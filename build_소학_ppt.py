# -*- coding: utf-8 -*-
"""
소학(小學) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '주자(朱子)가 엮은 인격의 교과서 · 조선 500년의 사람 만드는 양식',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.6),
                '小 學',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '소 학',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
                '習 與 性 成 — 습관이 본성을 이룬다',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '입교 · 명륜 · 경신 · 계고 · 가언 · 선행 — 내편 4 · 외편 2 · 총 6편',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 소학은 어떤 책인가'),
        ('Ⅱ.', '구조 — 내편 4 · 외편 2'),
        ('Ⅲ.', '교육 5대 원리'),
        ('Ⅳ.', '6편 깊이 읽기'),
        ('Ⅴ.', '핵심 개념 — 구용·구사·오륜'),
    ]
    items_right = [
        ('Ⅵ.', '명구와 일화'),
        ('Ⅶ.', '조선의 소학 — 500년의 운명'),
        ('Ⅷ.', '오늘 청소년에게'),
        ('Ⅸ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_what(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '소학(小學) — 어린 시절의 배움',
              '대학에 들어가기 전, 사람 됨됨이의 뼈대를 짓는 책')
    add_paragraphs(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.5), [
        ('· 「소(小)」 = 「격이 낮다」가 아니라 「어리고 근본적인 단계」',
         {'font_size': 18, 'space_before': 6}),
        ('· 「대학(大學)」이 「대인의 학(大人之學)」이라면, 소학은 「인간 되기의 바탕」',
         {'font_size': 18, 'space_before': 6}),
        ('· 옛 주(周)나라 교육 — 8세에 소학 입학 → 15세에 대학 입학',
         {'font_size': 18, 'space_before': 6}),
        ('· 이 오래된 교육 체계를 남송의 주희(朱熹)가 재구성한 책이 곧 『소학』',
         {'font_size': 18, 'space_before': 6}),
        ('· 새로 지은 책이 아니라, 옛 경전에서 가려 뽑아 엮은 「교육 안솔로지」',
         {'font_size': 18, 'space_before': 6}),
        ('· 조선 500년 — 서당의 필수 교재, 사림파의 정체성, 선비의 「운영체제」',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC1)
def i_author(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '편찬자와 성립 — 주희(朱熹)의 만년 작업',
              '1187년 · 주희 58세 · 제자 유청지와 함께')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 주희(朱熹, 1130~1200) — 남송의 대유, 성리학의 집대성자',
         {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● 유청지(劉清之, 자 자징) — 자료 수집·1차 정리를 주도',
         {'font_size': 18, 'space_before': 10}),
        ('● 1187년(남송 순희 14년) 완성 — 주희 만년의 교육적 결정판',
         {'font_size': 18, 'space_before': 10}),
        ('● 출전 — 『예기』·『의례』·『논어』·『맹자』·『효경』·『서경』·『시경』·『좌전』 등에서 가려 뽑음',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 주희의 문제의식 — 「대학(수기치인)의 학문은 깊지만, 일상의 기본 습관이 없으면 사상누각」',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 「소학은 배움의 뼈대를 만들고, 대학은 그 위에 근육을 붙인다」 — 주희의 선언',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC1)
def i_position(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '주희의 교육 체계 속 소학의 자리')
    rows = [
        ('8~14세',  '소학(小學)',       '일상 예절 · 기본 인성의 배움'),
        ('15세~',   '대학(大學)',       '격물·치지·성의·정심·수신·제가·치국·평천하'),
        ('이후',    '논어(論語)',       '배움의 근본 확립'),
        ('이후',    '맹자(孟子)',       '배움의 발휘'),
        ('이후',    '중용(中庸)',       '배움의 미묘에 이름'),
    ]
    for i, (age, name, desc) in enumerate(rows):
        y = Inches(2.5 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.8), Inches(0.7), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.8), Inches(0.7),
                    age, font_size=14, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(2.7), y, Inches(2.8), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(2.7), y, Inches(2.8), Inches(0.7),
                    name, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.8), y + Inches(0.05), Inches(7.2), Inches(0.65),
                    desc, font_size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.4),
                '「소학 없이 대학을 읽는 사람은 기초 체력 없이 역도에 뛰어드는 자와 같다」 — 주희',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC1)
def i_ages(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '나이별 교육 단계',
              '입교(立敎) 편이 제시하는 발달 단계별 학습')
    rows = [
        ('태교(胎敎)', '뱃속의 아기도 교육 대상 — 부모의 보고 듣는 것이 영향'),
        ('6세',  '숫자와 방위의 이름을 가르친다'),
        ('7세',  '남녀의 자리를 구별한다 (남녀칠세부동석의 원형)'),
        ('8세',  '집을 드나들거나 식사 때 반드시 어른 뒤에 따른다'),
        ('9세',  '날짜와 달력을 가르친다'),
        ('10세', '외부 스승에게 — 글·수학·쇄소응대진퇴·예악·사어의 기초'),
        ('13세', '음악을 배우고 시(詩)를 외우며, 예전(禮典)을 익힌다'),
        ('15세', '대학(大學)에 들어가 「대학의 도」를 본격 배운다'),
    ]
    for i, (age, desc) in enumerate(rows):
        y = Inches(2.2 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.2), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.2), Inches(0.5),
                    age, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.1), y, Inches(9.8), Inches(0.5),
                    desc, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅱ. 구조 ==============
SEC2 = 'Ⅱ. 구조'

@S(SEC2)
def ii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '소학의 구조 — 내편 4 · 외편 2 · 총 6편',
              '원리 → 실증 → 보충의 완성된 교육 설계')
    add_paragraphs(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.8), [
        ('● 내편(內篇) — 4편 — 「원리와 그 증명」', {'font_size': 19, 'bold': True, 'color': ACCENT}),
        ('     1. 입교(立敎) — 교육의 확립 (원리)', {'font_size': 16, 'space_before': 4}),
        ('     2. 명륜(明倫) — 인륜의 밝힘 (원리)', {'font_size': 16, 'space_before': 2}),
        ('     3. 경신(敬身) — 자기 몸의 공경 (원리)', {'font_size': 16, 'space_before': 2}),
        ('     4. 계고(稽古) — 옛 성현의 실천 사례 (실증)', {'font_size': 16, 'space_before': 2}),
        ('● 외편(外篇) — 2편 — 「보충과 확장」', {'font_size': 19, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     5. 가언(嘉言) — 한·당·송 명현의 가르침 (말씀)', {'font_size': 16, 'space_before': 4}),
        ('     6. 선행(善行) — 한·당·송 명현의 모범적 행실 (행동)', {'font_size': 16, 'space_before': 2}),
        ('● 내편은 「고대 경전」, 외편은 「근세 스승들」', {'font_size': 15, 'space_before': 14, 'color': SUB}),
    ])


@S(SEC2)
def ii_chart(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '6편 한 폭으로 보기')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.55), INK)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(1.7), Inches(0.55),
                '구분', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(2.4), Inches(2.3), Inches(2.2), Inches(0.55),
                '편명', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(4.6), Inches(2.3), Inches(2.0), Inches(0.55),
                '성격', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(6.6), Inches(2.3), Inches(6.1), Inches(0.55),
                '핵심 주제', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('내편', '1. 입교(立敎)', '원리', '교육의 확립 — 태교부터 15세 대학 입학까지'),
        ('내편', '2. 명륜(明倫)', '원리', '인륜(오륜)의 밝힘 — 부자·군신·부부·장유·붕우'),
        ('내편', '3. 경신(敬身)', '원리', '자기 몸의 공경 — 마음·몸가짐·옷·음식'),
        ('내편', '4. 계고(稽古)', '실증', '옛 성현이 실제로 어떻게 살았는가의 사례'),
        ('외편', '5. 가언(嘉言)', '보충', '한·당·송 명현의 아름다운 말씀'),
        ('외편', '6. 선행(善行)', '보충', '한·당·송 명현의 모범적 행실'),
    ]
    for i, (cat, name, kind, desc) in enumerate(rows):
        y = Inches(2.85 + i * 0.65)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.65), bg)
        add_textbox(slide, Inches(0.7), y, Inches(1.7), Inches(0.65),
                    cat, font_size=14, color=SUB, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.4), y, Inches(2.2), Inches(0.65),
                    name, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.6), y, Inches(2.0), Inches(0.65),
                    kind, font_size=14, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.7), y, Inches(6.0), Inches(0.65),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅲ. 교육 5대 원리 ==============
SEC3 = 'Ⅲ. 교육 5대 원리'

@S(SEC3)
def iii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '소학을 꿰뚫는 다섯 가지 원리',
              '6편을 관통하는 교육 철학')
    boxes = [
        ('1', '習與性成', '습여성성', '습관이 본성을 이룬다'),
        ('2', '灑掃應對進退', '쇄소응대진퇴', '일상이 곧 교육'),
        ('3', '先小學 後大學', '선소학 후대학', '순서가 있다'),
        ('4', '立 志', '입지', '뜻을 세움'),
        ('5', '以身敎', '이신교', '몸으로 가르친다'),
    ]
    for i, (num, han, kor, desc) in enumerate(boxes):
        y = Inches(2.4 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(0.85), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(0.85), Inches(0.7),
                    num, font_size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.7), y, Inches(3.8), Inches(0.7), PALE)
        add_textbox(slide, Inches(1.7), y, Inches(3.8), Inches(0.7),
                    han, font_size=22, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.7), y + Inches(0.05), Inches(2.4), Inches(0.65),
                    kor, font_size=15, color=SUB, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(8.2), y + Inches(0.05), Inches(4.8), Inches(0.65),
                    desc, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC3)
def iii_habit(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '원리 1 — 습여성성(習與性成)',
              '습관이 본성을 이룬다 — 사람은 만들어지는 것이다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 교육의 핵심은 「지식 전달」이 아니라 「습관 형성」',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 사람은 타고나는 것이 아니라 어린 시절의 반복으로 만들어진다',
         {'font_size': 18, 'space_before': 10}),
        ('● 좋은 습관이 몸에 배면 인격이 되고, 나쁜 습관이 몸에 배면 그 또한 인격이 된다',
         {'font_size': 18, 'space_before': 10}),
        ('● 출전 — 『서경·태갑』 「茲乃不義, 習與性成」', {'font_size': 17, 'color': SUB, 'space_before': 12, 'font_name': 'Batang'}),
        ('● 오늘날의 behavioral psychology · 「Atomic Habits」의 2,000년 선행 이론',
         {'font_size': 17, 'space_before': 12, 'color': SUB}),
        ('● 「어릴 때 들인 습관은 끝까지 간다」 — 소학 교육의 정언명령',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC3)
def iii_sweep(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '원리 2 — 쇄소응대진퇴(灑掃應對進退)',
              '일상이 곧 교육 — 거창함이 아니라 작은 행위에서 출발')
    cols = [
        ('灑 掃', '쇄소',
         '물 뿌리고 빗자루로 쓸기.\n청결과 책임의 습관.\n\n방을 치우는 아이가\n인생도 치울 수 있다.'),
        ('應 對', '응대',
         '부르면 대답하고\n물으면 답하기.\n\n말의 태도가 곧 인격.\n응답의 결이 사람을 만든다.'),
        ('進 退', '진퇴',
         '나아가고 물러나는\n몸가짐.\n\n어른 앞에서, 모임에서,\n언제 어디서 멈출 줄 아는 감각.'),
    ]
    for i, (han, kor, body) in enumerate(cols):
        x = Inches(0.7 + i * 4.2)
        add_filled_rect(slide, x, Inches(2.3), Inches(3.9), Inches(1.0), INK)
        add_textbox(slide, x, Inches(2.3), Inches(3.9), Inches(0.6),
                    han, font_size=28, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.85), Inches(3.9), Inches(0.4),
                    kor, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.3), Inches(3.6), Inches(3.5), Inches(3.5),
                       [(body, {'font_size': 14, 'color': INK})], line_spacing=1.4)


@S(SEC3)
def iii_order(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '원리 3 — 선소학 후대학(先小學 後大學)',
              '순서가 있다 — 기초 없는 심화는 공허')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 작은 배움을 먼저, 큰 배움을 나중',
         {'font_size': 18, 'space_before': 6}),
        ('● 수기(修己) 없는 치인(治人)은 재앙 — 인격 없는 권력은 흉기',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 「집 없이 지붕을 올리는 일과 같다」 — 퇴계 이황의 비유',
         {'font_size': 17, 'space_before': 12, 'color': SUB}),
        ('● 「소학 없이 사서삼경에 뛰어드는 건 빠른 게 아니라 빗나간 것」',
         {'font_size': 17, 'space_before': 12}),
        ('● 김굉필이 30세까지 소학만 파고든 이유 — 「아직 소학도 다 체득하지 못했다」',
         {'font_size': 17, 'space_before': 12, 'color': SUB}),
    ])


@S(SEC3)
def iii_will(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '원리 4 — 입지(立志)',
              '뜻을 세움 — 「성인으로 자기를 기약하라」')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(1.0), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(1.0),
                '立 志 以 聖 人 自 期',
                font_size=32, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(2.0), Inches(3.3), Inches(9.3), Inches(0.4),
                '뜻을 세우되 성인(聖人)으로 자기를 기약하라',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 소학은 「순한 애를 만드는 책」이 아니다 — 「큰 사람을 만드는 책」',
         {'font_size': 18, 'space_before': 8, 'bold': True, 'color': ACCENT}),
        ('● 뜻이 낮으면 평생 낮다 — 꿈의 크기가 삶의 크기',
         {'font_size': 18, 'space_before': 10}),
        ('● 율곡 『격몽요결』 — 16개 조의 제1조가 「입지(立志)」',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 「착한 아이가 되라」가 아닌 「역사에 이름을 남기는 사람이 되라」',
         {'font_size': 17, 'space_before': 10}),
    ])


@S(SEC3)
def iii_body(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '원리 5 — 이신교(以身敎)',
              '몸으로 가르친다 — 말이 아닌 삶으로')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(1.0), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(1.0),
                '以 身 敎, 不 以 言 敎',
                font_size=30, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(2.0), Inches(3.3), Inches(9.3), Inches(0.4),
                '몸으로 가르치되, 말로 가르치지 말라',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 부모·스승·어른이 먼저 실천해야 아이가 배운다',
         {'font_size': 18, 'space_before': 8}),
        ('● 「계고」·「선행」 편의 모든 사례 — 「이 사람은 이렇게 살았다」는 모범의 실증',
         {'font_size': 18, 'space_before': 10}),
        ('● 증자(曾子) — 아내가 농담으로 한 약속(돼지를 잡아주마)을 실제로 지킴',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 가장 강력한 교육은 옆에서 보고 자라는 것이다',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


# ============== Ⅳ. 6편 깊이 읽기 ==============
SEC4 = 'Ⅳ. 6편 깊이 읽기'

@S(SEC4)
def iv_lijiao(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC4} (1/6)', n, t)
    add_title(slide, '① 입교(立敎) — 교육의 확립',
              '내편 첫 편 — 「무엇을 어떻게 가르치는가」의 원리')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 핵심 명제 — 「習與性成」 — 습관이 본성을 이룬다',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 발달 단계별 교육 — 태교 → 6세 → 7세 → 8세 → 9세 → 10세 → 13세 → 15세',
         {'font_size': 17, 'space_before': 12}),
        ('● 「쇄소응대진퇴」 — 소학의 근본 행위 — 청소·인사·자리 지키기',
         {'font_size': 17, 'space_before': 10}),
        ('● 「예악사어서수(禮樂射御書數)」 — 육예(六藝)의 학습',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
        ('● 교육은 거창한 사상 전수가 아니라 「일상의 반복 훈련」',
         {'font_size': 17, 'space_before': 12, 'color': SUB}),
        ('● 「청소를 시키고, 인사를 시키고, 자리를 구분하게 하는 것이 교육의 뼈대」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC4)
def iv_minglun(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC4} (2/6)', n, t)
    add_title(slide, '② 명륜(明倫) — 인륜을 밝힘',
              '내편 둘째 편 — 오륜(五倫), 인간관계의 다섯 기둥')
    rows = [
        ('父子有親', '부자유친', '부모-자식', '친·효 — 문안, 봉양, 제사'),
        ('君臣有義', '군신유의', '군주-신하', '의·충 — 바른 일로 섬김, 간언'),
        ('夫婦有別', '부부유별', '남편-아내', '별·경 — 역할 구분, 상호 공경'),
        ('長幼有序', '장유유서', '어른-아이', '서·제 — 형우제공, 자리의 질서'),
        ('朋友有信', '붕우유신', '벗-벗',     '신·우 — 말의 신뢰, 선으로 인도'),
    ]
    for i, (han, kor, rel, desc) in enumerate(rows):
        y = Inches(2.4 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.7), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(2.7), Inches(0.7),
                    han, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, Inches(3.5), y, Inches(1.8), Inches(0.7), PALE)
        add_textbox(slide, Inches(3.5), y, Inches(1.8), Inches(0.7),
                    kor, font_size=14, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.5), y + Inches(0.05), Inches(2.0), Inches(0.65),
                    rel, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.6), y + Inches(0.05), Inches(5.3), Inches(0.65),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_jingshen(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC4} (3/6)', n, t)
    add_title(slide, '③ 경신(敬身) — 자기 몸을 공경하라',
              '내편 셋째 편 — 「내 몸은 부모의 끼친 몸」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 핵심 명제 — 「君子之自身, 父母之遺體」 — 내 몸은 부모의 끼친 몸',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 네 갈래로 자기를 다스린다', {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 心術(심술) — 마음가짐 — 구용(九容)·구사(九思)',
         {'font_size': 16, 'space_before': 6}),
        ('     · 威儀(위의) — 앉고 서고 걷는 몸가짐',
         {'font_size': 16, 'space_before': 4}),
        ('     · 衣服(의복) — 때·자리에 맞는 옷차림, 사치·남루 모두 경계',
         {'font_size': 16, 'space_before': 4}),
        ('     · 飮食(음식) — 식사의 예절, 식탐 경계, 절도',
         {'font_size': 16, 'space_before': 4}),
        ('● 작은 태도가 쌓여 한 사람의 품격이 된다 — 「body language의 동양 원형」',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
    ])


@S(SEC4)
def iv_jigu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC4} (4/6)', n, t)
    add_title(slide, '④ 계고(稽古) — 옛것을 상고하다',
              '내편 넷째 편 — 「원리의 실증 — 옛사람은 이렇게 했다」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 입교·명륜·경신의 원리를 「상고(上古)의 인물 사례」로 증명',
         {'font_size': 18, 'space_before': 6}),
        ('● 대표 사례', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('     · 순(舜) — 완악한 아버지·계모·동생 사이에서 효를 다함',
         {'font_size': 16, 'space_before': 6}),
        ('     · 맹모삼천(孟母三遷)·맹모단기(孟母斷機) — 교육 환경의 결정성',
         {'font_size': 16, 'space_before': 4}),
        ('     · 어린 공자 — 제기를 모아 제사 놀이를 한 일',
         {'font_size': 16, 'space_before': 4}),
        ('     · 증자(曾子) — 아내의 농담을 진실로 지켜 돼지를 잡음',
         {'font_size': 16, 'space_before': 4}),
        ('     · 동중서(董仲舒) — 3년간 장막을 내리고 책만 읽음',
         {'font_size': 16, 'space_before': 4}),
        ('● 「스토리텔링 교육」의 원형 — 추상이 아닌 구체로 가르친다',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC4)
def iv_jiayan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC4} (5/6)', n, t)
    add_title(slide, '⑤ 가언(嘉言) — 아름다운 말씀',
              '외편 첫 편 — 한·당·송 명현의 어록')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 송대 성리학자들의 명언을 중심으로 선별',
         {'font_size': 17, 'space_before': 6}),
        ('● 인용된 주요 인물', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('     · 주돈이(周敦頤, 염계) — 『태극도설』·『통서』',
         {'font_size': 15, 'space_before': 4}),
        ('     · 장재(張載, 횡거) — 「위천지입심, 위생민입명, 위왕성계절학, 위만세개태평」',
         {'font_size': 15, 'space_before': 4}),
        ('     · 정호·정이(明道·伊川) 형제 — 성리학의 뼈대',
         {'font_size': 15, 'space_before': 4}),
        ('     · 사마광(司馬光) — 『가범(家範)』의 가훈',
         {'font_size': 15, 'space_before': 4}),
        ('     · 주희 자신 — 자기 편지·어록 일부',
         {'font_size': 15, 'space_before': 4}),
        ('● 핵심 주제 — 자기 수양·독서·치가·관리의 자세·말의 신중함',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 대표 명구 — 「立志以聖人自期」 「讀書百遍義自見」',
         {'font_size': 15, 'space_before': 8, 'font_name': 'Batang'}),
    ])


@S(SEC4)
def iv_shanxing(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC4} (6/6)', n, t)
    add_title(slide, '⑥ 선행(善行) — 모범적 행실',
              '외편 둘째 편 — 가언이 「말」이라면 선행은 「행동」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 송대 유학자들의 일상 — 부모·형제·스승·벗·자식·이웃에게 어떻게 살았는가',
         {'font_size': 17, 'space_before': 6}),
        ('● 대표 일화', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('     · 범중엄(范仲淹) — 「先天下之憂而憂, 後天下之樂而樂」 · 의전(義田) 설립',
         {'font_size': 15, 'space_before': 4, 'font_name': 'Batang'}),
        ('     · 사마광(司馬光) — 어린 시절 장독대에 빠진 친구를 돌로 깨뜨려 구함',
         {'font_size': 15, 'space_before': 4}),
        ('     · 정이천(程伊川) — 눈 속에 스승 집 앞에 선 채 기다림 (「程門立雪」)',
         {'font_size': 15, 'space_before': 4, 'font_name': 'Batang'}),
        ('     · 황향(黃香) — 여름엔 부친 잠자리를 부채질, 겨울엔 이불을 데움 (「黃香溫席」)',
         {'font_size': 15, 'space_before': 4}),
        ('● 「행실의 데이터베이스」 — 「이럴 때 이런 어른은 이렇게 했다」',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
        ('● 오늘날의 case-based learning, 자기계발서 사례 학습의 완벽한 원형',
         {'font_size': 15, 'space_before': 8, 'color': SUB}),
    ])


# ============== Ⅴ. 핵심 개념 ==============
SEC5 = 'Ⅴ. 핵심 개념'

@S(SEC5)
def v_jiurong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '구용(九容) — 몸가짐의 아홉 가지',
              '경신(敬身) 편의 핵심 — 일거수일투족의 규범')
    rows = [
        ('足容重', '족용중', '발은 무겁게 — 경솔히 움직이지 않는다'),
        ('手容恭', '수용공', '손은 공손하게 — 함부로 휘젓지 않는다'),
        ('目容端', '목용단', '눈은 단정하게 — 흘낏거리지 않는다'),
        ('口容止', '구용지', '입은 다물고 — 까닭 없이 떠벌리지 않는다'),
        ('聲容靜', '성용정', '소리는 고요하게 — 큰 소리로 외치지 않는다'),
        ('頭容直', '두용직', '머리는 곧게 — 기울이지 않는다'),
        ('氣容肅', '기용숙', '기는 엄숙하게 — 거친 숨을 내쉬지 않는다'),
        ('立容德', '입용덕', '섰을 때는 덕스럽게 — 비스듬히 기대지 않는다'),
        ('色容莊', '색용장', '얼굴빛은 장엄하게 — 가벼이 변하지 않는다'),
    ]
    for i, (han, kor, desc) in enumerate(rows):
        col, row = i % 3, i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.3 + row * 1.5)
        add_filled_rect(slide, x, y, Inches(3.9), Inches(1.3), PALE)
        add_textbox(slide, x, y + Inches(0.1), Inches(3.9), Inches(0.45),
                    han, font_size=18, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(0.6), Inches(3.9), Inches(0.3),
                    kor, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.6), Inches(0.4),
                    desc, font_size=11, color=INK, align=PP_ALIGN.CENTER)


@S(SEC5)
def v_jiusi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '구사(九思) — 생각의 아홉 가지',
              '경신(敬身) 편의 또 다른 핵심 — 매 순간의 마음 점검')
    rows = [
        ('視思明', '시사명', '볼 때는 — 밝음을 생각하라'),
        ('聽思聰', '청사총', '들을 때는 — 총명함을 생각하라'),
        ('色思溫', '색사온', '얼굴빛은 — 따스함을 생각하라'),
        ('貌思恭', '모사공', '모습은 — 공손함을 생각하라'),
        ('言思忠', '언사충', '말은 — 충실함을 생각하라'),
        ('事思敬', '사사경', '일은 — 공경함을 생각하라'),
        ('疑思問', '의사문', '의심 날 때는 — 물음을 생각하라'),
        ('忿思難', '분사난', '화날 때는 — 어려움을 생각하라'),
        ('見得思義', '견득사의', '얻을 때는 — 의로움을 생각하라'),
    ]
    for i, (han, kor, desc) in enumerate(rows):
        col, row = i % 3, i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.3 + row * 1.5)
        add_filled_rect(slide, x, y, Inches(3.9), Inches(1.3), PALE)
        add_textbox(slide, x, y + Inches(0.1), Inches(3.9), Inches(0.45),
                    han, font_size=18, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(0.6), Inches(3.9), Inches(0.3),
                    kor, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.6), Inches(0.4),
                    desc, font_size=11, color=INK, align=PP_ALIGN.CENTER)


@S(SEC5)
def v_filial(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '효(孝) — 명륜의 절반을 차지하는 주제',
              '소학이 효를 가르치는 결')
    add_filled_rect(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.1), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.1),
                '身 體 髮 膚, 受 之 父 母\n不 敢 毁 傷, 孝 之 始 也',
                font_size=23, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.4),
                '몸·머리카락·살갗은 부모로부터 받은 것 / 감히 헐고 다치게 하지 않음이 효의 시작이다 — 『효경』',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.2), Inches(12.0), Inches(3.0), [
        ('● 「혼정신성(昏定晨省)」 — 저녁엔 잠자리를 봐드리고, 새벽엔 문안 인사',
         {'font_size': 17, 'space_before': 8, 'font_name': 'Batang'}),
        ('● 「부모의 뜻을 살핀다」 — 말씀하시기 전에 알아채는 것이 효의 본질',
         {'font_size': 17, 'space_before': 10}),
        ('● 효는 부모만이 아닌 모든 어른에 대한 「공경의 결」로 확장',
         {'font_size': 17, 'space_before': 10}),
        ('● 「兄友弟恭(형우제공)」 — 형은 사랑하고 아우는 공경 — 효의 자매편',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
    ])


@S(SEC5)
def v_friends(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '벗을 가리는 법 — 익자삼우 · 손자삼우',
              '명륜·붕우유신의 핵심 — 어떤 친구를 사귈 것인가')
    cols = [
        ('益者三友', '이로운 세 벗',
         '直 — 정직한 벗\n\n諒 — 신실한 벗\n\n多 聞 — 견문이 넓은 벗', ACCENT),
        ('損者三友', '해로운 세 벗',
         '便辟 — 비위 맞추는 벗\n\n善柔 — 줏대 없이 부드러운 벗\n\n便佞 — 말만 잘하는 벗', INK),
    ]
    for i, (han, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=26, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.85), Inches(5.9), Inches(0.4),
                    label, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.3), Inches(3.6), Inches(5.5), Inches(3.5),
                       [(body, {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER})],
                       line_spacing=1.5)
    add_textbox(slide, Inches(0.7), Inches(6.9), Inches(12.0), Inches(0.4),
                '— 『논어·계씨』 인용 (명륜 편)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅵ. 명구와 일화 ==============
SEC6 = 'Ⅵ. 명구와 일화'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC6)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC6} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=30, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=17, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 17, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('서경·태갑 / 소학 입교', '習 與 性 成',
     '습관이 본성을 이룬다',
     '소학 전체를 꿰뚫는 한 마디. 사람은 타고나는 것이 아니라 어린 시절의 반복으로 만들어진다.'),
    ('소학 입교', '灑 掃 應 對 進 退',
     '물 뿌리고 쓸기 · 응답하기 · 나아가고 물러나기',
     '소학의 근본 행위. 큰 도는 작은 행동에서 시작된다. 청소·인사·자리 지키기가 교육의 뼈대.'),
    ('효경 / 소학 명륜', '身 體 髮 膚, 受 之 父 母  不 敢 毁 傷, 孝 之 始 也',
     '몸과 머리카락은 부모로부터 받은 것이니, 헐고 다치지 않음이 효의 시작',
     '내 몸은 내 것이 아니라 부모의 끼친 몸이다 — 자기 관리가 효의 출발이라는 동양적 관점.'),
    ('소학 경신 — 구용', '足 容 重, 手 容 恭, 目 容 端, 口 容 止',
     '발은 무겁게, 손은 공손하게, 눈은 단정하게, 입은 다물고',
     '경신 편의 핵심. 율곡 『격몽요결』 「지신」장의 근간 — 조선 선비의 일거수일투족을 규정한 원전.'),
    ('소학 경신 — 구사', '見 得 思 義',
     '얻을 때는 의로움을 생각하라',
     '구사(九思)의 마지막 항목. 이익 앞에서 의를 잊지 않는 것 — 동양 윤리의 가장 어려운 한 줄.'),
    ('논어 계씨 / 소학 명륜', '益 者 三 友 — 直 · 諒 · 多 聞',
     '이로운 세 벗 — 정직·신실·견문 넓은 벗',
     '벗 사귐의 기준. 누구와 어울리느냐가 사람을 결정한다. 명륜의 「붕우유신」 핵심.'),
    ('장재(횡거) / 소학 가언', '爲 天 地 立 心  爲 生 民 立 命  爲 往 聖 繼 絕 學  爲 萬 世 開 太 平',
     '천지를 위해 마음을 세우고, 백성을 위해 명을 세우며, 옛 성인을 위해 끊긴 학문을 잇고, 만세를 위해 태평을 연다',
     '횡거 사구(四句) — 조선 선비의 정신적 좌우명. 입지(立志)의 최고 표현.'),
    ('범중엄 / 소학 선행', '先 天 下 之 憂 而 憂  後 天 下 之 樂 而 樂',
     '천하의 근심을 먼저 근심하고, 천하의 즐거움을 뒤에 즐긴다',
     '북송 명재상 범중엄의 좌우명. 「공인(公人)」의 자세 — 자기 즐거움을 뒤로 미루는 사람만이 공직에 적합하다.'),
    ('소학 가언', '立 志 以 聖 人 自 期',
     '뜻을 세우되 성인(聖人)으로 자기를 기약하라',
     '소학의 입지(立志) 가르침. 율곡 『격몽요결』 제1조 「입지」가 이 한 줄을 풀이.'),
    ('주희 / 소학 가언', '讀 書 百 遍, 義 自 見',
     '책을 백 번 읽으면 뜻이 저절로 드러난다',
     '주희가 강조한 독서법. 진리는 반복 속에서 스스로 모습을 드러낸다 — 동양 독서론의 핵심.'),
    ('정문입설(程門立雪) — 소학 선행', '雪 中 立 候, 程 門 三 尺',
     '눈 속에 서서 기다리니, 정문(程門)에 눈이 석 자',
     '정이천의 제자 양시(楊時)와 유작(游酢)이 잠든 스승을 깨우지 않으려 눈 속에 선 채 기다림 — 스승을 모시는 정성의 표상.'),
    ('황향온석(黃香溫席) — 소학 선행', '夏 扇 父 枕  冬 溫 父 衾',
     '여름엔 아버지 베개에 부채질을, 겨울엔 아버지 이불을 데움',
     '9세의 황향(黃香)이 아버지를 위해 한 일. 효의 가장 따스한 동양적 이미지로 24효(孝)의 하나.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅶ. 조선의 소학 ==============
SEC7 = 'Ⅶ. 조선의 소학'

@S(SEC7)
def vii_intro(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '조선 500년 — 한 권의 책이 한 학파를 만들다',
              '서당의 필수 교재 · 사림파의 정체성 · 선비의 운영체제')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 고려 말 — 안향·정몽주·이색을 통해 주자학과 함께 도입',
         {'font_size': 18, 'space_before': 4}),
        ('● 세종 — 왕실 교육의 기본 교재. 경연에서 신하들과 소학을 강론',
         {'font_size': 18, 'space_before': 10}),
        ('● 세조·성종 — 『소학언해(小學諺解)』 간행, 한문 모르는 이도 읽을 수 있게',
         {'font_size': 18, 'space_before': 10}),
        ('● 중종 — 조광조 주도의 「소학 부흥 운동」, 전국 유생 필수화 시도',
         {'font_size': 18, 'space_before': 10}),
        ('● 연산군(1505) — 갑자사화 후 소학을 금서로 지정',
         {'font_size': 18, 'space_before': 10, 'color': ACCENT, 'bold': True}),
        ('● 중종반정(1506) 후 복권 — 「금지되었다는 사실 자체가 소학의 권위를 더 높였다」',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC7)
def vii_curriculum(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '조선 서당의 학습 단계',
              '소학은 「유학의 세계로 들어가는 문」')
    rows = [
        ('5~7세',  '千 字 文', '천자문', '한자 1,000자와 기본 우주론'),
        ('7~9세',  '童 蒙 先 習', '동몽선습', '오륜의 간략한 해설 (박세무, 조선 자작)'),
        ('9~15세', '小 學', '소학', '일상·오륜·역사 사례의 종합', True),
        ('15세~',  '四 書', '사서', '대학 → 논어 → 맹자 → 중용'),
        ('이후',   '三 經', '삼경', '예기·춘추 → 과거 준비'),
    ]
    for i, row in enumerate(rows):
        age, han, kor, desc = row[0], row[1], row[2], row[3]
        highlighted = len(row) > 4 and row[4]
        y = Inches(2.4 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.5), Inches(0.7), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.7),
                    age, font_size=14, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        color = ACCENT if highlighted else INK
        add_filled_rect(slide, Inches(2.4), y, Inches(2.6), Inches(0.7), color)
        add_textbox(slide, Inches(2.4), y, Inches(2.6), Inches(0.7),
                    han, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.2), y + Inches(0.05), Inches(1.6), Inches(0.65),
                    kor, font_size=14, color=SUB, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.9), y + Inches(0.05), Inches(6.0), Inches(0.65),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC7)
def vii_kim(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '김굉필(金宏弼) — 「소학동자(小學童子)」',
              '소학만 읽어도 군자가 될 수 있음을 증명한 사람')
    add_filled_rect(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.1), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.1),
                '讀 小 學 書, 方 覺 從 前 非 人',
                font_size=26, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.4),
                '소학을 읽고 나서야 비로소 어제까지의 내가 사람이 아니었음을 깨달았다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.2), Inches(12.0), Inches(3.0), [
        ('● 한훤당 김굉필(1454~1504) — 20대 후반까지 소학만 깊이 파고듦',
         {'font_size': 17, 'space_before': 6}),
        ('● 자기 호를 「소학동자(小學童子) — 소학을 배우는 어린아이」라 부르게 함',
         {'font_size': 17, 'space_before': 8}),
        ('● 30세가 되어서야 비로소 육경(六經)으로 나아감',
         {'font_size': 17, 'space_before': 8}),
        ('● 사림파의 시조 — 그의 제자가 정여창·조광조, 그 학맥이 이언적·이황·이이로',
         {'font_size': 17, 'space_before': 8}),
        ('● 갑자사화(1504)에 처형되었으나 그가 뿌린 씨앗이 조선 성리학의 정통이 됨',
         {'font_size': 17, 'space_before': 8, 'color': SUB}),
    ])


@S(SEC7)
def vii_cho(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '조광조(趙光祖) — 소학을 정치 개혁의 무기로',
              '정암(靜菴) · 38세에 죽임당한 천재 개혁가')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 조광조(1482~1519) — 김굉필의 제자, 중종의 신임을 얻어 사림 정치 주도',
         {'font_size': 18, 'space_before': 4}),
        ('● 「지치주의(至治主義)」 개혁의 중심에 「소학 보급」을 놓음',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 전국 향교·서원의 소학 교육 의무화를 추진',
         {'font_size': 18, 'space_before': 10}),
        ('● 「온 나라를 소학의 원리로 교화한다」는 비전',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 기묘사화(1519)로 38세에 죽임당함 — 「소학이 정치의 뿌리가 될 수 있다」는 그의 시도는 꺾임',
         {'font_size': 17, 'space_before': 12}),
        ('● 그러나 그의 정신은 이후 300년간 조선 선비들의 표준이 됨',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC7)
def vii_four(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '조선 4대 유학자의 소학 평가')
    blocks = [
        ('퇴계 이황', '退溪 李滉 (1501~1570)',
         '「소학이 대학의 밑받침」 — 『성학십도』 제3도 「소학도(小學圖)」를 따로 그려 선조에게 올림.\n왕에게도 소학을 권한 담대한 진언.'),
        ('율곡 이이', '栗谷 李珥 (1536~1584)',
         '『격몽요결』 — 사실상 「조선판 소학 재편집」.\n『학교모범』 — 16개 조항의 제1조 입지(立志), 제2조 검신(檢身).\n「小學者, 立德之基, 敎化之本」.'),
        ('우암 송시열', '尤庵 宋時烈 (1607~1689)',
         '「만약 한 권의 책을 평생 읽으라 한다면 나는 소학을 택하겠다.」\n노년에도 소학을 손에서 놓지 않았다.'),
        ('다산 정약용', '茶山 丁若鏞 (1762~1836)',
         '두 아들에게 보낸 유배지의 편지 — 「소학의 구용·구사를 매일 실천하라」.\n실학자가 마지막에 인정한 소학의 보편성.\n『소학주관(小學珠串)』 저술.'),
    ]
    for i, (name, en, body) in enumerate(blocks):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.4)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(2.2), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.45),
                    name, font_size=17, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.35),
                    en, font_size=11, color=SUB, font_name='Batang')
        add_paragraphs(slide, x + Inches(0.2), y + Inches(0.95), Inches(5.5), Inches(1.2),
                       [(body, {'font_size': 12, 'color': INK})], line_spacing=1.35)


@S(SEC7)
def vii_routine(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '소학이 만든 조선 선비의 하루',
              '5~6년간 소학을 읽으면 몸에 배는 10가지 습관')
    items = [
        ('1', '새벽 기상 — 부모님께 문안 (혼정신성)'),
        ('2', '의관 정제 — 머리부터 발끝까지 예법대로'),
        ('3', '자기 방 청소 — 「쇄소」를 직접 실행'),
        ('4', '식사 예절 — 어른보다 먼저 수저를 들지 않음'),
        ('5', '독서 자세 — 단정히 앉아 책을 함부로 다루지 않음'),
        ('6', '언어 예절 — 상대에 따라 호칭·어미를 가림'),
        ('7', '외출 전후 — 부모에게 행선지·도착을 알림'),
        ('8', '벗 가리기 — 직·량·다문의 벗을 가까이'),
        ('9', '자기 성찰 — 하루에 세 번 자기를 돌아봄 (삼성)'),
        ('10', '뜻 세우기 — 「나는 성인을 기약한다」'),
    ]
    for i, (num, txt) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 0.9)
        add_filled_rect(slide, x, y, Inches(0.7), Inches(0.7), ACCENT)
        add_textbox(slide, x, y, Inches(0.7), Inches(0.7),
                    num, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + Inches(0.9), y + Inches(0.05), Inches(5.0), Inches(0.65),
                    txt, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅷ. 오늘 청소년에게 ==============
SEC8 = 'Ⅷ. 오늘 청소년에게'

@S(SEC8)
def viii_lack(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '현대 청소년의 결핍 지도',
              '왜 지금 다시 소학을 펼치는가')
    items = [
        ('습관의 결핍',     '스마트폰이 일상의 호흡을 빼앗는다'),
        ('관계의 서투름',   '디지털 소통의 증가, 오프라인 관계 예의의 약화'),
        ('자기 관리의 공백', '수면·식사·운동·정리정돈의 기본 루틴 붕괴'),
        ('목표의 부재',     '「무엇이 되고 싶은가」는 묻지만 「어떤 사람이 되고 싶은가」는 묻지 않음'),
        ('모범의 결핍',     '롤모델이 연예인·인플루언서로 집중, 역사 인물을 모름'),
    ]
    for i, (cat, desc) in enumerate(items):
        y = Inches(2.5 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.5), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(3.5), Inches(0.7),
                    cat, font_size=17, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.4), y + Inches(0.05), Inches(8.5), Inches(0.65),
                    desc, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC8)
def viii_fill(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '소학이 채워줄 수 있는 5가지')
    items = [
        ('1', '습관의 설계도',     '쇄소응대진퇴 = 「방 치우기·인사·자리 지키기」 — Atomic Habits의 동양 원조'),
        ('2', '관계의 문법',        '오륜 = 봉건이 아닌 「관계마다 달라지는 적절한 태도」 — 위세대·동료·친구를 가리는 감각'),
        ('3', '자기 관리 체크리스트', '구용·구사 = 오늘날의 morning routine · personal standards의 원형'),
        ('4', '역사 속 롤모델',      '계고·선행 = 「행동 모델의 데이터베이스」 — 맹모삼천·정문입설·황향온석'),
        ('5', '큰 뜻의 지평',         '입지(立志) = 「착한 아이가 되라」가 아닌 「역사에 이름을 남기는 사람이 되라」'),
    ]
    for i, (num, cat, desc) in enumerate(items):
        y = Inches(2.4 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(0.8), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(0.8), Inches(0.7),
                    num, font_size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.6), y, Inches(3.0), Inches(0.7), PALE)
        add_textbox(slide, Inches(1.6), y, Inches(3.0), Inches(0.7),
                    cat, font_size=15, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.8), y + Inches(0.05), Inches(8.2), Inches(0.65),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC8)
def viii_reinterpret(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '현대적 재해석 — 무엇은 바꿔 읽는가',
              '세부 규정이 아니라 근본 정신을 취한다')
    rows = [
        ('부부유별의 위계적 해석', '→', '상호 존중과 역할 분담'),
        ('남녀칠세부동석의 엄격 적용', '→', '적절한 경계와 예의'),
        ('효(孝)의 절대적 복종',   '→', '부모-자녀의 상호 돌봄'),
        ('계급·신분 전제',         '→', '모든 관계에 보편 적용되는 기본 예의'),
    ]
    for i, (old, arrow, new) in enumerate(rows):
        y = Inches(2.5 + i * 0.9)
        add_filled_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.7), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(5.5), Inches(0.7),
                    old, font_size=15, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.3), y, Inches(0.6), Inches(0.7),
                    arrow, font_size=20, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(7.0), y, Inches(5.9), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(7.0), y, Inches(5.9), Inches(0.7),
                    new, font_size=15, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.4),
                '「세부 규정」이 아니라 「근본 정신」 — 12세기의 산물을 21세기로 옮긴다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC8)
def viii_seven(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '오늘 우리가 소학에서 배울 7가지')
    items = [
        '1. 작은 것이 근본이다 — 큰 일은 작은 습관에서 태어난다',
        '2. 일상이 인격이다 — 특별한 순간이 아니라 매일의 반복이 당신을 만든다',
        '3. 관계에는 문법이 있다 — 모든 사람을 똑같이 대하는 게 평등이 아니다',
        '4. 몸이 배우면 마음이 따라온다 — 태도·자세를 바로 하면 마음도 정돈된다',
        '5. 옛사람의 이야기는 내 이야기의 거울이다 — 역사 속 선택에서 오늘을 배운다',
        '6. 뜻이 낮으면 일생이 낮다 — 꿈의 크기가 삶의 크기, 성인을 기약하라',
        '7. 순서를 지켜라 — 기본 없이 고급은 없다. 소학 없이 대학 없다',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.3 + i * 0.65)
        add_textbox(slide, Inches(0.7), y, Inches(12.0), Inches(0.55),
                    txt, font_size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅸ. 마무리 ==============
SEC9 = 'Ⅸ. 마무리'

@S(SEC9)
def ix_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '소학, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 소학은 주희가 1187년 엮은 「어린 시절의 배움 — 인간 되기의 바탕」.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 내편 4(입교·명륜·경신·계고) + 외편 2(가언·선행) — 총 6편.',
         {'font_size': 18, 'space_before': 8}),
        ('● 다섯 원리 — 습여성성 · 쇄소응대진퇴 · 선소학후대학 · 입지 · 이신교.',
         {'font_size': 18, 'space_before': 8}),
        ('● 구용·구사·오륜 — 일거수일투족의 동양적 매뉴얼.',
         {'font_size': 18, 'space_before': 8}),
        ('● 조선 500년 — 김굉필·조광조·이황·이이·송시열·정약용이 모두 그 길을 걸음.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「큰 사람은 작은 행위에서 만들어진다」는 책.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC9)
def ix_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.0),
                '習 與 性 成',
                font_size=100, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(3.6), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.6),
                '습 관 이  본 성 을  이 룬 다',
                font_size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
                '— 소학 입교(立敎) 편의 정언명령',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                '立 志 以 聖 人 自 期',
                font_size=24, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '뜻을 세우되 성인(聖人)으로 자기를 기약하라',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\소학_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
