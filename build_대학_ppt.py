# -*- coding: utf-8 -*-
"""
대학 발표자료 — 망라적 80장 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 1. 표지 ==============
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8),
                '大 學', font_size=130, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                'The Great Learning · 대학',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.75), Inches(12.3), Inches(0.5),
                '"대인지학(大人之學)" — 사서(四書)의 입덕지문(入德之門)',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
                '본래 『예기(禮記)』의 한 편 · 약 1,700자 · 주희(朱熹) 사서 독립 · 800년의 표준',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '"自天子以至於庶人, 壹是皆以修身爲本"',
                font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '— 천자로부터 서민에 이르기까지, 모두 한결같이 수신(修身)을 근본으로 삼는다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 2. 목차 ==============
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.8), Inches(0.7),
                '목 차', font_size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.15), Inches(12.8))
    items = [
        [('Ⅰ', '개요 — 대학이란 무엇인가'),
         ('Ⅱ', '시대 배경과 사상사적 위치'),
         ('Ⅲ', '구성 — 경 1장 + 전 10장'),
         ('Ⅳ', '경 1장 — 대학의 헌법'),
         ('Ⅴ', '삼강령 ① 明明德'),
         ('Ⅵ', '삼강령 ② 親民·新民'),
         ('Ⅶ', '삼강령 ③ 止於至善'),
         ('Ⅷ', '팔조목 개관'),
         ('Ⅸ', '팔조목 ① 格物'),
         ('Ⅹ', '팔조목 ② 致知'),
         ('Ⅺ', '팔조목 ③ 誠意')],
        [('Ⅻ', '팔조목 ④ 正心'),
         ('ⅩⅢ', '팔조목 ⑤ 修身'),
         ('ⅩⅣ', '팔조목 ⑥ 齊家'),
         ('ⅩⅤ', '팔조목 ⑦ 治國'),
         ('ⅩⅥ', '팔조목 ⑧ 平天下'),
         ('ⅩⅦ', '명구절 8선'),
         ('ⅩⅧ', '7대 핵심 메시지'),
         ('ⅩⅨ', '현대적 의의'),
         ('ⅩⅩ', '다른 고전과의 비교'),
         ('ⅩⅩⅠ', '마무리')],
    ]
    for col, group in enumerate(items):
        x = 0.7 + col * 6.4
        top = 1.5
        for num, title in group:
            add_textbox(slide, Inches(x), Inches(top), Inches(1.0), Inches(0.4),
                        num, font_size=15, bold=True, color=ACCENT)
            add_textbox(slide, Inches(x + 1.0), Inches(top), Inches(5.3), Inches(0.4),
                        title, font_size=15, color=INK)
            top += 0.5


# ============== Ⅰ. 개요 ==============
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '대학(大學)이란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                '"큰 배움" — 대인(大人)의 학(學)',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.5),
                '주희(朱熹)의 풀이 — "대인지학(大人之學)"',
                font_size=17, color=ACCENT, align=PP_ALIGN.CENTER)
    nums = [('약 1,700', '자(字)'), ('경 1장', '+ 전 10장'),
            ('삼강령', '+ 팔조목'), ('800', '년 표준')]
    for i, (n, lbl) in enumerate(nums):
        x = 0.6 + i * 3.05
        add_textbox(slide, Inches(x), Inches(4.0), Inches(2.9), Inches(1.0),
                    n, font_size=46, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(5.1), Inches(2.9), Inches(0.5),
                    lbl, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                '어린아이의 소학(小學)을 넘어 — "사람다운 어른이 되고 세상을 이롭게 하는 학문"',
                font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '고대 교육 — 8세 소학(灑掃應對) → 15세 대학(窮理正心修己治人)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_status(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '책의 위상 — 핵심 정보')
    rows = [
        ('서명',    '대학(大學)',                    '"큰 배움 · 대인지학"'),
        ('원출처',  '『예기(禮記)』 제42편',          '본래 독립 책이 아닌 한 편'),
        ('경문 저자','공자(孔子)의 말 → 증자(曾子) 기록', '전통적 견해'),
        ('전문 저자','증자의 뜻 → 문인(門人) 풀이',   '전 10장'),
        ('사서 독립','주희(朱熹, 1130~1200)',          '남송에서 사서로 묶음'),
        ('편집',    '대학장구(大學章句)',              '경 1장 + 전 10장 + 보망장'),
        ('영향',    '조선 500년 · 일본 막부기',        '동아시아 교육 표준'),
    ]
    top = 2.15
    for i, (tag, val, note) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.65), PALE)
        add_filled_rect(slide, Inches(2.85), Inches(top), Inches(10.0), Inches(0.65), bg)
        add_textbox(slide, Inches(0.55), Inches(top + 0.18), Inches(2.2), Inches(0.4),
                    tag, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.95), Inches(top + 0.05), Inches(4.5), Inches(0.5),
                    val, font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.5), Inches(top + 0.08), Inches(5.3), Inches(0.5),
                    note, font_size=13, color=SUB)
        top += 0.7


@S('Ⅰ. 개요')
def s_authors(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '성립과 저자 — 공자·증자·문인의 합작')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '경문(經文) 1장', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('공자(孔子)의 말을',
         {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('증자(曾子)가 기록한 것',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 단 205자에 압축',
         {'font_size': 14, 'space_before': 8}),
        ('• 삼강령·팔조목이 모두 담김',
         {'font_size': 14, 'space_before': 4}),
        ('• "헌법의 전문(前文)" 같은 성격',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '전문(傳文) 10장', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('증자의 뜻을',
         {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('문인(門人)이 풀어 기록',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 약 1,500자',
         {'font_size': 14, 'space_before': 8}),
        ('• 시경·서경·공자 어록 인용',
         {'font_size': 14, 'space_before': 4}),
        ('• 전 1~3장 = 삼강령 해설',
         {'font_size': 14, 'space_before': 4}),
        ('• 전 4~10장 = 팔조목 해설',
         {'font_size': 14}),
    ], line_spacing=1.3)


@S('Ⅰ. 개요')
def s_zhuxi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '주희(朱熹) — 사서로 독립시킨 800년의 결단')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(11.9), Inches(0.55),
                '남송(南宋) 주희 (1130~1200) — 『예기』에서 대학·중용을 독립시켜 사서로 묶음',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.3), Inches(11.9), Inches(3.5), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(3.45), Inches(11.3), Inches(3.3), [
        ('주희의 선언 — 대학은 "입덕지문(入德之門)"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"대학을 먼저 읽어 그 규모(規模)를 잡고',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 논어를 읽어 그 근본(根本)을 세우고',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 맹자를 읽어 그 발휘(發揮)를 보고',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 중용을 읽어 그 미묘(微妙)를 얻는다"',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('→ 사서 학습의 첫 책 — 800년간 동아시아 교육의 출발점',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.3)


@S('Ⅰ. 개요')
def s_daehakjanggu(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '대학장구(大學章句)와 보망장(補亡章)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('주희의 편집 — "대학장구(大學章句)"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('예기의 본문을 재구성하여 경 1장 + 전 10장으로 분류',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('조선 500년·일본 막부기 유학 교육의 표준본',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('보망장(補亡章) — 전 5장의 누락을 채움',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('주희는 전 5장(격물치지)이 누락됐다고 보고',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('자신이 직접 글을 지어 보충 → "보망장"',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"활연관통(豁然貫通)"의 명문이 여기서 나옴',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


# ============== Ⅱ. 시대·사상사 ==============
@S('Ⅱ. 시대·사상사')
def s_education_system(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대·사상사', page, total)
    add_title(slide, '고대 교육 체계 — 소학(小學)과 대학(大學)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.5),
                '소학(小學) — 8세부터', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"쇄소응대진퇴(灑掃應對進退)"',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('일상의 예절', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 청소(灑掃)', {'font_size': 14, 'space_before': 8}),
        ('• 응대(應對)', {'font_size': 14, 'space_before': 4}),
        ('• 진퇴(進退)', {'font_size': 14, 'space_before': 4}),
        ('• 육예(六藝)', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
        ('  예·악·사·어·서·수', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.9), Inches(0.5),
                '대학(大學) — 15세부터', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"궁리정심수기치인"',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('(窮理正心修己治人)',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 궁리(窮理) — 이치 궁구',
         {'font_size': 14, 'space_before': 8}),
        ('• 정심(正心) — 마음 바르게',
         {'font_size': 14, 'space_before': 4}),
        ('• 수기(修己) — 자기 닦음',
         {'font_size': 14, 'space_before': 4}),
        ('• 치인(治人) — 남을 다스림',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '8세 일상 예절 → 15세 어른의 학(學) — 단계적 성장의 교육 체계',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 시대·사상사')
def s_ipdeokjimun(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대·사상사', page, total)
    add_title(slide, '사서(四書)의 입문 — 입덕지문(入德之門)')
    books = [
        ('대학(大學)',  '規模 — 규모',  '전체의 틀',    '먼저 읽어 큰 그림을 잡는다', True),
        ('논어(論語)',  '根本 — 근본',  '인(仁)의 뿌리','두 번째로 근본을 세운다', False),
        ('맹자(孟子)',  '發揮 — 발휘',  '활달한 논변',  '세 번째로 발휘를 본다', False),
        ('중용(中庸)',  '微妙 — 미묘',  '형이상학',     '마지막으로 미묘함을 얻는다', False),
    ]
    top = 2.4
    for name, role, desc, when, is_first in books:
        c = ACCENT if is_first else SUB
        bg_right = PALE if is_first else RGBColor(0xFA, 0xFA, 0xFA)
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.95), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(2.5), Inches(0.5),
                    name, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(2.5), Inches(0.95), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.27), Inches(2.5), Inches(0.5),
                    role, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.8), Inches(top), Inches(7.0), Inches(0.95), bg_right)
        add_textbox(slide, Inches(6.0), Inches(top + 0.1), Inches(6.7), Inches(0.4),
                    desc, font_size=13, bold=True, color=INK)
        add_textbox(slide, Inches(6.0), Inches(top + 0.5), Inches(6.7), Inches(0.4),
                    when, font_size=12, color=SUB)
        top += 1.05
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '주희가 정한 800년의 학습 순서 — 대학이 곧 출발점',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 시대·사상사')
def s_joseon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대·사상사', page, total)
    add_title(slide, '조선 유학의 중심 — 대학의 영향')
    items = [
        ('퇴계 이황',  '성학십도(聖學十圖)',
         '대학의 삼강령·팔조목을 기본 골격으로 한 선조에게 올린 도설'),
        ('율곡 이이',  '성학집요(聖學輯要)',
         '대학의 구조로 쓴 제왕학(帝王學) — 선조 시대'),
        ('세종대왕',   '일신우일신(日日新)을 좌우명으로',
         '집무실에 걸어두고 매일 자기 점검'),
        ('정약용',     '대학공의(大學公議)',
         '실학의 시각에서 재해석한 대학'),
    ]
    top = 2.5
    for name, work, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    name, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(3.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.22), Inches(3.5), Inches(0.5),
                    work, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.8), Inches(top), Inches(6.0), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.0), Inches(top + 0.22), Inches(5.7), Inches(0.5),
                    desc, font_size=12, color=INK)
        top += 0.97
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"군주·관료·선비 모두 반드시 외워야 하는 책" — 조선 500년 교육의 중심',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 시대·사상사')
def s_east_asia(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대·사상사', page, total)
    add_title(slide, '동아시아 리더십 교육의 원전')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '"수기치인(修己治人)의 완성된 프레임 — 자기와 세상을 잇는 다리"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(3.4), [
        ('대학이 유학에서 차지하는 4가지 위치',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('① 사서(四書)의 입문 — 논어·맹자·중용의 뼈대',
         {'font_size': 15, 'space_before': 8}),
        ('② 수기치인의 완성된 프레임 — 개인의 덕과 세상의 다스림은 하나의 연속체',
         {'font_size': 15, 'space_before': 6}),
        ('③ 조선 유학의 중심 — 퇴계·율곡의 제왕학 모두 대학 골격',
         {'font_size': 15, 'space_before': 6}),
        ('④ 동아시아 리더십의 원전 — 한국·중국·일본·베트남 공통 교과서',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
    ], line_spacing=1.35)


# ============== Ⅲ. 구성 ==============
@S('Ⅲ. 구성')
def s_structure(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 구성', page, total)
    add_title(slide, '경(經) 1장 + 전(傳) 10장 — 헌법과 주석')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '경(經) 1장 — 총론', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('단 205자 안에 모든 것이 담겨 있다',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 삼강령(三綱領) 선언', {'font_size': 14, 'space_before': 8}),
        ('• 팔조목(八條目) 연쇄', {'font_size': 14, 'space_before': 6}),
        ('• 본말(本末) 선언', {'font_size': 14, 'space_before': 6}),
        ('• "수신위본"의 결론', {'font_size': 14, 'space_before': 6}),
        ('', {'font_size': 6}),
        ('"헌법의 전문(前文)" 같은 성격',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '전(傳) 10장 — 해설', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('약 1,500자 — 시·서·공자 어록 인용 해설',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 전1~3장 — 삼강령 해설', {'font_size': 14, 'space_before': 8}),
        ('  명명덕·신민·지어지선', {'font_size': 12, 'color': SUB}),
        ('• 전4장 — 본말 선언', {'font_size': 14, 'space_before': 6}),
        ('• 전5장 — 격물치지 (보망장)', {'font_size': 14, 'space_before': 6}),
        ('• 전6~10장 — 성의·정심~평천하', {'font_size': 14, 'space_before': 6}),
    ], line_spacing=1.3)


@S('Ⅲ. 구성')
def s_concentric(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 구성', page, total)
    add_title(slide, '동심원 구조 — 안에서 밖으로')
    # 동심원 시뮬레이션 — 8개 박스를 위계로 표현
    layers = [
        ('1', '격물 (格物)', '사물의 이치 궁구'),
        ('2', '치지 (致知)', '앎의 극진'),
        ('3', '성의 (誠意)', '뜻의 성실'),
        ('4', '정심 (正心)', '마음의 바름'),
        ('5', '수신 (修身)', '★ 축 — 모든 것의 근본'),
        ('6', '제가 (齊家)', '집을 가지런히'),
        ('7', '치국 (治國)', '나라를 다스림'),
        ('8', '평천하 (平天下)', '세계의 평화'),
    ]
    top = 2.0
    for i, (num, name, desc) in enumerate(layers):
        is_axis = (num == '5')
        c = ACCENT if is_axis else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_axis else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.6), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.12), Inches(0.7), Inches(0.4),
                    num, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(3.5), Inches(0.6), bg)
        add_textbox(slide, Inches(1.5), Inches(top + 0.13), Inches(3.4), Inches(0.4),
                    name, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.0), Inches(top), Inches(7.8), Inches(0.6), bg)
        add_textbox(slide, Inches(5.2), Inches(top + 0.13), Inches(7.5), Inches(0.4),
                    desc, font_size=14, color=INK,
                    bold=is_axis)
        top += 0.65
    add_textbox(slide, Inches(0.5), Inches(7.4), Inches(12.3), Inches(0.4),
                '1~4 = 내면 / 5 = 축(수신) / 6~8 = 외적 확장',
                font_size=11, color=SUB, align=PP_ALIGN.CENTER, bold=True)


@S('Ⅲ. 구성')
def s_brevity(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 구성', page, total)
    add_title(slide, '분량의 극적인 대비 — 간결함이 곧 권위')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(2.0),
                '205', font_size=160, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(4.5), Inches(5.5), Inches(0.5),
                '자(字)', font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(5.0), Inches(5.5), Inches(0.5),
                '— 경 1장 전체 분량',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER, bold=True)
    add_paragraphs(slide, Inches(7.0), Inches(2.4), Inches(6.0), Inches(4.5), [
        ('총 분량', {'bold': True, 'font_size': 20, 'color': ACCENT}),
        ('약 1,700자 — A4 세 장 분량', {'font_size': 16}),
        ('', {'font_size': 8}),
        ('그러나', {'bold': True, 'font_size': 20, 'color': ACCENT, 'space_before': 10}),
        ('동아시아 2,000년 지성사의 중심',
         {'font_size': 16, 'bold': True, 'color': INK}),
        ('', {'font_size': 8}),
        ('"간결함이 곧 권위였다"', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('손자병법(약 6,000자)보다도 짧지만', {'font_size': 14, 'color': SUB}),
        ('유학의 모든 것을 담은 압축의 미학', {'font_size': 14, 'color': SUB, 'bold': True}),
    ], line_spacing=1.3)


# ============== Ⅳ. 경 1장 ==============
@S('Ⅳ. 경 1장')
def s_gyeong_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '경 1장 — 205자의 "대학의 헌법"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('단 205자 안에 모든 것이 담겨 있다',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('삼강령 선언 + "지(止)"의 효능 + 본말 선언 +',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('팔조목 역순·순순 + "수신위본"의 결론',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('6개 원문의 흐름',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('① 삼강령 선언 — 大學之道, 在明明德, 在親民, 在止於至善',
         {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('② "지(止)"의 효능 — 知止 → 定 → 靜 → 安 → 慮 → 得',
         {'font_size': 13, 'align': PP_ALIGN.CENTER}),
        ('③ 본말 선언 — 物有本末, 事有終始',
         {'font_size': 13, 'align': PP_ALIGN.CENTER}),
        ('④⑤ 팔조목 — 천하에서 격물로(역순), 격물에서 천하로(순순)',
         {'font_size': 13, 'align': PP_ALIGN.CENTER}),
        ('⑥ 결론 — 自天子以至於庶人, 壹是皆以修身爲本',
         {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅳ. 경 1장')
def s_gyeong_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '원문 ① 삼강령 선언', '"대학의 도(道)는 어디에 있는가?"')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '大 學 之 道',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.7),
                '在 明 明 德   在 親 民   在 止 於 至 善',
                font_size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.5),
                '대학지도 · 재명명덕 · 재친민 · 재지어지선',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.6), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.6),
                '"대학의 도는 밝은 덕을 밝히는 데 있고',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.6),
                ' 백성을 새롭게 하는 데 있으며',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.6),
                ' 지극한 선에 머무는 데 있다"',
                font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"재(在)"의 반복이 리드미컬 — 셋이 병렬이 아니라 하나의 도(道)의 세 얼굴',
                font_size=12, color=SUB, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅳ. 경 1장')
def s_gyeong_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '원문 ② "지(止)"의 효능 — 5단계 내면 심리학')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.65),
                '知止而后有定, 定而后能靜, 靜而后能安, 安而后能慮, 慮而后能得',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('知止', '지지', '목표의식 — 지극한 선이 어디인지 안다'),
        ('定',   '정',   '결심 — 방향이 정해진다'),
        ('靜',   '정',   '평정 — 마음이 흔들리지 않는다'),
        ('安',   '안',   '안정 — 처한 상황 속에 편안하다'),
        ('慮',   '려',   '통찰 — 맑은 머리로 사유한다'),
        ('得',   '득',   '성취 — 마침내 얻는다'),
    ]
    top = 3.5
    for han, eum, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(1.0), Inches(0.4),
                    han, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(1.5), Inches(0.55), PALE)
        add_textbox(slide, Inches(1.7), Inches(top + 0.13), Inches(1.5), Inches(0.4),
                    eum, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.3), Inches(top), Inches(9.5), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.5), Inches(top + 0.13), Inches(9.2), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.6


@S('Ⅳ. 경 1장')
def s_gyeong_3(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '원문 ③ 본말의 선언 — 대학 전체의 방법론')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '物 有 本 末   事 有 終 始',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.9),
                '知 所 先 後   則 近 道 矣',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                '물유본말 · 사유종시 · 지소선후 · 즉근도의',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.95), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.6),
                '"사물에는 근본과 말단이 있고, 일에는 시작과 끝이 있다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.6),
                '"먼저 할 바와 나중 할 바를 알면 도(道)에 가까워진다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '→ 스티븐 코비의 "first things first"의 원형',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅳ. 경 1장')
def s_gyeong_4(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '원문 ④ 팔조목 역방향 — "천하"에서 "격물"로')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '"옛날에 천하에 밝은 덕을 밝히려는 자는…"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    chains = [
        ('欲明明德於天下', '천하에 밝은 덕을 밝히려', '先治其國', '먼저 그 나라를 다스렸다'),
        ('欲治其國',       '나라를 다스리려',          '先齊其家', '먼저 그 집을 가지런히 했다'),
        ('欲齊其家',       '집을 가지런히 하려',        '先修其身', '먼저 그 몸을 닦았다'),
        ('欲修其身',       '몸을 닦으려',              '先正其心', '먼저 그 마음을 바르게 했다'),
        ('欲正其心',       '마음을 바르게 하려',        '先誠其意', '먼저 그 뜻을 정성스럽게 했다'),
        ('欲誠其意',       '뜻을 정성스럽게 하려',      '先致其知', '먼저 그 앎을 극진히 했다'),
        ('致知',           '앎을 극진히 함은',          '在格物',   '사물의 이치를 궁구함에 있다'),
    ]
    top = 3.15
    for han_l, kor_l, han_r, kor_r in chains:
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.55), Inches(top + 0.05), Inches(2.2), Inches(0.5),
                    han_l, font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.85), Inches(top), Inches(3.2), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(2.95), Inches(top + 0.13), Inches(3.0), Inches(0.4),
                    kor_l, font_size=11, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(6.1), Inches(top + 0.13), Inches(0.4), Inches(0.4),
                    '→', font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.5), Inches(top), Inches(2.3), Inches(0.55), PALE)
        add_textbox(slide, Inches(6.55), Inches(top + 0.05), Inches(2.2), Inches(0.5),
                    han_r, font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(8.85), Inches(top), Inches(3.95), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(8.95), Inches(top + 0.13), Inches(3.8), Inches(0.4),
                    kor_r, font_size=11, color=INK, align=PP_ALIGN.CENTER)
        top += 0.6


@S('Ⅳ. 경 1장')
def s_gyeong_5(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '원문 ⑤ 팔조목 순방향 — "격물"에서 "천하"로')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '"사물이 궁구된 뒤에 앎이 이르고…"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    chains = [
        ('物格 → 知至',   '사물이 궁구된 뒤에 앎이 이르고'),
        ('知至 → 意誠',   '앎이 이른 뒤에 뜻이 성실해지고'),
        ('意誠 → 心正',   '뜻이 성실해진 뒤에 마음이 바르게 되고'),
        ('心正 → 身修',   '마음이 바르게 된 뒤에 몸이 닦이고'),
        ('身修 → 家齊',   '몸이 닦인 뒤에 집이 가지런해지고'),
        ('家齊 → 國治',   '집이 가지런해진 뒤에 나라가 다스려지고'),
        ('國治 → 天下平', '나라가 다스려진 뒤에 천하가 평안해진다'),
    ]
    top = 3.15
    for han, kor in chains:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(3.5), Inches(0.4),
                    han, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(8.5), Inches(0.55), PALE)
        add_textbox(slide, Inches(4.5), Inches(top + 0.13), Inches(8.2), Inches(0.4),
                    kor, font_size=14, color=INK)
        top += 0.6


@S('Ⅳ. 경 1장')
def s_dual_structure(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '팔조목의 이중 구조 — "방법과 목표가 서로를 끌고 밀어간다"')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '원문 ④ — 역순', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"천하 → 격물"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('목표에서 출발해', {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('방법으로 내려감', {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"무엇을 위해 무엇이 필요한가?"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '원문 ⑤ — 순순', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"격물 → 천하"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('방법에서 출발해', {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('목표로 올라감', {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"무엇을 하면 무엇에 이르는가?"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '같은 팔조목을 두 방향으로 — 문장 구조 자체가 "방법↔목표 상호작용"을 그린다',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 경 1장')
def s_gyeong_6(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '원문 ⑥ 결론 — 수신(修身)을 근본으로')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '自 天 子 以 至 於 庶 人',
                font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.9),
                '壹 是 皆 以 修 身 爲 本',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.5),
                '자천자이지어서인 · 일시개이수신위본',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.8), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
                '"천자(天子)로부터 서민(庶人)에 이르기까지',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.6),
                ' 모두 한결같이 수신(修身)을 근본으로 삼는다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '"일시(壹是)" — 한결같이, 신분과 무관하게 — 보편 교육의 선언',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅳ. 경 1장')
def s_universal(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 경 1장', page, total)
    add_title(slide, '"천자로부터 서민까지" — 보편 교육의 선언')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('수신이 모든 것의 축',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('팔조목 8개 중 가운데 자리인 "수신(修身)"이 대학의 축',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('앞 4조목(격물·치지·성의·정심) = 수신의 내적 준비',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('뒤 3조목(제가·치국·평천하) = 수신의 외적 확장',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"천자로부터 서민까지" — 혁명적 선언',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('유학은 사실상 신분을 뛰어넘는 보편 교육을 선언했다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('왕도 평민도 수신의 책임에서는 동등',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('2,400년 전의 이 한 문장은 매우 혁신적이었다',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


# ============== Ⅴ. 明明德 ==============
@S('Ⅴ. 明明德')
def s_myeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 明明德', page, total)
    add_title(slide, '明明德 — 내 안의 등불을 밝혀라')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '明\n德', font_size=130, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(4.2), Inches(2.4), Inches(8.7), Inches(4.5), [
        ('두 글자의 의미', {'bold': True, 'font_size': 18, 'color': ACCENT}),
        ('앞 明 = 동사 — "밝히다"', {'font_size': 16}),
        ('뒤 明 = 형용사 — "밝은"', {'font_size': 16}),
        ('德 = 덕(德)', {'font_size': 16}),
        ('', {'font_size': 6}),
        ('전체 의미', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('"이미 밝은 덕을 더 밝게 닦는다"',
         {'font_size': 18, 'bold': True, 'color': INK}),
        ('', {'font_size': 6}),
        ('핵심 전제', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('인간은 본래 밝은 덕(明德)을 타고난다',
         {'font_size': 15}),
        ('맹자의 성선설(性善說)과 이어지는 유학의 인간관',
         {'font_size': 14, 'color': SUB}),
    ], line_spacing=1.3)


@S('Ⅴ. 明明德')
def s_why_clouded(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 明明德', page, total)
    add_title(slide, '왜 덕이 가려지는가 — 그리고 교육의 본질')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('덕(德)이 가려지는 세 가지 원인',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('① 기질(氣質)의 치우침 — 타고난 편향',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('② 욕망의 침식 — 끝없는 욕심',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('③ 습관의 덧칠 — 반복된 부정적 패턴',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('교육의 본질에 대한 혁명적 정의',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"새 것을 채우는 것"이 아니라',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"가린 것을 걷어내는 것"',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('루소·페스탈로치·몬테소리까지 이어진 현대 교육사상의 동양적 원형',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅴ. 明明德')
def s_three_quotes(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 明明德', page, total)
    add_title(slide, '전 1장의 3가지 인용 — 명덕은 "끊임없이 닦는 일"')
    items = [
        ('서경 강고(康誥)',  '克 明 德',         '극명덕',
         '능히 덕을 밝혔다'),
        ('서경 태갑(太甲)',  '顧 諟 天 之 明 命', '고시천지명명',
         '이 하늘의 밝은 명을 돌아본다'),
        ('서경 제전(帝典)',  '克 明 峻 德',       '극명준덕',
         '능히 큰 덕을 밝혔다'),
    ]
    top = 2.4
    for source, han, eum, kor in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.8), Inches(1.05), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.32), Inches(2.8), Inches(0.5),
                    source, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.5), Inches(top), Inches(3.5), Inches(1.05), PALE)
        add_textbox(slide, Inches(3.5), Inches(top + 0.15), Inches(3.5), Inches(0.4),
                    han, font_size=18, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(3.5), Inches(top + 0.6), Inches(3.5), Inches(0.4),
                    eum, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.1), Inches(top), Inches(5.7), Inches(1.05),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.3), Inches(top + 0.35), Inches(5.4), Inches(0.5),
                    kor, font_size=15, color=INK)
        top += 1.18
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '한 번 밝혀 끝이 아니라 매일 닦아야 빛을 유지한다 — 자기성찰의 영구 작업',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅵ. 親民·新民 ==============
@S('Ⅵ. 親民·新民')
def s_two_interpretations(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 親民·新民', page, total)
    add_title(slide, '친민·신민 — 유학사의 큰 갈림길')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '주희의 新民 (신민)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"親"은 "新"의 오기(誤記)',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"백성을 새롭게 한다"',
         {'font_size': 19, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 자기 덕을 밝힌 뒤', {'font_size': 14, 'space_before': 8}),
        ('• 백성을 교화하여 새롭게 함', {'font_size': 14}),
        ('• 상하 관계 · 가르침의 관점', {'font_size': 14, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '왕양명의 親民 (친민)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('원문 그대로 "親"으로 읽음',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"백성과 친히 하나가 된다"',
         {'font_size': 19, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 백성과 친(親)히 만남', {'font_size': 14, 'space_before': 8}),
        ('• 함께 변화함', {'font_size': 14}),
        ('• 수평 관계 · 공감의 관점', {'font_size': 14, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '오늘날의 통합 — "자기의 변화가 타인의 변화로 이어지도록 한다"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅵ. 親民·新民')
def s_ilsin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 親民·新民', page, total)
    add_title(slide, '일신우일신(日日新) — 탕왕(湯王)의 반명(盤銘)')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '苟 日 新   日 日 新   又 日 新',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
                '구일신 일일신 우일신',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.75), Inches(12.3), Inches(0.4),
                '— 은(殷) 탕왕(湯王)이 세숫대야(盤)에 새긴 글',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '"진실로 하루가 새롭거든, 날로 새롭고, 또 날로 새로워라"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.5), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.75), Inches(11.3), Inches(1.3), [
        ('세숫대야에 새긴 글 — 매일 얼굴을 씻을 때마다 자기를 새롭게 하라',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('세종이 집무실에 걸어두었고, 이승훈·김구가 좌우명으로 삼았다',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅵ. 親民·新民')
def s_propagate(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 親民·新民', page, total)
    add_title(slide, '변화의 파급 — 작신민(作新民)·주수구방(周雖舊邦)')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '서경 — 作新民', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(3.1), Inches(5.9), Inches(1.0),
                '作 新 民',
                font_size=66, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(4.2), Inches(5.9), Inches(0.4),
                '작신민', font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(4.8), Inches(5.7), Inches(0.5),
                '"백성을 새롭게 하라"',
                font_size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(5.4), Inches(5.7), Inches(1.2),
                '내가 어제의 나에서 벗어나지 못하면\n내 주변도 어제와 같다',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '시경 — 周雖舊邦', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.1), Inches(5.9), Inches(0.7),
                '周 雖 舊 邦   其 命 維 新',
                font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.95), Inches(5.9), Inches(0.4),
                '주수구방 기명유신', font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.9), Inches(4.5), Inches(5.7), Inches(0.5),
                '"주(周)는 비록 옛 나라이나',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.9), Inches(5.0), Inches(5.7), Inches(0.5),
                ' 그 명(命)은 오직 새롭다"',
                font_size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.9), Inches(5.6), Inches(5.7), Inches(1.0),
                '오래된 조직도\n사명(命)은 새로울 수 있다',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅶ. 止於至善 ==============
@S('Ⅶ. 止於至善')
def s_jieojiseon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 止於至善', page, total)
    add_title(slide, '止於至善 — 정상에 머무는 능력')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '止', font_size=240, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(4.2), Inches(2.4), Inches(8.7), Inches(4.5), [
        ('"지(止)"의 의미', {'bold': True, 'font_size': 20, 'color': ACCENT}),
        ('일회적 도달이 아니라', {'font_size': 17}),
        ('"머물러 유지하는 상태(安住)"', {'font_size': 18, 'bold': True, 'color': INK}),
        ('', {'font_size': 8}),
        ('정상에 한 번 오르는 것이 아니라', {'font_size': 15, 'color': SUB, 'space_before': 10}),
        ('그 자리에 안주(安住)하는 능력', {'font_size': 15, 'color': SUB}),
        ('', {'font_size': 8}),
        ('지극한 선(至善)이란?', {'bold': True, 'font_size': 20, 'color': ACCENT, 'space_before': 10}),
        ('각자의 자리에서 추구할 최고의 모습', {'font_size': 15}),
        ('역할마다 "지선"이 다르다 — 다음 슬라이드',
         {'font_size': 14, 'color': SUB, 'bold': True}),
    ], line_spacing=1.3)


@S('Ⅶ. 止於至善')
def s_five_roles(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 止於至善', page, total)
    add_title(slide, '5가지 자리(役割)의 지선(至善) — 시경 인용')
    items = [
        ('군주(君)',  '仁', '인', '어짊'),
        ('신하(臣)',  '敬', '경', '공경'),
        ('자식(子)',  '孝', '효', '효도'),
        ('부모(父)',  '慈', '자', '자애'),
        ('벗(友)',    '信', '신', '믿음'),
    ]
    top = 2.3
    for role, han, eum, kor in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    role, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(3.4), Inches(top + 0.12), Inches(1.0), Inches(0.5),
                    '止 於', font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.6), Inches(top), Inches(1.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(4.6), Inches(top + 0.05), Inches(1.5), Inches(0.5),
                    han, font_size=34, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.3), Inches(top), Inches(6.5), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.5), Inches(top + 0.08), Inches(6.3), Inches(0.4),
                    f'{eum} — {kor}',
                    font_size=15, bold=True, color=INK)
        add_textbox(slide, Inches(6.5), Inches(top + 0.5), Inches(6.3), Inches(0.4),
                    f'{role}는 {kor}({han})에 머문다',
                    font_size=12, color=SUB)
        top += 0.97
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '한 사람이 군주이자 부모이자 벗 — 각 역할의 지선에 모두 머묾이 "지어지선"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅶ. 止於至善')
def s_bird(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 止於至善', page, total)
    add_title(slide, '"새조차 머물 곳을 안다" — 시경의 비유')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('緡 蠻 黃 鳥   止 于 丘 隅',
         {'font_size': 28, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('민만황조 지우구우 — "꾀꼬리도 자기 머무를 자리를 안다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.7),
                '"새조차 머물 곳을 아는데',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.7),
                ' 사람이 새만 못할까?"',
                font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
                '"지(止) 한 글자"가 대학 전체를 관통한다',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅷ. 팔조목 개관 ==============
@S('Ⅷ. 팔조목 개관')
def s_8items(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 팔조목 개관', page, total)
    add_title(slide, '팔조목 — 안에서 밖으로 펼쳐지는 동심원')
    items = [
        ('1', '格物', '격물', '사물의 이치를 궁구'),
        ('2', '致知', '치지', '앎을 극진히 함'),
        ('3', '誠意', '성의', '뜻을 성실히 함'),
        ('4', '正心', '정심', '마음을 바르게 함'),
        ('5', '修身', '수신', '몸을 닦음 — ★ 축'),
        ('6', '齊家', '제가', '집을 가지런히 함'),
        ('7', '治國', '치국', '나라를 다스림'),
        ('8', '平天下', '평천하', '천하를 평안케 함'),
    ]
    top = 2.2
    for num, han, eum, kor in items:
        is_axis = (num == '5')
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_axis else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(0.7), Inches(0.4),
                    num, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.5), Inches(0.55), bg)
        add_textbox(slide, Inches(1.4), Inches(top + 0.1), Inches(2.5), Inches(0.4),
                    han, font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.0), Inches(top), Inches(1.5), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.0), Inches(top + 0.13), Inches(1.5), Inches(0.4),
                    eum, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.6), Inches(top), Inches(7.2), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.8), Inches(top + 0.13), Inches(7.0), Inches(0.4),
                    kor, font_size=14, color=INK, bold=is_axis)
        top += 0.62


@S('Ⅷ. 팔조목 개관')
def s_3stages(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 팔조목 개관', page, total)
    add_title(slide, '知 · 體 · 用 — 팔조목의 3단계 구조')
    stages = [
        ('知 (지)',  '격물·치지 (1~2)',  '앎의 단계',     ACCENT,
         '대상에 다가가 이치를 파악한다'),
        ('體 (체)',  '성의·정심·수신 (3~5)', '수양의 단계', RGBColor(0xA0, 0x40, 0x40),
         '안에서 자기를 닦는다 — 핵심'),
        ('用 (용)',  '제가·치국·평천하 (6~8)', '실천의 단계', SUB,
         '밖으로 펼쳐 세상에 작용한다'),
    ]
    top = 2.4
    for tag, scope, role, color, desc in stages:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(1.4), color)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(2.0), Inches(0.5),
                    tag, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.8), Inches(top), Inches(3.5), Inches(1.4), PALE)
        add_textbox(slide, Inches(2.8), Inches(top + 0.32), Inches(3.5), Inches(0.4),
                    scope, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.8), Inches(top + 0.78), Inches(3.5), Inches(0.4),
                    role, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.5), Inches(top), Inches(6.3), Inches(1.4),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.7), Inches(top + 0.48), Inches(6.0), Inches(0.5),
                    desc, font_size=15, color=INK)
        top += 1.55


# ============== Ⅸ. 格物 ==============
@S('Ⅸ. 格物')
def s_gyeokmul(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 格物', page, total)
    add_title(slide, '格物 — 사물의 이치를 끝까지 파고든다')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '格', font_size=240, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(4.2), Inches(2.4), Inches(8.7), Inches(4.5), [
        ('"격(格)"의 두 해석', {'bold': True, 'font_size': 20, 'color': ACCENT}),
        ('', {'font_size': 6}),
        ('주희 — 이르다(至)', {'bold': True, 'font_size': 18, 'space_before': 8}),
        ('"사물에 이르러 그 이치를 궁구한다"',
         {'font_size': 15, 'color': SUB}),
        ('= 객관적 탐구', {'font_size': 14, 'color': INK}),
        ('', {'font_size': 6}),
        ('왕양명 — 바로잡다(正)', {'bold': True, 'font_size': 18, 'space_before': 10}),
        ('"사물의 뜻(事)을 바르게 한다"',
         {'font_size': 15, 'color': SUB}),
        ('= 마음의 부정을 바로잡는 것', {'font_size': 14, 'color': INK}),
        ('', {'font_size': 6}),
        ('두 해석이 유학사 800년의 분기점',
         {'font_size': 13, 'color': ACCENT, 'bold': True, 'space_before': 10}),
    ], line_spacing=1.3)


@S('Ⅸ. 格物')
def s_jzhu_yang(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 格物', page, total)
    add_title(slide, '주희 vs 왕양명 — 해석의 대립')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(5.9), Inches(0.6), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(0.5),
                '주희 — 정주학(程朱學)', font_size=18, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(6.9), Inches(2.2), Inches(5.9), Inches(0.6), SUB)
    add_textbox(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(0.5),
                '왕양명 — 양명학(陽明學)', font_size=18, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    rows = [
        ('격(格)의 풀이',  '이르다(至)',          '바로잡다(正)'),
        ('탐구 방향',       '밖 → 안 (대상 → 마음)','안 → 밖 (마음 → 세상)'),
        ('대상',             '천하 사물의 이치',     '마음의 부정(不正)'),
        ('방법',             '독서·관찰·추론',       '치양지(致良知)'),
        ('주된 영역',       '학문·과학적 탐구',     '내면 수양·실천'),
        ('영향',             '조선 주류 (퇴계·율곡)','조선 일부 (정제두 등)'),
    ]
    top = 3.0
    for tag, jzhu, yang in rows:
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.55), Inches(top + 0.13), Inches(2.4), Inches(0.4),
                    tag, font_size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.05), Inches(top), Inches(3.95), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.2), Inches(top + 0.13), Inches(3.75), Inches(0.4),
                    jzhu, font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.05), Inches(top), Inches(5.8), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.2), Inches(top + 0.13), Inches(5.6), Inches(0.4),
                    yang, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        top += 0.62


@S('Ⅸ. 格物')
def s_bomangjang(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 格物', page, total)
    add_title(slide, '주희의 보망장(補亡章) — 활연관통(豁然貫通)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.4), [
        ('"내 앎을 이루려 하면 사물에 나아가 그 이치를 궁구하는 데 있다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        (' …오래도록 힘을 쓰다가 하루아침에',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('豁 然 貫 通',
         {'font_size': 36, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('활연관통 — 활짝 통하다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('모든 사물의 겉과 속이 이르지 않는 바가 없고',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('내 마음의 전체와 큰 작용이 밝지 않은 바가 없게 된다"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.6), [
        ('현대적 의미',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('반복된 탐구가 어느 순간 질적 도약을 일으키는 학습의 구조',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('= 오늘의 "aha moment"·"insight" 이론',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


# ============== Ⅹ. 致知 ==============
@S('Ⅹ. 致知')
def s_chiji(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 致知', page, total)
    add_title(slide, '致知 — 앎을 극진히 한다')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.45), Inches(11.3), Inches(1.8), [
        ('"격물과 치지는 동전의 양면"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('격물(格物) = 대상에 다가가는 노력',
         {'font_size': 17, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('치지(致知) = 나의 앎이 극진해지는 결과',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.6), Inches(4.5), Inches(11.9), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.9), Inches(4.65), Inches(11.3), Inches(2.2), [
        ('현대 과학과의 만남 — "격물(格物)"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('19세기 일본의 니시 아마네는 "science"를 처음 번역할 때',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"격물학(格物學)"이라 했다',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('조선 실학자(홍대용·박지원·정약용)도 격물의 대상을 확장 → 실사구시(實事求是)',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅹ. 致知')
def s_hwarol(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 致知', page, total)
    add_title(slide, '활연관통(豁然貫通)의 순간 — 학습의 질적 도약')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '豁 然 貫 通',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                '활연관통 — "활짝 통하다"',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.6), Inches(4.7), Inches(11.9), Inches(2.3), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(4.85), Inches(11.3), Inches(2.1), [
        ('주희: "오랜 격물이 쌓인 어느 날',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        (' 불현듯 전체가 환히 꿰뚫리는 순간"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('축적된 격물이 어느 순간 통합적 이해로 도약',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"aha moment" · "사고의 임계점"의 동양적 정식',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


# ============== Ⅺ. 誠意 ==============
@S('Ⅺ. 誠意')
def s_seongui(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 誠意', page, total)
    add_title(slide, '誠意 — 자기기만 없는 진심')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '所 謂 誠 其 意 者   毋 自 欺 也',
                font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '소위성기의자 무자기야',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
                '— 전 6장', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '"이른바 뜻을 성실히 한다는 것은',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.6),
                ' 스스로를 속이지 않는 것이다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(6.1), Inches(11.9), Inches(1.0), PALE)
    add_textbox(slide, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.5),
                '"如惡惡臭 如好好色" — 나쁜 냄새를 싫어하듯 하고, 아름다운 것을 좋아하듯 한다',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.4),
                '여악악취 여호호색 — 감정과 행동이 일치하는 진정한 마음',
                font_size=12, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅺ. 誠意')
def s_mujagi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 誠意', page, total)
    add_title(slide, '毋自欺 — 가장 깊은 병, 자기기만')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '毋 自 欺',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                '무자기 — "스스로를 속이지 말라"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(2.1), [
        ('"자기(自欺)는 남을 속이는 것보다 더 깊은 병"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 악을 싫어한다면서도 실은 즐기는 마음',
         {'font_size': 15, 'space_before': 6}),
        ('• 선을 좋아한다면서도 실은 귀찮아하는 마음',
         {'font_size': 15, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('현대 — "authenticity(진정성)"의 가장 정밀한 동양적 정식',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅺ. 誠意')
def s_sindok(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 誠意', page, total)
    add_title(slide, '愼獨 — 혼자 있을 때를 삼가다')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '愼 獨',
                font_size=130, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                '신독 — "홀로 있을 때를 삼가다"',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(2.1), [
        ('남이 보지 않을 때의 나와 보는 때의 내가 같다면',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('뜻이 성실한 것이다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"사람이 자기를 보는 것이 폐와 간을 꿰뚫어 보는 것과 같으니',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 무엇이 유익하겠는가?"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


# ============== Ⅻ. 正心 ==============
@S('Ⅻ. 正心')
def s_jeongsim(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 正心', page, total)
    add_title(slide, '正心 — 마음의 4가지 치우침을 경계')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.9), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(11.9), Inches(0.55),
                '"마음이 ~한 바가 있으면 그 바름(正)을 얻지 못한다"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('忿 懥', '분치', '분노 · 노여움',     '판단을 흐리는 격노'),
        ('恐 懼', '공구', '두려움 · 불안',     '결정을 막는 위축'),
        ('好 樂', '호요', '좋아함 · 집착',     '평정을 무너뜨리는 편애'),
        ('憂 患', '우환', '근심 · 걱정',       '현재를 잠식하는 미래 불안'),
    ]
    top = 3.4
    for han, eum, role, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(0.8), PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.18), Inches(2.0), Inches(0.5),
                    han, font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.7), Inches(top), Inches(1.5), Inches(0.8),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(2.7), Inches(top + 0.23), Inches(1.5), Inches(0.4),
                    eum, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(3.5), Inches(0.8),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.3), Inches(top + 0.23), Inches(3.5), Inches(0.4),
                    role, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.9), Inches(top), Inches(4.9), Inches(0.8),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(8.1), Inches(top + 0.23), Inches(4.6), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.93


@S('Ⅻ. 正心')
def s_simbujae(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 正心', page, total)
    add_title(slide, '心不在焉 — 마음이 없으면…')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '心 不 在 焉',
                font_size=80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '심부재언 — "마음이 여기 없으면"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    items = [
        ('視 而 不 見', '시이불견', '보아도 보이지 않는다'),
        ('聽 而 不 聞', '청이불문', '들어도 들리지 않는다'),
        ('食 而 不 知 其 味', '식이불지기미', '먹어도 맛을 모른다'),
    ]
    top = 4.5
    for han, eum, kor in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(4.0), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(4.0), Inches(0.4),
                    han, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.7), Inches(top), Inches(2.5), Inches(0.65), PALE)
        add_textbox(slide, Inches(4.7), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    eum, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.3), Inches(top), Inches(5.5), Inches(0.65),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.5), Inches(top + 0.15), Inches(5.2), Inches(0.4),
                    kor, font_size=15, color=INK)
        top += 0.75
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '리더가 감정에 휘둘리면 팀 전체가 "보지 못하고 듣지 못한다"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅻ. 正心')
def s_jeongsim_modern(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 正心', page, total)
    add_title(slide, '현대 — 감정 조절(Emotional Regulation)의 고전')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('정심(正心) = 마음챙김(Mindfulness)의 동양적 원형',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"분치·공구·호요·우환"의 4가지 치우침은',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('현대 뇌과학의 "편도체 과잉 활성"과 정확히 일치',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 응용',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 다니엘 골먼 "감성지능(EQ)" — 자기 인식·자기 관리',
         {'font_size': 14, 'space_before': 6}),
        ('• 행동경제학 — 감정에 휩쓸린 의사결정의 오류',
         {'font_size': 14, 'space_before': 4}),
        ('• "분노했을 때 결정하지 말라"는 모든 리더십 원칙의 원조',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== XⅢ. 修身 ==============
@S('XⅢ. 修身')
def s_susin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 修身', page, total)
    add_title(slide, '修身 — 다섯 가지 편벽(辟)을 뛰어넘어라')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.9), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(11.9), Inches(0.55),
                '"사람은 ~ 하는 대상에 대해 편벽(辟)된다"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('親愛', '친애', '사랑하는 대상',     '결점을 보지 못함'),
        ('賤惡', '천오', '천하게 여겨 미워하는 대상', '장점을 보지 못함'),
        ('畏敬', '외경', '두려워하며 공경하는 대상', '비판을 못 함'),
        ('哀矜', '애긍', '불쌍히 여기는 대상', '엄격함을 못 가짐'),
        ('敖惰', '오타', '교만·게으름 부리는 대상', '존중을 못 함'),
    ]
    top = 3.4
    for han, eum, role, prob in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.5), Inches(0.65), PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.12), Inches(1.5), Inches(0.4),
                    han, font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.2), Inches(top), Inches(1.5), Inches(0.65),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(2.2), Inches(top + 0.15), Inches(1.5), Inches(0.4),
                    eum, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.8), Inches(top), Inches(4.5), Inches(0.65),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.95), Inches(top + 0.15), Inches(4.3), Inches(0.4),
                    role, font_size=13, color=INK)
        add_filled_rect(slide, Inches(8.4), Inches(top), Inches(4.4), Inches(0.65),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(8.55), Inches(top + 0.15), Inches(4.2), Inches(0.4),
                    f'→ {prob}', font_size=13, color=INK, bold=True)
        top += 0.75


@S('XⅢ. 修身')
def s_susin_balance(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 修身', page, total)
    add_title(slide, '愛而知其惡, 惡而知其美 — 수신의 경지')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '愛 而 知 其 惡',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.9),
                '惡 而 知 其 美',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                '애이지기악 · 오이지기미',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.9), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
                '"사랑하되 그의 결점을 알고,',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
                ' 미워하되 그의 장점을 안다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '→ 인지편향(confirmation bias) 극복 — 현대 360도 평가의 원형',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('XⅢ. 修身')
def s_susin_axis(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 修身', page, total)
    add_title(slide, '수신 — 팔조목 8개의 축(軸)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('"수신(修身)이 모든 것의 근본"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('팔조목 8개 중 가운데 자리인 "수신"이 대학의 축',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('앞 4조목(격물·치지·성의·정심) = 수신의 내적 준비',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('뒤 3조목(제가·치국·평천하) = 수신의 외적 확장',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('대학 전체가 결국 한 구절로 수렴',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"自天子以至於庶人 壹是皆以修身爲本"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"천자로부터 서민까지 한결같이 수신을 근본으로 삼는다"',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


# ============== XⅣ. 齊家 ==============
@S('XⅣ. 齊家')
def s_jega(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 齊家', page, total)
    add_title(slide, '齊家 — 집은 곧 축소된 나라')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '其 家 不 可 敎   而 能 敎 人 者   無 之',
                font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
                '기가불가교 이능교인자 무지',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.75), Inches(12.3), Inches(0.4),
                '— 전 9장', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '"자기 집도 가르치지 못하면서',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.6),
                ' 남을 가르칠 수 있는 자는 없다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                '→ 리더십의 진정한 시험대는 "가정" — 가족에게 보이는 모습이 진짜 얼굴',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('XⅣ. 齊家')
def s_hyo_je_ja(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 齊家', page, total)
    add_title(slide, '효(孝)·제(悌)·자(慈) — 가정의 3덕은 사회의 연습장')
    items = [
        ('孝', '효', '효도',  '所以事君 — 군주를 섬기는 방법',
         '가정에서의 孝 → 사회의 忠(충)'),
        ('悌', '제', '공경',  '所以事長 — 윗사람을 섬기는 방법',
         '가정에서의 悌 → 사회의 공경'),
        ('慈', '자', '자애',  '所以使衆 — 무리를 부리는 방법',
         '가정에서의 慈 → 사회의 사랑'),
    ]
    top = 2.4
    for han, eum, kor, principle, modern in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.5), Inches(1.2), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.35), Inches(1.5), Inches(0.6),
                    han, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.2), Inches(top), Inches(2.0), Inches(1.2), PALE)
        add_textbox(slide, Inches(2.2), Inches(top + 0.3), Inches(2.0), Inches(0.4),
                    eum, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.2), Inches(top + 0.7), Inches(2.0), Inches(0.4),
                    kor, font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(8.5), Inches(1.2),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.5), Inches(top + 0.15), Inches(8.2), Inches(0.45),
                    principle, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(4.5), Inches(top + 0.65), Inches(8.2), Inches(0.45),
                    modern, font_size=13, color=SUB)
        top += 1.35


@S('XⅣ. 齊家')
def s_one_family(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 齊家', page, total)
    add_title(slide, '一家仁 一國興仁 — 한 집의 인이 한 나라의 인을 일으킨다')
    items = [
        ('一 家 仁', '일가인', '一國興仁 — 한 나라가 어짊을 일으킨다'),
        ('一 家 讓', '일가양', '一國興讓 — 한 나라가 사양을 일으킨다'),
        ('一 人 貪 戾', '일인탐려', '一國作亂 — 한 나라에 난이 일어난다'),
    ]
    top = 2.5
    for han, eum, result in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.95), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.2), Inches(3.5), Inches(0.5),
                    han, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.6), Inches(3.5), Inches(0.3),
                    eum, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(8.5), Inches(0.95), PALE)
        add_textbox(slide, Inches(4.5), Inches(top + 0.27), Inches(8.2), Inches(0.5),
                    result, font_size=15, bold=True, color=INK)
        top += 1.1
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '"한 사람·한 집의 행위가 한 나라의 운명을 결정한다"',
                font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
                '요·순(仁)을 따랐고, 걸·주(暴)를 따랐다 — 조직 문화는 최상위 리더의 사적 행동에서 결정된다',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== XⅤ. 治國 ==============
@S('XⅤ. 治國')
def s_chiguk(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 治國', page, total)
    add_title(slide, '治國 — 제가(齊家)의 확장')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '"所謂治國必先齊其家者…"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(1.7), [
        ('"나라를 다스리려면 반드시 먼저 집을 가지런히 하라"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"자기 집도 가르치지 못하면서 남을 가르칠 수 있는 자가 없기 때문이다"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.4), [
        ('"君子不出家而成敎於國"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"군자는 집 밖으로 나가지 않고도 나라에 교화를 이룬다"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('XⅤ. 治國')
def s_lead_by_life(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 治國', page, total)
    add_title(slide, '"교화는 언어가 아니라 삶으로 전해진다"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('교화(敎化)는 말이 아니라 삶 자체로 이루어진다',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('정치 리더가 집안을 가지런히 하는 것이',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('이미 가장 강력한 정치 행위',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 리더십의 발견과 일치',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 변혁적 리더십(Transformational Leadership) — 행동으로의 모범',
         {'font_size': 14, 'space_before': 6}),
        ('• "Walk the talk" — 말한 대로 살아라',
         {'font_size': 14, 'space_before': 4}),
        ('• 조직 문화는 CEO의 사적 행동에서 결정된다',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== XⅥ. 平天下 ==============
@S('XⅥ. 平天下')
def s_pyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 平天下', page, total)
    add_title(slide, '平天下 — 혈구지도(絜矩之道)')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '絜 矩 之 道',
                font_size=80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '혈구지도',
                font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
                '— 전 10장', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.55), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8), [
        ('"목수의 곱자(矩)로 재는 것처럼',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        (' 자기를 기준 삼아 타인을 헤아리는 도(道)"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('역지사지(易地思之)의 가장 정밀한 정식 — 6방향의 원칙',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('XⅥ. 平天下')
def s_six_directions(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 平天下', page, total)
    add_title(slide, '혈구지도의 6방향 — "내가 싫었던 것은 남에게 하지 말라"')
    items = [
        ('上', '윗사람에게서 싫었던 것을', '아랫사람에게 시키지 말라'),
        ('下', '아랫사람에게서 싫었던 것으로', '윗사람을 섬기지 말라'),
        ('前', '앞사람에게서 싫었던 것으로', '뒷사람을 인도하지 말라'),
        ('後', '뒷사람에게서 싫었던 것으로', '앞사람을 따르지 말라'),
        ('右', '오른쪽에서 싫었던 것으로', '왼쪽에 건네지 말라'),
        ('左', '왼쪽에서 싫었던 것으로', '오른쪽에 건네지 말라'),
    ]
    top = 2.2
    for han, source, action in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.85), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(0.85), Inches(0.4),
                    han, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.55), Inches(top), Inches(5.6), Inches(0.65), PALE)
        add_textbox(slide, Inches(1.7), Inches(top + 0.15), Inches(5.4), Inches(0.4),
                    source, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.25), Inches(top), Inches(5.55), Inches(0.65),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(7.4), Inches(top + 0.15), Inches(5.35), Inches(0.4),
                    action, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        top += 0.72
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '칸트의 정언명령 · 예수의 황금률 · 공자의 기소불욕(己所不欲)과 본질이 같다',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('XⅥ. 平天下')
def s_deok_bon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 平天下', page, total)
    add_title(slide, '德本財末 — 덕은 근본, 재물은 말단')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '德 者 本 也   財 者 末 也',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
                '덕자본야 · 재자말야',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(3.9), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '"덕(德)은 근본이요, 재물(財)은 말단이다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.9), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.7), [
        ('"仁者以財發身 不仁者以身發財"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"어진 자는 재물로 몸(사람)을 일으키고',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 어질지 못한 자는 몸(사람)을 희생해 재물을 일으킨다"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('XⅥ. 平天下')
def s_jaechui_minsan(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 平天下', page, total)
    add_title(slide, '財聚則民散 — 재물을 모으면 백성이 흩어진다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '財 聚 則 民 散',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.9),
                '財 散 則 民 聚',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                '재취즉민산 · 재산즉민취',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.95), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.6),
                '"재물을 모으면 백성이 흩어지고',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.6),
                ' 재물을 흩으면 백성이 모인다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '→ 부의 재분배와 사회 안정의 2,400년 전 통찰',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('XⅥ. 平天下')
def s_iro_ui(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 平天下', page, total)
    add_title(slide, '國不以利爲利 — 의(義)를 이(利)로 삼다')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.5), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.65), Inches(12.3), Inches(0.6),
                '國 不 以 利 爲 利   以 義 爲 利 也',
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.4),
                '국불이리위리 이의위리야',
                font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.25), Inches(11.7), Inches(1.5), [
        ('"나라는 이(利)로써 이(利)를 삼지 않고',
         {'font_size': 18, 'align': PP_ALIGN.CENTER}),
        (' 의(義)로써 이(利)를 삼는다"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_textbox(slide, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.5),
                '"장기 신뢰 자본이 단기 이익을 이긴다"',
                font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
                '엔론·위워크·페이스북의 사례가 모두 "재물을 먼저 구하고 덕을 말단으로 본" 결과',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== XⅦ. 명구절 ==============
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=42):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=22, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('XⅦ. 명구절 (1/8)',
    '大 學 之 道\n在 明 明 德   在 親 民   在 止 於 至 善',
    '대학지도 · 재명명덕 재친민 재지어지선',
    '대학의 도는 밝은 덕을 밝히고, 백성을 새롭게 하고, 지극한 선에 머무는 데 있다',
    '경 1장 (삼강령 선언)', hanmun_size=26), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (2/8)',
    '知 止 而 后 有 定',
    '지지이후유정',
    '머무를 곳을 안 뒤에 정해짐(定)이 있다 — 5단계 내면 심리학의 출발',
    '경 1장 (지의 효능)', hanmun_size=52), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (3/8)',
    '物 有 本 末   事 有 終 始',
    '물유본말 · 사유종시',
    '사물에는 근본과 말단이, 일에는 시작과 끝이 있다 — 본말의 선언',
    '경 1장', hanmun_size=42), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (4/8)',
    '自 天 子 以 至 於 庶 人\n壹 是 皆 以 修 身 爲 本',
    '자천자이지어서인 · 일시개이수신위본',
    '천자로부터 서민에 이르기까지 모두 한결같이 수신을 근본으로 삼는다',
    '경 1장 (결론)', hanmun_size=28), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (5/8)',
    '苟 日 新   日 日 新   又 日 新',
    '구일신 일일신 우일신',
    '진실로 하루가 새롭거든, 날로 새롭고, 또 날로 새로워라',
    '전 2장 (탕왕의 반명)', hanmun_size=38), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (6/8)',
    '所 謂 誠 其 意 者\n毋 自 欺 也',
    '소위성기의자 · 무자기야',
    '뜻을 성실히 한다는 것은 스스로를 속이지 않는 것이다',
    '전 6장 (성의)', hanmun_size=32), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (7/8)',
    '絜 矩 之 道',
    '혈구지도',
    '자기를 기준 삼아 타인을 헤아리는 도 — 6방향의 역지사지',
    '전 10장 (평천하)', hanmun_size=92), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (8/8)',
    '德 者 本 也   財 者 末 也',
    '덕자본야 · 재자말야',
    '덕은 근본이요, 재물은 말단이다 — 삶의 우선순위',
    '전 10장 (평천하)', hanmun_size=42), 'XⅦ. 명구절'))


# ============== XⅧ. 7대 메시지 ==============
@S('XⅧ. 7대 메시지')
def s_message_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 7대 메시지', page, total)
    add_title(slide, '7대 핵심 메시지 ① ~ ③')
    items = [
        ('1', '사람은 본래 빛나는 존재다',
         '명명덕의 전제 — 교육은 빈 그릇 채움이 아니라 가려진 빛의 드러냄'),
        ('2', '지행(知行)은 하나의 과정이다',
         '격물·치지(知) + 성의·정심·수신(行) = 끊을 수 없는 연속체'),
        ('3', '수신(修身)이 모든 것의 축이다',
         '"自天子以至於庶人 壹是皆以修身爲本" — 대학의 결론'),
    ]
    top = 2.5
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(1.3), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.42), Inches(1.0), Inches(0.5),
                    num, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.8), Inches(top), Inches(11.0), Inches(1.3), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.22), Inches(10.6), Inches(0.5),
                    title, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.0), Inches(top + 0.75), Inches(10.6), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 1.5


@S('XⅧ. 7대 메시지')
def s_message_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 7대 메시지', page, total)
    add_title(slide, '7대 핵심 메시지 ④ ~ ⑤')
    items = [
        ('4', '변화는 안에서 밖으로 간다',
         '팔조목 순서는 예외 없음 — 내면 정돈 없이 외적 성과 추구는 "근본 버리고 말단 구하기"'),
        ('5', '역지사지가 정치의 기본이다',
         '혈구지도는 경영·정치·교육·가정 모든 영역의 기본값 — 6방향 원칙'),
    ]
    top = 2.7
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(1.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(1.0), Inches(0.5),
                    num, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.8), Inches(top), Inches(11.0), Inches(1.5), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.3), Inches(10.6), Inches(0.5),
                    title, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.0), Inches(top + 0.9), Inches(10.6), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 1.7


@S('XⅧ. 7대 메시지')
def s_message_3(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 7대 메시지', page, total)
    add_title(slide, '7대 핵심 메시지 ⑥ ~ ⑦')
    items = [
        ('6', '재물은 말단이다 (德本財末)',
         '재물을 목적으로 삼으면 사람을 잃는다 — 덕이 있어야 재물이 따른다'),
        ('7', '"머물러야(止)" 얻는다',
         '"知止 → 定 → 靜 → 安 → 慮 → 得" — 멈춤이 얻음의 어머니'),
    ]
    top = 2.7
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(1.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(1.0), Inches(0.5),
                    num, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.8), Inches(top), Inches(11.0), Inches(1.5), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.3), Inches(10.6), Inches(0.5),
                    title, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.0), Inches(top + 0.9), Inches(10.6), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 1.7


# ============== XⅨ. 현대 의의 ==============
@S('XⅨ. 현대 의의')
def s_modern_integration(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅨ. 현대 의의', page, total)
    add_title(slide, '현대 ① — 통합된 삶의 지도 (Work-Life Integration)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('현대인의 조각난 삶',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('직장의 자아 · 가정의 자아 · SNS의 자아 · 친구들 사이의 자아',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('서로 연결되지 않은 여러 자아를 돌아가며 산다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('대학이 제시하는 반대의 그림',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('격물 → 치지 → 성의 → 정심 → 수신 → 제가 → 치국 → 평천하',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('내면과 외면, 공과 사, 일과 가정이 하나의 연속체',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('Work-Life Balance → Work-Life Integration의 동양적 원형',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('XⅨ. 현대 의의')
def s_modern_authenticity(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅨ. 현대 의의', page, total)
    add_title(slide, '현대 ② — 진정성(Authenticity)의 회복')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('현대사회의 가장 큰 병폐 — 자기기만(自欺)',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('SNS의 필터링된 이미지',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('성공의 대본을 따르는 행동',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('본심을 숨긴 사회생활',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('대학의 처방 — 성의(誠意)와 신독(愼獨)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"악을 싫어하되 나쁜 냄새처럼 하라"',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"감정과 행동이 일치하는 삶"의 선언',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('진정성(authenticity)은 현대 리더십의 핵심 가치로 재발견되고 있다',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('XⅨ. 현대 의의')
def s_modern_empathy(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅨ. 현대 의의', page, total)
    add_title(slide, '현대 ③ — 혈구지도(絜矩) · 공감 능력의 정밀화')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('혈구지도는 단순한 "착하게 살자"가 아니다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('구체적으로 "내가 윗사람에게 싫었던 것을 아랫사람에게 하지 말라"는 식으로',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('방향과 대상을 명시 — 가장 정밀한 공감의 기술',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('오늘의 적용',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 면접관은 자기가 면접자였을 때 싫었던 질문을 기억해야 한다',
         {'font_size': 14, 'space_before': 6}),
        ('• 부모는 자기가 자식이었을 때 싫었던 말투를 기억해야 한다',
         {'font_size': 14, 'space_before': 4}),
        ('• 리더는 자기가 부하였을 때 싫었던 지시 방식을 기억해야 한다',
         {'font_size': 14, 'space_before': 4}),
    ], line_spacing=1.3)


@S('XⅨ. 현대 의의')
def s_modern_stop(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅨ. 현대 의의', page, total)
    add_title(slide, '현대 ④ — 지지(知止) · 멈춤의 기술')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"어디에서 멈출 것인가"를 모르는 채 질주하는 시대',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 스마트폰의 무한 스크롤 → 눈의 "지(止)" 없음',
         {'font_size': 14, 'space_before': 6}),
        ('• 24시간 접속 문화 → 마음의 "지(止)" 없음',
         {'font_size': 14, 'space_before': 4}),
        ('• 끝없는 성장 신화 → 삶의 "지(止)" 없음',
         {'font_size': 14, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('지지(知止)는 무기력한 포기가 아니다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"지극한 선"이라는 목표를 알기에 가능한 능동적 머묾',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('= 현대 마음챙김(Mindfulness)·Flow 이론의 동양적 원형',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('XⅨ. 현대 의의')
def s_modern_esg(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅨ. 현대 의의', page, total)
    add_title(slide, '현대 ⑤ — 덕본재말(德本財末)과 ESG')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"덕은 근본, 재물은 말단"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('덕이 있는 사람에게 재물이 따를 수는 있으나',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('재물을 뒤쫓다가 덕을 잃으면 결국 재물까지 잃는다',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 경영 — 덕본재말의 재발견',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• "장기 신뢰 자본이 단기 이익을 이긴다"',
         {'font_size': 14, 'space_before': 6}),
        ('• ESG(환경·사회·지배구조) 경영의 동양적 원조',
         {'font_size': 14, 'space_before': 4}),
        ('• 페이스북·엔론·위워크의 실패 = "재물을 먼저 구하고 덕을 말단으로"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== XX. 비교 ==============
@S('XX. 비교')
def s_compare_saseo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XX. 비교', page, total)
    add_title(slide, '사서(四書) 비교 — 대학의 위치')
    rows = [
        ('대학(大學)',  '規模',  '약 1,700자',  '체계의 뼈대 — 입덕지문',  True),
        ('논어(論語)',  '根本',  '약 16,000자', '인(仁)의 단편적 가르침',  False),
        ('맹자(孟子)',  '發揮',  '약 35,000자', '왕성한 논변',             False),
        ('중용(中庸)',  '微妙',  '약 3,500자',  '형이상학·중도',           False),
    ]
    top = 2.2
    add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(0.55), SUB)
    add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(2.5), Inches(0.4),
                '서명', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(3.05), Inches(top), Inches(2.5), Inches(0.55), SUB)
    add_textbox(slide, Inches(3.05), Inches(top + 0.1), Inches(2.5), Inches(0.4),
                '주희의 평', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(5.6), Inches(top), Inches(2.5), Inches(0.55), SUB)
    add_textbox(slide, Inches(5.6), Inches(top + 0.1), Inches(2.5), Inches(0.4),
                '분량', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(8.15), Inches(top), Inches(4.7), Inches(0.55), SUB)
    add_textbox(slide, Inches(8.15), Inches(top + 0.1), Inches(4.7), Inches(0.4),
                '특징', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    row_h = 0.85
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        c = ACCENT if row[4] else INK
        bg = RGBColor(0xFA, 0xE5, 0xE5) if row[4] else (
            RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE)
        add_filled_rect(slide, Inches(0.5), Inches(y), Inches(2.5), Inches(row_h), PALE)
        add_textbox(slide, Inches(0.55), Inches(y + 0.22), Inches(2.4), Inches(0.4),
                    row[0], font_size=15, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.05), Inches(y), Inches(2.5), Inches(row_h), bg)
        add_textbox(slide, Inches(3.05), Inches(y + 0.22), Inches(2.5), Inches(0.4),
                    row[1], font_size=15, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.6), Inches(y), Inches(2.5), Inches(row_h), bg)
        add_textbox(slide, Inches(5.6), Inches(y + 0.22), Inches(2.5), Inches(0.4),
                    row[2], font_size=14, color=c, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(8.15), Inches(y), Inches(4.7), Inches(row_h), bg)
        add_textbox(slide, Inches(8.3), Inches(y + 0.22), Inches(4.5), Inches(0.4),
                    row[3], font_size=13, color=c)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '대학이 사서 중 가장 짧지만 전체의 틀(規模)을 잡는 입문서',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('XX. 비교')
def s_compare_modern(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XX. 비교', page, total)
    add_title(slide, '대학과 현대 자기계발 도서')
    rows = [
        ('스티븐 코비',  '성공하는 사람들의 7가지 습관',
         '사적 영역(1~3) → 공적 영역(4~6) → 쇄신(7)',
         '= 팔조목의 안→밖 구조'),
        ('짐 콜린스',    '좋은 기업을 넘어 위대한 기업으로',
         '"5단계 리더십" — 겸손(謙) + 직업적 의지(志)',
         '= 성의·정심·수신의 현대판'),
        ('다니엘 골먼',  '감성지능(EQ)',
         '자기 인식·자기 관리·공감·관계 관리',
         '= 정심·수신·혈구·제가의 재명명'),
        ('칼 로저스',    '인간중심 상담',
         '"진정성(congruence)·무조건적 긍정"',
         '= 성의·명명덕의 심리학'),
    ]
    top = 2.3
    for author, book, principle, link in rows:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(1.05), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.12), Inches(2.5), Inches(0.4),
                    author, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(2.5), Inches(0.5),
                    book, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(9.6), Inches(1.05), PALE)
        add_textbox(slide, Inches(3.35), Inches(top + 0.12), Inches(9.3), Inches(0.45),
                    principle, font_size=13, bold=True, color=INK)
        add_textbox(slide, Inches(3.35), Inches(top + 0.55), Inches(9.3), Inches(0.45),
                    link, font_size=13, color=ACCENT, bold=True)
        top += 1.18


# ============== XXI. 마무리 ==============
@S('XXI. 마무리')
def s_one_page(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXI. 마무리', page, total)
    add_title(slide, '한 장으로 보는 대학')
    add_filled_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.15), Inches(12.3), Inches(0.5),
                '삼강령 (三綱領) = 목표',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.65), Inches(12.3), Inches(0.7),
                '明明德 → 親民(新民) → 止於至善',
                font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.5), RGBColor(0xA0, 0x40, 0x40))
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.5),
                '팔조목 (八條目) = 방법',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.7),
                '格物 → 致知 → 誠意 → 正心 → 修身 → 齊家 → 治國 → 平天下',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.4), [
        ('결론 한 문장',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"自天子以至於庶人 壹是皆以修身爲本"',
         {'font_size': 20, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"천자로부터 서민까지 모두 수신이 근본"',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 2}),
    ], line_spacing=1.25)


@S('XXI. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXI. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 대학')
    add_filled_rect(slide, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.9), PALE)
    add_paragraphs(slide, Inches(1.1), Inches(2.3), Inches(11.1), Inches(4.6), [
        ('"당신의 덕은 이미 당신 안에 있다(明明德)',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        (' 그것을 가린 먼지를 매일 닦고(修身)',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        (' 당신이 선 자리에서 지극한 선에 머물러라(止於至善)',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        (' 그러면 당신의 변화가 가족과 조직과 세상을 새롭게 할 것이다(新民)',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        (' 그 순서를 거스르지 말라"',
         {'font_size': 21, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('', {'font_size': 8}),
        ('— 2,400년 전의 책이지만 항상 현재형으로 답한다 —',
         {'font_size': 16, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
    ], line_spacing=1.2)


@S('XXI. 마무리')
def s_meaning(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXI. 마무리', page, total)
    add_title(slide, '대학이 당신에게 건네는 한 문장')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.6),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(4.3), [
        ('대학은',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('"지금 내가 어디에 서 있는가"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('를 물을 때마다',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('항상 현재형으로 답한다',
         {'font_size': 19, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 8}),
        ('어디로 가는지가 아니라',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 14}),
        ('지금 어디에 머무를지를 묻는 책',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('— 멈춤(止)이 곧 얻음(得)의 어머니 —',
         {'font_size': 18, 'bold': True, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 14}),
    ], line_spacing=1.25)


@S('XXI. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.5),
                '止 於 至 善',
                font_size=120, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '지 어 지 선', font_size=24, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.7),
                '"지극한 선에 머물러라"',
                font_size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '당신이 선 자리에서 — 군주는 仁에, 부모는 慈에, 벗은 信에',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
                '— 경 1장 · 대학의 첫 선언과 마지막 결론',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
                '감사합니다', font_size=26, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\대학_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')