# -*- coding: utf-8 -*-
"""
정관정요 발표자료 — 망라적 95장 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 1. 표지 ==============
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8),
                '貞 觀 政 要', font_size=92, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
                'Essentials of Government in the Zhenguan Era · 정관정요',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.3), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.55), Inches(12.3), Inches(0.5),
                '오긍(吳兢) 편 — 당 태종 정관 23년의 정치 문답집',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.4),
                '8세기 초 편찬 · 10권 40편 약 460장 · 동아시아 제왕학(帝王學)의 절대 표준',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.5),
                '"君者舟也, 庶人者水也. 水則載舟, 水則覆舟"',
                font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
                '— 임금은 배요, 백성은 물 · 물은 배를 띄우기도 뒤집기도 한다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 2. 목차 ==============
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.8), Inches(0.7),
                '목 차', font_size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.15), Inches(12.8))
    items = [
        [('Ⅰ', '개요 — 정관정요란 무엇인가'),
         ('Ⅱ', '시대 배경 — 수의 멸망과 현무문'),
         ('Ⅲ', '인물 — 당 태종과 정관의 군신'),
         ('Ⅳ', '구성 — 10권 40편'),
         ('Ⅴ', '군도(君道) — 군주의 도리'),
         ('Ⅵ', '정체(政體) — 정치의 근본'),
         ('Ⅶ', '간쟁(諫諍) — 간언의 정치'),
         ('Ⅷ', '三鏡 — 세 거울'),
         ('Ⅸ', '군신(君臣)의 도'),
         ('Ⅹ', '임현(任賢) — 인재 등용'),
         ('Ⅺ', '민본(民本) — 백성이 근본')],
        [('Ⅻ', '권5 — 5대 핵심 덕목'),
         ('ⅩⅢ', '권6 — 9대 자기 점검'),
         ('ⅩⅣ', '문화·민생 정책'),
         ('ⅩⅤ', '군사·외교'),
         ('ⅩⅥ', '신종(愼終) — 끝을 삼가라'),
         ('ⅩⅦ', '명구절 10선'),
         ('ⅩⅧ', '후대의 수용'),
         ('ⅩⅨ', '다른 고전과의 비교'),
         ('ⅩⅩ', '마무리')],
    ]
    for col, group in enumerate(items):
        x = 0.7 + col * 6.4
        top = 1.5
        for num, title in group:
            add_textbox(slide, Inches(x), Inches(top), Inches(0.9), Inches(0.4),
                        num, font_size=15, bold=True, color=ACCENT)
            add_textbox(slide, Inches(x + 0.9), Inches(top), Inches(5.4), Inches(0.4),
                        title, font_size=15, color=INK)
            top += 0.5


# ============== Ⅰ. 개요 ==============
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '정관정요(貞觀政要)란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                '당 태종(唐太宗) 이세민의 정관(貞觀) 연간 정치를 정리한 책',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.5),
                '"정관(貞觀)"은 태종의 연호 (627~649) — 곧 "정관 시대의 정치 요체"',
                font_size=15, color=ACCENT, align=PP_ALIGN.CENTER)
    nums = [('10', '권(卷)'), ('40', '편(篇)'), ('약 460', '장(章)'), ('1,300', '년 영향')]
    for i, (n, lbl) in enumerate(nums):
        x = 0.6 + i * 3.05
        add_textbox(slide, Inches(x), Inches(4.0), Inches(2.9), Inches(1.0),
                    n, font_size=50, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(5.15), Inches(2.9), Inches(0.5),
                    lbl, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                '동아시아 제왕학(帝王學)의 절대 표준 — 한국·중국·일본·베트남 군주의 필독서',
                font_size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '편년체(編年體)가 아니라 주제별 — 현대 경영 핸드북의 원조 구조',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_status(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '책의 위상 — 핵심 정보')
    rows = [
        ('서명',       '정관정요(貞觀政要)',         '"정관 시대의 정치 요체"'),
        ('정관(貞觀)', '당 태종 이세민의 연호',       '627~649년, 23년'),
        ('편자',       '오긍(吳兢, 670?~749)',       '당 현종(玄宗) 때의 사관(史官)'),
        ('편찬 시기',  '8세기 초',                    '태종 사후 약 60~70년'),
        ('분량',       '10권 40편',                   '약 460장의 일화·문답'),
        ('성격',       '동아시아 제왕학의 절대 표준',  '주제별 편찬 → 사례집 형식'),
        ('영향',       '1,300년 군주의 필독서',       '한국·중국·일본·베트남'),
    ]
    top = 2.15
    for i, (tag, val, note) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.65), PALE)
        add_filled_rect(slide, Inches(2.85), Inches(top), Inches(10.0), Inches(0.65), bg)
        add_textbox(slide, Inches(0.55), Inches(top + 0.18), Inches(2.2), Inches(0.4),
                    tag, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.95), Inches(top + 0.05), Inches(4.5), Inches(0.5),
                    val, font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.5), Inches(top + 0.08), Inches(5.3), Inches(0.5),
                    note, font_size=13, color=SUB)
        top += 0.7


@S('Ⅰ. 개요')
def s_4_uniqueness(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '왜 동양 제왕학의 최고봉인가 — 네 가지 독특성')
    items = [
        ('1', '실재 명군의 기록',  '요·순·우 같은 신화적 성왕이 아니라\n실제로 재위한 당 태종(627~649)의 통치 그대로'),
        ('2', '실패도 함께 기록',  '고구려 원정 실패·태자 폐위·후반 흔들림까지\n숨기지 않는 "흔들리며 배운 군주"의 드라마'),
        ('3', '쿠데타 군주의 속죄', '현무문의 변으로 형제를 죽이고 즉위 →\n수 양제의 멸망을 반면교사로 삼은 자기 교정'),
        ('4', '주제별 사례집',     '편년이 아닌 덕목·정책·병폐별 분류 →\n오늘의 경영 핸드북과 같은 실용 구조'),
    ]
    top = 2.2
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.9), Inches(1.05), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.3), Inches(0.9), Inches(0.5),
                    num, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(2.7), Inches(1.05), PALE)
        add_textbox(slide, Inches(1.7), Inches(top + 0.3), Inches(2.7), Inches(0.5),
                    title, font_size=17, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.6), Inches(top), Inches(8.2), Inches(1.05),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.8), Inches(top + 0.13), Inches(8.0), Inches(0.85),
                    desc, font_size=14, color=INK)
        top += 1.2


@S('Ⅰ. 개요')
def s_ogyung(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '편자 오긍(吳兢, 670?~749)', '— 당 현종 때의 사관(史官)')
    lines = [
        ('신분', {'bold': True, 'font_size': 18, 'color': ACCENT}),
        ('  측천무후 시기부터 활동한 사관 — 황실 일거수일투족의 기록자',
         {'font_size': 16}),
        ('작업', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('  정관 23년의 실록·기거주(起居注)·군신 대화록을 오랜 세월에 걸쳐',
         {'font_size': 16}),
        ('  편년이 아닌 "주제별"로 재분류·편집',
         {'font_size': 16, 'bold': True}),
        ('편찬 목적', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('  당 현종(재위 712~756)을 위시한 후대 군주가',
         {'font_size': 16}),
        ('  "정관의 치세"를 본받게 하려는 진헌(進獻)의 책',
         {'font_size': 16}),
        ('역사적 평가', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('  사관으로서의 정직 — 태종의 실패까지 그대로 기록',
         {'font_size': 16}),
    ]
    add_paragraphs(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.5),
                   lines, line_spacing=1.4)


@S('Ⅰ. 개요')
def s_irony(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '편찬 목적과 역설(逆說)')
    add_filled_rect(slide, Inches(0.6), Inches(2.2), Inches(11.9), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.33), Inches(11.9), Inches(0.5),
                '오긍의 의도 — "후대 군주에게 정관의 모범을 본받게 하자"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.2), Inches(11.9), Inches(1.6), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(3.4), Inches(11.3), Inches(1.4), [
        ('그러나 이 책을 가장 먼저 받은 현종(玄宗) 자신이',
         {'font_size': 18, 'align': PP_ALIGN.CENTER}),
        ('후반기에 초심을 잃고 안사(安史)의 난(755~763)을 부르는 비극을 자초',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.4)
    add_rule(slide, Inches(5.5), Inches(5.1), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
                '"끝을 삼가라(愼終如始)"는 책의 가장 큰 교훈을',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
                '책을 가장 먼저 받은 군주가 가장 먼저 어긴 역사의 아이러니',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                '→ 이 역설이 오히려 "신종여시(愼終如始)" 메시지의 설득력을 강화',
                font_size=13, color=SUB, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 시대 배경 ==============
@S('Ⅱ. 시대 배경')
def s_sui_fall(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대 배경', page, total)
    add_title(slide, '수(隋)의 멸망 — 태종의 평생 반면교사',
              '정관정요에서 가장 자주 언급되는 역사적 교훈')
    rows = [
        ('581~604', '문제(文帝) 양견(楊堅)', '검소·근면, 명군',     '중국 통일, 30년 안정'),
        ('604~618', '양제(煬帝) 양광(楊廣)', '사치·방종, 토목광',   '14년 만에 멸망 — 호위군에게 살해'),
    ]
    top = 2.5
    for era, ruler, char, result in rows:
        bg = PALE if '문제' in ruler else RGBColor(0xFA, 0xE5, 0xE5)
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(11.9), Inches(1.6), bg)
        add_textbox(slide, Inches(0.8), Inches(top + 0.2), Inches(2.5), Inches(0.5),
                    era, font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(3.4), Inches(top + 0.2), Inches(4.0), Inches(0.5),
                    ruler, font_size=20, bold=True, color=INK)
        add_textbox(slide, Inches(0.8), Inches(top + 0.8), Inches(5.0), Inches(0.4),
                    f'• {char}', font_size=15, color=SUB)
        add_textbox(slide, Inches(6.5), Inches(top + 0.8), Inches(6.2), Inches(0.4),
                    f'→ {result}', font_size=15, color=INK, bold=True)
        top += 1.8
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '정관정요의 숨은 구조 — "수서(隋書)에 대한 반서(反書)"',
                font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '정관 시대의 안정은 수 양제의 반대를 모든 면에서 실행한 결과',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅱ. 시대 배경')
def s_yangdi_4_evils(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대 배경', page, total)
    add_title(slide, '수 양제의 4대 악정(惡政)', '— 태종이 반복적으로 환기한 반면교사')
    evils = [
        ('1', '토목 광기',     '낙양 동도 건설(5년·수십만 명)\n대운하 건설(수백만 명·기아·사망)\n강도(江都) 호화 행궁'),
        ('2', '고구려 원정',   '612·613·614 세 차례\n113만 대군 동원(당시 세계 최대)\n모두 실패, 막대한 인명 손실'),
        ('3', '간언 무시',     '충신 유의(劉毅) 처형\n간신 우세기(虞世基)만 곁에 둠\n비위 맞춤의 정치'),
        ('4', '멸망',          '백성의 기근·전쟁 피해 폭발\n전국 반란 발발\n618년 강도에서 호위군에 살해'),
    ]
    top = 2.3
    box_w = 3.0
    gap = 0.13
    start_x = (13.333 - (box_w * 4 + gap * 3)) / 2
    for i, (num, title, desc) in enumerate(evils):
        x = start_x + i * (box_w + gap)
        add_filled_rect(slide, Inches(x), Inches(top), Inches(box_w), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(x), Inches(top + 0.13), Inches(box_w), Inches(0.5),
                    f'{num}. {title}', font_size=17, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x), Inches(top + 0.75), Inches(box_w), Inches(3.0), PALE)
        add_textbox(slide, Inches(x + 0.15), Inches(top + 0.95), Inches(box_w - 0.3),
                    Inches(2.7), desc, font_size=13, color=INK, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '→ 정관정요의 검약·구간·정벌·안변 편이 모두 이 4가지에 대한 응답',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 시대 배경')
def s_xuanwumen(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대 배경', page, total)
    add_title(slide, '현무문의 변(玄武門之變, 626년)', '— 태종 즉위의 원죄(原罪)')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(2.0), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.35), Inches(12.3), Inches(0.5),
                '권력 투쟁의 구도', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    figures = [
        ('당 고조 이연(李淵)', '건국자, 결단력 약함'),
        ('태자 이건성(李建成)', '첫째 — 위징을 모사로'),
        ('이세민(李世民)',     '둘째 — 방현령·두여회를 모사로'),
        ('제왕 이원길(李元吉)', '셋째 — 이건성 편'),
    ]
    for i, (name, role) in enumerate(figures):
        x = 0.6 + (i % 2) * 6.1
        y = 3.0 + (i // 2) * 0.55
        add_textbox(slide, Inches(x), Inches(y), Inches(3.5), Inches(0.4),
                    name, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(x + 3.5), Inches(y), Inches(2.5), Inches(0.4),
                    role, font_size=12, color=SUB)
    add_filled_rect(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.5), RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.55), Inches(11.7), Inches(2.3), [
        ('626년 7월 2일 새벽 — 현무문(玄武門)',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('이세민이 선제 공격하여 형 이건성과 동생 이원길을 살해',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('형제 2명의 자녀 10명도 모두 처형 — 후환 방지',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('고조는 압박에 못 이겨 양위 → 626년 8월 즉위, 정관 원년 시작',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('태종의 평생 죄의식 — "나는 형을 죽이고 즉위한 자다',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('그러니 나는 반드시 최고의 치세를 만들어야 한다"',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅱ. 시대 배경')
def s_zhenguan_age(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대 배경', page, total)
    add_title(slide, '정관의 치(貞觀之治, 627~649)',
              '— 역사상 가장 완성된 23년')
    rows = [
        ('민생', '"도불습유·야불폐호" — 길에 떨어진 물건을 줍지 않고 문을 잠그지 않는다'),
        ('법제', '당률(唐律) 편찬 — 이후 1,000년 동아시아 법제의 표준'),
        ('외교', '돌궐 평정(630), 서역 확대, 문성공주의 토번 화친'),
        ('문화', '『오경정의』·『진서』·『수서』 등 국가 편찬 사업'),
        ('제도', '삼성육부제 완성, 과거제 확대, 조·용·조 세제 정비'),
        ('영토', '동쪽 고구려부터 서쪽 페르시아 국경까지 — "천가한(天可汗)" 칭호'),
        ('인구', '수말 1,000만 이하 → 정관 말 약 3,000만 회복'),
    ]
    top = 2.25
    for i, (tag, desc) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(1.8), Inches(0.62), PALE)
        add_filled_rect(slide, Inches(2.35), Inches(top), Inches(10.5), Inches(0.62), bg)
        add_textbox(slide, Inches(0.5), Inches(top + 0.17), Inches(1.8), Inches(0.4),
                    tag, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.5), Inches(top + 0.17), Inches(10.3), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.67


@S('Ⅱ. 시대 배경')
def s_late_drift(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 시대 배경', page, total)
    add_title(slide, '정관 후반의 흔들림 — "신종(愼終)"의 배경')
    timeline = [
        ('정관 11년경 (637)',  '점차 사치 경향 시작'),
        ('정관 13년 (639)',    '궁궐 증축·대규모 행차 증가'),
        ('정관 13년경',         '위징의 「신종소(愼終疏)」 진언 — 십점불극종(十漸不克終)'),
        ('정관 17년 (643)',    '위징 사망 → 직간자(直諫者) 부재'),
        ('정관 19년 (645)',    '고구려 원정 실패'),
        ('정관 23년 (649)',    '태종 붕어, 향년 52세 — 만년에 「제범(帝範)」 12편 친저'),
    ]
    top = 2.3
    for era, event in timeline:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.3), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.17), Inches(3.3), Inches(0.4),
                    era, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.05), Inches(top), Inches(8.8), Inches(0.65), PALE)
        add_textbox(slide, Inches(4.25), Inches(top + 0.17), Inches(8.6), Inches(0.4),
                    event, font_size=15, color=INK)
        top += 0.72
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"명군도 후반에 흔들린다" — 정관정요의 마지막 장이 신종편인 까닭',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅲ. 인물 ==============
@S('Ⅲ. 인물')
def s_taizong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '당 태종 이세민(唐太宗 李世民, 598~649)',
              '— 정관 정치의 주인공')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '太\n宗', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('연보', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  598년 출생 — 618년 당 건국 참여 (20세 대장군)',
         {'font_size': 14}),
        ('  626년 현무문의 변 → 28세 즉위',
         {'font_size': 14}),
        ('  627~649년 정관 23년 → 649년 붕어, 향년 52세',
         {'font_size': 14}),
        ('성격', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  문무겸비 — 초원의 전사이면서 한학에 밝음',
         {'font_size': 14}),
        ('  자기 인식이 강함 — 약점을 알고 끊임없이 점검',
         {'font_size': 14}),
        ('  감정 기복 있음 — 위징에게 분노, 그러나 결국 수용',
         {'font_size': 14}),
        ('  실패에 솔직 — 고구려 원정 실패·후반 흔들림 인정',
         {'font_size': 14}),
        ('  인재 발굴의 천재 — 적도 쓰고, 능력자를 과감히 기용',
         {'font_size': 14}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅲ. 인물')
def s_weizheng(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '위징(魏徵, 580~643)', '— 직간(直諫)의 대명사, 태종의 거울')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '魏\n徵', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('연보', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  수말 혼란기 여러 세력 전전 → 태자 이건성의 모사',
         {'font_size': 14}),
        ('  626년 현무문 후 포로 → 태종이 등용',
         {'font_size': 14}),
        ('  정관 17년간 간의대부·비서감·시중 역임',
         {'font_size': 14}),
        ('  643년 사망, 태종이 친히 묘지명을 씀',
         {'font_size': 14}),
        ('업적', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  평생 200여 차례 직간(直諫)',
         {'font_size': 14, 'bold': True}),
        ('  「간태종십사소」·「신종소」 등 명문 상소',
         {'font_size': 14}),
        ('  『정관예(貞觀禮)』·『수서(隋書)』 편찬 주도',
         {'font_size': 14}),
        ('리더십 자질', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  강직·냉정·지혜·헌신 — 군주를 위해 목숨 걸고 간언',
         {'font_size': 14}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅲ. 인물')
def s_empress(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '황후 장손씨(長孫皇后, 601~636)',
              '— 정관 정치의 숨은 기둥')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '長孫\n皇后', font_size=64, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('역할', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  총명·인자·겸손 — 태종의 감정 조절자',
         {'font_size': 14}),
        ('  직접 정치 개입은 피하되, 결정적 순간에 개입',
         {'font_size': 14}),
        ('대표 일화', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  ① 위징 보호 — "임금이 현명해야 신하가 곧다(君明臣直)"',
         {'font_size': 14}),
        ('  ② 사치 거절 — "한 나라의 어머니가 사치하면 모든 부녀자가 따른다"',
         {'font_size': 14}),
        ('  ③ 임종 유언 — "제 친정(장손씨 가문)을 결코 요직에 들이지 마소서"',
         {'font_size': 14}),
        ('의의', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  한(漢)대 외척 정치의 비극을 반면교사 삼은 자기 절제',
         {'font_size': 14}),
        ('  636년 36세에 별세 — 태종의 후반 흔들림이 시작된 분기점',
         {'font_size': 14, 'color': SUB}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅲ. 인물')
def s_fang(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '방현령(房玄齡, 579~648)', '— 행정·기획의 달인 · 정관 조정의 축')
    lines = [
        ('역할', {'bold': True, 'font_size': 20, 'color': ACCENT}),
        ('  "꾀(謀)"의 대가 — 여러 선택지를 능숙하게 제시',
         {'font_size': 17}),
        ('  창업기부터 태종과 함께한 최고 행정가',
         {'font_size': 17}),
        ('  재상 자리에 20여 년 — 정관 조정의 축',
         {'font_size': 17}),
        ('', {'font_size': 8}),
        ('창업 논의 발언 (정체편)', {'bold': True, 'font_size': 20, 'color': ACCENT, 'space_before': 10}),
        ('"창업이 어렵습니다. 천하가 어지러울 때 여러 영웅이 다투니',
         {'font_size': 17}),
        (' 이를 제압하는 일이 쉽지 않습니다."',
         {'font_size': 17, 'bold': True}),
        ('  → 창업 경험자의 현실적 증언',
         {'font_size': 14, 'color': SUB}),
    ]
    add_paragraphs(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.5),
                   lines, line_spacing=1.4)


@S('Ⅲ. 인물')
def s_du(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '두여회(杜如晦, 585~630)', '— 결단(斷)의 달인')
    lines = [
        ('특징', {'bold': True, 'font_size': 20, 'color': ACCENT}),
        ('  방현령이 짠 꾀(謀) 중 최선을 즉시 선택',
         {'font_size': 17}),
        ('  결단력의 화신 — 우유부단함의 정반대',
         {'font_size': 17}),
        ('  태종이 가장 의지한 참모 중 하나',
         {'font_size': 17}),
        ('', {'font_size': 8}),
        ('비극', {'bold': True, 'font_size': 20, 'color': ACCENT, 'space_before': 10}),
        ('  정관 4년(630) 젊은 나이로 요절',
         {'font_size': 17}),
        ('  태종이 매우 아쉬워함 — 만년에도 자주 그를 회상',
         {'font_size': 17}),
        ('', {'font_size': 8}),
        ('태종의 인물평', {'bold': True, 'font_size': 20, 'color': ACCENT, 'space_before': 10}),
        ('  "방(房)은 꾀를 짜고, 두(杜)는 결단한다 — 房謀杜斷"',
         {'font_size': 17, 'bold': True, 'color': INK}),
    ]
    add_paragraphs(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.5),
                   lines, line_spacing=1.4)


@S('Ⅲ. 인물')
def s_fangmo_dudan(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '방모두단(房謀杜斷) — 강점의 조합')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '房 謀 杜 斷',
                font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '방모두단 — "방은 꾀를 짜고, 두는 결단한다"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.1), Inches(2.3), color=RULE, weight=1.5)
    # 좌우 분할
    add_filled_rect(slide, Inches(0.7), Inches(4.4), Inches(5.9), Inches(2.6), PALE)
    add_textbox(slide, Inches(0.7), Inches(4.55), Inches(5.9), Inches(0.5),
                '房玄齡 — 謀 (꾀)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(5.2), Inches(5.5), Inches(1.7),
                '"문제 해결 5가지 방법"을 제시\n선택지를 풍부하게 만들어내는 자',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(6.8), Inches(4.4), Inches(5.9), Inches(2.6), PALE)
    add_textbox(slide, Inches(6.8), Inches(4.55), Inches(5.9), Inches(0.5),
                '杜如晦 — 斷 (결단)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.0), Inches(5.2), Inches(5.5), Inches(1.7),
                '"그 중 최선 1가지"를 골라 즉시 실행\n결단하여 실행하는 자',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '강점이 다른 두 사람의 조합 — 현대 팀빌딩 이론의 원형',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅲ. 인물')
def s_wanggui(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '왕규(王珪) — 인물 평가의 대가',
              '간의대부(諫議大夫) — 자기 인식이 인재 안목')
    add_filled_rect(slide, Inches(0.6), Inches(2.2), Inches(11.9), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.33), Inches(11.9), Inches(0.5),
                '태종: "신하들의 장단(長短)을 평해 보라"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rows = [
        ('방현령', '큰일을 맡길 만함'),
        ('두여회', '결단을 맡길 만함'),
        ('위 징', '잘못을 바로잡아 드릴 만함'),
        ('자기 자신', '옳고 그름을 가려 말하는 데 약간의 재주가 있을 뿐'),
    ]
    top = 3.2
    for name, desc in rows:
        is_self = '자기' in name
        bg_color = RGBColor(0xFA, 0xE5, 0xE5) if is_self else PALE
        add_filled_rect(slide, Inches(0.7), Inches(top), Inches(2.5), Inches(0.7), bg_color)
        add_textbox(slide, Inches(0.7), Inches(top + 0.2), Inches(2.5), Inches(0.4),
                    name, font_size=17, bold=True,
                    color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.3), Inches(top), Inches(9.5), Inches(0.7),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.5), Inches(top + 0.2), Inches(9.1), Inches(0.4),
                    desc, font_size=16, color=INK)
        top += 0.85
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '자기 자리를 정확히 아는 자만이 남의 자리도 정확히 본다',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅲ. 인물')
def s_li_yu(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '이정(李靖) · 우세남(虞世南) — 군사와 문예')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '이정(李靖, 571~649)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(3.0), Inches(5.9), Inches(0.4),
                '— 최고의 군사 전략가', font_size=14, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.6), Inches(5.3), Inches(3.0), [
        ('• 돌궐 평정(630) — 힐리가한 생포', {'font_size': 15}),
        ('  돌궐 와해 → 태종 "천가한" 칭호', {'font_size': 14, 'color': SUB}),
        ('', {'font_size': 6}),
        ('• 토욕혼 정벌(635)', {'font_size': 15, 'space_before': 6}),
        ('• 토번(吐蕃) 외교의 군사적 기반', {'font_size': 15}),
        ('', {'font_size': 6}),
        ('• 말수 적고 권력 다툼에 무관여', {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ('• 정관 후기 조용히 은퇴', {'font_size': 14, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '우세남(虞世南, 558~638)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.0), Inches(5.9), Inches(0.4),
                '— 학자·서예가', font_size=14, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.6), Inches(5.3), Inches(3.0), [
        ('• 태종의 문학·예악 자문', {'font_size': 15}),
        ('• 당대 최고 서예가 중 하나', {'font_size': 15}),
        ('', {'font_size': 6}),
        ('태종의 "오절(五絶)" 칭송', {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('  덕행 · 충직 · 박학 · 문장 · 서한(書翰)',
         {'font_size': 14, 'color': INK}),
        ('', {'font_size': 6}),
        ('• 사후 태종이 직접 추도', {'font_size': 14, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)


@S('Ⅲ. 인물')
def s_minister_overview(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 인물', page, total)
    add_title(slide, '정관의 핵심 군신 — 한눈에 보기')
    rows = [
        ('당 태종', '598~649', '군주',     '문무겸비 · 자기 점검의 화신'),
        ('위징',    '580~643', '간의대부', '200여 회 직간 · "거울 같은 신하"'),
        ('장손황후','601~636', '황후',     '감정 조절 동반자 · 외척 절제'),
        ('방현령',  '579~648', '재상',     '꾀(謀)의 대가 · 20여 년 행정 총괄'),
        ('두여회',  '585~630', '재상',     '결단(斷)의 대가 · 정관 4년 요절'),
        ('왕규',    '?~639',   '간의대부', '인물 평가의 대가'),
        ('이정',    '571~649', '명장',     '돌궐 평정 · 천가한의 군사 기반'),
        ('우세남',  '558~638', '학자',     '"오절(五絶)" · 문예 자문'),
        ('장손무기','594?~659','재상',     '황후의 오빠 · 태종 처남이자 동지'),
    ]
    top = 2.0
    for i, (name, era, role, char) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(12.3), Inches(0.55), bg)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(2.0), Inches(0.4),
                    name, font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.7), Inches(top + 0.13), Inches(1.8), Inches(0.4),
                    era, font_size=13, color=SUB)
        add_textbox(slide, Inches(4.5), Inches(top + 0.13), Inches(2.0), Inches(0.4),
                    role, font_size=13, color=INK)
        add_textbox(slide, Inches(6.6), Inches(top + 0.13), Inches(6.2), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.58


# ============== Ⅳ. 구성 ==============
@S('Ⅳ. 구성')
def s_structure(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 구성', page, total)
    add_title(slide, '10권 40편 개관 — 주제별 편찬의 완성')
    rows = [
        ('권1', '1~2편',   '군주의 근본 (군도·정체)'),
        ('권2', '3~5편',   '인재와 간언 (임현·구간·납간)'),
        ('권3', '6~8편',   '군신·관제 (군신감계·택관·봉건)'),
        ('권4', '9~12편',  '후계자 교육 4편'),
        ('권5', '13~17편', '5대 핵심 덕목 (인의·충의·효우·공평·성신)'),
        ('권6', '18~26편', '9대 자기 점검 병폐'),
        ('권7', '27~29편', '문화 정책 (숭유학·문사·예악)'),
        ('권8', '30~34편', '민생·법·경제 (무농·형법·사령·공부·변흥망)'),
        ('권9', '35~36편', '군사·외교 (정벌·안변)'),
        ('권10','37~40편', '군주의 일상과 마무리 (행행·전렵·재상·신종)'),
    ]
    top = 2.0
    for kwon, scope, theme in rows:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.2), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(1.2), Inches(0.4),
                    kwon, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.9), Inches(top), Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(1.9), Inches(top + 0.1), Inches(2.0), Inches(0.4),
                    scope, font_size=14, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(4.1), Inches(top + 0.1), Inches(8.7), Inches(0.4),
                    theme, font_size=14, color=INK)
        top += 0.55


@S('Ⅳ. 구성')
def s_circle(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 구성', page, total)
    add_title(slide, '동심원 구조 — "자기"에서 출발해 "자기"로 귀결')
    layers = [
        ('1~2',   '군주의 근본 — 자기'),
        ('3~5',   '인재와 간언 — 사람'),
        ('6~8',   '군신·관제 — 구조'),
        ('9~12',  '후계 — 계승'),
        ('13~17', '5대 덕목 — 덕'),
        ('18~26', '9대 병폐 — 자기 점검'),
        ('27~29', '문화 — 장기 정당성'),
        ('30~34', '민생·법·경제 — 제도'),
        ('35~36', '군사·외교 — 외부'),
        ('37~39', '일상 — 절제'),
        ('40',    '신종 — 자기의 마무리'),
    ]
    # 2열로 표시
    top_start = 2.05
    for i, (scope, role) in enumerate(layers):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = top_start + row * 0.55
        is_start_or_end = (scope == '1~2' or scope == '40')
        c = ACCENT if is_start_or_end else INK
        bg = PALE if is_start_or_end else RGBColor(0xFA, 0xFA, 0xFA)
        add_filled_rect(slide, Inches(x), Inches(y), Inches(6.0), Inches(0.5), bg)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.1), Inches(1.6), Inches(0.4),
                    f'[{scope}]', font_size=13, bold=True, color=c,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 1.8), Inches(y + 0.1), Inches(4.1), Inches(0.4),
                    role, font_size=13, color=c, bold=is_start_or_end)
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                '자기(1·2) → 사람·구조·계승·덕·점검·문화·제도·외부·일상 → 자기 마무리(40)',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 구성')
def s_part_1_5(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 구성', page, total)
    add_title(slide, '권1~5 — 주요 편의 핵심')
    rows = [
        ('1',  '군도(君道)',     '군주의 도리 — 십사소 10항목'),
        ('2',  '정체(政體)',     '창업과 수성·군주민수·겸청즉명'),
        ('3',  '임현(任賢)',     '현자 임용 — 인물 평가'),
        ('4',  '구간(求諫)',     '간언을 구하라 — 얼굴빛 관리'),
        ('5',  '납간(納諫)',     '간언을 받아들이라 — 매·"이 촌놈" 일화'),
        ('6~8', '군신·택관·봉건', '군신은 한 몸·소수정예·구조적 지속'),
        ('9~12','후계 교육 4편',  '태자제왕·존경사부·교계태자·규간태자'),
        ('13', '인의(仁義)',     '덕으로 품으면 복종 — 사형수 390명 일화'),
        ('14', '충의(忠義)',     '충은 복종이 아니라 바로잡기'),
        ('15', '효우(孝友)',     '"후대 군주는 결코 나를 본받지 말라"'),
        ('16', '공평(公平)',     '삼복주(三覆奏) 제도화 — 친인척 절제'),
        ('17', '성신(誠信)',     '"民無信不立" — 신용 자산'),
    ]
    top = 2.0
    for num, name, desc in rows:
        add_textbox(slide, Inches(0.7), Inches(top), Inches(0.8), Inches(0.4),
                    num, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.6), Inches(top), Inches(2.8), Inches(0.4),
                    name, font_size=15, bold=True, color=INK)
        add_textbox(slide, Inches(4.5), Inches(top), Inches(8.3), Inches(0.4),
                    desc, font_size=13, color=SUB)
        top += 0.42


@S('Ⅳ. 구성')
def s_part_6_10(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 구성', page, total)
    add_title(slide, '권6~10 — 주요 편의 핵심')
    rows = [
        ('18', '검약(儉約)',     '"사치는 모든 악정의 출발"'),
        ('19', '겸양(謙讓)',     '"가득 참을 경계하고 비어 있음을 기뻐하라"'),
        ('20', '인측(仁惻)',     '재해·기근 때 직접 살피고 구호'),
        ('21', '신소호(愼所好)', '"임금이 좋아하면 천하가 따른다" — 매 일화의 사상적 근거'),
        ('22', '신언어(愼言語)', '"천자의 말은 사관이 기록한다" — 영구 기록'),
        ('23', '두참사(杜讒邪)', '"비밀 정보는 1/3로 깎아서 듣는다"'),
        ('24', '회과(悔過)',     '"내가 틀렸다"고 공개 인정하는 위엄'),
        ('25/26', '사종/탐비',  '사치와 방종·탐욕과 비루함 경계'),
        ('27~29', '문화 3편',   '숭유학·문사(史官 직필)·예악'),
        ('30~34', '민생·법·경제', '무농·형법(당률)·사령·공부·변흥망'),
        ('35~36', '군사·외교',  '정벌·안변 — 兵者凶器'),
        ('37~40', '일상과 마무리', '행행·전렵·재상·신종 — 신종여시'),
    ]
    top = 2.0
    for num, name, desc in rows:
        add_textbox(slide, Inches(0.7), Inches(top), Inches(0.9), Inches(0.4),
                    num, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.7), Inches(top), Inches(3.0), Inches(0.4),
                    name, font_size=15, bold=True, color=INK)
        add_textbox(slide, Inches(4.8), Inches(top), Inches(8.0), Inches(0.4),
                    desc, font_size=13, color=SUB)
        top += 0.42


# ============== Ⅴ. 군도(君道) ==============
@S('Ⅴ. 군도')
def s_seonjeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 군도', page, total)
    add_title(slide, '군주의 근본 — 先正其身, 然後能正人')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '先 正 其 身   然 後 能 正 人',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '선정기신 연후능정인',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 군도편 제1 (위징의 답)', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.7), Inches(4.8), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(4.95), Inches(11.3), Inches(1.8), [
        ('태종: "군주의 도(道)는 무엇이 가장 먼저인가?"',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('위징: "먼저 제 몸을 바르게 한 뒤에야 남을 바르게 할 수 있습니다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('→ 修己治人(수기치인)의 원조 — 군주의 가장 큰 적은 외부가 아니라 자기 욕망',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅴ. 군도')
def s_halgo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 군도', page, total)
    add_title(slide, '할고담복(割股啖腹) — 자기파괴적 착취의 경계')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '割 股 啖 腹   腹 飽 而 身 斃',
                font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '할고담복 복포이신폐',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 군도편 제1', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
                '"임금이 백성을 수탈하여 자기 몸을 받드는 것은',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
                '자기 허벅지 살을 베어 자기 배를 채우는 것과 같다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.6),
                '"배는 부를지언정 몸은 죽는다"',
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '→ 착취 경제의 자기 파괴성 — 2,000년 전의 통찰',
                font_size=13, color=SUB, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅴ. 군도')
def s_10thoughts_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 군도', page, total)
    add_title(slide, '간태종십사소(諫太宗十思疏) ① — 1~5항',
              '위징이 올린 상소 — 태종이 병풍에 써 붙여 매일 본 체크리스트')
    items = [
        ('1', '見可欲, 則思知足',         '견가욕 즉사지족',
         '보고 싶은 것이 있으면 만족할 줄 앎을 생각'),
        ('2', '將有作, 則思知止',         '장유작 즉사지지',
         '토목을 일으키려 하면 그칠 줄을 생각'),
        ('3', '念高危, 則思謙沖',         '염고위 즉사겸충',
         '높은 자리를 생각하면 겸허히 낮춤을 생각'),
        ('4', '懼滿盈, 則思江海下百川',  '구만영 즉사강해하백천',
         '가득 참이 무서우면 강이 낮은 곳으로 흐르는 것을'),
        ('5', '盤遊, 則思三驅以爲度',    '반유 즉사삼구이위도',
         '사냥의 즐거움이 있으면 삼면 몰이의 옛 법을'),
    ]
    top = 2.3
    for num, han, eum, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(5.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.05), Inches(5.3), Inches(0.5),
                    han, font_size=17, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), Inches(top + 0.48), Inches(5.3), Inches(0.4),
                    eum, font_size=12, color=SUB)
        add_filled_rect(slide, Inches(7.0), Inches(top), Inches(5.8), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.15), Inches(top + 0.22), Inches(5.6), Inches(0.5),
                    desc, font_size=13, color=INK)
        top += 0.95


@S('Ⅴ. 군도')
def s_10thoughts_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 군도', page, total)
    add_title(slide, '간태종십사소 ② — 6~10항')
    items = [
        ('6',  '憂懈怠, 則思愼始而敬終',     '우해태 즉사신시이경종',
         '게으름이 두려우면 처음을 삼가고 끝을 공경함을'),
        ('7',  '慮壅蔽, 則思虛心以納下',     '여옹폐 즉사허심이납하',
         '막힌 말이 두려우면 허심탄회하게 아랫사람을'),
        ('8',  '想讒邪, 則思正身以黜惡',     '상참사 즉사정신이출악',
         '참소가 두려우면 제 몸을 바르게 하여 악을 물리침'),
        ('9',  '恩所加, 則思無因喜以謬賞',   '은소가 즉사무인희이류상',
         '은혜를 베풀려면 사사로운 기쁨으로 잘못된 상을'),
        ('10', '罰所及, 則思無以怒而濫刑',   '벌소급 즉사무이노이남형',
         '벌을 내리려면 노여움으로 형을 남용하지 않을 것을'),
    ]
    top = 2.3
    for num, han, eum, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(5.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.05), Inches(5.3), Inches(0.5),
                    han, font_size=16, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), Inches(top + 0.48), Inches(5.3), Inches(0.4),
                    eum, font_size=12, color=SUB)
        add_filled_rect(slide, Inches(7.0), Inches(top), Inches(5.8), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.15), Inches(top + 0.22), Inches(5.6), Inches(0.5),
                    desc, font_size=13, color=INK)
        top += 0.95


@S('Ⅴ. 군도')
def s_10thoughts_structure(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 군도', page, total)
    add_title(slide, '십사소의 3층 구조 — 자기·조직·타자의 점검')
    layers = [
        ('1~5번', '자기 욕망의 제어',
         '소유욕·확장욕·자만·포만감·여가까지\n다섯 가지 개인의 욕망을 한 번씩 점검',
         ACCENT),
        ('6~7번', '리더십의 유지',
         '게으름의 위험·정보 차단의 위험\n초심 유지와 개방성',
         RGBColor(0xA0, 0x40, 0x40)),
        ('8~10번','공정한 상벌',
         '참소·편향된 보상·감정적 처벌\n관계를 다루는 공정의 3원칙',
         SUB),
    ]
    top = 2.4
    for tag, title, desc, color in layers:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(1.2), color)
        add_textbox(slide, Inches(0.6), Inches(top + 0.42), Inches(2.5), Inches(0.5),
                    tag, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.3), Inches(top), Inches(3.0), Inches(1.2), PALE)
        add_textbox(slide, Inches(3.3), Inches(top + 0.42), Inches(3.0), Inches(0.5),
                    title, font_size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.5), Inches(top), Inches(6.3), Inches(1.2),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.7), Inches(top + 0.18), Inches(6.0), Inches(1.0),
                    desc, font_size=14, color=INK)
        top += 1.4
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '리더가 자기 → 조직 → 타자 세 차원에서 검토해야 할 열 가지 순간',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅵ. 정체(政體) ==============
@S('Ⅵ. 정체')
def s_chang_su(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 정체', page, total)
    add_title(slide, '창업과 수성(創業守成) — 어느 것이 더 어려운가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.35), Inches(12.3), Inches(0.5),
                '태종: "제왕의 사업에서 창업(創業)이 어려운가, 수성(守成)이 어려운가?"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.2), Inches(5.9), Inches(3.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(3.35), Inches(5.9), Inches(0.5),
                '방현령 — 創業', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.95), Inches(5.3), Inches(2.5), [
        ('"창업이 어렵습니다"', {'font_size': 16, 'bold': True}),
        ('', {'font_size': 6}),
        ('"천하가 어지러울 때', {'font_size': 14, 'space_before': 6}),
        (' 여러 영웅이 다투니', {'font_size': 14}),
        (' 이를 제압하는 일이 쉽지 않습니다"', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('— 창업 경험자의 증언', {'font_size': 13, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(3.2), Inches(5.9), Inches(3.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(3.35), Inches(5.9), Inches(0.5),
                '위징 — 守成', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.95), Inches(5.3), Inches(2.5), [
        ('"수성이 더 어렵습니다"', {'font_size': 16, 'bold': True}),
        ('', {'font_size': 6}),
        ('"제왕은 간난(艱難) 속에', {'font_size': 14, 'space_before': 6}),
        (' 천하를 얻고', {'font_size': 14}),
        (' 편안함 속에 잃습니다"', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('— 안정기의 통찰', {'font_size': 13, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                '태종: "창업의 어려움은 이미 지나갔다. 수성의 어려움은 여러분과 함께 삼가겠다"',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅵ. 정체')
def s_chang_su_meaning(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 정체', page, total)
    add_title(slide, '"수성기의 적은 교만·안일·사치" — 성공의 역설')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(1.7), [
        ('創 業 難   守 成 亦 難',
         {'font_size': 36, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('창업도 어렵지만 수성도 어렵다',
         {'font_size': 16, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
                '현대 경영학의 같은 발견', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(2.0), [
        ('• "성공한 기업의 혁신 실패" — Innovator\'s Dilemma (크리스텐슨)',
         {'font_size': 15}),
        ('• 빌 게이츠 — "성공은 나쁜 교사다(Success is a lousy teacher)"',
         {'font_size': 15, 'space_before': 6}),
        ('• 스타트업 성공 후 안주 단계가 가장 위험',
         {'font_size': 15, 'space_before': 6}),
        ('• "성공한 순간이 가장 위험한 순간"이라는 통찰의 동양적 원형',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
    ], line_spacing=1.4)


@S('Ⅵ. 정체')
def s_boat_water(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 정체', page, total)
    add_title(slide, '군주민수(君舟民水) — 배와 물')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.0),
                '君 者 舟 也   庶 人 者 水 也',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.0),
                '水 則 載 舟   水 則 覆 舟',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.5),
                '군자주야 서인자수야 · 수즉재주 수즉복주',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
                '— 정체편 제2 (순자의 말을 태종이 자주 인용)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
                '"임금은 배, 백성은 물 — 물은 배를 띄우기도 하고 뒤집기도 한다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '태종이 태자에게 배를 타며 반복 훈계 — 민본의 절대 원칙',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅵ. 정체')
def s_gyeongcheong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 정체', page, total)
    add_title(slide, '兼聽則明 偏信則暗 — 두루 들으면 밝다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0),
                '兼 聽 則 明   偏 信 則 暗',
                font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
                '겸청즉명 편신즉암',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
                '— 정체편 제2 (위징의 답)', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.7), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
                '두루 들으면 밝고, 한쪽만 믿으면 어둡다',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.9), Inches(11.9), Inches(1.2), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(6.05), Inches(11.3), Inches(1.0), [
        ('현대의 다양성 경영(Diversity Management)의 원조',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('한쪽 정보에만 의존하지 않는 의사결정 시스템',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅶ. 간쟁(諫諍) ==============
@S('Ⅶ. 간쟁')
def s_ganjaeng(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 간쟁', page, total)
    add_title(slide, '간쟁(諫諍) — 정관 정치의 심장')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(11.9), Inches(0.5),
                '"사람이 제 얼굴을 보려면 거울이 반드시 있어야 한다',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(2.83), Inches(11.9), Inches(0.5),
                '임금이 자기의 허물을 알려면 충신이 반드시 있어야 한다"',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.4), Inches(5.9), Inches(3.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(3.55), Inches(5.9), Inches(0.5),
                '구간(求諫) — 간언을 구하라', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(4.2), Inches(5.3), Inches(2.5), [
        ('"임금이 간언을 구할 때는', {'font_size': 15, 'bold': True}),
        (' 반드시 너그러운 낯을 지어야 한다"', {'font_size': 15}),
        ('', {'font_size': 6}),
        ('"엄한 얼굴로 \'말하라\' 해봐야', {'font_size': 14, 'color': SUB, 'space_before': 6}),
        (' 누가 바른말을 하겠는가?"', {'font_size': 14, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(3.4), Inches(5.9), Inches(3.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(3.55), Inches(5.9), Inches(0.5),
                '납간(納諫) — 간언을 받아들이라', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(4.2), Inches(5.3), Inches(2.5), [
        ('"화가 나도 참는 근육"을 기르는 일', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 6}),
        ('황후 · 위징 · 자기 자신', {'font_size': 14, 'space_before': 6}),
        ('— 세 겹의 장치로 참아낸다', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('"참는 제도"가 아니라', {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ('"참는 습관"이 명군을 만든다', {'font_size': 14, 'color': SUB}),
    ], line_spacing=1.3)


@S('Ⅶ. 간쟁')
def s_face(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 간쟁', page, total)
    add_title(slide, '구간(求諫) — 얼굴빛부터 관리하라')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(1.7), [
        ('"간언은 제도보다 분위기의 문제"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('리더의 미간·말투·눈빛이 조직의 정보 흐름을 결정한다',
         {'font_size': 17, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(4.55), Inches(12.3), Inches(0.5),
                '현대 응용', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(2.0), [
        ('• CEO가 회의에서 먼저 웃는 얼굴로 질문', {'font_size': 15}),
        ('• "내가 틀렸을 수 있다"는 자기 오류 인정 문화',
         {'font_size': 15, 'space_before': 6}),
        ('• 아마존 6-pager·침묵의 독서 — 감정 개입 전 문서 검토',
         {'font_size': 15, 'space_before': 6}),
        ('• 구글 "심리적 안전감(psychological safety)" — Aristotle 연구',
         {'font_size': 15, 'space_before': 6}),
    ], line_spacing=1.4)


@S('Ⅶ. 간쟁')
def s_200_remonstrations(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 간쟁', page, total)
    add_title(slide, '위징의 200여 회 직간(直諫)')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '200',
                font_size=180, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
                '여 회의 직간 — 정관 17년 동안',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.2), Inches(2.3), color=RULE, weight=1.5)
    add_paragraphs(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(2.0), [
        ('일부는 목숨을 걸고 — 다른 신하라면 처형되었을 발언도 서슴지 않음',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('가장 유명한 두 상소 — 「간태종십사소」·「신종소」',
         {'font_size': 15, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('태종이 병풍에 써 붙여 매일 점검',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅶ. 간쟁')
def s_first_meeting(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 간쟁', page, total)
    add_title(slide, '일화 ① 첫 대면 (626년) — "왜 우리 형제를 이간했는가"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('현무문의 변 직후 — 태종이 포로 위징을 심문',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('태종: "너는 왜 우리 형제 사이를 이간질했는가?"',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('위징: "만약 태자(이건성)가 일찍 내 말을 들었다면',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('       오늘 같은 화는 없었을 것입니다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.4), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.1), [
        ('좌우 신하들: "저놈을 죽이라"',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('태종은 오히려 그를 간의대부(諫議大夫)로 임명',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('→ "적이었던 자도 재능이 있으면 쓴다"의 전형',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('이 한 장면이 정관 23년의 군신 관계를 결정',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅶ. 간쟁')
def s_country_bumpkin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 간쟁', page, total)
    add_title(slide, '일화 ② "이 촌놈을 죽이겠다" — 황후의 개입')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.7),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.4), [
        ('어느 날 위징의 신랄한 간언 후, 태종이 내실에서 분통:',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"이 촌놈(위징)을 내가 반드시 죽여야겠다!"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.8), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5), [
        ('황후 장손씨가 정장(正裝)으로 나와 절을 올림',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"임금이 현명해야 신하가 곧은 말을 합니다(君明臣直)"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"위징이 곧은 말을 할 수 있는 것은 폐하가 현명하기 때문이니',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('   제가 어찌 축하하지 않겠습니까"',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('태종의 분노가 풀림 — 황후의 지혜가 위징과 태종 사이의 가교',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅶ. 간쟁')
def s_hawk(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 간쟁', page, total)
    add_title(slide, '일화 ③ 매를 품에 숨긴 일화 — 자기 절제의 극치')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.8), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(2.5), [
        ('태종이 희귀한 매(鷹)를 팔뚝에 앉혀 놀고 있는데',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('멀리서 위징이 오는 것이 보임',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('태종은 급히 매를 품 안에 숨김',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('위징은 의도적으로 긴 보고를 올림',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('보고가 끝날 무렵 — 태종의 품에서 매가 질식해 죽어 있었다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('태종은 아무 말 없이 웃기만 함',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
                '"군주의 마음이 취미에 기우는 것"을 스스로 경계',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6),
                '위징은 그 미묘한 심리를 알고 행동으로 제어',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                '→ 정관 정치의 "자기 절제"의 극치를 보여주는 일화',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅶ. 간쟁')
def s_tribute(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 간쟁', page, total)
    add_title(slide, '위징 사망 후 — "세 거울" 추도 (643년)')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.2), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.9), [
        ('643년 위징이 죽자, 태종은 통곡하며:',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"以銅爲鑒, 可以正衣冠;',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 以古爲鑒, 可以知興替;',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 以人爲鑒, 可以明得失."',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"今魏徵沒, 朕亡一鑒矣."',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.3)
    add_rule(slide, Inches(5.5), Inches(5.6), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.6),
                '"이제 위징이 죽으니, 짐이 거울 하나를 잃었다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '정관정요 전체에서 가장 유명한 한 구절',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅷ. 三鏡 ==============
@S('Ⅷ. 三鏡')
def s_3mirrors(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 三鏡', page, total)
    add_title(slide, '三鏡 — 세 개의 거울')
    mirrors = [
        ('銅', '동', '구리(銅) 거울',  '正衣冠',  '의관(衣冠)을 바로 함',  '일상의 자기 점검'),
        ('古', '고', '옛일의 거울',     '知興替',  '흥망(興替)을 앎',         '역사의 거시적 패턴'),
        ('人', '인', '사람의 거울',     '明得失',  '득실(得失)을 밝힘',      '실시간 직언자'),
    ]
    for i, (han, eum, name, role_han, role_eum, mean) in enumerate(mirrors):
        x = 0.5 + i * 4.3
        add_filled_rect(slide, Inches(x), Inches(2.3), Inches(4.1), Inches(4.7), PALE)
        add_textbox(slide, Inches(x), Inches(2.45), Inches(4.1), Inches(1.4),
                    han, font_size=120, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.15), Inches(4.1), Inches(0.5),
                    name, font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 0.2), Inches(4.85), Inches(3.7), Inches(0.5),
                        ACCENT)
        add_textbox(slide, Inches(x + 0.2), Inches(4.95), Inches(3.7), Inches(0.4),
                    role_han, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(5.45), Inches(3.7), Inches(0.4),
                    role_eum, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(5.95), Inches(3.7), Inches(0.8),
                    mean, font_size=14, color=INK, align=PP_ALIGN.CENTER)


@S('Ⅷ. 三鏡')
def s_3mirrors_meaning(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 三鏡', page, total)
    add_title(slide, '세 거울의 의미 — 현대 리더십에의 적용')
    items = [
        ('銅鑒 — 일상 점검',  '분기별 성과 리뷰 · 자기 평가 · KPI · OKR',
         '눈에 보이는 자기 모습 — 즉각적 피드백'),
        ('古鑒 — 역사 패턴',   '업계 역사 · 선례 분석 · 케이스 스터디',
         '같은 실수를 반복하지 않기 위한 거시적 시야'),
        ('人鑒 — 실시간 직언', '솔직히 직언하는 멘토 · 사외이사 · 코치',
         '가장 희귀하고 가장 중요한 거울'),
    ]
    top = 2.3
    for tag, modern, meaning in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(1.3), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.45), Inches(3.5), Inches(0.5),
                    tag, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(8.5), Inches(1.3), PALE)
        add_textbox(slide, Inches(4.5), Inches(top + 0.12), Inches(8.2), Inches(0.5),
                    modern, font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(4.5), Inches(top + 0.7), Inches(8.2), Inches(0.5),
                    meaning, font_size=14, color=INK)
        top += 1.5


@S('Ⅷ. 三鏡')
def s_human_mirror(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 三鏡', page, total)
    add_title(slide, '人鏡 — 가장 희귀한 거울')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '人 鑒',
                font_size=200, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.7), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.7), Inches(5.0), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.7), [
        ('"위징 같은 사람 한 명"의 가치가 조직의 운명을 좌우한다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('태종은 위징 사후 흔들리기 시작했다 — 인경 부재의 대가',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('고구려 원정 실패 후 "위징이 있었다면…" 한탄',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅸ. 군신(君臣)의 도 ==============
@S('Ⅸ. 군신의 도')
def s_one_body(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 군신의 도', page, total)
    add_title(slide, '군신은 한 몸 — 머리와 팔다리')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(1.7), [
        ('군주와 신하는 한 몸의 머리와 팔다리',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('머리가 밝아야 팔다리가 건실하고',
         {'font_size': 17, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('팔다리가 건실해야 머리가 안전하다',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.5),
                '— 군신감계편 제6의 핵심 사상', font_size=14, color=SUB,
                align=PP_ALIGN.CENTER, bold=True)
    add_filled_rect(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.7),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(1.0), Inches(5.35), Inches(11.3), Inches(1.5), [
        ('현대 적용',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('조직의 부패는 "말하지 않는 신하"에서 시작',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('리더는 나쁜 소식을 가져오는 자를 먼저 보호해야 정보가 흐른다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅸ. 군신의 도')
def s_true_loyalty(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 군신의 도', page, total)
    add_title(slide, '진정한 충(忠) — 복종이 아니라 바로잡기')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '잘못된 충(忠)', font_size=18, bold=True, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.2), Inches(5.3), Inches(3.5), [
        ('"맹목적 복종"', {'font_size': 18, 'bold': True}),
        ('', {'font_size': 6}),
        ('임금이 옳건 그르건', {'font_size': 15, 'color': SUB, 'space_before': 6}),
        ('일단 "예"라고 답한다', {'font_size': 15, 'color': SUB}),
        ('', {'font_size': 6}),
        ('수 양제의 우세기(虞世基)', {'font_size': 14, 'color': ACCENT, 'space_before': 6}),
        ('— 비위만 맞추다 나라를 망쳤다', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '진정한 충(忠)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.2), Inches(5.3), Inches(3.5), [
        ('"바로잡기(正君)"', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('', {'font_size': 6}),
        ('임금이 잘못하면', {'font_size': 15, 'space_before': 6}),
        ('목숨 걸고 직언한다', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 6}),
        ('위징의 200여 회 직간', {'font_size': 14, 'color': ACCENT, 'space_before': 6}),
        ('— 정관정요가 정의하는 진짜 충신', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)


@S('Ⅸ. 군신의 도')
def s_sui_huse(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 군신의 도', page, total)
    add_title(slide, '양제와 우세기 — 정관의 반면교사')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.8),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.5), [
        ('태종이 되짚는다 — 군신감계편:',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"양제(煬帝)는 스스로 성군이라 자처하고 간언을 싫어했다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 우세기(虞世基)는 양제의 비위만 맞췄다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 그래서 도적이 성 아래 몰려왔는데도 양제는 알지 못했다"',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"임금이 혼미하면 신하가 아첨한다 — 두렵다(可懼也)"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '현대 응용', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.5), [
        ('• "예스맨"으로만 채워진 임원진은 조직 멸망의 신호',
         {'font_size': 15}),
        ('• 다양한 목소리가 사라진 순간이 가장 위험',
         {'font_size': 15, 'space_before': 6}),
    ], line_spacing=1.3)


# ============== Ⅹ. 임현(任賢) ==============
@S('Ⅹ. 임현')
def s_get_people(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 임현', page, total)
    add_title(slide, '為政之要 唯在得人 — 정치의 요체는 사람을 얻는 것')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '為 政 之 要   唯 在 得 人',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '위정지요 유재득인',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 택관편 제7', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
                '"정치의 요체는 오직 사람을 얻는 데 있다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.1), [
        ('"官在得人, 不在員多"', {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('관리는 사람을 얻는 데 있지 수가 많은 데 있지 않다 — 소수정예의 원조',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅹ. 임현')
def s_use_enemy(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 임현', page, total)
    add_title(slide, '"적이었던 자도 재능이 있으면 쓴다"')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(1.7), [
        ('대표 사례 — 위징(魏徵)',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('형 이건성의 모사였다가 → 태종의 최고 간관(諫官)으로',
         {'font_size': 17, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"적의 책사라도 재능이 있으면 등용"의 가장 극적인 예',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
                '현대 사례', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(2.0), [
        ('• 디즈니의 픽사 인수(2006) 후 존 래시터 중용 — 적대적 M&A 후 상대 회사 임원 기용',
         {'font_size': 14}),
        ('• 창업 경쟁자 출신을 CTO로 영입하는 글로벌 IT 기업 사례',
         {'font_size': 14, 'space_before': 6}),
        ('• "능력 vs 충성"의 갈등에서 능력을 택하는 결단력',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
    ], line_spacing=1.4)


@S('Ⅹ. 임현')
def s_team_strength(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 임현', page, total)
    add_title(slide, '강점의 조합 — 정관 핵심팀의 구성')
    rows = [
        ('위 징',    '직간(直諫)',   '잘못을 바로잡는 거울'),
        ('방현령',   '꾀(謀)',       '큰일을 맡기는 기획자'),
        ('두여회',   '결단(斷)',     '결정을 내리는 실행자'),
        ('왕 규',    '시비 판별',     '옳고 그름을 가리는 평가자'),
        ('이 정',    '군사(軍事)',   '돌궐을 평정한 명장'),
        ('우세남',   '문예(文藝)',   '문학·서예의 조언자'),
    ]
    top = 2.2
    for i, (name, strength, role) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(0.65), PALE)
        add_textbox(slide, Inches(0.5), Inches(top + 0.17), Inches(2.5), Inches(0.4),
                    name, font_size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.05), Inches(top), Inches(3.5), Inches(0.65), bg)
        add_textbox(slide, Inches(3.2), Inches(top + 0.17), Inches(3.3), Inches(0.4),
                    strength, font_size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.6), Inches(top), Inches(6.2), Inches(0.65), bg)
        add_textbox(slide, Inches(6.8), Inches(top + 0.17), Inches(5.9), Inches(0.4),
                    role, font_size=15, color=INK)
        top += 0.7
    add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
                '서로 다른 강점을 가진 6인의 조합 → 정관의 안정 23년',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅹ. 임현')
def s_wanggui_insight(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 임현', page, total)
    add_title(slide, '왕규의 자기 인식 — "자기를 아는 자가 남도 안다"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(3.2), [
        ('태종: "신하들의 장단을 평해 보라. 그리고 자네 자신은 어떤가?"',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('왕규(王珪): "방현령은 큰일을 맡길 만하고',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('             두여회는 결단을 맡길 만하고',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('             위징은 잘못을 바로잡아 드릴 만합니다',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('             저는 그저 옳고 그름을 가려 말하는 데',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('             약간의 재주가 있을 뿐입니다"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
                '"자기 인식이 곧 인재 안목"',
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                '자기 자리를 정확히 아는 자만이 남의 자리도 정확히 본다',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅺ. 민본 ==============
@S('Ⅺ. 민본')
def s_minbon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 민본', page, total)
    add_title(slide, '民惟邦本 — 백성이 나라의 근본')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0),
                '民 惟 邦 本   本 固 邦 寧',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '민유방본 본고방녕',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
                '— 정관정요 곳곳, 태종이 자주 인용 (서경 출전)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.6), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.6),
                '백성이 오직 나라의 근본이니, 근본이 굳으면 나라가 편안하다',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.1), [
        ('세금·형벌·전쟁·토목 모든 결정에서',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('"백성의 물이 거칠어지는가"를 본다 — 민본의 절대 원칙',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅺ. 민본')
def s_taizi_teach(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 민본', page, total)
    add_title(slide, '태자에게 배를 타며 — 군주민수의 산교육')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.7), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(2.4), [
        ('태종이 태자와 함께 배를 타며:',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"물이 배를 띄우지만 뒤집기도 한다',
         {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 백성이 곧 물이니 두려워해야 한다"',
         {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.6), Inches(5.2), Inches(11.9), Inches(1.8),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.9), Inches(5.35), Inches(11.3), Inches(1.6), [
        ('밥그릇 앞에서 (또 다른 일화):',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"네가 이 밥이 어디서 나왔는지 아느냐?',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        (' 농부가 봄·여름·가을 땀 흘려 지은 것이다',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 그 수고를 알아야 이 밥을 먹을 자격이 있다"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅺ. 민본')
def s_tightrope(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 민본', page, total)
    add_title(slide, '戰戰兢兢 如履薄冰 — 살얼음 위를 걷듯')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '戰 戰 兢 兢   如 履 薄 冰',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '전전긍긍 여리박빙',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 시경 출전 · 태종이 자기 마음가짐을 표현할 때 자주 인용',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
                '"두렵고 조심하기를 마치 살얼음을 밟듯이"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.3), PALE)
    add_textbox(slide, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.5),
                '천하의 군주가 가질 마음 — 자만이 아니라 조심',
                font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
                '최고의 자리일수록 가장 두려워해야 한다는 통찰',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅻ. 권5 — 5대 핵심 덕목 ==============
@S('Ⅻ. 5대 덕목')
def s_inui(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 5대 덕목', page, total)
    add_title(slide, '인의(仁義) — 사형수 390명의 귀가 일화')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(2.2), [
        ('태종이 전국의 사형수 390명을',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        ('임시 귀가시켜 설을 쇠게 하고, 이듬해 다시 돌아와 형을 받게 함',
         {'font_size': 17, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('— 모두 약속대로 돌아왔고, 태종이 이들을 특사(特赦) —',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
                '懷 之 以 德   則 無 不 服 者',
                font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
                '회지이덕 즉무불복자 — "덕으로 품으면 복종하지 않는 자가 없다"',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                '— 인의편 제13', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '신뢰받은 사람은 신뢰로 응답한다 — 경영학 연구와 일치',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅻ. 5대 덕목')
def s_chungui(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 5대 덕목', page, total)
    add_title(slide, '충의(忠義) — 충(忠)의 재정의')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.5),
                '"충(忠)은 복종이 아니라 바로잡기다"',
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(3.3), Inches(11.9), Inches(1.6), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(3.45), Inches(11.3), Inches(1.4), [
        ('위징 같은 직간자(直諫者)가 진정한 충신',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('비위 맞춤은 아첨일 뿐, 진짜 충은 임금을 바른 길로 이끄는 것',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.5),
                '현대 응용 — "충성의 정의 전환"', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.5), [
        ('• 복종하는 자보다 바른말을 하는 자가 조직을 위한다',
         {'font_size': 15, 'bold': True}),
        ('• "좋은 직원"의 정의 — Yes-man이 아니라 무엇이 잘못되었는지 알리는 사람',
         {'font_size': 14, 'space_before': 6}),
    ], line_spacing=1.4)


@S('Ⅻ. 5대 덕목')
def s_hyo_woo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 5대 덕목', page, total)
    add_title(slide, '효우(孝友) — 태종의 자기 고백')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('"임금이 효·우를 갖추지 못하면',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        (' 어찌 천하에 덕화(德化)를 베풀겠는가"',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('그러나 태종 자신은 형 이건성, 동생 이원길을 살해한 자',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('이 모순이 효우편에 묘한 긴장을 만든다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('태종의 공개 선언',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"후대 군주는 결코 나를 본받지 말라"',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('→ 리더가 자기 과오를 인정·공표할 때 조직 전체에 정직의 문화가 퍼진다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅻ. 5대 덕목')
def s_gongpyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 5대 덕목', page, total)
    add_title(slide, '공평(公平) — 삼복주(三覆奏) 제도화')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(1.7), [
        ('일화 — 장온고(張蘊古) 처형의 후회',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('총애하던 장온고가 중죄 → 태종이 분노하여 즉시 처형',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('뒤에 깊이 후회: "내가 너무 성급히 처형했다"',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.6), Inches(4.5), Inches(11.9), Inches(2.4),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_textbox(slide, Inches(0.6), Inches(4.65), Inches(11.9), Inches(0.5),
                '→ 三 覆 奏 제도 제정 — 사형 판결의 3중 재심',
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(5.25), Inches(11.5), Inches(1.6), [
        ('① 1차 법관 판결 → ② 2차 중앙 심리 재심 → ③ 3차 황제 직속 재심',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('+ 집행 당일 마지막 확인 → 한 번이라도 이의 있으면 즉시 중단',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('이 제도는 당률에 편입 → 이후 천 년 동아시아 법제의 표준',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)


@S('Ⅻ. 5대 덕목')
def s_seongsin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 5대 덕목', page, total)
    add_title(slide, '성신(誠信) — 民無信不立')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '民 無 信 不 立',
                font_size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '민무신불립',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
                '— 성신편 제17 (논어 인용)', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.6), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.6),
                '"백성에게 신의(信)가 없으면 나라가 서지 않는다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.95), Inches(11.3), Inches(1.1), [
        ('현대 적용 — 리더의 말 한마디가 조직의 신용 자산',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('약속과 어긋나면 반드시 공개 정정 — 일관성의 정치',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== XⅢ. 권6 — 9대 자기 점검 ==============
@S('XⅢ. 9대 병폐')
def s_9_disorders(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 9대 병폐', page, total)
    add_title(slide, '권6 — 9대 병폐 자기 점검 체크리스트',
              '군주가 빠지기 쉬운 개인적 함정 9가지 — 현대 리더의 일일 점검')
    items = [
        ('18', '검약(儉約)',   '사치 경계 — 모든 악정의 출발'),
        ('19', '겸양(謙讓)',   '가득 참을 경계하라'),
        ('20', '인측(仁惻)',   '재난 시 직접 살핀다'),
        ('21', '신소호(愼所好)','좋아하는 것을 삼가라'),
        ('22', '신언어(愼言語)','말은 영구 기록이다'),
        ('23', '두참사(杜讒邪)','참소를 1/3로 깎아 듣는다'),
        ('24', '회과(悔過)',   '잘못을 공개 인정한다'),
        ('25', '사종(奢縱)',   '사치와 방종을 경계'),
        ('26', '탐비(貪鄙)',   '작은 이익을 탐하지 말라'),
    ]
    top = 2.3
    for i, (num, name, desc) in enumerate(items):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.3
        y = top + row * 1.45
        add_filled_rect(slide, Inches(x), Inches(y), Inches(4.1), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(x), Inches(y + 0.13), Inches(4.1), Inches(0.4),
                    f'{num}. {name}', font_size=15, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x), Inches(y + 0.6), Inches(4.1), Inches(0.7), PALE)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.7), Inches(3.8), Inches(0.5),
                    desc, font_size=13, color=INK, align=PP_ALIGN.CENTER)


@S('XⅢ. 9대 병폐')
def s_geomyak(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 9대 병폐', page, total)
    add_title(slide, '검약(儉約) — 사치는 모든 악정의 출발')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.7),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.4), [
        ('"수 양제는 낙양·강도 궁전을 짓다 나라를 잃었다"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('— 정관정요가 반복적으로 환기하는 반면교사',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.4), [
        ('정관정요의 통찰',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('사치는 단순한 재정 낭비가 아니라',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"군주의 마음이 기울어지는 표지"',
         {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('황후 장손씨도 보석·장식을 거절하며 솔선수범',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"한 나라의 어머니가 사치하면 모든 부녀자가 따라 합니다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('XⅢ. 9대 병폐')
def s_sinhoho(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 9대 병폐', page, total)
    add_title(slide, '신소호(愼所好) — "좋아하는 것을 삼가라"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"임금이 매(鷹)를 좋아하면 천하가 매를 바치고',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 임금이 글씨를 좋아하면 천하가 글씨를 바친다"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.3), [
        ('리더의 취향 = 조직 자원 배분의 암묵적 신호',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('좋아하는 것일수록 공개하지 말라',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('현대 사례 — 잡스가 커피를 좋아한다고 알려지자',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('모든 회의에 다양한 커피가 등장하고 정작 의제는 뒤로 밀림',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('XⅢ. 9대 병폐')
def s_sineoneo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 9대 병폐', page, total)
    add_title(slide, '신언어(愼言語) — "말을 삼가라"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.2), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.9), [
        ('"천자(天子)의 말 한마디는 사관(史官)이 기록한다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 후대가 보는 것이니 어찌 삼가지 않겠는가"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(2.1), [
        ('리더의 말은 영구 기록 — 무심코 한 농담도 정책이 된다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('현대 사례 — SNS 시대의 CEO 트윗 하나가 주가를 움직인다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('일론 머스크의 트윗 사건 — 2,000년 전의 통찰이 그대로',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('XⅢ. 9대 병폐')
def s_duchamsa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 9대 병폐', page, total)
    add_title(slide, '두참사(杜讒邪) — "참소를 막아라"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.3), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0), [
        ('"참소(讒訴)는 원래 교묘해서 처음에는 알아채기 어렵다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 그러므로 간언이 오면 공개된 자리에서 검증하고',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 비밀리에 올라온 고자질은 삼분의 일로 깎아서 듣는다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.95), Inches(11.7), Inches(2.0), [
        ('"비밀 정보는 과대평가되기 쉽다"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('투명한 검증 프로세스가 조직을 살린다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('익명 제보·비밀 보고의 신뢰성 검증 절차 필수',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('XⅢ. 9대 병폐')
def s_hoegwa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 9대 병폐', page, total)
    add_title(slide, '회과(悔過) — "잘못을 뉘우쳐라"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.3), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0), [
        ('태종은 여러 차례 공개적으로 자기 실책을 인정',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"내가 틀렸다"고 말할 수 있는 것이 임금의 위엄',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.95), Inches(11.7), Inches(2.0), [
        ('현대 사례',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('워런 버핏의 연례 주주서한 — "올해 내 실수" 코너 상시',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('제프 베조스 — "Disagree and commit" 후 실패하면 공개 인정',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('→ 리더의 자기 오류 인정이 조직의 학습 문화를 만든다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


# ============== XⅣ. 문화·민생 ==============
@S('XⅣ. 문화·민생')
def s_sunghyu(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 문화·민생', page, total)
    add_title(slide, '숭유학(崇儒學) — 홍문관과 18학사',
              '문화는 장식이 아니라 통치의 기반')
    items = [
        ('홍문관(弘文館)', '태종이 설치한 학술 기관 — 정관 정치의 두뇌'),
        ('18학사(十八學士)', '당대 최고 학자 18명을 궁에 두고 — 태종이 직접 강론에 참여'),
        ('공자 배향',        '공자를 국가 제사의 정식 대상으로'),
        ('『오경정의(五經正義)』', '경전 표준 주석서 — 후대 1,000년의 교과서'),
        ('국가 편찬 사업',     '『진서(晉書)』·『수서(隋書)』 등 정사 편찬'),
    ]
    top = 2.5
    for tag, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.17), Inches(3.5), Inches(0.4),
                    tag, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(8.5), Inches(0.7), PALE)
        add_textbox(slide, Inches(4.5), Inches(top + 0.17), Inches(8.2), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.8
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                'CEO의 "학습하는 조직" 구축 — 베조스 매일 독서, 게이츠 Think Week의 원형',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('XⅣ. 문화·민생')
def s_historian(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 문화·민생', page, total)
    add_title(slide, '사관(史官)의 직필(直筆) — "임금도 자기 기록을 보지 말라"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.3),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0), [
        ('태종이 자기 즉위 과정(현무문의 변)을 기록한',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('실록을 보자고 청함',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('사관: "군주가 사서(史書)를 보면 직필(直筆)이 어려워집니다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('태종은 처음엔 불쾌했으나 결국 받아들임',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.95), Inches(11.7), Inches(2.0), [
        ('기록의 독립성 — 리더도 자기 기록을 손대지 않는 원칙',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"권력이 기록을 마음대로 고치면 역사는 거짓이 된다"',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('현대의 회계 감사·외부 이사·언론 자유의 원형',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('XⅣ. 문화·민생')
def s_munong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 문화·민생', page, total)
    add_title(slide, '무농(務農) — "농업은 나라의 근본(國之本)"')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '農 — 國 之 本',
                font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(3.6), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.05), Inches(11.7), Inches(1.8), [
        ('태종의 친경(親耕) — 매년 봄 황제가 직접 밭을 갈고 씨를 뿌림',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('농번기에는 요역(徭役) 동원을 피함',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('— 백성의 일정에 국가가 맞춘다 —',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_textbox(slide, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.5),
                '균전제(均田制) — 모든 농민에게 일정 면적의 땅을 지급',
                font_size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
                '조용조(租·庸·調) — 곡식·노동·옷감을 정해진 비율로 거둠',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('XⅣ. 문화·민생')
def s_dangryul(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 문화·민생', page, total)
    add_title(slide, '형법(刑法) — 당률(唐律)과 관대한 법')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '寬 簡 之 法 — 관대하고 평이한 법',
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rows = [
        ('5형(五刑) 체계',  '笞(태)·杖(장)·徒(도)·流(유)·死(사)'),
        ('정상참작',         '나이·신분·환경 고려'),
        ('약자 보호',         '노인·어린이·여성·장애인 특별 조항'),
        ('삼복주(三覆奏)',    '사형은 3중 재심 + 집행 당일 확인'),
        ('수의 가혹한 법 폐지', '엄벌주의를 인정주의로 전환'),
    ]
    top = 3.3
    for tag, desc in rows:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.6), PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.15), Inches(3.5), Inches(0.4),
                    tag, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(8.5), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.5), Inches(top + 0.15), Inches(8.2), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.68
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '→ 당률은 이후 1,000년 동아시아(한국·일본·베트남) 법제의 뼈대',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== XⅤ. 군사·외교 ==============
@S('XⅤ. 군사·외교')
def s_byeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 군사·외교', page, total)
    add_title(slide, '兵者凶器 — "병은 흉기라 부득이할 때만 쓴다"')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '兵 者 凶 器   不 得 已 用 之',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '병자흉기 부득이용지',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 정벌편 제35 (노자·손자병법과 통하는 사상)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
                '"전쟁(兵)은 흉기(凶器)이니 부득이할 때만 쓴다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.1), [
        ('손자병법의 "不戰而屈人之兵"과 같은 정신',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('정관의 외교는 "정복보다 회유·교역·혼인 동맹"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('XⅤ. 군사·외교')
def s_dolgwol(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 군사·외교', page, total)
    add_title(slide, '돌궐 평정(630) — 천가한(天可汗) 칭호')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('이정(李靖) 장군의 기습 작전 → 힐리가한(頡利可汗) 생포',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('항복자 포용 — 힐리가한을 죽이지 않고 황궁 근처에 거주',
         {'font_size': 16, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('다른 추장들은 당의 장군으로 받아들임',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
                '天   可   汗',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.5),
                '천가한 — "하늘의 가한(可汗·황제)"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.5),
                '북방·서역의 여러 부족이 태종에게 바친 존칭',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
                '— "정복 후 포용"이 만든 동아시아 최초의 다민족 평화 체제',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('XⅤ. 군사·외교')
def s_munseong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 군사·외교', page, total)
    add_title(slide, '문성공주(文成公主)의 토번 화친 (641년)',
              '— 전쟁 없는 외교의 모범')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('토번(吐蕃·티베트)의 송첸감포(松贊干布)가 결혼 요청',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('태종이 종실의 딸 문성공주를 파견',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('→ 전쟁 없이 수백 년 평화 —',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('당 문화가 토번으로 전파되는 계기',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('현대 적용 — 소프트 파워(Soft Power)의 원조',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('무력이 아니라 문화·교류·동맹으로 영향력을 확장',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('전쟁 비용보다 외교 비용이 훨씬 싸다는 통찰',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('XⅤ. 군사·외교')
def s_goguryeo_fail(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 군사·외교', page, total)
    add_title(slide, '고구려 원정 실패(645) — "위징이 있었다면…"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('정관 19년(645) — 태종이 직접 고구려 원정',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('안시성(安市城) 공략 실패 → 철군',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('직간자 위징은 이미 2년 전에 사망 (643년)',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.2), [
        ('철군하며 태종이 한탄:',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"위징이 살아 있었다면',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 나를 이런 원정으로 보내지 않았을 것이다"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('→ 명군도 자만에 빠진다. 직간자 부재의 위험을 자기 경험으로 입증',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


# ============== XⅥ. 신종 ==============
@S('XⅥ. 신종')
def s_daily(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 신종', page, total)
    add_title(slide, '권10 — 군주의 일상 절제')
    items = [
        ('37 행행(行幸)', '"임금 한 번 움직이면 10만 가구가 고달프다"',
         '행차·순행 최소화 — 통치자의 발걸음 자체가 사회적 비용'),
        ('38 전렵(畋獵)', '취미는 재앙의 시작',
         '사냥도 삼면 몰이(三驅)의 옛 예를 지켜 과하지 않게'),
        ('39 재상(災祥)', '"상서를 좋아하는 임금은 재이가 온다"',
         '가뭄·홍수를 하늘의 경고로 받아 자기 반성의 계기로'),
        ('40 신종(愼終)', '정관정요의 대미(大尾)',
         '위징의 「신종소」 — 십점불극종(十漸不克終)의 경고'),
    ]
    top = 2.3
    for tag, quote, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.8), Inches(1.05), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.3), Inches(2.8), Inches(0.5),
                    tag, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.6), Inches(top), Inches(9.2), Inches(1.05), PALE)
        add_textbox(slide, Inches(3.75), Inches(top + 0.1), Inches(8.9), Inches(0.45),
                    quote, font_size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(3.75), Inches(top + 0.55), Inches(8.9), Inches(0.45),
                    desc, font_size=13, color=SUB)
        top += 1.2


@S('XⅥ. 신종')
def s_sinjongso_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 신종', page, total)
    add_title(slide, '「신종소(愼終疏)」 ① — 십점불극종 1~5',
              '정관 13년경 위징이 올린 가장 강력한 상소 — 열 가지 점진적 흔들림')
    items = [
        ('1', '검소 → 사치',     '궁궐 증축·토목 욕망 — 사치의 시작'),
        ('2', '간언 환영 → 회피', '귀에 거슬리는 말을 피함 — 정보 차단'),
        ('3', '근면 → 게으름',    '새벽 조회 → 사냥장에서 시간 소비'),
        ('4', '신뢰 → 의심',      '신하에게 권한 위임 → 작은 일까지 개입'),
        ('5', '인의 → 처벌',      '인의로 다스림 → 처벌을 즐기는 경향'),
    ]
    top = 2.5
    for num, change, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.75), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.18), Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(4.8), Inches(0.75), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.18), Inches(4.6), Inches(0.5),
                    change, font_size=16, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.3), Inches(top), Inches(6.5), Inches(0.75),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.5), Inches(top + 0.18), Inches(6.2), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 0.85


@S('XⅥ. 신종')
def s_sinjongso_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 신종', page, total)
    add_title(slide, '「신종소」 ② — 십점불극종 6~10')
    items = [
        ('6',  '욕망 제어 → 합리화',    '"이 정도는 괜찮다" — 자기 기준 완화'),
        ('7',  '능력 판단 → 편애',      '좋아하는 사람을 편애'),
        ('8',  '공사 분리 → 혼동',      '개인적 호오가 정책에 영향'),
        ('9',  '신하 수고 감사 → 당연시', '겸양 상실'),
        ('10', '백성 두려움 → 자만',     '"나의 덕으로 다스려진다"의 자만'),
    ]
    top = 2.4
    for num, change, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.75), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.18), Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(4.8), Inches(0.75), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.18), Inches(4.6), Inches(0.5),
                    change, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.3), Inches(top), Inches(6.5), Inches(0.75),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.5), Inches(top + 0.18), Inches(6.2), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 0.85
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                '태종의 반응: "위징이 없었다면 나는 이 변화를 보지 못했을 것이다"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('XⅥ. 신종')
def s_sinjong_yeosi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 신종', page, total)
    add_title(slide, '愼終如始 — 정관정요의 최종 결론')
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
                '愼 終 如 始',
                font_size=140, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.5),
                '신종여시',
                font_size=24, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.55), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.7),
                '"처음처럼 끝까지 삼가면, 실패할 일이 없다(則無敗事)"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.85), Inches(11.9), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(6.0), Inches(11.3), Inches(1.1), [
        ('태종은 이 상소를 받고 "매일 보며 자신의 흔들림을 점검하겠다"고 공개 다짐',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('40편 460여 일화가 결국 이 한 줄로 수렴 — 정관정요의 압축된 결론',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== XⅦ. 명구절 ==============
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=44):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=22, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('XⅦ. 명구절 (1/10)',
    '先 正 其 身   然 後 能 正 人',
    '선정기신 연후능정인',
    '먼저 제 몸을 바르게 한 뒤에야 남을 바르게 할 수 있다',
    '군도편 제1 (위징의 답)', hanmun_size=36), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (2/10)',
    '君 者 舟 也   庶 人 者 水 也\n水 則 載 舟   水 則 覆 舟',
    '군자주야 서인자수야 · 수즉재주 수즉복주',
    '임금은 배요 백성은 물 — 물은 배를 띄우기도 뒤집기도 한다',
    '정체편 제2 (순자 인용)', hanmun_size=26), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (3/10)',
    '兼 聽 則 明   偏 信 則 暗',
    '겸청즉명 편신즉암',
    '두루 들으면 밝고, 한쪽만 믿으면 어둡다',
    '정체편 제2 (위징의 답)', hanmun_size=48), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (4/10)',
    '以 銅 爲 鑒   可 以 正 衣 冠\n以 古 爲 鑒   可 以 知 興 替\n以 人 爲 鑒   可 以 明 得 失',
    '이동위감 가이정의관 · 이고위감 가이지흥체 · 이인위감 가이명득실',
    '구리로 거울 삼으면 의관, 옛일로 거울 삼으면 흥망, 사람으로 거울 삼으면 득실을 안다',
    '구간편 — 위징 추도 (643년)', hanmun_size=22), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (5/10)',
    '為 政 之 要   唯 在 得 人',
    '위정지요 유재득인',
    '정치의 요체는 오직 사람을 얻는 데 있다',
    '택관편 제7', hanmun_size=42), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (6/10)',
    '民 無 信 不 立',
    '민무신불립',
    '백성에게 신의(信)가 없으면 나라가 서지 않는다',
    '성신편 제17 (논어 인용)', hanmun_size=70), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (7/10)',
    '兵 者 凶 器   不 得 已 用 之',
    '병자흉기 부득이용지',
    '병(兵)은 흉기(凶器)라 부득이할 때만 쓴다',
    '정벌편 제35', hanmun_size=40), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (8/10)',
    '愼 終 如 始   則 無 敗 事',
    '신종여시 즉무패사',
    '처음처럼 끝까지 삼가면 실패할 일이 없다',
    '신종편 제40 (위징 「신종소」)', hanmun_size=42), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (9/10)',
    '魏 徵 沒   朕 亡 一 鑒 矣',
    '위징몰 짐망일감의',
    '위징이 죽으니, 짐이 거울 하나를 잃었다',
    '구간편 — 태종의 추도 (643년)', hanmun_size=42), 'XⅦ. 명구절'))

SLIDES.append((make_quote_slide('XⅦ. 명구절 (10/10)',
    '創 業 難   守 成 亦 難',
    '창업난 수성역난',
    '창업도 어렵지만, 수성(守成)도 어렵다',
    '정체편 제2', hanmun_size=52), 'XⅦ. 명구절'))


# ============== XⅧ. 후대 ==============
@S('XⅧ. 후대 수용')
def s_korea(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 후대 수용', page, total)
    add_title(slide, '한국 — 세종과 정조의 정관정요')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '세종(1397~1450)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(3.0), Inches(5.9), Inches(0.4),
                '— "조선의 정관"', font_size=14, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.6), Inches(5.3), Inches(3.1), [
        ('• 정관정요를 통치의 기본 텍스트로', {'font_size': 14}),
        ('• 경연(經筵) 1,898회 중 핵심 교재', {'font_size': 14, 'space_before': 6}),
        ('• 집현전·한글·과학·음악·법전', {'font_size': 14, 'space_before': 6}),
        ('• 황희·맹사성·신숙주 등', {'font_size': 14, 'space_before': 6}),
        ('  — 위징·방현령·두여회 구도의 재현', {'font_size': 13, 'color': SUB}),
        ('• 정관의 치를 능가하는 성과', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '정조(1752~1800)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.0), Inches(5.9), Inches(0.4),
                '— 정관정요의 학문적 재해석', font_size=14, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.6), Inches(5.3), Inches(3.1), [
        ('• 정관정요 연구에 깊이 투자', {'font_size': 14}),
        ('• 『홍재전서(弘齋全書)』에 분석 수록', {'font_size': 14, 'space_before': 6}),
        ('• 군신감계 정신으로 당쟁 조율', {'font_size': 14, 'space_before': 6}),
        ('  — 노론·소론·남인 다당제', {'font_size': 13, 'color': SUB}),
        ('• 초계문신제(抄啟文臣制) 운영', {'font_size': 14, 'space_before': 6}),
        ('  — 홍문관 18학사의 조선판', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)


@S('XⅧ. 후대 수용')
def s_japan(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 후대 수용', page, total)
    add_title(slide, '일본 — 도쿠가와 이에야스의 수용',
              '"치도(治道)의 스승"으로 숭배')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.7), Inches(2.43), Inches(11.9), Inches(0.5),
                '도쿠가와 이에야스(徳川家康, 1543~1616)',
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.3), Inches(11.5), Inches(3.7), [
        ('• 정관정요를 "치도(治道)의 스승"으로 숭배',
         {'font_size': 17}),
        ('', {'font_size': 8}),
        ('• 에도 시대 300년 평화(1603~1867)의 사상적 기초',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('', {'font_size': 8}),
        ('• 임제종 선사(臨濟禪師) 고젠(狐禪)에게 강론을 청함',
         {'font_size': 17, 'space_before': 10}),
        ('', {'font_size': 8}),
        ('• 막부 관료 교육에 정관정요를 정식 편입',
         {'font_size': 17, 'space_before': 10}),
        ('', {'font_size': 8}),
        ('• "창업수성" 개념을 도쿠가와 가문 승계의 핵심 원칙으로',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ], line_spacing=1.4)


@S('XⅧ. 후대 수용')
def s_vietnam(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 후대 수용', page, total)
    add_title(slide, '베트남·동아시아 한자 문화권 전체로')
    items = [
        ('베트남 리(李) 왕조',     '1009~1225',   '궁정 필독서로 정관정요 채택'),
        ('베트남 쩐(陳) 왕조',     '1225~1400',   '왕세자 교육 교과서로'),
        ('류큐(琉球) 왕국',         '14~19세기',   '책봉 외교와 함께 정관정요 수용'),
        ('동아시아 한자 문화권',    '8세기 이후',  '한국·중국·일본·베트남 공통 경전'),
    ]
    top = 2.5
    for region, era, role in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(4.0), Inches(0.9), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(4.0), Inches(0.4),
                    region, font_size=16, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(4.0), Inches(0.4),
                    era, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.8), Inches(top), Inches(8.0), Inches(0.9), PALE)
        add_textbox(slide, Inches(5.0), Inches(top + 0.27), Inches(7.8), Inches(0.5),
                    role, font_size=15, color=INK)
        top += 1.05
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '1,300년 동안 동아시아 군주들이 공유한 단 하나의 정치 교과서',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('XⅧ. 후대 수용')
def s_modern_biz(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 후대 수용', page, total)
    add_title(slide, '현대 기업 경영과의 접점')
    rows = [
        ('수성(守成)의 어려움', 'Innovator\'s Dilemma — 성공한 기업의 혁신 실패'),
        ('세 거울(三鏡)',        '360도 피드백 — 다각도 평가'),
        ('임현(任賢)',           '강점 기반 팀 빌딩 (Gallup)'),
        ('구간·납간',            '심리적 안전감 — Google Aristotle 연구'),
        ('창업수성',             '블리츠스케일링 vs 장기 경영의 균형'),
        ('권6 9병폐',            'CEO 성과 평가 체크리스트'),
        ('삼복주',               '다중 검증 시스템 (Checks & Balances)'),
        ('겸청즉명',             '다양성 경영(Diversity Management)'),
        ('신종여시',             '아마존 베조스 "Day 1" 정신'),
    ]
    top = 2.0
    for i, (concept, modern) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(4.5), Inches(0.5), PALE)
        add_filled_rect(slide, Inches(5.05), Inches(top), Inches(7.8), Inches(0.5), bg)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(4.4), Inches(0.4),
                    concept, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(5.2), Inches(top + 0.1), Inches(7.6), Inches(0.4),
                    modern, font_size=13, color=INK)
        top += 0.55


# ============== XⅨ. 비교 ==============
@S('XⅨ. 비교')
def s_compare(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅨ. 비교', page, total)
    add_title(slide, '정관정요 · 한비자 · 군주론(마키아벨리) 비교')
    rows = [
        ('인간관',     '덕 + 이기심의 혼재',      '이기심 중심',           '이기적 본성'),
        ('통치 기반',  '덕치(德治) + 제도',        '법(法)·술(術)·세(勢)',  '권력·공포'),
        ('군주상',     '덕 있는 조율자',          '엄밀한 시스템 설계자',   '결단의 권력자'),
        ('신하상',     '거울 같은 동반자',         '경계 대상',              '도구'),
        ('변화관',     '초심 유지(愼終)',         '현실 적응',              '단기 실용'),
        ('시점',       '안정기 치세',             '혼란기 통제',            '혼란기 권력 획득'),
        ('대표 덕목',  '신종여시(愼終如始)',      '엄밀한 법치',            'virtù(역량)'),
    ]
    top = 1.95
    headers = ['항목', '정관정요', '한비자', '군주론']
    widths = [2.0, 3.6, 3.6, 3.2]
    x = 0.5
    for i, (h, w) in enumerate(zip(headers, widths)):
        color = SUB if i == 0 else (ACCENT if i == 1 else SUB)
        add_filled_rect(slide, Inches(x), Inches(top), Inches(w), Inches(0.55), color)
        add_textbox(slide, Inches(x), Inches(top + 0.1), Inches(w), Inches(0.4),
                    h, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w + 0.05
    row_h = 0.65
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE
        x = 0.5
        add_filled_rect(slide, Inches(x), Inches(y), Inches(widths[0]), Inches(row_h), PALE)
        add_textbox(slide, Inches(x + 0.05), Inches(y + 0.15), Inches(widths[0] - 0.1),
                    Inches(0.4), row[0], font_size=13, bold=True, color=INK,
                    align=PP_ALIGN.CENTER)
        x += widths[0] + 0.05
        for i in range(3):
            add_filled_rect(slide, Inches(x), Inches(y), Inches(widths[i + 1]),
                            Inches(row_h), bg)
            color = ACCENT if i == 0 else INK
            add_textbox(slide, Inches(x + 0.1), Inches(y + 0.15),
                        Inches(widths[i + 1] - 0.2), Inches(0.4),
                        row[i + 1], font_size=13, color=color,
                        bold=(i == 0), align=PP_ALIGN.CENTER)
            x += widths[i + 1] + 0.05


# ============== XX. 마무리 ==============
@S('XX. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XX. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 정관정요')
    add_filled_rect(slide, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.9), PALE)
    add_paragraphs(slide, Inches(1.1), Inches(2.3), Inches(11.1), Inches(4.6), [
        ('현무문의 변으로 형을 죽이고 즉위한 당 태종이',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('수 양제의 멸망을 반면교사 삼아',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('위징·방현령·두여회·장손황후와 함께',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('정관 23년간 만들어 낸',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"정관의 치(治)"의 모든 비결을',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('사관 오긍이 10권 40편으로 정리한',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('동아시아 제왕학의 절대 표준',
         {'font_size': 24, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('', {'font_size': 4}),
        ('— 자기 욕망을 다스리는 리더, 반대 의견을 받아내는 조직 —',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('"끝을 삼가는 용기"의 23년 기록',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.2)


@S('XX. 마무리')
def s_last_words(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XX. 마무리', page, total)
    add_title(slide, '태종의 임종 유언', '— 제왕의 어려움은 시작이 아니라 끝에 있다')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.55), Inches(11.7), Inches(4.2), [
        ('"제왕의 어려움은 시작이 아니라 끝에 있다',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        (' 나는 후반에 이르러 많은 실수를 했다"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"너는 내 전반기를 본받고',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 12}),
        (' 후반기는 거울로 삼으라"',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('— 태종의 마지막 가르침 (649년, 향년 52세)',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 12}),
        ('', {'font_size': 8}),
        ('자기 실패까지 정직하게 물려준 명군의 마지막 말',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.25)


@S('XX. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.5),
                '愼 終 如 始',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '신종여시', font_size=24, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
                '"처음처럼 끝까지 삼가면, 실패할 일이 없다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
                '— 신종편 제40, 위징 「신종소」',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '정관 23년의 비결, 그리고 1,300년 동안 이어진 한 줄의 가르침',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
                '감사합니다', font_size=26, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\정관정요_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
