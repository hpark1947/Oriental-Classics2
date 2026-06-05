# -*- coding: utf-8 -*-
"""
안씨가훈(顔氏家訓) 발표자료 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '동아시아 가훈서의 시조 · 4왕조를 살아낸 안지추의 절실한 증언',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '顔 氏 家 訓',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '안 씨 가 훈',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '不 教 而 愛   非 愛 也',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '사랑한다고 해서 가르치지 않는 것은 — 사랑이 아니다 — 교자편',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '안지추(顔之推, 531~591경) · 수(隋) 개황 연간 · 7권 20편 · 1,400년 가훈서의 모범',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 안씨가훈은 어떤 책인가'),
        ('Ⅱ.', '저자 안지추 — 격동기의 산증인'),
        ('Ⅲ.', '편찬 동기와 시대 배경'),
        ('Ⅳ.', '7권 20편의 5단계 구조'),
        ('Ⅴ.', '핵심 편 12 깊이 읽기'),
    ]
    items_right = [
        ('Ⅵ.', '책 전체를 관통하는 5대 사상'),
        ('Ⅶ.', '명구 12선'),
        ('Ⅷ.', '한국과 동아시아 수용사'),
        ('Ⅸ.', '오늘 다시 펼치는 이유'),
        ('Ⅹ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 안씨가훈',
              '동아시아 가훈서의 시조 — 古今家訓之祖')
    rows = [
        ('서명', '顔氏家訓 — 안씨 가문의 가르침'),
        ('저자', '안지추(顔之推, 531~591경) — 4왕조를 섬긴 학자'),
        ('편찬', '수(隋) 개황 연간(581~600년경), 6세기 말 완성'),
        ('분량', '7권 20편 — 가정·교육·학문·처세·종교·죽음까지'),
        ('성격', '체험에서 우러난 절실한 증언 — 추상이 아닌 사례'),
        ('위상', '古今家訓之祖 — 1,400년 가훈서의 시조이자 정점'),
        ('정신', '실용 합리주의 + 유·불·도 균형'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.8), Inches(0.45), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.8), Inches(0.45),
                    k, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.7), y, Inches(10.3), Inches(0.45),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_first(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '「가훈(家訓)」이라는 단어 자체의 출처',
              '안씨가훈이 만든 보통명사')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 오늘 우리가 쓰는 「**가훈**」이라는 단어 — 이 책에서 시작', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   이 책 이전에도 자식에게 남긴 글은 있었으나', {'font_size': 14, 'color': SUB}),
        ('   20편 7권으로 체계화된 본격 가훈서는 — 이 책이 처음', {'font_size': 14, 'color': SUB}),
        ('● 명·청대 학자들의 일치된 평가 — 「**古今家訓之祖**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('   고금 가훈의 비조(鼻祖) — 모든 후대 가훈서의 모범', {'font_size': 14, 'color': SUB}),
        ('● 한국 율곡의 「**격몽요결**」, 정약용의 「**가계(家誡)**」도 영향', {'font_size': 17, 'space_before': 12}),
        ('● 1,400년 동안 — 동아시아 가훈서의 표준 형식', {'font_size': 16, 'color': SUB, 'space_before': 12}),
    ])


@S(SEC1)
def i_essence(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '안씨가훈의 정수 — 한 마디로',
              '천하를 가르치려는 책이 아니라 — 내 자손에게')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.6), Inches(1.8), [
        ('내가 이 책을 짓는 것은 세상을 바로잡기 위함이 아니다.', {'font_size': 17, 'color': INK}),
        ('다만 내 집안을 가지런히 하고,', {'font_size': 17, 'color': INK}),
        ('자손들을 이끌고자 함이다.', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('— 서치편(序致篇)', {'font_size': 13, 'color': SUB, 'space_before': 8}),
    ])
    add_paragraphs(slide, Inches(0.7), Inches(4.7), Inches(12.0), Inches(2.5), [
        ('● 천하를 가르치는 거대한 책이 아니라', {'font_size': 16, 'space_before': 6}),
        ('● 내 자손에게 절실하게 말하는 책', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('● 이 절실함이 — 1,400년 사랑받은 비결', {'font_size': 16, 'color': SUB, 'space_before': 6}),
        ('● 「**진정한 보편은 — 절실한 개별에서 나온다**」', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅱ. 저자 안지추 ==============
SEC2 = 'Ⅱ. 저자 안지추'

@S(SEC2)
def ii_life(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '안지추(顔之推, 531~591경)',
              '4왕조를 섬긴 격동기의 산증인')
    rows = [
        ('531년경', '남조 양(梁) 강릉(江陵, 호북성)에서 출생 — 명문 안씨 가문'),
        ('552년경', '양 원제(元帝)의 막료로 출사'),
        ('554년', '서위(西魏)가 강릉 함락 — 양 원제 살해 — 망국 체험'),
        ('555년', '북주(北周)에 포로처럼 끌려감'),
        ('556년', '북제(北齊)로 망명 — 황문시랑 등 요직'),
        ('577년', '북제 멸망 — 북주에 귀속'),
        ('581년', '수(隋) 건국 — 4번째 왕조를 섬김'),
        ('591년경', '수에서 학사(學士)로 활동하다 사망'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.52)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.5), Inches(0.42), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.42),
                    k, font_size=13, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.4), y, Inches(10.5), Inches(0.42),
                    v, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC2)
def ii_witness(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '4왕조를 섬긴 한 인간',
              '양 → 북제 → 북주 → 수')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 자기 의지가 아닌 — 격동의 시대의 운명', {'font_size': 17, 'space_before': 6}),
        ('● 직접 겪은 일들', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   고향 강릉의 함락과 약탈(554)', {'font_size': 14, 'color': SUB}),
        ('   포로로 북방에 끌려가는 굴욕', {'font_size': 14, 'color': SUB}),
        ('   남방 한족과 북방 호족 문화 충돌을 양쪽에서 모두 경험', {'font_size': 14, 'color': SUB}),
        ('   남조 귀족 사회의 타락과 북조의 거친 풍토를 비교', {'font_size': 14, 'color': SUB}),
        ('● 이 모든 격동의 체험이 안씨가훈의 모든 줄에 배어 있음', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('● 추상이 아니라 — 살아 있는 증언', {'font_size': 16, 'color': ACCENT, 'space_before': 6}),
    ])


@S(SEC2)
def ii_thought(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '안지추의 사상',
              '유가 토대 + 불교 신앙 + 도가 비판')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 유가(儒家)를 근본으로 — 가정 윤리·정치·교육', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('● 불교(佛敎)를 깊이 신앙 — 「귀심편(歸心篇)」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   6조시대 사대부의 전형 — 「유교는 살아 있는 동안의 도리, 불교는 영혼의 평안」', {'font_size': 14, 'color': SUB}),
        ('● 도가의 현묘한 청담(淸談)·황로술 비판', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   양생은 받아들이되 — 신선술에 빠지지는 말라', {'font_size': 14, 'color': SUB}),
        ('● 「**실용적 합리주의**」 — 빈말이 아닌 실제로 도움 되는 가르침', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
    ])


# ============== Ⅲ. 편찬 동기 ==============
SEC3 = 'Ⅲ. 편찬 동기'

@S(SEC3)
def iii_motive(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '왜 이 책을 썼는가',
              '서치편(序致篇)이 밝히는 동기')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 격동기에 가문이 흩어지면 가르침이 끊긴다 — 입에서 입으로 전하기 어렵다', {'font_size': 17, 'space_before': 6}),
        ('● 자기 평생의 시행착오를 — 자손이 반복하지 않도록', {'font_size': 17, 'space_before': 10}),
        ('● 양 귀족 사회의 타락을 직접 보았기에 — 자손이 같은 길을 안 가게', {'font_size': 17, 'space_before': 10}),
        ('● 사대부의 가풍이 흔들리는 시대에 — 전통의 핵심을 압축해 남김', {'font_size': 17, 'space_before': 10}),
        ('● 「**천하를 가르치는 책이 아니라, 내 자손에게 절실히 말하는 책**」', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅳ. 20편 구조 ==============
SEC4 = 'Ⅳ. 20편 구조'

@S(SEC4)
def iv_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '7권 20편의 5단계',
              '요람에서 무덤까지')
    rows = [
        ('1단계 (1~5)', '가정 윤리 — 서치·교자·형제·후취·치가'),
        ('2단계 (6~7)', '사회 진출 — 풍조·모현'),
        ('3단계 (8~11)', '학문과 직업 — 면학·문장·명실·섭무'),
        ('4단계 (12~16)', '처세와 양생 — 성사·지족·계병·양생·귀심'),
        ('5단계 (17~20)', '학문 심화와 죽음 — 서증·음사·잡예·종제'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.6 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.8), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.8), Inches(0.55),
                    k, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.7), y, Inches(9.2), Inches(0.55),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5),
                '한 인간의 전 생애를 망라하는 구조 — from cradle to grave',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_list(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '20편 한 폭으로',
              '각 편의 핵심 주제')
    left = [
        '1. 서치(序致) — 저술 동기',
        '2. 교자(敎子) — 자식 교육',
        '3. 형제(兄弟) — 형제 우애',
        '4. 후취(後娶) — 재혼·계모',
        '5. 치가(治家) — 가정 운영',
        '6. 풍조(風操) — 예절·관습',
        '7. 모현(慕賢) — 어진 이를',
        '8. 면학(勉學) — 학문 (백미)',
        '9. 문장(文章) — 글쓰기',
        '10. 명실(名實) — 명성과 실질',
    ]
    right = [
        '11. 섭무(涉務) — 실무 능력',
        '12. 성사(省事) — 일을 줄임',
        '13. 지족(止足) — 만족할 줄 앎',
        '14. 계병(誡兵) — 무력 경계',
        '15. 양생(養生) — 건강',
        '16. 귀심(歸心) — 불교 신앙',
        '17. 서증(書證) — 고증',
        '18. 음사(音辭) — 음운학',
        '19. 잡예(雜藝) — 다양한 기예',
        '20. 종제(終制) — 죽음·장례',
    ]
    for i, txt in enumerate(left):
        add_textbox(slide, Inches(0.7), Inches(2.4 + i * 0.42), Inches(6.0), Inches(0.4),
                    txt, font_size=14, color=INK)
    for i, txt in enumerate(right):
        add_textbox(slide, Inches(7.0), Inches(2.4 + i * 0.42), Inches(6.0), Inches(0.4),
                    txt, font_size=14, color=INK)


# ============== Ⅴ. 핵심 편 12 깊이 읽기 ==============
SEC5 = 'Ⅴ. 핵심 편 12 깊이 읽기'

def make_chapter_slide(idx_total, chapter_name, hanja_quote, korean_quote, lesson):
    @S(SEC5)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, SEC5, n, t)
        add_title(slide, f'{idx_total} — {chapter_name}', '안지추의 절실한 한 마디')
        add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.2), PALE)
        add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.2),
                    hanja_quote,
                    font_size=22, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.8),
                    korean_quote,
                    font_size=16, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.8), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.8),
                    lesson,
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


chapters = [
    ('1/12', '교자편(敎子篇) — 자식 교육', '不 教 而 愛   非 愛 也',
     '사랑한다고 해서 가르치지 않는 것은 — 사랑이 아니다',
     '교육은 태어날 때부터(태교) — 어렸을 때 사랑하면서도 위엄을 잃지 말 것. 친애(親愛)와 엄숙(嚴肅)의 균형. 너무 사랑해서 잘못을 못 보고 넘기면 자식을 망친다.'),
    ('2/12', '형제편(兄弟篇) — 형제 우애', '兄 弟 者   父 母 之 遺 體',
     '형제는 — 부모가 남긴 같은 몸',
     '부모는 한 분, 형제는 부모가 남긴 같은 핏줄. 부모 사후에는 형제가 곧 부모의 연장. 형제 사이가 무너지면 가문이 무너진다.'),
    ('3/12', '치가편(治家篇) — 가정 운영', '儉 而 不 嗇   而 亦 不 奢',
     '검소하되 인색하지 않고, 그렇다고 사치하지도 않는다',
     '가정의 근본은 검소(儉素). 너무 인색해도 안 되고 사치해도 안 됨 — 중도(中道). 노복(奴僕)을 다루는 법 — 인간적 존엄을 잊지 말 것.'),
    ('4/12', '풍조편(風操篇) — 예절', '禮 之 大 體   出 入 應 對',
     '예의 큰 체(體)는 — 들고 남과 응대에 있다',
     '남북의 다른 풍습 비교 — 안지추가 양쪽을 다 살았기에 가능. 휘(諱)의 예법, 결혼·상장의 예절. 「예절은 작은 것 같지만 인격의 외피」.'),
    ('5/12', '모현편(慕賢篇) — 어진 이를 본받음', '見 賢 而 慕   人 之 大 道',
     '어진 이를 보고 흠모하는 것이 사람의 큰 길',
     '가까이 있는 사람은 평범해 보여도 — 시간이 지나면 그 가치가 드러남. 어진 이를 만나면 반드시 본받으라. 비방하는 말을 함부로 듣지 말라 — 시기심에서 나오는 경우가 많다.'),
    ('6/12', '면학편(勉學篇) — 학문 (이 책의 백미)', '人 生 在 世   會 當 有 業',
     '사람이 세상에 났다면 — 마땅히 학업이 있어야 한다',
     '학문하지 않은 자의 비참한 말로(末路). 격동기에 학문은 유일하게 빼앗기지 않는 자산. 「**늙어서라도 학문을 시작하라 — 어두운 밤에 등불을 든 것 같다**」. 학문은 위인지학(爲人之學)이 아닌 위기지학(爲己之學).'),
    ('7/12', '문장편(文章篇) — 글쓰기', '文 章 當 以 理 致 為 心 腎',
     '글은 — 이치의 도달을 심장과 콩팥으로 삼아야',
     '화려한 수사보다 내실(內實)이 중요. 옛사람의 좋은 문장을 익혀 자기 것으로. 「**글은 그 사람이다 — 글쓰기는 인격 수양의 일부**」.'),
    ('8/12', '명실편(名實篇) — 명성과 실질', '名 之 與 實   猶 形 之 與 影',
     '명성과 실질은 — 형체와 그림자와 같다',
     '명성을 좇지 말라 — 명성은 결과로 따라오는 것. 헛된 평판에 휘둘리는 자의 어리석음. 진짜 실력을 갖춘 자는 명성을 따라가지 않아도 따라온다. 「**명(名)은 실(實)의 그림자**」.'),
    ('9/12', '섭무편(涉務篇) — 실무 능력', '士 君 子   不 可 不 涉 務',
     '사군자(士君子)는 — 실무에 관여하지 않을 수 없다',
     '책상물림 학자가 아니라 실제 일을 처리할 줄 아는 사람이 되라. 6조 사대부들이 청담에 빠져 실무를 모른다고 비판. 농사·의술·법률·산수 등 실용 학문의 가치. 「**고담만 하고 일을 못하면 위기 때 한 치도 못 움직인다**」.'),
    ('10/12', '지족편(止足篇) — 만족할 줄 앎', '少 欲 知 足   是 名 富 貴',
     '욕심을 적게 하고 만족할 줄 앎 — 이것이 부귀의 이름',
     '욕심에는 끝이 없다 — 욕심이 곧 화의 근원. 자기 분수를 알고 만족하면 위태롭지 않다. 노자의 「**지족불욕(知足不辱)**」과 공명. 「**만족은 가장 큰 부유함**」.'),
    ('11/12', '귀심편(歸心篇) — 불교 신앙', '佛 家 之 教   亦 是 善 道',
     '불가의 가르침도 — 선한 도이다',
     '안지추 자신의 불교 신앙 고백. 인과응보·윤회의 가르침. 유교만으로 설명되지 않는 인생의 깊은 차원에 대한 답을 불교에서. 6조 사대부의 종교 의식의 전형. 「**유교는 살아 있는 동안의 도리, 불교는 영혼의 평안**」.'),
    ('12/12', '종제편(終制篇) — 죽음과 장례', '吾 死 之 後   葬 以 葦 席',
     '내가 죽은 후에 — 갈자리로 장사 지내라',
     '자신의 유언과 장례 지침 — 검소한 장례를 명시. 화려한 묘·후한 부장품 거부. 죽음을 자연스럽게 받아들이는 자세. 「**죽음 또한 가르침이다 — 검소하고 담담하게**」.'),
]

for tag, name, hj, kr, ls in chapters:
    make_chapter_slide(tag, name, hj, kr, ls)


# ============== Ⅵ. 5대 사상 ==============
SEC6 = 'Ⅵ. 5대 사상'

@S(SEC6)
def vi_experience(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 1 — 체험에서 우러난 절실함',
              '추상이 아닌 사례의 가르침')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 모든 가르침이 — 안지추 자신이 직접 보고 겪고 깨달은 것에서', {'font_size': 17, 'space_before': 6}),
        ('● 추상적 도덕 강의가 아니라', {'font_size': 17, 'space_before': 10}),
        ('● 「**내가 이런 사람을 보았는데 그 결과는 이러했다**」 — 구체적 사례', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 1,400년 동안 사랑받은 비결 — 절실함의 힘', {'font_size': 17, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC6)
def vi_balance(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 2 — 유·불·도의 균형',
              '6세기 사대부 정신의 종합')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 유교를 근본으로 — 윤리·정치·교육', {'font_size': 17, 'space_before': 6}),
        ('● 불교의 인과응보·내세관으로 — 죽음과 영혼', {'font_size': 17, 'space_before': 10}),
        ('● 도가의 양생술은 부분적 수용하되 — 황로술·신선술은 비판', {'font_size': 17, 'space_before': 10}),
        ('● 「**유교는 살아 있는 동안, 불교는 영혼의 평안, 도가는 신체 보양**」', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 12}),
    ])


@S(SEC6)
def vi_home(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 3 — 가정에서 시작해 가정으로 끝',
              '인생은 결국 한 가정의 이야기')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 첫 편(서치) — 자기 가문의 내력으로 시작', {'font_size': 17, 'space_before': 6}),
        ('● 마지막 편(종제) — 자기 장례를 명시하며 끝', {'font_size': 17, 'space_before': 10}),
        ('● 모든 인생이 — 한 가정에서 출발해 한 가정으로 돌아간다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 책의 형식 자체가 — 그 진리를 보여준다', {'font_size': 17, 'space_before': 10}),
    ])


@S(SEC6)
def vi_survival(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 4 — 격동기의 생존 매뉴얼',
              '4왕조의 격변에서 살아남는 법')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 표면은 가훈서 — 본질은 격동기 생존 매뉴얼', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('● 무엇을 배우고(면학)', {'font_size': 17, 'space_before': 8}),
        ('● 무엇을 줄이고(성사)', {'font_size': 17, 'space_before': 6}),
        ('● 무엇에 만족하고(지족)', {'font_size': 17, 'space_before': 6}),
        ('● 무엇을 피하라(계병)', {'font_size': 17, 'space_before': 6}),
        ('● 안지추가 격변 속에서 살아남으며 깨달은 모든 것의 압축', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC6)
def vi_middle(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 5 — 중도(中道)와 균형',
              '모든 분야에서 양극단을 피하라')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 너무 사랑하지도, 너무 엄하지도 말라 (교자)', {'font_size': 17, 'space_before': 6}),
        ('● 너무 인색하지도, 사치하지도 말라 (치가)', {'font_size': 17, 'space_before': 8}),
        ('● 양생을 하되 신선술에 빠지지 말라 (양생)', {'font_size': 17, 'space_before': 8}),
        ('● 학문을 하되 청담에 빠지지 말라 (섭무)', {'font_size': 17, 'space_before': 8}),
        ('● 「**모든 분야에서 — 양극단의 어느 쪽으로도 기울지 말라**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
    ])


# ============== Ⅶ. 명구 12선 ==============
SEC7 = 'Ⅶ. 명구 12선'

def make_quote(idx_total, hanja, korean, comment):
    @S(SEC7)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, SEC7, n, t)
        add_title(slide, f'명구 {idx_total}', '안씨가훈 — 한 마디의 절실함')
        add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.2), PALE)
        add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.2),
                    hanja,
                    font_size=24, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(1.0),
                    korean,
                    font_size=18, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.8), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.8),
                    comment,
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


quotes = [
    ('1/12', '不 教 而 愛   非 愛 也',
     '가르치지 않고 사랑하는 것 — 사랑이 아니다',
     '교자편. 자식을 망치는 사랑 — 잘못을 못 보고 넘기는 것은 자식을 죽이는 것.'),
    ('2/12', '幼 而 學 者   如 日 出 之 光',
     '어려서 배운 자는 — 해 뜨는 빛과 같다',
     '면학편. 어린 시절의 학문이 — 평생을 비추는 빛.'),
    ('3/12', '老 而 學 者   如 秉 燭 夜 行',
     '늙어서 배운 자도 — 어두운 밤에 등불을 들고 가는 것과 같다',
     '면학편. 학문에 늦은 때는 없다 — 평생 학습의 정수.'),
    ('4/12', '少 欲 知 足   是 名 富 貴',
     '욕심을 적게 하고 만족할 줄 앎 — 이것이 부귀',
     '지족편. 노자 「지족불욕」과 공명. 만족이 곧 부유함.'),
    ('5/12', '名 之 與 實   猶 形 之 與 影',
     '명성과 실질은 — 형체와 그림자와 같다',
     '명실편. 명성은 결과로 따라온다 — 그림자를 좇지 말라.'),
    ('6/12', '兄 弟 者   父 母 之 遺 體',
     '형제는 — 부모가 남긴 같은 몸',
     '형제편. 부모 사후 형제는 곧 부모의 연장 — 우애의 정수.'),
    ('7/12', '一 家 之 親   生 於 父 母',
     '한 집안의 친애 — 부모에게서 생긴다',
     '서치편. 모든 친애의 출발은 부모.'),
    ('8/12', '婚 姻 之 道   合 二 姓 之 好',
     '혼인의 도 — 두 성씨의 우호를 합한다',
     '풍조편. 혼인은 개인이 아닌 가문의 결합.'),
    ('9/12', '人 生 在 世   會 當 有 業',
     '사람이 세상에 났다면 — 마땅히 학업이 있어야',
     '면학편. 격동기에 학문은 유일하게 빼앗기지 않는 자산.'),
    ('10/12', '安 危 在 出 令   存 亡 在 所 任',
     '안위는 명령을 내림에 있고, 존망은 맡기는 데 있다',
     '치가편·섭무편 통합 정신. 위정자의 책임.'),
    ('11/12', '禮 為 教 本',
     '예(禮)는 교육의 근본',
     '풍조편. 예가 작아 보여도 인격의 외피 — 작은 예의에서 큰 도리.'),
    ('12/12', '吾 死 之 後   葬 以 葦 席',
     '내가 죽은 후에 — 갈자리로 장사 지내라',
     '종제편. 검소한 장례 — 죽음 또한 가르침이다.'),
]

for tag, hj, kr, cm in quotes:
    make_quote(tag, hj, kr, cm)


# ============== Ⅷ. 한국·동아시아 수용 ==============
SEC8 = 'Ⅷ. 한국·동아시아 수용'

@S(SEC8)
def viii_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '한국에서의 수용',
              '율곡 격몽요결·정약용 가계의 모범')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 고려·조선 사대부의 가훈 작성 모범', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('● 율곡 이이의 『**격몽요결(擊蒙要訣)**』', {'font_size': 17, 'space_before': 10}),
        ('   어린이 교육서 — 안씨가훈의 정신을 한국적으로', {'font_size': 14, 'color': SUB}),
        ('● 정약용의 『**가계(家誡)**』', {'font_size': 17, 'space_before': 10}),
        ('   유배지에서 자식들에게 보낸 편지 — 가훈서의 한국적 정점', {'font_size': 14, 'color': SUB}),
        ('● 조선 사대부 가문의 가훈·가법(家法) 문화의 원천', {'font_size': 17, 'space_before': 10}),
    ])


@S(SEC8)
def viii_value(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '학술적 가치',
              '단순한 가훈서를 넘어')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 6조시대의 거의 유일한 1차 사료', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   정사(正史)로 보충 불가능한 사대부 일상', {'font_size': 14, 'color': SUB}),
        ('● 음운학·문헌학의 초기 보고 (서증·음사편)', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   중국 음운학의 초기 자료로 학술적 가치 큼', {'font_size': 14, 'color': SUB}),
        ('● 6조 불교의 사대부 수용 양상 자료 (귀심편)', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 일본·베트남에서도 사대부 교양서로 정착', {'font_size': 17, 'space_before': 10}),
    ])


# ============== Ⅸ. 오늘 다시 펼치는 이유 ==============
SEC9 = 'Ⅸ. 오늘 다시 펼치는 이유'

@S(SEC9)
def ix_today1(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '1 — 학문이 가장 마지막에 남는다',
              '면학편의 영원한 가르침')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 부귀·권력·재산은 — 한순간에 사라질 수 있다', {'font_size': 17, 'space_before': 6}),
        ('● 그러나 자기 안에 쌓은 학문과 인격은 — 누구도 빼앗아 갈 수 없다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 격동기일수록 — 이 진실은 더 분명', {'font_size': 17, 'space_before': 10}),
        ('● 오늘 — 「**평생 학습**」이 가장 깊은 보험', {'font_size': 17, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today2(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '2 — 평범한 가르침이 가장 깊다',
              '거창함 없이 정직한 답')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 안씨가훈에는 거창한 형이상학이 없다', {'font_size': 17, 'space_before': 6}),
        ('● 자식 잘 가르치는 법 · 형제 사이 잘 지내는 법 · 글 잘 쓰는 법', {'font_size': 17, 'space_before': 10}),
        ('● 누구나 평생 부딪히는 일들에 대한 — 정직한 답', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 그래서 어떤 시대 어떤 사람에게도 — 살아 있다', {'font_size': 17, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today3(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '3 — 내 오늘 하루가 가문의 운명',
              '책임감의 정수')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 서치편(인생 시작)에서 종제편(죽음)까지', {'font_size': 17, 'space_before': 6}),
        ('● 한 사람의 인생이 — 결국 자기가 속한 가문의 한 챕터', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 나의 행위는 — 내 자손에게 이어진다', {'font_size': 17, 'space_before': 10}),
        ('● 「**나의 오늘 하루가 — 가문의 운명을 결정한다**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 가장 깊은 책임감 — 4왕조를 살아낸 안지추의 깊은 자각', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


# ============== Ⅹ. 마무리 ==============
SEC10 = 'Ⅹ. 마무리'

@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '안씨가훈이 일러주는 7가지',
              '한 폭으로 정리')
    items = [
        '가르치지 않고 사랑하는 것은 사랑이 아니다 — 不教而愛 非愛也',
        '늙어서라도 학문을 시작하라 — 老而學者 如秉燭夜行',
        '욕심을 적게 하고 만족하라 — 少欲知足 是名富貴',
        '명성을 좇지 말고 실질을 닦으라 — 名實之分',
        '실무 능력을 갖춰라 — 청담만 하지 말라',
        '모든 분야에서 중도(中道) — 양극단을 피하라',
        '검소함은 죽음에도 — 葬以葦席',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.5 + i * 0.6)
        add_textbox(slide, Inches(0.9), y, Inches(0.6), Inches(0.5),
                    f'{i+1}.', font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.5),
                    txt, font_size=16, color=INK)


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                '안씨가훈 — 4왕조를 살아낸 자가 자손에게 절실히 남긴 20편',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '不 教 而 愛',
                font_size=140, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
                '가르치지 않고 사랑하는 것은 — 사랑이 아니다',
                font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.8), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
                '체험에서 우러난 절실한 가르침 — 1,400년 가훈서의 시조',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)


# ============== 빌드 ==============
total = len(SLIDES)
for i, (fn, sec) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    fn(slide, i, total)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\안씨가훈.pptx'
prs.save(out_path)
print(f'생성 완료: {out_path}  슬라이드 수: {total}')
