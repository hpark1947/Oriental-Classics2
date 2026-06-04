# -*- coding: utf-8 -*-
"""
논어(論語) 발표자료 — 전면 보강판 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
20편 각 1장 깊이 읽기 · 공자 생애 · 핵심 제자 · 사상 7기둥 · 후대 영향 망라
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '동아시아 2,500년의 정신적 DNA · 사서(四書)의 으뜸',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '論 語',
                font_size=120, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '논 어',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
                '學 而 時 習 之 不 亦 說 乎  — 배우고 때때로 익히면 기쁘지 아니한가',
                font_size=17, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '공자(孔子, BC 551~479)와 제자들의 언행록 · 20편 482장 약 16,000자',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '半 部 論 語 治 天 下 — 논어 반 권으로 천하를 다스린다 (송 조보趙普)',
                font_size=14, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 논어란 무엇인가'),
        ('Ⅱ.', '편찬 — 400년의 완성'),
        ('Ⅲ.', '공자의 생애 — 한 인간의 드라마'),
        ('Ⅳ.', '공문사과와 핵심 제자'),
        ('Ⅴ.', '20편의 구조'),
        ('Ⅵ.', '20편 각 편 깊이 읽기'),
    ]
    items_right = [
        ('Ⅶ.', '핵심 사상 7기둥'),
        ('Ⅷ.', '명구 16선'),
        ('Ⅸ.', '한국·일본·서양 수용사'),
        ('Ⅹ.', '오늘 우리에게'),
        ('Ⅺ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 논어')
    rows = [
        ('서명',  '논어(論語) — 「토론하여 정리한 말씀」 (論 토론 + 語 말씀)'),
        ('저자',  '공자와 제자들 — 공자 사후 제자들이 수 차례 편찬'),
        ('시대',  '춘추시대 말기 (BC 6~5세기)'),
        ('분량',  '20편 482장 · 약 16,000자'),
        ('성격',  '공자와 제자들의 언행록(語錄) · 단편적 어록 모음'),
        ('위상',  '사서(四書)의 으뜸 — 논어 · 맹자 · 대학 · 중용'),
        ('영향',  '동아시아 2,500년의 정신적 DNA · 한국·일본·베트남의 교양 기초'),
        ('상징어', '반부논어치천하(半部論語治天下) — 송 조보(趙普)의 말'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_why(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '왜 논어인가 — 동양 경전의 원점',
              '반부논어치천하 — 논어 반 권으로 천하를 다스린다')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.0), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.0),
                '半 部 論 語 治 天 下',
                font_size=44, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.2), Inches(12.0), Inches(0.4),
                '북송의 재상 조보(趙普)가 한 말 — 논어가 얼마나 깊고 실천적인 책인가',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 노신(魯迅) — 「중국 사상의 뼈대는 공자에게서 나왔다」',
         {'font_size': 17, 'space_before': 6}),
        ('● 빌 게이츠 — 인생의 책으로 꼽은 동양 고전',
         {'font_size': 17, 'space_before': 10}),
        ('● 피터 드러커 — 「공자의 인(仁)과 서(恕)는 현대 경영의 근본 원리다」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 한국 — 조선 500년 사대부의 「제1 교양서」 · 일본 — 에도 막부의 정치 이념',
         {'font_size': 17, 'space_before': 10}),
        ('● 16,000자의 작은 책이 2,500년의 동아시아를 만들었다',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC1)
def i_unique(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '논어가 다른 고전과 다른 점',
              '「체계적 저작」이 아니다 — 그래서 더 사람 냄새 나는 책')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1. 체계적 저작이 아닌 「대화록」',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('     『한비자』·『손자병법』은 체계 논문, 논어는 단편적 어록의 모음',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 2. 공자 자신이 쓴 책이 아니다 — 「述而不作」',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     공자는 「전할 뿐 짓지 않는다」 — 제자들이 스승의 말과 행동을 기록·편찬',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 3. 인간 공자의 약점까지 기록',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     화내는 공자·우는 공자·유머하는 공자 — 완벽한 성인이 아닌 흔들리며 가르친 스승',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 4. 해석 가능성의 무한함 — 「因材施敎」',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     같은 「인(仁)」을 제자마다 다르게 정의 — 독자마다 다른 논어를 읽는다',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
    ])


# ============== Ⅱ. 편찬 ==============
SEC2 = 'Ⅱ. 편찬'

@S(SEC2)
def ii_compile(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '편찬 과정 — 400년에 걸친 완성',
              '한 사람이 한 번에 쓴 책이 아니다')
    rows = [
        ('1차', '공자 사후 (BC 479~)', '중궁·자유·자하 등 직계 제자',  '기본 자료 정리'),
        ('2차', '증자(曾子) 사후',       '유자·민자건 등',                 '보충과 정리'),
        ('3차', '전국시대 (BC 400~300)', '맹자 시기 전후',                  '첨가와 보충'),
        ('4차', '전한 말 (BC 1세기)',   '장우(張禹)',                      '현행 판본의 기초 — 후한에서 정본화'),
    ]
    for i, (step, era, who, what) in enumerate(rows):
        y = Inches(2.5 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.0), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.0), Inches(0.8),
                    step, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.9), y, Inches(3.2), Inches(0.8), PALE)
        add_textbox(slide, Inches(1.9), y, Inches(3.2), Inches(0.8),
                    era, font_size=14, color=SUB, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.3), y + Inches(0.1), Inches(4.0), Inches(0.65),
                    who, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(9.5), y + Inches(0.1), Inches(3.4), Inches(0.65),
                    what, font_size=14, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
                '「논어」라는 이름은 전한 경제·무제 시기(BC 157~87)에 정착',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅲ. 공자의 생애 ==============
SEC3 = 'Ⅲ. 공자의 생애'

@S(SEC3)
def iii_profile(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '공자(孔子) — 한 인간의 기본 정보',
              '가난한 무관의 아들이 동아시아의 스승이 되기까지')
    rows = [
        ('이름',     '공구(孔丘), 자(字) 중니(仲尼)'),
        ('생몰',     'BC 551 ~ BC 479 (72~73세)'),
        ('출신',     '노(魯)나라 (현 중국 산동성 곡부시)'),
        ('부친',     '숙량흘(叔梁紇) — 노나라 하급 무관 (공자 3세에 사망)'),
        ('모친',     '안징재(顔徵在) — 홀어머니 밑에서 가난하게 성장'),
        ('아들',     '공리(孔鯉) — 공자보다 먼저 사망'),
        ('손자',     '자사(子思) — 『중용』 저자'),
        ('제자',     '약 3,000명, 뛰어난 자 72명 (칠십이현 七十二賢)'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.1), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC3)
def iii_timeline(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '공자가 직접 남긴 인생 연표 — 위정편 4장',
              '六 단어가 한국어의 「나이」를 가리키는 말이 되다')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.7),
                '吾 十 有 五 而 志 于 學  三 十 而 立  四 十 而 不 惑',
                font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(12.0), Inches(0.7),
                '五 十 而 知 天 命  六 十 而 耳 順  七 十 而 從 心 所 欲 不 踰 矩',
                font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    rows = [
        ('15', '志于學', '지학', '학문에 뜻을 둠'),
        ('30', '而立',   '이립', '자립함 — 독립적 사상 확립'),
        ('40', '不惑',   '불혹', '미혹되지 않음 — 의심·흔들림 없음'),
        ('50', '知天命', '지천명', '천명을 앎 — 소명·한계 자각'),
        ('60', '耳順',   '이순', '귀가 순해짐 — 어떤 말도 편안히 수용'),
        ('70', '從心',   '종심', '마음대로 해도 법도 어긋나지 않음'),
    ]
    for i, (age, han, kor, desc) in enumerate(rows):
        y = Inches(4.0 + i * 0.5)
        add_textbox(slide, Inches(0.7), y, Inches(0.8), Inches(0.4),
                    age, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.7), y, Inches(1.8), Inches(0.4),
                    han, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(3.6), y, Inches(1.5), Inches(0.4),
                    kor, font_size=13, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.2), y, Inches(7.6), Inches(0.4),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC3)
def iii_drama(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '공자의 인생 — 4단계 드라마',
              '소년의 가난 → 학문과 관직 → 14년 주유 → 노년의 교육')
    rows = [
        ('소년기 (0~19세)', '가난한 성장',
         '3세에 부친 사망 · 「吾少也賤 故多能鄙事」 — 어려서 천했기에 비천한 일을 많이 할 줄 안다.\n그럼에도 15세에 학문에 뜻을 둠(志于學).'),
        ('청·장년기 (20~50세)', '학문과 관직',
         '노 계손씨 가문 하급 관리.\n주(周) 낙양에서 노자를 만났다는 전승 · 제(齊)에서 「소(韶) 음악」 듣고 3개월간 고기 맛을 잊음.\n50세 전후 노나라 대사구(大司寇) 임명 — 정치적 절정.'),
        ('천하주유 (55~68세)', '14년의 방랑',
         '노에서 정치 개혁 실패 → 위·조·송·정·진·채 6개국 주유.\n광(匡)의 위기 · 진채(陳蔡)의 절량 · 정나라 「상갓집 개」.\n어느 제후도 등용하지 않음 — 그러나 가르침은 더 깊어짐.'),
        ('노년 귀향 (68~73세)', '교육과 저술',
         '68세 노로 귀환 — 후학 양성에 전념.\n『춘추』 편찬 · 『시경·서경·역경·예기』 정리.\n안연 요절(71세) → 「天喪予!」 · 자로 전사(72세).\nBC 479년 봄 사망, 향년 73세.'),
    ]
    for i, (stage, title, body) in enumerate(rows):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.4)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(2.2), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.45),
                    stage, font_size=14, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.4),
                    title, font_size=13, color=SUB, bold=True)
        add_paragraphs(slide, x + Inches(0.2), y + Inches(0.95), Inches(5.5), Inches(1.2),
                       [(body, {'font_size': 11, 'color': INK})], line_spacing=1.4)


# ============== Ⅳ. 공문사과와 핵심 제자 ==============
SEC4 = 'Ⅳ. 공문사과와 핵심 제자'

@S(SEC4)
def iv_four(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '공문사과(孔門四科) — 제자 분류',
              '선진편 3장 — 공자의 강점 기반 교육관')
    cols = [
        ('德 行', '덕행',  '인격·도덕 실천',  '안연 · 민자건\n염백우 · 중궁'),
        ('言 語', '언어',  '외교·웅변',       '재아 · 자공'),
        ('政 事', '정사',  '정치·행정',       '염유 · 자로'),
        ('文 學', '문학',  '학문·경전',       '자유 · 자하'),
    ]
    for i, (han, kor, role, who) in enumerate(cols):
        x = Inches(0.7 + i * 3.1)
        add_filled_rect(slide, x, Inches(2.3), Inches(2.9), Inches(1.0), ACCENT)
        add_textbox(slide, x, Inches(2.3), Inches(2.9), Inches(0.6),
                    han, font_size=26, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.85), Inches(2.9), Inches(0.4),
                    kor, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, x, Inches(3.4), Inches(2.9), Inches(0.6), PALE)
        add_textbox(slide, x, Inches(3.4), Inches(2.9), Inches(0.6),
                    role, font_size=13, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.2), Inches(4.2), Inches(2.6), Inches(2.5),
                       [(who, {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER})],
                       line_spacing=1.6)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '공자는 제자의 「강점」을 파악해 길을 안내 — 현대 strengths-based leadership의 원조',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC4)
def iv_disciples(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '핵심 제자 4인 — 논어를 빛낸 인물들')
    rows = [
        ('顔淵 안연 (BC 521~481)', '안빈낙도의 화신',
         '「賢哉回也! 一簞食一瓢飮 在陋巷 不改其樂」\n공자의 가장 사랑받은 제자 · 41세 요절\n「天喪予! — 하늘이 나를 버렸다!」 (공자의 통곡)'),
        ('子路 자로 (BC 542~480)', '용기의 화신',
         '공자보다 9세 어린 맏형격 제자 · 호위대장 같은 역할\n「知之爲知之 不知爲不知 是知也」(공자의 가르침)\n위(衛) 반란에서 「君子死 冠不免」 외치며 전사'),
        ('子貢 자공 (BC 520?~456?)', '외교·경제의 천재',
         '언어과 대표 · 거상(巨商) · 사기 화식열전에도 등장\n공자 사후 6년간 시묘살이 — 제자 중 가장 길게\n「夫子之不可及也 猶天之不可階而升也」'),
        ('曾子 증자 (BC 505~435)', '효도의 화신 · 도통의 핵심',
         '46세 어린 최연소급 제자 · 『대학』·『효경』 저자\n「吾日三省吾身」 — 매일 세 가지로 자기를 반성\n공자 → 증자 → 자사 → 맹자의 도통(道統) 핵심'),
    ]
    for i, (name, label, body) in enumerate(rows):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.4)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(2.2), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.45),
                    name, font_size=15, bold=True, color=ACCENT, font_name='Batang')
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.35),
                    label, font_size=12, color=SUB, bold=True)
        add_paragraphs(slide, x + Inches(0.2), y + Inches(0.95), Inches(5.5), Inches(1.2),
                       [(body, {'font_size': 11, 'color': INK})], line_spacing=1.4)


# ============== Ⅴ. 20편 구조 ==============
SEC5 = 'Ⅴ. 20편의 구조'

@S(SEC5)
def v_table(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '20편 한 폭으로 — 상론(1~10) + 하론(11~20)',
              '각 편명은 첫 두 글자에서 따옴 — 체계가 아닌 어록 모음')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(5.9), Inches(0.4), INK)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(5.9), Inches(0.4),
                '上 論 (상론, 1~10편)', font_size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_filled_rect(slide, Inches(7.0), Inches(2.2), Inches(5.7), Inches(0.4), INK)
    add_textbox(slide, Inches(7.0), Inches(2.2), Inches(5.7), Inches(0.4),
                '下 論 (하론, 11~20편)', font_size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    upper = [
        ('1', '學而 학이',   '학문의 출발'),
        ('2', '爲政 위정',   '정치의 도와 효'),
        ('3', '八佾 팔일',   '예의 본질'),
        ('4', '里仁 이인',   '인(仁)의 거처'),
        ('5', '公冶長 공야장', '인물 평가'),
        ('6', '雍也 옹야',   '안빈낙도'),
        ('7', '述而 술이',   '공자의 자기 진술'),
        ('8', '泰伯 태백',   '성왕과 군자'),
        ('9', '子罕 자한',   '공자의 가르침'),
        ('10', '鄕黨 향당',  '공자의 일상'),
    ]
    lower = [
        ('11', '先進 선진',  '제자들의 인물상'),
        ('12', '顔淵 안연',  '인과 군자의 정의'),
        ('13', '子路 자로',  '정치의 실무'),
        ('14', '憲問 헌문',  '인물·정치 평론'),
        ('15', '衛靈公 위령공', '군자의 도'),
        ('16', '季氏 계씨',  '삼우삼락'),
        ('17', '陽貨 양화',  '인간 본성'),
        ('18', '微子 미자',  '은자와 공자'),
        ('19', '子張 자장',  '제자들의 어록'),
        ('20', '堯曰 요왈',  '천명과 마무리'),
    ]
    for i, (no, name, desc) in enumerate(upper):
        y = Inches(2.62 + i * 0.42)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(5.9), Inches(0.4), bg)
        add_textbox(slide, Inches(0.7), y, Inches(0.6), Inches(0.4),
                    no, font_size=12, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.3), y, Inches(2.0), Inches(0.4),
                    name, font_size=12, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(3.3), y, Inches(3.3), Inches(0.4),
                    desc, font_size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    for i, (no, name, desc) in enumerate(lower):
        y = Inches(2.62 + i * 0.42)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(7.0), y, Inches(5.7), Inches(0.4), bg)
        add_textbox(slide, Inches(7.0), y, Inches(0.7), Inches(0.4),
                    no, font_size=12, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.7), y, Inches(2.1), Inches(0.4),
                    name, font_size=12, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(9.8), y, Inches(2.9), Inches(0.4),
                    desc, font_size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅵ. 20편 깊이 읽기 ==============
SEC6 = 'Ⅵ. 20편 깊이 읽기'

def make_chapter_slide(num, total, name_han, name_kor, headline,
                        original, modern, theme, point):
    @S(SEC6)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC6} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=26, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    headline, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.4), [
            (original, {'font_size': 18, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 편의 핵심', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(theme, {'font_size': 14, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 오늘에의 적용', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 15, 'color': INK})], line_spacing=1.35)


CHAPTERS = [
    ('學而 학이', '학이', '제1편 — 학문의 출발 · 논어의 첫 편',
     '學 而 時 習 之  不 亦 說 乎',
     '배우고 때때로 익히면 기쁘지 아니한가',
     '논어의 첫 구절 — 「학(學)」으로 시작 / 「유붕자원방래(有朋自遠方來)」 — 벗이 멀리서 옴 / 「인부지이불온(人不知而不慍)」 — 남이 알아주지 않아도 / 증자의 「吾日三省吾身」 — 매일 자기 점검',
     '공부의 기쁨은 외부 평가가 아닌 내부 변화에서 — 「위기지학(爲己之學)」의 출발점.'),
    ('爲政 위정', '위정', '제2편 — 정치의 도와 효',
     '爲 政 以 德  譬 如 北 辰  居 其 所 而 衆 星 共 之',
     '덕으로 정치를 하면, 마치 북극성이 제자리에 있고 뭇 별이 함께 도는 것과 같다',
     '「위정이덕(爲政以德)」 — 덕치의 정의 / 공자의 인생 연표(15~70) / 「知之爲知之」 — 자로에게 / 효(孝)에 대한 다양한 답 (인재시교)',
     '리더십은 「누가 명령하는가」가 아니라 「누구를 따르고 싶은가」 — 덕이 사람을 끌어당기는 힘.'),
    ('八佾 팔일', '팔일', '제3편 — 예(禮)의 본질',
     '人 而 不 仁  如 禮 何   人 而 不 仁  如 樂 何',
     '사람이 인하지 않다면 예를 어찌하며, 사람이 인하지 않다면 악을 어찌하랴',
     '「八佾舞於庭 是可忍也 孰不可忍也」 — 계씨의 무례에 분노 / 「祭如在 祭神如神在」 — 정성의 예 / 「樂而不淫 哀而不傷」 — 관저편 평어',
     '형식만의 예는 의미 없다 — 「내용」으로서의 인(仁)이 먼저. 진정성 없는 의례는 죽은 예.'),
    ('里仁 이인', '이인', '제4편 — 인(仁)의 거처',
     '里 仁 爲 美   君 子 喩 於 義  小 人 喩 於 利',
     '인한 마을에 사는 것이 아름답다 / 군자는 의로움에 밝고, 소인은 이익에 밝다',
     '「朝聞道 夕死可矣」 — 아침에 도를 들으면 저녁에 죽어도 좋다 / 「君子喩於義 小人喩於利」 — 군자와 소인의 결정적 차이 / 「見賢思齊」 — 어진 이를 보면 같아질 것을 생각',
     '환경이 사람을 만든다 — 누구와 함께 있느냐가 자기 인격의 절반. 어떤 가치에 밝은가가 군자/소인을 가른다.'),
    ('公冶長 공야장', '공야장', '제5편 — 인물 평가',
     '聞 一 以 知 十  聞 一 以 知 二',
     '안회는 하나를 들으면 열을 알지만, 자공은 하나를 들으면 둘을 알 뿐',
     '인물 평가의 모음편 / 자공과 안연의 비교 / 「敏而好學 不恥下問」 — 공문자의 시호 풀이 / 「老者安之 朋友信之 少者懷之」 — 공자의 포부',
     '인물을 평가할 때 강점에 주목 — 누구를 닮고 싶은가의 안목이 곧 자기 길의 안목.'),
    ('雍也 옹야', '옹야', '제6편 — 안빈낙도(安貧樂道)',
     '賢 哉 回 也  一 簞 食 一 瓢 飮  不 改 其 樂',
     '훌륭하도다 안회여 — 한 그릇 밥 한 표주박 물에도 그 즐거움을 바꾸지 않는다',
     '안연의 안빈낙도 / 「知之者不如好之者 好之者不如樂之者」 — 아는 자 < 좋아하는 자 < 즐기는 자 / 「中庸之爲德也 其至矣乎」 — 중용의 덕',
     '즐길 줄 아는 자가 가장 강하다 — 「열정」 위에 「즐거움」. 외부 조건이 흔들어도 무너지지 않는 내면.'),
    ('述而 술이', '술이', '제7편 — 공자의 자기 진술',
     '述 而 不 作  信 而 好 古',
     '전할 뿐 짓지 않으며, 옛것을 믿고 좋아한다',
     '공자가 자기를 그린 편 — 「不作」의 겸손 / 「三人行 必有我師焉」 — 세 사람이면 반드시 스승 / 「飯疏食 飮水 曲肱而枕之 樂亦在其中矣」 — 거친 밥에도 즐거움 / 「發憤忘食 樂以忘憂」',
     '겸손과 자족 — 새것을 짓는 게 아니라 옛것을 잇는다는 자세. 거친 밥에도 즐거움이 있다는 마음.'),
    ('泰伯 태백', '태백', '제8편 — 성왕과 군자',
     '士 不 可 以 不 弘 毅  任 重 而 道 遠',
     '선비는 넓고 굳세지 않으면 안 된다 — 짐은 무겁고 갈 길은 멀다',
     '증자의 「임중도원(任重道遠)」 / 「興於詩 立於禮 成於樂」 — 시에서 일으키고 예에서 서고 악에서 완성 / 「不在其位 不謀其政」 — 그 자리에 있지 않으면 그 일을 도모하지 않는다',
     '인생의 짐과 길은 늘 무겁고 멀다 — 그 무게를 감당하는 자가 「선비(士)」. 자기 책임의 자각.'),
    ('子罕 자한', '자한', '제9편 — 공자의 가르침과 자취',
     '子 絕 四 — 毋 意  毋 必  毋 固  毋 我',
     '공자에게 네 가지가 없었다 — 억측·고집·완고·자기 중심',
     '「子在川上曰 逝者如斯夫 不舍晝夜」 — 흐르는 물처럼 / 「歲寒然後知松柏之後彫也」 — 추워진 뒤에야 송백을 안다 / 「絕四」 — 네 가지 없음 / 「後生可畏」 — 후배가 두렵다',
     '비움의 미학 — 억측·고집·완고·자기 중심을 비울 때 비로소 보인다. 시간이 모든 것을 가려준다.'),
    ('鄕黨 향당', '향당', '제10편 — 공자의 일상',
     '食 不 厭 精  膾 不 厭 細',
     '밥은 정한 것을 싫어하지 않고, 회는 가는 것을 싫어하지 않는다',
     '공자의 일상 생활 묘사 — 식사·옷·말투·움직임의 디테일 / 「食不語 寢不言」 — 먹을 때와 잘 때 말 없음 / 「鄕人飮酒 杖者出 斯出矣」 — 노인을 먼저',
     '인격은 거창한 곳이 아니라 일상의 작은 동작에서 — 향당편은 공자의 「살아 있는 모습」을 가장 생생히 보여준다.'),
    ('先進 선진', '선진', '제11편 — 제자들의 인물상',
     '未 知 生  焉 知 死',
     '삶도 모르거늘 어찌 죽음을 알랴',
     '공문사과(德行·言語·政事·文學) 출전 / 자로의 죽음에 대한 답 — 「未能事人 焉能事鬼」 / 안연의 요절에 「天喪予!」 / 자로·증석·염유·공서화의 「각자의 뜻」',
     '죽음을 묻기 전에 삶을 먼저 — 형이상학적 사변보다 현재의 사람과 일에 충실. 「지금 여기」의 윤리.'),
    ('顔淵 안연', '안연', '제12편 — 인(仁)과 군자의 정의',
     '克 己 復 禮 爲 仁   己 所 不 欲  勿 施 於 人',
     '자기를 이기고 예로 돌아감이 인이다 / 자기가 바라지 않는 바를 남에게 베풀지 말라',
     '안연의 인(仁) 질문 → 克己復禮 / 중궁의 인(仁) 질문 → 己所不欲 勿施於人 / 사마우의 인(仁) 질문 → 訒(말의 신중함) / 「君君臣臣父父子子」 — 정명(正名)',
     '인(仁)의 정의가 사람마다 다른 까닭 — 「因材施敎」. 같은 진리도 그 사람에게 가장 절실한 말로 전해진다.'),
    ('子路 자로', '자로', '제13편 — 정치의 실무',
     '君 子 和 而 不 同   小 人 同 而 不 和',
     '군자는 화합하되 같지 않고, 소인은 같되 화합하지 않는다',
     '「화이부동(和而不同)」 — 군자와 소인의 결정적 차이 / 「其身正 不令而行」 — 자기 몸이 바르면 명령 없어도 행해진다 / 「先之勞之」 — 솔선수범',
     '진정한 화합은 「같음」이 아닌 「다름의 조화」. 다양성이 곧 힘 — 동질화는 동조 압력일 뿐.'),
    ('憲問 헌문', '헌문', '제14편 — 인물·정치 평론',
     '不 患 人 之 不 己 知  患 其 不 能 也',
     '남이 나를 알아주지 않음을 근심하지 말고, 자기의 능력 없음을 근심하라',
     '「不患人之不己知」 — 학이편의 변주 / 「以直報怨 以德報德」 — 곧음으로 원망에 갚고, 덕으로 덕에 갚는다 / 관중·자산 등 역대 인물 평론',
     '평가의 주체를 자기로 옮긴다 — 외부 인정에 끌려다니지 않는 자유. 자기 능력의 키움이 진짜 답.'),
    ('衛靈公 위령공', '위령공', '제15편 — 군자의 도',
     '己 所 不 欲  勿 施 於 人 — 一 言 而 可 以 終 身 行 之 者',
     '자기가 바라지 않는 바를 남에게 베풀지 말라 — 평생 한 마디로 행할 수 있는 것',
     '자공의 「평생 한 글자」 질문 → 「恕」와 「己所不欲 勿施於人」 / 「君子求諸己 小人求諸人」 — 군자는 자기에게서 구하고 소인은 남에게서 구한다 / 「君子固窮 小人窮斯濫矣」',
     '도덕의 가장 짧은 원칙 — 황금률(Golden Rule)의 동양적 음각. 「하지 말라」의 자기 절제로 시작.'),
    ('季氏 계씨', '계씨', '제16편 — 삼우삼락',
     '益 者 三 友 — 友 直  友 諒  友 多 聞',
     '이로운 세 벗 — 곧은 벗·신실한 벗·견문 넓은 벗',
     '익자삼우(益者三友) / 손자삼우(損者三友) — 便辟·善柔·便佞 / 익자삼락(益者三樂) — 절예악·도인지선·다현우 / 군자에게 「三戒」 — 색·투·득',
     '관계의 질이 인생의 질 — 누구와 어울리는가가 곧 나의 인격. 시기별 경계(三戒)는 자기 관리의 매뉴얼.'),
    ('陽貨 양화', '양화', '제17편 — 인간 본성과 시대 비판',
     '性 相 近 也  習 相 遠 也',
     '본성은 서로 가깝지만, 습관이 서로 멀게 한다',
     '인간 본성에 대한 가장 짧은 명제 / 「色厲而內荏」 — 겉은 엄하지만 속이 약함 / 「鄕原 德之賊也」 — 향원은 덕의 도적 / 「飽食終日 無所用心 難矣哉」',
     '사람의 차이는 타고난 본성이 아니라 「습관(習)」이 만든다 — 교육과 환경의 결정적 힘. 향원의 위선 경계.'),
    ('微子 미자', '미자', '제18편 — 은자와 공자',
     '鳥 獸 不 可 與 同 群  吾 非 斯 人 之 徒 與  而 誰 與',
     '새와 짐승과는 무리 지을 수 없으니, 내가 이 사람들과 함께하지 않으면 누구와 함께하랴',
     '은자(隱者)들 — 장저·걸닉·접여 / 그들이 공자를 「풍자」하고 공자가 답함 / 「鳥獸不可與同群」 — 공자의 인본주의 / 「天下有道 丘不與易也」 — 천하에 도가 있다면 내가 바꾸려 하지 않을 것',
     '세상을 떠나는 자에 대한 공자의 답 — 「인간 가운데 있어야 한다」. 도피보다 참여, 좌절보다 시도.'),
    ('子張 자장', '자장', '제19편 — 제자들의 어록',
     '博 學 而 篤 志  切 問 而 近 思  仁 在 其 中 矣',
     '널리 배우고 뜻을 돈독히 하며, 절실하게 묻고 가까이 생각하면 인은 그 가운데 있다',
     '공자의 말이 아니라 자장·자하·증자·자공의 어록 / 자하의 「博學篤志 切問近思」 — 성균관대학 교훈, 『근사록』 서명 출전 / 자공의 「夫子之不可及也 猶天之不可階而升也」',
     '스승의 말을 잇는 제자들의 목소리 — 가르침의 진정성은 「누가 이어 말하는가」로 입증된다.'),
    ('堯曰 요왈', '요왈', '제20편 — 천명과 마무리',
     '不 知 命  無 以 爲 君 子 也',
     '천명을 모르면 군자가 될 수 없다',
     '논어의 마지막 편 / 요·순·우·탕·문·무 — 성왕들의 말씀 인용 / 「不知命 不知禮 不知言」 — 천명·예·말의 세 가지를 알아야 / 학이편의 「학(學)」으로 시작하여 요왈편의 「명(命)」으로 끝남',
     '논어의 끝은 다시 시작 — 학문의 길은 결국 천명의 자각에 이른다. 수미상관(首尾相關)의 정점.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅶ. 핵심 사상 7기둥 ==============
SEC7 = 'Ⅶ. 핵심 사상 7기둥'

def make_concept_slide(num, total, han, kor, source, principle, today):
    @S(SEC7)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC7} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{han}  ({kor})',
                    font_size=30, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    source, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.7), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.6),
                       [(principle, {'font_size': 14, 'color': INK})], line_spacing=1.45)
        add_textbox(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.4),
                    '◆ 오늘에의 함의', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.3), Inches(12.0), Inches(3.0),
                       [(today, {'font_size': 15, 'color': INK})], line_spacing=1.5)


CONCEPTS = [
    ('仁', '인', '논어 전체의 중심 — 109회 등장',
     '· 「克己復禮爲仁」 — 자기를 이기고 예로 돌아감 (안연편)\n· 「己所不欲 勿施於人」 — 자기가 원치 않는 것을 남에게 베풀지 말라\n· 「仁者愛人」 — 인이란 사람을 사랑함\n· 「巧言令色 鮮矣仁」 — 꾸민 말과 얼굴엔 인이 드물다',
     '동양 도덕의 가장 큰 그릇 — 「사람다움」.\n현대 윤리학·EQ·사회적 감수성의 동양적 원천.\n공자가 제자마다 다르게 정의한 「열린 개념」.'),
    ('義', '의', '이익에 대비되는 「올바름」',
     '· 「君子喩於義 小人喩於利」 — 군자는 의에 밝고 소인은 이익에 밝다\n· 「見利思義」 — 이익을 보면 의를 생각하라\n· 「不義而富且貴 於我如浮雲」 — 의롭지 않은 부귀는 뜬구름 같다',
     '이익과 도덕의 분기점 — 「의(義)냐 이(利)냐」가 군자/소인을 가른다.\n현대 비즈니스 윤리·ESG의 동양 원형.\n「견리사의」가 모든 결정의 저울.'),
    ('禮', '예', '사회 규범과 질서의 형식',
     '· 「人而不仁 如禮何」 — 사람이 인하지 않다면 예를 어찌하랴\n· 「克己復禮爲仁」 — 자기를 이기고 예로 돌아감이 인\n· 「非禮勿視 非禮勿聽 非禮勿言 非禮勿動」 — 사물(四勿)',
     '예는 형식이 아니라 「인의 표현」 — 내용과 형식의 통합.\n현대의 매너·프로토콜·기업 윤리 강령의 근원.\n인 없는 예는 죽은 의례.'),
    ('知 / 學', '지·학', '앎과 배움의 윤리',
     '· 「知之爲知之 不知爲不知 是知也」 — 아는 것을 안다, 모르는 것을 모른다 함이 앎\n· 「學而時習之 不亦說乎」 — 배우고 익힘의 기쁨\n· 「知之者不如好之者 好之者不如樂之者」 — 아는 자 < 좋아하는 자 < 즐기는 자',
     '지적 겸손이 진짜 지혜 — 소크라테스 「너 자신을 알라」와 통한다.\n외부 인정이 아닌 내부 변화로서의 학문 — 「위기지학」.\n즐기는 자가 가장 강하다.'),
    ('孝 / 悌', '효·제', '관계 윤리의 출발',
     '· 「孝弟也者 其爲仁之本」 — 효와 제는 인의 근본\n· 「父在觀其志 父沒觀其行」 — 부모 생전에는 뜻을, 사후에는 행을 본다\n· 「事父母幾諫」 — 부모를 섬기되 은근히 간한다',
     '모든 관계의 시작은 부모-자식 — 가까운 데서 먼 데로 확장되는 윤리.\n「가까운 사람에게 막대하지 않기」의 동양적 원형.\n효(孝)는 절대 복종이 아닌 「깊은 배려」.'),
    ('忠 / 恕', '충·서', '공자 사상의 한 줄 요약',
     '· 「夫子之道 忠恕而已矣」 — 공자의 도는 충과 서일 뿐 (증자, 이인편)\n· 「忠」 — 마음을 다하는 것 (盡心)\n· 「恕」 — 같이 헤아리는 것 (推己及人)\n· 「己所不欲 勿施於人」 — 서(恕)의 정의 (위령공편)',
     '서양 「황금률」의 동양 음각 버전 — 「하지 말라」의 자기 절제.\n현대 협상·관계의 가장 짧은 원칙.\n공자 사상 전체를 두 글자로 요약 — 모든 도덕의 코드.'),
    ('君 子', '군자', '논어의 이상적 인간상 — 107회 등장',
     '· 「君子和而不同 小人同而不和」 — 군자는 화합하되 같지 않다\n· 「君子求諸己 小人求諸人」 — 군자는 자기에게서 구한다\n· 「君子坦蕩蕩 小人長戚戚」 — 군자는 평탄하고 소인은 늘 근심\n· 「君子喩於義 小人喩於利」',
     '논어 전체의 인격적 도착점 — 「군자」.\n태생이 아닌 「됨」의 인격 — 누구나 군자가 될 수 있다.\n현대 인격·리더십의 동양 표준.'),
]

for i, c in enumerate(CONCEPTS, 1):
    make_concept_slide(i, len(CONCEPTS), *c)


# ============== Ⅷ. 명구 16선 ==============
SEC8 = 'Ⅷ. 명구 16선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC8)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC8} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=16, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 16, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('학이편 첫 구절', '學 而 時 習 之  不 亦 說 乎',
     '배우고 때때로 익히면 기쁘지 아니한가',
     '논어 첫 문장. 학문의 본질을 「기쁨」으로 정의 — 외부 평가가 아닌 내부 변화. 「위기지학」의 출발.'),
    ('학이편 · 증자', '吾 日 三 省 吾 身',
     '나는 하루에 세 가지로 나를 반성한다',
     '증자의 일과 — 충(忠)·신(信)·습(習) 세 가지로 자기 점검. 매일 자기 점검 루틴의 가장 짧고 강력한 원전.'),
    ('위정편 · 공자의 인생 연표', '三 十 而 立  四 十 而 不 惑  五 十 而 知 天 命',
     '서른에 자립하고, 마흔에 미혹되지 않고, 쉰에 천명을 안다',
     '이립·불혹·지천명 — 한국어에서 나이를 가리키는 말이 된 6단어의 원전. 인생 연표의 동양적 표준.'),
    ('위정편', '知 之 爲 知 之  不 知 爲 不 知  是 知 也',
     '아는 것을 안다, 모르는 것을 모른다 하는 것이 아는 것이다',
     '공자가 자로에게. 지적 겸손이 곧 지혜의 시작 — 「너 자신을 알라」의 동양 음각 버전.'),
    ('이인편', '朝 聞 道  夕 死 可 矣',
     '아침에 도를 들으면 저녁에 죽어도 좋다',
     '진리에 대한 가장 뜨거운 갈망. 안위가 아니라 「도(道)」를 위해 사는 삶 — 모든 학자의 좌우명.'),
    ('이인편', '君 子 喩 於 義  小 人 喩 於 利',
     '군자는 의로움에 밝고, 소인은 이익에 밝다',
     '군자와 소인의 결정적 차이 — 무엇에 「밝은가」가 사람을 가른다. 비즈니스 윤리의 동양적 원천.'),
    ('이인편 · 증자가 풀이', '夫 子 之 道  忠 恕 而 已 矣',
     '선생님의 도는 충과 서일 뿐이다',
     '공자 사상 전체를 두 글자로 요약. 충(忠) — 마음을 다함 / 서(恕) — 같이 헤아림. 모든 도덕의 코드.'),
    ('옹야편 · 안연 찬탄', '賢 哉 回 也  一 簞 食  一 瓢 飮  不 改 其 樂',
     '훌륭하도다 안회여 — 한 그릇 밥, 한 표주박 물에도 그 즐거움을 바꾸지 않는다',
     '안빈낙도(安貧樂道)의 화신 안연에 대한 공자의 찬탄. 외부 조건과 무관한 내면의 강건함.'),
    ('옹야편', '知 之 者 不 如 好 之 者  好 之 者 不 如 樂 之 者',
     '아는 자는 좋아하는 자만 못하고, 좋아하는 자는 즐기는 자만 못하다',
     '학습·일·삶의 3단계 위계. 「즐기는 자가 가장 강하다」 — 현대 몰입(flow) 이론의 2,500년 선행 명제.'),
    ('술이편', '三 人 行  必 有 我 師 焉',
     '세 사람이 가면 그 가운데 반드시 내 스승이 있다',
     '겸손과 학습 의지의 결정판. 모든 사람에게서 배운다는 자세 — 「스승은 도처에 있다」.'),
    ('자한편 · 절사(絕四)', '子 絕 四 — 毋 意  毋 必  毋 固  毋 我',
     '공자에게 네 가지가 없었다 — 억측·고집·완고·자기 중심',
     '공자의 인격 자화상. 비움의 미학 — 네 가지를 비울 때 비로소 본다. 현대 열린 마음의 동양 원형.'),
    ('자한편', '逝 者 如 斯 夫  不 舍 晝 夜',
     '흘러가는 것이 이와 같구나 — 밤낮을 가리지 않는다',
     '공자가 강가에서 한 말. 시간의 무상함에 대한 가장 시적인 표현. 현재에 충실하라는 동양적 명제.'),
    ('안연편', '克 己 復 禮  爲 仁',
     '자기를 이기고 예로 돌아감이 인이다',
     '안연이 인을 물은 답. 인(仁)의 가장 핵심 정의. 「자기 절제 + 사회적 규범」의 통합.'),
    ('안연편 / 위령공편', '己 所 不 欲  勿 施 於 人',
     '자기가 바라지 않는 바를 남에게 베풀지 말라',
     '서(恕)의 정의. 황금률의 동양 음각. 「평생 한 글자로 행할 만한 것」 — 도덕의 가장 짧은 원칙.'),
    ('자로편', '君 子 和 而 不 同  小 人 同 而 不 和',
     '군자는 화합하되 같지 않고, 소인은 같되 화합하지 않는다',
     '진정한 화합은 「다름의 조화」. 동질화 압력에 대한 가장 오래된 비판. 다양성의 가치를 선언.'),
    ('요왈편 · 논어의 마지막', '不 知 命  無 以 爲 君 子 也',
     '천명을 모르면 군자가 될 수 없다',
     '논어의 마지막 구절. 학이편의 「학(學)」으로 시작해 요왈편의 「명(命)」으로 끝난다 — 수미상관의 정점.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅸ. 한국·일본·서양 수용사 ==============
SEC9 = 'Ⅸ. 한국·일본·서양 수용사'

@S(SEC9)
def ix_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '한국 — 조선 500년의 「제1 교양서」',
              '사림의 정체성 · 과거 시험의 핵심 · 현대까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 고려 — 안향이 원에서 주자학과 함께 도입',
         {'font_size': 17, 'space_before': 4}),
        ('● 조선 — 사서삼경의 첫머리 · 과거 시험(생원·진사)의 핵심',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 퇴계 이황 — 『논어석의』 · 율곡 이이 — 『논어석의』',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 정약용 — 『논어고금주』 — 가장 방대한 조선 시대 논어 주석서',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 경연(經筵) — 왕에게 가장 자주 진강된 책',
         {'font_size': 17, 'space_before': 10}),
        ('● 현대 — 「오늘 우리도 매일 쓰는 표현」 — 학이시습·온고지신·삼인행필유아사 등',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 한국인의 「민족 어휘」에 가장 깊이 새겨진 외래 고전',
         {'font_size': 16, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


@S(SEC9)
def ix_japan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '일본 — 에도 막부의 정치 이념',
              '오규 소라이·이토 진사이의 비판적 독해부터 시부사와 에이이치까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 오진 천황 시기 백제 왕인(王仁)이 『논어』를 전했다는 전승',
         {'font_size': 17, 'space_before': 4}),
        ('● 에도 막부 — 도쿠가와의 통치 이념으로 채택 (하야시 라잔)',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 이토 진사이(伊藤仁齋) · 오규 소라이(荻生徂徠) — 주자 비판하고 「원전 논어」로 회귀',
         {'font_size': 17, 'space_before': 10}),
        ('● 시부사와 에이이치(澁澤榮一) — 「논어와 주판」',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 일본 자본주의의 아버지 — 「도덕(논어) + 경제(주판)」의 결합',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 메이지 이후 — 일본 기업 윤리의 사상적 기반',
         {'font_size': 17, 'space_before': 10}),
    ])


@S(SEC9)
def ix_west(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '서양 — 「Confucius」가 된 공자',
              '예수회 선교사부터 빌 게이츠·피터 드러커까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 17세기 — 예수회 선교사 마테오 리치 → 「Confucius」로 라틴어 명명',
         {'font_size': 17, 'space_before': 4}),
        ('● 1687 — 『Confucius Sinarum Philosophus』 (파리) — 논어 최초 라틴어 완역',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 볼테르 · 라이프니츠 — 공자를 「이상적 철학자」로 격찬 — 계몽주의에 영향',
         {'font_size': 17, 'space_before': 10}),
        ('● 헤겔 — 비판적이었으나 동양 사유의 대표로 인정',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 현대 — 빌 게이츠의 「인생의 책」 · 피터 드러커 「인(仁)과 서(恕)는 현대 경영의 근본」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 영문 번역 50종 이상 — 동양 고전 중 가장 많이 번역·인용된 책',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


# ============== Ⅹ. 오늘 우리에게 ==============
SEC10 = 'Ⅹ. 오늘 우리에게'

@S(SEC10)
def x_modern(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '논어가 오늘 우리에게 일러주는 12가지')
    items = [
        '1. 공부는 외부 평가가 아닌 내부 변화 — 學而時習之',
        '2. 매일 자기를 세 가지로 점검하라 — 三省吾身',
        '3. 모르는 것을 모른다 함이 진짜 앎 — 知之爲知之 不知爲不知',
        '4. 의(義)와 이(利)의 갈림에서 의를 보라 — 君子喩於義',
        '5. 자기가 원치 않는 것을 남에게 베풀지 말라 — 己所不欲 勿施於人',
        '6. 즐기는 자가 가장 강하다 — 樂之者',
        '7. 세 사람이면 반드시 내 스승이 있다 — 三人行必有我師',
        '8. 다름의 조화 — 和而不同',
        '9. 평가를 자기에게 돌리라 — 不患人之不己知',
        '10. 한 덕목에 고착되지 말라 — 絕四(억측·고집·완고·자기 중심)',
        '11. 진리에 대한 갈망 — 朝聞道 夕死可矣',
        '12. 천명을 안다는 것 — 不知命 無以爲君子',
    ]
    for i, txt in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 0.7)
        add_textbox(slide, x, y, Inches(6.0), Inches(0.6),
                    txt, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅺ. 마무리 ==============
SEC11 = 'Ⅺ. 마무리'

@S(SEC11)
def xi_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '논어, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 공자와 제자들의 언행록 — 20편 482장 약 16,000자.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 사서(四書)의 으뜸 — 동아시아 2,500년의 정신적 DNA.',
         {'font_size': 18, 'space_before': 8}),
        ('● 학(學)으로 시작해 명(命)으로 끝나는 수미상관의 책.',
         {'font_size': 18, 'space_before': 8}),
        ('● 7대 사상 기둥 — 인 · 의 · 예 · 지/학 · 효/제 · 충/서 · 군자.',
         {'font_size': 18, 'space_before': 8}),
        ('● 학이시습·온고지신·이립불혹·지천명·종심·삼인행필유아사 — 일상에 살아 있는 표현.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「흔들리며 가르친 한 스승의 살아 있는 모습」.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC11)
def xi_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.1),
                '學 而 時 習 之',
                font_size=84, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.9),
                '不 亦 說 乎',
                font_size=60, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.4), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.6),
                '배 우 고  때 때 로  익 히 면  기 쁘 지  아 니 한 가',
                font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
                '— 논어 제1편 학이(學而), 첫 구절',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                '論  語',
                font_size=22, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\논어_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
