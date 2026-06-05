# -*- coding: utf-8 -*-
"""
명심보감(明心寶鑑) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '마음을 밝혀주는 보배로운 거울 · 800년 동아시아의 인격 교과서',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.6),
                '明 心 寶 鑑',
                font_size=92, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '명 심 보 감',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
                '爲善者 天報之以福 — 선을 행하면 하늘이 복으로 갚는다',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '추적(秋適, 고려) · 범립본(范立本, 명) — 유·불·도 삼교 융합 격언집',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 명심보감은 어떤 책인가'),
        ('Ⅱ.', '편찬 — 추적과 범립본'),
        ('Ⅲ.', '삼교합일 — 유·불·도의 융합'),
        ('Ⅳ.', '20편의 구조 — 다섯 군(群)'),
        ('Ⅴ.', '20편 한 폭으로 보기'),
    ]
    items_right = [
        ('Ⅵ.', '핵심 편 깊이 읽기'),
        ('Ⅶ.', '명구 18선'),
        ('Ⅷ.', '조선과 동아시아의 명심보감'),
        ('Ⅸ.', '오늘 다시 펼치는 이유'),
        ('Ⅹ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '명심보감(明心寶鑑) — 이름이 곧 교육 철학',
              '마음을 밝혀 주는 보배로운 거울')
    add_filled_rect(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.0), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.0),
                '明 心  =  마음을 밝히다     寶 鑑  =  보배로운 거울',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 사람의 마음은 본래 거울 — 욕망·편견·분노의 먼지로 가려지기 쉽다',
         {'font_size': 18, 'space_before': 8}),
        ('● 매일 선현의 한 구절을 읽어 그 먼지를 닦아내는 책',
         {'font_size': 18, 'space_before': 10}),
        ('● 즉석 깨달음이 아닌 「축적된 거울」 — 외워두면 어느 순간 나를 비추는 한 줄',
         {'font_size': 18, 'space_before': 10}),
        ('● 이름이 곧 독서법 — 「매일의 한 구절」 습관을 만드는 책',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC1)
def i_outline(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 명심보감')
    rows = [
        ('서명',     '명심보감(明心寶鑑) — 「마음을 밝혀 주는 보배로운 거울」'),
        ('원편자',   '추적(秋適, 1246~1317) — 고려 충렬왕 때 예문관제학'),
        ('증보자',   '범립본(范立本) — 명 홍무 26년(1393)'),
        ('편수',     '현행 상·하 2권 20편 (증보판 최대 24~25편)'),
        ('성격',     '처세·수양 격언집 · 유·불·도 삼교 융합'),
        ('교육 위상', '조선 서당 — 천자문 다음, 소학 앞에 놓인 기초 교재'),
        ('국제 영향', '한·중·일·베트남, 1592년 스페인어 번역(마닐라)'),
        ('현존 최고본', '1454년 「신간대자명심보감」 — 청주고인쇄박물관 소장'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_position(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '조선 서당 교육에서의 자리',
              '천자문 다음, 소학 앞 — 「짧은 금언으로 한문과 인격을 동시에」')
    rows = [
        ('5~6세', '千 字 文',    '천자문', '한자 1,000자와 기본 세계관'),
        ('7~8세', '明 心 寶 鑑', '명심보감', '금언 모음으로 한문 문리와 생활 윤리', True),
        ('9~15세', '小 學',       '소학', '일상 예절과 인격의 체계'),
        ('15세~',  '四 書 三 經', '사서삼경', '본격 유학의 세계'),
    ]
    for i, row in enumerate(rows):
        age, han, kor, desc = row[0], row[1], row[2], row[3]
        highlighted = len(row) > 4 and row[4]
        y = Inches(2.5 + i * 1.0)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.5), Inches(0.8), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.8),
                    age, font_size=14, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        color = ACCENT if highlighted else INK
        add_filled_rect(slide, Inches(2.4), y, Inches(3.0), Inches(0.8), color)
        add_textbox(slide, Inches(2.4), y, Inches(3.0), Inches(0.8),
                    han, font_size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.6), y + Inches(0.05), Inches(1.8), Inches(0.7),
                    kor, font_size=15, color=SUB, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.5), y + Inches(0.05), Inches(5.4), Inches(0.7),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_thought(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '핵심 사상 네 기둥')
    boxes = [
        ('三 敎 合 一', '삼교합일', '유·불·도의 융합 — 한 권 안의 「지혜의 뷔페」'),
        ('勸 善 懲 惡', '권선징악', '선에는 복, 악에는 화 — 인과응보의 단순하고 강력한 감각'),
        ('修 己 治 人', '수기치인', '자기 수양에서 가정·사회로 — 동심원적 확장'),
        ('安 分 知 足', '안분지족', '분수를 알고 만족함 — 욕망의 절제와 참된 행복'),
    ]
    for i, (han, kor, desc) in enumerate(boxes):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(6.0), Inches(2.0), PALE)
        add_textbox(slide, x, y + Inches(0.25), Inches(6.0), Inches(0.7),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(1.0), Inches(6.0), Inches(0.4),
                    kor, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.3), y + Inches(1.4), Inches(5.4), Inches(0.55),
                    desc, font_size=13, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 편찬 ==============
SEC2 = 'Ⅱ. 편찬 — 두 사람의 책'

@S(SEC2)
def ii_chu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '원편자 — 추적(秋適, 1246~1317)',
              '고려 충렬왕 때 예문관제학 · 노당(露堂)')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1305년 전후 — 왕자와 유생들의 기초 교육을 위해 편찬',
         {'font_size': 18, 'space_before': 4}),
        ('● 경전·제자백가·역사서에서 「마음을 닦는 금언」 수백 구절을 가려 뽑아 19편으로 엮음',
         {'font_size': 18, 'space_before': 10}),
        ('● 시대 배경 — 고려 말 원(元) 간섭기, 몽골풍 성행, 가치관 붕괴',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 불교의 세속화·신유학의 도입 초기 — 「무엇을 기준으로 살 것인가」의 공백',
         {'font_size': 17, 'space_before': 8, 'color': SUB}),
        ('● 추적의 선택 — 「무거운 경전을 늘어놓기 전에, 짧고 깊은 격언으로 마음의 방향부터」',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 편찬 철학 — 「이론 이전에 태도, 사상 이전에 감각, 체계 이전에 습관」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC2)
def ii_fan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '증보자 — 범립본(范立本, 1393)',
              '명 홍무 26년 — 대중 교화의 격언집으로 확장')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 명 태조 주원장의 「육유(六諭)」 반포 시기 — 대중 교화가 시대적 과제',
         {'font_size': 18, 'space_before': 4}),
        ('● 추적본을 저본으로 증보 — 상·하 2권 20편(혹은 24~25편)으로 확장',
         {'font_size': 18, 'space_before': 10}),
        ('● 추적본이 「왕실·유생용 정선본」이라면, 범립본은 「민간·아동용 대중본」',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 조선에서는 범립본 계통의 20편본이 주로 유통됨',
         {'font_size': 17, 'space_before': 10}),
        ('● 한·중·일·베트남·필리핀 — 동아시아 한문 문화권의 대중 교육 베스트셀러',
         {'font_size': 17, 'space_before': 10}),
        ('● 1592년 마닐라 — 스페인어 번역 「Beng Sim Po Cam」 — 동양 고전 유럽어 번역의 효시 중 하나',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC2)
def ii_versions(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '주요 판본 — 700년의 책의 물리적 역사')
    rows = [
        ('1454', '신간대자명심보감(新刊大字明心寶鑑)',
         '현존 최고(最古) 간본 · 청주고인쇄박물관 소장 · 큰 글자(大字)로 아동 학습용'),
        ('조선 전·중기', '조선 목판본 10여 종',
         '국립중앙도서관·규장각·각 대학 도서관 소장'),
        ('조선 후기~일제', '중간본(重刊本)·석판본',
         '편의 수가 20·23·24·25편으로 달라지며 다양한 증보판이 유통'),
        ('1592', '스페인어 번역본 (마닐라)',
         '예수회 선교사가 화교를 통해 번역 — 「Beng Sim Po Cam」'),
    ]
    for i, (year, name, desc) in enumerate(rows):
        y = Inches(2.4 + i * 1.05)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.85),
                    year, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(2.9), y + Inches(0.05), Inches(10.0), Inches(0.8), [
            (name, {'font_size': 15, 'bold': True, 'color': INK}),
            (desc, {'font_size': 13, 'color': SUB, 'space_before': 4}),
        ], line_spacing=1.25)
    add_textbox(slide, Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.4),
                '「정본(正本)이 없는 책」 — 지역·시대·대상에 따라 자유롭게 변용된 「살아 있는 책」',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅲ. 삼교합일 ==============
SEC3 = 'Ⅲ. 삼교합일'

@S(SEC3)
def iii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '명심보감의 가장 큰 특징 — 삼교합일(三敎合一)',
              '유·불·도의 가르침이 한 책에 자연스럽게 섞이다')
    cols = [
        ('儒 유', '사회의 윤리',
         '효(孝) · 오륜 · 수기치인\n\n출전 — 논어·맹자·예기·시경·\n서경·주역·주돈이·정명도·주희',
         ACCENT),
        ('道 도', '자연과 분수',
         '순명(順命) · 안분(安分)\n무위(無爲)의 지혜\n\n출전 — 장자·열자·태상감응편·\n문창제군 등 도교 경전',
         INK),
        ('佛 불', '인과와 자비',
         '인과응보 · 겸손 · 용서\n참음의 덕\n\n출전 — 여러 불경의 경구,\n선종 조사의 말씀',
         SUB),
    ]
    for i, (han, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 4.2)
        add_filled_rect(slide, x, Inches(2.3), Inches(3.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(3.9), Inches(0.6),
                    han, font_size=28, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.85), Inches(3.9), Inches(0.4),
                    label, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.3), Inches(3.6), Inches(3.5), Inches(3.5),
                       [(body, {'font_size': 13, 'color': INK})], line_spacing=1.4)


@S(SEC3)
def iii_meaning(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '왜 한 권 안에 섞였는가',
              '삼교 융합의 의미')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 대중 교육의 현실주의 — 어린이와 서민에게 학파 논쟁은 사치',
         {'font_size': 18, 'space_before': 6}),
        ('● 실생활에 도움이 되는 가르침이면 어느 학파든 가져다 쓴다',
         {'font_size': 18, 'space_before': 10}),
        ('● 동아시아의 「3중 구조 사회」의 반영',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 유학 — 국가 체제 · 불교 — 개인 구원 · 도교 — 자연·운명',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 편찬자의 개방성 — 「사람에게 좋은 말이면 어디서 왔든 다 좋다」',
         {'font_size': 17, 'space_before': 12}),
        ('● 하나의 이데올로기로 훈육하지 않는 책 — 「이념 교육」이 아닌 「지혜 축적」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 「지혜의 뷔페」 — 아이는 자기에게 맞는 구절을 고르고, 자라며 새 의미를 발견',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅳ. 20편의 다섯 군 ==============
SEC4 = 'Ⅳ. 20편 — 다섯 군(群)'

@S(SEC4)
def iv_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '20편을 다섯 군으로 묶다',
              '세계관 → 관계 → 자기 수양 → 사회 → 실천 매뉴얼')
    boxes = [
        ('군 1', '세계관의 뼈대', '계선 · 천명 · 순명', '선악과 천명 — 인과의 질서'),
        ('군 2', '관계의 시작', '효행', '부모-자식 — 모든 관계의 출발'),
        ('군 3', '내면의 다스림', '정기·안분·존심·계성·근학·훈자', '자기 수양 — 명심보감의 심장부'),
        ('군 4', '세상을 사는 지혜', '성심 상·하·입교·치정·치가·안의', '성찰에서 가정·사회·국가로'),
        ('군 5', '실천 매뉴얼', '준례 · 언어 · 교우 · 부행', '예·말·친구·가정의 구체 처방'),
    ]
    for i, (g, label, scope, desc) in enumerate(boxes):
        y = Inches(2.3 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.2), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.2), Inches(0.8),
                    g, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(2.0), y, Inches(3.0), Inches(0.8), PALE)
        add_textbox(slide, Inches(2.0), y, Inches(3.0), Inches(0.8),
                    label, font_size=15, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(5.2), y + Inches(0.05), Inches(7.7), Inches(0.75), [
            (scope, {'font_size': 14, 'bold': True, 'color': ACCENT}),
            (desc,  {'font_size': 13, 'color': SUB, 'space_before': 4}),
        ], line_spacing=1.25)


@S(SEC4)
def iv_flow(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '교육의 흐름 — 「세계관 → 사람 → 나 → 세상 → 일상」',
              '왜 이 순서인가')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 먼저 「세상에 질서가 있다」는 감각을 심는다 (1~3편)',
         {'font_size': 18, 'space_before': 6}),
        ('     선이 복을 부르고 악이 화를 부른다 — 우연이 아닌 인과의 세계',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 모든 관계의 시작인 「부모와 자식」 (4편)',
         {'font_size': 18, 'space_before': 10}),
        ('     이 관계에서 예의·감사를 배워야 다른 관계를 맺을 수 있다',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 어떤 상황에서도 변하지 않는 「자기 자신을 다스리는 능력」 (5~10편)',
         {'font_size': 18, 'space_before': 10}),
        ('     명심보감의 심장부 — 여섯 편이 모두 내면의 관리에 집중',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 「수신 → 제가 → 치국」으로 확장된 사회의 지혜 (11~16편)',
         {'font_size': 18, 'space_before': 10}),
        ('     『대학』의 수신제가치국평천하가 짧은 금언으로 구현',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 마지막 — 「예·말·친구·가정」의 실천 매뉴얼 (17~20편)',
         {'font_size': 18, 'space_before': 10}),
    ])


# ============== Ⅴ. 20편 한 폭으로 보기 ==============
SEC5 = 'Ⅴ. 20편 한 폭으로'

@S(SEC5)
def v_table1(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '20편 전체 (1) — 1~10편',
              '세계관 · 관계 · 자기 수양')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.5), INK)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(1.5), Inches(0.5),
                'No.', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(2.2), Inches(2.2), Inches(3.5), Inches(0.5),
                '편명', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(5.7), Inches(2.2), Inches(7.0), Inches(0.5),
                '핵심 주제', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('1', '繼善篇 계선편', '선을 이어가라 — 선행을 끊임없이 실천'),
        ('2', '天命篇 천명편', '천명을 두려워하라 — 선과 천도의 관계'),
        ('3', '順命篇 순명편', '운명에 순응하되 선은 쌓는다'),
        ('4', '孝行篇 효행편', '부모에 대한 효도의 도리'),
        ('5', '正己篇 정기편', '자기를 바르게 — 청렴과 자기반성'),
        ('6', '安分篇 안분편', '분수를 알고 만족하라'),
        ('7', '存心篇 존심편', '마음을 보존하라 — 겸손과 용서'),
        ('8', '戒性篇 계성편', '성질을 경계하라 — 분노를 참는 덕'),
        ('9', '勤學篇 근학편', '부지런히 배워라'),
        ('10', '訓子篇 훈자편', '자녀를 잘 가르쳐라'),
    ]
    for i, (no, name, desc) in enumerate(rows):
        y = Inches(2.7 + i * 0.42)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.42), bg)
        add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.42),
                    no, font_size=13, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.2), y, Inches(3.5), Inches(0.42),
                    name, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.9), y, Inches(6.7), Inches(0.42),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC5)
def v_table2(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '20편 전체 (2) — 11~20편',
              '성찰 · 사회 · 실천 매뉴얼')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.5), INK)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(1.5), Inches(0.5),
                'No.', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(2.2), Inches(2.2), Inches(3.5), Inches(0.5),
                '편명', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(5.7), Inches(2.2), Inches(7.0), Inches(0.5),
                '핵심 주제', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('11', '省心篇 上 성심편 상', '마음의 성찰 (일반) — 인간 본성과 처세'),
        ('12', '省心篇 下 성심편 하', '마음의 성찰 (사회) — 사람 사귐과 분별'),
        ('13', '立敎篇 입교편', '가르침을 세움 — 삼강오륜의 정립'),
        ('14', '治政篇 치정편', '정치의 도 — 애민(愛民)'),
        ('15', '治家篇 치가편', '집을 다스림 — 근검과 화목'),
        ('16', '安義篇 안의편', '인륜의 의리 — 부부·형제·붕우'),
        ('17', '遵禮篇 준례편', '예의를 지켜라'),
        ('18', '言語篇 언어편', '말을 삼가라'),
        ('19', '交友篇 교우편', '벗을 가려라'),
        ('20', '婦行篇 부행편', '부인의 네 가지 덕(사덕 四德)'),
    ]
    for i, (no, name, desc) in enumerate(rows):
        y = Inches(2.7 + i * 0.42)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.42), bg)
        add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.42),
                    no, font_size=13, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.2), y, Inches(3.5), Inches(0.42),
                    name, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.9), y, Inches(6.7), Inches(0.42),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅵ. 핵심 편 깊이 읽기 ==============
SEC6 = 'Ⅵ. 핵심 편 깊이 읽기'

def make_chapter_slide(num, total, name_han, name_kor, headline, original, modern, theme, point):
    @S(SEC6)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC6} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=28, bold=True, color=INK)
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    headline, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.4), [
            (original, {'font_size': 18, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 편 전체의 가르침', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(theme, {'font_size': 16, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 음미할 지점', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 16, 'color': INK})], line_spacing=1.35)


CHAPTERS = [
    ('계선편(繼善篇)', '계선편', '제1편 — 선을 이어가라 · 명심보감 전체의 서두',
     '爲善者 天報之以福  爲不善者 天報之以禍',
     '선을 행하면 하늘이 복으로 갚고, 악을 행하면 화로 갚는다',
     '명심보감의 첫 구절. 선악의 인과응보라는 전체 기조를 단번에 세운다.',
     '복잡한 논증 없이 「선=복, 악=화」의 단순한 등식 — 아이의 마음에 박는 「생활의 디폴트 설정」.'),
    ('효행편(孝行篇)', '효행편', '제4편 — 부모 은혜의 무한함과 효의 실천',
     '父兮生我 母兮鞠我  欲報之德 昊天罔極',
     '아버지는 나를 낳으시고 어머니는 나를 기르시니, 그 은덕을 갚고자 하나 하늘처럼 끝이 없다',
     '『시경·소아·요아』에서 가져온 구절. 부모의 은혜는 갚을 길 없을 만큼 깊다.',
     '효는 「의무」이전에 「감사의 자연스러운 표현」. 모든 인간관계의 시작이 여기에 있다.'),
    ('정기편(正己篇)', '정기편', '제5편 — 자기를 바르게 하라',
     '見人之善而尋己之善  見人之惡而尋己之惡',
     '남의 선을 보면 자기의 선을 찾고, 남의 악을 보면 자기의 악을 찾는다',
     '남을 거울 삼아 자기를 본다 — 정기(正己)의 핵심 자세.',
     '비판은 거울처럼 자기에게 돌리고, 칭찬은 자기를 점검하는 자료로 쓴다.'),
    ('안분편(安分篇)', '안분편', '제6편 — 분수를 알고 만족하라',
     '知足者 貧賤亦樂  不知足者 富貴亦憂',
     '만족할 줄 아는 자는 빈천해도 즐겁고, 만족할 줄 모르는 자는 부귀해도 근심한다',
     '행복은 가진 것의 크기가 아니라 만족할 줄 아는 마음에서 온다.',
     '도교적 무위·유교적 절제·불교적 욕망 비움이 한 줄에 녹은 「안분지족」의 정수.'),
    ('계성편(戒性篇)', '계성편', '제8편 — 성질을 경계하라',
     '忍一時之忿  免百日之憂',
     '한순간의 분을 참으면 백일의 근심을 면한다',
     '분노의 격발 한 번이 백일의 후회를 부른다 — 감정 관리의 동양적 격언.',
     '「인(忍)」 한 글자가 명심보감 전체의 키워드 중 하나. 어른의 자제력은 한순간에 결정된다.'),
    ('근학편(勤學篇)', '근학편', '제9편 — 부지런히 배워라',
     '玉不琢 不成器  人不學 不知道',
     '옥은 다듬지 않으면 그릇이 되지 않고, 사람은 배우지 않으면 도를 모른다',
     '재능이 아무리 좋아도 갈고 닦지 않으면 쓰임을 얻지 못한다 — 『예기·학기』 인용.',
     '명심보감이 학습을 가르치는 방식 — 「공부하라」 명령이 아닌 「옥이 다듬어지듯」의 비유.'),
    ('훈자편(訓子篇)', '훈자편', '제10편 — 자녀를 잘 가르쳐라',
     '黃金滿籯  不如敎子一經',
     '황금이 한 광주리에 가득해도 자식에게 경전 한 권을 가르치는 것만 못하다',
     '재산이 아니라 가르침이 진정한 유산 — 자녀 교육의 가치.',
     '돈을 물려준 부모와 책을 읽힌 부모 — 그 차이가 자녀의 평생을 가른다.'),
    ('성심편(省心篇)', '성심편', '제11~12편 — 마음의 성찰 · 명심보감의 분량의 절반',
     '路遙知馬力  日久見人心',
     '길이 멀어야 말의 힘을 알고, 날이 오래되어야 사람의 마음을 안다',
     '사람을 판단하기에 가장 좋은 도구는 시간이다.',
     '성심편 상·하는 처세·인간관계·세상 보는 눈의 백과 — 어른의 매뉴얼.'),
    ('치정편(治政篇)', '치정편', '제14편 — 정치의 도 · 애민(愛民)',
     '愛民如子  如子視父',
     '백성을 자식처럼 사랑하면 백성이 (다스리는 자를) 아비처럼 본다',
     '정치의 요체 — 한쪽이 자식처럼 사랑할 때 다른 쪽이 아비처럼 본다.',
     '「公則生明 廉則生威」 — 공정함에서 밝음이, 청렴함에서 위엄이 나온다. 공직자의 두 기둥.'),
    ('치가편(治家篇)', '치가편', '제15편 — 집을 다스리는 근본',
     '勤儉  治家之本',
     '부지런함과 검소함이 집안 다스림의 근본이다',
     '집안의 흥망은 거창한 일이 아닌 「부지런함과 검소함」 두 마디로 정해진다.',
     '근(勤)과 검(儉) — 동양 가정 경영의 두 키워드. 사치와 게으름이 들어오면 가문은 무너진다.'),
    ('언어편(言語篇)', '언어편', '제18편 — 말을 삼가라',
     '口是傷人斧  言是割舌刀',
     '입은 사람을 상하게 하는 도끼요, 말은 혀를 베는 칼이다',
     '말은 도끼이자 칼 — 함부로 휘두르면 남도 다치고 결국 자기도 다친다.',
     '「一言半句, 重値千金」 — 짧은 말 한마디가 천금. 말의 무게는 말한 사람의 무게다.'),
    ('교우편(交友篇)', '교우편', '제19편 — 벗을 가려라',
     '相識滿天下  知心能幾人',
     '아는 이는 천하에 가득해도 마음을 아는 이는 몇이나 되겠는가',
     '「아는 사람」과 「마음을 아는 사람」의 큰 차이 — 진정한 우정의 희소함.',
     '「오래 사귀어야 그 사람의 마음을 안다」 — 빠른 친밀과 진짜 우정의 구별.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅶ. 명구 18선 ==============
SEC7 = 'Ⅶ. 명구 18선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC7)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC7} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=30, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=17, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 17, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('계선편 — 공자', '爲 善 者 天 報 之 以 福  爲 不 善 者 天 報 之 以 禍',
     '선을 행하는 자에게 하늘은 복으로 갚고, 악을 행하는 자에게 화로 갚는다',
     '명심보감의 첫 구절. 「선=복, 악=화」 — 단순하지만 가장 강력한 동양 윤리의 디폴트.'),
    ('계선편 — 유비', '勿 以 善 小 而 不 爲  勿 以 惡 小 而 爲 之',
     '선이 작다 하여 하지 않으면 안 되고, 악이 작다 하여 행해서는 안 된다',
     '유비가 아들 유선에게 남긴 유훈. 도덕에는 크고 작음이 없다 — 작은 행동이 인격을 만든다.'),
    ('계선편 — 주역', '積 善 之 家  必 有 餘 慶',
     '선을 쌓은 집안에는 반드시 남는 경사가 있다',
     '『주역·곤괘 문언전』의 명구. 선행의 효과는 개인을 넘어 가문 전체에 미치며 후대까지 이어진다.'),
    ('계선편', '一 日 不 念 善  諸 惡 皆 自 起',
     '하루라도 선을 생각하지 않으면 온갖 악이 저절로 일어난다',
     '선한 마음은 의식적으로 유지해야 한다 — 방심하면 악이 저절로 자라난다는 경고. 불교의 「염(念)의 철학」.'),
    ('순명편 — 논어', '死 生 有 命  富 貴 在 天',
     '삶과 죽음은 명이요, 부귀는 하늘에 있다',
     '『논어』 인용. 운명의 영역과 노력의 영역을 구별하는 동양적 지혜 — 인사를 다하고 천명을 기다린다.'),
    ('효행편 — 시경', '父 兮 生 我  母 兮 鞠 我  欲 報 之 德  昊 天 罔 極',
     '아버지는 나를 낳으시고 어머니는 나를 기르시니, 그 은덕은 갚고자 해도 하늘처럼 끝이 없다',
     '『시경·소아·요아』 인용. 효의 출발은 「의무」가 아닌 「갚을 길 없는 은혜에 대한 자각」.'),
    ('정기편', '見 人 之 善 而 尋 己 之 善  見 人 之 惡 而 尋 己 之 惡',
     '남의 선을 보면 자기의 선을 찾고, 남의 악을 보면 자기의 악을 찾는다',
     '비판도 칭찬도 거울처럼 자기에게 돌린다 — 자기 점검의 핵심 자세.'),
    ('안분편', '知 足 者 貧 賤 亦 樂  不 知 足 者 富 貴 亦 憂',
     '만족할 줄 아는 자는 빈천해도 즐겁고, 만족할 줄 모르는 자는 부귀해도 근심한다',
     '행복은 가진 것의 크기가 아니라 만족할 줄 아는 마음의 크기에서 온다.'),
    ('안분편', '知 足 常 足  終 身 不 辱',
     '만족할 줄 알면 늘 만족하고, 평생 욕을 당하지 않는다',
     '욕망을 채우려는 사람은 욕을 당하고, 만족할 줄 아는 사람은 평생 자유롭다.'),
    ('계성편', '忍 一 時 之 忿  免 百 日 之 憂',
     '한순간의 분을 참으면 백일의 근심을 면한다',
     '「忍(인)」 한 글자가 백일의 인생을 결정한다. 분노 관리의 가장 짧고 강력한 격언.'),
    ('근학편 — 예기', '玉 不 琢 不 成 器  人 不 學 不 知 道',
     '옥은 다듬지 않으면 그릇이 되지 않고, 사람은 배우지 않으면 도를 모른다',
     '재능은 다듬어야 빛난다 — 학문이 사람을 사람답게 만든다는 『예기·학기』의 명제.'),
    ('근학편 — 주자', '少 年 易 老 學 難 成  一 寸 光 陰 不 可 輕',
     '소년은 쉽게 늙고 학문은 이루기 어렵다, 한 치의 시간이라도 가벼이 하지 말라',
     '주자(주희)의 「권학문(勸學文)」 — 시간의 한정성과 학문의 무한성 사이의 긴장.'),
    ('훈자편', '黃 金 滿 籯  不 如 敎 子 一 經',
     '황금이 한 광주리에 가득해도 자식에게 경전 한 권을 가르치는 것만 못하다',
     '재산은 무너지지만 가르침은 평생 — 자녀에게 무엇을 남길지에 대한 명심보감의 답.'),
    ('성심편', '水 至 淸 則 無 魚  人 至 察 則 無 徒',
     '물이 너무 맑으면 물고기가 살지 않고, 사람이 너무 살피면 친구가 없다',
     '엄격함의 한계 — 완벽주의는 관계를 끊는다. 어른의 관용에 대한 동양적 통찰.'),
    ('성심편', '路 遙 知 馬 力  日 久 見 人 心',
     '길이 멀어야 말의 힘을 알고, 날이 오래되어야 사람의 마음을 안다',
     '사람을 판단하는 가장 정확한 도구는 시간 — 빠른 친밀과 진짜 우정의 구별.'),
    ('성심편', '疑 人 莫 用  用 人 勿 疑',
     '의심스러우면 쓰지 말고, 쓴 사람은 의심하지 말라',
     '용인(用人)의 가장 짧은 원칙. 한 번 결정했으면 끝까지 — 의심하며 부리는 일은 양쪽 모두를 망친다.'),
    ('준례편', '若 要 人 重 我  無 過 我 重 人',
     '남이 나를 존중하기를 바라거든, 내가 남을 존중하는 것만 한 것이 없다',
     '존중의 원리 — 자기가 받고 싶은 대접을 먼저 남에게 행하라. 「황금률」의 동양적 표현.'),
    ('언어편', '口 是 傷 人 斧  言 是 割 舌 刀',
     '입은 사람을 상하게 하는 도끼요, 말은 혀를 베는 칼이다',
     '말의 무서움 — 함부로 휘두른 말은 남도 베고 결국 자기 혀도 벤다. 언어편 한 권의 정수.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅷ. 조선과 동아시아 ==============
SEC8 = 'Ⅷ. 조선과 동아시아'

@S(SEC8)
def viii_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '조선 — 「천자문 다음, 소학 앞」 자리의 책',
              '서당 교육의 표준 교재 · 500년 민족 정신의 토양')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 조선 초부터 서당 표준 교재 — 천자문을 마친 7~8세 아동이 처음 만나는 「뜻 있는 글」',
         {'font_size': 18, 'space_before': 4}),
        ('● 1454년 「신간대자명심보감(新刊大字明心寶鑑)」 — 큰 글자로 인쇄된 아동 학습용 판본',
         {'font_size': 17, 'space_before': 12}),
        ('● 사대부·서민·여성까지 — 신분을 가리지 않은 전 계층의 교양서',
         {'font_size': 17, 'space_before': 10}),
        ('● 임진왜란 후 — 일본 출판인들이 가져가 무로마치 후기·에도 초기 교육서로 활용',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 일제 강점기에도 서당의 마지막 교재로 살아남음',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 「사상의 책」보다 「감각의 책」 — 한 줄 한 줄이 평생의 판단 순간마다 떠오르는 나침반',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC8)
def viii_world(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '동아시아와 세계로 — 1592년 마닐라의 기적',
              '동양 고전이 처음으로 유럽 언어로 옮겨지다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1592년 — 스페인 도미니코회 선교사 후안 코보(Juan Cobo)',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('     필리핀 마닐라에서 화교 사회를 통해 명심보감을 접함',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 「Beng Sim Po Cam (寶鑑)」 — 명심보감을 스페인어·중국어로 대역 번역',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     동양 고전이 유럽어로 번역된 최초의 사례 중 하나',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 1595년 — 펠리페 2세에게 헌상 — 마드리드 국립도서관 소장',
         {'font_size': 17, 'space_before': 12}),
        ('● 일본·베트남·필리핀까지 — 동아시아 한자 문화권 전체의 「민중 교과서」',
         {'font_size': 17, 'space_before': 10}),
        ('● 추적이 고려에서 엮은 책이 4세기를 건너 스페인 왕실에 도달한 여정',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC8)
def viii_decline(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '왜 오늘 명심보감의 존재감은 약해졌는가',
              '근대 교육이 잃어버린 「매일의 한 구절」 습관')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 근대 학제의 도입 — 한문 교재가 교과서에서 밀려남',
         {'font_size': 17, 'space_before': 4}),
        ('● 「사상의 체계」 중심 교육 — 「격언의 축적」이라는 옛 방식이 낯설어짐',
         {'font_size': 17, 'space_before': 10}),
        ('● 일부 구절의 보수성(부행편의 사덕 등)이 시대와 맞지 않는다는 평가',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 그러나 — 「매일의 한 구절」이라는 독서법 자체는 여전히 유효',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 자기계발서·명상 앱·morning quote — 모두 명심보감 독서법의 현대적 변주',
         {'font_size': 17, 'space_before': 10}),
        ('● 「머리로 배우는 책」이 아닌 「마음에 담는 책」 — 그 방식이 오늘 다시 필요해진 이유',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


# ============== Ⅸ. 오늘 다시 펼치는 이유 ==============
SEC9 = 'Ⅸ. 오늘 다시 펼치는 이유'

@S(SEC9)
def ix_modern(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '오늘의 결핍 지도 — 명심보감이 채울 수 있는 자리')
    items = [
        ('정보 과잉, 지혜 결핍', '하루 수천 개의 짧은 글을 보지만 마음에 남는 한 줄은 드물다'),
        ('감정 제어의 어려움',   '계성편이 가르치는 「忍一時之忿」 — 한순간의 분이 인생을 흔든다'),
        ('비교 경쟁의 피로',     '안분편 「知足者貧賤亦樂」 — 비교에서 자기 자신으로 시선을 돌리기'),
        ('말의 가벼움',         '언어편 「口是傷人斧」 — 디지털 시대에 가장 위험해진 영역'),
        ('우정과 신뢰의 모호함', '교우편 「日久見人心」 — 시간이 진짜 친구를 가려준다'),
        ('자녀 교육의 혼란',     '훈자편 「黃金滿籯 不如敎子一經」 — 무엇을 물려줄 것인가'),
    ]
    for i, (cat, desc) in enumerate(items):
        y = Inches(2.4 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.8), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(3.8), Inches(0.6),
                    cat, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.7), y + Inches(0.05), Inches(8.2), Inches(0.55),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC9)
def ix_howread(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '명심보감 오늘 읽는 법 — 「매일의 한 구절」',
              '책의 이름이 곧 독서법')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 처음부터 끝까지 읽지 않는다 — 처음 펼친 자리부터 한 구절씩',
         {'font_size': 18, 'space_before': 6}),
        ('● 마음에 닿는 구절은 메모 — 직접 손으로 적어두면 더 오래 남는다',
         {'font_size': 18, 'space_before': 10}),
        ('● 하루의 시작 또는 끝에 한 구절 — 책상 위·휴대폰 잠금화면도 좋다',
         {'font_size': 18, 'space_before': 10}),
        ('● 자녀와 함께 — 일주일에 한 구절씩 외우고 그 의미를 이야기 나눈다',
         {'font_size': 18, 'space_before': 10}),
        ('● 시간이 지나며 같은 구절의 새로운 의미가 보인다 — 「축적되는 거울」의 효과',
         {'font_size': 18, 'space_before': 10}),
        ('● 800년 전 추적의 의도 그대로 — 「이론이 아닌 감각, 체계가 아닌 습관」',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC9)
def ix_seven(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '명심보감이 오늘 우리에게 일러주는 7가지')
    items = [
        '1. 선과 악은 결국 갚아진다 — 큰 시야로 보라',
        '2. 작은 선도 외면하지 말고, 작은 악도 가볍게 보지 말라',
        '3. 만족은 가진 것이 아니라 마음의 결에서 온다',
        '4. 한순간의 분이 백일의 후회 — 「忍」 한 글자가 어른을 만든다',
        '5. 남을 보고 자기를 돌아보라 — 비판도 칭찬도 거울이다',
        '6. 말은 도끼와 칼 — 짧은 말 한마디가 천금이자 흉기',
        '7. 진짜 친구는 시간만이 가려준다 — 빨리 친해진 사람을 빨리 믿지 말라',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.3 + i * 0.65)
        add_textbox(slide, Inches(0.7), y, Inches(12.0), Inches(0.55),
                    txt, font_size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅹ. 마무리 ==============
SEC10 = 'Ⅹ. 마무리'

@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '명심보감, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 명심보감은 고려 추적이 1305년 엮고, 명 범립본이 1393년 증보한 격언집.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 20편 — 세계관 → 관계 → 자기 수양 → 사회 → 실천 매뉴얼.',
         {'font_size': 18, 'space_before': 8}),
        ('● 유·불·도 삼교가 자연스럽게 녹은 「지혜의 뷔페」.',
         {'font_size': 18, 'space_before': 8}),
        ('● 조선 500년 — 천자문 다음, 소학 앞 — 서당의 표준 교재.',
         {'font_size': 18, 'space_before': 8}),
        ('● 1592년 — 마닐라에서 스페인어로 번역, 동양 고전의 첫 유럽어 번역.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「마음을 매일 닦는 거울」 — 800년의 인격 교과서.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.1),
                '明 心 寶 鑑',
                font_size=92, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(3.6), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.6),
                '마 음 을  밝 혀 주 는  보 배 로 운  거 울',
                font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.5),
                '爲 善 者 天 報 之 以 福',
                font_size=24, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
                '선을 행하는 자에게 하늘은 복으로 갚는다',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '— 명심보감 계선편 첫 구절',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\명심보감.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
