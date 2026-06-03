# -*- coding: utf-8 -*-
"""
순자 발표자료 재작성 스크립트
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.italic = opts.get('italic', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=32, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=15, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ---------- 1. 표지 ----------
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.6),
                '荀 子', font_size=110, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                'Xunzi · 순자', font_size=24, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
                '순황(荀況)의 32편 — 선진 유가의 마지막 거대 체계',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
                '전국시대 말기 (BC 313~238경) · 32편 · 성악설·예론·천론의 체계',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ---------- 2. 목차 ----------
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.8), Inches(0.7),
                '목 차', font_size=36, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.3), Inches(12.8))
    items = [
        ('Ⅰ', '개요 — 순자란 무엇인가'),
        ('Ⅱ', '32편의 구성'),
        ('Ⅲ', '핵심 사상 ① — 성악설(性惡說)'),
        ('Ⅳ', '핵심 사상 ② — 예론(禮論)과 천론(天論)'),
        ('Ⅴ', '핵심 사상 ③ — 수양론과 정명론'),
        ('Ⅵ', '명구절 10선'),
        ('Ⅶ', '순자의 구조적 특징'),
        ('Ⅷ', '현대적 의의'),
        ('Ⅸ', '다른 사상가와의 비교'),
        ('Ⅹ', '마무리'),
    ]
    top = 1.7
    for num, title in items:
        add_textbox(slide, Inches(1.2), Inches(top), Inches(1.0), Inches(0.4),
                    num, font_size=22, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.4), Inches(top), Inches(10.0), Inches(0.4),
                    title, font_size=20, color=INK)
        top += 0.5


# ---------- Ⅰ. 개요 ----------
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '순자(荀子)란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.5),
                '전국시대 말기 유학자 순황의 사상을 담은 유가 철학서',
                font_size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.5),
                '논어·맹자가 어록 중심이라면, 순자는 최초의 체계적 논설문',
                font_size=17, color=ACCENT, align=PP_ALIGN.CENTER)
    nums = [('32', '편(篇)'), ('3대', '선진 유학 저작'), ('BC 313', '~238경')]
    for i, (n, lbl) in enumerate(nums):
        x = 1.8 + i * 3.5
        add_textbox(slide, Inches(x), Inches(4.4), Inches(3.0), Inches(1.0),
                    n, font_size=58, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(5.6), Inches(3.0), Inches(0.5),
                    lbl, font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '논어·맹자·순자 — 선진 유학의 3대 저작',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER, bold=True)


@S('Ⅰ. 개요')
def s_xunhuang(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '순황(荀況, BC 313~238경)',
              '— 자(字)는 경(卿), 후세 사람들이 "순경(荀卿)"으로 존칭')
    lines = [
        ('출신', {'bold': True, 'font_size': 18, 'color': ACCENT}),
        ('  조(趙)나라 출신 — 북방 학풍',
         {'font_size': 18}),
        ('활동', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 12}),
        ('  제(齊)나라 직하학궁(稷下學宮)에서 세 차례 제주(祭酒, 학장) 역임',
         {'font_size': 18}),
        ('  당대 최고 학자로 인정받은 사상가',
         {'font_size': 14, 'color': SUB}),
        ('만년', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 12}),
        ('  초(楚)나라 난릉령(蘭陵令) 역임 후 저술에 전념',
         {'font_size': 18}),
        ('  난릉에서 생을 마침',
         {'font_size': 14, 'color': SUB}),
    ]
    add_paragraphs(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.5),
                   lines, line_spacing=1.4)


@S('Ⅰ. 개요')
def s_disciples(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '직하학궁과 두 제자', '— 유가에서 길러낸 법가의 거장들')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(1.2), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(11.9), Inches(0.5),
                '직하학궁(稷下學宮)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(3.05), Inches(11.9), Inches(0.4),
                '전국시대 제(齊)나라가 운영한 동아시아 최초의 국립 종합 학술 기관',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    # 두 제자
    pupils = [
        ('이사', '李斯', '진(秦)나라 승상 — 진시황의 천하통일을 보좌',
         '분서(焚書)·문자 통일 등 통치 제도화'),
        ('한비자', '韓非子', '법가 사상의 집대성자 — 한비자(韓非子) 저술',
         '법(法)·술(術)·세(勢)의 통합 이론'),
    ]
    for i, (name, hanmun, role, contrib) in enumerate(pupils):
        x = 0.6 + i * 6.1
        add_filled_rect(slide, Inches(x), Inches(3.8), Inches(5.8), Inches(3.0), PALE)
        add_textbox(slide, Inches(x), Inches(3.95), Inches(5.8), Inches(0.8),
                    hanmun, font_size=44, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.85), Inches(5.8), Inches(0.5),
                    name, font_size=20, bold=True, color=INK,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(5.45), Inches(5.4), Inches(0.5),
                    role, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(5.95), Inches(5.4), Inches(0.8),
                    contrib, font_size=14, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '순자가 법가 사상에 결정적 영향을 미친 까닭 — 그의 사상사적 양면성',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_three_innovations(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '사상사적 위치 — 세 가지 혁신')
    innovations = [
        ('1', '인성론',  '맹자의 성선설에 대해 성악설(性惡說)을 주장',
         '유가 내부 논쟁 촉발 — 인성에 대한 깊은 사유의 시작'),
        ('2', '예론',    '예(禮)를 사회 질서의 근본 원리로 체계화',
         '추상적 덕(德)이 아닌 제도적 장치로서의 예'),
        ('3', '천론',    '하늘(天)을 자연 현상으로 규정',
         '동아시아 합리주의 전통의 수립 — 미신·신비주의에 대한 비판'),
    ]
    top = 2.3
    for num, field, claim, impact in innovations:
        add_filled_rect(slide, Inches(0.7), Inches(top), Inches(1.0), Inches(1.4), ACCENT)
        add_textbox(slide, Inches(0.7), Inches(top + 0.45), Inches(1.0), Inches(0.6),
                    num, font_size=36, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.9), Inches(top), Inches(2.5), Inches(1.4), PALE)
        add_textbox(slide, Inches(1.9), Inches(top + 0.5), Inches(2.5), Inches(0.5),
                    field, font_size=22, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(4.6), Inches(top + 0.15), Inches(8.3), Inches(0.5),
                    claim, font_size=17, bold=True, color=INK)
        add_textbox(slide, Inches(4.6), Inches(top + 0.75), Inches(8.3), Inches(0.5),
                    impact, font_size=14, color=SUB)
        top += 1.55


# ---------- Ⅱ. 구성 ----------
@S('Ⅱ. 구성')
def s_structure_overview(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 구성', page, total)
    add_title(slide, '32편의 구성 개관 — 네 가지 영역')
    blocks = [
        ('핵심 사상편', '1~9편',   '권학·수신·왕제 등',  '학문과 수양·왕제론',  ACCENT),
        ('경세 실천편', '10~16편', '부국·왕패·의병 등',  '경제·정치·군사',        RGBColor(0xA0, 0x40, 0x40)),
        ('철학 심화편', '17~23편', '천론·예론·정명·성악', '천관·예론·인식·언어',  RGBColor(0x70, 0x40, 0x60)),
        ('부록·잡편',  '24~32편', '군자·성상·격언 등',   '운문·일화·잠언',        SUB),
    ]
    top = 2.3
    for tag, scope, contents, theme, color in blocks:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.4), Inches(1.0), color)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(2.4), Inches(0.5),
                    tag, font_size=18, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.55), Inches(2.4), Inches(0.4),
                    scope, font_size=13,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(9.5), Inches(1.0), PALE)
        add_textbox(slide, Inches(3.4), Inches(top + 0.13), Inches(9.1), Inches(0.4),
                    contents, font_size=15, bold=True, color=INK)
        add_textbox(slide, Inches(3.4), Inches(top + 0.55), Inches(9.1), Inches(0.4),
                    theme, font_size=13, color=SUB)
        top += 1.15


@S('Ⅱ. 구성')
def s_part1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 구성', page, total)
    add_title(slide, '1~16편 — 핵심 사상편 + 경세 실천편')
    rows = [
        ('1',  '권학(勸學)',   '학문의 권면 — 순자 사상의 출발점'),
        ('2',  '수신(修身)',   '인격 수양의 기본'),
        ('3',  '불구(不苟)',   '구차하게 살지 않는 군자의 태도'),
        ('4',  '영욕(榮辱)',   '명예와 치욕의 분별'),
        ('5',  '비상(非相)',   '관상술 비판'),
        ('6',  '비십이자(非十二子)', '당대 12사상가 비판'),
        ('7',  '중니(仲尼)',   '공자와 중궁의 평가'),
        ('8',  '유효(儒效)',   '진정한 유학의 효용'),
        ('9',  '왕제(王制)',   '왕도 정치 제도론'),
        ('10', '부국(富國)',   '경제론 — 부국의 길'),
        ('11', '왕패(王覇)',   '왕도와 패도의 분별'),
        ('12', '군도(君道)',   '군주의 도리'),
        ('13', '신도(臣道)',   '신하의 도리'),
        ('14', '치사(致士)',   '인재 등용'),
        ('15', '의병(議兵)',   '군사론 — 손자병법과 다른 관점'),
        ('16', '강국(彊國)',   '강국이 되는 길'),
    ]
    top = 2.0
    # 2열 배치
    for i, (num, name, desc) in enumerate(rows):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = top + row * 0.6
        add_textbox(slide, Inches(x), Inches(y), Inches(0.5), Inches(0.4),
                    num, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.5), Inches(y), Inches(1.8), Inches(0.4),
                    name, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(x + 2.4), Inches(y), Inches(3.9), Inches(0.4),
                    desc, font_size=12, color=SUB)


@S('Ⅱ. 구성')
def s_part2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 구성', page, total)
    add_title(slide, '17~32편 — 철학 심화편 + 부록·잡편')
    rows = [
        ('17', '천론(天論)',   '하늘은 자연 — 순자 철학의 정점'),
        ('18', '정론(正論)',   '잘못된 통설의 교정'),
        ('19', '예론(禮論)',   '예의 기원과 의미'),
        ('20', '악론(樂論)',   '음악의 도덕적·사회적 기능'),
        ('21', '해폐(解蔽)',   '인식의 폐단을 푸는 인식론'),
        ('22', '정명(正名)',   '이름과 실재 — 언어 철학'),
        ('23', '성악(性惡)',   '성악설의 본 장'),
        ('24', '군자(君子)',   '군자에 대한 잠언'),
        ('25', '성상(成相)',   '리듬을 가진 운문체 격언'),
        ('26', '부(賦)',       '부 형식의 다섯 편'),
        ('27', '대략(大略)',   '여러 짧은 잠언 모음'),
        ('28', '유좌(宥坐)',   '공자의 일화 모음'),
        ('29', '자도(子道)',   '효(孝)와 도(道)'),
        ('30', '법행(法行)',   '본받아 행할 일들'),
        ('31', '애공(哀公)',   '애공과 공자의 문답'),
        ('32', '요문(堯問)',   '요(堯)의 물음 — 마무리'),
    ]
    top = 2.0
    for i, (num, name, desc) in enumerate(rows):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = top + row * 0.6
        add_textbox(slide, Inches(x), Inches(y), Inches(0.5), Inches(0.4),
                    num, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.5), Inches(y), Inches(1.8), Inches(0.4),
                    name, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(x + 2.4), Inches(y), Inches(3.9), Inches(0.4),
                    desc, font_size=12, color=SUB)


# ---------- Ⅲ. 성악설 ----------
@S('Ⅲ. 성악설')
def s_seongak(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 성악설', page, total)
    add_title(slide, '성악설(性惡說) — 본성은 악하다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '人 之 性 惡   其 善 者 僞 也',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '인지성악 기선자위야',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 성악편 제23', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
                '사람의 본성은 악하다',
                font_size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6),
                '선한 것은 인위(僞)다',
                font_size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '※ 맹자의 성선설(性善說)에 대한 정면 반박',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅲ. 성악설')
def s_seong_wi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 성악설', page, total)
    add_title(slide, '성(性) vs 위(僞) — 두 개념의 구분')
    # 좌
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '性 — 성(性)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(3.1), Inches(5.9), Inches(1.4),
                '性',
                font_size=130, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(4.8), Inches(5.3), Inches(1.9), [
        ('타고난 본성 — 자연적·본능적', {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('이기적 욕망, 감각적 욕구', {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('배고프면 먹고 싶고, 추우면 따뜻하고 싶다', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    # 우
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '僞 — 위(僞)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.1), Inches(5.9), Inches(1.4),
                '僞',
                font_size=130, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(4.8), Inches(5.3), Inches(1.9), [
        ('인위적 노력 — 의식적·후천적', {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('예의(禮義)·사양(辭讓)·문화(文化)', {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('"거짓"이 아니라 "人 + 爲(사람의 행위)"', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '인간이 인간다워지는 것은 "성"이 아니라 "위"의 산물이다',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅲ. 성악설')
def s_hwaseong_giwi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 성악설', page, total)
    add_title(slide, '화성기위(化性起僞) — 순자 사상의 핵심 공식')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '化 性 起 僞',
                font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '화성기위 — "본성을 변화시켜 인위를 일으킨다"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    # 흐름도: 性 → 僞 → 禮義 → 法度
    steps = [('性', '본성'), ('僞', '인위'), ('禮義', '예의'), ('法度', '법도')]
    box_w = 2.4
    gap = 0.5
    total_w = box_w * 4 + gap * 3
    start_x = (13.333 - total_w) / 2
    for i, (han, kor) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_filled_rect(slide, Inches(x), Inches(4.6), Inches(box_w), Inches(1.5), PALE)
        add_textbox(slide, Inches(x), Inches(4.75), Inches(box_w), Inches(0.8),
                    han, font_size=44, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(5.65), Inches(box_w), Inches(0.4),
                    kor, font_size=15, color=INK, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_x = x + box_w
            add_textbox(slide, Inches(arrow_x), Inches(5.05), Inches(0.5), Inches(0.5),
                        '▶', font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '성인이 본성을 변화시켜 인위를 일으키고, 인위가 예의를 낳고, 예의가 법도를 제정한다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅲ. 성악설')
def s_dojiin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 성악설', page, total)
    add_title(slide, '도지인가이위우(塗之人可以爲禹) — 교육 낙관론')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '塗 之 人   可 以 爲 禹',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '도지인 가이위우 — "길가의 사람도 우(禹) 임금이 될 수 있다"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 성악편 제23', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.7), Inches(4.8), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(4.95), Inches(11.3), Inches(1.8), [
        ('성악설의 비관 → 교육 낙관론으로 전환',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('본성이 악하더라도, 학문과 수양을 거치면',
         {'font_size': 17, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('누구나 성왕(聖王)이 될 수 있다',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('맹자의 성선설과 다른 출발점 — 같은 평등적 결론',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)


# ---------- Ⅳ. 예론과 천론 ----------
@S('Ⅳ. 예·천론')
def s_yelron_origin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 예·천론', page, total)
    add_title(slide, '예(禮)의 기원 — 왜 예가 필요한가')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.42), Inches(11.9), Inches(0.5),
                '禮 起 於 何 也',
                font_size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    # 4단계 흐름
    flow = [
        ('욕망',     '인간에게는 욕망이 있다'),
        ('충족 불가', '욕망을 다 충족할 수는 없다'),
        ('추구',     '그래도 사람은 끝없이 추구한다'),
        ('다툼',     '한도가 없으면 서로 다투게 된다'),
        ('禮 필요',  '그래서 예(禮)가 일어났다'),
    ]
    top = 3.3
    for i, (tag, desc) in enumerate(flow):
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.4), Inches(0.55),
                        ACCENT if i == 4 else PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.12), Inches(2.4), Inches(0.4),
                    tag, font_size=16, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF) if i == 4 else INK,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(9.5), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.4), Inches(top + 0.13), Inches(9.1), Inches(0.4),
                    desc, font_size=15, color=INK)
        top += 0.65
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '— 예론편 제19 · 욕망 조절의 사회적 필요에서 예가 발생했다',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER, bold=True)


@S('Ⅳ. 예·천론')
def s_three_functions(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 예·천론', page, total)
    add_title(slide, '예의 세 가지 기능 — 양(養)·별(別)·문(文)')
    funcs = [
        ('養', '양', '욕망을 적절히 충족시킴',
         '禮者 養也 — 예는 기름(충족)이다.\n욕망의 억압이 아니라 절도 있는 만족.'),
        ('別', '별', '사회적 분별을 세움',
         '귀천·노소·친소의 분별이 있어야\n사회 질서가 유지된다.'),
        ('文', '문', '문화적 형식을 부여함',
         '단순한 본능적 행위에\n인간다운 의미와 형식을 더한다.'),
    ]
    for i, (han, kor, role, desc) in enumerate(funcs):
        x = 0.5 + i * 4.3
        add_filled_rect(slide, Inches(x), Inches(2.3), Inches(4.1), Inches(4.7), PALE)
        add_textbox(slide, Inches(x), Inches(2.5), Inches(4.1), Inches(1.4),
                    han, font_size=120, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.2), Inches(4.1), Inches(0.5),
                    kor, font_size=20, bold=True, color=INK,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.8), Inches(4.1), Inches(0.4),
                    role, font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(5.4), Inches(3.7), Inches(1.5),
                    desc, font_size=13, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '禮者 人道之極也 — 예는 인간 도의 극치다',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 예·천론')
def s_cheonron(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 예·천론', page, total)
    add_title(slide, '천론(天論) — 하늘은 자연이다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '天 行 有 常',
                font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.9),
                '不 爲 堯 存   不 爲 桀 亡',
                font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
                '천행유상 불위요존 불위걸망',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.4),
                '— 천론편 제17', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
                '하늘의 운행에는 일정한 법칙이 있어,',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                '요(堯) 때문에 존재하지도, 걸(桀) 때문에 없어지지도 않는다',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)


@S('Ⅳ. 예·천론')
def s_je_cheonmyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 예·천론', page, total)
    add_title(slide, '제천명이용지(制天命而用之) — 자연을 활용하라')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '制 天 命 而 用 之',
                font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '제천명이용지 — "천명(天命)을 제어하여 이용하라"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
                '— 천론편 제17', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.55), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.7), Inches(4.85), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.05), Inches(11.3), Inches(1.7), [
        ('자연에 대한 수동적 경외 → 적극적 활용으로의 전환',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('하늘을 신비화하지 않고 자연 법칙으로 인식할 때',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('비로소 인간은 자연을 이해하고 활용할 수 있다',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('중국 사상사에서 가장 혁신적인 합리주의적 자연관',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)


# ---------- Ⅴ. 수양론과 정명론 ----------
@S('Ⅴ. 수양·정명')
def s_jeok(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 수양·정명', page, total)
    add_title(slide, '적(積)의 사상 — 축적의 힘')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '積 土 成 山   風 雨 興 焉',
                font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '적토성산 풍우흥언',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
                '— 권학편 제1', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
                '흙을 쌓아 산을 이루면 바람과 비가 일어난다',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.4), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.2), [
        ('不 積 蹞 步   無 以 至 千 里',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('부적규보 무이지천리 — "반 걸음을 쌓지 않으면 천 리에 이를 수 없다"',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅴ. 수양·정명')
def s_cheongchulam(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 수양·정명', page, total)
    add_title(slide, '청출어람(靑出於藍) — 학문의 가능성')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '靑 取 之 於 藍   而 靑 於 藍',
                font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '청취지어람 이청어람',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
                '— 권학편 제1 (순자가 쓴 원문이 사자성어의 출전)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '푸른색은 쪽(藍)에서 취하였으나, 쪽보다 더 푸르다',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6),
                '학문과 수양을 통해 본래의 자기를 넘어설 수 있다는 통찰',
                font_size=17, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                '※ 후대에 "제자가 스승을 능가한다"는 의미로 확장 사용',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅴ. 수양·정명')
def s_heoiljeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 수양·정명', page, total)
    add_title(slide, '허일이정(虛壹而靜) — 인식론의 세 조건', '— 해폐편 제21')
    items = [
        ('虛', '허', '마음을 비움', '기존 지식·선입견에 얽매이지 않는 개방성'),
        ('壹', '일', '집중',       '분산되지 않는 전일적(專一的) 주의력'),
        ('靜', '정', '고요함',     '감정에 흔들리지 않는 마음의 안정'),
    ]
    for i, (han, eum, role, desc) in enumerate(items):
        x = 0.5 + i * 4.3
        add_filled_rect(slide, Inches(x), Inches(2.3), Inches(4.1), Inches(4.0), PALE)
        add_textbox(slide, Inches(x), Inches(2.5), Inches(4.1), Inches(1.4),
                    han, font_size=130, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.2), Inches(4.1), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.75), Inches(4.1), Inches(0.5),
                    role, font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(5.3), Inches(3.7), Inches(1.2),
                    desc, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '대청명(大淸明) — 사물을 있는 그대로 파악하는 인식의 경지',
                font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅴ. 수양·정명')
def s_jeongmyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 수양·정명', page, total)
    add_title(slide, '정명론(正名論) — 이름은 약속이다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '名 無 固 宜   約 之 以 命',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '명무고의 약지이명',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
                '— 정명편 제22', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '이름에는 본래부터 정해진 적합함이 없으며,',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.6),
                '사회적 약속(約)으로 비로소 그 의미가 정해진다',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.85), PALE)
    add_textbox(slide, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
                '約 定 俗 成  — 약정속성', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
                '현대 언어 철학의 규약주의(conventionalism)와 통하는 선구적 통찰',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ---------- Ⅵ. 명구절 ----------
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=46):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=22, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('Ⅵ. 명구절 (1/10)',
    '學 不 可 以 已',
    '학 불 가 이 이',
    '배움은 그만둘 수 없다',
    '권학편 제1 — 순자 32편의 첫 문장', hanmun_size=72), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (2/10)',
    '靑 取 之 於 藍   而 靑 於 藍',
    '청취지어람 이청어람',
    '푸른색은 쪽(藍)에서 취하였으나 쪽보다 푸르다 — 청출어람(靑出於藍)의 출전',
    '권학편 제1', hanmun_size=34), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (3/10)',
    '積 土 成 山   風 雨 興 焉',
    '적토성산 풍우흥언',
    '흙을 쌓아 산을 이루면 바람과 비가 일어난다 — 축적의 힘',
    '권학편 제1', hanmun_size=36), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (4/10)',
    '木 受 繩 則 直   金 就 礪 則 利',
    '목수승즉직 금취려즉리',
    '나무는 먹줄을 받으면 곧아지고, 쇠는 숫돌에 갈면 날카로워진다 — 교육의 효과',
    '권학편 제1', hanmun_size=32), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (5/10)',
    '天 行 有 常\n不 爲 堯 存   不 爲 桀 亡',
    '천행유상 · 불위요존 불위걸망',
    '하늘의 운행에는 법칙이 있어 요(堯) 때문에 있지도, 걸(桀) 때문에 없지도 않다',
    '천론편 제17', hanmun_size=28), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (6/10)',
    '君 者 舟 也   庶 人 者 水 也\n水 則 載 舟   水 則 覆 舟',
    '군자주야 서인자수야 · 수즉재주 수즉복주',
    '군주는 배요, 백성은 물 — 물은 배를 띄우기도 하고 뒤집기도 한다',
    '왕제편 제9 (재주복주, 載舟覆舟의 출전)', hanmun_size=24), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (7/10)',
    '人 之 性 惡   其 善 者 僞 也',
    '인지성악 기선자위야',
    '사람의 본성은 악하며, 선한 것은 인위(僞)다',
    '성악편 제23 — 성악설의 본 명제', hanmun_size=36), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (8/10)',
    '塗 之 人   可 以 爲 禹',
    '도지인 가이위우',
    '길가의 사람도 우(禹) 임금이 될 수 있다 — 누구나 성인이 될 수 있다는 평등적 수양론',
    '성악편 제23', hanmun_size=44), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (9/10)',
    '禮 起 於 何 也\n求 而 無 度 量 分 界   則 不 能 不 爭',
    '예기어하야 · 구이무도량분계 즉불능부쟁',
    '예는 어디서 일어났는가 — 욕망 추구에 한도가 없으면 다툼을 피할 수 없기에',
    '예론편 제19', hanmun_size=22), 'Ⅵ. 명구절'))

SLIDES.append((make_quote_slide('Ⅵ. 명구절 (10/10)',
    '虛 壹 而 靜   謂 之 大 淸 明',
    '허일이정 위지대청명',
    '비우고 집중하며 고요하면 대청명(大淸明)의 경지에 이른다',
    '해폐편 제21', hanmun_size=32), 'Ⅵ. 명구절'))


# ---------- Ⅶ. 구조 ----------
@S('Ⅶ. 구조')
def s_essay_form(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 구조', page, total)
    add_title(slide, '유가 최초의 체계적 논설문')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(5.7), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.5), Inches(5.7), Inches(0.5),
                '논어 · 맹자의 형식', font_size=18, bold=True, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(1.0), Inches(3.2), Inches(5.1), Inches(3.5), [
        ('어록 · 대화 중심', {'font_size': 17, 'bold': True}),
        ('', {'font_size': 6}),
        ('단편적 발언과 일화', {'font_size': 15, 'space_before': 6}),
        ('체계적 논증 없이', {'font_size': 15}),
        ('주제별로 흩어진 구성', {'font_size': 15}),
    ], line_spacing=1.4)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.7), Inches(4.5), PALE)
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.7), Inches(0.5),
                '순자의 형식', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(3.2), Inches(5.1), Inches(3.5), [
        ('비유 → 논증 → 결론의 논설문', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('', {'font_size': 6}),
        ('하나의 주제를 끝까지', {'font_size': 15, 'space_before': 6}),
        ('논리적으로 전개', {'font_size': 15}),
        ('각 편이 독립된 논문 구조', {'font_size': 15}),
    ], line_spacing=1.4)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '중국 산문 발전사의 결정적 전환점 — 후대 사상 논변의 표준 형식',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅶ. 구조')
def s_metaphors(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 구조', page, total)
    add_title(slide, '비유의 풍부함 — 네 가지 대표 비유')
    metaphors = [
        ('積土成山',   '적토성산',   '축적의 힘',                   '흙을 쌓아 산이 되면 비바람이 인다'),
        ('靑出於藍',   '청출어람',   '학문의 발전',                  '쪽보다 푸른 청 — 본래를 넘어섬'),
        ('蓬生麻中',   '봉생마중',   '환경의 영향',                  '쑥도 삼밭에서 자라면 곧아진다'),
        ('載舟覆舟',   '재주복주',   '민심의 결정적 중요성',         '물(백성)은 배(군주)를 띄우기도 뒤집기도'),
    ]
    top = 2.3
    for han, eum, theme, desc in metaphors:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.6), Inches(1.0), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.15), Inches(2.6), Inches(0.5),
                    han, font_size=22, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.62), Inches(2.6), Inches(0.4),
                    eum, font_size=12,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.4), Inches(top), Inches(2.6), Inches(1.0), PALE)
        add_textbox(slide, Inches(3.4), Inches(top + 0.32), Inches(2.6), Inches(0.5),
                    theme, font_size=17, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.2), Inches(top), Inches(6.5), Inches(1.0),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.4), Inches(top + 0.32), Inches(6.2), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 1.13


@S('Ⅶ. 구조')
def s_flow(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 구조', page, total)
    add_title(slide, '32편의 논리 흐름')
    flow = [
        ('학문론',   '1편',     '권학 — 모든 것의 출발은 배움'),
        ('수양론',   '2~4편',   '수신·불구·영욕 — 인격 형성'),
        ('비판론',   '5~6편',   '비상·비십이자 — 사이비 학설의 비판'),
        ('유학론',   '7~8편',   '중니·유효 — 진정한 유학의 효용'),
        ('정치론',   '9~16편',  '왕제·부국·왕패·의병 — 경세 실천'),
        ('철학론',   '17~23편', '천론·예론·정명·성악 — 사상의 정점'),
        ('부록',     '24~32편', '군자·성상·격언 — 보충 자료'),
    ]
    top = 2.3
    for i, (name, scope, desc) in enumerate(flow):
        add_filled_rect(slide, Inches(0.7), Inches(top), Inches(2.0), Inches(0.6),
                        PALE)
        add_textbox(slide, Inches(0.7), Inches(top + 0.1), Inches(2.0), Inches(0.4),
                    name, font_size=16, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.9), Inches(top + 0.1), Inches(1.5), Inches(0.4),
                    scope, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(4.6), Inches(top + 0.1), Inches(8.4), Inches(0.4),
                    desc, font_size=15, color=INK)
        if i < len(flow) - 1:
            add_textbox(slide, Inches(1.5), Inches(top + 0.55), Inches(0.5), Inches(0.3),
                        '▼', font_size=10, color=SUB, align=PP_ALIGN.CENTER)
        top += 0.65


# ---------- Ⅷ. 현대적 의의 ----------
def make_modern_slide(title, kor_subtitle, lines):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, 'Ⅷ. 현대적 의의', page, total)
        add_title(slide, title, kor_subtitle)
        add_paragraphs(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(4.6),
                       lines, line_spacing=1.45, font_size=18)
    return renderer


SLIDES.append((make_modern_slide(
    '현대 ① — 현실주의적 인간관',
    '성악설이 현대 사회과학의 출발 가설과 통한다',
    [
        ('성악설(性惡說)', {'font_size': 22, 'bold': True, 'color': ACCENT}),
        ('  인간을 이상화하지 않고 욕망과 이기심을 인정',
         {'font_size': 16, 'space_before': 4}),
        ('', {'font_size': 8}),
        ('현대적 연결', {'font_size': 20, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('• 경제학의 합리적 행위자(rational agent) 전제',
         {'font_size': 16, 'space_before': 4}),
        ('• 정치학의 견제와 균형 원리',
         {'font_size': 16}),
        ('• 제도 설계의 출발점이 되는 인간관',
         {'font_size': 16}),
    ]), 'Ⅷ. 현대적 의의'))

SLIDES.append((make_modern_slide(
    '현대 ② — 제도주의(Institutionalism)',
    '예론(禮論) — 개인의 덕성이 아닌 제도의 힘',
    [
        ('예(禮) = 사회 제도', {'font_size': 22, 'bold': True, 'color': ACCENT}),
        ('  개인의 도덕적 결심이 아니라',
         {'font_size': 16, 'space_before': 4}),
        ('  사회적 장치로 질서를 유지',
         {'font_size': 16}),
        ('', {'font_size': 8}),
        ('현대적 연결', {'font_size': 20, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('• 신제도주의 경제학(New Institutional Economics)',
         {'font_size': 16, 'space_before': 4}),
        ('• "사람을 바꾸기보다 룰을 바꿔라"는 조직 설계의 원칙',
         {'font_size': 16}),
    ]), 'Ⅷ. 현대적 의의'))

SLIDES.append((make_modern_slide(
    '현대 ③ — 과학적 세계관',
    '천론(天論) — 자연을 신비화하지 않는 합리주의',
    [
        ('하늘 = 자연 법칙', {'font_size': 22, 'bold': True, 'color': ACCENT}),
        ('  미신·재이설(災異說)에 대한 비판',
         {'font_size': 16, 'space_before': 4}),
        ('  자연에 대한 능동적·실용적 태도',
         {'font_size': 16}),
        ('', {'font_size': 8}),
        ('현대적 연결', {'font_size': 20, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('• 과학적 세계관의 동아시아적 원형',
         {'font_size': 16, 'space_before': 4}),
        ('• 環境工學·應用科學의 사상적 기반',
         {'font_size': 16}),
    ]), 'Ⅷ. 현대적 의의'))

SLIDES.append((make_modern_slide(
    '현대 ④ — 인지심리학과 비판적 사고',
    '해폐(解蔽) — 인지 편향을 깨는 지혜',
    [
        ('해폐(解蔽)', {'font_size': 22, 'bold': True, 'color': ACCENT}),
        ('  한쪽 측면에 가려져 큰 이치에 어두워지는 것',
         {'font_size': 16, 'space_before': 4}),
        ('  마음을 비우고 집중하여 고요할 때 바로 본다',
         {'font_size': 15, 'color': SUB}),
        ('', {'font_size': 8}),
        ('현대적 연결', {'font_size': 20, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('• 인지편향(cognitive bias) 연구와 직결',
         {'font_size': 16, 'space_before': 4}),
        ('• 비판적 사고(critical thinking) 교육의 동양적 원형',
         {'font_size': 16}),
    ]), 'Ⅷ. 현대적 의의'))

SLIDES.append((make_modern_slide(
    '현대 ⑤ — 성장 마인드셋',
    '적(積)의 사상 — 작은 축적이 위대한 성취로',
    [
        ('적토성산(積土成山)', {'font_size': 22, 'bold': True, 'color': ACCENT}),
        ('  반 걸음이 쌓여 천 리에 이른다',
         {'font_size': 16, 'space_before': 4}),
        ('  점진적 발전론 — 노력과 시간의 누적',
         {'font_size': 16}),
        ('', {'font_size': 8}),
        ('현대적 연결', {'font_size': 20, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('• Carol Dweck의 성장 마인드셋(Growth Mindset)',
         {'font_size': 16, 'space_before': 4}),
        ('• 습관 형성·복리(複利) 효과의 동양적 원형',
         {'font_size': 16}),
    ]), 'Ⅷ. 현대적 의의'))

SLIDES.append((make_modern_slide(
    '현대 ⑥ — 언어 철학과 기호학',
    '정명론(正名論) — 이름은 사회적 약속이다',
    [
        ('약정속성(約定俗成)', {'font_size': 22, 'bold': True, 'color': ACCENT}),
        ('  이름과 사물의 관계는 본질이 아니라 규약',
         {'font_size': 16, 'space_before': 4}),
        ('  사회적 합의로 의미가 정해진다',
         {'font_size': 16}),
        ('', {'font_size': 8}),
        ('현대적 연결', {'font_size': 20, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('• 소쉬르의 기호의 자의성(arbitrariness of sign)',
         {'font_size': 16, 'space_before': 4}),
        ('• 분석철학의 규약주의(conventionalism)',
         {'font_size': 16}),
    ]), 'Ⅷ. 현대적 의의'))


# ---------- Ⅸ. 비교 ----------
@S('Ⅸ. 비교')
def s_compare(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 비교', page, total)
    add_title(slide, '공자 · 맹자 · 순자 — 유가 삼대 사상가 비교')
    rows = [
        ('인성론',     '성상근(性相近)',  '성선설(性善說)', '성악설(性惡說)'),
        ('핵심 덕목',  '인(仁)',           '인의(仁義)',     '예(禮)'),
        ('수양 방법',  '학·극기복례',     '확충(擴充)',     '화성기위·적(積)'),
        ('천(天)관',   '도덕적 천',       '의지적 천',      '자연적 천'),
        ('정치론',     '덕치·정명',       '왕도·역성혁명',  '왕도·예법 병용'),
        ('문체',       '간결한 어록',     '웅변적 논변',    '체계적 논설문'),
        ('후대 영향',  '유가 정통',       '성리학',         '법가·예학(禮學)'),
    ]
    top = 1.95
    add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.2), Inches(0.55), SUB)
    add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(2.2), Inches(0.4),
                '항목', font_size=15, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    for i, name in enumerate(['공자(논어)', '맹자', '순자']):
        cx = 0.5 + 2.2 + i * 3.45
        color = ACCENT if name == '순자' else SUB
        add_filled_rect(slide, Inches(cx), Inches(top), Inches(3.45), Inches(0.55), color)
        add_textbox(slide, Inches(cx), Inches(top + 0.1), Inches(3.45), Inches(0.4),
                    name, font_size=16, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    row_h = 0.62
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        add_filled_rect(slide, Inches(0.5), Inches(y), Inches(2.2), Inches(row_h), PALE)
        add_textbox(slide, Inches(0.55), Inches(y + 0.12), Inches(2.1), Inches(0.4),
                    row[0], font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        for i in range(3):
            cx = 0.5 + 2.2 + i * 3.45
            add_filled_rect(slide, Inches(cx), Inches(y), Inches(3.45), Inches(row_h), bg)
            text_color = ACCENT if i == 2 else INK
            add_textbox(slide, Inches(cx + 0.05), Inches(y + 0.12),
                        Inches(3.35), Inches(0.4),
                        row[i + 1], font_size=14, color=text_color,
                        bold=(i == 2), align=PP_ALIGN.CENTER)


# ---------- Ⅹ. 마무리 ----------
@S('Ⅹ. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 순자')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.5), PALE)
    add_paragraphs(slide, Inches(1.1), Inches(2.6), Inches(11.1), Inches(4.1), [
        ('인간의 본성이 악하다고 진단하면서도',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('학문과 수양을 통해 누구나 성인이 될 수 있다고 확신한',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('— 합리주의적 사상가 —',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('하늘을 자연 법칙으로',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('예를 문명의 근본으로',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('이름을 사회적 약정으로',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('— 합리주의와 제도주의의 선구자 —',
         {'font_size': 20, 'bold': True, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
    ], line_spacing=1.2)


@S('Ⅹ. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
                '塗 之 人   可 以 爲 禹',
                font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                '도지인 가이위우', font_size=22, color=SUB,
                align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '"길가의 보통 사람도 우(禹) 임금이 될 수 있다"',
                font_size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
                '— 성악편 제23 · 본성이 악해도 누구나 성인에 이를 수 있다',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
                '감사합니다', font_size=28, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ---------- 빌드 ----------
total_pages = len(SLIDES)
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\순자_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
