# -*- coding: utf-8 -*-
"""
부모은중경(父母恩重經) 발표자료 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '동아시아 불교의 효경(孝經) · 어머니 사랑의 가장 절절한 경전',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '父 母 恩 重 經',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '부 모 은 중 경 · 대 방 편 불 보 은 경',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '父 兮 生 我   母 兮 鞠 我',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '아버지가 나를 낳으시고, 어머니가 나를 기르시니  — 시경',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '한역자 미상 · 십대은(十大恩) · 1,000년 동아시아인의 눈물을 적신 효경',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 두 경전의 만남'),
        ('Ⅱ.', '부모은중경의 정체'),
        ('Ⅲ.', '십대은(十大恩) — 부모의 열 가지 은혜'),
        ('Ⅳ.', '보은의 어려움'),
        ('Ⅴ.', '대방편불보은경 — 부처의 효'),
    ]
    items_right = [
        ('Ⅵ.', '두 경전이 함께 가르치는 것'),
        ('Ⅶ.', '한국에서의 위상 — 조선의 효경'),
        ('Ⅷ.', '명구 10선'),
        ('Ⅸ.', '오늘 우리가 배워야 할 교훈'),
        ('Ⅹ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 두 경전',
              '불교의 효 사상을 입체적으로 보여주는 한 쌍')
    rows = [
        ('부모은중경', '불설대보부모은중경 — 짧고 감성적'),
        ('대방편불보은경', '7권 9품 — 길고 교리적'),
        ('주체', '부모은중경 = 부모(어머니) / 대방편불보은경 = 부처'),
        ('성격', '동아시아 불교의 효경(孝經) — 효 사상의 집대성'),
        ('정신', '부모의 은혜는 측량 불가 — 십대은(十大恩)'),
        ('위상', '한국 1,000년 — 가장 많이 간행된 불교 경전'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.5), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.5), Inches(0.5),
                    k, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.4), y, Inches(9.5), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_essence(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '두 경전의 핵심 — 한 마디로',
              '부모 사랑의 가장 절절한 표현')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.6), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.7),
                '父 母 恩 重  難 報 之 德',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(12.0), Inches(0.5),
                '부모의 은혜는 무겁고, 갚기 어려운 덕이다',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.4),
                '어머니가 자식을 위해 베푸는 열 가지 은혜 — 십대은',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.5), Inches(12.0), Inches(2.5), [
        ('● 부모은중경 — 어머니 시선의 사랑 묘사', {'font_size': 16, 'space_before': 6}),
        ('● 대방편불보은경 — 부처(자식) 시선의 보은 이야기', {'font_size': 16, 'space_before': 6}),
        ('● 두 경전이 함께 — 불교 효 사상의 입체적 정수', {'font_size': 16, 'space_before': 6}),
        ('● 「**갚을 수 없음의 자각이 — 진짜 효의 시작**」', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅱ. 부모은중경의 정체 ==============
SEC2 = 'Ⅱ. 부모은중경의 정체'

@S(SEC2)
def ii_apocrypha(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '「위경(僞經)」이지만 「진짜 경(經)」',
              '인도 원전 없이 — 동아시아에서 형성된 자생 경전')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.8), [
        ('● 위경 논란의 이유', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   산스크리트어·빨리어 원전이 발견되지 않음', {'font_size': 14, 'color': SUB}),
        ('   인도 불교 전통에 없는 표현·논리가 다수', {'font_size': 14, 'color': SUB}),
        ('   중국 당대(唐代) 학자가 작성한 것으로 추정', {'font_size': 14, 'color': SUB}),
        ('● 그런데도 1,000년 사랑받은 이유', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('   메시지 자체가 너무도 강력하고 진실 — 어머니 사랑의 보편적 진실', {'font_size': 14, 'color': SUB}),
        ('   종교를 넘어 인간 보편의 정서에 닿음', {'font_size': 14, 'color': SUB}),
        ('● 「**동아시아 불교의 자생적 창작**」', {'font_size': 16, 'color': ACCENT, 'space_before': 12}),
        ('   인도 불교에 없던 「부모-자식 관계의 깊이」를 동아시아의 눈으로', {'font_size': 14, 'color': SUB}),
    ])


@S(SEC2)
def ii_dabhang(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '대방편불보은경 — 정통 인도 경전',
              '7권 9품 — 부처 자신의 효 이야기')
    rows = [
        ('정식 명칭', '대방편불보은경(大方便佛報恩經)'),
        ('성립', '후한~서진 시기 한역 (인도 원전 있는 정통)'),
        ('분량', '7권 9품 — 부모은중경에 비해 방대'),
        ('주체', '부처(석가모니) 자신의 전생 효행'),
        ('핵심 편', '효양품(孝養品) — 부처의 전생 효행 이야기'),
        ('메시지', '부처가 부처가 된 것은 — 효의 결과'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.2), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.2), Inches(0.5),
                    k, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.1), y, Inches(9.8), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅲ. 십대은 ==============
SEC3 = 'Ⅲ. 십대은'

@S(SEC3)
def iii_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '십대은 한 폭으로',
              '어머니가 자식에게 베푸는 열 가지 은혜')
    items = [
        '1. 회탐수호은(懷耽守護恩) — 잉태하여 지켜 주심',
        '2. 임산수고은(臨産受苦恩) — 해산의 고통',
        '3. 생자망우은(生子忘憂恩) — 자식 낳고 모든 근심 잊음',
        '4. 인고토감은(咽苦吐甘恩) — 쓴 것 삼키고 단 것 먹여줌',
        '5. 회건취습은(廻乾就濕恩) — 진자리 마른자리 가려 누임',
        '6. 유포양육은(乳哺養育恩) — 젖을 먹여 키워줌',
        '7. 세탁부정은(洗濯不淨恩) — 더러운 것을 씻어줌',
        '8. 원행억념은(遠行憶念恩) — 멀리 떠난 자식 걱정',
        '9. 위조악업은(爲造惡業恩) — 자식 위해 악업도 마다 않음',
        '10. 구경연민은(究竟憐愍恩) — 끝까지 가엾이 여기심',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.4 + i * 0.45)
        add_textbox(slide, Inches(0.9), y, Inches(12.0), Inches(0.4),
                    txt, font_size=14, color=INK)


def make_eun_slide(idx_total, hanja, korean, story, lesson):
    @S(SEC3)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, SEC3, n, t)
        add_title(slide, f'{idx_total} — {korean}', '어머니 사랑의 한 모습')
        add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.0), PALE)
        add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.0),
                    hanja,
                    font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_paragraphs(slide, Inches(0.7), Inches(3.6), Inches(12.0), Inches(2.0), [
            ('어머니의 모습', {'font_size': 14, 'bold': True, 'color': SUB}),
            (story, {'font_size': 15, 'color': INK}),
        ], line_spacing=1.4)
        add_filled_rect(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.0), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.85), Inches(11.6), Inches(0.9),
                    f'오늘 우리에게 — {lesson}',
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


euns = [
    ('1/10', '懷 耽 守 護 恩', '회탐수호은 — 잉태하여 지키신 은혜',
     '뱃속에 자식을 품은 열 달 — 한순간도 잊지 않고 보호하셨다. 한 점의 정(精)이 자식의 시작, 어머니의 몸이 곧 자식의 첫 집.',
     '나는 어머니 뱃속에서 열 달을 보호받았다 — 그것이 내 생명의 시작'),
    ('2/10', '臨 産 受 苦 恩', '임산수고은 — 해산의 고통을 견디신 은혜',
     '출산의 고통은 산을 옮기는 것보다 무겁고 — 뼈와 살이 찢기는 진통. 자식이 무사히 나오는 것만이 어머니의 유일한 소원.',
     '내가 태어난 것 자체가 — 어머니의 가장 큰 고통의 결과'),
    ('3/10', '生 子 忘 憂 恩', '생자망우은 — 자식 낳고 모든 근심을 잊으신 은혜',
     '그 모든 고통이 — 자식의 첫 울음 한 번에 사라진다. 자식의 얼굴을 보는 순간 모든 것을 잊는다.',
     '사랑이 고통을 압도하는 신비 — 어머니 사랑의 정수'),
    ('4/10', '咽 苦 吐 甘 恩', '인고토감은 — 쓴 것은 삼키고 단 것은 먹여 주신 은혜',
     '자기 입에 들어간 좋은 것도 자식 입에 넣어 주고 — 쓴 것·맛없는 것은 자기가 먹는다.',
     '자기 것을 자식의 것으로 여기는 — 무조건적 사랑'),
    ('5/10', '廻 乾 就 濕 恩', '회건취습은 — 진자리 마른자리 가려 누이신 은혜',
     '자식의 오줌으로 젖은 자리는 자기가 자고 — 마른 자리는 자식에게 양보. 춥고 불편한 것은 자기, 편안한 것은 자식.',
     '자기 불편을 — 자식의 편안함으로 바꾸는 사랑'),
    ('6/10', '乳 哺 養 育 恩', '유포양육은 — 젖을 먹여 키워 주신 은혜',
     '자기 피와 살에서 짜낸 젖 — 자식이 클수록 어머니의 몸은 야위어 간다. 어머니의 몸이 곧 자식의 양식.',
     '내 몸의 한 부분이 — 어머니 몸에서 왔다는 자각'),
    ('7/10', '洗 濯 不 淨 恩', '세탁부정은 — 더러운 것을 깨끗이 씻어 주신 은혜',
     '자식의 더러운 옷·기저귀·몸을 손수 씻음. 손이 트고 갈라져도 멈추지 않음.',
     '자기 손의 더러움을 — 자식의 깨끗함으로 바꾸는 노동'),
    ('8/10', '遠 行 憶 念 恩', '원행억념은 — 자식이 멀리 떠난 뒤에도 걱정하신 은혜',
     '자식이 집을 떠나도 한순간도 잊지 못함. 비 오면 비 걱정, 추우면 추위 걱정. 자식의 거리는 어머니의 마음에 닿지 않는다.',
     '나는 어디에 있든 — 어머니의 마음 안에 있다'),
    ('9/10', '爲 造 惡 業 恩', '위조악업은 — 자식 위해 악업도 마다 않으신 은혜',
     '자식 굶기지 않으려 거짓말도 하고 — 자식 입히려 도둑질도 마다 않음. 자기 영혼의 죄를 짊어져서라도 자식을 지키는 사랑.',
     '자식을 위해서라면 — 자기 업(業)을 더하면서까지의 사랑'),
    ('10/10', '究 竟 憐 愍 恩', '구경연민은 — 끝까지 가엾이 여기고 사랑하신 은혜',
     '자식이 백 살이 되어도 어머니에게는 어린아이. 어머니가 죽는 그 순간까지 자식 걱정. 죽어서도 끝나지 않는 사랑.',
     '가장 마지막까지 — 가장 깊은 은혜 — 죽음으로도 끊을 수 없다'),
]

for tag, hj, kr, st, ls in euns:
    make_eun_slide(tag, hj, kr, st, ls)


# ============== Ⅳ. 보은의 어려움 ==============
SEC4 = 'Ⅳ. 보은의 어려움'

@S(SEC4)
def iv_difficulty(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '「이렇게까지 해도 다 갚지 못한다」',
              '보은의 어려움 — 충격적 비유')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(3.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.6), Inches(2.8), [
        ('자식이 부모를 왼쪽 어깨에 메고, 오른쪽 어깨에 메고', {'font_size': 16, 'color': INK}),
        ('수미산을 백천 번 돌아 — 살갗이 닳고 뼈가 부서지고', {'font_size': 16, 'color': INK}),
        ('골수가 흘러내려도 — 부모의 은혜를 다 갚지 못한다.', {'font_size': 16, 'bold': True, 'color': ACCENT}),
        ('', {'font_size': 8}),
        ('자식이 부모를 위해 칼로 자기 살을 베고', {'font_size': 16, 'color': INK}),
        ('수억 겁 동안 부모를 위해 자기 살을 보시해도', {'font_size': 16, 'color': INK}),
        ('부모의 은혜를 다 갚지 못한다.', {'font_size': 16, 'bold': True, 'color': ACCENT}),
    ])
    add_textbox(slide, Inches(0.7), Inches(5.7), Inches(12.0), Inches(1.4),
                '이 무서운 묘사는 — 「갚을 수 없으니 포기하라」가 아니라\n「그만큼 깊으니 평생 갚으려고 노력하라」는 의미',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S(SEC4)
def iv_two_paths(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '보은의 두 길',
              '세속적 효행 + 불교적 보은')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.8), [
        ('● 세속적 효행 — 유교적 효와 일치', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   부모 봉양 · 부모 마음 편안케 · 살아 계실 때 정성 · 제사 정성', {'font_size': 14, 'color': SUB}),
        ('● 불교적 보은 — 부모은중경의 독자적 가르침', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('   부모를 위해 경전을 만들거나 베껴 쓰기', {'font_size': 14, 'color': SUB}),
        ('   부모의 이름으로 불상을 조성하거나 절을 지음', {'font_size': 14, 'color': SUB}),
        ('   부모의 명복을 위해 재(齋)를 올림', {'font_size': 14, 'color': SUB}),
        ('   부모를 도(道)에 들게 함이 가장 큰 효', {'font_size': 14, 'bold': True, 'color': ACCENT}),
        ('● 「**부모 봉양은 작은 효, 부모를 도에 들게 함은 큰 효**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
    ])


# ============== Ⅴ. 대방편불보은경 ==============
SEC5 = 'Ⅴ. 대방편불보은경'

@S(SEC5)
def v_buddha_filial(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '부처의 전생 효행 이야기',
              '효는 깨달음의 출발점')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.8), [
        ('● 발단 — 외도들의 비방', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   「부처는 어머니를 떠나 출가했으니 불효자다」', {'font_size': 14, 'color': SUB}),
        ('● 부처의 답 — 전생 효행 이야기', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('   수다나(須大拏) 태자 — 자기 살을 베어 부모를 살린 전생', {'font_size': 14, 'color': SUB}),
        ('   삼십삼천 — 돌아가신 어머니 마야 부인을 위해 천상에 올라 설법', {'font_size': 14, 'color': SUB}),
        ('   살타 태자 — 굶주린 호랑이에게 몸을 던진 보살행', {'font_size': 14, 'color': SUB}),
        ('● 핵심 — 「**부처가 부처가 된 것은 — 효의 결과**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('   효는 단순한 가족 윤리가 아닌 — 깨달음의 출발점', {'font_size': 14, 'color': SUB}),
    ])


@S(SEC5)
def v_filial_universe(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '효의 우주적 차원',
              '대방편불보은경의 깊이')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 효는 — 가족 안의 도덕에 머물지 않는다', {'font_size': 17, 'space_before': 6}),
        ('● 효는 — 불성(佛性)의 발현이다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 효는 — 깨달음으로 가는 가장 깊은 길', {'font_size': 17, 'space_before': 10}),
        ('● 부모은중경(어머니 시선) + 대방편불보은경(자식 시선)', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   두 경전이 합쳐져 — 효의 입체적 정수', {'font_size': 14, 'color': SUB}),
    ])


# ============== Ⅵ. 두 경전이 가르치는 것 ==============
SEC6 = 'Ⅵ. 두 경전이 가르치는 것'

@S(SEC6)
def vi_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '두 경전의 통합 메시지',
              '효의 모든 차원')
    items = [
        '부모의 은혜는 측량할 수 없다 — 십대은과 전생담',
        '그 은혜에 보답하는 것은 사람의 가장 근본 도리',
        '봉양의 효 (살아 계실 때 정성) — 작은 효',
        '부모를 도(道)에 들게 함 — 큰 효',
        '갚을 수 없음의 자각이 — 진짜 효의 시작',
        '효는 가족 윤리가 아닌 — 영적 책임',
        '불교가 유교 효의 보완 — 부모 시선 + 자식 시선',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.5 + i * 0.6)
        add_textbox(slide, Inches(0.9), y, Inches(0.6), Inches(0.5),
                    f'{i+1}.', font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.5),
                    txt, font_size=16, color=INK)


# ============== Ⅶ. 한국 위상 ==============
SEC7 = 'Ⅶ. 한국 위상'

@S(SEC7)
def vii_korea1(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '조선의 효경(孝經)',
              '1,000년 한국인의 눈물을 적신 경전')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 고려 말~조선 말까지 — 약 80여 종 간행', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   가장 많이 간행된 불교 경전 중 하나', {'font_size': 14, 'color': SUB}),
        ('● 한글 창제 후 언해본(諺解本) 다수 — 일반 백성도 읽음', {'font_size': 17, 'space_before': 10}),
        ('● 사찰뿐 아니라 민간 가정의 필수 경전', {'font_size': 17, 'space_before': 10}),
        ('● 유교 사회였지만 — 불교 박해 시대에도 폭넓게 유통', {'font_size': 17, 'space_before': 10}),
        ('   부모 사랑이라는 보편 정서 앞에서 종교 차이가 무의미', {'font_size': 14, 'color': SUB}),
    ])


@S(SEC7)
def vii_jeongjo(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '정조와 용주사본',
              '왕의 효심이 만든 가장 아름다운 판본')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 정조(正祖)가 1796년 용주사 창건과 함께 간행을 직접 명함', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('● 김홍도(金弘道)의 그림이 첨가된 변상도(變相圖)', {'font_size': 17, 'space_before': 10}),
        ('● 미술사적·인쇄사적 가치 매우 큼', {'font_size': 17, 'space_before': 10}),
        ('● 정조 자신의 사도세자에 대한 깊은 효심', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   「아버지에게 못 다한 효를 — 이 경전으로 갚는다」', {'font_size': 14, 'color': SUB}),
        ('● 21~22장의 판화 — 글을 모르는 백성도 십대은의 의미를 깨달음', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


# ============== Ⅷ. 명구 10선 ==============
SEC8 = 'Ⅷ. 명구 10선'

def make_quote(idx_total, hanja, korean, comment):
    @S(SEC8)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, SEC8, n, t)
        add_title(slide, f'명구 {idx_total}', '부모은중경 — 한 줄에 담긴 사랑')
        add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.3), PALE)
        add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.3),
                    hanja,
                    font_size=24, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(1.0),
                    korean,
                    font_size=18, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.7), Inches(5.1), Inches(12.0), Inches(1.6), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.1), Inches(11.6), Inches(1.6),
                    comment,
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


quotes = [
    ('1/10', '父 母 恩 重   難 報 之 德',
     '부모의 은혜는 무겁고, 갚기 어려운 덕이다',
     '부모은중경의 핵심. 이 한 마디가 모든 가르침의 정수.'),
    ('2/10', '父 兮 生 我   母 兮 鞠 我',
     '아버지가 나를 낳으시고, 어머니가 나를 기르셨다',
     '시경 인용. 부모 은혜의 가장 단순한 표현.'),
    ('3/10', '哀 哀 父 母   生 我 劬 勞',
     '슬프고 슬프다 부모님이여, 나를 낳아 고생하셨도다',
     '시경 인용. 부모의 노고에 대한 자식의 가장 깊은 탄식.'),
    ('4/10', '母 行 路 苦   父 行 役 勞',
     '어머니는 길 가는 고통, 아버지는 일하는 수고',
     '부모가 자식을 위해 겪는 모든 노고의 압축.'),
    ('5/10', '昊 天 罔 極',
     '하늘처럼 끝이 없다',
     '시경 인용. 부모 은혜의 무한함 — 끝이 없는 하늘처럼.'),
    ('6/10', '生 我 養 我   出 入 腹 我',
     '나를 낳고 기르시며, 들고 나며 나를 안고 다니심',
     '어머니의 일상 — 매 순간 자식과 함께.'),
    ('7/10', '欲 報 之 德   昊 天 罔 極',
     '그 은혜에 보답하고자 하나, 하늘처럼 끝이 없도다',
     '시경 인용. 갚을 수 없음의 자각이 — 진짜 효의 시작.'),
    ('8/10', '風 樹 之 嘆',
     '나무는 고요하고자 하나 바람이 그치지 않고, 자식은 봉양하고자 하나 부모는 기다려 주지 않는다',
     '효의 시간을 놓치지 말라 — 오늘 미루지 말라는 절박한 호소.'),
    ('9/10', '小 孝 奉 養   大 孝 道 行',
     '봉양은 작은 효, 부모를 도(道)에 들게 함은 큰 효',
     '효의 두 차원 — 봉양 + 영적 책임.'),
    ('10/10', '不 報 父 母 恩   是 為 大 失',
     '부모의 은혜를 갚지 않는 것 — 그것이 가장 큰 잘못이다',
     '효를 잊는 것이 — 사람으로서 가장 큰 결함.'),
]

for tag, hj, kr, cm in quotes:
    make_quote(tag, hj, kr, cm)


# ============== Ⅸ. 오늘 우리에게 ==============
SEC9 = 'Ⅸ. 오늘 우리에게'

@S(SEC9)
def ix_today1(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '1 — 부모 살아 계실 때 하라',
              '풍수지탄(風樹之嘆)의 절박함')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.6), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.6), Inches(1.4), [
        ('나무는 고요하고자 하나 바람이 그치지 않고,', {'font_size': 17, 'color': INK}),
        ('자식은 봉양하고자 하나 부모는 기다려 주지 않는다.', {'font_size': 17, 'bold': True, 'color': ACCENT}),
    ])
    add_paragraphs(slide, Inches(0.7), Inches(4.3), Inches(12.0), Inches(2.5), [
        ('● 부모은중경의 가장 무거운 가르침', {'font_size': 17, 'space_before': 6}),
        ('● 「**오늘 미루지 말라**」 — 절박한 호소', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 오늘의 전화 한 통 · 안부 한 마디 · 잠깐의 방문이 — 곧 효', {'font_size': 17, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today2(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '2 — 갚을 수 없음의 자각이 진짜 효',
              '의무가 아닌 마음의 자세')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 효를 「**의무**」로만 보면 무거운 짐이 된다', {'font_size': 17, 'space_before': 6}),
        ('● 「**이 은혜는 다 갚을 수 없다**」는 자각에서 시작하면', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 평생의 작은 행위들이 — 모두 의미 있는 보답이 된다', {'font_size': 17, 'space_before': 10}),
        ('● 다 갚지 못해도 좋다 — 갚으려는 마음 자체가 효', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 매일의 작은 마음 씀 — 그것이 부모은중경의 정수', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today3(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '3 — 자기 자식에게 보여주는 효',
              '효는 가르치는 게 아니라 보여주는 것')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 부모은중경 가르침의 가장 깊은 실천', {'font_size': 17, 'space_before': 6}),
        ('● 자기 자식이 보는 데서 — 자기 부모에게 효를 다하는 것', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 부모 사랑하는 모습을 보지 못한 자식은 — 효를 모른다', {'font_size': 17, 'space_before': 10}),
        ('● 자기가 보여 준 만큼만 — 자식이 따라 한다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 「**효는 가르치는 것이 아니라 — 보여 주는 것**」', {'font_size': 16, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅹ. 마무리 ==============
SEC10 = 'Ⅹ. 마무리'

@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '부모은중경이 일러주는 7가지',
              '한 폭으로 정리')
    items = [
        '부모의 은혜는 헤아릴 수 없다 — 十大恩',
        '갚을 수 없음의 자각이 — 진짜 효의 시작',
        '오늘 미루지 말라 — 風樹之嘆',
        '작은 효 (봉양) + 큰 효 (부모를 도에 들게 함)',
        '자기 자식에게 보여주는 효',
        '무조건적 사랑의 모범으로서의 부모',
        '하늘처럼 끝이 없다 — 昊天罔極',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.5 + i * 0.6)
        add_textbox(slide, Inches(0.9), y, Inches(0.6), Inches(0.5),
                    f'{i+1}.', font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.5),
                    txt, font_size=16, color=INK)


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                '부모은중경 — 어머니 사랑을 가장 절절하게 그린 경전',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '昊 天 罔 極',
                font_size=140, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
                '하늘처럼 끝이 없다 — 부모의 은혜',
                font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.8), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
                '오늘, 부모님께 전화 한 통이라도 드리는 것이 — 곧 효',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)


# ============== 빌드 ==============
total = len(SLIDES)
for i, (fn, sec) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    fn(slide, i, total)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\부모은중경.pptx'
prs.save(out_path)
print(f'생성 완료: {out_path}  슬라이드 수: {total}')
