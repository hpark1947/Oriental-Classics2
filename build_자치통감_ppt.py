# -*- coding: utf-8 -*-
"""
자치통감(資治通鑑) 발표자료 — 전면 보강판 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
16편 각 부분 깊이 읽기 · 사마광 생애 · 후대 영향 망라
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '제왕학의 절대 교본 · 동양 1,000년의 「판단의 학교」',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '資 治 通 鑑',
                font_size=88, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '자 치 통 감',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '鑑 於 往 事  有 資 於 治 道',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '사마광(司馬光, 1019~1086) · 19년 편찬 · 294권 · 1,362년',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '주(周) 위열왕 BC 403 → 후주(後周) 세종 AD 959 · 「臣光曰」 218회의 직접 평론',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 자치통감이란'),
        ('Ⅱ.', '편자 사마광 — 19년의 침잠'),
        ('Ⅲ.', '편찬 과정 — 「叢目 → 長編 → 定本」'),
        ('Ⅳ.', '16편의 흐름 — 1,362년 한눈에'),
        ('Ⅴ.', '16편 각 부분 깊이 읽기'),
        ('Ⅵ.', '사관 — 「資治」와 「鑑」'),
    ]
    items_right = [
        ('Ⅶ.', '「臣光曰」 218회의 직접 평론'),
        ('Ⅷ.', '명장면 10선'),
        ('Ⅸ.', '명구 12선'),
        ('Ⅹ.', '후대 영향과 한국 수용'),
        ('Ⅺ.', '오늘 우리에게'),
        ('Ⅻ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 자치통감')
    rows = [
        ('편자',     '사마광(司馬光, 1019~1086) — 송 명재상·사학자'),
        ('보조',     '유서(劉恕)·유반(劉攽)·범조우(范祖禹)·사마강(사마광 子)'),
        ('편찬',     '1066~1084 · 19년 — 영종이 명, 신종이 책명 하사'),
        ('분량',     '294권 + 목록 30권 + 고이(考異) 30권 · 약 300만 자'),
        ('시대 폭',  '주 위열왕 23년(BC 403) ~ 후주 세종 6년(AD 959) · 1,362년'),
        ('체제',     '편년체(編年體) · 16개 왕조'),
        ('서명 뜻',  '「治에 도움 되도록(資) 옛일을 통합적으로(通) 거울 삼음(鑑)」'),
        ('사평',     '「신광왈(臣光曰)」 218회 — 사마광의 직접 평론'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.1), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '서명의 뜻 — 「資 · 治 · 通 · 鑑」',
              '신종(神宗)이 직접 명명·御製序 작성')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5),
                '鑑 於 往 事  有 資 於 治 道',
                font_size=32, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.4),
                '옛일을 거울 삼아 다스림의 도에 도움이 있게 한다 — 신종 御製序',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    cols = [
        ('資', '자', '도움이 되다'),
        ('治', '치', '다스림'),
        ('通', '통', '통하다 — 통사(通史)'),
        ('鑑', '감', '거울 — 以史爲鑑'),
    ]
    for i, (han, kor, desc) in enumerate(cols):
        x = Inches(0.7 + i * 3.1)
        y = Inches(4.5)
        add_filled_rect(slide, x, y, Inches(2.9), Inches(1.0), ACCENT)
        add_textbox(slide, x, y, Inches(2.9), Inches(1.0),
                    han, font_size=44, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, y + Inches(1.0), Inches(2.9), Inches(0.4),
                    kor, font_size=14, color=SUB,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.5), Inches(2.9), Inches(0.4),
                    desc, font_size=13, color=INK,
                    align=PP_ALIGN.CENTER)


@S(SEC1)
def i_unique(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '사기 vs 자치통감 — 두 정점의 비교',
              '「인간의 역사」 vs 「통치의 역사」')
    cols = [
        ('史 記', '사기 (사마천)',
         '기전체(紀傳體)\n인물 중심 — 본기·세가·열전\n\nBC 2600 ~ BC 100\n130권 약 52만 자\n\n「인간의 역사」\n태사공왈', INK),
        ('資 治 通 鑑', '자치통감 (사마광)',
         '편년체(編年體)\n시간 중심 — 연·월·일 순\n\nBC 403 ~ AD 959 (1,362년)\n294권 약 300만 자\n\n「통치의 역사」\n신광왈 (218회)', ACCENT),
    ]
    for i, (han, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    label, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.5), Inches(3.6), Inches(4.9), Inches(3.5),
                       [(body, {'font_size': 14, 'color': INK})], line_spacing=1.5)


# ============== Ⅱ. 사마광 ==============
SEC2 = 'Ⅱ. 사마광'

@S(SEC2)
def ii_profile(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '사마광(司馬光) — 송 명재상·사학자',
              '시호 「文正」 — 신하 최고의 시호')
    rows = [
        ('자(字)',   '군실(君實)'),
        ('호(號)',   '우수(迂叟) — 자호 「우원한 노인」'),
        ('시호',     '문정(文正) — 신하가 받을 수 있는 최고의 시호'),
        ('생몰',     '1019~1086 (향년 68세)'),
        ('출생',     '산서성 하현(夏縣)'),
        ('관직',     '한림학사·어사중승·추밀부사·재상'),
        ('대표 일화', '7세 때 「사마광 격옹(擊甕)」 — 친구 구함의 결단'),
        ('별명',     '「司馬君子」 — 군자다움의 화신'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.1), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC2)
def ii_geyong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '7세의 결단 — 「사마광 격옹(司馬光擊甕)」',
              '한 사람의 평생 성격을 어린 시절 한 장면으로 압축')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 사마광이 7세 때 친구들과 정원에서 놂',
         {'font_size': 18, 'space_before': 4}),
        ('● 한 아이가 큰 물독에 빠짐 — 다른 아이들은 모두 놀라 도망',
         {'font_size': 18, 'space_before': 12, 'color': SUB}),
        ('● 어린 사마광 — 큰 돌을 들어 독을 깨뜨려 친구 구함',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 「순간의 판단과 결단」의 표상 — 「**격옹구우(擊甕救友)**」',
         {'font_size': 17, 'space_before': 12, 'font_name': 'Batang'}),
        ('● 송대 이래 중국 어린이 교과서의 단골 이야기',
         {'font_size': 16, 'space_before': 10}),
        ('● 오늘날에도 중국·한국·일본 초등 교재에 「司馬光擊甕」으로 실림',
         {'font_size': 16, 'space_before': 10, 'color': SUB, 'font_name': 'Batang'}),
        ('● 한 인물의 평생 성격을 어린 시절의 한 장면으로 압축한 동양적 인물 묘사',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC2)
def ii_wangan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '왕안석과의 대립 — 「19년 침잠」의 토대',
              '정치적 좌절이 곧 자치통감의 탄생 조건')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1069년 — 신종(神宗) 즉위, 왕안석(王安石)을 등용해 신법(新法) 추진',
         {'font_size': 17, 'space_before': 4}),
        ('● 청묘법·모역법·시역법·보갑법·균수법 — 부국강병의 거대한 개혁',
         {'font_size': 16, 'space_before': 8, 'color': SUB}),
        ('● 사마광의 비판 — 「개혁의 의도는 좋으나 실행 방식이 백성을 해친다」',
         {'font_size': 17, 'space_before': 12}),
        ('● 두 거인은 본래 친구 — 「**여왕개보서(與王介甫書)**」 세 편이 동양 정치 서신의 명문',
         {'font_size': 16, 'space_before': 10, 'color': SUB, 'font_name': 'Batang'}),
        ('● 1071년 — 사마광이 정치적으로 패배 → 변경 떠나 낙양(洛陽) 칩거',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 15년간 통감 편찬에만 몰두 — 자호 「우수(迂叟)」',
         {'font_size': 17, 'space_before': 10}),
        ('● 사마천이 궁형을 사기로 승화시켰듯, 사마광은 정치 패배를 통감으로 승화',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
        ('● 1086년 — 통감 완성 2년 후 재상 복귀 8개월 만에 사망',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅲ. 편찬 과정 ==============
SEC3 = 'Ⅲ. 편찬 과정'

@S(SEC3)
def iii_three(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '3단계 작업법 — 「叢目 → 長編 → 定本」',
              '현대 역사학의 「실증 → 종합 → 해석」을 11세기에 체계화')
    rows = [
        ('1단계', '叢目 총목 (자료 카드)',
         '각 분담 학자가 자기 시대의 모든 사료에서 사건을 카드로 만듦.\n출전 명확히 — 카드 하나가 한 사건.'),
        ('2단계', '長編 장편 (풍성한 초고)',
         '총목 카드들을 시간 순으로 배열해 풍성한 초고로.\n모든 자료 일단 포함 — 모순이 있으면 「考異」로 따로.'),
        ('3단계', '定本 정본 (사마광의 최종 정리)',
         '사마광이 장편을 모두 읽고 직접 압축·정리.\n중복 제거·문장 통일·「臣光曰」 사평 추가.\n300만 자 장편이 약 100만 자 정본으로.'),
    ]
    for i, (step, title, body) in enumerate(rows):
        y = Inches(2.4 + i * 1.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.2), Inches(1.35), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.2), Inches(1.35),
                    step, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(2.1), y, Inches(3.4), Inches(1.35), PALE)
        add_textbox(slide, Inches(2.1), y, Inches(3.4), Inches(1.35),
                    title, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_paragraphs(slide, Inches(5.7), y + Inches(0.15), Inches(7.2), Inches(1.2),
                       [(body, {'font_size': 13, 'color': INK})], line_spacing=1.4)


@S(SEC3)
def iii_team(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '보조 편찬자들 — 「사학 드림팀」',
              '당대 최고의 전문가 4인의 분업')
    rows = [
        ('劉恕 유서', '1032~1078', '위진남북조', '가장 광범위한 사료 박학자 · 평생 친구'),
        ('劉攽 유반', '1023~1089', '양한(兩漢)',  '한대(漢代) 연구의 전문가'),
        ('范祖禹 범조우', '1041~1098', '당(唐)',  '당사 연구의 권위자'),
        ('司馬康 사마강', '? ~ ?',    '정리·교정',  '사마광의 친아들 — 부친의 손발'),
    ]
    for i, (name, life, era, role) in enumerate(rows):
        y = Inches(2.4 + i * 1.0)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.8), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(2.8), Inches(0.4),
                    name, font_size=16, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), y + Inches(0.45), Inches(2.8), Inches(0.4),
                    life, font_size=12, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(3.7), y, Inches(2.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.7), y, Inches(2.5), Inches(0.85),
                    era, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.4), y + Inches(0.1), Inches(6.5), Inches(0.7),
                    role, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC3)
def iii_kaoyi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '考異(고이) — 사료 비판의 정점',
              '독일 랑케보다 800년 앞선 실증 사학의 정수')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 『통감고이(通鑑考異) 30권』 — 자치통감의 핵심 부록',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 같은 사건에 여러 사료가 다른 기록 → 어떤 자료를 채택했고 왜 버렸는지 일일이 밝힘',
         {'font_size': 17, 'space_before': 12}),
        ('● 현대 「사료 비판(source criticism)」의 11세기 체계화',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 독일 랑케(Leopold von Ranke)의 실증주의 역사학보다 800년 앞섬',
         {'font_size': 16, 'space_before': 10}),
        ('● 모택동의 평 — 「고이가 통감의 가장 위대한 부분이다」',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
        ('     「한 줄을 쓰기 위해 사마광이 얼마나 많은 자료를 읽었는지를 보여준다」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 사마광 — 「사료가 다르면 옳고 그름을 가려야 한다. 가리지 않으면 사학이 아니다」',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


# ============== Ⅳ. 16편의 흐름 ==============
SEC4 = 'Ⅳ. 16편의 흐름'

@S(SEC4)
def iv_table(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '16편 1,362년 — 한 폭으로 보기',
              '주 위열왕 BC 403 → 후주 세종 AD 959')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.45), INK)
    headers = [('No.', 0.7), ('편명', 2.4), ('권수', 0.9), ('시대', 3.3), ('핵심 사건', 4.7)]
    x = Inches(0.7)
    for label, w in headers:
        add_textbox(slide, x, Inches(2.2), Inches(w), Inches(0.45),
                    label, font_size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(w)
    rows = [
        ('1', '周紀 주기', '5', 'BC 403~256', '삼가분진·상앙·소진장의'),
        ('2', '秦紀 진기', '3', 'BC 255~207', '진시황·분서갱유·항우유방'),
        ('3', '漢紀 한기', '60', 'BC 206~AD 219', '초한쟁패·문경의 치·왕망·삼국 정립'),
        ('4', '魏紀 위기', '10', 'AD 220~265', '백제성 탁고·제갈량 출사표·사마의'),
        ('5', '晉紀 진기', '40', 'AD 266~419', '팔왕의 난·영가의 난·비수대전'),
        ('6', '宋紀 송기', '16', 'AD 420~478', '원가의 치·종친 살육'),
        ('7', '齊紀 제기', '10', 'AD 479~501', '영명체·동혼후 광기'),
        ('8', '梁紀 양기', '22', 'AD 502~556', '양 무제 48년·후경의 난'),
        ('9', '陳紀 진기', '10', 'AD 557~588', '옥수후정화·수의 통일'),
        ('10', '隋紀 수기', '8', 'AD 589~617', '개원의 치·대운하·고구려 정벌'),
        ('11', '唐紀 당기', '81', 'AD 618~906', '정관·측천·개원·안사·번진환관'),
        ('12', '後梁紀 후량', '6', 'AD 907~922', '주온의 당 시해·골육상잔'),
        ('13', '後唐紀 후당', '8', 'AD 923~935', '이존욱·석경당'),
        ('14', '後晉紀 후진', '6', 'AD 936~946', '兒皇帝·연운16주 양도'),
        ('15', '後漢紀 후한', '4', 'AD 947~950', '가장 짧은 왕조·황포가신'),
        ('16', '後周紀 후주', '5', 'AD 951~959', '곽위·세종 시영의 30년 계획'),
    ]
    for i, (no, name, vol, era, event) in enumerate(rows):
        y = Inches(2.65 + i * 0.29)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.29), bg)
        add_textbox(slide, Inches(0.7), y, Inches(0.7), Inches(0.29),
                    no, font_size=11, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.4), y, Inches(2.4), Inches(0.29),
                    name, font_size=11, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(3.8), y, Inches(0.9), Inches(0.29),
                    vol, font_size=11, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.7), y, Inches(3.3), Inches(0.29),
                    era, font_size=10, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(8.0), y, Inches(4.7), Inches(0.29),
                    event, font_size=10, color=INK,
                    anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅴ. 16편 깊이 읽기 ==============
SEC5 = 'Ⅴ. 16편 깊이 읽기'

def make_chapter_slide(num, total, han, kor, vol_period, key_event,
                        original, modern, sapyeong):
    @S(SEC5)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC5} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{han}  ({kor})',
                    font_size=26, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    vol_period, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.6), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.5), [
            (original, {'font_size': 17, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.4),
                    '◆ 핵심 사건·인물', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.2), Inches(12.0), Inches(1.2),
                       [(key_event, {'font_size': 14, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.4),
                    '◆ 신광왈(臣光曰)·시대의 교훈', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.3),
                       [(sapyeong, {'font_size': 14, 'color': INK})], line_spacing=1.35)


CHAPTERS = [
    ('周 紀 주기', '주기', '권1~5 · 5권 · BC 403~256 · 148년 — 전국시대',
     '天 子 之 職 莫 大 於 禮  禮 莫 大 於 分  分 莫 大 於 名',
     '천자의 직책은 예보다 크지 않고, 예는 분(分), 분은 명(名)보다 크지 않다 — 권1 첫 사평',
     '삼가분진(三家分晉) · 위 문후·상앙의 변법 · 소진장의 합종연횡 · 인상여염파 문경지교 · 범저 원교근공',
     '권1 첫 「臣光曰」 — 정명(正名)의 위대한 선언. 「예가 무너진 시작이 곧 천하 혼란의 시작」. 자치통감 전체 사관의 출발.'),
    ('秦 紀 진기', '진기', '권6~8 · 3권 · BC 255~207 · 49년 — 통일과 멸망',
     '亡 國 不 可 以 復 存  死 者 不 可 以 復 生',
     '망한 나라는 되살릴 수 없고, 죽은 자는 다시 살릴 수 없다 — 사마광이 진의 빠른 멸망을 평하며',
     '여불위 「기화가거」 · 한비자의 죽음 · 형가의 진왕 암살 · 진 통일 BC 221 · 분서갱유 · 진승오광 · 항우·유방',
     '「정복은 무력으로, 통치는 인덕(仁德)으로 — 시황제가 이 한 마디를 몰랐다」. 빠른 영광이 곧 빠른 멸망.'),
    ('漢 紀 한기', '한기', '권9~68 · 60권 · BC 206~AD 219 · 425년 — 한 왕조 전체',
     '夫 運 籌 帷 幄 之 中  決 勝 千 里 之 外',
     '장막 안에서 계책을 짜 천 리 밖의 승부를 결정한다 — 한고조의 장량 평',
     '초한쟁패 · 한고조 삼걸(장량·소하·한신) · 한신 토사구팽 · 문경의 치 · 한 무제·죄기조 · 사마천 사기 · 왕망의 가면 · 광무제·후한 · 십상시·황건적·동탁·관도·적벽',
     '60권에 한 왕조 전체. 「창업·수성·부패·멸망」의 완벽한 사이클. 송이 가장 두려워한 「외척·환관」의 가장 큰 거울.'),
    ('魏 紀 위기', '위기', '권69~78 · 10권 · AD 220~265 · 46년 — 삼국 정립',
     '從 此 以 往  至 于 軍 門  則 將 軍 制 之',
     '지금부터 군문에 이르기까지는 장군이 제어한다 — 사마광의 「위 정통」 변호',
     '백제성 탁고 · 제갈량 「출사표」·5차 북벌·오장원 · 사마의 10년 잠복·고평릉 사변 · 司馬昭之心 路人皆知 · 樂不思蜀 · 위·진 찬탈의 시작',
     '사마광 「위(魏) 정통」 선언 — 「**실질이 명분에 앞선다**」. 후에 주희가 「촉 정통」으로 바꿈 — 동양 사학사 가장 큰 논쟁.'),
    ('晉 紀 진기', '진기', '권79~118 · 40권 · AD 266~419 · 155년 — 짧은 통일·긴 분열',
     '何 不 食 肉 糜',
     '어찌 고기죽을 먹지 않느냐 — 진 혜제의 무지의 결정판',
     '진 무제 통일과 사치 · 혜제 「何不食肉糜」 · 팔왕의 난 16년 · 영가의 난·낙양 함락 · 5호16국 · 왕도 「王與馬 共天下」 · 비수대전·사안의 침착 · 유유의 진 찬탈',
     '40권의 비극. 「자기 종친의 내전이 외적보다 무섭다」·「검소함은 한 세대를 못 넘긴다」. 가장 슬픈 부분.'),
    ('宋 紀 송기', '송기', '권119~134 · 16권 · AD 420~478 · 59년 — 남조 첫 왕조',
     '寄 奴',
     '기노 — 유유의 어린 시절 별명 / 가난에서 황제로',
     '유유의 가난·송 건국 · 문제 「원가의 치」 24년 · 친아들 유소의 부친 시해 · 효무제·전폐제 종친 살육 · 소도성의 찬탈',
     '「검소함은 한 세대를 못 넘긴다」·「가족 안의 의심이 곧 왕조의 무덤」. 남조 첫 왕조부터 살육의 패턴.'),
    ('齊 紀 제기', '제기', '권135~144 · 10권 · AD 479~501 · 23년 — 가장 짧은 남조',
     '永 明 體',
     '영명체 — 시문의 황금기 / 그러나 정치는 잔혹',
     '소도성 검소 · 무제 「영명의 치」·심약·사조의 영명체 시문 · 명제 종친 살육 · 동혼후의 「금련 위 미녀」 · 소연의 양 건국',
     '23년에 4단계 패턴 압축 — 「죄가 빨리 자기를 죽인다」. 시문은 자라되 정치는 잔혹.'),
    ('梁 紀 양기', '양기', '권145~166 · 22권 · AD 502~556 · 55년 — 학문 군주의 비극',
     '自 我 得 之  自 我 失 之  亦 復 何 恨',
     '내가 얻었고, 내가 잃었으니, 또 무슨 한이 있겠는가 — 양 무제의 마지막',
     '양 무제 소연 48년 · 소명문선·심약·종영 · 만년 불교 광신 · 4차 사신·황제 보살 · 후경의 난 · 사찰에서 굶어 죽은 황제',
     '「자비롭되 어리석지 말라」·「한 번의 잘못된 자비가 48년의 영광을 지운다」. 「自我得之 自我失之」.'),
    ('陳 紀 진기', '진기', '권167~176 · 10권 · AD 557~588 · 33년 — 남조 마지막',
     '玉 樹 後 庭 花',
     '옥수후정화 — 진 후주가 직접 지은 망국 음악',
     '진패선의 건국 · 문제·선제의 짧은 안정 · 진 후주의 「옥수후정화」·금련경양정 · 수 양견의 통일 (AD 589) · 280년 분열 종언',
     '「사치에 빠진 황제는 자기 발밑의 위기를 보지 못한다」. 두목의 「商女不知亡國恨」의 출전.'),
    ('隋 紀 수기', '수기', '권177~184 · 8권 · AD 589~617 · 29년 — 짧은 통일',
     '何 罪 之 有  乎',
     '내가 무슨 죄가 있느냐 — 강도에서 시해당하기 직전 수 양제',
     '양견 통일·개황의 치·3성6부·과거제 · 양광 시해·즉위 · 대운하·동도 낙양·만리장성 재건 · 고구려 3차 정벌·살수대첩 · 강도에서 친위대에 시해',
     '「큰 공을 한 시대에 다 하면 그 시대가 무너진다」. 진(秦)과 같은 패턴 — 통일 후 짧은 멸망.'),
    ('唐 紀 당기', '당기', '권185~265 · 81권 · AD 618~906 · 290년 — 자치통감 최대',
     '兼 聽 則 明  偏 信 則 暗',
     '두루 들으면 밝고, 치우치게 믿으면 어둡다 — 위징이 당 태종에게',
     '현무문의 변 · 태종 「정관의 치」·위징·방현령·두여회·이정 · 측천무후·무자비 · 현종 개원의 치 · 양귀비·이임보·안록산 · 안사의 난·번진·환관 · 황소의 난·주온 찬탈',
     '81권의 거대한 흐름. 「겸청·인재·여색 절제·변경 군권 통제·환관 금지」 5대 원리. 송이 가장 두려워한 「번진과 환관」의 가장 큰 거울.'),
    ('後 梁 紀 후량기', '후량기', '권266~271 · 6권 · AD 907~922 · 16년',
     '全 忠 → 弑 君',
     '「全忠(완전한 충성)」이라는 이름을 받은 자가 황제를 시해',
     '주온(주전충)의 당 시해 (AD 904) · 후량 건국 (AD 907) · 백마역의 화·사대부 학살 · 친아들에게 시해됨 · 16년 4황제 가족 살육 · 이존욱의 후당에 멸망',
     '「이름과 행동의 정반대」가 가장 큰 비극을 부른다. 한 가문의 자기 파괴.'),
    ('後 唐 紀 후당기', '후당기', '권272~279 · 8권 · AD 923~935 · 13년',
     '燕 雲 十 六 州',
     '연운16주 — 석경당이 거란에 양도한 영토',
     '이존욱 후량 멸·후당 건국 · 「의자(伶人) 정치」·흥교문의 변에 시해 · 명종의 짧은 안정 · 석경당이 거란에 의존·후진 건국 · 연운16주 양도',
     '「창업과 수성의 다름」·「영토를 판 자의 죄」. 사마광이 가장 길게 한탄한 연운16주 양도 — 송 200년 짊어진 슬픔.'),
    ('後 晉 紀 후진기', '후진기', '권280~285 · 6권 · AD 936~946 · 11년',
     '兒 皇 帝',
     '아황제 — 석경당이 거란 야율덕광을 「부황제」로 모시며 자처',
     '석경당 「兒皇帝」의 굴종 · 연운16주 양도 · 매년 비단 30만 필 · 출제의 거란 결렬 · 두중위 배신 · 거란의 변경 함락 · 야율덕광의 짧은 중원 통치',
     '「한 사람의 권력욕이 천 년의 비극」. 가장 굴욕적인 왕조. 송이 200년 회복하지 못한 슬픔의 원천.'),
    ('後 漢 紀 후한기', '후한기', '권286~289 · 4권 · AD 947~950 · 4년 — 가장 짧음',
     '黃 袍 加 身',
     '황포가신 — 군사들이 황색 도포로 곽위를 황제로 만듦',
     '유지원 거란 회군 후 즉위·1년 통치·사망 · 18세 은제 즉위 · 측근에 휘둘려 곽위 가족 학살 · 곽위 분노·진군·은제 살해 · 곽위 「황포가신」으로 후주 건국',
     '4권으로 끝나는 가장 짧은 부분. 「신하의 가족을 죽인 황제는 자기 가족도 잃는다」. 송 태조 「진교병변」의 예행 연습.'),
    ('後 周 紀 후주기', '후주기', '권290~294 · 5권 · AD 951~959 · 9년 — 자치통감 마지막',
     '十 年 開 拓  十 年 養 民  十 年 致 太 平',
     '10년 정복·10년 부국·10년 화평 — 후주 세종의 30년 계획',
     '곽위 검소·시영을 후계자로 · 세종 시영 「30년 계획」 · 고평전투 · 회창 폐불 · 회남 평정 · 연운16주 회복 시도 · 39세 요절 · 자치통감 AD 959에서 끝남',
     '오대 가장 영명한 황제. 「한 사람의 수명이 한 시대를 결정한다」. 「현재 왕조의 역사는 쓰지 않는다」 — 사마광의 마지막 원칙.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅵ. 사관 ==============
SEC6 = 'Ⅵ. 사관 — 「資治」와 「鑑」'

@S(SEC6)
def vi_jian(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '「以史爲鑑」 — 거울론',
              '당 태종의 세 가지 거울에서 책명까지')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(2.2), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.6),
                '以 銅 爲 鏡  可 以 正 衣 冠',
                font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(2.9), Inches(12.0), Inches(0.6),
                '以 古 爲 鏡  可 以 知 興 替',
                font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.6),
                '以 人 爲 鏡  可 以 明 得 失',
                font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(0.4),
                '구리·옛것·사람 — 세 가지 거울 (당 태종이 위징 사후 한 탄식)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.5), [
        ('● 자치통감 「**鑑**」의 사상적 원천',
         {'font_size': 17, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 사마광 — 「**옛일에서 군주에게 유익한 거울을 보여드리려 했을 뿐**」',
         {'font_size': 17, 'space_before': 10}),
        ('● 신종 御製序 — 「**鑑於往事 有資於治道**」가 곧 책명',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
    ])


@S(SEC6)
def vi_zhengtong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '정통론(正統論) — 「실질로 본 정통」',
              '사마광 vs 주희의 영원한 논쟁')
    cols = [
        ('사마광', '실질로 본 정통',
         '위(魏)를 정통\n— 천하의 7할을 차지\n\n「**실질이 명분에 앞선다**」\n\n현실적·합리적\n사학의 안목', INK),
        ('주희', '명분으로 본 정통',
         '촉(蜀)을 정통\n— 한 황실 후예 유비\n\n「**명분이 실질에 앞선다**」\n\n도덕적·도통적\n도학의 안목', ACCENT),
    ]
    for i, (name, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    name, font_size=24, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, Inches(2.9), Inches(5.9), Inches(0.4),
                    label, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.5), Inches(3.6), Inches(4.9), Inches(3.5),
                       [(body, {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER})],
                       line_spacing=1.6)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '주희의 『자치통감강목』 — 사마광의 위 정통을 촉 정통으로 변경 · 동양 사학사 최대 논쟁',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC6)
def vi_principles(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사마광의 7대 정치 원리',
              '218개의 「신광왈」을 종합하면 드러나는 7가지')
    items = [
        ('1', '正 名', '정명',  '이름이 흔들리면 모든 것이 흔들린다 — 권1 첫 사평'),
        ('2', '民 本', '민본',  '백성이 흩어지면 왕조가 망한다 — 民惟邦本'),
        ('3', '人 材', '인재',  '누구를 등용하는가가 시대의 운명 — 蓄人才者 國家大本'),
        ('4', '兼 聽', '겸청',  '두루 듣는 자가 밝다 — 兼聽則明 偏信則暗'),
        ('5', '賞 罰', '상벌',  '공정함이 곧 권위 — 信賞必罰'),
        ('6', '節 制', '절제',  '황제 자신이 절제하지 못하면 모든 것이 무너진다'),
        ('7', '守 成', '수성',  '창업보다 어려운 지킴 — 진짜 시험은 2~3대'),
    ]
    for i, (num, han, kor, desc) in enumerate(items):
        y = Inches(2.3 + i * 0.65)
        add_filled_rect(slide, Inches(0.7), y, Inches(0.7), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(0.7), Inches(0.55),
                    num, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.5), y, Inches(1.6), Inches(0.55), PALE)
        add_textbox(slide, Inches(1.5), y, Inches(1.6), Inches(0.55),
                    han, font_size=16, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(3.3), y + Inches(0.05), Inches(1.5), Inches(0.5),
                    kor, font_size=13, color=SUB, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.9), y + Inches(0.05), Inches(8.0), Inches(0.5),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅶ. 신광왈 ==============
SEC7 = 'Ⅶ. 「臣光曰」'

@S(SEC7)
def vii_first(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '가장 유명한 「臣光曰」 — 권1 첫 사평',
              '삼가분진 사건에 대한 약 1,000자의 긴 사평')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.8), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.6),
                '天 子 之 職 莫 大 於 禮',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(2.9), Inches(12.0), Inches(0.6),
                '禮 莫 大 於 分',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.6),
                '分 莫 大 於 名',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.3), Inches(12.0), Inches(3.0), [
        ('● 천자의 직책은 예(禮)보다 큰 것이 없고, 예는 분(分, 신분 질서)보다 큰 것이 없으며, 분은 명(名, 명분)보다 큰 것이 없다',
         {'font_size': 14, 'space_before': 4, 'color': SUB}),
        ('● 「**예 → 분 → 명**」의 위계가 무너지면 천하가 무너진다',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 공자의 「**정명(正名)**」 사상의 송학적 부활',
         {'font_size': 17, 'space_before': 10}),
        ('● 자치통감 294권 전체의 사관(史觀)을 한 권 첫 페이지에서 명백히 선언',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC7)
def vii_jianting(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '두루 들음 — 「兼聽則明 偏信則暗」',
              '당 태종과 위징의 대화 · 자치통감의 가장 자주 인용된 명제')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.8), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.8),
                '兼 聽 則 明   偏 信 則 暗',
                font_size=44, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(0.4),
                '두루 들으면 밝고, 치우치게 믿으면 어둡다',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.7), Inches(12.0), Inches(2.5), [
        ('● 당 태종: 「人主何以爲明 何以爲暗」(군주가 어떻게 밝고 어떻게 어두워지는가)?',
         {'font_size': 15, 'space_before': 6, 'font_name': 'Batang'}),
        ('● 위징: 「**兼聽則明 偏信則暗**」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 자치통감 전체에서 가장 자주 인용된 정치 원리',
         {'font_size': 16, 'space_before': 12}),
        ('● 사마광 — 「**여러 신하의 의견을 두루 듣는 군주가 밝고, 한 사람의 말만 믿는 군주는 어둡다**」',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC7)
def vii_warning(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '사마광이 송(宋)에 던진 5대 경고',
              '218개 「신광왈」의 가장 깊은 메시지')
    items = [
        ('1', '번진(藩鎭)을 통제하라',     '당 멸망의 첫째 원인 — 안사의 난·번진 할거'),
        ('2', '환관에게 권력을 주지 말라',  '후한·당 멸망의 둘째 원인 — 십상시·감로의 변'),
        ('3', '외척을 견제하라',          '한 멸망의 원인 — 왕망의 30년 가면'),
        ('4', '이민족에 영토를 팔지 말라', '석경당의 연운16주 양도 — 송이 200년 짊어진 슬픔'),
        ('5', '창업의 잔혹을 수성에 이어가지 말라', '진·수의 짧은 멸망 — 인덕 없는 통일'),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(2.3 + i * 0.9)
        add_filled_rect(slide, Inches(0.7), y, Inches(0.8), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(0.8), Inches(0.7),
                    num, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.7), y, Inches(4.7), Inches(0.7), PALE)
        add_textbox(slide, Inches(1.7), y, Inches(4.7), Inches(0.7),
                    title, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.6), y + Inches(0.1), Inches(6.3), Inches(0.6),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅷ. 명장면 ==============
SEC8 = 'Ⅷ. 명장면'

@S(SEC8)
def viii_scenes(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '자치통감 명장면 10선',
              '1,362년의 거대한 흐름에서 가장 극적인 장면들')
    items = [
        ('1', 'BC 403 · 三家分晉',        '권1 첫 사평의 출발 — 정명의 붕괴'),
        ('2', 'BC 206 · 鴻門宴',          '항우와 유방의 운명을 가른 연회'),
        ('3', 'BC 202 · 垓下歌',          '항우의 「力拔山兮氣蓋世」 마지막 노래'),
        ('4', 'AD 9 · 왕망의 30년 가면',  '한실 찬탈·신(新) 건국·15년 멸망'),
        ('5', 'AD 23 · 곤양 전투',         '광무제 1만이 왕망 40만 격파'),
        ('6', 'AD 223 · 백제성 탁고',      '유비가 제갈량에게 천하를 맡김'),
        ('7', 'AD 383 · 비수대전',         '동진 8만이 전진 80만 격파 — 「草木皆兵」'),
        ('8', 'AD 626 · 玄武門의 변',      '이세민이 형제 사살·당 태종 즉위'),
        ('9', 'AD 755 · 安史의 난',        '당의 결정적 분기점 — 양귀비 죽음'),
        ('10', 'AD 936 · 연운16주 양도',   '석경당이 거란에 영토 양도 — 사마광 최대 한탄'),
    ]
    for i, (num, title, desc) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 0.95)
        add_filled_rect(slide, x, y, Inches(0.7), Inches(0.8), ACCENT)
        add_textbox(slide, x, y, Inches(0.7), Inches(0.8),
                    num, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, x + Inches(0.8), y, Inches(2.5), Inches(0.8), PALE)
        add_textbox(slide, x + Inches(0.8), y, Inches(2.5), Inches(0.8),
                    title, font_size=12, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x + Inches(3.4), y + Inches(0.1), Inches(2.7), Inches(0.6),
                    desc, font_size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅸ. 명구 12선 ==============
SEC9 = 'Ⅸ. 명구 12선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC9)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC9} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=16, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 16, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('권1 첫 사평 · 삼가분진', '天 子 之 職 莫 大 於 禮  禮 莫 大 於 分  分 莫 大 於 名',
     '천자의 직책은 예, 예는 분, 분은 명보다 크지 않다',
     '자치통감 전체의 첫 정치 원리 — 「정명(正名)」의 송학적 부활. 권1 첫 「臣光曰」.'),
    ('한고조 자평 · 한기', '夫 運 籌 帷 幄 之 中  決 勝 千 里 之 外',
     '장막 안에서 책략을 짜 천 리 밖에서 승부를 결정한다',
     '장량(張良)에 대한 한고조 평. 「**내가 능력은 그들만 못하나, 그들을 쓸 수 있어 천하를 얻었다**」 — 인재 등용술의 정점.'),
    ('한신의 한탄 · 한기', '狡 兎 死  良 狗 烹  高 鳥 盡  良 弓 藏',
     '교활한 토끼가 죽으면 좋은 사냥개를 삶고, 높은 새가 사라지면 좋은 활을 감춘다',
     '「**토사구팽**」의 출전. 공신의 영원한 운명 — 송 태조의 「배주석병권」과 정반대의 길.'),
    ('백제성 탁고 · 위기', '君 才 十 倍 曹 丕  必 能 安 國  終 定 大 事',
     '그대의 재주는 조비의 10배이니 반드시 나라를 안정시키고 대업을 이룰 것이오',
     '유비의 제갈량 탁고. 「**천하를 맡길 만큼 신뢰한 군주와 죽음까지 충성을 지킨 신하**」 — 동양 군신 관계의 정점.'),
    ('제갈량 출사표 · 위기', '鞠 躬 盡 瘁  死 而 後 已',
     '몸을 굽혀 마음을 다하고, 죽은 뒤에야 그칩니다',
     '제갈량의 「출사표」. 동양 충신의 정수 — 후한~송대까지 가장 자주 인용된 글.'),
    ('진 혜제 · 진기晉', '何 不 食 肉 糜',
     '어찌 고기죽을 먹지 않느냐',
     '백성이 굶주린다는 보고에 진 혜제의 답. 「**부유한 자의 무지**」의 영원한 표상.'),
    ('당 태종 · 당기', '以 銅 爲 鏡 可 以 正 衣 冠  以 古 爲 鏡 可 以 知 興 替  以 人 爲 鏡 可 以 明 得 失',
     '구리·옛것·사람 — 세 가지 거울로 의관·흥망·득실을 안다',
     '위징 사후 당 태종의 탄식. 「**내가 거울 하나를 잃었다**」 — 자치통감 책명 「鑑」의 원천.'),
    ('위징 · 당기', '兼 聽 則 明  偏 信 則 暗',
     '두루 들으면 밝고, 치우치게 믿으면 어둡다',
     '당 태종에게 위징의 답. 자치통감 전체에서 가장 자주 인용된 정치 원리.'),
    ('당 태종 · 당기', '凡 創 業 之 難 已 往  守 成 之 難 方 當 與 諸 公 愼 之',
     '창업의 어려움은 지났으나, 수성의 어려움은 지금부터 신중해야 한다',
     '정관(貞觀)의 치를 만든 자세. **창업과 수성 중 어느 것이 어려운가**의 토론의 답.'),
    ('한신 인용·통감 곳곳', '當 斷 不 斷  反 受 其 亂',
     '결단해야 할 때 결단하지 않으면 오히려 화를 입는다',
     '한기 누경의 말. **결단의 타이밍** — 망설임이 오히려 더 큰 혼란을 부른다.'),
    ('서경 인용·통감 곳곳', '民 惟 邦 本  本 固 邦 寧',
     '백성이 곧 나라의 근본이니, 근본이 굳어야 나라가 평안하다',
     '신광왈에서 반복 환기되는 동양 민본주의의 절대 명제. 사마광이 가장 자주 강조한 원리.'),
    ('신종 御製序', '鑑 於 往 事  有 資 於 治 道',
     '옛일을 거울 삼아 다스림의 도에 도움이 있게 한다',
     '신종의 「자치통감」 명명 사유. 「**자치통감**」 네 글자의 풀이. 책 전체의 정신을 한 줄로 압축.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅹ. 후대 영향 ==============
SEC10 = 'Ⅹ. 후대 영향'

@S(SEC10)
def x_china(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '중국 — 호삼성·주희·왕부지·모택동',
              '900년의 주석과 비평')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 남송 호삼성(胡三省, 1230~1302) — 30년 작업의 표준 주석 『資治通鑑音注』',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('     · 송이 망한 뒤 원 시대에 작업 — 망국의 슬픔을 학문에 쏟음',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 남송 주희(朱熹) — 『資治通鑑綱目』으로 압축·재편',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 사마광의 「위 정통」을 「촉 정통」으로 변경 — 도통론 적용',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 명말청초 왕부지(王夫之) — 『讀通鑑論』 — 청 초의 깊은 비평',
         {'font_size': 17, 'space_before': 12}),
        ('● 모택동(毛澤東) — 「**17번 읽었다**」 · 베갯머리 책',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 「**내가 정치를 배운 것은 자치통감에서다**」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 「**통감 안에는 정치의 모든 것이 있다. 신광왈을 특히 깊이 읽어야**」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
    ])


@S(SEC10)
def x_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '한국 — 조선 경연의 핵심 텍스트',
              '권근·정도전·세종·영조·정조의 책')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 고려 중기 도입 → 조선 초 사서삼경과 함께 사대부 필독서',
         {'font_size': 17, 'space_before': 4}),
        ('● 태종·세종 — 경연에서 자주 진강',
         {'font_size': 17, 'space_before': 10}),
        ('● 권근·정도전 — 통감의 정명·민본 사상으로 조선 건국 이념 형성',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 영조 — 경연에서 통감을 가장 자주 진강',
         {'font_size': 17, 'space_before': 10}),
        ('● 정조 — 통감을 깊이 연구 · 『통감강목속편』 편찬 관심',
         {'font_size': 17, 'space_before': 10}),
        ('● 조선 시대는 통감 원문 294권보다 주희의 『자치통감강목』이 더 자주 읽힘',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 권중달 역 『자치통감』 전31권 — 한국에서 가장 권위 있는 완역',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC10)
def x_japan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '일본 — 막부의 정치 교과서',
              '도쿠가와·시부사와 에이이치·도쿠토미 소호')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 헤이안 시대부터 한학자들이 통감 학습',
         {'font_size': 17, 'space_before': 4}),
        ('● 에도 막부 — 도쿠가와가 통감강목을 정치 교과서로',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 시부사와 에이이치(澁澤榮一) — 메이지 자본주의의 아버지 · 통감 애독',
         {'font_size': 17, 'space_before': 12}),
        ('● 도쿠토미 소호(德富蘇峰) — 통감을 평생의 책으로',
         {'font_size': 17, 'space_before': 10}),
        ('● 일본 한학자가 작성한 『통감대전』·『통감의해』 등 주석서 다수',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
        ('● 통감의 「**거울론**」이 메이지 유신 사상의 한 토대',
         {'font_size': 16, 'space_before': 10, 'color': ACCENT, 'bold': True}),
    ])


# ============== Ⅺ. 오늘 우리에게 ==============
SEC11 = 'Ⅺ. 오늘 우리에게'

@S(SEC11)
def xi_modern(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '자치통감이 오늘 우리에게 일러주는 10가지')
    items = [
        '1. 이름이 흔들리면 질서가 무너진다 — 정명(正名)의 첫 원리',
        '2. 두루 들어야 밝다 — 兼聽則明 偏信則暗',
        '3. 창업보다 수성이 어렵다 — 「유종의 미」의 어려움',
        '4. 인재가 곧 시대다 — 누구를 곁에 두는가',
        '5. 민이 흩어지면 모든 것이 끝난다 — 民惟邦本',
        '6. 절정에서 내려올 줄 알아야 한다 — 안사의 난·당 현종',
        '7. 위선의 정치는 반드시 무너진다 — 왕망의 30년 가면',
        '8. 상벌이 흔들리면 신뢰가 사라진다 — 신상필벌',
        '9. 결단의 순간을 놓치면 화를 입는다 — 當斷不斷',
        '10. 사례는 옷을 바꿔 다시 등장한다 — 鑑於往事',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.3 + i * 0.45)
        add_textbox(slide, Inches(0.7), y, Inches(12.0), Inches(0.4),
                    txt, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅻ. 마무리 ==============
SEC12 = 'Ⅻ. 마무리'

@S(SEC12)
def xii_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '자치통감, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 사마광이 19년에 걸쳐 1,362년의 중국사를 294권 약 300만 자로 정리.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 16편 — 주·진·한·위·진·송·제·양·진·수·당·후량·후당·후진·후한·후주.',
         {'font_size': 18, 'space_before': 8}),
        ('● 방법론 — 叢目·長編·定本의 3단계 작업법 + 「考異」의 사료 비판.',
         {'font_size': 18, 'space_before': 8}),
        ('● 「臣光曰」 218회 — 사마광의 직접 정치 평론 · 7대 원리 (정명·민본·인재·겸청·상벌·절제·수성).',
         {'font_size': 18, 'space_before': 8}),
        ('● 후대 — 호삼성·주희·왕부지·모택동의 베갯머리 책 · 조선 경연·일본 막부의 교과서.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「鑑於往事 有資於治道」 — 동양 1,000년의 「판단의 학교」.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC12)
def xii_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
                '臣 光 曰',
                font_size=84, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(3.5), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.6),
                '신 광 이  말 씀 드 립 니 다',
                font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.5),
                '鑑 於 往 事  有 資 於 治 道',
                font_size=28, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
                '옛일을 거울 삼아 다스림의 도에 도움이 있게 한다',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '— 신종(神宗) 어제서(御製序) · 1084년 12월',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '資  治  通  鑑',
                font_size=22, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\자치통감.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
