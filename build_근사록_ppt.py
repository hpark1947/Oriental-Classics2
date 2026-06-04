# -*- coding: utf-8 -*-
"""
근사록(近思錄) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '신유학(性理學)의 종합 입문서 · 동아시아 800년 사대부 정신 교육의 정수',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '近 思 錄',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '근 사 록',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '切 問 而 近 思  仁 在 其 中',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '절실하게 묻고 가까이 생각하면, 인(仁)은 그 가운데 있다  — 『논어』 자장편',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '주희(朱熹) · 여조겸(呂祖謙) 공편 · 1175 한천정사 · 14편 622조 · 북송 4자(四子)',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 근사록이란'),
        ('Ⅱ.', '편찬 — 1175년 한천정사의 10일'),
        ('Ⅲ.', '인용된 북송 4자(四子)'),
        ('Ⅳ.', '14편의 동심원 구조'),
        ('Ⅴ.', '14편 각 편 깊이 읽기'),
    ]
    items_right = [
        ('Ⅵ.', '근사(近思)의 5가지 학습 원리'),
        ('Ⅶ.', '명구 12선'),
        ('Ⅷ.', '한국에서의 수용 — 800년의 책'),
        ('Ⅸ.', '오늘 다시 펼치는 이유'),
        ('Ⅹ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '서명의 뜻 — 「가까이서 생각함」',
              '『논어』 자장편에서 따온 학문의 정신')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5),
                '博 學 而 篤 志   切 問 而 近 思   仁 在 其 中 矣',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(3.5), [
        ('● 자하(子夏)의 말 — 『논어』 「자장(子張)」편',
         {'font_size': 17, 'space_before': 6}),
        ('● 「널리 배우고(博學) · 뜻을 돈독히 하며(篤志) · 절실하게 묻고(切問) · 가까이 생각하라(近思)」',
         {'font_size': 17, 'space_before': 10}),
        ('● 「인(仁)이 그 가운데 있다」 — 인격의 완성은 바로 이 네 가지에 있다',
         {'font_size': 17, 'space_before': 10}),
        ('● 「近思」 — 「내 몸 가까이 있는 일부터 생각하라」',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 고원하고 현묘한 이론에 매달리지 말고, 자기 몸·가정·일상에 절실한 문제부터 성찰',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 책 이름 자체가 「실천 가능하고 절실한 학문」의 선언',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 근사록')
    rows = [
        ('편자',     '주희(朱熹, 1130~1200) · 여조겸(呂祖謙, 1137~1181, 호 동래東萊)'),
        ('편찬 시기',  '1175년 한천정사(寒泉精舍) 회동 · 1178년 최종 간행'),
        ('분량',     '14권(편) · 622조의 어록·문장 발췌'),
        ('인용 대상',  '북송 4자 — 주돈이(염계)·정호(명도)·정이(이천)·장재(횡거)'),
        ('성격',     '정주학(程朱學) = 성리학(性理學)의 종합 입문서'),
        ('교육 위상',  '사서삼경 학습 전 단계의 「입문 정본」'),
        ('한국 수용',  '고려 말 도입(1370 진주 초간) · 조선 사림의 필독서'),
        ('정조 어제서', '1794년 정조가 직접 어제서(御製序)를 지어 그 중요성을 강조'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_position(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '교육 체계 속 근사록의 자리',
              '소학 다음, 사서삼경 앞 — 성리학으로 들어가는 문')
    rows = [
        ('8~14세',  '소학(小學)',        '일상 예절과 인격의 체계'),
        ('이후',    '근사록(近思錄)',     '정주학(성리학)의 종합 입문', True),
        ('이후',    '사서(四書)',        '대학 → 논어 → 맹자 → 중용'),
        ('이후',    '삼경(三經)',        '시경·서경·역경 — 본격 유학'),
        ('이후',    '북송 4자 본 저술',   '태극도설·정몽·이정전서 등 원전'),
    ]
    for i, row in enumerate(rows):
        age, han, desc = row[0], row[1], row[2]
        highlighted = len(row) > 3 and row[3]
        y = Inches(2.5 + i * 0.9)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.8), Inches(0.75), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.8), Inches(0.75),
                    age, font_size=14, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        color = ACCENT if highlighted else INK
        add_filled_rect(slide, Inches(2.7), y, Inches(2.8), Inches(0.75), color)
        add_textbox(slide, Inches(2.7), y, Inches(2.8), Inches(0.75),
                    han, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.8), y + Inches(0.05), Inches(7.2), Inches(0.7),
                    desc, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '여조겸의 발문 — 「이 책은 초학자(初學者)를 위한 입문서이지 도학의 완성이 아니다」',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 편찬 ==============
SEC2 = 'Ⅱ. 편찬'

@S(SEC2)
def ii_meet(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '1175년 한천정사(寒泉精舍)의 10일',
              '두 거장이 머리를 맞대고 신유학의 입문서를 짓다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1175년 4월(남송 효종 순희 2년) — 여조겸이 주희가 있던 한천정사(복건성 건양) 방문',
         {'font_size': 17, 'space_before': 4}),
        ('● 약 10일간 함께 머무르며 북송 4자의 방대한 저술을 함께 읽음',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 편찬 동기 — 정주학을 공부하려는 학자들이 마주한 세 가지 어려움',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 북송 4자의 저술이 너무 방대 — 어디서부터 읽어야 할지 막막',
         {'font_size': 15, 'color': SUB, 'space_before': 6}),
        ('     · 어록과 문장이 흩어져 있어 체계적 학습이 어려움',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('     · 입문자가 본질적이고 핵심적인 가르침에 빨리 도달할 길이 필요',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 두 사람의 분담', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 주희 — 전체 기획, 사상적 골격, 분류 체계',
         {'font_size': 15, 'space_before': 6}),
        ('     · 여조겸 — 자료 정리, 편집 협조, 발문 작성',
         {'font_size': 15, 'space_before': 4}),
    ])


@S(SEC2)
def ii_zhu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '두 편자')
    cols = [
        ('朱 熹 주희', '1130~1200',
         '남송 대유학자\n성리학의 집대성자\n\n사서집주(四書集注)·\n중용장구·시집전·서집전\n\n동아시아 성리학의 표준\n— 「주자(朱子)」'),
        ('呂 祖 謙 여조겸', '1137~1181',
         '호 동래(東萊)\n금화학파(金華學派)\n\n역사학·문헌학에 정통\n『동래박의』·『대사기』\n\n주희·장식과 「동남삼현(東南三賢)」\n— 향년 45세 요절'),
    ]
    for i, (han, lifespan, body) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), ACCENT if i == 0 else INK)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=26, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.85), Inches(5.9), Inches(0.4),
                    lifespan, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.4), Inches(3.6), Inches(5.1), Inches(3.5),
                       [(body, {'font_size': 14, 'color': INK})], line_spacing=1.5)


# ============== Ⅲ. 북송 4자 ==============
SEC3 = 'Ⅲ. 북송 4자(四子)'

@S(SEC3)
def iii_four(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '북송 4자 — 정주학의 네 뿌리',
              '주희가 종합·체계화한 성리학의 원천')
    rows = [
        ('周敦頤 주돈이', '濂溪 염계', '1017~1073', '태극(太極)·무극(無極) · 우주생성론의 창시',
         '『태극도설』·『통서』'),
        ('程顥 정호',    '明道 명도', '1032~1085', '인(仁)·천리(天理) · 만물 일체의 인',
         '어록(이정전서)'),
        ('程頤 정이',    '伊川 이천', '1033~1107', '이기(理氣)·격물치지·경(敬) — 정주학의 뼈대',
         '『역전(易傳)』·어록'),
        ('張載 장재',    '橫渠 횡거', '1020~1077', '기(氣) 일원론 · 민포물여(民胞物與)',
         '『정몽』·『서명(西銘)』'),
    ]
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.5), INK)
    headers = [('학자', 2.3), ('호', 1.5), ('생몰', 1.5), ('핵심 사상', 4.5), ('대표 저술', 2.2)]
    x = Inches(0.7)
    for label, w in headers:
        add_textbox(slide, x, Inches(2.3), Inches(w), Inches(0.5),
                    label, font_size=13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(w)
    for i, (name, alias, life, thought, book) in enumerate(rows):
        y = Inches(2.8 + i * 0.95)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.9), bg)
        add_textbox(slide, Inches(0.7), y, Inches(2.3), Inches(0.9),
                    name, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(3.0), y, Inches(1.5), Inches(0.9),
                    alias, font_size=13, color=SUB, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.5), y, Inches(1.5), Inches(0.9),
                    life, font_size=11, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.0), y + Inches(0.1), Inches(4.4), Inches(0.7),
                    thought, font_size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(10.5), y, Inches(2.2), Inches(0.9),
                    book, font_size=11, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC3)
def iii_four_concepts(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '4자가 남긴 4가지 핵심 개념')
    boxes = [
        ('太 極', '주돈이', '우주 만물이 펼쳐지는 근원\n「無極而太極」 — 무극이 곧 태극'),
        ('仁 / 萬 物 一 體', '정호', '만물이 한 몸 — 「一草一木 皆有春意」\n인(仁)이 곧 천리(天理)'),
        ('格 物 致 知 / 敬', '정이', '사물에 나아가 이치를 궁구함\n경(敬) = 「主一無適」'),
        ('民 胞 物 與', '장재', '백성은 내 동포, 만물은 내 짝\n「爲生民立命 為萬世開太平」'),
    ]
    for i, (han, who, desc) in enumerate(boxes):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(6.0), Inches(2.0), PALE)
        add_textbox(slide, x, y + Inches(0.2), Inches(6.0), Inches(0.7),
                    han, font_size=22, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(0.9), Inches(6.0), Inches(0.4),
                    who, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.3), y + Inches(1.3), Inches(5.4), Inches(0.65),
                    desc, font_size=12, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅳ. 14편의 동심원 구조 ==============
SEC4 = 'Ⅳ. 14편의 동심원 구조'

@S(SEC4)
def iv_circle(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '안에서 밖으로 · 작은 데서 큰 데로',
              '대학(大學)의 수신·제가·치국·평천하와 정확히 대응')
    rows = [
        ('단계 1',  '권1',     '도체(道體)',         '형이상학적 토대 — 「세상에는 일관된 원리가 있다」'),
        ('단계 2',  '권2~5',  '위학·치지·존양·극기', '자기 수양 — 학문 동기·인식·마음·자기 통제'),
        ('단계 3',  '권6',     '가도(家道)',         '가정 윤리 — 「집 안의 모습이 진짜 인격」'),
        ('단계 4',  '권7~10', '출처·치체·치법·정사', '사회·정치 — 진퇴·정치 본체·제도·실무'),
        ('단계 5',  '권11~12', '교학·경계',         '교육과 자기 점검 — 적용과 반성'),
        ('단계 6',  '권13~14', '변이단·관성현',     '분별과 모범 — 무엇을 거를까 / 누구를 본받을까'),
    ]
    for i, (step, vol, name, desc) in enumerate(rows):
        y = Inches(2.3 + i * 0.78)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.3), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.3), Inches(0.65),
                    step, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(2.1), y, Inches(1.2), Inches(0.65), PALE)
        add_textbox(slide, Inches(2.1), y, Inches(1.2), Inches(0.65),
                    vol, font_size=14, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.4), y, Inches(3.2), Inches(0.65),
                    name, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(6.7), y, Inches(6.2), Inches(0.65),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅴ. 14편 각 편 깊이 읽기 ==============
SEC5 = 'Ⅴ. 14편 각 편 깊이 읽기'

def make_chapter_slide(num, total, vol_no, name_han, name_kor, items_count,
                        topic, principle, today):
    @S(SEC5)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC5} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{vol_no}  {name_han}  ({name_kor})',
                    font_size=26, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    f'{items_count}조 · {topic}', font_size=14, color=SUB)
        add_textbox(slide, Inches(0.7), Inches(2.1), Inches(12.0), Inches(0.4),
                    '◆ 핵심 명제', font_size=14, bold=True, color=ACCENT)
        add_filled_rect(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(1.7), PALE)
        add_paragraphs(slide, Inches(0.9), Inches(2.6), Inches(11.6), Inches(1.6),
                       [(principle, {'font_size': 15, 'color': INK})], line_spacing=1.45)
        add_textbox(slide, Inches(0.7), Inches(4.4), Inches(12.0), Inches(0.4),
                    '◆ 오늘 우리가 배울 것', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(today, {'font_size': 16, 'color': INK})], line_spacing=1.45)


CHAPTERS = [
    ('권1', '道體 도체', '도체', '51조', '형이상학 — 세상과 나는 어떻게 생겼는가',
     '· 우주 만물은 「태극(太極)」에서 음양·오행·만물로 펼쳐진다\n· 천(天)의 도는 음양, 지(地)의 도는 강유, 인(人)의 도는 인의(仁義)\n· 이(理)와 기(氣) — 모든 존재는 이(원리)와 기(질료)의 결합',
     '「내가 사는 세상에는 일관된 원리가 있다」는 사고 방식.\n작은 일을 처리할 때도 「이 일의 근본 원리는 무엇인가」를 묻는 습관.'),
    ('권2', '爲學 위학', '위학', '111조', '학문 동기 — 14편 중 가장 분량이 많은 실질 중심',
     '· 학문의 목적은 「성인(聖人)이 되는 것」 — 단순한 지식 축적이 아니다\n· 「입지(立志)」 — 뜻을 세움이 모든 학문의 시작\n· 「위기지학(爲己之學) — 위인지학(爲人之學)이 아니다」 — 자기를 위한 공부',
     '「공부의 목적은 자기 변화이지 스펙·인정이 아니다」.\n시험·취업·인정을 위한 공부는 자기를 소진시키고, 자기 변화를 위한 공부는 평생의 동력이 된다.'),
    ('권3', '致知 치지', '치지', '78조', '인식론 — 어떻게 진짜 앎에 도달하는가',
     '· 격물(格物) — 사물 하나하나의 이치를 궁구함\n· 궁리(窮理) — 그 이치의 끝까지 탐구함\n· 활연관통(豁然貫通) — 누적된 격물이 어느 순간 전체 깨달음으로 통합됨',
     '「앎은 누적과 도약의 결합」.\n자료 검색·요약에 만족하지 말고 「왜 그런가」를 끝까지 묻는 자세.\n한 분야를 깊이 파면 다른 분야가 보이는 통섭(通涉)의 경험.'),
    ('권4', '存養 존양', '존양', '70조', '마음 수양 — 마음을 어떻게 지키는가',
     '· 주일무적(主一無適) — 마음이 한 곳에 집중, 흐트러지지 않음 = 경(敬)\n· 정중유동 동중유정(靜中有動 動中有靜) — 고요함 속 깨어 있음, 분주함 속 고요함\n· 정좌(靜坐) — 마음을 가라앉혀 본성을 자각',
     '「바쁠수록 마음의 중심을 잃지 말라」.\n스마트폰·SNS의 시대 — 「경(敬, 하나에 깨어 있음)」이 가장 절박한 능력.\n명상·집중력 훈련·디지털 디톡스의 동양적 원형.'),
    ('권5', '克己 극기', '극기', '41조', '자기 통제 — 사욕을 어떻게 다스리는가',
     '· 천리(天理)와 인욕(人慾)의 대립 — 도덕의 결정적 갈림길\n· 분노·욕심·태만 등 인심의 병통(病痛) 분석\n· 「한 번 끊으면 평생 자유로워지는 결단」의 힘',
     '「자기 약점에 대한 정직한 인식과 결단」.\n게으름·욕심·분노·자존심 — 외면하지 말고 정면으로 보고 한 번에 끊어내는 결단.\n중독·게으름·분노 조절 — 모든 자기 통제의 원리.'),
    ('권6', '家道 가도', '가도', '22조', '가정 윤리 — 가정은 어떻게 다스리는가',
     '· 가정은 사회 윤리의 출발점이자 모범\n· 부모의 솔선수범이 가장 강력한 가정 교육\n· 부부 간 「상경여빈(相敬如賓)」 — 서로 손님처럼 공경\n· 형제의 우애가 효(孝)의 자연스러운 확장',
     '「바깥의 인격은 집 안 모습이 진짜다」.\n사회에서 아무리 훌륭해도 가정에서 무너지면 그 인격은 가짜.\n「가족이라서 막대해도 된다」는 사고를 버리는 것.'),
    ('권7', '出處 출처', '출처', '39조', '진퇴 의리 — 나아갈 때와 물러설 때',
     '· 「의(義)에 맞으면 나아가고, 의에 맞지 않으면 물러난다」\n· 시류에 영합하기 위해 나아가지 않는다\n· 시중(時中) — 때에 맞는 처신\n· 받을 것과 안 받을 것의 분명한 구분',
     '「기회 앞에서 본질을 묻는 안목」.\n모든 기회를 다 잡으려는 사람은 결국 큰 기회를 놓친다.\n이직·창업·승진 — 「내 가치와 맞는가」를 묻는 자세. 거절할 줄 아는 용기.'),
    ('권8', '治體 치체', '치체', '25조', '정치 본체 — 정치의 근본은 무엇인가',
     '· 다스리는 자는 먼저 자기 자신을 다스려야 한다 (수신 = 치국의 근본)\n· 정치의 근본은 민심이다\n· 「올바름(正)으로 다스리는 것」이 모든 정치 기술에 우선\n· 인재를 알아보고 등용하는 안목이 군주의 첫째 자질',
     '「리더십의 본질은 기교가 아니라 자기 정립(正)」.\n회사 CEO든 작은 팀장이든 — 자기 자신부터 바르게 서지 않으면 사람을 움직일 수 없다.\n부하에게 요구하기 전에 자기에게 먼저 묻는 습관.'),
    ('권9', '治法 치법', '치법', '27조', '제도·법 — 어떤 제도로 다스리는가',
     '· 정치는 「사람(治人)」이 먼저고, 「법(治法)」이 그 다음\n· 그러나 좋은 법·제도가 없으면 좋은 사람도 효과를 내지 못한다\n· 「예(禮)는 부드러운 다스림, 형(刑)은 마지못한 보충」\n· 학교 제도·인재 선발·토지 제도의 이상',
     '「사람을 바꾸는 것은 결국 시스템이고, 시스템을 만드는 것은 결국 사람이다」.\n조직 운영에서 「사람을 갈아 넣지 말고 시스템을 바꿔라」 + 「시스템도 결국 사람이 만든다」 — 양면을 함께 보는 안목.'),
    ('권10', '政事 정사', '정사', '64조', '실무 — 실제로 어떻게 일을 처리하는가',
     '· 사세(事勢)를 정확히 보라 — 명분만으로 일이 되지 않는다\n· 작은 일을 소홀히 하지 말라\n· 백성의 처지를 먼저 헤아려라\n· 결정과 책임을 회피하지 말라',
     '「이상은 높게, 실무는 정확하게」.\n도덕적 이상만 외치는 사람은 일을 망치고, 실무 기교만 따지는 사람은 방향을 잃는다.\n큰 그림을 잃지 않으면서 작은 디테일까지 챙기는 균형.'),
    ('권11', '敎學 교학', '교학', '21조', '교육 — 어떻게 가르치고 어떻게 배우게 하는가',
     '· 솔선수범이 모든 가르침의 근본\n· 학생의 수준에 맞춰 가르쳐야 한다 (수기隨機 교육)\n· 교학상장(敎學相長) — 가르치는 일은 자기 공부의 연장\n· 잘못을 꾸짖되 기를 꺾지 말라',
     '「가장 잘 가르치는 사람은 자기가 가장 잘 배우고 있는 사람」.\n부모로서 자식을, 선배로서 후배를, 상사로서 부하를 이끌 때 — 「말로 가르치지 말고 모습으로 보여라」.'),
    ('권12', '警戒 경계', '경계', '33조', '자기 점검 — 어떤 함정을 피해야 하는가',
     '· 자만(自滿) — 조금 알았다고 다 안 줄 알고 멈춘다\n· 태만(怠慢) — 작심삼일, 지속의 부재\n· 편견(偏見) — 자기 생각에 갇힘\n· 명리(名利)에 휘둘림 — 본래 목적을 잊음\n· 외적 평판에의 집착 — 위인지학으로의 타락',
     '「공부와 수양이 망가지는 것은 외부 적이 아니라 자기 안의 함정」.\n위학편이 길을 제시했다면, 경계편은 그 길에서 빠지기 쉬운 함정들을 짚어 준다 — 자기 점검의 체크리스트.'),
    ('권13', '辨異端 변이단', '변이단', '14조', '비판적 사고 — 그릇된 학설과 어떻게 구별하는가',
     '· 이단의 비판은 그 사상가 개인의 비방이 아니라 학설의 오류 분석\n· 외형이 비슷해도 근본 전제와 실천 방향이 다르면 다른 학문\n· 인의(仁義)·인륜(人倫)을 부정하면 결국 사회를 해친다\n· 「가짜 같음」이 가장 위험하다 — 진리에 가까울수록 분별이 어렵다',
     '단순한 종교 비판이 아닌 「비판적 사고력」.\n표면만 보지 말고 근본 전제를 묻고, 결과적 방향을 보는 안목.\n정보·이념·자기계발서 홍수 속 — 「이 주장의 전제는, 받아들이면 어디로 가는가」.'),
    ('권14', '觀聖賢 관성현', '관성현', '26조', '모범 — 어떤 모습을 본받아야 하는가',
     '· 성인은 추상적 이상이 아니라 구체적 사람이었다\n· 그들의 언행·처세·교육·정치가 모두 모범\n· 공자의 온화함, 맹자의 호연지기, 안연의 고요함 — 성인도 다양한 인격\n· 성현의 모습을 마음에 새기면 자기 행동의 척도가 된다',
     '「모방할 수 있는 구체적 모범 없이는 변화도 없다」.\n추상적 이상은 동기 부여가 약하지만, 구체적 인격의 살아 있는 모습은 사람을 바꾼다.\n누구를 마음의 스승으로 삼느냐가 그 사람의 미래를 결정.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅵ. 5가지 학습 원리 ==============
SEC6 = 'Ⅵ. 5가지 학습 원리'

@S(SEC6)
def vi_five(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '근사(近思)의 정신 — 다섯 글자에 압축된 학습관',
              '논어 자장편 + 정주학의 14편을 관통하는 원리')
    boxes = [
        ('博 學', '박학', '널리 배우라', '14편이 우주론에서 정치·교육·이단 비판까지 망라하듯, 학문은 좁아서는 안 된다'),
        ('篤 志', '독지', '뜻을 돈독히 하라', '위학편이 강조하듯, 모든 학문의 시작은 「왜 공부하는가」에 대한 분명한 답'),
        ('切 問', '절문', '절실하게 물으라', '치지편의 격물치지가 가르치듯, 표면에 만족하지 말고 끝까지 묻는 태도'),
        ('近 思', '근사', '가까이 생각하라', '책 이름 자체 — 거창한 이론이 아닌, 자기 몸·가정·일상의 절실한 문제부터'),
        ('力 行', '역행', '힘써 실천하라', '극기·존양·정사편이 강조하듯, 앎은 행함으로 완성된다'),
    ]
    for i, (han, kor, label, desc) in enumerate(boxes):
        y = Inches(2.3 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.8),
                    han, font_size=24, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, Inches(2.9), y, Inches(2.5), Inches(0.8), PALE)
        add_paragraphs(slide, Inches(2.9), y + Inches(0.1), Inches(2.5), Inches(0.7), [
            (label, {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
            (kor,   {'font_size': 12, 'color': SUB, 'space_before': 2, 'align': PP_ALIGN.CENTER}),
        ], line_spacing=1.2)
        add_textbox(slide, Inches(5.6), y + Inches(0.1), Inches(7.3), Inches(0.7),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅶ. 명구 12선 ==============
SEC7 = 'Ⅶ. 명구 12선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC7)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC7} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=17, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 17, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('논어 자장편 / 근사록 서명의 출전', '博 學 而 篤 志  切 問 而 近 思',
     '널리 배우고 뜻을 돈독히 하며, 절실하게 묻고 가까이 생각하라',
     '자하(子夏)의 말. 「근사(近思)」가 곧 책 이름이 되었다. 「인(仁)이 그 가운데 있다」 — 인격 완성의 사대 명제.'),
    ('도체 · 주돈이 태극도설', '無 極 而 太 極',
     '무극이 곧 태극이다',
     '주돈이가 우주생성론의 첫 명제로 던진 한 줄. 태극의 펼쳐짐이 곧 음양·오행·만물. 「세상에 일관된 원리가 있다」는 사고의 출발.'),
    ('도체 · 정호의 명제', '一 草 一 木  皆 有 春 意',
     '한 포기 풀, 한 그루 나무에도 모두 봄의 뜻이 있다',
     '정호(명도)의 인(仁) 사상 — 만물 일체. 우주 만물이 같은 생명력으로 통한다는 생명 감수성의 명제.'),
    ('도체 · 장재 서명', '民 吾 同 胞  物 吾 與 也',
     '백성은 내 동포요, 만물은 내 짝이다',
     '장재(횡거)의 『서명(西銘)』 — 「민포물여(民胞物與)」의 원전. 동아시아 인본주의·자연관의 최고 표현.'),
    ('위학 · 정주학의 학문관', '爲 己 之 學  非 爲 人 之 學',
     '자기를 위한 공부지 남에게 보이기 위한 공부가 아니다',
     '학문의 본질은 인정·평가·스펙이 아니라 자기 변화. 「위기지학(爲己之學)」 — 정주학의 학습 윤리 핵심.'),
    ('치지 · 정이의 인식론', '今 日 格 一 物  明 日 格 一 物  豁 然 貫 通',
     '오늘 한 사물을 격(格)하고 내일 또 한 사물을 격하면, 어느 순간 활연히 관통한다',
     '격물치지(格物致知)의 누적과 도약. 한 분야를 깊이 파면 어느 순간 전체가 보이는 통섭의 원리.'),
    ('존양 · 경(敬)의 정의', '主 一 無 適',
     '하나에 집중하여 다른 데로 가지 않음',
     '경(敬)의 가장 짧은 정의. 마음 챙김(mindfulness)의 동양 원형 — 디지털 시대에 가장 절실한 능력.'),
    ('극기 · 논어 안연편 인용', '克 己 復 禮  天 下 歸 仁',
     '자기를 이기고 예로 돌아가면 천하가 인(仁)으로 돌아온다',
     '극기복례(克己復禮) — 공자의 인(仁) 정의. 정주학 자기 통제의 출발점 — 「한 번 끊는 결단」.'),
    ('가도 · 부부의 도', '相 敬 如 賓',
     '서로 손님처럼 공경한다',
     '부부의 도 — 가장 가까운 사람에게 가장 정중할 때 인격이 완성된다. 가족이라서 막대해도 된다는 사고의 정면 부정.'),
    ('치체 · 군주의 도', '其 身 正  不 令 而 行',
     '자기 몸이 바르면, 명령하지 않아도 행해진다',
     '논어 자로편 인용 — 리더십의 본질은 기교가 아닌 자기 정립. 부하에게 요구하기 전에 자기에게 먼저 묻는 습관.'),
    ('교학 · 예기 학기 인용', '敎 學 相 長',
     '가르치고 배움이 서로 자라게 한다',
     '가르치는 일은 자기 공부의 연장. 「가장 잘 가르치는 사람은 자기가 가장 잘 배우는 사람」.'),
    ('장재 · 횡거 사구(四句)', '為 天 地 立 心  為 生 民 立 命  為 往 聖 繼 絕 學  為 萬 世 開 太 平',
     '천지를 위해 마음을 세우고, 백성을 위해 명을 세우며, 옛 성인의 끊긴 학문을 잇고, 만세를 위해 태평을 연다',
     '장재의 횡거사구(橫渠四句) — 조선 선비의 정신적 좌우명. 입지(立志)의 가장 웅장한 표현.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅷ. 한국 수용 ==============
SEC8 = 'Ⅷ. 한국 수용'

@S(SEC8)
def viii_intro(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '고려 말에서 조선까지 — 800년의 책',
              '신유학 도입의 첫 번째 텍스트')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 고려 말 — 안향 등 신유학 도입자들이 원나라에서 수입',
         {'font_size': 18, 'space_before': 4}),
        ('● 1370년 — 진주에서 첫 간행 (이인민이 진주목사로 부임할 때)',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 조선 초 — 사서집주·주자가례와 함께 「신유학 핵심 3권」 중 하나',
         {'font_size': 18, 'space_before': 12}),
        ('● 사서삼경 학습 전 단계의 「성리학 입문 정본」',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 퇴계 이황·율곡 이이를 비롯한 모든 성리학자의 필수 학습서',
         {'font_size': 17, 'space_before': 12}),
        ('● 향약·서원·교육 제도 — 조선 사회 시스템의 사상적 기반',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC8)
def viii_scholars(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '조선의 주요 주석본 — 한국 학자들의 깊이 있는 독해')
    rows = [
        ('퇴계 이황', '退溪 李滉 (1501~1570)',
         '『근사록석의(近思錄釋疑)』\n어려운 부분에 대한 해석.\n조선 성리학의 표준 주석으로 자리잡음.'),
        ('성호 이익', '星湖 李瀷 (1681~1763)',
         '『근사록질서(近思錄疾書)』\n실학적 관점의 재해석.\n전통 주석에서 한 걸음 나간 비판적 독해.'),
        ('한주 정엽', '寒洲 鄭曄 (1563~1625)',
         '근사록에 대한 강설·주해.\n조선 중기 성리학 깊이의 한 단면.'),
        ('기타 학자', '권벌·이덕홍·송시열·정조 등',
         '경연(經筵)에서 왕에게 진강.\n중종·영조·정조가 모두 깊이 중시.\n정조는 1794년 어제서(御製序) 직접 작성.'),
    ]
    for i, (name, en, body) in enumerate(rows):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.4)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(2.2), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.45),
                    name, font_size=17, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.35),
                    en, font_size=11, color=SUB, font_name='Batang')
        add_paragraphs(slide, x + Inches(0.2), y + Inches(0.95), Inches(5.5), Inches(1.2),
                       [(body, {'font_size': 12, 'color': INK})], line_spacing=1.4)


@S(SEC8)
def viii_kingscholar(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '경연(經筵)의 책 — 왕이 읽은 근사록',
              '중종·영조·정조가 모두 깊이 중시한 텍스트')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 중종 시기 — 권벌(權橃)의 진강(進講)이 유명',
         {'font_size': 18, 'space_before': 4}),
        ('● 영조 — 권벌 후손 권만(權萬)을 불러 직접 열람',
         {'font_size': 18, 'space_before': 12}),
        ('● 정조 — 1794년 어제서(御製序)를 지어 그 중요성을 강조',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     「학문에 들어가는 첫 걸음의 진정한 길잡이」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 사대부 — 신유학 입문, 사서삼경의 토대로 학습',
         {'font_size': 17, 'space_before': 14}),
        ('● 향약·서원·과거(科擧) — 모든 교육·정치 제도의 사상적 토대',
         {'font_size': 17, 'space_before': 10}),
        ('● 조선 5백 년 도덕 교육의 이론적 토대',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


# ============== Ⅸ. 오늘 다시 펼치는 이유 ==============
SEC9 = 'Ⅸ. 오늘 다시 펼치는 이유'

@S(SEC9)
def ix_lack(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '오늘의 결핍 지도 — 근사록이 채울 수 있는 자리')
    items = [
        ('이론 과잉, 실천 부재', '추상적 사상은 많지만 「내 일상에 어떻게 적용할 것인가」가 약하다'),
        ('정보 폭발, 집중력 붕괴', '존양편의 「主一無適」 — 한 가지에 집중하는 시간이 사라지고 있다'),
        ('비교 경쟁의 피로',     '위학편의 「위기지학」 — 남에게 보이기 위한 공부가 아니라 자기 변화의 공부'),
        ('관계의 형식주의',     '가도편 「相敬如賓」 — 가장 가까운 사람에게 가장 정중하기를 잊어버렸다'),
        ('롤모델의 부재',       '관성현편 「구체적 모범」 — 추상적 이상이 아닌 구체적 인격이 사람을 바꾼다'),
        ('비판적 사고력 약화',   '변이단편 — 표면이 아닌 전제와 결과를 묻는 분별의 안목'),
    ]
    for i, (cat, desc) in enumerate(items):
        y = Inches(2.4 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.5), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(3.5), Inches(0.6),
                    cat, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.4), y + Inches(0.05), Inches(8.5), Inches(0.55),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC9)
def ix_howread(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '근사록을 오늘 어떻게 읽을까',
              '여조겸·주희의 당부 그대로 — 「가까이서, 절실하게」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 전체를 통독하지 말고, 권 단위로 한 편씩 천천히 읽는다',
         {'font_size': 17, 'space_before': 4}),
        ('● 입문에 좋은 권 — 권2 위학(왜 공부하는가) → 권5 극기(자기 통제) → 권6 가도(가정)',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 「14편을 다 읽기」가 아닌 「한 구절을 마음에 새기기」가 목표',
         {'font_size': 17, 'space_before': 10}),
        ('● 책 이름 그대로 — 「가까이서 생각」: 내 삶에 적용 가능한 한 구절씩',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 자녀와 함께 — 일주일에 한 조(條)씩 읽고 「오늘 우리 가족엔?」 이야기 나눔',
         {'font_size': 16, 'space_before': 10}),
        ('● 주희·여조겸의 당부 — 「입문서이지 결론이 아니다, 갇히지 말라」',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


# ============== Ⅹ. 마무리 ==============
SEC10 = 'Ⅹ. 마무리'

@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '근사록, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 주희와 여조겸이 1175년 한천정사에서 10일간 함께 엮은 신유학 종합 입문서.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 북송 4자(주돈이·정호·정이·장재) 어록 622조를 14편으로 분류.',
         {'font_size': 18, 'space_before': 8}),
        ('● 동심원 구조 — 도체 → 자기 수양 → 가도 → 사회·정치 → 교육·반성 → 분별·모범.',
         {'font_size': 18, 'space_before': 8}),
        ('● 다섯 학습 원리 — 박학·독지·절문·근사·역행.',
         {'font_size': 18, 'space_before': 8}),
        ('● 조선 800년 — 입문 정본·경연의 책·과거의 토대.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「학문을 가까이서 시작해 멀리까지 도달하게 하는」 책.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
                '切 問 而 近 思',
                font_size=80, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.9),
                '仁 在 其 中',
                font_size=64, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.1), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '절 실 하 게  묻 고  가 까 이  생 각 하 면',
                font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.6),
                '인 (仁) 은  그  가 운 데  있 다',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '— 『논어』 자장(子張)편 · 자하(子夏)의 말',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '近  思  錄',
                font_size=20, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\근사록_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
