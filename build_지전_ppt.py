# -*- coding: utf-8 -*-
"""
지전(智典) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
중국 3,000년 역사의 지혜 — 창세·치세·패권·진퇴의 4편
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '지혜로 역사를 논하고, 역사로 지혜를 말한다 · 3,000년 중국사의 지혜 경전',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '智 典',
                font_size=130, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '지 전',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
                '以 智 論 史  以 史 說 智',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '냉성금(冷成金, 1962~2021) 저 · 장연 역 · 김영사 (한국 2003 출간, 전4권)',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '춘추전국 → 양한 → 수당송원 → 명청 — 창세·치세·패권·진퇴의 변증법',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 지전이란 무엇인가'),
        ('Ⅱ.', '방법론 — 이지논사·제자백가 융합'),
        ('Ⅲ.', '4편의 구조 — 창세·치세·패권·진퇴'),
        ('Ⅳ.', '1편 춘추전국 — 창세의 책략'),
        ('Ⅴ.', '2편 전한·후한 — 치세의 지모'),
    ]
    items_right = [
        ('Ⅵ.', '3편 수·당·송·원 — 패권의 쟁투'),
        ('Ⅶ.', '4편 명·청 — 진퇴의 법칙'),
        ('Ⅷ.', '4편을 관통하는 5대 원리'),
        ('Ⅸ.', '명장면·명구 12선'),
        ('Ⅹ.', '오늘 우리에게 · 마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 지전')
    rows = [
        ('원제',   '知典 (2001 중국) — 후에 『三千年来谁著史 (3천년 동안 누가 역사를 썼나)』로 재출간'),
        ('저자',   '냉성금(冷成金, 렁청진, 1962~2021) — 중국인민대학 문학원 교수'),
        ('역자',   '장연 · 한국판 출판사 김영사'),
        ('한국 출간', '2003년 (전2권 세트) → 이후 전4권으로 확장'),
        ('성격',   '중국 3,000년 역사 속 지혜와 권모술수를 분석한 「지혜의 경전」'),
        ('방법론', '以智論史 以史說智 — 지혜로 역사를 논하고, 역사로 지혜를 말한다'),
        ('출전',   '『좌전』·『사기』·『전국책』·『자치통감』 등에서 100여 편의 일화'),
        ('주제',   '제후·재상·책략가·종횡가들의 지혜와 지략 · 경세 철학'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.1), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_author(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '저자 냉성금(冷成金) — 「역사 속 지혜의 발굴자」',
              '중국인민대학 문학원 교수 · 소동파 연구의 권위자')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1962년 산동성(山東省) 출생 — 2021년 3월 9일 사망 (59세)',
         {'font_size': 17, 'space_before': 4}),
        ('● 중국인민대학(中國人民大學) 문학원 교수 · 박사생 지도교수',
         {'font_size': 17, 'space_before': 10}),
        ('● 중국 소식(蘇軾) 연구회 전 부회장 — 소동파 연구의 권위자',
         {'font_size': 17, 'space_before': 10}),
        ('● 대표 저작', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 『智典』 / 『三千年来谁著史』 — 본 발표의 원전',
         {'font_size': 14, 'space_before': 4}),
        ('     · 『读史有智慧』 — 「역사를 읽는 지혜」',
         {'font_size': 14, 'space_before': 4}),
        ('     · 『有一种境界叫苏东坡』 — 「소동파라는 경지」',
         {'font_size': 14, 'space_before': 4}),
        ('     · 『논어의 정신』·『중국권지(中國權智)』',
         {'font_size': 14, 'space_before': 4}),
        ('● 「역사를 살아 있는 지혜로 전환시키는」 글쓰기로 중국·한국 독서계의 사랑을 받음',
         {'font_size': 15, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


@S(SEC1)
def i_what(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '지전 — 단순한 역사책이 아니다',
              '"지혜의 경전" — 100여 편의 역사 이야기로 지혜를 길어내는 책')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「지혜의 경전(智 + 典)」 — 이름 자체가 책의 성격',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 『좌전』·『사기』·『전국책』·『자치통감』 등에서 100여 편의 이야기를 가려 뽑음',
         {'font_size': 17, 'space_before': 10}),
        ('● 제후·재상·책략가·종횡가들의 지혜와 지략·경세 철학을 「이야기식」으로 풀이',
         {'font_size': 17, 'space_before': 10}),
        ('● 유가(儒) · 법가(法) · 도가(道) · 종횡가(縱橫) · 병가(兵) · 음양가(陰陽) 통합',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 1편 = 「천하를 어떻게 얻는가」 · 2편 = 「천하를 어떻게 지키는가」',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 창업(創業)과 수성(守成)의 변증법 — 두 권이 유기적 완결 구조',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅱ. 방법론 ==============
SEC2 = 'Ⅱ. 방법론'

@S(SEC2)
def ii_method(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '이지논사(以智論史)·이사설지(以史說智)',
              '단순 사건 나열이 아닌 「지혜의 원리」 추출')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.7), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.7),
                '以 智 論 史   以 史 說 智',
                font_size=32, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(12.0), Inches(0.7),
                '지혜로 역사를 논하고, 역사로 지혜를 말한다',
                font_size=18, color=SUB,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(0.7), Inches(4.2), Inches(12.0), Inches(3.0), [
        ('● 「以智論史」 — 역사 속 인물의 판단·행동에 담긴 「지혜의 원리」를 분석',
         {'font_size': 17, 'space_before': 6}),
        ('● 「以史說智」 — 추상적 지혜를 「구체적 역사 사건」으로 입증',
         {'font_size': 17, 'space_before': 10}),
        ('● 사기·자치통감이 「무엇이 일어났는가」를 묻는다면,',
         {'font_size': 17, 'space_before': 12}),
        ('     지전은 「왜 그렇게 행동했는가, 우리는 무엇을 배울 것인가」를 묻는다',
         {'font_size': 17, 'space_before': 8, 'bold': True, 'color': ACCENT}),
        ('● 「역사책 + 자기계발서 + 철학서」의 결합 — 동양 「응용 역사학」의 정점',
         {'font_size': 16, 'space_before': 14, 'color': SUB, 'bold': True}),
    ])


@S(SEC2)
def ii_baijia(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '제자백가(諸子百家)의 융합적 해석',
              '한 사건을 6가지 학파의 눈으로 동시에 본다')
    boxes = [
        ('儒 유가', '德治 덕치', '인(仁)·의(義)로 마음을 정복', ACCENT),
        ('法 법가', '法·勢·術', '법으로 질서, 세로 권위, 술로 신하 다스림', INK),
        ('道 도가', '無爲 무위', '억지로 하지 않고 자연의 리듬을 따른다', SUB),
        ('縱橫 종횡가', '合縱連橫', '외교·변론으로 천하를 움직임', ACCENT),
        ('兵 병가', '不戰勝', '싸우지 않고 이김 · 부전승의 철학', INK),
        ('陰陽 음양가', '陰陽五行', '천지의 리듬·기운의 흐름 읽기', SUB),
    ]
    for i, (han, label, desc, color) in enumerate(boxes):
        col, row = i % 3, i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(3.9), Inches(0.8), color)
        add_textbox(slide, x, y, Inches(3.9), Inches(0.8),
                    han, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, x, y + Inches(0.9), Inches(3.9), Inches(0.5), PALE)
        add_textbox(slide, x, y + Inches(0.9), Inches(3.9), Inches(0.5),
                    label, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.5), Inches(0.5),
                    desc, font_size=12, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅲ. 4편의 구조 ==============
SEC3 = 'Ⅲ. 4편의 구조'

@S(SEC3)
def iii_four(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '4편의 시대와 주제 — 「창세 → 치세 → 패권 → 진퇴」',
              '약 2,700년 역사의 변증법적 구조')
    rows = [
        ('1편', '春秋戰國 춘추전국',  'BC 770~221',  '創世 창세의 책략',  '천하를 어떻게 얻는가'),
        ('2편', '前漢·後漢 전한·후한', 'BC 206~AD 220','治世 치세의 지모',  '천하를 어떻게 지키는가'),
        ('3편', '隋唐宋元 수당송원',  'AD 581~1368', '覇權 패권의 쟁투',  '제왕의 기상과 인성의 약점'),
        ('4편', '明淸 명·청',          'AD 1368~1912', '進退 진퇴의 법칙',  '나아가고 물러설 때'),
    ]
    for i, (vol, era, period, theme, msg) in enumerate(rows):
        y = Inches(2.4 + i * 1.0)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.0), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.0), Inches(0.85),
                    vol, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.9), y, Inches(3.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.9), y + Inches(0.05), Inches(3.0), Inches(0.4),
                    era, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(1.9), y + Inches(0.45), Inches(3.0), Inches(0.4),
                    period, font_size=12, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.1), y + Inches(0.1), Inches(3.5), Inches(0.7),
                    theme, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(8.7), y + Inches(0.1), Inches(4.2), Inches(0.7),
                    msg, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
                '4편이 모여 「창업과 수성·도전과 절제」의 인생 사이클을 그려낸다',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅳ. 1편 춘추전국 ==============
SEC4 = 'Ⅳ. 1편 춘추전국'

@S(SEC4)
def iv_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '1편 — 창세(創世)의 책략',
              '주 왕실의 권위가 무너진 550년 · 새로운 질서를 세우는 책략들')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● BC 770 주의 동천 → BC 221 진의 통일 · 약 550년의 대분열기',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 「예악정벌이 천자에게서 나온다」는 원칙 붕괴 — 권력의 진공',
         {'font_size': 17, 'space_before': 10}),
        ('● 제자백가의 출현 — 사상과 정치 전략의 동시 폭발',
         {'font_size': 17, 'space_before': 10}),
        ('● 4부 구성', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 제1부 정치의 지혜 — 덕치와 권모, 양모와 음모',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제2부 전략적 전술 — 이(利)를 중심으로 한 합종연횡',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제3부 인품과 운명 — 개인의 성품과 역사적 운명의 관계',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제4부 현자와 국가 — 인재 등용술과 국가 흥망',
         {'font_size': 14, 'space_before': 4}),
        ('● 핵심 메시지 — 「누가 새로운 질서를 세울 것인가」',
         {'font_size': 16, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


@S(SEC4)
def iv_guanzhong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '관중과 제환공 — 「원수를 재상으로」',
              '유가적 용인(用人)의 극치 — 사사로운 원한을 넘는 대의')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 관중(管仲)이 공자규(公子糾)를 위해 소백(훗날 제환공)에게 활을 쏨 — 사구(射鉤)',
         {'font_size': 16, 'space_before': 4}),
        ('● 살아 돌아온 제환공이 즉위 — 환공은 관중을 「원수」로 여김',
         {'font_size': 16, 'space_before': 10}),
        ('● 그러나 포숙아(鮑叔牙)의 추천 — 「관중이 있는 나라가 천하를 얻습니다」',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 환공의 결단 — 사사로운 원한을 접고 관중을 재상으로 등용, 「중부(仲父)」로 모심',
         {'font_size': 16, 'space_before': 10}),
        ('● 관중의 정치 — 「倉廩實而知禮節 衣食足而知榮辱」 (창고가 차야 예절을 안다)',
         {'font_size': 15, 'space_before': 10, 'font_name': 'Batang'}),
        ('● 결과 — 「九合諸侯」 아홉 차례 회맹·존왕양이로 춘추 첫 패자가 됨',
         {'font_size': 16, 'space_before': 10}),
        ('● 「生我者父母 知我者鮑子也」 — 관중의 회고 · 「지인(知人)」이 곧 위대한 능력',
         {'font_size': 15, 'space_before': 10, 'color': SUB, 'bold': True, 'font_name': 'Batang'}),
    ])


@S(SEC4)
def iv_shangyang(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '상앙의 변법 — 「나무 하나로 천하의 신뢰를 사다」',
              '사문입목(徙木立信) — 법가 사상의 가장 극적인 실현')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 위(衛) 출신 상앙 — 위 혜왕이 등용도 살해도 안 함 → 진(秦) 효공에 등용',
         {'font_size': 16, 'space_before': 4}),
        ('● 세 번의 면담 — 帝道·王道엔 졸고 覇道(부국강병)엔 무릎을 침',
         {'font_size': 16, 'space_before': 10}),
        ('● 변법의 첫 장애물 — 백성의 불신 (「이전 정권들이 늘 약속을 어겼다」)',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 「사문입목(徙木立信)」 — 도성 남문에 나무를 세우고',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     · 「북문으로 옮기는 자에게 10금」 → 아무도 안 함 → 50금으로 인상',
         {'font_size': 14, 'space_before': 6}),
        ('     · 한 사람이 옮기자 즉시 50금 지급 → 신뢰 확립 → 변법 시행',
         {'font_size': 14, 'space_before': 4}),
        ('● 진을 부강케 해 통일의 기초를 놓음 — 그러나 상앙 자신은 거열형(車裂)으로 죽음',
         {'font_size': 16, 'space_before': 12}),
        ('● 「제도를 만든 자는 종종 그 제도로 죽는다」 — 법가의 비극',
         {'font_size': 15, 'space_before': 10, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC4)
def iv_hengjong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '합종연횡 — 종횡가의 천하 외교 시대',
              '소진의 합종(合縱) vs 장의의 연횡(連橫)')
    cols = [
        ('蘇 秦', '소진', '합종(合縱)',
         '낙양 평민 → 6국 동맹\n6국 재상의 인을 동시에 차다\n진은 15년간 함곡관을 못 나옴'),
        ('張 儀', '장의', '연횡(連橫)',
         '「내 혀가 아직 있느냐?」\n진을 중심으로 6국을 차례로 분열\n6국 합종을 무너뜨려 진 통일의 토대'),
    ]
    for i, (han, kor, label, body) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), ACCENT if i == 0 else INK)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    f'{kor} — {label}', font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.4), Inches(3.6), Inches(5.1), Inches(3.5),
                       [(body, {'font_size': 14, 'color': INK})], line_spacing=1.5)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '귀곡선생(鬼谷子)의 두 제자 — 같은 스승, 정반대의 길 · 외교 전략의 영원한 두 원형',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅴ. 2편 전한·후한 ==============
SEC5 = 'Ⅴ. 2편 전한·후한'

@S(SEC5)
def v_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '2편 — 치세(治世)의 지모',
              '천하를 지키는 법 · BC 206~AD 220 약 420년')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 진(秦) 통일 15년 만에 멸망 → 초한쟁패 → 한(漢) 건국 → 왕망의 찬탈 → 광무제 중흥',
         {'font_size': 17, 'space_before': 4}),
        ('● 1편이 「창세(창업)」였다면 2편은 「치세(수성)」 — 관점의 근본적 전환',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 난세의 영웅 → 태평성대의 통치자',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 전쟁터의 전략 → 궁정의 정치술',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 3부 구성', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 제1부 영웅과 시대 — 초한쟁패, 영웅의 조건',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제2부 성세와 난세 — 왕조의 흥망성쇠 순환 원리',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제3부 선비와 정치 — 지식인과 권력의 관계',
         {'font_size': 14, 'space_before': 4}),
    ])


@S(SEC5)
def v_chuhan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '초한쟁패 — 항우 vs 유방, 영웅의 조건',
              '중국사 가장 극적인 대결 · 「누가 진짜 영웅인가」')
    cols = [
        ('項 羽', '항우', '비극의 영웅',
         '거록(鉅鹿) 「파부침주(破釜沈舟)」\n역발산기개세(力拔山氣蓋世)\n\n약점 — 독단·아집·인재 경시\n범증·한신·진평을 다 잃음\n\n해하(垓下) 사면초가 → 오강에서 자결\n「내가 무슨 면목으로 강동 부로를…」', INK),
        ('劉 邦', '유방', '난세의 승자',
         '패현의 정장(亭長) 출신\n\n약법삼장(約法三章) — 민심 장악\n장량·소하·한신 3걸을 모두 안음\n\n「내 능력이 그들만 못하나,\n그들을 쓰는 능력이 있다」\n「의심나면 쓰지 말고, 쓰면 의심 말라」', ACCENT),
    ]
    for i, (han, kor, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    f'{kor} — {label}', font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.4), Inches(3.6), Inches(5.1), Inches(3.5),
                       [(body, {'font_size': 13, 'color': INK})], line_spacing=1.5)


@S(SEC5)
def v_three_genius(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '한고조의 「삼걸(三傑)」 — 인재가 천하를 만든다',
              '장량·소하·한신 — 유방이 항우를 이긴 진짜 이유')
    rows = [
        ('張 良 장량',  '책사(策士)',
         '귀신 같은 책략의 천재 · 「운주유악(運籌帷幄) 결승천리(決勝千里)」\n홍문연·약법삼장의 핵심 조언자 · 황석공에게 받은 삼략의 책사'),
        ('蕭 何 소하',  '재상(宰相)',
         '관중의 안정·군량의 확보 · 한신을 알아본 「소하월하추한신(蕭何月下追韓信)」\n「내가 없이 한신을 추천하면 누구도 안 했을 것」'),
        ('韓 信 한신',  '명장(名將)',
         '국사무쌍(國士無雙)의 명장 · 사면초가·배수진의 천재 전략가\n그러나 후에 「토사구팽」으로 죽음 — 천하 평정 후 명장의 운명'),
    ]
    for i, (name, role, body) in enumerate(rows):
        y = Inches(2.3 + i * 1.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(1.35), ACCENT)
        add_textbox(slide, Inches(0.7), y + Inches(0.15), Inches(2.4), Inches(0.5),
                    name, font_size=17, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), y + Inches(0.7), Inches(2.4), Inches(0.5),
                    role, font_size=12, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(3.3), y + Inches(0.15), Inches(9.6), Inches(1.2),
                       [(body, {'font_size': 13, 'color': INK})], line_spacing=1.4)


@S(SEC5)
def v_wangmang(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '왕망의 찬탈과 광무제 — 「수성의 어려움」',
              '한 왕조의 사이클 — 위기와 중흥의 시소')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 왕망(王莽, AD 9~23) — 외척 출신 · 「겸손한 군자」의 가면을 쓴 찬탈자',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('     · 30년에 걸친 신중한 정치적 등반 — 「위장된 덕」의 극치',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 신(新) 왕조 건국 — 그러나 무리한 개혁으로 15년 만에 멸망',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 광무제 유수(劉秀, AD 25~57) — 한실 종친 출신 후한의 중흥조',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 「유수가 곤양(昆陽)에서 1만으로 왕망의 40만을 격파」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 「공신을 죽이지 않고도 천하를 지킨」 드문 황제 — 한신과 정반대',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 교훈 — 「수성(守成)은 창업(創業)보다 어렵다」 — 당 태종이 위징에게 물은 그 질문',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


# ============== Ⅵ. 3편 수당송원 ==============
SEC6 = 'Ⅵ. 3편 수당송원'

@S(SEC6)
def vi_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '3편 — 패권(覇權)의 쟁투',
              '제왕의 기상과 인성의 약점 · AD 581~1368 약 780년')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 수·당·송·원 — 통일·분열·재통일·이민족 정복의 격동기',
         {'font_size': 17, 'space_before': 4}),
        ('● 3편 구성', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 제1편 제왕의 기상 — 풍모와 인재 전략',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제2편 인성의 약점을 이겨라 — 인간 본성의 함정과 극복',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제3편 역사의 성패로 영웅을 논하다 — 성패의 기준',
         {'font_size': 14, 'space_before': 4}),
        ('● 핵심 인물 — 당 태종·측천무후·송 태조·왕안석·주원장 등',
         {'font_size': 17, 'space_before': 12}),
        ('● 핵심 메시지 — 「큰 인물도 작은 약점에 무너진다」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 인성의 약점 — 의심·시기·교만·노여움·태만 등을 어떻게 다스리는가',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC6)
def vi_taizong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '당 태종 이세민 — 「제왕의 기상」의 표상',
              '정관(貞觀)의 치 · 위징을 「인간 거울」로 삼은 군주')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「以銅爲鏡可正衣冠 以古爲鏡可知興替 以人爲鏡可明得失」',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     구리로 거울 삼으면 의관을, 옛것으로 거울 삼으면 흥망을, 사람으로 거울 삼으면 득실을',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 위징(魏徵) — 원수 진영(이건성) 출신을 등용 → 200여 회 간언',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 위징 사후 — 「내가 거울 하나를 잃었다」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 정관(貞觀)의 치 — 「창고가 차고 거리에 도둑이 없으며」',
         {'font_size': 17, 'space_before': 12}),
        ('● 그러나 만년에 — 고구려 정벌 실패·태자 폐립 등 「큰 군주도 늙으면 흔들린다」',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 『정관정요(貞觀政要)』 — 당 태종과 신하들의 대화를 정리한 동양 제왕학의 최고봉',
         {'font_size': 15, 'space_before': 10, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC6)
def vi_wuzetian(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '측천무후 — 중국 유일의 여황제 · 인성의 양면',
              '권력을 위한 무자비함과 통치의 탁월함의 공존')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 측천무후(則天武后, 624~705) — 중국 역사상 유일한 여황제',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 태종의 후궁 → 고종의 황후 → 무주(武周) 황제 (재위 690~705)',
         {'font_size': 17, 'space_before': 10}),
        ('● 권력 장악의 무자비함', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 친자식까지 의심하면 가차없이 제거 · 혹리(酷吏)를 이용한 공포 정치',
         {'font_size': 14, 'space_before': 4}),
        ('● 그러나 통치의 탁월함', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 과거 제도의 확대 — 「전시(殿試)·무거(武擧)」 신설',
         {'font_size': 14, 'space_before': 4}),
        ('     · 적인걸(狄仁傑) 등 명재상 등용 · 농업 진흥 · 경제 안정',
         {'font_size': 14, 'space_before': 4}),
        ('     · 자기 사후 묘비를 무자비(無字碑)로 — 「공과는 후세가 판단하라」',
         {'font_size': 14, 'space_before': 4}),
        ('● 「선악을 한 사람 안에서 다 보여준」 역사의 가장 복잡한 인물',
         {'font_size': 15, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


@S(SEC6)
def vi_song(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '송 태조와 「배주석병권(杯酒釋兵權)」',
              '술 한 잔으로 무인의 권력을 거두다 · 가장 우아한 정치 술')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 송 태조 조광윤(趙匡胤) — 「진교병변」으로 후주를 찬탈해 송 건국',
         {'font_size': 17, 'space_before': 4}),
        ('● 가장 큰 두려움 — 「내가 한 짓을 부하 장수가 또 할까봐」',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 한고조의 한신 처형, 명 태조의 공신 숙청과는 다른 길',
         {'font_size': 17, 'space_before': 10}),
        ('● 杯酒釋兵權(배주석병권) — 술 한 잔으로 병권을 풀다',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     · 부하 장수들을 술자리에 모아 「부귀를 누리며 안락하게 사세요」',
         {'font_size': 14, 'space_before': 6}),
        ('     · 다음날 모두 사직 → 거대한 부와 영지로 보상',
         {'font_size': 14, 'space_before': 4}),
        ('● 한 방울의 피도 흘리지 않고 군권을 황제에게 집중 — 「우아한 권력 정리」의 모범',
         {'font_size': 16, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 그러나 부작용 — 무력 약화로 송은 끊임없이 북방 이민족에 시달림',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅶ. 4편 명·청 ==============
SEC7 = 'Ⅶ. 4편 명·청'

@S(SEC7)
def vii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '4편 — 진퇴(進退)의 법칙',
              '나아갈 때와 물러설 때 · AD 1368~1912 약 540년')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 명(明) 1368 건국 → 청(淸) 1644 입관 → 1912 신해혁명',
         {'font_size': 17, 'space_before': 4}),
        ('● 통일 제국 후기의 지혜 — 「세상이 안정될수록 진퇴가 어렵다」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 3편 구성', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 제1편 개국과 고권(固權) — 나라를 세우고 권력을 굳히기',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제2편 대재(大才)와 소재(小才) — 큰 인재와 작은 인재의 차이',
         {'font_size': 14, 'space_before': 4}),
        ('     · 제3편 정치는 대지혜이다 — 정치의 본질과 진퇴의 법칙',
         {'font_size': 14, 'space_before': 4}),
        ('● 핵심 인물 — 주원장·장거정·강희제·옹정제·건륭제·증국번·이홍장 등',
         {'font_size': 17, 'space_before': 12}),
        ('● 핵심 메시지 — 「물러날 줄 아는 자가 가장 멀리 간다」',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC7)
def vii_zhuyuanzhang(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '주원장 — 거지에서 황제로, 가장 잔혹한 개국조',
              '명의 개국과 「공신 숙청」의 극단')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 주원장(朱元璋, 1328~1398) — 빈농 출신 · 거지·승려 → 황제',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 원나라를 무너뜨리고 명 건국 — 한족 회복의 황제',
         {'font_size': 17, 'space_before': 10}),
        ('● 그러나 가장 잔혹한 개국조 — 「누구도 믿지 않는다」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 호유용(胡惟庸)의 옥 — 3만여 명 처형 · 승상 제도 폐지',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 남옥(藍玉)의 옥 — 또 1만 5천여 명 처형 · 공신 거의 전멸',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 「토사구팽」의 극단적 실현 — 한신·팽월·영포의 운명을 대규모로',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 송 태조의 「배주석병권」과의 정반대 — 같은 문제, 정반대의 해법',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
        ('● 「창업의 잔혹함」과 「수성의 안정」 사이의 가장 극적인 사례',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC7)
def vii_zengguofan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '증국번(曾國藩) — 청 말의 「물러설 줄 아는 자」',
              '태평천국 진압의 영웅 · 그러나 권력을 스스로 내려놓다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 증국번(曾國藩, 1811~1872) — 청 말 한족 출신 명재상',
         {'font_size': 17, 'space_before': 4}),
        ('● 태평천국의 난(1851~1864) 진압의 핵심 인물 — 상군(湘軍) 창설',
         {'font_size': 17, 'space_before': 10}),
        ('● 천하의 영웅이 되었으나 — 스스로 「병권을 거두고 군을 해산」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 한신·악비의 운명을 알았기에 — 절정에서 내려옴',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 「中天而退」 — 한가운데서 물러난다',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 자기 일기·가서(家書)를 후세에 남겨 — 동양 자기 수양의 모범',
         {'font_size': 16, 'space_before': 10}),
        ('● 4편의 정점 — 「진퇴의 법칙」의 살아 있는 사례',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
        ('● 마오쩌둥·장개석 모두 증국번을 「자기가 배우려는 모범」으로 꼽음',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC7)
def vii_kanglongqian(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '강·옹·건 — 청의 130년 황금기',
              '강희(61) → 옹정(13) → 건륭(60) — 3대가 만든 「대청 성세」')
    rows = [
        ('康熙 강희', '재위 61년 (1661~1722)',
         '8세에 즉위 · 오삼계의 난·정성공 평정 · 러시아·갈단 격퇴\n학자형 군주 — 「강희자전」 편찬 · 유학 진흥'),
        ('雍正 옹정', '재위 13년 (1722~1735)',
         '잔혹한 정리자 — 형제와 신하를 가차없이 숙청\n그러나 부패 척결·재정 정비 — 청의 130년 성세를 가능케 함\n「큰 일을 위해 작은 일을 잔혹하게」'),
        ('乾隆 건륭', '재위 60년 (1735~1795)',
         '청의 절정 — 영토 최대·인구 폭증\n그러나 만년 — 화신(和珅) 등 부패 만연\n「절정이 곧 쇠퇴의 시작」'),
    ]
    for i, (name, period, body) in enumerate(rows):
        y = Inches(2.3 + i * 1.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(1.35), ACCENT)
        add_textbox(slide, Inches(0.7), y + Inches(0.15), Inches(2.4), Inches(0.5),
                    name, font_size=17, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), y + Inches(0.7), Inches(2.4), Inches(0.5),
                    period, font_size=11, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(3.3), y + Inches(0.15), Inches(9.6), Inches(1.2),
                       [(body, {'font_size': 13, 'color': INK})], line_spacing=1.4)


# ============== Ⅷ. 5대 원리 ==============
SEC8 = 'Ⅷ. 4편 관통 5대 원리'

@S(SEC8)
def viii_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '4편을 관통하는 5대 지혜 원리',
              '냉성금이 100여 편의 일화에서 길어낸 보편 원리')
    items = [
        ('1', '創 業 / 守 成', '창업과 수성',  '천하를 얻는 지혜와 지키는 지혜는 다르다 — 1편↔2편'),
        ('2', '知 人 / 用 人', '지인과 용인',  '인재를 알아보는 안목이 곧 군주의 가장 큰 능력 — 관중·삼걸'),
        ('3', '進 退 之 道',   '진퇴의 도',    '나아갈 때 가는 자보다 물러설 때 멈출 줄 아는 자가 멀리 간다'),
        ('4', '德 治 / 法 治', '덕치와 법치',  '인(仁)으로 마음을 정복하면서도 법(法)으로 질서를 세운다'),
        ('5', '人 性 / 大 義', '인성과 대의',  '인성의 약점(의심·시기·교만)을 다스려 대의를 이루는 자가 영웅'),
    ]
    for i, (num, han, kor, desc) in enumerate(items):
        y = Inches(2.3 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(0.7), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(0.7), Inches(0.7),
                    num, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.5), y, Inches(3.0), Inches(0.7), PALE)
        add_textbox(slide, Inches(1.5), y, Inches(3.0), Inches(0.7),
                    han, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.7), y + Inches(0.05), Inches(2.2), Inches(0.6),
                    kor, font_size=13, color=SUB, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.0), y + Inches(0.05), Inches(5.9), Inches(0.6),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅸ. 명장면·명구 12선 ==============
SEC9 = 'Ⅸ. 명장면·명구 12선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC9)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC9} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=16, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 16, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('관중 · 1편', '生 我 者 父 母  知 我 者 鮑 子 也',
     '나를 낳아준 것은 부모이지만, 나를 알아준 것은 포숙아이다',
     '제환공의 원수 관중을 추천한 포숙아 — 「지인(知人)」이 곧 위대한 능력. 인재를 알아보는 안목이 곧 권력의 가장 큰 자산.'),
    ('관자 · 1편', '倉 廩 實 而 知 禮 節  衣 食 足 而 知 榮 辱',
     '창고가 차야 예절을 알고, 입고 먹는 것이 족해야 영욕을 안다',
     '관중의 실용 정치 철학. 도덕 이전에 경제 안정 — 맹자의 「항산항심(恒産恒心)」과 짝.'),
    ('상앙 · 1편', '徙 木 立 信',
     '나무를 옮겨 신뢰를 세우다',
     '상앙의 변법. 「작은 약속을 어김없이 지킴」으로써 큰 개혁의 신뢰를 쌓는다 — 현대 조직 변화 관리의 동양 원형.'),
    ('소진 · 1편', '寧 爲 鷄 口  無 爲 牛 後',
     '닭의 부리가 될지언정 소의 꼬리가 되지 말라',
     '소진이 한 선혜왕을 합종으로 설득할 때 — 작은 나라의 주인이 될지언정 큰 나라의 종이 되지 말라.'),
    ('항우 · 2편', '力 拔 山 兮  氣 蓋 世',
     '힘은 산을 뽑고 기개는 세상을 덮건만',
     '해하 사면초가의 밤 · 항우의 최후 노래(垓下歌). 「힘」만으로는 천하를 못 얻는다 — 동양 비극의 정점.'),
    ('유방 · 2편', '約 法 三 章',
     '법 세 가지로 약속한다',
     '관중 입성 후 유방의 선언 — 살인자死·상해자 도둑은 처벌, 나머지 가혹한 법은 모두 폐지. 민심 장악의 결정타.'),
    ('한고조 · 2편', '運 籌 帷 幄  決 勝 千 里',
     '장막 안에서 책략을 짜 천 리 밖에서 승부를 결정한다',
     '유방의 장량(張良) 평. 책사의 가치 — 직접 싸우지 않고도 천 리 밖의 전투를 결정짓는 능력.'),
    ('한신 · 2편', '兎 死 狗 烹',
     '토끼가 죽으면 사냥개를 삶는다',
     '한신의 최후 한탄. 「국사무쌍」의 명장도 천하 평정 후엔 위협 — 공신의 영원한 운명.'),
    ('당 태종 · 3편', '以 人 爲 鏡  可 明 得 失',
     '사람으로 거울을 삼으면 득실을 밝힐 수 있다',
     '위징(魏徵)에 대한 평. 「인간 거울」을 곁에 두는 군주가 큰 군주 — 「내가 거울 하나를 잃었다」 한 위징 사후의 탄식.'),
    ('송 태조 · 3편', '杯 酒 釋 兵 權',
     '술 한 잔으로 병권을 풀다',
     '송 태조 조광윤이 부하 장수들을 술자리에 모아 군권을 거둠 — 한 방울 피도 흘리지 않는 「우아한 권력 정리」의 모범.'),
    ('주원장 · 4편', '狡 兎 死  良 狗 烹',
     '교활한 토끼가 죽으면 좋은 사냥개를 삶는다',
     '주원장의 호유용·남옥의 옥 — 한신의 운명을 「국가 규모」로 실행. 송 태조의 배주석병권과 정반대의 길.'),
    ('증국번 · 4편', '中 天 而 退',
     '한가운데서 물러난다',
     '태평천국 평정 후 증국번이 스스로 군권을 내려놓음. 4편 「진퇴의 법칙」의 정점 — 「가장 잘 나갈 때 멈출 줄 아는 자」가 가장 멀리 간다.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅹ. 오늘 우리에게 + 마무리 ==============
SEC10 = 'Ⅹ. 오늘 우리에게'

@S(SEC10)
def x_today(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '지전이 오늘 우리에게 일러주는 10가지')
    items = [
        '1. 지혜로 역사를 논하고, 역사로 지혜를 말하라 — 以智論史 以史說智',
        '2. 창업과 수성은 다른 지혜 — 천하를 얻는 법과 지키는 법이 같지 않다',
        '3. 인재를 알아보는 안목이 가장 큰 능력 — 知人이 用人에 앞선다',
        '4. 가장 위태로운 적은 외부가 아니라 자기 안의 약점이다',
        '5. 사사로운 원한을 넘어 대의를 보는 자가 큰 그릇이다 — 관중·환공',
        '6. 작은 약속을 어김없이 지켜야 큰 신뢰가 선다 — 사문입목',
        '7. 「토사구팽」을 알았다면 절정에서 내려와라 — 한신과 증국번',
        '8. 한 방울 피 없이 권력을 정리하는 길도 있다 — 배주석병권',
        '9. 가장 큰 권력은 가장 큰 두려움을 만든다 — 주원장의 잔혹함',
        '10. 진퇴(進退)의 도 — 「中天而退」, 한가운데서 물러설 줄 아는 자',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.3 + i * 0.45)
        add_textbox(slide, Inches(0.7), y, Inches(12.0), Inches(0.4),
                    txt, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '지전, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 중국 인민대학 냉성금 교수가 100여 편의 역사 일화로 엮은 「지혜의 경전」.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 4편 — 춘추전국(창세) · 양한(치세) · 수당송원(패권) · 명청(진퇴).',
         {'font_size': 18, 'space_before': 8}),
        ('● 방법론 — 以智論史 以史說智 · 단순 역사가 아닌 「지혜의 원리」 추출.',
         {'font_size': 18, 'space_before': 8}),
        ('● 제자백가 통합 — 유·법·도·종횡·병·음양가의 6눈으로 한 사건을 본다.',
         {'font_size': 18, 'space_before': 8}),
        ('● 100여 명의 인물 — 관중·상앙·소진·장의·항우·유방·한신·당 태종·송 태조·주원장·증국번.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「창업과 수성·진퇴와 인성」의 영원한 변증법을 보여주는 책.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
                '以 智 論 史',
                font_size=84, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.9),
                '以 史 說 智',
                font_size=60, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.3), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
                '지 혜 로  역 사 를  논 하 고',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
                '역 사 로  지 혜 를  말 한 다',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '— 냉성금, 『지전(智典) · 3,000년 동안 누가 역사를 썼나』',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '智  典',
                font_size=22, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\지전.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
