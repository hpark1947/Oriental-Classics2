# -*- coding: utf-8 -*-
"""
육도삼략(六韜三略) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '제왕의 병서 · 손자병법의 결정적 짝 · 무경칠서의 두 자리',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '六 韜  ·  三 略',
                font_size=86, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '육 도  ·  삼 략',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '天 下 非 一 人 之 天 下  乃 天 下 之 天 下',
                font_size=17, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '천하는 한 사람의 천하가 아니라 천하 사람들의 천하다',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '강태공(姜太公) · 황석공(黃石公) — 두 책의 합본 · 육도 60편 + 삼략 3편',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 한 권이 아닌 두 권'),
        ('Ⅱ.', '두 책의 전승 — 강태공과 황석공'),
        ('Ⅲ.', '무경칠서에서의 자리'),
        ('Ⅳ.', '육도 6권의 흐름 — 文에서 武로'),
        ('Ⅴ.', '육도 6권 깊이 읽기'),
        ('Ⅵ.', '삼략 3편 깊이 읽기'),
    ]
    items_right = [
        ('Ⅶ.', '제왕학 7대 원칙'),
        ('Ⅷ.', '손자병법과의 비교'),
        ('Ⅸ.', '명구 20선'),
        ('Ⅹ.', '후대 영향 — 한국·일본'),
        ('Ⅺ.', '오늘 우리에게'),
        ('Ⅻ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_two(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '육도삼략은 「한 권」이 아니다',
              '서로 다른 두 책 — 시대도 저자도 문체도 다른 합본')
    cols = [
        ('六 韜', '육도',
         '주(周) 문왕·무왕 ↔ 강태공(姜太公)의 문답\n\n6권 60편 · 수만 자\n\n종합 통치 + 군사 전략서\n전국시대 말~한대 편찬 추정',
         ACCENT),
        ('三 略', '삼략',
         '진말(秦末) 황석공(黃石公)이\n장량(張良)에게 전한 비전서\n\n상·중·하 3편 · 약 3,800자\n\n리더십·정치 철학서\n전한 말~후한 편찬 추정',
         INK),
    ]
    for i, (han, kor, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    kor, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.3), Inches(3.6), Inches(5.5), Inches(3.5),
                       [(body, {'font_size': 15, 'color': INK})], line_spacing=1.5)


@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '이름의 뜻',
              '「韜(도)」와 「略(략)」 — 두 글자가 가르치는 것')
    cols = [
        ('韜 도',
         '활이나 검을 싸는 주머니\n→ 「감춤·깊이 숨긴 것」\n\n「육도」 = 여섯 가지 심오한\n비밀의 가르침',
         '문(文)·무(武)·용(龍)·\n호(虎)·표(豹)·견(犬)'),
        ('略 략',
         '간략한 요약 · 핵심 전략\n→ 「큰 줄기로 가려 뽑은 것」\n\n「삼략」 = 세 단계의 전략\n— 짧지만 응축의 극치',
         '상략(上略)·중략(中略)·\n하략(下略)'),
    ]
    for i, (han, body, foot) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), ACCENT)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(1.0),
                    han, font_size=36, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_paragraphs(slide, x + Inches(0.3), Inches(3.5), Inches(5.5), Inches(2.5),
                       [(body, {'font_size': 15, 'color': INK})], line_spacing=1.5)
        add_filled_rect(slide, x, Inches(6.1), Inches(5.9), Inches(0.9), PALE)
        add_textbox(slide, x, Inches(6.1), Inches(5.9), Inches(0.9),
                    foot, font_size=14, color=ACCENT, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_status(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 육도삼략')
    rows = [
        ('편자(전승)',  '육도 — 태공망 여상(呂尙) · 삼략 — 황석공(黃石公)'),
        ('실제 편찬',   '전국시대 말 ~ 진한(秦漢) 교체기'),
        ('편수',        '육도 6권 60편 + 삼략 3편'),
        ('전체 분량',    '수만 자 — 손자병법(약 6,000자)의 여러 배'),
        ('성격',        '종합 통치서 + 군사 전략 + 리더십 철학'),
        ('주된 독자',    '「장수」가 아닌 「군주(왕)」'),
        ('위상',        '송대 무경칠서(武經七書) 7종 중 2종'),
        ('수용 범위',    '중국·한국·일본·베트남 — 동아시아 군사학의 정전'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_thought(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '핵심 사상 네 기둥')
    boxes = [
        ('天 下 公 有', '천하공유', '천하는 한 사람의 것이 아니다 — 이익을 나누는 자가 얻는다'),
        ('得 賢 者 昌', '득현자창', '인재 중심주의 — 사람을 얻으면 흥, 잃으면 망'),
        ('柔 能 制 剛', '유능제강', '부드러움이 굳셈을 제어한다 — 유와 강의 조화'),
        ('師 出 以 義', '사출이의', '군대는 의(義)로써 나선다 — 의롭지 않은 전쟁은 패한다'),
    ]
    for i, (han, kor, desc) in enumerate(boxes):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(6.0), Inches(2.0), PALE)
        add_textbox(slide, x, y + Inches(0.25), Inches(6.0), Inches(0.7),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(1.0), Inches(6.0), Inches(0.4),
                    kor, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.3), y + Inches(1.4), Inches(5.4), Inches(0.55),
                    desc, font_size=13, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 전승 ==============
SEC2 = 'Ⅱ. 두 책의 전승'

@S(SEC2)
def ii_taegong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '강태공(姜太公) — 위수의 낚시꾼, 주 800년의 설계자',
              '70세까지 학문에 전념 · 문왕을 만나 왕사(王師)가 되다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 이름 — 여상(呂尙), 성은 강(姜) · BC 11세기 은말주초',
         {'font_size': 17, 'space_before': 4}),
        ('● 위수(渭水) 가에서 낚시 — 문왕을 만나 등용 (BC 1100년경)',
         {'font_size': 17, 'space_before': 8}),
        ('● 문왕의 점괘 — 「용도 호랑이도 아닌, 公侯를 얻으리라」',
         {'font_size': 16, 'space_before': 8, 'color': SUB}),
        ('● 「우리 선왕 太公이 바라던(望) 분」 — 「太公望」의 유래',
         {'font_size': 16, 'space_before': 8, 'color': SUB, 'font_name': 'Batang'}),
        ('● 무왕을 도와 목야(牧野)에서 은 주왕을 멸 — 주나라 800년 왕조의 토대',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 제(齊)나라 초대 제후 — 동방의 새 왕조를 건설',
         {'font_size': 17, 'space_before': 8}),
        ('● 중국 최초의 「王師(왕의 스승)」의 원형 — 후대 모든 책사의 모범',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC2)
def ii_huang(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '황석공(黃石公) — 하비교의 노인, 한 왕조의 비밀 스승',
              '신발 세 번의 시험 · 장량에게 병서를 전수하다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 진(秦) 말 · BC 3세기 — 장량이 하비(下邳)에 숨어 있을 때',
         {'font_size': 17, 'space_before': 4}),
        ('● 다리 위의 노인 — 신발을 다리 아래로 던지며 「젊은이, 주워오라」',
         {'font_size': 17, 'space_before': 8}),
        ('● 새벽 약속을 세 번 시험 — 장량이 끝까지 화내지 않고 통과',
         {'font_size': 17, 'space_before': 8, 'color': SUB}),
        ('● 「태공망의 병법(곧 삼략)」을 전수 — 「이를 읽으면 王師가 된다」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 13년 뒤 곡성산(穀城山)의 황석(黃石) — 「그것이 나다」라는 수수께끼',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 장량은 이 병법으로 유방을 도와 한(漢) 400년 왕조를 세움',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 「전승의 서사」가 책의 성격을 결정 — 「왕조를 세우는 병법」의 메시지',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
    ])


@S(SEC2)
def ii_kingdom(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '두 전승이 주고자 한 메시지',
              '「우리는 왕조를 세우는 병법이다」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 손자병법 — 「전장의 병법」 · 이미 결정된 전쟁에서 어떻게 이길 것인가',
         {'font_size': 17, 'space_before': 6}),
        ('● 육도삼략 — 「천하 창건의 병법」 · 어떻게 왕조를 세우고 다스릴 것인가',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 강태공의 육도 → 주(周) 800년',
         {'font_size': 17, 'space_before': 14}),
        ('● 황석공의 삼략 → 한(漢) 400년',
         {'font_size': 17, 'space_before': 8}),
        ('● 즉, 두 책은 합하여 「1,200년의 왕조를 만든 비전서」의 신화를 갖는다',
         {'font_size': 17, 'space_before': 14, 'bold': True, 'color': ACCENT}),
        ('● 「왕이 이 책을 만나면 왕조를 열고, 장수가 만나면 천하를 얻는다」',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
    ])


# ============== Ⅲ. 무경칠서 ==============
SEC3 = 'Ⅲ. 무경칠서'

@S(SEC3)
def iii_table(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '무경칠서(武經七書) — 동아시아 군사학의 7대 정전',
              '송 신종(1043) 공식 선정 · 육도·삼략이 두 자리 차지')
    rows = [
        ('손자병법(孫子兵法)',  '손무(孫武)',    '전쟁의 원리 — 부전승의 철학'),
        ('오자병법(吳子兵法)',  '오기(吳起)',    '군사 행정 · 장수와 사기'),
        ('六 韜  육도',           '태공망(太公望)', '종합 통치 · 인재 · 전략 · 전술', True),
        ('三 略  삼략',           '황석공(黃石公)', '리더십 · 강유 조화 · 정당한 전쟁', True),
        ('사마법(司馬法)',      '전양저(田穰苴)', '군례(軍禮) · 군의 윤리'),
        ('울료자(尉繚子)',      '울료(尉繚)',    '법치 군사론 · 엄정한 군기'),
        ('이위공문대(李衛公問對)', '이정(李靖)',  '당대 전략 문답 · 정·기의 변용'),
    ]
    for i, row in enumerate(rows):
        name, author, role = row[0], row[1], row[2]
        highlight = len(row) > 3 and row[3]
        y = Inches(2.3 + i * 0.6)
        bg = PALE if not highlight else ACCENT
        text_color = INK if not highlight else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(3.7), Inches(0.55), bg)
        add_textbox(slide, Inches(0.7), y, Inches(3.7), Inches(0.55),
                    name, font_size=15, bold=True, color=text_color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.6), y, Inches(2.2), Inches(0.55),
                    author, font_size=14, color=SUB if not highlight else ACCENT, bold=highlight,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.0), y, Inches(5.9), Inches(0.55),
                    role, font_size=14, color=INK if not highlight else ACCENT, bold=highlight,
                    anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
                '「군주가 읽어야 할 두 권」이 곧 육도와 삼략',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅳ. 육도 6권의 흐름 ==============
SEC4 = 'Ⅳ. 육도 6권의 흐름'

@S(SEC4)
def iv_flow(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '육도 6권 — 「文에서 武로」',
              '통치 → 전략 → 장수·참모 → 전술 → 기동')
    boxes = [
        ('文 韜', '문도', '12편', '문치(文治)와 국가 통치의 원리'),
        ('武 韜', '무도',  '5편', '국가 전략과 모략 — 문벌(文伐) 12계'),
        ('龍 韜', '용도', '13편', '참모 조직과 장수 선발 — 왕익(王翼) 72인'),
        ('虎 韜', '호도', '12편', '기본 전술 — 힘과 정공법의 호(虎)'),
        ('豹 韜', '표도',  '8편', '특수 지형 전투 — 민첩의 표(豹)'),
        ('犬 韜', '견도', '10편', '병종별 운용·기동 — 조직의 견(犬)'),
    ]
    for i, (han, kor, num, desc) in enumerate(boxes):
        y = Inches(2.3 + i * 0.78)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.7), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.7), Inches(0.65),
                    han, font_size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, Inches(2.5), y, Inches(1.5), Inches(0.65), PALE)
        add_textbox(slide, Inches(2.5), y, Inches(1.5), Inches(0.65),
                    kor, font_size=15, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.1), y, Inches(1.4), Inches(0.65),
                    num, font_size=14, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.6), y + Inches(0.05), Inches(7.3), Inches(0.6),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_message(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '육도가 일러주는 가장 큰 메시지',
              '「병법은 전술로 시작하지 않는다. 통치에서 시작한다.」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 전반 3권(문·무·용) — 통치 · 전략 · 인사',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('     「전쟁 이전」에 결정되는 것들 — 60편 중 30편이 여기에',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 후반 3권(호·표·견) — 전술 · 지형 · 기동',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     「전쟁 자체」의 기술 — 실전 매뉴얼',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 손자가 「벌모 → 벌교 → 벌병 → 벌성」 순으로 우선순위를 말한다면,',
         {'font_size': 17, 'space_before': 14}),
        ('     육도는 그 순서를 책 자체의 구조로 보여준다 — 「文 → 武」',
         {'font_size': 17, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 통치가 없으면 전략이 없고, 전략이 없으면 전술도 의미 없다',
         {'font_size': 16, 'space_before': 14, 'color': SUB}),
    ])


# ============== Ⅴ. 육도 6권 깊이 읽기 ==============
SEC5 = 'Ⅴ. 육도 6권 깊이 읽기'

def make_chapter_slide(num, total, name_han, name_kor, headline, keys, principle, point):
    @S(SEC5)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC5} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=28, bold=True, color=INK)
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    headline, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_textbox(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.4),
                    '◆ 핵심 편(編)', font_size=13, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0),
                       [(keys, {'font_size': 15, 'color': INK})], line_spacing=1.4)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 권 전체의 가르침', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(principle, {'font_size': 16, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 음미할 지점', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 16, 'color': INK})], line_spacing=1.35)


CHAPTERS = [
    ('文 韜 (문도)', '12편', '제1권 — 문치(文治)와 제왕학의 정수',
     '문사(文師) · 영허(盈虛) · 국무(國務) · 대례(大禮) · 명전(明傳) · 육수(六守) · 수토 · 수국 · 상현(上賢) · 거현(擧賢) · 상벌 · 병도(兵道)',
     '60편 중 첫 12편이 「순수 제왕학」 — 정치·경제·인사가 통째 들어 있다.\n천하공유론 · 육수(六守) · 상현·거현 · 영허(盈虛)의 이치.',
     '「天下非一人之天下」 — 천하는 한 사람의 것이 아니라는 단언이 첫 편 첫 구절. 손자병법에는 없는 영역.'),
    ('武 韜 (무도)',  '5편', '제2권 — 국가 전략과 모략',
     '발계(發啓) · 문계(文啓) · 문벌(文伐) · 순계(順啓) · 삼의(三疑)',
     '무력 이전에 「문(文)으로 치는 법」 12가지 — 외교·간첩·이간·미인계·재물 회유·적신 포섭.\n손자의 「벌모·벌교」를 극단까지 구체화한 버전.',
     '문벌(文伐) 12계 — 현대 하이브리드 전쟁·정보전·기업 경쟁 전략의 원형. 칼을 뽑기 전에 먼저 머리를 쓴다.'),
    ('龍 韜 (용도)', '13편', '제3권 — 장수론과 참모 조직',
     '왕익(王翼) · 논장 · 선장 · 입장 · 장위 · 여군 · 음부 · 음서 · 군세 · 기병 · 오음 · 병징 · 농기',
     '왕익(王翼) — 군주의 참모를 72개 전문 직종으로 정밀 분류.\n복심·모사·천문·지리·병법·통량(군수)·법산(회계)·유사(외교)·술사(기술) 등.',
     '2,500년 전에 이미 「전문가 사단(cross-functional team)」의 설계도. 「혼자 다 하려는 왕이 가장 위험하다」.'),
    ('虎 韜 (호도)', '12편', '제4권 — 정규전의 전술 · 힘과 정공법의 호(虎)',
     '군용(軍用) · 삼진(三陣) · 질전 · 필출(포위 돌파) · 군략 · 임경 · 동해 · 금고 · 약지 · 변동 · 오운산병 · 오운택병',
     '기초 전술의 체계 — 무기 규격·진형·신호·돌파·진영 유지의 매뉴얼.\n「호(虎)」가 상징하는 것은 정면에서의 위력.',
     '전장의 「표준 작전 절차(SOP)」가 정밀하게 짜여 있다. 「용맹」이 아닌 「준비」가 호의 진정한 힘.'),
    ('豹 韜 (표도)',  '8편', '제5권 — 특수 지형 전투 · 민첩성의 표(豹)',
     '임전(숲) · 돌전(돌격) · 산병 · 화전(불) · 허루(빈 진영 활용) · 금림(매복) · 돌출 · 산공(산 공격)',
     '숲·산·물·불 — 다양한 지형과 상황별 실전 매뉴얼.\n예측을 깨는 기습과 변칙의 모음.',
     '정공법(虎)이 만능이 아니다 — 지형이 달라지면 전법도 달라진다. 「변칙의 미학」의 매뉴얼화.'),
    ('犬 韜 (견도)', '10편', '제6권 — 병종별 운용·기동 · 조직성의 견(犬)',
     '분합 · 무봉 · 연거(전차 훈련) · 연승(기병 훈련) · 무차사 · 연전 · 보기(보병·기병) · 전보 · 기전 · 전기',
     '보병·기병·전차 — 병종별 훈련과 편제의 구체적 매뉴얼.\n각 병종의 강점을 어떻게 결합할 것인가의 합동 작전론.',
     '「견(犬)」 = 조직성. 사냥개가 무리 지어 움직일 때 호랑이도 무너뜨린다 — 단일 병종이 아닌 「결합」의 위력.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅵ. 삼략 3편 깊이 읽기 ==============
SEC6 = 'Ⅵ. 삼략 3편 깊이 읽기'

@S(SEC6)
def vi_shang(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC6} (1/3)', n, t)
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                '上 略  (상략)',
                font_size=28, bold=True, color=INK, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                '리더의 덕 · 강유 조화 · 군심 장악',
                font_size=14, color=SUB)
    add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.8), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.8),
                '柔 能 制 剛  弱 能 制 強\n柔 者 德 也  剛 者 賊 也',
                font_size=26, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 「부드러움이 굳셈을 제어하고, 약함이 강함을 제어한다」',
         {'font_size': 17, 'space_before': 6}),
        ('● 「유는 덕이요, 강은 해로움이다」 — 동양 리더십의 핵심 명제',
         {'font_size': 17, 'space_before': 10}),
        ('● 그러나 결정적 명제 — 「能柔能剛, 其國彌光 / 純剛純強, 其國必亡」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     「유할 수도 강할 수도 있는 자」가 진짜 리더 — 마키아벨리의 「여우와 사자」보다 2,000년 앞섬',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 상벌의 공정 · 군심(軍心) 장악 · 「병사에게는 후히, 자기에게는 절약」',
         {'font_size': 16, 'space_before': 12}),
    ])


@S(SEC6)
def vi_zhong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC6} (2/3)', n, t)
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                '中 略  (중략)',
                font_size=28, bold=True, color=INK, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                '권변(權變)과 시대 읽기 — 시대마다 다른 리더십',
                font_size=14, color=SUB)
    add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5),
                '夫 人 衆 一  則 勝   不 一  則 敗',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(3.5), [
        ('● 「사람의 마음이 하나가 되면 이기고, 하나가 되지 못하면 진다」',
         {'font_size': 17, 'space_before': 6}),
        ('● 삼황(三皇)·오제(五帝)·삼왕(三王)·오패(五覇) — 통치 스타일의 시대 변천',
         {'font_size': 17, 'space_before': 10}),
        ('● 「세상이 바뀌면 법도 바뀐다」 — 경직된 원칙주의에 대한 경고',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 전권 위임(權變) — 「전쟁 중에는 임금의 명이라도 장수가 거역할 수 있다」',
         {'font_size': 16, 'space_before': 10}),
        ('● 시대 판독의 능력 — 같은 원칙이라도 어느 시대엔 약이 되고 어느 시대엔 독이 된다',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC6)
def vi_xia(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, f'{SEC6} (3/3)', n, t)
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                '下 略  (하략)',
                font_size=28, bold=True, color=INK, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                '도덕의 본질 · 정당한 전쟁론',
                font_size=14, color=SUB)
    add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5),
                '師 出 以 義',
                font_size=44, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(3.5), [
        ('● 「군대는 의(義)로써 출동한다」 — 의롭지 않은 전쟁은 반드시 패한다',
         {'font_size': 17, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 道·德·仁·義·禮 — 五者一體也 — 다섯은 한 몸이다',
         {'font_size': 17, 'space_before': 12, 'font_name': 'Batang'}),
        ('● 「得賢者昌, 失賢者亡」 — 현인 등용이 흥망을 가른다',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
        ('● 민심과 천도의 연결 — 백성의 뜻이 곧 하늘의 뜻',
         {'font_size': 17, 'space_before': 10}),
        ('● 삼략 = 「군주의 철학서」 — 군사 실무가 거의 없는 도덕·정치 텍스트',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
    ])


# ============== Ⅶ. 제왕학 7대 원칙 ==============
SEC7 = 'Ⅶ. 제왕학 7대 원칙'

@S(SEC7)
def vii_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '육도삼략은 어떻게 「제왕학」인가',
              '병서의 외피를 쓴 통치의 교과서 — 7대 원칙')
    items = [
        ('1', '천하공유(天下公有)',  '천하는 한 사람의 것이 아니다 — 이익을 나누는 자가 얻는다'),
        ('2', '육수(六守)',         '군주가 지켜야 할 인·의·충·신·용·모 + 6가지 검증법'),
        ('3', '상현·거현(上賢·擧賢)', '인재가 전부다 — 「6적·7해」를 가려내는 안목'),
        ('4', '입장(立將) 의례',     '도끼를 건네는 위임의 의식 — 「군문을 넘으면 장군이 제어한다」'),
        ('5', '왕익(王翼) 72인',    '전문 영역별 참모진 — 「혼자 다 하려는 왕이 가장 위험」'),
        ('6', '강유 조화(剛柔調和)', '덕만으로도 힘만으로도 안 된다 — 전환할 수 있는 자가 진짜'),
        ('7', '민본(民本)',         '백성의 마음이 모든 결정의 저울 — 「분육지심(分肉之心)」'),
    ]
    for i, (num, name, desc) in enumerate(items):
        y = Inches(2.3 + i * 0.65)
        add_filled_rect(slide, Inches(0.7), y, Inches(0.7), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(0.7), Inches(0.55),
                    num, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.5), y, Inches(3.5), Inches(0.55), PALE)
        add_textbox(slide, Inches(1.5), y, Inches(3.5), Inches(0.55),
                    name, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.2), y + Inches(0.05), Inches(7.7), Inches(0.5),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC7)
def vii_gongyou(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '제왕학 ① 천하공유(天下公有)',
              '문도·문사편 — 육도 전체의 정치 철학의 원기(元氣)')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.6), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.8),
                '天 下 非 一 人 之 天 下  乃 天 下 之 天 下 也',
                font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(12.0), Inches(0.8),
                '同 天 下 之 利 者 則 得 天 下   擅 天 下 之 利 者 則 失 天 下',
                font_size=18, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(3.0), [
        ('● 천하는 천하 사람들의 것 — 한 사람의 사사로운 소유가 아니다',
         {'font_size': 17, 'space_before': 6}),
        ('● 천하의 이익을 함께 나누는 자가 천하를 얻고, 독점하는 자는 잃는다',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 맹자의 「여민동락(與民同樂)」, 장자의 「천하위공(天下爲公)」과 공명',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 현대 환언 — 「공유 가치 창출(Shared Value)」·「이해관계자 자본주의」의 2,000년 선행 이론',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC7)
def vii_liushou(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '제왕학 ② 육수(六守) — 인재 검증의 6가지 시험',
              '문도·육수편 — 현대 HR의 다면 평가(360°)의 고전 원형')
    rows = [
        ('富 之 以 觀 其 無 犯', '부유하게 해서 범법 여부를 본다'),
        ('貴 之 以 觀 其 無 驕', '귀하게 해서 교만 여부를 본다'),
        ('付 之 以 觀 其 無 背', '일을 맡겨 배반 여부를 본다'),
        ('使 之 以 觀 其 無 隱', '부려서 숨기는 바가 없는지를 본다'),
        ('危 之 以 觀 其 無 恐', '위기에 처하게 해서 두려움 없는지를 본다'),
        ('事 之 以 觀 其 無 竊', '재물 다루게 해서 절도 없는지를 본다'),
    ]
    for i, (han, kor) in enumerate(rows):
        y = Inches(2.3 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.6), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(5.5), Inches(0.6),
                    han, font_size=17, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(6.5), y + Inches(0.05), Inches(6.5), Inches(0.55),
                    kor, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
                '「사람을 쓰기 전에 시험하고, 시험할 때는 여섯 각도로 본다」',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S(SEC7)
def vii_wangyi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '제왕학 ③ 왕익(王翼) — 72인의 참모 조직',
              '용도·왕익편 — 2,500년 전의 「전문가 사단」 설계도')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 군주의 참모를 72개 전문 직종으로 정밀 분류',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 대표 직종', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 복심(腹心) — 책사  /  모사(謀士) — 전략가  /  천문 · 지리 · 병법',
         {'font_size': 14, 'space_before': 6}),
        ('     · 통량(通糧) — 군수  /  분위(奮威) — 사기 진작  /  유사(遊士) — 외교',
         {'font_size': 14, 'space_before': 4}),
        ('     · 법산(法算) — 회계  /  술사(術士) — 기술자  /  방술사 — 의료',
         {'font_size': 14, 'space_before': 4}),
        ('     · 자재사 — 인사  /  권사·유선사 — 정보  /  이목·아자 — 요원  /  우사 — 통역',
         {'font_size': 14, 'space_before': 4}),
        ('● 현대 환언 — 「T자형 조직」·「Cross-functional team」의 동양 원조',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
        ('● 핵심 — 「왕은 혼자 다 하지 않는다. 혼자 다 할 수 있다고 생각하는 왕이 가장 위험하다」',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC7)
def vii_lijiang(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '제왕학 ④ 입장(立將) — 도끼를 건네는 의례',
              '용도·입장편 — 위임은 말이 아닌 의식(ritual)으로 완성된다')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5),
                '從 此 以 往  至 于 軍 門  則 將 軍 制 之',
                font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 「지금부터 군문에 이르기까지는 장군이 제어한다」 — 왕의 선언',
         {'font_size': 17, 'space_before': 6}),
        ('● 왕이 직접 도끼를 들어 장군에게 건넨다 — 전권 위임의 의식화',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 그 경계를 넘은 뒤에는 군주의 명도 장수의 판단보다 우선하지 않는다',
         {'font_size': 17, 'space_before': 10}),
        ('● 손자의 「將在外, 君命有所不受」가 육도에서는 의례화된 제도로 등장',
         {'font_size': 15, 'space_before': 10, 'color': SUB, 'font_name': 'Batang'}),
        ('● 현대 환언 — 「취임식·킥오프·권한 위임 문서」의 의식적 설계',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC7)
def vii_balance(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '제왕학 ⑤ 강유 조화 — 마키아벨리보다 2,000년 앞선 「여우와 사자」',
              '삼략·상략 — 덕만으로도 힘만으로도 안 된다')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(2.0), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.9),
                '能 柔 能 剛  其 國 彌 光',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.1), Inches(12.0), Inches(0.9),
                '純 剛 純 強  其 國 必 亡',
                font_size=24, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.5), Inches(12.0), Inches(3.0), [
        ('● 「유할 수도 강할 수도 있으면 나라가 빛나고, 순전히 강하기만 하면 반드시 망한다」',
         {'font_size': 17, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 한 스타일에 고착된 리더는 반드시 실패 — 부드러움과 단호함을 전환할 수 있는 자가 살아남는다',
         {'font_size': 16, 'space_before': 12}),
        ('● 마키아벨리의 「여우와 사자」(군주론, 1532)의 동양 원형 — 2,000년 앞섬',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 현대 환언 — 「situational leadership · leadership agility」의 고전적 정식',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅷ. 손자병법과의 비교 ==============
SEC8 = 'Ⅷ. 손자병법과의 비교'

@S(SEC8)
def viii_compare(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '손자병법 vs 육도삼략 — 한 폭으로 비교')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.5), INK)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(2.6), Inches(0.5),
                '항목', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(3.3), Inches(2.2), Inches(4.7), Inches(0.5),
                '손자병법', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(8.0), Inches(2.2), Inches(4.7), Inches(0.5),
                '육도삼략', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('시대',     '춘추 말 (BC 5c)',         '전국 말~진한 (BC 3~2c)'),
        ('분량',     '13편 · 6,000자',          '60편+3편 · 수만 자'),
        ('주된 독자', '장수',                    '군주 + 장수'),
        ('중심 관심', '전략·전술·조직',          '통치·정치·전략·전술 총괄'),
        ('전쟁관',   '「부전승 최선」',           '「사출이의(師出以義)」'),
        ('민심',     '도(道)로 간략히 언급',     '전반의 근본 ─ 천하공유'),
        ('장수론',   '5덕(智信仁勇嚴)',          '5덕(勇智仁信忠)+6수+8관찰법'),
        ('임명식',   '간단 언급',               '도끼 건네기 의례 상세'),
        ('참모 조직', '거의 없음',                '왕익 72인의 정밀 편성'),
        ('정신',     '냉철한 현실주의',         '덕과 현실의 결합'),
        ('문체',     '짧고 날카로움 · 압축',     '길고 풍부함 · 서사적'),
    ]
    for i, (k, s, l) in enumerate(rows):
        y = Inches(2.72 + i * 0.4)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.4), bg)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.4),
                    k, font_size=13, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.3), y, Inches(4.7), Inches(0.4),
                    s, font_size=13, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(8.0), y, Inches(4.7), Inches(0.4),
                    l, font_size=13, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC8)
def viii_pair(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '경쟁자가 아닌 「짝」 — 서로를 보완하는 두 책',
              '한나라 건국의 비밀 — 「삼략의 참모 + 손자의 장수」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 손자만 읽으면 — 뛰어난 전술가는 되나, 왜 싸워야 하는지·싸운 뒤 어떻게 통치할지를 놓친다',
         {'font_size': 17, 'space_before': 6}),
        ('● 육도삼략만 읽으면 — 통치 철학은 배우나, 실제 전장의 정밀한 승부 기술이 약해진다',
         {'font_size': 17, 'space_before': 10}),
        ('● 둘 다 읽으면 — 「왕의 눈」과 「장수의 눈」을 동시에 가진다',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 역사적 증언', {'font_size': 17, 'space_before': 14, 'bold': True, 'color': ACCENT}),
        ('     · 장량 ← 황석공의 삼략 → 한고조 유방의 책사',
         {'font_size': 15, 'space_before': 6}),
        ('     · 한신 ← 손자병법의 실전 대가 → 한고조의 장수',
         {'font_size': 15, 'space_before': 4}),
        ('     · 즉, 한나라 건국 = 「삼략의 참모(장량) + 손자의 장수(한신)」의 결합',
         {'font_size': 15, 'space_before': 6, 'color': SUB, 'bold': True}),
        ('● 조선의 이순신 — 손자의 「지피지기」와 육도의 「사출이의·민본」이 『난중일기』에 나란히 흐른다',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅸ. 명구 20선 ==============
SEC9 = 'Ⅸ. 명구 20선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC9)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC9} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=17, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 17, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('육도·문도·문사', '天 下 非 一 人 之 天 下  乃 天 下 之 天 下 也',
     '천하는 한 사람의 천하가 아니라 천하 사람들의 천하다',
     '육도 60편 전체의 첫 명제. 「공유 가치 창출(Shared Value)」의 2,000년 선행 이론.'),
    ('육도·문도·문사', '同 天 下 之 利 者 則 得 天 下  擅 天 下 之 利 者 則 失 天 下',
     '천하의 이익을 함께 나누는 자는 천하를 얻고, 독점하는 자는 잃는다',
     '권력의 가장 냉정한 현실주의 — 「혼자 다 가지면 결국 모두 잃는다」.'),
    ('육도·문도·영허', '盈 則 虛  虛 則 盈',
     '가득 차면 비고, 비면 다시 찬다',
     '주역의 「물극필반」과 공명 — 성공의 정점이 곧 몰락의 시작. 잘 나갈 때 「다음의 허」를 준비하라.'),
    ('육도·문도·육수', '富 之 以 觀 其 無 犯  貴 之 以 觀 其 無 驕',
     '부유하게 해서 범법 여부를 보고, 귀하게 해서 교만 여부를 본다',
     '인재 검증의 6가지 시험 중 첫 두 가지. 「부귀를 주어도 흔들리지 않는 자」가 진짜.'),
    ('육도·문도·거현', '得 賢 者 昌  失 賢 者 亡',
     '현인을 얻는 자는 흥하고 잃는 자는 망한다',
     '제왕학의 가장 짧은 명제. 그러나 더 중요한 것은 6적·7해를 가려내는 안목이다.'),
    ('육도·문도·병도', '全 勝 不 鬪  大 兵 無 創',
     '완전한 승리는 싸우지 않고, 큰 군대는 상처 입지 않는다',
     '손자의 「부전이굴인지병(不戰而屈人之兵)」과 같은 정신 — 최고의 승리는 전쟁 자체가 없는 것.'),
    ('육도·용도·논장', '勇 智 仁 信 忠 — 將 之 五 才',
     '용맹·지혜·인자·신실·충성 — 장수의 다섯 재질',
     '손자가 智信仁勇嚴이라면 육도는 勇智仁信忠 — 「충(忠)」이 「엄(嚴)」을 대체. 군주의 시각에서 본 장수론.'),
    ('육도·용도·입장', '將 受 命  則 忘 其 家  臨 軍 約 束  則 忘 其 親  援 枹 鼓  則 忘 其 身',
     '명을 받으면 집을 잊고, 군에 약속하면 부모를 잊고, 북채를 잡으면 자기 몸을 잊는다',
     '장수의 「세 번의 잊음」. 직책 앞에서 한 단계씩 자기를 비우는 의식.'),
    ('육도·용도·입장', '從 此 以 往  至 于 軍 門  則 將 軍 制 之',
     '지금부터 군문에 이르기까지는 장군이 제어한다',
     '왕이 도끼를 건네며 하는 선언 — 위임은 의식으로 완성된다. 현대 취임식·권한위임의 원형.'),
    ('육도·무도·문벌', '文 伐 十 二',
     '문(文)으로 치는 열두 가지',
     '무력 이전의 외교·간첩·이간·미인계·재물 회유·적신 포섭 — 손자의 「벌모·벌교」를 극한까지 구체화.'),
    ('삼략·상략', '柔 能 制 剛  弱 能 制 強',
     '부드러움이 굳셈을 제어하고, 약함이 강함을 제어한다',
     '동양 리더십의 핵심 명제. 노자 「상선약수(上善若水)」와 짝이 되는 군사적 표현.'),
    ('삼략·상략', '柔 者 德 也  剛 者 賊 也',
     '유는 덕이요, 강은 해로움이다',
     '그러나 이것이 전부가 아니다 — 다음 명제가 결정적. 「유만으로도 안 된다」.'),
    ('삼략·상략', '能 柔 能 剛  其 國 彌 光   純 剛 純 強  其 國 必 亡',
     '유할 수도 강할 수도 있으면 나라가 빛나고, 순전히 강하기만 하면 반드시 망한다',
     '마키아벨리의 「여우와 사자」보다 2,000년 앞선 강유 조화론. 진짜 리더는 「전환할 수 있는 자」.'),
    ('삼략·상략', '將 軍 者  國 之 命 也',
     '장군은 나라의 운명이다',
     '장수 한 사람의 무게 — 한 자리에 사람을 잘못 앉히면 한 나라가 흔들린다.'),
    ('삼략·상략', '良 將 之 養 士  不 易 於 身',
     '좋은 장수가 병사를 기르는 것은 자기 몸을 기르는 것과 같다',
     '리더십의 가장 따뜻한 한 줄 — 부하를 자기 몸으로 본다.'),
    ('삼략·상략', '賞 祿 有 功  通 志 於 衆',
     '공이 있는 자에게 상과 녹을 주고, 뜻을 무리에게 통하게 하라',
     '동기 부여의 두 축 — 「보상의 공정」과 「비전의 공유」.'),
    ('삼략·중략', '夫 人 衆 一  則 勝   不 一  則 敗',
     '사람의 마음이 하나가 되면 이기고, 하나가 되지 못하면 진다',
     '조직의 승패는 결국 「마음의 하나됨」에서 갈린다 — 전략·자원보다 앞서는 변수.'),
    ('삼략·하략', '師 出 以 義',
     '군대는 의(義)로써 출동한다',
     '의롭지 않은 전쟁은 반드시 실패한다 — 명분 없는 공격에 대한 가장 짧은 경고. 「대의」가 모든 경쟁의 토대.'),
    ('삼략·하략', '道 · 德 · 仁 · 義 · 禮  五 者 一 體',
     '도·덕·인·의·예 — 이 다섯은 한 몸이다',
     '동양 도덕철학의 한 줄 요약. 다섯은 따로 떨어진 것이 아닌 같은 뿌리의 다른 가지.'),
    ('육도·문도·문사', '取 天 下 者  若 逐 野 獸  而 天 下 皆 有 分 肉 之 心',
     '천하를 얻는 것은 들짐승을 쫓는 것과 같다 — 천하 사람이 모두 고기를 나눌 마음이 있어야 한다',
     '천하 얻기의 가장 솔직한 비유 — 혼자서는 안 된다. 모두가 「나도 이익을 본다」고 느낄 때만 따라온다.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅹ. 후대 영향 ==============
SEC10 = 'Ⅹ. 후대 영향'

@S(SEC10)
def x_china(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '중국 — 한 왕조를 세운 책에서 무묘(武廟)의 신주로',
              '한대의 표준 병학 → 당대 무묘 → 송대 무경칠서')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 한대 — 장량이 삼략으로 유방을 도와 한 400년 왕조를 세움',
         {'font_size': 17, 'space_before': 4}),
        ('● 당대(唐) — 무묘(武廟)에 강태공을 모시고 추앙 (731년)',
         {'font_size': 17, 'space_before': 10}),
        ('● 송대 — 신종(神宗) 시기 무경칠서 공식 선정 (1043) — 무과(武科) 필수 교재',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 삼국지연의 — 하후무가 「육도삼략을 익혔다」 자랑하다 제갈량에게 참패하는 일화',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 1972년 산둥 은작산(銀雀山) 한묘 — 육도 죽간 출토 → 진본임을 입증',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 청대까지 무관·관료의 필독서 — 「병법서이자 정치서」로 활용',
         {'font_size': 16, 'space_before': 10}),
    ])


@S(SEC10)
def x_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '한국 — 이순신의 애독서, 율곡의 강설',
              '조선 무과의 강경 과목 · 충무공의 책장에 있던 책')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 고려 — 도입 후 왕실·무신의 필독서',
         {'font_size': 18, 'space_before': 4}),
        ('● 조선 — 무경칠서의 하나로 무과 시험의 강경(講經) 과목',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 율곡 이이 — 선조에게 올린 『육도삼략 강설(講說)』이 남아 있음',
         {'font_size': 18, 'space_before': 10}),
        ('● 충무공 이순신 — 손자병법과 함께 평생의 애독서',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     『난중일기』에 손자의 「지피지기」와 육도의 「사출이의·민본」이 나란히 흐른다',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 정약용 — 『목민심서』에 「장수의 다섯 재능」 등 육도 개념 인용',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC10)
def x_japan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '일본 — 무가의 필독서 · 다이묘들의 좌우명',
              '도쿠가와 이에야스가 평생 옆에 두었다고 전하는 책')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 헤이안·가마쿠라 시대 — 무가(武家)의 필독서로 도입',
         {'font_size': 18, 'space_before': 4}),
        ('● 센고쿠(戰國) 시대 — 다이묘들의 필독서',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 다케다 신겐(武田信玄) — 「풍림화산(風林火山)」의 사상적 토대 중 하나',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('     · 우에스기 겐신(上杉謙信) — 「義」의 정치 — 삼략 「사출이의」 정신',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 도쿠가와 이에야스 — 평생 옆에 두었다고 전해지는 책',
         {'font_size': 18, 'space_before': 14, 'bold': True, 'color': ACCENT}),
        ('● 완구 기업 「반다이(BANDAI)」 — 사명이 육도 용도 왕익편 「萬代不易(만대불역)」에서 유래',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
        ('● 베트남 리(李)·쩐(陳) 왕조의 군사 교육 — 동아시아 한자 문화권 전반의 영향',
         {'font_size': 16, 'space_before': 10}),
    ])


# ============== Ⅺ. 오늘 우리에게 ==============
SEC11 = 'Ⅺ. 오늘 우리에게'

@S(SEC11)
def xi_modern(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '경영·리더십 — 육도삼략의 현대적 활용')
    rows = [
        ('인사 관리',  '육수(六守) 6가지 검증법 — 다면 평가(360°)의 고전 원형'),
        ('조직 설계',   '왕익(王翼) 72인 — Cross-functional team의 동양 원조'),
        ('인재 검증',  '6적·7해 — 「진짜 같은 가짜」를 거르는 안목'),
        ('리더십',      '오재(五才) · 강유 조화 · 솔선수범(將不言渴)'),
        ('경쟁 전략', '문벌(文伐) 12계 — 비즈니스 경쟁·정보전 매뉴얼'),
        ('산업 정책', '삼보(三寶: 농·공·상) — 산업 균형 발전론'),
        ('비전 공유', '여중동호(與衆同好) — 미션·비전 정렬'),
        ('위임 의례', '입장(立將) — 취임식·킥오프·권한 위임 문서의 설계'),
    ]
    for i, (cat, desc) in enumerate(rows):
        y = Inches(2.3 + i * 0.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.0), Inches(0.45), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(3.0), Inches(0.45),
                    cat, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.9), y + Inches(0.02), Inches(9.0), Inches(0.42),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC11)
def xi_ten(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '오늘 머리에 새겨야 할 10가지')
    items = [
        '1. 이익을 나누지 않는 자는 결국 잃는다 — 천하공유(天下公有)',
        '2. 사람이 전부다 — 그러나 더 어려운 것은 「가짜를 거르는 안목」',
        '3. 소프트 파워가 최종 승자 — 그러나 「전환할 수 있어야」 한다',
        '4. 혼자 다 하려는 리더가 가장 위험하다 — 전문가에게 권한을',
        '5. 위임은 의식(ritual)으로 완성된다 — 「오늘부터 너의 영역」 선언',
        '6. 명분 없는 공격은 반드시 실패 — 사출이의(師出以義)',
        '7. 정점에서 내려올 준비를 한다 — 영허(盈虛)의 순환',
        '8. 진짜 같은 가짜를 보는 눈 — 6적·7해의 경계',
        '9. 한 가지 스타일에 고착되지 말라 — 강유 조화의 이중 모드',
        '10. 마음이 하나가 되면 이긴다 — 부인중일 즉승(夫人衆一 則勝)',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.3 + i * 0.45)
        add_textbox(slide, Inches(0.7), y, Inches(12.0), Inches(0.4),
                    txt, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC11)
def xi_read(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '육도삼략 독서 순서 제안')
    rows = [
        ('1단계', '입문 — 삼략 3편',
         '짧고 강렬한 3,800자. 「강유 조화」·「사출이의」의 핵심 명제부터.'),
        ('2단계', '제왕학 — 육도 문도(文韜) 12편',
         '정치·민본·인재의 핵심 — 손자에 없는 영역.'),
        ('3단계', '전략 — 육도 무도(武韜) 5편',
         '문벌 12계를 중심으로 — 무력 이전의 모략.'),
        ('4단계', '장수론 — 육도 용도(龍韜) 13편',
         '왕익·논장·입장 편이 핵심 — 인사·조직·위임.'),
        ('5단계', '전술 — 호도·표도·견도',
         '실전 매뉴얼. 현대 독자는 가볍게 지나가도 좋다.'),
    ]
    for i, (step, title, desc) in enumerate(rows):
        y = Inches(2.3 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.4), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.4), Inches(0.8),
                    step, font_size=17, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(2.3), y + Inches(0.05), Inches(10.5), Inches(0.8), [
            (title, {'font_size': 16, 'bold': True, 'color': ACCENT}),
            (desc,  {'font_size': 13, 'color': SUB, 'space_before': 4}),
        ], line_spacing=1.25)


# ============== Ⅻ. 마무리 ==============
SEC12 = 'Ⅻ. 마무리'

@S(SEC12)
def xii_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '육도삼략, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 두 책의 합본 — 육도(강태공·60편) + 삼략(황석공·3편).',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 손자병법이 「전장의 병법」이라면, 육도삼략은 「천하 창건의 병법」.',
         {'font_size': 18, 'space_before': 8}),
        ('● 무경칠서의 두 자리 — 「군주가 읽어야 할 두 권」.',
         {'font_size': 18, 'space_before': 8}),
        ('● 네 사상 — 천하공유 · 득현자창 · 유능제강 · 사출이의.',
         {'font_size': 18, 'space_before': 8}),
        ('● 일곱 제왕학 원칙 — 천하공유·육수·상현·입장·왕익·강유·민본.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「병서의 외피를 쓴 제왕학」 — 1,200년 왕조의 비밀.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC12)
def xii_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.1),
                '天 下 之 天 下',
                font_size=92, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(3.6), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.6),
                '천 하 는  천 하 사 람 들 의  것',
                font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.5),
                '同 利 者 得  擅 利 者 失',
                font_size=24, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
                '이익을 나누는 자가 얻고, 독점하는 자는 잃는다',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '— 육도 문도·문사 (六韜 文韜·文師)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\육도삼략_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
