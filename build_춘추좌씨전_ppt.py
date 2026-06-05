# -*- coding: utf-8 -*-
"""
춘추좌씨전(春秋左氏傳) 발표자료 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '중국 최초의 본격 편년체 역사서 · 춘추 255년 격동의 기록',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.6),
                '春 秋 左 氏 傳',
                font_size=100, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.6),
                '춘 추 좌 씨 전  ·  좌 전',
                font_size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.0), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
                '夫 民   神 之 主 也',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                '무릇 백성은 — 신(神)의 주인이다 — 환공 6년',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '좌구명(左丘明) 전승 · 노 12공 255년 · 春秋三傳 중 사실 중심의 정점',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 좌전이란 무엇인가'),
        ('Ⅱ.', '저자와 시대'),
        ('Ⅲ.', '춘추 삼전(三傳)과의 관계'),
        ('Ⅳ.', '7대 사상 — 좌전의 중심축'),
        ('Ⅴ.', '12공의 시대 한 폭으로'),
    ]
    items_right = [
        ('Ⅵ.', '명장면 12선'),
        ('Ⅶ.', '명구·고사성어 16선'),
        ('Ⅷ.', '문학적 성취 — 서사의 정점'),
        ('Ⅸ.', '한국·동아시아 수용'),
        ('Ⅹ.', '오늘 우리에게'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 좌전',
              '중국 최초의 본격 편년체 역사서')
    rows = [
        ('정식 명칭', '春秋左氏傳 — 줄여서 「좌전(左傳)」'),
        ('저자', '좌구명(左丘明) — 노나라 사관 (전통적 견해)'),
        ('편집 완성', '전국 중기(BC 4세기) 전후로 추정'),
        ('내용 시대', '노 12공 — BC 722~468 약 255년'),
        ('성격', '춘추 경문의 해설서 — 사실 중심 서사'),
        ('실질 분량', '춘추 경문의 10배 이상'),
        ('위상', '편년체 역사서의 시조 · 동양 서사 문학의 정점'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.45), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.45),
                    k, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.0), Inches(0.45),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_position(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '좌전의 위치',
              '춘추 + 좌전 = 동양 역사학의 출발')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 공자가 노나라 역사를 재편집한 경전 — 『춘추(春秋)』', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   BC 722 은공 원년 ~ BC 481 애공 14년 — 242년 기록', {'font_size': 14, 'color': SUB}),
        ('● 춘추의 압축적 경문에 — 사실의 살을 붙인 책 = 좌전', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('● 「**경(經)을 풀이한 전(傳)**」이면서도', {'font_size': 17, 'space_before': 12}),
        ('● 춘추 경문에 없는 사건까지 폭넓게 다룸', {'font_size': 17, 'space_before': 6}),
        ('● 「**중국 최초의 본격 편년체 역사서**」로 평가받음', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('● 사기(기전체)와 함께 — 동양 역사 서술의 양대 모형', {'font_size': 16, 'color': SUB, 'space_before': 6}),
    ])


# ============== Ⅱ. 저자와 시대 ==============
SEC2 = 'Ⅱ. 저자와 시대'

@S(SEC2)
def ii_zuo(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '좌구명(左丘明)',
              '공자가 부끄러워한 바를 같이 부끄러워한 사관')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 전통적 견해 — 노(魯)나라 사관 좌구명', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   공자와 동시대 혹은 조금 뒤의 인물', {'font_size': 14, 'color': SUB}),
        ('● 유일한 직접 기록 — 논어 공야장편', {'font_size': 17, 'space_before': 12}),
        ('   「**左 丘 明 恥 之  丘 亦 恥 之**」 — 「좌구명이 부끄러워한 바를 나(공자)도 부끄러워한다」', {'font_size': 14, 'color': SUB}),
        ('● 현대 학계 — 전국 중기(BC 4세기) 전후에 지금 형태로 편집·증보', {'font_size': 17, 'space_before': 12}),
        ('   하나의 저자가 아니라 — 좌구명 전승의 사관 집단의 누적', {'font_size': 14, 'color': SUB}),
    ])


@S(SEC2)
def ii_period(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '좌전이 다루는 시대',
              'BC 722 ~ BC 468 — 약 255년 춘추시대')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 노 12공 — 은공·환공·장공·민공·희공·문공·선공', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('● ... 성공·양공·소공·정공·애공', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('● BC 722년(은공 원년) ~ BC 468년(애공 27년) 전후', {'font_size': 17, 'space_before': 12}),
        ('● 약 255년간 춘추시대 제후국들의 흥망과 분쟁', {'font_size': 17, 'space_before': 10}),
        ('● 다룬 나라 — 주·제·진·초·진·정·송·위·진·채·조·오·월 등 거의 모든 제후국', {'font_size': 16, 'color': SUB, 'space_before': 10}),
        ('● 노나라를 중심축으로 — 천하 전체를 조망', {'font_size': 16, 'color': SUB, 'space_before': 6}),
    ])


# ============== Ⅲ. 춘추 삼전 ==============
SEC3 = 'Ⅲ. 춘추 삼전'

@S(SEC3)
def iii_three(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '춘추 삼전(三傳)',
              '같은 경전의 세 가지 해석')
    rows = [
        ('좌씨전(左傳)', '사실·사건 중심의 서사(敍事) — 풍부한 역사 기록·인물 묘사'),
        ('공양전(公羊傳)', '의리(義理)·포폄(褒貶) 중심의 문답 — 한대 금문경학의 본류'),
        ('곡량전(穀梁傳)', '예(禮)·명분 중심의 짧은 해설 — 공양전과 비슷하나 간결'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.7 + i * 1.0)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.0), Inches(0.8), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(3.0), Inches(0.8),
                    k, font_size=17, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(3.9), y, Inches(9.1), Inches(0.8),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.8),
                '좌전이 사실 중심의 정점 — 가장 풍부한 서사와 인물 묘사',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅳ. 7대 사상 ==============
SEC4 = 'Ⅳ. 7대 사상'

@S(SEC4)
def iv_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '좌전을 관통하는 7대 사상축',
              '읽을 때 놓치지 말아야 할 키')
    items = [
        '禮의 질서 — 인간·국가 행위의 최고 평가 기준',
        '義·信의 중시 — 외교·군사에서 의와 신을 저버리면 패한다',
        '民本(민본) — 「**夫民, 神之主也**」 백성은 신의 주인',
        '德과 天命 — 천명은 덕 있는 자에게 옮아간다',
        '직필(直筆)의 사관 정신 — 동호직필(董狐直筆)',
        '포폄(褒貶)의 춘추필법 — 한 글자로 옳고 그름',
        '현실주의 치국관 — 형·정·덕의 균형, 부국강병',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.5 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.6), Inches(0.55),
                    f'{i+1}.', font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.55),
                    txt, font_size=16, color=INK)


@S(SEC4)
def iv_min_ben(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '핵심 명제 — 民本',
              '夫民 神之主也 — 백성은 신의 주인')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.4),
                '夫 民   神 之 主 也',
                font_size=32, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.5),
                '무릇 백성은 — 신(神)의 주인이다 (환공 6년 사효의 말)',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.7), Inches(12.0), Inches(2.5), [
        ('● 신을 섬기는 제사도 — 사실은 백성을 위한 것', {'font_size': 16, 'space_before': 6}),
        ('● 제사·천명보다 — 민심이 앞선다', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('● 서경 「**民惟邦本**」, 맹자 「**民貴君輕**」과 함께 — 동양 민본주의의 정점', {'font_size': 16, 'space_before': 6}),
        ('● BC 8세기에 이미 — 「**백성 우선**」의 정치 사상이 확립', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC4)
def iv_zhibi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '동호직필(董狐直筆)',
              '권세를 두려워하지 않는 사관 정신')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 진(晉) 영공이 무도하자 조돈(趙盾)이 망명', {'font_size': 16, 'space_before': 6}),
        ('● 사촌 조천(趙穿)이 영공을 시해', {'font_size': 16, 'space_before': 6}),
        ('● 태사 동호(董狐)가 「**趙盾弑其君**」(조돈이 그 임금을 시해했다) 적음', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 조돈의 항의에 — 「**경계를 넘지 않고 돌아와 역적을 안 쳤으니 그대가 시해**」', {'font_size': 16, 'space_before': 6}),
        ('● 공자의 칭찬 — 「**옛날의 양사(良史)다, 법에 따라 적고 숨기지 않았다**」', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 사관의 직필 — 좌전 전체의 윤리적 척추', {'font_size': 16, 'color': ACCENT, 'space_before': 8}),
    ])


# ============== Ⅴ. 12공 한 폭 ==============
SEC5 = 'Ⅴ. 12공의 시대'

@S(SEC5)
def v_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '노 12공의 시대 흐름',
              '255년의 격동을 한 폭으로')
    rows = [
        ('은공 (722~712)', '춘추 개막 — 정백극단어언·대의멸친'),
        ('환공 (711~694)', '쿠데타 즉위·수갈지전·문강의 비극'),
        ('장공 (693~662)', '제 환공 패업·관중 등용·조귀논전'),
        ('민공 (661~660)', '경보불사 노난미이'),
        ('희공 (659~627)', '제 환공 규구지회·진 문공 망명·성복대전'),
        ('문공 (626~609)', '진 패권 흔들림·초의 도전'),
        ('선공 (608~591)', '초 장왕 패업·문정경중·동호직필'),
        ('성공 (590~573)', '안지전·언릉전 - 진초 대결'),
        ('양공 (572~542)', '예의 시대 마감 가까워짐'),
        ('소공 (541~510)', '계씨의 전횡·공자 출생기'),
        ('정공 (509~495)', '협곡지회·공자 활동기'),
        ('애공 (494~468)', '오월쟁패·좌전 마감'),
    ]
    for i, (k, v) in enumerate(rows[:6]):
        y = Inches(2.5 + i * 0.4)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.32), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.32),
                    k, font_size=12, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.85), y, Inches(4.3), Inches(0.32),
                    v, font_size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    for i, (k, v) in enumerate(rows[6:]):
        y = Inches(2.5 + i * 0.4)
        add_filled_rect(slide, Inches(7.2), y, Inches(2.0), Inches(0.32), PALE)
        add_textbox(slide, Inches(7.2), y, Inches(2.0), Inches(0.32),
                    k, font_size=12, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(9.35), y, Inches(3.7), Inches(0.32),
                    v, font_size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅵ. 명장면 12선 ==============
SEC6 = 'Ⅵ. 명장면 12선'

def make_scene_slide(idx_total, title_name, period, story, lesson):
    @S(SEC6)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, SEC6, n, t)
        add_title(slide, f'{idx_total} — {title_name}', period)
        add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(3.0), [
            ('일화', {'font_size': 14, 'bold': True, 'color': SUB}),
            (story, {'font_size': 15, 'color': INK}),
        ], line_spacing=1.45)
        add_filled_rect(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.2), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.85), Inches(11.6), Inches(1.1),
                    f'중심 사상 — {lesson}',
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


scenes = [
    ('1/12', '정백극단어언(鄭伯克段於鄢)', '은공 원년 — 좌전의 첫 장면',
     '정 장공의 어머니 무강이 동생 공숙단을 편애. 단이 반역하자 장공이 「**多行不義 必自斃**」(많이 쌓아두면 스스로 무너진다)며 기다리다 쳤다. 어머니와 절교했다가 영고숙의 간언으로 황천(黃泉)에서 화해.',
     '혈친과 예의 충돌 — 끝내 효를 회복해야. 「많이 쌓이면 스스로 무너진다」의 자멸 법칙.'),
    ('2/12', '대의멸친(大義滅親)', '은공 4년',
     '위 주우가 환공을 죽이고 찬탈. 석작의 아들 석후가 일당이었다. 석작이 자식 석후를 진(陳)에 유인해 죽임. 「**공의 의가 사의 정을 끊는다**」.',
     '공의(公義)가 사정(私情)을 끊는다 — 신하의 충성이 혈육의 정보다 우선.'),
    ('3/12', '관포지교(管鮑之交)', '장공 8~9년',
     '제 환공이 즉위. 적이었던 관중을 — 포숙의 추천으로 재상 삼음. 「**자기보다 나은 자를 알아보는 안목**」 — 관포지교의 원형.',
     '용인(用人)의 대의 — 원수라도 재능 있으면 쓴다. 패업의 토대는 사람.'),
    ('4/12', '조귀논전(曹劌論戰)', '장공 10년 장작전투',
     '제가 노를 침공. 평민 조귀가 장공에게 「**송사를 정성껏 재판하는 것**(忠之屬也)」이 싸울 만한 근거라 답. 제군이 세 번 북 친 뒤 응전해 대승. 「**一鼓作氣 再而衰 三而竭**」.',
     '민심·송사의 공정이 국가의 진정한 힘. 전쟁 승패는 병력 아닌 사기의 운용.'),
    ('5/12', '여희의 난·중이 망명', '희공 4~5, 23년',
     '진 헌공의 여희가 태자 신생을 모함. 신생이 「**효를 지키려**」 자결. 동생 중이(重耳)는 19년간 8나라를 떠돌며 호언·조최·개자추 등과 고난. 진 목공 도움으로 귀국 → 진 문공.',
     '고난이 패자를 기른다. 효와 형세의 비극적 충돌. 결국 덕 있는 자가 천명을 받는다.'),
    ('6/12', '퇴피삼사·성복대전(退避三舍·城濮之戰)', '희공 28년',
     '중이가 망명 시 초 성왕에게 「**전장에서 만나면 90리(三舍) 물러나겠다**」 약속. 성복에서 정말 90리 물러난 뒤 화공으로 초를 대파. 「**의로 움직이니 천하가 부끄럽게 여기지 않는다**」.',
     '신(信)의 힘 — 약속은 생사의 전장에서도 지킨다. 신의가 전술적 우위를 만든다.'),
    ('7/12', '개자추의 불언록(不言祿)', '희공 24년',
     '진 문공이 논공행상할 때 개자추가 자기 공을 말하지 않고 어머니와 산에 숨음. 문공이 산을 태우자 끝내 불에 타 죽음. 「**천록은 내 공이 아니라 하늘의 뜻**」.',
     '공을 탐하지 않는 겸허. 위정자가 녹을 빠뜨리면 사람이 상한다 — 쌍방의 교훈.'),
    ('8/12', '효산지전(殽山之戰)', '희공 33년',
     '진(秦) 목공이 정나라 정벌 원정군이 효산 협곡에서 진(晉) 양공 복병에 전멸. 돌아온 장수 맹명시를 목공이 「**이는 내 잘못**」 자책하며 다시 기용.',
     '군주의 자책이 신하의 죽음을 막고 국력을 회복. 「**과오를 남 탓으로 돌리지 않는 자가 결국 이긴다**」.'),
    ('9/12', '문정경중(問鼎輕重)', '선공 3년',
     '초 장왕이 주 왕실 국경에서 열병 — 왕손만에게 「**구정(九鼎)의 경중을 묻는다**」. 왕손만의 답 — 「**在德不在鼎**」(정은 덕에 있지 크기에 있지 않다). 장왕이 말없이 물러남.',
     '권력의 정당성은 — 무게가 아니라 덕에 있다. 좌전 최고의 정치 명언.'),
    ('10/12', '동호직필(董狐直筆)', '선공 2년',
     '진 영공이 무도, 조돈이 망명. 사촌 조천이 영공 시해. 태사 동호가 「**趙盾弑其君**」 적음. 조돈 항의에 「**그대가 정경으로 도망쳤으되 경계를 넘지 않고 돌아와 역적을 안 쳤다**」. 공자: 「**옛 양사다**」.',
     '사관의 직필 — 권세 아닌 법과 명분에 따라 기록. 좌전 전체의 윤리적 척추.'),
    ('11/12', '협곡지회(夾谷之會)', '정공 10년',
     '제·노 회맹에서 제가 협박 시도. 공자(孔子)가 노 정공의 곁에서 「**예에 어긋난 음악**」을 지적해 제 경공이 굴복. 잃었던 노나라 땅을 돌려받음.',
     '예와 명분의 힘. 공자가 실제 정치에서 보여준 외교 능력 — 「예로 적을 굴복시킨다」.'),
    ('12/12', '오월쟁패(吳越爭覇)', '애공 — 좌전의 마지막',
     '오 부차 vs 월 구천. 구천의 와신상담 — 쓸개를 핥으며 패전의 치욕을 잊지 않음. 마지막에 월이 오를 멸함. 좌전이 마감되는 시점의 격동.',
     '와신상담(臥薪嘗膽) — 굴욕을 잊지 않는 인내. 천하의 패권은 끝없이 옮겨간다.'),
]

for tag, name, period, st, ls in scenes:
    make_scene_slide(tag, name, period, st, ls)


# ============== Ⅶ. 명구·고사성어 16선 ==============
SEC7 = 'Ⅶ. 명구·고사성어 16선'

def make_quote(idx_total, hanja, korean, comment):
    @S(SEC7)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, SEC7, n, t)
        add_title(slide, f'명구 {idx_total}', '좌전이 만든 한 마디')
        add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.2), PALE)
        add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.2),
                    hanja,
                    font_size=24, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(1.0),
                    korean,
                    font_size=18, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.8), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.8),
                    comment,
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


quotes = [
    ('1/16', '多 行 不 義   必 自 斃',
     '많이 행하는 불의는 — 스스로 무너진다',
     '은공 원년. 정 장공이 동생 단의 반역에 대해 한 말. 「자멸의 법칙」 — 좌전 전체의 한 키.'),
    ('2/16', '大 義 滅 親',
     '큰 의(義)가 — 친(親)을 없앤다',
     '은공 4년. 석작이 자식 석후를 죽이며. 「**공의가 사정을 끊는다**」 — 신하 윤리의 정수.'),
    ('3/16', '夫 民   神 之 主 也',
     '무릇 백성은 — 신의 주인이다',
     '환공 6년. 민본 사상의 가장 깊은 표현 — 제사보다 민심이 앞선다.'),
    ('4/16', '管 鮑 之 交',
     '관중과 포숙의 사귐',
     '장공 시기. 「자기보다 나은 자를 알아보는 안목」 — 인재 등용과 우정의 정수.'),
    ('5/16', '一 鼓 作 氣   再 而 衰   三 而 竭',
     '한 번 북 치면 사기 오르고, 두 번이면 쇠하고, 세 번이면 다한다',
     '장공 10년 조귀논전. 군대 사기 운용의 정수 — 「**기세를 한 번에 몰아라**」.'),
    ('6/16', '退 避 三 舍',
     '90리를 물러난다',
     '희공 28년 성복대전. 진 문공의 약속 — 신(信)의 힘이 전술 우위가 됨.'),
    ('7/16', '結 草 報 恩',
     '풀을 묶어 은혜를 갚는다',
     '선공 15년. 위과(魏顆)가 아버지 첩을 살려 보내자 — 그 첩의 죽은 아버지 혼이 풀을 묶어 적장을 넘어뜨림. 은혜는 죽어서도 갚는다.'),
    ('8/16', '問 鼎 輕 重',
     '구정(九鼎)의 경중을 묻는다',
     '선공 3년. 초 장왕의 야망 — 그러나 왕손만의 「**在德不在鼎**」으로 물러남.'),
    ('9/16', '在 德 不 在 鼎',
     '덕에 있지, 정(鼎)에 있지 않다',
     '선공 3년 왕손만. 권력의 정당성은 — 크기가 아니라 덕에 있다. 좌전 최고의 정치 명언.'),
    ('10/16', '董 狐 直 筆',
     '동호의 곧은 붓',
     '선공 2년. 권세를 두려워하지 않는 사관 정신 — 사관 직필의 영원한 모범.'),
    ('11/16', '輔 車 相 依   脣 亡 齒 寒',
     '수레와 바퀴는 서로 의지하고, 입술이 없으면 이가 시리다',
     '희공 5년. 우(虞)·괵(虢) 두 나라의 운명 — 가까운 자가 망하면 자기도 망한다.'),
    ('12/16', '玩 物 喪 志',
     '사물에 빠지면 뜻을 잃는다',
     '여러 편에 산재. 한 가지에 너무 빠지면 큰 뜻을 잃는다는 경계.'),
    ('13/16', '臥 薪 嘗 膽',
     '섶에 누워 쓸개를 핥는다',
     '월 구천의 인내 — 패전의 치욕을 잊지 않는 자세. 마지막에 오를 멸한 동력.'),
    ('14/16', '城 下 之 盟',
     '성 아래에서의 맹약',
     '굴욕적 항복의 맹약 — 적이 성 아래까지 와서 강요한 조약. 굴욕 외교의 대명사.'),
    ('15/16', '東 道 主',
     '동쪽 길의 주인',
     '희공 30년. 정나라가 진(秦) 사신의 주인 노릇 — 「**손님을 영접하는 자**」의 어원.'),
    ('16/16', '篳 路 藍 縷',
     '거친 수레와 누더기 옷',
     '선공 12년. 초나라 조상의 개척 정신 — 처음의 어려움을 견딘 자가 큰 일을 이룬다.'),
]

for tag, hj, kr, cm in quotes:
    make_quote(tag, hj, kr, cm)


# ============== Ⅷ. 문학적 성취 ==============
SEC8 = 'Ⅷ. 문학적 성취'

@S(SEC8)
def viii_narrative(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '서사의 정점',
              '동양 산문 문학의 모범')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 인물 묘사의 생생함 — 정 장공·관중·중이·동호', {'font_size': 17, 'space_before': 6}),
        ('● 대화의 정교함 — 짧은 한 마디로 인격을 드러냄', {'font_size': 17, 'space_before': 10}),
        ('● 전쟁 묘사의 박진감 — 성복대전·효산지전·언릉지전', {'font_size': 17, 'space_before': 10}),
        ('● 외교 변설의 예리함 — 촉지무·자산·안영', {'font_size': 17, 'space_before': 10}),
        ('● 사기·자치통감의 서사 — 모두 좌전을 모범으로 함', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
    ])


@S(SEC8)
def viii_diplomacy(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '외교 변설의 모범',
              '좌전이 전국책의 토대')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 촉지무(燭之武)가 진(秦)을 물러가게 함 — 희공 30년', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   「**亡鄭以陪鄰 鄰之厚 君之薄也**」 — 정을 멸하면 진(晉)만 강해진다', {'font_size': 14, 'color': SUB}),
        ('● 자산(子産)의 정나라 외교 — 양공·소공', {'font_size': 17, 'space_before': 10}),
        ('   강대국 사이에서 약소국이 살아남는 외교의 정수', {'font_size': 14, 'color': SUB}),
        ('● 안영(晏嬰)의 제나라 외교 — 양공·소공', {'font_size': 17, 'space_before': 10}),
        ('   「**오이 종자가 다르다(橘化爲枳)**」 — 환경의 차이를 말함', {'font_size': 14, 'color': SUB}),
        ('● 「**한 마디 말이 천하를 움직인다**」 — 좌전의 외교 미학', {'font_size': 16, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅸ. 한국·동아시아 수용 ==============
SEC9 = 'Ⅸ. 한국·동아시아 수용'

@S(SEC9)
def ix_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '한국에서의 위상',
              '삼경(三經) 중 춘추의 정수로')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 고려·조선 사대부의 필독서 — 사서삼경의 「**춘추**」 중 핵심', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('● 좌전을 통해 — 춘추 경문의 의미를 읽음', {'font_size': 17, 'space_before': 10}),
        ('● 조선 학자들의 좌전 주해 다수 — 김장생·송시열 등', {'font_size': 17, 'space_before': 10}),
        ('● 동호직필 — 한국 사관 정신의 가장 깊은 토대', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   조선 사관의 「**史官筆**」 정신이 — 좌전 동호의 직접 후예', {'font_size': 14, 'color': SUB}),
        ('● 김부식 『**삼국사기**』·정인지 『**고려사**』 — 좌전 서사의 영향', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC9)
def ix_influence(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '후대 영향',
              '동양 역사 서술의 시조')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 사기(司馬遷) — 좌전의 서사 정신을 기전체로', {'font_size': 17, 'space_before': 6}),
        ('● 자치통감(司馬光) — 좌전과 같은 편년체로', {'font_size': 17, 'space_before': 10}),
        ('● 전국책 — 좌전의 외교 변설의 직접 후예', {'font_size': 17, 'space_before': 10}),
        ('● 일본 — 헤이안 시대부터 사대부 필독서', {'font_size': 17, 'space_before': 10}),
        ('● 좌전 없이 — 동양 역사 서술은 성립하지 않는다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
    ])


# ============== Ⅹ. 오늘 우리에게 ==============
SEC10 = 'Ⅹ. 오늘 우리에게'

@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '좌전이 일러주는 7가지',
              '한 폭으로 정리')
    items = [
        '백성은 신의 주인 — 民本의 정수',
        '천명은 덕에 옮아간다 — 在德不在鼎',
        '신(信)을 지키는 자가 결국 이긴다 — 退避三舍',
        '많이 쌓이면 스스로 무너진다 — 多行不義 必自斃',
        '한 번에 기세를 몰아라 — 一鼓作氣',
        '사관의 직필 — 권세에 굽히지 않는 기록',
        '와신상담의 인내 — 굴욕을 잊지 않는 자세',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.5 + i * 0.6)
        add_textbox(slide, Inches(0.9), y, Inches(0.6), Inches(0.5),
                    f'{i+1}.', font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.5),
                    txt, font_size=16, color=INK)


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                '춘추좌씨전 — 춘추 255년 격동을 살아 있는 서사로',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '在 德 不 在 鼎',
                font_size=120, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
                '덕에 있지 정(鼎)에 있지 않다 — 권력 정당성의 정수',
                font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.8), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
                '동호의 직필 정신이 — 동양 역사학 2,500년을 만들었다',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)


# ============== 빌드 ==============
total = len(SLIDES)
for i, (fn, sec) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    fn(slide, i, total)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\춘추좌씨전.pptx'
prs.save(out_path)
print(f'생성 완료: {out_path}  슬라이드 수: {total}')
