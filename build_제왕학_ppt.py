# -*- coding: utf-8 -*-
"""
제왕학(帝王學) 발표자료 — 망라적 85장 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 1. 표지 ==============
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8),
                '帝 王 學', font_size=110, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                'The Way of the Emperor · 제왕학',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.5),
                '동양 2,000년 흥망사가 가르치는 리더의 학문',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
                '니와 슌페이 『중국고사에서 배우는 제왕학』(2001) — 한 권으로 압축된 동양 리더십',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                '"리더는 자기 능력으로 일하는 자가 아니라, 사람의 능력을 끌어내는 자다"',
                font_size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                '— 한 고조 유방, 천하 통일 후의 자기 평가',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 2. 목차 ==============
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.8), Inches(0.7),
                '목 차', font_size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.15), Inches(12.8))
    items = [
        [('Ⅰ', '개요 — 제왕학이란 무엇인가'),
         ('Ⅱ', '2,000년 발전사'),
         ('Ⅲ', '핵심 텍스트 7종'),
         ('Ⅳ', '동서 비교 — 마키아벨리 vs 동양'),
         ('Ⅴ', '결정적 인물 8인'),
         ('Ⅵ', '7대 핵심 원리')],
        [('Ⅶ', '8대 리더 자질'),
         ('Ⅷ', '명구·명장면 10선'),
         ('Ⅸ', '한국의 제왕학 — 세종·정조'),
         ('Ⅹ', '현대적 의의'),
         ('Ⅺ', '마무리')],
    ]
    for col, group in enumerate(items):
        x = 0.7 + col * 6.4
        top = 1.6
        for num, title in group:
            add_textbox(slide, Inches(x), Inches(top), Inches(1.0), Inches(0.5),
                        num, font_size=17, bold=True, color=ACCENT)
            add_textbox(slide, Inches(x + 1.0), Inches(top), Inches(5.3), Inches(0.5),
                        title, font_size=17, color=INK)
            top += 0.7


# ============== Ⅰ. 개요 ==============
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '제왕학(帝王學)이란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                '"한 나라의 군주가 갖추어야 할 학문"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.5),
                '현대적 의미 — "한 조직의 최고 책임자·리더가 지녀야 할 자질과 능력의 총체"',
                font_size=15, color=ACCENT, align=PP_ALIGN.CENTER)
    items = [
        ('CEO·임원',       '기업 경영의 도(道)'),
        ('정치 지도자',     '국가 통치의 도(道)'),
        ('관리자·팀장',     '조직 운영의 도(道)'),
        ('부모·교사',       '가정·학교의 도(道)'),
        ('자기 인생',       '자기 주재(主宰)의 도(道)'),
    ]
    top = 4.0
    for tag, desc in items:
        add_filled_rect(slide, Inches(0.7), Inches(top), Inches(3.0), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.7), Inches(top + 0.13), Inches(3.0), Inches(0.4),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.8), Inches(top), Inches(8.9), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.0), Inches(top + 0.13), Inches(8.7), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.62


@S('Ⅰ. 개요')
def s_book(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '『중국고사에서 배우는 제왕학』 — 책 정보')
    rows = [
        ('서명',    '중국고사에서 배우는 제왕학',          '"옛 고사로 풀어낸 제왕의 학문"'),
        ('저자',    '니와 슌페이(니와슈운페이)',            '일본의 동양고전 연구자'),
        ('옮긴이',  '이규은',                                '한국어 번역'),
        ('출판',    '도서출판 삶과꿈',                       '2001년 3월 (현재 절판)'),
        ('성격',    '리더십 압축서',                         '2,000년 흥망사가 응축'),
        ('출처',    '사기·정관정요·자치통감·한비자·육도삼략', '중국 고전 종합'),
        ('한국 의미','외환위기 이후 리더십 논쟁기 출간',     '동양 리더십의 가교'),
    ]
    top = 2.15
    for i, (tag, val, note) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.65), PALE)
        add_filled_rect(slide, Inches(2.85), Inches(top), Inches(10.0), Inches(0.65), bg)
        add_textbox(slide, Inches(0.55), Inches(top + 0.18), Inches(2.2), Inches(0.4),
                    tag, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.95), Inches(top + 0.05), Inches(4.5), Inches(0.5),
                    val, font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.5), Inches(top + 0.08), Inches(5.3), Inches(0.5),
                    note, font_size=13, color=SUB)
        top += 0.7


@S('Ⅰ. 개요')
def s_niwa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '저자 니와 슌페이 — 일본 동양고전 리더십의 흐름')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '"기업 경영의 답을 동양 고전에서 찾자" — 1970년대 이후 일본 경영의 흐름',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rows = [
        ('야마모토 시치헤이', '제왕학 — 정관정요에서 배우는 리더의 자격', '정관정요 중심'),
        ('모리야 히로시',    '중국 고전 인간학·군주론',                   '한비자·손자·관자 종합'),
        ('다케우치 미노루',  '중국 고전 통섭',                            '학술적 깊이'),
        ('시바 료타로',      '역사 소설로 재구성',                        '대중적 보급'),
        ('니와 슌페이',      '중국고사에서 배우는 제왕학 (본서)',         '중국 고사 종합 정리'),
    ]
    top = 3.3
    for i, (name, book, char) in enumerate(rows):
        is_this = '니와' in name
        color = ACCENT if is_this else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_this else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.8), Inches(0.65), color)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(2.8), Inches(0.4),
                    name, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.5), Inches(top), Inches(5.8), Inches(0.65), bg)
        add_textbox(slide, Inches(3.65), Inches(top + 0.13), Inches(5.6), Inches(0.4),
                    book, font_size=12, bold=True, color=INK)
        add_filled_rect(slide, Inches(9.4), Inches(top), Inches(3.4), Inches(0.65), bg)
        add_textbox(slide, Inches(9.55), Inches(top + 0.13), Inches(3.2), Inches(0.4),
                    char, font_size=12, color=SUB)
        top += 0.72
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '공통점 — 개별 인물의 결정적 일화에서 보편 원리 도출, 현대 경영에 직접 적용',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_korea_2001(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '2001년 한국 — 책이 번역된 시점의 의미')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('1997년 외환위기 직후 — 한국 사회의 충격',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"우리에게 필요한 리더는 어떤 사람인가?"',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('— 사회 전반에 절박하게 떠오른 질문',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('당시 한국 사회의 화두',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 김대중-노무현-이명박 시기의 리더십 논쟁',
         {'font_size': 14, 'space_before': 6}),
        ('• 기업의 "한국형 경영" 본질에 대한 질문',
         {'font_size': 14, 'space_before': 4}),
        ('• 일본의 리더십 책이 대거 번역·소개된 시기',
         {'font_size': 14, 'space_before': 4}),
        ('• 한국인이 자기 동양 전통에서 답을 찾는 가교',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅰ. 개요')
def s_why_now(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '왜 지금 제왕학을 읽는가')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('얇지만 가볍지 않은 책',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('짧은 페이지 안에 2,000년의 흥망사가 응축',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('그 응축이 묻는 것은 결국 한 가지 질문',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"사람을 이끄는 자는 어떤 사람이어야 하는가?"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('이 질문은 옛 황제만의 것이 아니다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('CEO·팀장·교사·부모·자기 자신 — 우리 모두 어떤 의미에서 리더',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('그러므로 제왕학은 황제만의 학문이 아니라 모두의 학문',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


# ============== Ⅱ. 2,000년 역사 ==============
@S('Ⅱ. 역사')
def s_history_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 역사', page, total)
    add_title(slide, '춘추전국시대 (BC 770~221) — 제왕학의 형성기',
              '"제왕학의 모든 원형이 만들어진 시대"')
    items = [
        ('공자(孔子)',     'BC 551~479',   '인(仁)·예(禮)의 덕치(德治)'),
        ('맹자(孟子)',     'BC 372~289',   '왕도정치·민본주의'),
        ('순자(荀子)',     'BC 298~238',   '예(禮)와 법(法)의 결합'),
        ('한비자(韓非子)', 'BC 280~233',   '법·술·세 — 법가 제왕학의 결정'),
        ('관자(管子)',     'BC 723~645',   '부국강병의 실용주의'),
        ('강태공·황석공',  'BC 11세기 등',  '육도삼략 — 전승 병법·통치학'),
    ]
    top = 2.3
    for name, era, idea in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(3.0), Inches(0.4),
                    name, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(2.5), Inches(0.65), PALE)
        add_textbox(slide, Inches(3.7), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    era, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.3), Inches(top), Inches(6.5), Inches(0.65),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.5), Inches(top + 0.13), Inches(6.2), Inches(0.4),
                    idea, font_size=13, color=INK)
        top += 0.72


@S('Ⅱ. 역사')
def s_history_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 역사', page, total)
    add_title(slide, '진한(秦漢) 시기 (BC 221~AD 220) — 제왕학의 통합',
              '"유가 + 법가의 결합 + 역사의 거울"')
    items = [
        ('진(秦)의 통일',     'BC 221',        '한비자 법가의 전면 적용 → 14년 만에 멸망'),
        ('"법치만으로는 안 된다"', '교훈',         '진의 실패가 후대 제왕학의 출발점'),
        ('한 무제의 결단',     'BC 141~87',     '유가(儒家)를 국가 이념으로 채택'),
        ('사마천 『사기』',     'BC 91 완성',    '인물 중심의 역사 서술 — 후대 모든 제왕학의 원형'),
        ('유향의 정치 일화',   'BC 1세기',      '『전국책』·『설원』·『신서』 — 정치 사례집'),
    ]
    top = 2.4
    for name, era, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(3.5), Inches(0.4),
                    name, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(3.5), Inches(0.3),
                    era, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.2), Inches(top), Inches(8.6), Inches(0.85), PALE)
        add_textbox(slide, Inches(4.4), Inches(top + 0.27), Inches(8.3), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.95
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '이 시기에 "외유내법(外儒內法)" — 겉은 유교, 속은 법가 — 패러다임 형성',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 역사')
def s_history_3(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 역사', page, total)
    add_title(slide, '당대(唐代, 618~907) — 제왕학의 황금기',
              '"정관(貞觀)의 치 — 중국 역사상 가장 모범적 통치"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('당 태종 이세민 (599~649)',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('정관(貞觀)의 치(治) — 627~649 · 23년의 황금기',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('위징·방현령·두여회 등 명신의 활약',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('오긍(吳兢) 『정관정요(貞觀政要)』',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('동양 제왕학의 정점이자 표준 텍스트',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('이후 모든 동양 제왕학은',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"정관의 치를 어떻게 재현할 것인가"의 변주',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅱ. 역사')
def s_history_4(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 역사', page, total)
    add_title(slide, '송대(宋代, 960~1279) — 제왕학의 학문화',
              '"단순 사례집에서 체계적 학문으로"')
    items = [
        ('사마광(司馬光)',  '1019~1086',
         '『자치통감(資治通鑑)』 — 1,084년 편찬 · 16기 294권',
         '1,362년의 통사 — "다스림에 도움이 되는 거울"'),
        ('주희(朱熹)',      '1130~1200',
         '『자치통감강목(綱目)』',
         '자치통감의 도덕적 재해석 — 성리학적 시각'),
        ('진덕수(眞德秀)',  '1178~1235',
         '『대학연의(大學衍義)』 — 43권',
         '『대학』 + 『자치통감』 결합 — 황제 교과서'),
    ]
    top = 2.4
    for name, era, work, char in items:
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(1.4), ACCENT)
        add_textbox(slide, Inches(0.5), Inches(top + 0.32), Inches(2.5), Inches(0.5),
                    name, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.5), Inches(top + 0.8), Inches(2.5), Inches(0.4),
                    era, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.1), Inches(top), Inches(9.7), Inches(1.4), PALE)
        add_textbox(slide, Inches(3.3), Inches(top + 0.2), Inches(9.4), Inches(0.5),
                    work, font_size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(3.3), Inches(top + 0.75), Inches(9.4), Inches(0.5),
                    char, font_size=13, color=INK)
        top += 1.55


@S('Ⅱ. 역사')
def s_history_5(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 역사', page, total)
    add_title(slide, '명청(明淸) 시기 (1368~1912) — 제왕학의 표준화',
              '"국가 표준 교육 체계로 정착"')
    items = [
        ('구준(丘濬)',         '1421~1495',  '『대학연의보』 160권',     '진덕수 책의 확장 — 명대 표준'),
        ('장거정(張居正)',     '1525~1582',  '명대 재상의 황제 교육',    '경연 제도의 정립'),
        ('강희제(康熙帝)',     '1654~1722',  '청대 황금기의 명군',       '정관정요·자치통감 평생 학습'),
        ('옹정제(雍正帝)',     '1678~1735',  '강희의 학습 전통 계승',     '제왕학 친저 다수'),
        ('건륭제(乾隆帝)',     '1711~1799',  '청 최성기',                  '60년 재위 · 어전 강의 일상화'),
    ]
    top = 2.4
    for name, era, work, char in items:
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(0.75), ACCENT)
        add_textbox(slide, Inches(0.5), Inches(top + 0.07), Inches(2.5), Inches(0.4),
                    name, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.5), Inches(top + 0.42), Inches(2.5), Inches(0.3),
                    era, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.1), Inches(top), Inches(4.5), Inches(0.75), PALE)
        add_textbox(slide, Inches(3.25), Inches(top + 0.18), Inches(4.3), Inches(0.4),
                    work, font_size=13, bold=True, color=ACCENT)
        add_filled_rect(slide, Inches(7.7), Inches(top), Inches(5.1), Inches(0.75),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.85), Inches(top + 0.18), Inches(4.9), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.85


# ============== Ⅲ. 핵심 텍스트 ==============
@S('Ⅲ. 핵심 텍스트')
def s_texts_overview(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 핵심 텍스트', page, total)
    add_title(slide, '핵심 텍스트 7종 — "이 책들에서 모든 것이 나왔다"')
    items = [
        ('정관정요',  '오긍 · 8세기 초',     '동양 제왕학의 절대 표준 — 40편'),
        ('자치통감',  '사마광 · 1084',       '1,362년의 통사 — 16기 294권'),
        ('대학연의',  '진덕수 · 1264',       '황제 교과서 — 43권 · 세종 100독'),
        ('한비자',    '한비 · BC 3세기',     '법가 제왕학의 결정 — 55편'),
        ('육도삼략',  '강태공·황석공 전승',  '병법 제왕학 — 무경칠서'),
        ('사기',      '사마천 · BC 91',      '인물 중심 역사 — 130권 · 모든 일화의 원천'),
        ('기타',      '논어·맹자·전국책 등',  '동양 정치사상의 보조 자료'),
    ]
    top = 2.0
    for tag, era, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    tag, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(3.0), Inches(0.55), PALE)
        add_textbox(slide, Inches(3.35), Inches(top + 0.13), Inches(2.8), Inches(0.4),
                    era, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.3), Inches(top), Inches(6.5), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.5), Inches(top + 0.13), Inches(6.2), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.62


@S('Ⅲ. 핵심 텍스트')
def s_jeonggwan(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 핵심 텍스트', page, total)
    add_title(slide, '정관정요(貞觀政要) — 동양 제왕학의 절대 표준')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '오긍(吳兢, 670~749) · 8세기 초 · 10권 40편',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.7), [
        ('당 태종과 신하들의 정치 문답을 주제별로 분류',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('군도·정체·임현·납간·군신감계·택관·성신·검약·인측·신소호·신언어',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.4), [
        ('"1,300년 동안 동아시아 모든 군주의 필수 학습서"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('한국·일본·베트남·청대 강희제·건륭제 모두 평생 학습',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅲ. 핵심 텍스트')
def s_jachi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 핵심 텍스트', page, total)
    add_title(slide, '자치통감(資治通鑑) — 1,362년의 거울')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '사마광(司馬光, 1019~1086) · 1084 완성 · 16기 294권 · 편년체',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.3), [
        ('주 위열왕(BC 403) ~ 후주 세종(959) — 1,362년의 통사',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('편년체(編年體) — 시간 순서로 서술',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"자치통감" — 책 이름의 의미',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('資(자) — 도움이 되다 · 治(치) — 다스림',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('通(통) — 역대를 관통 · 鑑(감) — 거울',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('"다스림에 도움이 되고 역대에 통하여 거울이 되는 책" — 송 신종이 명명',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅲ. 핵심 텍스트')
def s_daehakyeon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 핵심 텍스트', page, total)
    add_title(slide, '대학연의(大學衍義) — 황제의 표준 교과서')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '진덕수(眞德秀, 1178~1235) · 1264 경연 진강 · 43권',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.6), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.4), [
        ('"『대학』의 8조목을 자치통감 등의 사례로 풀어낸 책"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('격물·치지·성의·정심·수신·제가·치국·평천하 + 역사 사례',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.9),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.7), [
        ('조선에서의 위상 — "임금이 반드시 공부해야 하는 책"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('세종이 100번 이상 완독 · 숙종·정조도 자주 인용',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('정관정요 = 사례집 / 대학연의 = 체계적 이론서 (양대 기둥)',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅲ. 핵심 텍스트')
def s_hanbi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 핵심 텍스트', page, total)
    add_title(slide, '한비자(韓非子) — 법가 제왕학의 결정')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '한비(韓非, BC 280~233) · 한(韓)나라 공자 · 55편',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('法 (법)', '명문화된 객관적 규칙',     '공개·평등·일관성'),
        ('術 (술)', '군주가 신하를 다루는 통치술', '은밀한 관리 기법'),
        ('勢 (세)', '군주의 위엄과 권세',           '지위에서 나오는 힘'),
    ]
    top = 3.4
    for han, role, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.0), Inches(0.5),
                    han, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.7), Inches(top), Inches(4.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(2.85), Inches(top + 0.22), Inches(4.3), Inches(0.5),
                    role, font_size=15, bold=True, color=ACCENT)
        add_filled_rect(slide, Inches(7.3), Inches(top), Inches(5.5), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.5), Inches(top + 0.22), Inches(5.2), Inches(0.5),
                    desc, font_size=13, color=INK)
        top += 0.97
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
                '"진시황이 한비자의 책을 읽고 ‘이 사람을 만나 죽어도 여한이 없다’"',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '진의 통일 사상적 토대 — 법가 제왕학의 결정',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅲ. 핵심 텍스트')
def s_yukdo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 핵심 텍스트', page, total)
    add_title(slide, '육도삼략(六韜三略) — 병법 제왕학')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '『육도(六韜)』', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('저자(전승) — 강태공(姜太公)', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 6}),
        ('6편 60장', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 8}),
        ('• 문도(文韜) — 정치', {'font_size': 13, 'space_before': 4}),
        ('• 무도(武韜) — 군사 일반', {'font_size': 13}),
        ('• 용도(龍韜) — 장수의 도', {'font_size': 13}),
        ('• 호도(虎韜) — 전투의 道', {'font_size': 13}),
        ('• 표도(豹韜) — 야전', {'font_size': 13}),
        ('• 견도(犬韜) — 보병·기병', {'font_size': 13}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '『삼략(三略)』', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('저자(전승) — 황석공(黃石公) → 장량', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 6}),
        ('3편', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 8}),
        ('• 상략(上略) — 큰 그림', {'font_size': 13, 'space_before': 4}),
        ('• 중략(中略) — 군주의 도', {'font_size': 13}),
        ('• 하략(下略) — 책략과 응용', {'font_size': 13}),
        ('', {'font_size': 6}),
        ('전쟁 + 군주의 통치 원리 종합', {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ('무경칠서(武經七書) 중 둘', {'font_size': 13, 'color': SUB, 'bold': True}),
    ], line_spacing=1.3)


@S('Ⅲ. 핵심 텍스트')
def s_sagi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 핵심 텍스트', page, total)
    add_title(slide, '사기(史記) — 인물로 읽는 역사 · 제왕학의 보고(寶庫)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '사마천(司馬遷, BC 145?~86?) · 130권 · 기전체(紀傳體)',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rows = [
        ('본기(本紀)',  '12편',  '제왕의 일대기 — "역사를 움직인 자"'),
        ('세가(世家)',  '30편',  '제후의 일대기'),
        ('열전(列傳)',  '70편',  '신하·장군·학자의 일대기 — 인물 중심'),
        ('표(表)',      '10편',  '연표·계보'),
        ('서(書)',      '8편',   '제도사 — 천문·역법·예악'),
    ]
    top = 3.4
    for tag, scope, desc in rows:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    tag, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(1.8), Inches(0.55), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.13), Inches(1.8), Inches(0.4),
                    scope, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.1), Inches(top), Inches(7.7), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.3), Inches(top + 0.13), Inches(7.4), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.62
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"역사는 인물의 이야기다" — 동양 사관의 원형, 제왕학의 거의 모든 일화의 출처',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅳ. 동서 비교 ==============
@S('Ⅳ. 동서 비교')
def s_compare_west(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 동서 비교', page, total)
    add_title(slide, '동양 제왕학 vs 마키아벨리 『군주론』 — 본질적 차이')
    rows = [
        ('초점',      '권력 유지 기술',         '군주 인격의 완성'),
        ('인간관',    '본래 악함 → 두려움 필요', '본래 선함 → 덕(德)으로 감화'),
        ('핵심',      '"사랑보다 두려움"',      '"두려움보다 신뢰가 강하다"'),
        ('도덕',      '결과가 수단을 정당화',    '수단도 도덕적이어야 함'),
        ('시간 척도', '단기 권력 유지',          '장기 안정과 후세 평가'),
        ('모범 인물', '체사레 보르자',          '요·순·당 태종'),
    ]
    top = 2.0
    add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(0.55), SUB)
    add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(2.5), Inches(0.4),
                '항목', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(3.05), Inches(top), Inches(4.8), Inches(0.55), SUB)
    add_textbox(slide, Inches(3.05), Inches(top + 0.1), Inches(4.8), Inches(0.4),
                '마키아벨리 군주론', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(7.9), Inches(top), Inches(4.9), Inches(0.55), ACCENT)
    add_textbox(slide, Inches(7.9), Inches(top + 0.1), Inches(4.9), Inches(0.4),
                '동양 제왕학', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    row_h = 0.8
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(y), Inches(2.5), Inches(row_h), PALE)
        add_textbox(slide, Inches(0.55), Inches(y + 0.22), Inches(2.4), Inches(0.4),
                    row[0], font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.05), Inches(y), Inches(4.8), Inches(row_h), bg)
        add_textbox(slide, Inches(3.2), Inches(y + 0.22), Inches(4.5), Inches(0.4),
                    row[1], font_size=13, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.9), Inches(y), Inches(4.9), Inches(row_h), bg)
        add_textbox(slide, Inches(8.05), Inches(y + 0.22), Inches(4.6), Inches(0.4),
                    row[2], font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 동서 비교')
def s_compare_essence(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 동서 비교', page, total)
    add_title(slide, '핵심 차이 — "권력 기술" vs "좋은 사람"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('동양 제왕학은',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('"좋은 사람이 좋은 리더가 된다"는 확신 위에 서 있다',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('이것이 마키아벨리와의 결정적 차이',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"덕(德)·인(仁)·예(禮)·법(法)을 종합한 통치의 도(道)"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 덕(德) — 군주의 인격',
         {'font_size': 14, 'space_before': 6}),
        ('• 인(仁) — 백성에 대한 사랑',
         {'font_size': 14, 'space_before': 4}),
        ('• 예(禮) — 관계의 질서',
         {'font_size': 14, 'space_before': 4}),
        ('• 법(法) — 명확한 규칙',
         {'font_size': 14, 'space_before': 4}),
        ('이 넷을 종합한 "도(道)"가 동양 제왕학의 본질',
         {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅴ. 결정적 인물 8인 ==============
# 한 고조 유방
@S('Ⅴ. 인물 — 유방')
def s_yubang_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 유방', page, total)
    add_title(slide, '한 고조 유방(漢高祖 劉邦, BC 256~195)',
              '— 최고 지도자의 모범 · 평민에서 황제로')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '劉\n邦', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('출신', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  강소성 패현(沛縣)의 농민 출신 · 작은 정장(亭長)',
         {'font_size': 14}),
        ('  술과 여자를 좋아하고 일은 게을렀던 보통 사람',
         {'font_size': 13, 'color': SUB}),
        ('거병', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  BC 209년 진의 폭정에 봉기 · 60세에 황제',
         {'font_size': 14}),
        ('통일', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  항우(項羽)와 4년 초한전쟁 끝에 BC 202년 천하 통일',
         {'font_size': 14}),
        ('  한(漢)나라 건국',
         {'font_size': 14, 'bold': True}),
        ('명장면', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  ① 약법삼장(BC 206) · ② 자기 평가(BC 202) · ③ 백등산 탈출',
         {'font_size': 13, 'color': SUB}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅴ. 인물 — 유방')
def s_yubang_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 유방', page, total)
    add_title(slide, '유방의 자기 평가 — "사람을 알아보고 쓸 줄 아는 능력"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                'BC 202년, 천하 통일 후 신하들 앞에서',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('책략(運籌)',  '내가 장량(張良)보다 못하다'),
        ('내정(行政)',  '내가 소하(蕭何)보다 못하다'),
        ('군사(軍事)',  '내가 한신(韓信)보다 못하다'),
    ]
    top = 3.3
    for tag, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.7), PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.15), Inches(3.0), Inches(0.4),
                    tag, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(9.1), Inches(0.7),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.9), Inches(top + 0.15), Inches(8.8), Inches(0.4),
                    desc, font_size=15, color=INK)
        top += 0.78
    add_filled_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.3),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.95), Inches(11.7), Inches(1.1), [
        ('"그러나 이 세 사람을 알아보고 쓸 줄 아는 것 — 내가 천하를 얻은 까닭"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"항우는 범증 한 사람조차 활용하지 못했으니, 이것이 그가 나에게 진 까닭"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# 당 태종
@S('Ⅴ. 인물 — 당 태종')
def s_taizong_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 당 태종', page, total)
    add_title(slide, '당 태종 이세민(唐太宗 李世民, 599~649)',
              '— 정관(貞觀)의 치(治)의 명군')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '太\n宗', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('어두운 시작 — 현무문(玄武門)의 변(626)', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  형 이건성·동생 이원길을 직접 살해',
         {'font_size': 13}),
        ('  아버지를 강제 양위하게 하고 즉위',
         {'font_size': 13}),
        ('  → 즉위 자체가 형제 살해의 도덕적 부담 위',
         {'font_size': 13, 'color': SUB}),
        ('평생의 자기 절제로 갚음', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  위징의 직간 수용, 자기 분노 다스림, 백성 살핌',
         {'font_size': 13}),
        ('정관의 치 (627~649)', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  중국 역사상 가장 모범적 통치 시대',
         {'font_size': 14, 'bold': True}),
        ('위징과의 관계', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  본래 적이었던 자를 간의대부로 발탁',
         {'font_size': 13}),
        ('  평생 200여 차례 직간 — "세 거울"의 일화',
         {'font_size': 13, 'color': SUB}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅴ. 인물 — 당 태종')
def s_taizong_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 당 태종', page, total)
    add_title(slide, '"세 거울(三鏡)" — 위징 사후 태종의 통곡')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(3.2), [
        ('위징(643년 사망) 사후 태종이 통곡하며',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"以銅爲鑑 可以正衣冠"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('구리로 거울을 삼으면 의관을 바로잡을 수 있고',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"以古爲鑑 可以知興替"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('옛것을 거울로 삼으면 흥망성쇠를 알 수 있고',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"以人爲鑑 可以明得失"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('사람을 거울로 삼으면 잘잘못을 알 수 있다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.5),
                '"이제 위징이 죽으니 나는 거울 하나를 잃었구나"',
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
                '비판하는 신하를 보배로 여긴 명군의 진수',
                font_size=13, color=SUB, bold=True, align=PP_ALIGN.CENTER)


# 진 시황
@S('Ⅴ. 인물 — 진시황')
def s_qin_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 진시황', page, total)
    add_title(slide, '진 시황(秦始皇, BC 259~210) — 반면교사',
              '"무력과 법으로 얻은 권력은 덕(德)으로 지키지 못하면 무너진다"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '始\n皇', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('위대한 능력', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  중국 최초의 황제 · 6국 통일 (BC 221)',
         {'font_size': 14}),
        ('  도량형·문자·화폐·도로 통일',
         {'font_size': 14}),
        ('  만리장성 건설 · 거대한 능력의 군주',
         {'font_size': 14}),
        ('자기 절제의 실패', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  분서갱유(BC 213~212) — 지식인 학살',
         {'font_size': 14}),
        ('  만리장성·아방궁·여산릉 동시 추진 — 70~150만 동원',
         {'font_size': 14}),
        ('  불로장생 추구 — 방사들에게 속아 수은 마심',
         {'font_size': 14}),
        ('  법가 일변도 — 가혹한 형벌만 의존',
         {'font_size': 14}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅴ. 인물 — 진시황')
def s_qin_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 진시황', page, total)
    add_title(slide, '결과 — 통일 14년 만의 멸망')
    timeline = [
        ('BC 221', '6국 통일 · 황제 등극',                 False),
        ('BC 213~212', '분서갱유 · 사상 통제',              False),
        ('BC 210', '시황 사망 (50세) · 환관 조고 농단',    True),
        ('BC 209', '진승(陳勝)·오광(吳廣)의 농민 봉기',     True),
        ('BC 206', '유방이 함양 입성 → 진 멸망',           True),
    ]
    top = 2.3
    for era, event, is_fall in timeline:
        c = ACCENT if is_fall else SUB
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.75), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.18), Inches(2.5), Inches(0.4),
                    era, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_fall else PALE
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(9.6), Inches(0.75), bg)
        add_textbox(slide, Inches(3.4), Inches(top + 0.18), Inches(9.3), Inches(0.4),
                    event, font_size=14, color=INK)
        top += 0.85
    add_filled_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_textbox(slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.5),
                '"무력과 법으로 얻은 권력은 덕(德)으로 지키지 못하면 무너진다"',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# 항우
@S('Ⅴ. 인물 — 항우')
def s_hangwoo_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 항우', page, total)
    add_title(slide, '항우(項羽, BC 232~202) — 재능은 있으나 사람을 잃은 비극',
              '"천하 제일의 무인이 천하 제일의 외톨이가 되어 죽다"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '項\n羽', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('천하 명장의 능력', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  초나라 명문가 출신 · 무력 천하 제일',
         {'font_size': 14}),
        ('  거록의 전투(BC 207) — 진의 주력군 격파',
         {'font_size': 14}),
        ('  한때 천하 패권을 손에 쥠',
         {'font_size': 14}),
        ('결정적 실수 — 홍문연(鴻門宴, BC 206)', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  함양 진격 후 유방을 죽일 결정적 기회',
         {'font_size': 14}),
        ('  책사 범증이 거듭 "죽이라!" 신호',
         {'font_size': 14}),
        ('  항우는 결단을 못 내려 유방을 살려 보냄',
         {'font_size': 14, 'bold': True, 'color': ACCENT}),
        ('  → 이 한 번의 우유부단이 4년 후 자기 죽음의 원인',
         {'font_size': 13, 'color': SUB}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅴ. 인물 — 항우')
def s_hangwoo_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 항우', page, total)
    add_title(slide, '비참한 최후 — 해하(垓下)·사면초가·패왕별희 (BC 202)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('해하(垓下)에 한 군에 포위',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('사면초가(四面楚歌) — 사방에서 초나라 노래',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('사랑하는 우희(虞姬)와의 이별 노래 — 우희 자결',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('力 拔 山 兮 氣 蓋 世',
         {'font_size': 26, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('역발산혜기개세 — "힘은 산을 뽑고 기개는 세상을 덮는다"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('時 不 利 兮 騅 不 逝',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('시불리혜추불서 — "때가 이롭지 않으니 명마(騅)도 가지 않는구나"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('오강(烏江)에서 부끄러워 자결 (31세)',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


# 강태공
@S('Ⅴ. 인물 — 강태공')
def s_taegong_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 강태공', page, total)
    add_title(slide, '강태공(姜太公) — 인재는 나이·출신을 따지지 않는다',
              '"위수 강가의 70 노인 → 천하 재상"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('주 문왕(文王) 시대 (BC 11세기)',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('70 노인이 위수 강가에서 곧은 낚시(미늘 없는 낚시)로 시간을 보내고 있었다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('사람들이 비웃었으나 그는 답했다:',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"나는 물고기를 낚는 것이 아니라 때(時)를 낚고 있다"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('사냥 나온 문왕이 그를 만나 대화하니 천하의 책략이 그의 머릿속에',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('문왕: "내 할아버지(太公)가 바라던 사람이다" → 태공망(太公望)',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('문왕 사후 무왕을 도와 은(殷) 주왕을 멸하고 주(周) 건국 (BC 1046)',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅴ. 인물 — 강태공')
def s_taegong_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 강태공', page, total)
    add_title(slide, '강태공이 남긴 가르침 — 『육도(六韜)』')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('인재는 나이·출신·외모와 관계없다',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('70 노인의 머리에 천하의 책략이 들어 있다는 것을',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('알아본 문왕의 눈이 곧 주(周) 건국의 시작',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('현대적 교훈 — 외형보다 실력',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 학벌·간판·외모에 가려진 진짜 능력을 보라',
         {'font_size': 14, 'space_before': 6}),
        ('• 무명·은퇴·노년의 사람 안에 보석이 있을 수 있다',
         {'font_size': 14, 'space_before': 4}),
        ('• "때를 낚는 자"를 알아보는 안목이 곧 리더의 자질',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


# 장량
@S('Ⅴ. 인물 — 장량')
def s_jangryang_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 장량', page, total)
    add_title(slide, '장량(張良, BC 250?~186) — 참모의 정수',
              '"한 줌의 곡식과 영구히 벼슬 없는 한가함을 원할 뿐"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '張\n良', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('출신', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  한(韓)나라 5대 재상의 후손',
         {'font_size': 14}),
        ('  진이 한을 멸한 후 진시황 암살 시도 (BC 218, 박랑사)',
         {'font_size': 14, 'color': SUB}),
        ('병법 전수', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  황석공(黃石公)에게서 『삼략(三略)』을 받음 (전설)',
         {'font_size': 14}),
        ('한 고조의 최고 책사', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  유방을 만나 천하 통일의 1등 공신',
         {'font_size': 14}),
        ('  홍문연에서 유방을 살린 책략',
         {'font_size': 14}),
        ('  결정적 전투의 모든 책략이 그의 머리에서',
         {'font_size': 14, 'bold': True}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅴ. 인물 — 장량')
def s_jangryang_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 장량', page, total)
    add_title(slide, '공성신퇴(功成身退) — 떠날 때를 아는 자만이 살아남는다')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('천하 통일 후 — 유방이 권력과 부를 주려 함',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('장량은 모든 것을 사양',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('"한 줌의 곡식과 영구히 벼슬 없는 한가함을 원할 뿐"',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('산속에 들어가 신선술을 닦으며 천수(天壽)를 누림',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('한신과 정확히 반대되는 처신',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('장량 — 떠나서 살아남음',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('한신 — 머물러서 토사구팽',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"공을 세우는 능력보다 공을 세운 후 처신이 더 어렵다"',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)


# 한신
@S('Ⅴ. 인물 — 한신')
def s_hansin_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 한신', page, total)
    add_title(slide, '한신(韓信, BC ?~196) — 재능과 비극의 양면',
              '"백만 대군을 다다익선(多多益善)으로 부린 천재"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '韓\n信', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('처음', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  회음(淮陰)의 가난한 청년',
         {'font_size': 14}),
        ('  빨래터 노파의 밥을 얻어먹음',
         {'font_size': 14, 'color': SUB}),
        ('  과하지욕(跨下之辱) — 깡패의 가랑이 사이로',
         {'font_size': 14, 'color': SUB}),
        ('발탁', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  처음 항우 밑에 있다가 알아주지 않자 유방에게로',
         {'font_size': 14}),
        ('  유방도 못 알아봐 떠나려 함',
         {'font_size': 14}),
        ('  소하월하추한신(蕭何月下追韓信)',
         {'font_size': 14, 'bold': True, 'color': ACCENT}),
        ('  → 마침내 대장군 임명',
         {'font_size': 13, 'color': SUB}),
        ('천재성', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  "다다익선(多多益善)" — 백만 대군을 자유자재로',
         {'font_size': 14}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅴ. 인물 — 한신')
def s_hansin_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 한신', page, total)
    add_title(slide, '토사구팽(兎死狗烹) — 공을 세운 후 처신의 어려움')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('한신은 자기 공을 자랑하고 자기를 낮출 줄 모름',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('유방: "그대는 얼마나 군대를 부릴 수 있는가?"',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('한신: "다다익선(多多益善) — 많을수록 좋다"',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('유방: "나는?" 한신: "10만 정도"',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('유방의 의심을 산 결정적 답변 → BC 196년 모반의 누명으로 처형',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"狡 兎 死 走 狗 烹"',
         {'font_size': 28, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('교토사 주구팽 — "날쌘 토끼가 죽으면 사냥개가 삶긴다"',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"高 鳥 盡 良 弓 藏" — 높이 나는 새가 다 잡히면 좋은 활은 거두어진다',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# 위징
@S('Ⅴ. 인물 — 위징')
def s_wijing_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 위징', page, total)
    add_title(slide, '위징(魏徵, 580~643) — 직간(直諫)의 영원한 모범',
              '"평생 200여 차례 태종에게 직접 비판"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '魏\n徵', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('험난한 출세', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  당 건국기에 여러 군주를 거침',
         {'font_size': 14}),
        ('  태자 이건성(태종의 형)의 모사(謀士)',
         {'font_size': 14, 'color': SUB}),
        ('  이건성에게 "이세민을 죽이라"고 권한 인물',
         {'font_size': 14, 'color': SUB}),
        ('현무문 후 — 태종은 살리고 곁에 둠', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  "왜 내 형제 사이를 이간했느냐"',
         {'font_size': 14}),
        ('  위징: "태자께서 제 말을 들으셨다면 오늘의 화를 당하지 않으셨을 것"',
         {'font_size': 13, 'color': SUB}),
        ('  → 두려움 없는 답에 태종 감복',
         {'font_size': 13, 'color': SUB}),
        ('  → 간의대부(諫議大夫)로 발탁',
         {'font_size': 14, 'bold': True}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅴ. 인물 — 위징')
def s_wijing_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 위징', page, total)
    add_title(slide, '"이 시골 늙은이를 죽이고 말겠다!" — 황후의 지혜')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('위징이 너무 심하게 간하여 태종이 화가 나 침전에서 황후에게',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"이 시골 늙은이(村翁)를 죽이고 말겠다!"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('황후 장손씨가 정장(正裝)으로 절을 올리며',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"君 明 臣 直" — 군명신직',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"임금이 현명해야 신하가 곧은 말을 합니다',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 위징이 곧은 말을 하는 것은 폐하가 명군이시기 때문이니',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        (' 감축드립니다"',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


# 8인 한눈에
@S('Ⅴ. 인물 — 한눈에')
def s_8people(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 한눈에', page, total)
    add_title(slide, '결정적 인물 8인 — 한눈에 보기')
    rows = [
        ('한 고조 유방',   'BC 256~195',  '평민 황제 — 인재 활용의 모범'),
        ('당 태종 이세민', '599~649',     '정관의 치 — 비판 수용의 명군'),
        ('진 시황',         'BC 259~210',  '반면교사 — 무력만의 통치'),
        ('항 우',           'BC 232~202',  '재능은 있으나 사람을 잃은 비극'),
        ('강 태 공',         'BC 11세기',   '인재는 나이를 따지지 않음'),
        ('장 량',           'BC 250?~186', '참모의 정수 — 공성신퇴'),
        ('한 신',           'BC ?~196',    '천재와 토사구팽'),
        ('위 징',           '580~643',     '직간의 영원한 모범'),
    ]
    top = 2.0
    for i, (name, era, char) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(3.0), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.5), Inches(top + 0.13), Inches(3.0), Inches(0.4),
                    name, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.6), Inches(top), Inches(2.5), Inches(0.55), PALE)
        add_textbox(slide, Inches(3.6), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    era, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.2), Inches(top), Inches(6.6), Inches(0.55), bg)
        add_textbox(slide, Inches(6.4), Inches(top + 0.13), Inches(6.3), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.6


@S('Ⅴ. 인물 — 한눈에')
def s_lesson(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 인물 — 한눈에', page, total)
    add_title(slide, '8인의 교훈 — 두 갈래의 길')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.5),
                '성공의 길', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('• 유방 — 사람을 알아보고 씀', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 4}),
        ('• 당 태종 — 비판을 받아들임', {'font_size': 15, 'bold': True, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('• 강태공 — 때를 낚다', {'font_size': 15, 'bold': True, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('• 장량 — 공을 이루고 물러섬', {'font_size': 15, 'bold': True, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('• 위징 — 두려움 없이 직간', {'font_size': 15, 'bold': True, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.9), Inches(0.5),
                '실패의 길', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('• 진 시황 — 자기 욕망을 못 다스림', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 4}),
        ('• 항 우 — 자기 능력만 믿음', {'font_size': 15, 'bold': True, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('• 한 신 — 떠날 때를 모름', {'font_size': 15, 'bold': True, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('공 통 점', {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 8}),
        ('  자기 정립의 실패', {'font_size': 14, 'color': SUB}),
        ('  사람을 못 알아봄', {'font_size': 14, 'color': SUB}),
        ('  비판을 안 들음', {'font_size': 14, 'color': SUB}),
    ], line_spacing=1.3)


# ============== Ⅵ. 7대 핵심 원리 ==============
@S('Ⅵ. 7대 원리')
def s_7principles_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 7대 원리', page, total)
    add_title(slide, '7대 핵심 원리 — 동심원 구조',
              '"1번 수신 없이 나머지 6가지는 의미가 없다"')
    items = [
        ('7', '위기의 결단력', '결정적 순간의 판단'),
        ('6', '민심이 천심',   '권력의 정당성'),
        ('5', '공성신퇴',       '떠날 때를 알기'),
        ('4', '법과 덕의 균형', '통치의 두 축'),
        ('3', '간언 수용',       '비판을 듣는 그릇'),
        ('2', '인재 활용',       '사람을 쓰는 능력'),
        ('1', '수신(修身)',     '자기 정립 — 모든 것의 출발점'),
    ]
    top = 2.0
    for num, name, desc in items:
        is_base = (num == '1')
        c = ACCENT if is_base else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_base else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.6), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(0.7), Inches(0.4),
                    num, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(3.5), Inches(0.6), bg)
        add_textbox(slide, Inches(1.5), Inches(top + 0.13), Inches(3.4), Inches(0.4),
                    name, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.0), Inches(top), Inches(7.8), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.2), Inches(top + 0.13), Inches(7.5), Inches(0.4),
                    desc, font_size=14, color=INK, bold=is_base)
        top += 0.7
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '1번부터 7번까지 순서가 있는 단계 — 모두 연결된 한 그림',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def make_principle_slide(num, title_kor, title_subtitle, original, principle_text,
                          example, opposite, modern):
    @S('Ⅵ. 7대 원리')
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, 'Ⅵ. 7대 원리', page, total)
        add_title(slide, f'원리 {num} — {title_kor}', title_subtitle)
        add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.5), Inches(2.33), Inches(12.3), Inches(0.55),
                    principle_text, font_size=16, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(0.5), Inches(3.1), Inches(5.9), Inches(1.3), PALE)
        add_textbox(slide, Inches(0.5), Inches(3.25), Inches(5.9), Inches(0.4),
                    '대표 사례', font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(5.5), Inches(0.6),
                    example, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.9), Inches(3.1), Inches(5.9), Inches(1.3),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(6.9), Inches(3.25), Inches(5.9), Inches(0.4),
                    '반대 사례', font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(7.1), Inches(3.7), Inches(5.5), Inches(0.6),
                    opposite, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.4),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
                    '오늘의 적용', font_size=16, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_paragraphs(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.7),
                       modern, line_spacing=1.4, font_size=14)
    return renderer


make_principle_slide('1', '수신(修身)',
    '"자기를 다스리는 자만이 남을 다스릴 수 있다"',
    None,
    '"자기 분노·욕망·게으름·편견을 못 다스리는 자가 어떻게 남을 다스리겠는가"',
    '당 태종 — 위징의 비판 앞에서 자기 분노를 다스림',
    '진 시황 — 불로장생·과시 욕망 못 다스려 왕조 멸망',
    [
        ('• 매일 자기 점검 시간 확보', {'font_size': 14}),
        ('• 분노 폭발 전 0.5초 멈추기', {'font_size': 14, 'space_before': 4}),
        ('• "내가 가장 다스려야 할 사람은 누구인가?" → 나 자신', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_principle_slide('2', '인재 활용',
    '"사람을 알아보고 쓸 줄 아는 능력이 리더의 첫째 자질"',
    None,
    '"혼자 잘하려는 리더는 망한다 — 각 분야 최고를 알아보고 적재적소에"',
    '유방 — 장량·소하·한신 三傑을 알아보고 활용',
    '항우 — 충신 범증 한 사람조차 활용 못 함',
    [
        ('• "내 옆 사람은 무엇을 잘하는가"를 묻기', {'font_size': 14}),
        ('• 자기보다 잘난 사람을 두려워하지 말고 영입', {'font_size': 14, 'space_before': 4}),
        ('• 人事(인사)가 萬事(만사) — 사람 쓰기가 모든 일의 시작', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_principle_slide('3', '간언 수용',
    '"비판하는 신하를 두려워 말고 곁에 두라"',
    None,
    '"칭찬만 듣는 군주는 반드시 망한다 — 직간을 받아들일 그릇이 곧 군주의 그릇"',
    '당 태종 — 위징의 200여 차례 직간 수용 · "세 거울"',
    '수 양제 — 직간자 모두 죽이거나 내쫓아 멸망',
    [
        ('• 자기 곁에 솔직히 말해 주는 사람이 몇 명인가 점검', {'font_size': 14}),
        ('• 부하의 비판에 화내지 말고 감사 ("말해 주어 고맙다")', {'font_size': 14, 'space_before': 4}),
        ('• "예스맨"으로 둘러싸이는 순간 추락 시작', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_principle_slide('4', '법(法)과 덕(德)의 균형',
    '"법으로 다스리고 덕으로 이끈다"',
    None,
    '"한비자의 법치만으로도, 공자의 덕치만으로도 부족 — 둘이 함께"',
    '한 고조 — 약법삼장(約法三章) · "법은 단순하게 분명하게"',
    '진(秦) — 가혹한 법치 일변도로 멸망',
    [
        ('• 명확한 규칙(法) + 따뜻한 인정(德)의 균형', {'font_size': 14}),
        ('• "원칙은 분명하게, 적용은 인간적으로"', {'font_size': 14, 'space_before': 4}),
        ('• 처벌은 신속·분명하게, 격려는 자주', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_principle_slide('5', '공성신퇴(功成身退)',
    '"공을 세웠으면 물러설 줄 알라"',
    None,
    '"토사구팽 — 큰 공을 세운 신하는 군주의 의심을 받기 쉽다"',
    '장량 — 천하 통일 후 모든 권력 사양·산속으로',
    '한신 — 다다익선 자랑하다 모반 누명·처형',
    [
        ('• 큰 성공·승진·인정 받았을 때 더욱 자기를 낮춤', {'font_size': 14}),
        ('• "내가 한 것이 아니라 우리가 한 것"', {'font_size': 14, 'space_before': 4}),
        ('• 떠날 때를 아는 자가 진짜 능력자', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_principle_slide('6', '민심이 천심',
    '"백성의 마음을 잃으면 천하를 잃는다"',
    None,
    '"무력·법·재력으로 권력을 얻을 수 있어도, 그것을 지키는 것은 오직 민심"',
    '한 고조의 군기 — 백성 재산 보호로 민심 얻음',
    '항우의 약탈 — 백성에게 잔인하여 민심 잃음',
    [
        ('• "조직 구성원의 마음"을 알지 못하는 리더는 결국 외톨이', {'font_size': 14}),
        ('• 부하의 처지·고충을 알고 헤아리는 자세', {'font_size': 14, 'space_before': 4}),
        ('• "내가 옳다"가 아니라 "사람들이 어떻게 받아들이는가"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_principle_slide('7', '위기의 결단력',
    '"위기 때의 판단이 진짜 리더를 결정한다"',
    None,
    '"평상시에는 누구나 리더로 보이지만, 위기에서 진짜 리더가 가려진다"',
    '주유 — 적벽대전 결단 · 유방 — 백등산 탈출',
    '항우 — 홍문연에서 유방을 죽일 결단 못 내림',
    [
        ('• 결정적 순간 망설이는 리더가 가장 위험', {'font_size': 14}),
        ('• "60~70% 정보로 결단하고, 나머지는 실행하며 보정"', {'font_size': 14, 'space_before': 4}),
        ('• 결단 후에는 책임지는 자세', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])


@S('Ⅵ. 7대 원리')
def s_7integ(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 7대 원리', page, total)
    add_title(slide, '7대 원리 통합 — "리더십은 기술이 아니라 인격"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(4.4), [
        ('"리더십은 기술이 아니라 인격이다"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('인격은',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('① 자기 다스림에서 시작해',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('② 사람을 알아보고',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('③ 비판을 듣고',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('④ 균형을 잡고',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('⑤ 물러설 줄 알며',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('⑥ 민심을 헤아리고',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('⑦ 위기 때 결단할 줄 아는',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('— 7가지 그릇의 종합 —',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
    ], line_spacing=1.2)


# ============== Ⅶ. 8대 자질 ==============
@S('Ⅶ. 8대 자질')
def s_8q_overview(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 8대 자질', page, total)
    add_title(slide, '8대 리더 자질 — 한눈에 보기',
              '"1번이 없으면 모든 것이 무너지고, 8번이 없으면 모든 것이 변질된다"')
    items = [
        ('1', '자기 정립(正)',     '자기를 이긴 자만이 남을 이긴다'),
        ('2', '인재 안목',         '사람 알아보는 능력이 첫째 자질'),
        ('3', '위임의 용기',       '알아봤으면 믿고 맡겨라'),
        ('4', '비판 수용력',       '직간하는 자가 진짜 친구'),
        ('5', '결단력',            '위기에 진짜 리더가 가려진다'),
        ('6', '공정함',            '호오에 휘둘리는 순간 신뢰 잃는다'),
        ('7', '장기적 안목',       '5년 후를 보는 자가 5년을 산다'),
        ('8', '물러섬의 지혜',     '잘 떠나는 것이 가장 어려운 일'),
    ]
    top = 2.0
    for num, name, summary in items:
        is_axis = (num in ['1', '8'])
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_axis else PALE
        c = ACCENT if is_axis else SUB
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.6), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(0.7), Inches(0.4),
                    num, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(3.0), Inches(0.6), bg)
        add_textbox(slide, Inches(1.5), Inches(top + 0.13), Inches(2.9), Inches(0.4),
                    name, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.5), Inches(top), Inches(8.3), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.7), Inches(top + 0.13), Inches(8.0), Inches(0.4),
                    summary, font_size=14, color=INK)
        top += 0.65


def make_quality_slide(num, name, summary, example, opposite, diagnose):
    @S('Ⅶ. 8대 자질')
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, 'Ⅶ. 8대 자질', page, total)
        add_title(slide, f'자질 {num} — {name}', summary)
        add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(2.3), PALE)
        add_textbox(slide, Inches(0.5), Inches(2.45), Inches(5.9), Inches(0.4),
                    '대표 사례', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.7), Inches(2.9), Inches(5.5), Inches(1.5),
                    example, font_size=13, color=INK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(2.3),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(6.9), Inches(2.45), Inches(5.9), Inches(0.4),
                    '반대 사례', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(7.1), Inches(2.9), Inches(5.5), Inches(1.5),
                    opposite, font_size=13, color=INK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.4),
                    '진단 질문 — 자기에게 물어보라', font_size=16, bold=True,
                    color=ACCENT, align=PP_ALIGN.CENTER)
        add_paragraphs(slide, Inches(0.8), Inches(5.45), Inches(11.7), Inches(1.5),
                       diagnose, line_spacing=1.4, font_size=14)
    return renderer


make_quality_slide('1', '자기 정립(正)',
    '자기 약점을 알고 다스리는 능력 — 모든 리더십의 근본',
    '당 태종 — 자기 분노를 위징 비판 앞에서 다스림\n유방 — "나는 장량보다 못하다" 자기 한계 인정',
    '항우 — 자기 능력을 자만하여 멸망\n진 시황 — 자기 욕망 못 다스려 왕조 멸망',
    [
        ('• "지난 한 달 동안 내가 가장 후회하는 행동은?"', {'font_size': 14}),
        ('• "내가 다른 사람을 가장 화나게 한 순간과 그 원인은?"', {'font_size': 14, 'space_before': 4}),
        ('• "내가 가장 약한 부분(분노·게으름·자만·편견·욕심)은?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_quality_slide('2', '인재 안목',
    '사람을 알아보는 눈 — 실력을 꿰뚫어 보는 능력',
    '유방 — 한신을 처음 만나 알아봄 (소하 추천 수용)\n주 문왕 — 70 노인 강태공을 한눈에\n유비 — 제갈량 삼고초려',
    '항우 — 한신이 자기 밑에 있을 때 알아보지 못함',
    [
        ('• "내 주변에 진짜 능력 있는 사람이 누구인가?"', {'font_size': 14}),
        ('• "그 능력이 무엇인지 구체적으로 말할 수 있는가?"', {'font_size': 14, 'space_before': 4}),
        ('• "나는 그를 알아보고 그에 합당한 기회를 주었는가?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_quality_slide('3', '위임의 용기',
    '알아본 사람을 진짜로 쓰는 능력 — "알아보고도 못 맡기면 안목이 무의미"',
    '유방 — 한신에게 백만 대군 지휘권을 통째로 줌\n당 태종 — 위징에게 직간의 자유 보장\n세종 — 황희·맹사성에게 정사 위임',
    '항우 — 범증의 책략을 듣지도 따르지도 않음\n현대 마이크로매니저 — 시키고도 끝까지 간섭',
    [
        ('• "내가 다른 사람에게 진짜로 맡긴 일이 있는가?"', {'font_size': 14}),
        ('• "맡기고도 자꾸 들여다보고 있지 않은가?"', {'font_size': 14, 'space_before': 4}),
        ('• "부하의 실패를 한 번이라도 감싸 준 적이 있는가?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_quality_slide('4', '비판 수용력',
    '듣기 싫은 말을 듣는 그릇 — "이것이 없으면 예스맨에 둘러싸인다"',
    '당 태종 — 황후의 지혜로 위징 비판 수용\n"君明臣直 — 임금이 현명해야 신하가 곧다"',
    '수 양제 — 직간자 모두 죽이고 멸망\n현대 폐쇄적 CEO — 비판 듣고 보복',
    [
        ('• "최근 누가 나에게 솔직한 비판을 했는가?"', {'font_size': 14}),
        ('• "그때 내 반응은 어땠는가?"', {'font_size': 14, 'space_before': 4}),
        ('• "내 곁에 진짜 직언해 주는 사람이 몇 명인가?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_quality_slide('5', '결단력',
    '결정적 순간의 판단 — "위기에서 진짜 리더가 가려진다"',
    '유방의 관중 진격\n주유의 적벽대전 결단\n이순신의 명량해전 ("12척으로 133척")',
    '항우의 홍문연 — 결정적 기회에 망설여 놓침\n현대 "더 검토해 보자"의 무한 반복',
    [
        ('• "최근 한 달 내가 내린 가장 어려운 결단은?"', {'font_size': 14}),
        ('• "결단을 미루다가 후회한 일이 있는가?"', {'font_size': 14, 'space_before': 4}),
        ('• "결단 후 자꾸 흔들리지 않는가?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_quality_slide('6', '공정함',
    '호오에 휘둘리지 않는 평등한 대우 — "감정의 노예가 되지 않는 자기 통제"',
    '제갈량의 읍참마속(泣斬馬謖)\n— 자기가 가장 아끼던 마속을 울며 처형\n당 태종 — 정치적 입장이 달랐던 위징을 중용',
    '현대 사내 정치 — 친한 사람에게 좋은 자리\n싫은 사람에게 부당한 처우',
    [
        ('• "내가 편애하는 사람이 있는가?"', {'font_size': 14}),
        ('• "내가 부당하게 박하게 대하는 사람이 있는가?"', {'font_size': 14, 'space_before': 4}),
        ('• "내 결정의 기준은 객관적인가, 감정적인가?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_quality_slide('7', '장기적 안목',
    '눈앞 이익에 흔들리지 않는 큰 그림 — "5년 후를 보는 자가 5년을 산다"',
    '세종 — 한글 창제 (사대부 반대에도)\n당 태종 — 단기 업적 욕심 자제',
    '수 양제 — 대운하·만리장성 동시 추진으로 멸망\n현대 분기 실적 매몰',
    [
        ('• "내가 최근 한 결정 중 단기 이익에 휘둘린 것은?"', {'font_size': 14}),
        ('• "5년 후 이 결정이 어떻게 평가될까?"', {'font_size': 14, 'space_before': 4}),
        ('• "장기적으로 옳은 길과 단기적으로 편한 길 중 무엇을?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])

make_quality_slide('8', '물러섬의 지혜',
    '시대와 자기 자리를 아는 것 — "잘 떠나는 것이 잘 시작하는 것보다 어렵다"',
    '장량 — 한 고조 통일 후 산속으로\n범려 — 월왕 구천 패업 후 미련 없이 떠남\n이순신 — "내 죽음을 적에게 알리지 말라"',
    '한신 — 공을 세우고도 떠나지 못해 토사구팽\n권력 말기까지 매달리는 정치인·CEO',
    [
        ('• "내 후임자가 누구인지 분명한가?"', {'font_size': 14}),
        ('• "내가 없어도 조직이 굴러가는가?"', {'font_size': 14, 'space_before': 4}),
        ('• "이 자리가 내 정체성의 전부는 아닌가?"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ])


@S('Ⅶ. 8대 자질')
def s_8q_check(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 8대 자질', page, total)
    add_title(slide, '8대 자질 자가 진단 체크리스트')
    items = [
        '오늘 내가 가장 다스려야 할 사람이 누구인지 알고 있다 — 자기 정립',
        '내 옆 사람의 강점을 구체적으로 말할 수 있다 — 인재 안목',
        '한 사람에게 진짜로 권한과 책임을 맡긴 일이 있다 — 위임',
        '내 곁에 솔직히 직언해 주는 사람이 1명 이상 있다 — 비판 수용',
        '최근 한 달 내가 책임지고 내린 결단이 있다 — 결단력',
        '나의 결정 기준이 객관적이라 자신할 수 있다 — 공정함',
        '내 결정의 5년 후 결과를 한 번이라도 자문해 보았다 — 장기 안목',
        '나의 후임자가 누구인지 분명히 그릴 수 있다 — 물러섬',
    ]
    top = 2.2
    for item in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.55), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(0.55), Inches(0.4),
                    '☐', font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.25), Inches(top), Inches(11.55), Inches(0.55), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.13), Inches(11.3), Inches(0.4),
                    item, font_size=14, color=INK)
        top += 0.62


@S('Ⅶ. 8대 자질')
def s_8q_integ(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 8대 자질', page, total)
    add_title(slide, '가장 중요한 한 가지 — 첫째 조건의 절대성')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('8가지 중 단 하나만 고르라면',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('첫 번째 — 자기 정립(正)',
         {'font_size': 24, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('"리더가 자기를 다스리지 못하면 다른 7가지는 모두 의미가 없다"',
         {'font_size': 15, 'bold': True, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('자기 분노를 못 다스리는 리더가 어떻게 공정하겠는가?',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('자기 욕심을 못 다스리는 리더가 어떻게 위임하겠는가?',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('자기 자만을 못 다스리는 리더가 어떻게 비판을 받아들이겠는가?',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"수신(修身)이 곧 치국(治國)의 근본"',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.25)


# ============== Ⅷ. 명구·명장면 10선 ==============
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=42):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=20, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('Ⅷ. 명구 (1/10)',
    '夫 運 籌 帷 帳 之 中\n決 勝 千 里 之 外   吾 不 如 子 房',
    '부 운주유장지중 결승천리지외 오불여자방',
    '군영에서 책략을 짜내어 천 리 밖의 승리를 결정짓는 것은 내가 장량보다 못하다',
    '한 고조 유방, BC 202년 자기 평가 (사기 고조본기)', hanmun_size=22), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (2/10)',
    '約 法 三 章',
    '약 법 삼 장',
    '법은 단지 세 가지 — 살인은 사형, 상해·도둑은 처벌, 그 외 진의 법은 모두 폐지',
    '한 고조 유방, BC 206년 함양 입성 후', hanmun_size=110), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (3/10)',
    '狡 兎 死   走 狗 烹\n高 鳥 盡   良 弓 藏',
    '교토사 주구팽 · 고조진 양궁장',
    '날쌘 토끼가 죽으면 사냥개가 삶기고, 높이 나는 새가 다 잡히면 좋은 활은 거두어진다',
    '한신, BC 196년 처형 직전의 한탄 (사기 회음후열전)', hanmun_size=26), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (4/10)',
    '功 成 身 退',
    '공 성 신 퇴',
    '공을 이루었으면 몸을 물린다 — 장량이 한 고조 통일 후 권력 사양',
    '장량의 처세 · 노자 9장에서 유래', hanmun_size=110), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (5/10)',
    '以 銅 爲 鑑   可 以 正 衣 冠\n以 古 爲 鑑   可 以 知 興 替\n以 人 爲 鑑   可 以 明 得 失',
    '이동위감 가이정의관 · 이고위감 가이지흥체 · 이인위감 가이명득실',
    '구리 거울로 의관을, 옛것을 거울로 흥망을, 사람을 거울로 득실을 안다 — 三鏡',
    '당 태종, 위징 사후의 통곡 (정관정요)', hanmun_size=20), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (6/10)',
    '四 面 楚 歌\n力 拔 山 兮   氣 蓋 世',
    '사면초가 · 역발산혜 기개세',
    '사방에서 초나라 노래 · "힘은 산을 뽑고 기개는 세상을 덮건만"',
    '항우, BC 202년 해하의 마지막 노래', hanmun_size=28), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (7/10)',
    '多 多 益 善   跨 下 之 辱',
    '다다익선 · 과하지욕',
    '"많을수록 좋다" — 한신의 군사 천재성 · "가랑이 사이의 모욕" — 한신의 인내',
    '한신의 두 면모 (사기 회음후열전)', hanmun_size=40), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (8/10)',
    '吾 釣 非 釣 魚   釣 時 也',
    '오 조비조어 조시야',
    '나는 물고기를 낚는 것이 아니라 때(時)를 낚고 있다 — 강태공의 곧은 낚시',
    '강태공, 위수 강가에서', hanmun_size=38), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (9/10)',
    '見 可 欲   則 思 知 足\n…   罰 所 及   則 思 無 以 怒 而 濫 刑',
    '견가욕 즉사지족 … 벌소급 즉사무이노이남형',
    '十思 — 위징이 태종에게 올린 「諫太宗十思疏」 · 매일 점검할 10가지 자기 점검',
    '위징, 정관 11년 「諫太宗十思疏」', hanmun_size=18), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (10/10)',
    '君 君   臣 臣   父 父   子 子',
    '군 군 신 신 부 부 자 자',
    '임금은 임금답게, 신하는 신하답게, 부모는 부모답게, 자식은 자식답게 — 정명(正名)',
    '공자, 논어 안연편 · 동양 제왕학의 근본 원리', hanmun_size=44), 'Ⅷ. 명구'))


# ============== Ⅸ. 한국의 제왕학 ==============
@S('Ⅸ. 한국 제왕학')
def s_kr_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 한국 제왕학', page, total)
    add_title(slide, '한국 제왕학 — 우리 역사의 두 빛',
              '"세종과 정조 — 동양 제왕학을 자기 살로 만든 군주"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('한국 제왕학의 특수성',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('성리학적 도덕주의가 강하게 결합',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"임금이 곧 도덕적 모범이어야 한다"는 강한 윤리 의식',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    items = [
        ('경연(經筵) — 임금이 신하들과 학문을 토론하는 자리',  '독특한 학습 제도'),
        ('사관(史官)의 일거수일투족 기록',                       '임금의 자기 경계심 극대화'),
        ('신하의 직간을 자기 도덕 점검의 도구로',                '단지 받아들임이 아님'),
        ('대학연의 — 임금이 반드시 공부해야 하는 책',           '세종 100번 완독'),
    ]
    top = 4.5
    for tag, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(7.0), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.8), Inches(top + 0.13), Inches(6.7), Inches(0.4),
                    tag, font_size=13, bold=True, color=WHITE)
        add_filled_rect(slide, Inches(7.7), Inches(top), Inches(5.1), Inches(0.55), PALE)
        add_textbox(slide, Inches(7.9), Inches(top + 0.13), Inches(4.9), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.62


@S('Ⅸ. 한국 제왕학')
def s_sejong_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 한국 제왕학', page, total)
    add_title(slide, '세종(世宗, 1397~1450) — 경연(經筵)의 임금',
              '"즉위 32년 동안 경연 1,898회 — 한국사 어느 임금보다도 많은"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '世\n宗', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('학습의 화신', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  경연 1,898회 — 매주 한두 번 꼴',
         {'font_size': 14, 'bold': True}),
        ('경연에서 다룬 책들', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  대학연의 · 정관정요 · 자치통감',
         {'font_size': 14}),
        ('  근사록 · 심경 · 성리대전 · 소학 · 국조보감',
         {'font_size': 13, 'color': SUB}),
        ('  동양 제왕학의 모든 핵심 텍스트가 망라',
         {'font_size': 13, 'color': SUB, 'bold': True}),
        ('대학연의 100번 완독', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  진덕수 『대학연의』를 평생 깊이 학습',
         {'font_size': 14}),
        ('  세종 통치의 사상적 토대',
         {'font_size': 14, 'bold': True}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅸ. 한국 제왕학')
def s_sejong_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 한국 제왕학', page, total)
    add_title(slide, '세종의 결실 — 한글·과학·민본 정치')
    items = [
        ('1443년',  '훈민정음 창제',          '"백성을 위한" 정치의 정수'),
        ('1420년',  '집현전 설치',            '인재 양성의 모범'),
        ('과학',     '장영실 등 과학자 활용', '인재 안목의 모범 — 노비 출신도 발탁'),
        ('법제',     '노비에게도 출산휴가',    '동시대 어느 나라에도 없던 민본 정책'),
        ('농업',     '농사직설 · 측우기',     '백성 삶의 직접적 개선'),
    ]
    top = 2.4
    for era, work, char in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.75), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.18), Inches(2.5), Inches(0.4),
                    era, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(3.5), Inches(0.75), PALE)
        add_textbox(slide, Inches(3.35), Inches(top + 0.18), Inches(3.3), Inches(0.4),
                    work, font_size=14, bold=True, color=ACCENT)
        add_filled_rect(slide, Inches(6.8), Inches(top), Inches(6.0), Inches(0.75),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.0), Inches(top + 0.18), Inches(5.7), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.85
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"세종은 동양 제왕학의 모든 원리를 한국에서 가장 완벽하게 구현한 군주"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅸ. 한국 제왕학')
def s_jeongjo_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 한국 제왕학', page, total)
    add_title(slide, '정조(正祖, 1752~1800) — 학자 군주',
              '"내가 더 이상 경들에게 배울 것이 없으니, 내가 직접 가르치겠다"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '正\n祖', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('사도세자의 비극', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  뒤주 사건(1762) — 11세에 직접 목격',
         {'font_size': 14, 'color': SUB}),
        ('  깊은 상처를 평생 학문·정치로 승화',
         {'font_size': 14}),
        ('초계문신제(抄啟文臣制, 1781)', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  경연을 폐지하고 오히려 자신이 가르치는 자리로 전환',
         {'font_size': 14}),
        ('  젊고 재능 있는 신하를 뽑아 정조가 직접 교육',
         {'font_size': 14, 'bold': True}),
        ('규장각(奎章閣, 1776)', {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 8}),
        ('  즉위 직후 국왕 직속 학술·정책 연구 기관',
         {'font_size': 14}),
        ('  정약용·박제가·이덕무 등 실학자 등용',
         {'font_size': 14, 'bold': True}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅸ. 한국 제왕학')
def s_jeongjo_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 한국 제왕학', page, total)
    add_title(slide, '『홍재전서(弘齋全書)』 — 한국사 가장 학문적인 임금')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('정조는 단지 학습자가 아니라 저술가였다',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('『홍재전서(弘齋全書)』 — 184권',
         {'font_size': 20, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('한국사에서 가장 학문적인 임금의 친저',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"학문 자체가 통치"라는 정관정요·대학연의 정신의 가장 깊은 실천',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 수원 화성 건설 — 실학과 제왕학의 결합',
         {'font_size': 14, 'space_before': 6}),
        ('• 탕평책 — 당파를 넘어선 인재 등용',
         {'font_size': 14, 'space_before': 4}),
        ('• 신해통공(辛亥通共, 1791) — 시전 상인의 독점권 폐지',
         {'font_size': 14, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅸ. 한국 제왕학')
def s_kr_essence(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 한국 제왕학', page, total)
    add_title(slide, '한국 제왕학의 결론 — 세종·정조의 공통점')
    items = [
        ('학습', '경연 1,898회 / 초계문신제',
         '"임금이 가장 많이 공부하는 사람"'),
        ('자기 절제',  '신하의 직간을 자기 도덕 점검의 도구로',
         '"성리학적 도덕주의가 결합된 한국 제왕학의 특수성"'),
        ('인재',     '집현전 / 규장각',
         '"노비 출신 장영실, 서얼 정약용까지 등용"'),
        ('민본',     '훈민정음 / 탕평·통공',
         '"백성을 위한" 정치의 가장 강력한 두 사례'),
        ('기록',     '실록 · 일성록 · 승정원일기',
         '"임금의 자기 경계심 극대화 — 세계사적 기록 문화"'),
    ]
    top = 2.3
    for tag, what, why in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.0), Inches(0.5),
                    tag, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.7), Inches(top), Inches(4.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(2.85), Inches(top + 0.22), Inches(4.3), Inches(0.5),
                    what, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.3), Inches(top), Inches(5.5), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.45), Inches(top + 0.22), Inches(5.3), Inches(0.5),
                    why, font_size=13, color=INK)
        top += 0.97


# ============== Ⅹ. 현대 의의 ==============
@S('Ⅹ. 현대 의의')
def s_modern_personality(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 현대 의의', page, total)
    add_title(slide, '현대 ① — "리더십은 기술이 아니라 인격의 표현"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.5),
                '서양 경영학', font_size=20, bold=True, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('리더십 = 스킬(skill)', {'font_size': 16, 'bold': True}),
        ('', {'font_size': 6}),
        ('• 의사소통 기술', {'font_size': 14, 'space_before': 8}),
        ('• 의사결정 모델', {'font_size': 14, 'space_before': 4}),
        ('• 동기부여 이론', {'font_size': 14, 'space_before': 4}),
        ('• 협상 전략 …', {'font_size': 14, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.9), Inches(0.5),
                '동양 제왕학', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('리더십 = 인격(德)의 표현', {'font_size': 16, 'bold': True, 'color': ACCENT}),
        ('', {'font_size': 6}),
        ('"좋은 리더가 되려면', {'font_size': 14, 'space_before': 8}),
        (' 좋은 사람이 되어야 한다"', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 6}),
        ('"자기 분노를 못 다스리는 사람이', {'font_size': 13, 'color': SUB, 'space_before': 6}),
        (' 어떻게 부하의 분노를 다스리겠는가?"', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)


@S('Ⅹ. 현대 의의')
def s_modern_relationship(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 현대 의의', page, total)
    add_title(slide, '현대 ② — "리더십은 혼자가 아니라 관계의 함수"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('동양 제왕학의 모든 명장면 — 두 사람의 관계',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('유방-장량 · 태종-위징 · 문왕-강태공 · 유비-제갈량 · 세종-황희',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('— 혼자 잘난 리더는 없다 —',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"누구를 곁에 두는가가 곧 리더의 운명을 결정한다"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('항우 — 천하 명장이었지만 범증 한 사람을 못 지킴',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('유방 — 평범한 농민이었지만 삼걸을 알아보고 함께함',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('차이가 운명을 갈랐다',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅹ. 현대 의의')
def s_modern_vessel(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 현대 의의', page, total)
    add_title(slide, '현대 ③ — "리더십은 능력이 아니라 그릇의 문제"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('"재능 있는 자는 천하에 많다 — 그러나 천하를 얻는 자는 드물다"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('차이는 그릇의 크기에 있다',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('• 칭찬을 들어도 자만하지 않는 그릇',
         {'font_size': 15, 'bold': True}),
        ('• 비판을 들어도 화내지 않는 그릇',
         {'font_size': 15, 'bold': True, 'space_before': 4}),
        ('• 위기에도 흔들리지 않는 그릇',
         {'font_size': 15, 'bold': True, 'space_before': 4}),
        ('• 공을 이루어도 물러설 줄 아는 그릇',
         {'font_size': 15, 'bold': True, 'space_before': 4, 'color': ACCENT}),
    ], line_spacing=1.3)


@S('Ⅹ. 현대 의의')
def s_modern_history(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 현대 의의', page, total)
    add_title(slide, '현대 ④ — "역사는 반복된다, 그래서 옛 사례에서 배운다"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"2,000년 전 한 고조와 항우의 차이는',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 오늘날 어떤 두 CEO의 차이와 다르지 않다"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('인간 본성과 권력의 역학은 변하지 않기 때문',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"자만하는 자는 망하고',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 사람을 잃는 자는 무너지고',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 자기를 못 다스리는 자는 추락하고',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 떠날 때를 모르는 자는 비참해진다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('— 이 패턴은 어느 시대 어느 조직에서도 똑같이 반복된다 —',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅹ. 현대 의의')
def s_modern_heart(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 현대 의의', page, total)
    add_title(slide, '현대 ⑤ — "리더십의 궁극은 사람의 마음"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"모든 권력의 시작과 끝은 사람의 마음"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('무력으로 얻을 수 있어도 마음으로만 지킬 수 있다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('권력의 진짜 원천',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('• 부하의 신뢰',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 동료의 존경',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('• 사회의 사랑',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('돈으로 살 수 없지만 모든 권력의 진짜 원천',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('이를 망각한 모든 권력자는 결국 무너졌다 — 역사는 한결같이 증언',
         {'font_size': 13, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅹ. 현대 의의')
def s_modern_everyone(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 현대 의의', page, total)
    add_title(slide, '현대 ⑥ — "우리 모두가 어떤 의미에서 리더이다"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '"제왕학은 황제만의 학문이 아니라 모두의 학문"',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('CEO·임원',     '회사 구성원에 대한 리더'),
        ('부서장·팀장',  '팀에 대한 리더'),
        ('선생님',       '학생에 대한 리더'),
        ('부모',         '자녀에 대한 리더'),
        ('선배',         '후배에 대한 리더'),
        ('자영업자',     '직원·고객에 대한 리더'),
        ('자기 자신',    '자기 인생에 대한 리더'),
    ]
    top = 3.2
    for tag, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(3.0), Inches(0.4),
                    tag, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(9.1), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.9), Inches(top + 0.1), Inches(8.8), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.55


# ============== Ⅺ. 마무리 ==============
@S('Ⅺ. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 제왕학')
    add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(11.9), Inches(5.0), PALE)
    add_paragraphs(slide, Inches(1.1), Inches(2.2), Inches(11.1), Inches(4.7), [
        ('강태공의 위수 낚시에서 진시황의 분서갱유까지',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('한 고조의 약법삼장에서 당 태종의 정관의 치까지',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('장량의 공성신퇴에서 한신의 토사구팽까지',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('— 2,000년 중국 흥망사의 결정적 인물 일화를 압축하여 —',
         {'font_size': 15, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 6}),
        ('"사람을 이끄는 자에게 진정 필요한 것은',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 12}),
        (' 자기 인격 · 인재 안목 · 비판 수용력 · 결단력의',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 8가지 자질"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('— 짧지만 무게 있는 동양 리더십의 정수 —',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
    ], line_spacing=1.2)


@S('Ⅺ. 마무리')
def s_key_message(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 마무리', page, total)
    add_title(slide, '가장 중요한 한 마디')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.6),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(4.3), [
        ('한 고조 유방이 자기 입으로 한 말',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"리더는 자기 능력으로 일하는 자가 아니라',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 4}),
        (' 사람의 능력을 끌어내는 자다"',
         {'font_size': 24, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('', {'font_size': 8}),
        ('모든 동양 제왕학이 결국 이 한 줄로 모인다',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 14}),
        ('', {'font_size': 4}),
        ('"좋은 사람이 좋은 리더가 된다"',
         {'font_size': 20, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('— 2,000년 흥망사가 가르치는 단 하나의 진실 —',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅺ. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.5),
                '修 身 治 國',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '수 신 치 국', font_size=24, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.7),
                '"수신(修身)이 곧 치국(治國)의 근본이다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '— 동양 제왕학 2,000년이 가르치는 단 하나의 진실 —',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '사람은 누구나 어떤 의미에서 리더 — 그래서 제왕학은 모두의 학문',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
                '감사합니다', font_size=26, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\제왕학.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')