# -*- coding: utf-8 -*-
"""
법구경(法句經, Dhammapada) 발표자료 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '진리의 길 · 마음의 경전 · 동방의 성서',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '法 句 經',
                font_size=130, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '담 마 빠 다  (Dhammapada)',
                font_size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '心 爲 法 本   心 尊 心 使',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '마음은 모든 법의 근본이요, 마음이 주(主)가 되어 모든 일을 시킨다  — 제1게송',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '편자 법구(法救) · 빨리어 26품 423게송 / 한역 39품 752게송',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 법구경은 어떤 책인가'),
        ('Ⅱ.', '서명과 편자'),
        ('Ⅲ.', '성립과 전승의 두 갈래'),
        ('Ⅳ.', '26품의 5단계 구조'),
        ('Ⅴ.', '핵심 게송 16선'),
    ]
    items_right = [
        ('Ⅵ.', '7대 핵심 사상'),
        ('Ⅶ.', '왜 동방의 성서인가'),
        ('Ⅷ.', '한국과 동아시아 수용'),
        ('Ⅸ.', '오늘 다시 펼치는 이유'),
        ('Ⅹ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 법구경', '부처의 직설(直說)에 가장 가까운 게송 모음')
    rows = [
        ('원제', 'Dhammapada (담마빠다) — 진리의 길'),
        ('편자', '법구(法救, Dharmatrāta) — 부처 열반 후 약 300년 인도 학자'),
        ('빨리어본', '26품 423게송 — 표준본'),
        ('한역본', '39품 752게송 — 오나라 유기난 등 한역'),
        ('성격', '시(詩) 형식의 짧은 가르침 모음 — 산문 아닌 게송'),
        ('위상', '불교 경전 중 가장 오래되고 가장 사랑받는 책'),
        ('정신', '마음이 모든 것의 근본 — 心爲法本'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.8), Inches(0.45), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.8), Inches(0.45),
                    k, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.7), y, Inches(10.3), Inches(0.45),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_essence(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '법구경의 정수', '한 줄을 외워 평생 사색하는 경전')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.4),
                '心 爲 法 本   心 尊 心 使',
                font_size=32, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.5),
                '마음은 모든 법의 근본이요, 마음이 주(主)가 되어 모든 일을 시킨다',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.7), Inches(12.0), Inches(2.5), [
        ('● 짧고 강렬한 시 — 한 게송이 평생의 사색 거리', {'font_size': 16, 'space_before': 6}),
        ('● 일상의 비유 — 꽃·코끼리·수레·그림자·강·등불', {'font_size': 16, 'space_before': 6}),
        ('● 종교를 초월한 보편 지혜 — 신자 아니어도 자기계발서로', {'font_size': 16, 'space_before': 6}),
        ('● 「**마음을 다스리는 법**」의 가장 오래된 시집', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC1)
def i_translation(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '슬기바다 11권 — 한명숙 옮김본',
              '깨달음과 그 실천의 울림')
    rows = [
        ('시리즈', '동양고전 슬기바다 11권 (홍익출판사)'),
        ('옮긴이', '한명숙 — 동양고전 전문 번역가'),
        ('출판일', '2005년 4월 11일'),
        ('형태', '양장본 412쪽 — 시리즈 중 분량이 큰 편'),
        ('한국의 다른 번역', '법정 『진리의 말씀』(1984) · 전재성 · 일아 · 김서리'),
        ('의의', '학술 번역이 아닌 — 일반 독자가 한 게송씩 음미하는 형식'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.0), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅱ. 서명과 편자 ==============
SEC2 = 'Ⅱ. 서명과 편자'

@S(SEC2)
def ii_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, 'Dhammapada — 「진리의 길」',
              '서명 자체가 책의 본질')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● Dhamma(담마) — 법(法), 진리, 부처의 가르침', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● Pada(빠다) — 길, 발자취, 말씀, 시구(詩句)', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 합치면 — 「**진리의 길**」, 「**법의 말씀**」, 「**가르침의 짧은 구절들**」', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 한역(漢譯) — Dhamma=法, Pada=句 → 「**법구경(法句經)**」', {'font_size': 16, 'space_before': 10}),
        ('● 「**진리의 짧은 구절을 모은 경전**」 — 이름이 곧 본질', {'font_size': 16, 'color': SUB, 'space_before': 10}),
        ('● 긴 설법이 아니라 — 한 줄, 한 게송으로 가슴을 치는 말씀들', {'font_size': 16, 'color': SUB, 'space_before': 6}),
    ])


@S(SEC2)
def ii_editor(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '편자 법구(다르마트라타)',
              '부처 열반 후 약 300년 — 흩어진 게송을 모으다')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 다르마트라타(Dharmatrāta) — 한역명 「**법구**」', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   「법(다르마)을 구원하는 자」라는 뜻의 이름', {'font_size': 14, 'color': SUB}),
        ('● 부처(BC 5세기) 열반 후 약 300년경 — BC 3~2세기 인도 학자', {'font_size': 17, 'space_before': 10}),
        ('   그가 직접 부처의 말씀을 들은 것은 아님', {'font_size': 14, 'color': SUB}),
        ('● 여러 경로로 전해진 부처의 게송들을 — 모으고 분류', {'font_size': 17, 'space_before': 10}),
        ('● 「**편자**」 — 저자가 아니라 — **모으고 정리한 사람**', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 게송 자체는 — 부처의 직설(直說)에 가장 가까운 형식', {'font_size': 16, 'color': SUB, 'space_before': 6}),
    ])


# ============== Ⅲ. 성립과 전승 ==============
SEC3 = 'Ⅲ. 성립과 전승'

@S(SEC3)
def iii_two_streams(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '두 갈래의 전승',
              '남방 빨리어 vs 북방 산스크리트')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.8), [
        ('● 남방 상좌부(上座部) 전승 — 빨리어 담마빠다', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('   26품 423게송 — 가장 권위 있는 표준본', {'font_size': 14, 'color': SUB}),
        ('   스리랑카·미얀마·태국 등에서 전승', {'font_size': 14, 'color': SUB}),
        ('● 북방 설일체유부(說一切有部) 전승 — 산스크리트 우다나바르가', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('   33품 약 1,000게송 — 빨리어본보다 확장', {'font_size': 14, 'color': SUB}),
        ('   인도 북부·중앙아시아·티베트에서 전승', {'font_size': 14, 'color': SUB}),
        ('● 한역 법구경은 — 후자 계통의 모태', {'font_size': 16, 'color': ACCENT, 'space_before': 12}),
    ])


@S(SEC3)
def iii_han_translation(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '한역(漢譯)의 역사',
              '삼국 오나라에서 시작된 동아시아 전승')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 삼국시대 오(吳)나라(229~280)', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('   유기난(維祇難)·축장염(竺將炎)·축율염(竺律炎) 공동 한역', {'font_size': 14, 'color': SUB}),
        ('● 한역본 — 빨리어본 26품 + 13품 추가 = **39품 752게송**', {'font_size': 17, 'space_before': 10}),
        ('● 고려대장경(팔만대장경) 수록 — 한국 정착', {'font_size': 17, 'space_before': 10}),
        ('● 한·중·일·베트남 — 모두 한역본을 통해 법구경 수용', {'font_size': 17, 'space_before': 10}),
        ('● 빨리어 원전은 19세기 막스 뮬러 영역(英譯) 이후 서양에 본격 알려짐', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


# ============== Ⅳ. 26품 구조 ==============
SEC4 = 'Ⅳ. 26품 구조'

@S(SEC4)
def iv_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '26품의 큰 흐름',
              '마음의 본질에서 깨달음의 길까지')
    rows = [
        ('1단계 (1~3품)', '마음의 본질 — 쌍서·불방일·심'),
        ('2단계 (4~7품)', '두 모습 — 꽃·우암·명철·나한'),
        ('3단계 (8~11품)', '행위의 결과 — 술천·악행·도장·노모'),
        ('4단계 (12~18품)', '처세와 행복 — 애신·세속·술불·안녕·호희·분노·진구'),
        ('5단계 (19~26품)', '해탈의 길 — 봉지·도행·...·범지'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.6 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.2), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(3.2), Inches(0.55),
                    k, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.1), y, Inches(8.8), Inches(0.55),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5),
                '마음의 본질부터 깨달음의 길까지 — 단계적 안내',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_chapters(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '26품 한 폭으로',
              '각 품의 핵심 주제')
    left = [
        '1. 쌍서품 (Yamaka) — 마음에서',
        '2. 불방일품 — 깨어 있음',
        '3. 심품 — 마음 다스림',
        '4. 화품 — 꽃의 비유',
        '5. 우암품 — 어리석은 자',
        '6. 명철품 — 지혜로운 자',
        '7. 나한품 — 깨달은 자',
        '8. 술천품 — 수의 비유',
        '9. 악행품 — 악업의 결과',
        '10. 도장품 — 비폭력',
        '11. 노모품 — 무상함과 늙음',
        '12. 애신품 — 자기 책임',
        '13. 세속품 — 세상을 보는 눈',
    ]
    right = [
        '14. 술불품 — 부처의 위대함',
        '15. 안녕품 — 참된 행복',
        '16. 호희품 — 집착의 위험',
        '17. 분노품 — 분노 다스림',
        '18. 진구품 — 마음의 때',
        '19. 봉지품 — 법대로 사는 자',
        '20. 도행품 — 팔정도',
        '21. 광연품 — 다양한 가르침',
        '22. 지옥품 — 지옥의 길',
        '23. 상유품 — 코끼리의 비유',
        '24. 애욕품 — 갈애의 끊음',
        '25. 사문품 — 수행자의 길',
        '26. 범지품 — 진정한 거룩한 자',
    ]
    for i, txt in enumerate(left):
        add_textbox(slide, Inches(0.7), Inches(2.4 + i * 0.36), Inches(6.0), Inches(0.32),
                    txt, font_size=13, color=INK)
    for i, txt in enumerate(right):
        add_textbox(slide, Inches(7.0), Inches(2.4 + i * 0.36), Inches(6.0), Inches(0.32),
                    txt, font_size=13, color=INK)


@S(SEC4)
def iv_first_chapter(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '첫 품 「쌍서품」의 첫 게송',
              '법구경 전체의 머리 — 모든 메시지의 압축')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(3.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.6), Inches(2.8), [
        ('마음은 모든 일의 근본이다.', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('마음이 주(主)가 되어 모든 일을 시키나니,', {'font_size': 16, 'color': INK}),
        ('나쁜 마음으로 말하거나 행동하면', {'font_size': 16, 'color': INK}),
        ('괴로움이 그를 따른다, 수레바퀴가 소를 따르듯이.', {'font_size': 16, 'color': INK, 'space_before': 4}),
        ('', {'font_size': 8}),
        ('선한 마음으로 말하거나 행동하면', {'font_size': 16, 'color': INK}),
        ('즐거움이 그를 따른다, 그림자가 그를 따르듯이.', {'font_size': 16, 'color': INK}),
    ])
    add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.5),
                '— 법구경 제1, 2 게송 (쌍서품)',
                font_size=13, color=SUB, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.8),
                '이 한 게송 안에 — 마음·행위·결과의 모든 가르침이 들어 있다',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅴ. 핵심 게송 16선 ==============
SEC5 = 'Ⅴ. 핵심 게송 16선'

def make_gatha_slide(section, idx_total, hanja, korean, comment):
    @S(section)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, section, n, t)
        add_title(slide, f'게송 {idx_total}', f'한 줄을 평생 사색하는 시(詩)')
        add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.4), PALE)
        add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.4),
                    hanja,
                    font_size=24, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(1.4),
                    korean,
                    font_size=18, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.4), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.5), Inches(11.6), Inches(1.4),
                    comment,
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


gathas = [
    ('1/16', '心 爲 法 本   心 尊 心 使',
     '마음은 모든 법의 근본이요, 마음이 주가 되어 모든 일을 시킨다',
     '제1게송. 법구경 전체의 머리. 모든 행위·말·결과의 근원은 마음.'),
    ('2/16', '不 放 逸 不 死   放 逸 是 死 路',
     '불방일(不放逸)은 죽지 않는 길, 방일(放逸)은 죽음의 길',
     '제21게송 불방일품. 깨어 있음은 살아 있음, 게으름은 죽음.'),
    ('3/16', '心 動 難 攝   輕 躁 難 持',
     '마음은 움직임이 빠르고 다잡기 어렵다, 가벼이 들떠 붙잡기 어렵다',
     '심품. 마음의 변덕스러운 본성 — 다스리기 어렵지만 다스려야 한다.'),
    ('4/16', '無 病 最 利   知 足 最 富',
     '병 없음이 최고의 이익, 만족할 줄 앎이 최고의 부유함',
     '안녕품. 건강과 만족 — 모든 부귀보다 깊은 행복.'),
    ('5/16', '從 親 生 憂   從 親 生 怖',
     '사랑하는 데서 근심이 생기고, 사랑하는 데서 두려움이 생긴다',
     '호희품. 집착이 곧 고통의 뿌리 — 사랑 자체가 아니라 집착하는 사랑.'),
    ('6/16', '勝 怒 從 慈   勝 惡 從 善',
     '분노를 이기는 것은 자비, 악을 이기는 것은 선',
     '분노품. 분노에 분노로 갚지 말라 — 자비만이 분노를 그치게 한다.'),
    ('7/16', '心 為 法 本   行 為 結 果',
     '마음이 원인이요, 행위가 결과다 — 모든 업의 도리',
     '쌍서·심품 통합. 마음의 의도가 업, 업이 운명.'),
    ('8/16', '欲 河 難 渡   愛 海 難 越',
     '욕망의 강은 건너기 어렵고, 갈애의 바다는 넘기 어렵다',
     '애욕품. 모든 괴로움의 뿌리는 갈애(taṇhā).'),
    ('9/16', '一 切 皆 苦   一 切 無 我',
     '모든 것은 괴로움이요, 모든 것은 무아(無我)다',
     '봉지품. 삼법인의 정수 — 苦·無常·無我.'),
    ('10/16', '心 是 諸 法 本',
     '마음이 모든 법의 근본이다',
     '쌍서·심품. 마음을 보는 것이 모든 수행의 시작.'),
    ('11/16', '勿 以 善 小 而 不 為',
     '선이 작다고 하지 않지 말라 — 작은 선이 큰 선이 된다',
     '술천품. 한 방울의 물이 항아리를 채우듯, 작은 선이 큰 선이 된다.'),
    ('12/16', '己 為 己 主   己 為 己 護',
     '자기가 자기의 주인이요, 자기가 자기의 보호자다',
     '애신품. 자기 책임의 정수 — 누구도 나를 대신 구원하지 않는다.'),
    ('13/16', '生 而 不 殺   不 起 殺 念',
     '살아 있는 것을 죽이지 말고 — 죽이려는 마음도 일으키지 말라',
     '도장품. 모든 생명은 죽음을 두려워한다 — 자기를 헤아려 남을.'),
    ('14/16', '猶 如 老 牛   人 老 亦 然',
     '늙은 소처럼 — 사람도 늙는다',
     '노모품. 무상함의 자각 — 누구나 늙고 병들고 죽는다.'),
    ('15/16', '非 因 出 身   行 為 為 梵',
     '태어남이 아니라 행위로 바라문이 된다',
     '범지품. 카스트 제도에 대한 부처의 혁명적 선언.'),
    ('16/16', '一 句 之 法   勝 千 無 義',
     '의미 있는 한 게송이 — 의미 없는 천 마디보다 낫다',
     '술천품. 압축의 가치 — 짧고 진실한 한 마디.'),
]

for tag, hj, kr, cm in gathas:
    make_gatha_slide(SEC5, tag, hj, kr, cm)


# ============== Ⅵ. 7대 핵심 사상 ==============
SEC6 = 'Ⅵ. 7대 핵심 사상'

@S(SEC6)
def vi_mind(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 1 — 모든 것은 마음에서',
              '心爲法本 — 법구경 전체의 머리')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 마음 = 모든 행위의 근원', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● 행위 = 결과를 가져옴', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 결과 = 그림자처럼 따라옴 (선 → 즐거움, 악 → 괴로움)', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 「**마음을 보라, 거기서 모든 것이 시작된다**」', {'font_size': 17, 'color': ACCENT, 'space_before': 12}),
        ('● 마음을 다스리지 못하면 — 어떤 노력도 어긋난다', {'font_size': 16, 'color': SUB, 'space_before': 8}),
    ])


@S(SEC6)
def vi_karma(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 2 — 업(業, Karma)의 법칙',
              '행한 대로 받는 인과의 절대성')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 행한 대로 받는다 — 인과의 절대성', {'font_size': 18, 'space_before': 6}),
        ('● 그러나 업의 핵심은 「**의도(意圖)**」', {'font_size': 18, 'space_before': 10}),
        ('● 부처: 「**내가 업이라 하는 것은 의도이다**」', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 마음의 의도가 곧 업, 업이 곧 운명', {'font_size': 18, 'space_before': 10}),
        ('● 「**자기 책임의 선언**」 — 누구도 대신 구원하지 않는다', {'font_size': 17, 'color': ACCENT, 'space_before': 12}),
    ])


@S(SEC6)
def vi_awakening(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 3 — 불방일(不放逸)의 정신',
              '깨어 있음이 곧 생명')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 깨어 있음(覺)이 곧 생명', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● 방일(放逸, 게으름)은 죽음의 길', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 「**불방일은 죽지 않는 길, 방일은 죽음의 길**」 (제21게송)', {'font_size': 17, 'space_before': 10}),
        ('● 매 순간 의식적으로 살라 — 관성에 끌려가지 말라', {'font_size': 17, 'color': ACCENT, 'space_before': 12}),
        ('● 오늘 — 마음 챙김(mindfulness)의 가장 오래된 정수', {'font_size': 16, 'color': SUB, 'space_before': 8}),
    ])


@S(SEC6)
def vi_two_paths(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 4 — 두 가지 길의 분별',
              '어리석은 자(우암) vs 지혜로운 자(명철)')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 어리석은 자 — 욕망에 끌려가는 자', {'font_size': 17, 'space_before': 6}),
        ('● 지혜로운 자 — 자기 마음의 주인이 되는 자', {'font_size': 17, 'space_before': 10}),
        ('● 매 순간 우리가 두 길 중 하나를 선택한다는 자각', {'font_size': 17, 'space_before': 10}),
        ('● 「**선악**」이 아닌 「**지혜와 어리석음**」의 분별', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('● 종교적 선악 판단이 아닌 — 「**어떤 길이 행복으로 이끄는가**」', {'font_size': 16, 'color': SUB, 'space_before': 8}),
    ])


@S(SEC6)
def vi_compassion(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 5 — 자비와 비폭력',
              '모든 생명은 죽음을 두려워한다')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 모든 생명은 죽음을 두려워한다 — 자기를 헤아려 남을 헤아려라', {'font_size': 17, 'space_before': 6}),
        ('● 분노에 분노로 갚지 말라', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 분노는 사랑으로만 그친다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 진정한 거룩한 자(범지)는 — 폭력을 떠난 자', {'font_size': 17, 'space_before': 10}),
        ('● 「**비폭력(Ahiṃsā)**」 — 간디·마틴 루서 킹의 동양적 원천', {'font_size': 16, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC6)
def vi_craving(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 6 — 갈애의 끊음 = 해탈',
              '모든 괴로움의 뿌리')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 모든 괴로움의 뿌리는 갈애(渴愛, taṇhā)', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('   더 가지려는 마음, 잃지 않으려는 마음 — 곧 고통', {'font_size': 14, 'color': SUB}),
        ('● 갈애를 뽑으면 — 평안에 이른다', {'font_size': 18, 'space_before': 10}),
        ('● 집착하지 않음이 곧 자유', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 욕망을 「**부정**」하라가 아니라 — 욕망에 「**끌려가지 말라**」', {'font_size': 17, 'color': ACCENT, 'space_before': 12}),
        ('● 욕망을 가진 채 자유로울 수 있는 길 — 그게 갈애 끊음', {'font_size': 16, 'color': SUB, 'space_before': 8}),
    ])


@S(SEC6)
def vi_brahmin(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 7 — 진정한 거룩한 자',
              '출신이 아닌 행위로 결정된다')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 출신·계급·복장이 아니라 — 마음의 청정함이 거룩함의 기준', {'font_size': 17, 'space_before': 6}),
        ('● 「**태어남이 아니라 행위로 바라문이 된다**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 2,500년 전 인도 카스트 사회에서 부처가 외친 가장 혁명적 선언', {'font_size': 17, 'space_before': 10}),
        ('● 「**진정한 고귀함은 — 자기를 다스린 자**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 오늘 — 「**평등**」 사상의 가장 오래된 종교적 정식', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


# ============== Ⅶ. 동방의 성서 ==============
SEC7 = 'Ⅶ. 동방의 성서'

@S(SEC7)
def vii_eastern_bible(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '왜 「동방의 성서」라 불리는가',
              '종교를 초월한 보편 지혜')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 종교적 차이를 초월한 보편 지혜', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   불교 신자가 아니어도 — 윤리·심리학·자기계발서로 읽힌다', {'font_size': 14, 'color': SUB}),
        ('● 어느 게송 하나만 골라도 — 평생 사색할 가치', {'font_size': 17, 'space_before': 10}),
        ('● 톨스토이·헤르만 헤세·하이데거 — 서양 사상가들이 깊이 영향', {'font_size': 17, 'space_before': 10}),
        ('● 19세기 막스 뮬러의 영역 이후 — 서양에서 가장 많이 읽힌 불교 경전', {'font_size': 17, 'space_before': 10}),
        ('● 「**마음을 다스리는 법**」 — 종교 초월한 인류 공통 과제', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC7)
def vii_appeal(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '법구경의 형식적 매력',
              '왜 2,500년 사랑받는가')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 짧고 강렬한 시(詩) 형식', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   외우기 쉽고 — 기억에 남는다', {'font_size': 14, 'color': SUB}),
        ('● 누구나 이해할 수 있는 평이함 — 농부도 학자도', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   추상 형이상학이 아닌 — 일상의 비유(꽃·코끼리·수레·강)', {'font_size': 14, 'color': SUB}),
        ('● 부처의 인격적 매력 — 위협 아닌 권유', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   「이 길로 가라」가 아니라 — 「이 길이 행복한 길이다」', {'font_size': 14, 'color': SUB}),
    ])


# ============== Ⅷ. 한국·동아시아 수용 ==============
SEC8 = 'Ⅷ. 한국 수용'

@S(SEC8)
def viii_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '한국에서의 위상',
              '가장 사랑받는 불교 경전')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 한국 불교에서 — 반야심경·금강경과 함께 가장 널리 읽힘', {'font_size': 17, 'space_before': 6}),
        ('● 신자뿐 아니라 일반 교양서로도 폭넓게 읽힘', {'font_size': 17, 'space_before': 10}),
        ('● 법정 스님 『**진리의 말씀**』(1984) — 가장 유명한 한국어 번역', {'font_size': 17, 'space_before': 10}),
        ('● 전재성·일아·김서리·한명숙 등 다수 번역본', {'font_size': 17, 'space_before': 10}),
        ('● 한국 시인·소설가들의 정신적 자양분 (한용운·법정 등)', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC8)
def viii_modern(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '현대의 부활',
              '명상·마음 챙김·자기계발의 가장 깊은 원천')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 현대 명상·마음 챙김(mindfulness) 운동의 동양적 원천', {'font_size': 17, 'space_before': 6}),
        ('● 「**마음을 다스린다**」 — 자기계발·심리학의 가장 깊은 정수', {'font_size': 17, 'space_before': 10}),
        ('● 「**갈애를 끊어라**」 — 미니멀리즘·미디어 다이어트의 정신', {'font_size': 17, 'space_before': 10}),
        ('● 「**분노는 자비로**」 — 비폭력 갈등 해결의 정수', {'font_size': 17, 'space_before': 10}),
        ('● 「**한 게송 하루**」 — 일상 명상 텍스트로 가장 적합', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅸ. 오늘 다시 펼치는 이유 ==============
SEC9 = 'Ⅸ. 오늘 다시 펼치는 이유'

@S(SEC9)
def ix_today1(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '1 — 마음을 보라',
              '거기서 모든 것이 시작된다')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 제1게송이 책 전체를 말한다', {'font_size': 17, 'space_before': 6}),
        ('● 모든 행위·말·결과의 근원은 — 마음', {'font_size': 17, 'space_before': 10}),
        ('● 외부 환경을 바꾸기 전에 — 자기 마음을 본다', {'font_size': 17, 'space_before': 10}),
        ('● 마음을 다스리지 못하면 — 어떤 노력도 어긋난다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 「**오늘 내 마음은 어디에 있는가?**」 — 매일 묻는 질문', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today2(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '2 — 한 줄을 평생 사색하라',
              '법구경의 가장 좋은 읽기 방법')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 법구경은 — 빨리 읽고 끝내는 책이 아니다', {'font_size': 17, 'space_before': 6}),
        ('● 한 게송을 가슴에 품고 — 그 의미를 살아내는 책', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 「**하루에 한 게송**」 — 1년이면 365게송, 평생이면 평생', {'font_size': 17, 'space_before': 10}),
        ('● 한 줄의 시(詩)가 — 자기 일상의 거울이 된다', {'font_size': 17, 'space_before': 10}),
        ('● 「**느리게 읽기(slow reading)**」의 가장 오래된 모델', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today3(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '3 — 갈애를 끊어 자유로워라',
              '오늘 가장 절실한 가르침')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 더 가지려는 마음 — 끝없는 욕망의 시대', {'font_size': 17, 'space_before': 6}),
        ('● SNS·광고·소비주의가 갈애를 부채질', {'font_size': 17, 'space_before': 10}),
        ('● 법구경의 답 — 「**욕망에 끌려가지 말라**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 「**무병이 최고의 이익, 지족이 최고의 부유**」', {'font_size': 17, 'space_before': 10}),
        ('● 미니멀리즘·디지털 디톡스의 가장 깊은 동양적 원천', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


# ============== Ⅹ. 마무리 ==============
SEC10 = 'Ⅹ. 마무리'

@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '법구경이 일러주는 7가지',
              '한 폭으로 정리')
    items = [
        '마음을 보라 — 거기서 모든 것이 시작된다',
        '오늘 한 행위가 곧 내일의 자기 — 업의 자기 책임',
        '깨어 있음이 살아 있음 — 불방일',
        '어리석은 자가 아닌 지혜로운 자의 길',
        '분노는 자비로만 그친다 — 비폭력',
        '갈애를 끊으면 자유로워진다 — 집착 없음',
        '진정한 고귀함은 출신이 아닌 행위 — 평등',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.5 + i * 0.6)
        add_textbox(slide, Inches(0.9), y, Inches(0.6), Inches(0.5),
                    f'{i+1}.', font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.5),
                    txt, font_size=16, color=INK)


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                '법구경 — 부처의 직설에 가장 가까운 진리의 길',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '心 爲 法 本',
                font_size=140, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
                '마음이 모든 법의 근본 — 제1게송의 정수',
                font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.8), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
                '한 게송을 가슴에 품고 평생 사색하라 — 동방의 성서',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)


# ============== 빌드 ==============
total = len(SLIDES)
for i, (fn, sec) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    fn(slide, i, total)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\법구경.pptx'
prs.save(out_path)
print(f'생성 완료: {out_path}  슬라이드 수: {total}')
