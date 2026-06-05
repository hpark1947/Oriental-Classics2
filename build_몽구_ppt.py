# -*- coding: utf-8 -*-
"""
몽구(蒙求) 발표자료 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '동아시아 최초의 어린이 역사·인물 학습서 · 1,300년 동아시아 교양의 원천',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '蒙 求',
                font_size=130, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '몽 구',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '匪 我 求 童 蒙  童 蒙 求 我',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '내가 어린아이를 찾아가는 게 아니라, 어린아이가 나를 찾아온다  — 『주역』 몽괘',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '이한(李翰, 당 천보 5년 746년경) · 596구 2,384자 · 4언 운문 대구의 600여 인물 일화',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 몽구는 어떤 책인가'),
        ('Ⅱ.', '서명의 유래 — 주역 몽괘에서'),
        ('Ⅲ.', '저자 이한과 편찬 배경'),
        ('Ⅳ.', '8자 1구 4자 대구의 구조'),
        ('Ⅴ.', '명일화 16선'),
    ]
    items_right = [
        ('Ⅵ.', '5대 교육 사상'),
        ('Ⅶ.', '후세 사자성어의 원전'),
        ('Ⅷ.', '한국·일본 수용사'),
        ('Ⅸ.', '오늘 다시 펼치는 이유'),
        ('Ⅹ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 몽구', '동아시아 최초의 본격적 어린이 역사 학습서')
    rows = [
        ('서명', '蒙求 (몽구) — 「어리석은 어린이가 배움을 구한다」'),
        ('저자', '이한(李翰) — 당 천보 5년(746) 무렵 한림학사'),
        ('분량', '596구 2,384자 · 4자 1구의 대구로 600여 인물 일화'),
        ('시대 범위', '춘추전국 ~ 남북조 (약 1,200년)'),
        ('성격', '어린이용 역사·인물 학습서 · 운문 + 대구 + 일화'),
        ('위상', '천자문 다음, 삼자경 이전 — 동아시아 어린이 교양 원형'),
        ('정신', '童蒙求我 — 자발적 배움 · 인물 중심 역사관'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.8), Inches(0.45), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(1.8), Inches(0.45),
                    k, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.7), y, Inches(10.3), Inches(0.45),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_position(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '동아시아 어린이 교재의 가교',
              '천자문 → 몽구 → 삼자경의 흐름에서 인물·역사의 입문서')
    add_filled_rect(slide, Inches(0.7), Inches(2.4), Inches(12.0), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.4), Inches(12.0), Inches(1.4),
                '천자문(6세기) → 몽구(8세기) → 백가성·삼자경(13세기)',
                font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 천자문 — 한자 1,000자의 입문', {'font_size': 16, 'bold': True, 'color': ACCENT}),
        ('   글자를 익히는 단계 — 의미보다 문자 자체', {'font_size': 14, 'color': SUB}),
        ('● 몽구 — 한자로 인물·역사·도덕을 학습', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   600여 인물의 결정적 한 장면을 4자로 압축 — 외우면서 역사 인식 형성', {'font_size': 14, 'color': SUB}),
        ('● 삼자경·백가성 — 도덕 명제·성씨의 정리', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   몽구가 정착시킨 「운문 + 일화」의 형식이 후대에 계승', {'font_size': 14, 'color': SUB}),
    ])


@S(SEC1)
def i_essence(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '몽구의 핵심 — 한 마디로',
              '「사람의 한 행동이 그를 영원히 정의한다」')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.6), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.6),
                '匡 衡 鑿 壁  孫 敬 閉 戶',
                font_size=32, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(0.5),
                '광형은 벽을 뚫어 글을 읽었고, 손경은 문을 닫아걸고 학문에 몰두했다',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.9), Inches(12.0), Inches(2.5), [
        ('● 600여 인물 × 결정적 한 장면 — 4자로 압축', {'font_size': 16, 'space_before': 6}),
        ('● 두 인물을 짝지어 한 단위 — 비교와 종합 사고', {'font_size': 16, 'space_before': 6}),
        ('● 운율과 대구로 어린이가 외우기 쉽게 — 평생 입에 남음', {'font_size': 16, 'space_before': 6}),
        ('● 한 행동이 그 사람을 영원히 정의한다 — 자기 4자에 대한 자각', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC1)
def i_translation(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '슬기바다 7권 — 유동환 옮김본',
              '한국에서의 본격 일반 독자용 번역')
    rows = [
        ('시리즈', '동양고전 슬기바다 7권 (홍익출판사)'),
        ('옮긴이', '유동환'),
        ('출판일', '2005년 4월 18일'),
        ('형태', '양장본 328쪽 · 152×223mm'),
        ('편집 원칙', '원문 · 독음 · 번역 · 해설을 함께 제공'),
        ('의의', '일반 독자가 596구 전체를 따라 읽을 수 있도록 정리한 본격 한국어 번역본'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.0), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅱ. 서명의 유래 ==============
SEC2 = 'Ⅱ. 서명의 유래'

@S(SEC2)
def ii_meng_gua(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '주역 몽괘에서 따온 이름',
              '책 이름 자체가 교육 철학의 선언')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.5),
                '匪 我 求 童 蒙   童 蒙 求 我',
                font_size=30, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.5),
                '내가 어린아이를 찾아가는 게 아니라, 어린아이가 나를 찾아온다 — 『주역』 몽괘 단사',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.7), Inches(12.0), Inches(2.5), [
        ('● 蒙(몽) — 어리석음·어두움·아직 깨어나지 않은 어린이', {'font_size': 16, 'space_before': 6}),
        ('● 求(구) — 구하다·청하다·찾다', {'font_size': 16, 'space_before': 6}),
        ('● 「몽구」 — 어린이가 스승을 찾아 배움을 청한다', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 책 이름 자체가 — 자발적 배움의 유가적 교육관 선언', {'font_size': 15, 'color': SUB, 'space_before': 6}),
    ])


@S(SEC2)
def ii_voluntary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '자발적 배움(童蒙求我)의 정신',
              '강요가 아니라 호기심에서 시작하는 학습')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 가르침은 일방적 주입이 아니다', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   배우는 자가 스스로 동기를 갖고 청해야 한다 — 동양 교육관의 정수', {'font_size': 14, 'color': SUB}),
        ('● 운율과 대구라는 형식의 매력', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   어린이가 스스로 외우고 싶어지도록 — 형식이 곧 동기', {'font_size': 14, 'color': SUB}),
        ('● 인물 일화의 흥미', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   추상 도덕이 아닌 구체 사람 — 아이의 호기심에 직접 호소', {'font_size': 14, 'color': SUB}),
        ('● 두 인물 짝짓기의 비교 사고', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   왜 둘이 함께 거론되는가? — 단순 암기를 넘어 종합 사고로', {'font_size': 14, 'color': SUB}),
    ])


# ============== Ⅲ. 저자와 편찬 ==============
SEC3 = 'Ⅲ. 저자와 편찬'

@S(SEC3)
def iii_li_han(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '저자 이한(李翰) — 당 한림학사',
              '당 대종 시대 최고 문장가의 어린이 교과서')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 이한(李翰) — 당 천보(天寶) 연간 한림학사', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   한림학사 = 황제 조서 작성과 학문을 담당한 최고위 문신', {'font_size': 14, 'color': SUB}),
        ('   일부 문헌은 「이한(李瀚)」 — 청대 고증으로 동일 인물 확인', {'font_size': 14, 'color': SUB}),
        ('● 편찬 시기 — 천보 5년(746) 무렵', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   학계 일부는 758~765년 사이로 보기도 함', {'font_size': 14, 'color': SUB}),
        ('● 황제 헌상 — 요주자사 이량(李良)의 「**천몽구표(薦蒙求表)**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   이로써 천하에 알려짐 — 국가 공인 어린이 교과서로 격상', {'font_size': 14, 'color': SUB}),
        ('● 결과 — 중국 역사상 가장 오래·가장 널리 읽힌 어린이 교재', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC3)
def iii_intent(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '편찬 의도',
              '왜 600여 인물을 8자 한 단위로 묶었는가')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 흩어진 역사 인물 일화를 효율적으로 암기시키기 위함', {'font_size': 17, 'space_before': 6}),
        ('● 운율과 대구의 형식 미학에 기대어 — 어린이가 스스로 외우고 싶게', {'font_size': 17, 'space_before': 10}),
        ('● 한 단위(8자)에 두 인물 — 비슷하거나 대조되는 행적을 함께', {'font_size': 17, 'space_before': 10}),
        ('● 1,200년의 중국사를 600여 인물로 압축 — 「압축된 역사 백과」', {'font_size': 17, 'space_before': 10}),
        ('● 추상 도덕이 아닌 구체 사례로 효·근면·충절·지혜를 가르침', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅳ. 구조 ==============
SEC4 = 'Ⅳ. 8자 1구 구조'

@S(SEC4)
def iv_form(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '경이로운 압축의 형식',
              '8자 한 단위 = 4자 + 4자, 두 인물 일화의 짝')
    rows = [
        ('총 글자 수', '2,384자 (이본은 2,484자 설)'),
        ('구절 수', '596구 (4자 1구 기준)'),
        ('한 단위', '8자 = 4자 + 4자 — 두 인물 일화'),
        ('운율', '4언 운문 · 짝수 구절마다 운을 맞춤'),
        ('압운 변화', '8구마다 환운(換韻) — 단조롭지 않게'),
        ('구성', '인물 이름 + 그를 정의하는 결정적 한 장면'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.2), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.2), Inches(0.5),
                    k, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.1), y, Inches(9.8), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_first_lines(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '첫머리의 8자 한 단위들',
              '몽구가 어떻게 시작되는지 — 두 인물의 짝')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.9), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.9),
                '王 戎 簡 要   裴 楷 淸 通',
                font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.3), Inches(12.0), Inches(0.4),
                '왕융은 간결하고 요령 있었으며, 배해는 청아하고 통달하였다 — 몽구 첫 구절',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.9), PALE)
    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.9),
                '孔 明 臥 龍   呂 望 非 熊',
                font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(4.9), Inches(12.0), Inches(0.4),
                '공명(제갈량)은 누워 있는 용, 여망(강태공)은 곰이 아니었다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.9), PALE)
    add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.9),
                '楊 震 關 西   丁 寬 易 東',
                font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4),
                '양진은 관서의 공자라 불렸고, 정관은 동쪽으로 역(易)을 전했다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC4)
def iv_pairing(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '두 인물 짝짓기의 깊이',
              '비교 사고를 통한 종합적 역사 인식')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 비슷한 덕목의 두 인물', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   광형착벽 + 손경폐호 — 둘 다 「**가난 속의 면학**」', {'font_size': 14, 'color': SUB}),
        ('   차윤낭형 + 손강영설 — 둘 다 「**빛이 없는 가운데 책 읽기**」', {'font_size': 14, 'color': SUB}),
        ('● 대조되는 두 인물', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   왕융간요 + 배해청통 — 「**간결**」 vs 「**청아**」 두 미덕', {'font_size': 14, 'color': SUB}),
        ('● 같은 시대의 두 인물', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   공명와룡 + 여망비웅 — 천하의 두 책사가 발탁된 순간', {'font_size': 14, 'color': SUB}),
        ('● 어린이가 스스로 묻게 한다 — 「왜 둘이 함께인가?」', {'font_size': 16, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅴ. 명일화 16선 ==============
SEC5 = 'Ⅴ. 명일화 16선'

def make_anecdote_slide(section, idx_total, four_chars, korean_title, story, lesson):
    @S(section)
    def _(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, section, n, t)
        add_title(slide, f'명일화 {idx_total} — {korean_title}',
                  f'몽구의 한 구절이 만든 사자성어')
        add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.1), PALE)
        add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.1),
                    four_chars,
                    font_size=32, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_paragraphs(slide, Inches(0.7), Inches(3.6), Inches(12.0), Inches(2.0), [
            ('일화', {'font_size': 14, 'bold': True, 'color': SUB}),
            (story, {'font_size': 15, 'color': INK}),
        ], line_spacing=1.35)
        add_filled_rect(slide, Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.0), PALE)
        add_textbox(slide, Inches(0.9), Inches(5.95), Inches(11.6), Inches(0.9),
                    f'배움 — {lesson}',
                    font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return _


anecdotes = [
    ('1/16', '匡 衡 鑿 壁', '광형착벽 — 광형이 벽을 뚫다',
     '한나라 광형(匡衡)이 어렸을 때 가난해 등잔불이 없었다. 이웃집 등불 빛이 새어 들도록 벽을 뚫어 책을 읽었다.',
     '환경의 한계를 핑계 삼지 않는 학습 의지'),
    ('2/16', '孫 敬 閉 戶', '손경폐호 — 손경이 문을 닫다',
     '한나라 손경(孫敬)은 학문에 몰두하기 위해 문을 닫아걸고 외부와 단절하여 글에만 전념했다.',
     '집중을 위한 자기 격리'),
    ('3/16', '車 胤 囊 螢', '차윤낭형 — 차윤이 반딧불을 모으다',
     '진(晉)나라 차윤(車胤)은 가난해 등잔 기름을 살 수 없어, 여름밤 반딧불을 모아 그 빛으로 책을 읽었다.',
     '없음을 탓하지 않는 창의적 면학 — 형설지공의 한 축'),
    ('4/16', '孫 康 映 雪', '손강영설 — 손강이 눈빛에 비추다',
     '진나라 손강(孫康)은 겨울밤 흰 눈에 반사된 달빛으로 책을 읽었다. 차윤과 함께 「**형설지공**」의 원전.',
     '계절의 자연을 학습에 활용하는 지혜'),
    ('5/16', '孔 明 臥 龍', '공명와룡 — 공명이 누운 용이다',
     '제갈량(諸葛亮)이 남양 융중에 은거할 때, 사람들이 그를 「잠긴 용(臥龍)」이라 불렀다. 유비가 삼고초려해 모셨다.',
     '진짜 인재는 알아보는 이가 있어야 일어난다'),
    ('6/16', '呂 望 非 熊', '여망비웅 — 여망은 곰이 아니었다',
     '주 문왕이 위수(渭水) 강가에서 점을 쳐 「범도 곰도 아닌 것을 얻으리라」 했는데, 만난 자가 강태공(여망)이었다.',
     '인재는 곰처럼 뚜렷한 위세 없이 평범한 모습으로 온다'),
    ('7/16', '陶 母 截 髮', '도모절발 — 도간 어머니의 머리카락',
     '진(晉) 도간(陶侃)의 어머니는 손님을 대접할 형편이 안 되자, 자기 머리카락을 잘라 팔아 술과 안주를 마련했다.',
     '자식의 인맥을 위해 자기를 다 내어주는 어머니의 사랑'),
    ('8/16', '孟 母 三 遷', '맹모삼천 — 맹자 어머니의 세 번 이사',
     '맹자의 어머니가 처음엔 묘지 근처, 다음엔 시장 근처, 마지막엔 서당 근처로 세 번 이사해 자식의 환경을 가다듬었다.',
     '환경이 곧 교육 — 三遷之敎의 원전'),
    ('9/16', '丁 蘭 刻 木', '정란각목 — 정란이 나무에 새기다',
     '한나라 정란(丁蘭)은 부모가 일찍 돌아가시자, 나무로 부모의 모습을 새겨 살아 계신 듯 봉양했다.',
     '효는 부모의 죽음으로 끝나지 않는다'),
    ('10/16', '孟 嘉 落 帽', '맹가낙모 — 맹가의 모자',
     '동진 맹가(孟嘉)가 환온의 연회에서 바람에 모자가 날아갔는데도 알아채지 못했다. 환온이 그를 풍류의 사람이라 칭송했다.',
     '풍류(風流)의 원전 — 작은 일에 흔들리지 않는 자연스러움'),
    ('11/16', '陳 蕃 一 室', '진번일실 — 진번의 한 방',
     '후한 진번(陳蕃)이 15세에 한 방을 어지럽혀 두자, 손님이 「**왜 청소하지 않는가**」 물었다. 진번이 「**대장부는 천하를 청소해야지 어찌 한 방을 청소하리오**」 답했다.',
     '큰 뜻을 품은 자의 자존 — 다만 큰 일은 작은 일에서 시작한다는 후세의 반론도'),
    ('12/16', '楊 震 四 知', '양진사지 — 양진의 네 가지가 안다',
     '후한 양진(楊震)이 밤에 뇌물을 받지 않으며 「**하늘이 알고 신이 알고 그대가 알고 내가 안다(天知 神知 子知 我知)**」고 했다.',
     '아무도 보지 않을 때의 자기 단속 — 慎獨의 원형'),
    ('13/16', '王 戎 觀 李', '왕융관리 — 왕융이 자두를 보다',
     '왕융이 어렸을 때 길가에 열매가 가득한 자두나무를 보고도 따러 가지 않았다. 「**열매가 많은데 길가에 있다면 분명 쓴 자두**」.',
     '겉만 보고 달려들지 않는 분별력'),
    ('14/16', '謝 安 高 潔', '사안고결 — 사안의 고결함',
     '동진 사안(謝安)이 동산(東山)에 20년 은거하다 나라가 위태로워지자 출사했다. 비수대전에서 부견을 격파한 뒤에도 흐트러지지 않았다.',
     '은거의 깊이가 곧 출사의 무게 — 평정한 위기 대응'),
    ('15/16', '蘇 武 持 節', '소무지절 — 소무가 부절을 들다',
     '한 무제의 사신 소무(蘇武)가 흉노에 19년 억류되었으나 부절을 한 번도 손에서 놓지 않았다. 결국 돌아왔다.',
     '사명에 대한 끝없는 충성과 인내'),
    ('16/16', '伯 牙 絶 絃', '백아절현 — 백아가 줄을 끊다',
     '거문고의 명인 백아(伯牙)는 자기 음악을 알아주던 종자기(鍾子期)가 죽자, 거문고 줄을 끊고 다시는 연주하지 않았다.',
     '나를 알아주는 한 사람의 가치 — 知音의 원전'),
]

for tag, fc, kt, st, ls in anecdotes:
    make_anecdote_slide(SEC5, tag, fc, kt, st, ls)


# ============== Ⅵ. 5대 교육 사상 ==============
SEC6 = 'Ⅵ. 5대 교육 사상'

@S(SEC6)
def vi_voluntary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 1 — 자발적 배움',
              '童蒙求我 — 어린이가 스스로 찾아오는 학습')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 강요가 아닌 호기심에 의한 학습', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   책 이름 자체가 그 정신을 선언 — 「**어린이가 나를 찾아온다**」', {'font_size': 14, 'color': SUB}),
        ('● 운율과 대구의 매력 — 형식이 동기를 만든다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   어린이가 외우고 싶어지도록 — 노래처럼 흘러가는 8자', {'font_size': 14, 'color': SUB}),
        ('● 인물 일화의 흥미 — 추상이 아닌 사람', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   광형이 벽을 뚫었다, 차윤이 반딧불을 모았다 — 그림이 그려진다', {'font_size': 14, 'color': SUB}),
        ('● 오늘날 — 「**내적 동기 학습**」, 「**게이미피케이션**」의 원조', {'font_size': 16, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC6)
def vi_person(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 2 — 인물 중심의 역사',
              '사건이 아니라 사람이 역사를 만든다')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 1,200년의 중국사를 600여 인물로 — 「**역사 = 사람의 선택**」', {'font_size': 17, 'space_before': 6}),
        ('● 각 인물의 결정적 한 행동을 4자로 압축', {'font_size': 17, 'space_before': 10}),
        ('● 어린이가 인물의 행적을 따라 살게 됨 — 모범의 내면화', {'font_size': 17, 'space_before': 10}),
        ('● 「**한 행동이 그를 영원히 정의한다**」 — 자기 4자의 자각', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 사기·자치통감의 인물 중심 사관의 어린이판 — 동일 정신', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC6)
def vi_concrete(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 3 — 추상이 아닌 사례',
              '덕목을 구체적 인물 일화로 환원')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.8), [
        ('● 「**효란 무엇인가**」를 설명하지 않는다', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   대신 — 정란이 부모의 모습을 나무에 새겨 모셨다고 말한다', {'font_size': 14, 'color': SUB}),
        ('● 「**근면이란 무엇인가**」를 설명하지 않는다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   대신 — 광형이 벽을 뚫고, 차윤이 반딧불을 모았다고 말한다', {'font_size': 14, 'color': SUB}),
        ('● 「**충절이란 무엇인가**」를 설명하지 않는다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   대신 — 소무가 19년 부절을 들고 돌아왔다고 말한다', {'font_size': 14, 'color': SUB}),
        ('● 「**스토리텔링 학습**」의 가장 오래된 동양적 원형', {'font_size': 16, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC6)
def vi_pairing(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 4 — 비교·종합 사고',
              '두 인물 짝짓기가 만드는 비판적 사고')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 어린이가 자연스럽게 묻게 된다 — 「왜 둘이 함께인가?」', {'font_size': 17, 'space_before': 6}),
        ('● 비슷한가? 대조되는가? 같은 시대인가? — 비교 사고', {'font_size': 17, 'space_before': 10}),
        ('● 단순 암기가 아닌 — 종합적 역사 인식으로의 도약', {'font_size': 17, 'space_before': 10}),
        ('● 오늘날 「**Compare & Contrast Learning**」의 동양적 원형', {'font_size': 16, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC6)
def vi_balance(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '사상 5 — 권선징악의 균형',
              '모범과 반면교사를 모두 보인다')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 대부분은 본받을 모범 사례 — 면학·효·충절·지혜', {'font_size': 17, 'space_before': 6}),
        ('● 그러나 간신·소인의 일화도 포함됨', {'font_size': 17, 'space_before': 10}),
        ('● 「**반면교사**」로서의 경계 역할', {'font_size': 17, 'space_before': 10}),
        ('● 어린이가 선악 양면을 보고 — 스스로 판단하게 한다', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 일방적 교화가 아닌 — 균형 잡힌 도덕 인식의 형성', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


# ============== Ⅶ. 사자성어 원전 ==============
SEC7 = 'Ⅶ. 사자성어 원전'

@S(SEC7)
def vii_idioms(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '몽구가 만든 사자성어 8선',
              '오늘 한국어에 살아 있는 표현들')
    rows = [
        ('螢雪之功 (형설지공)', '차윤낭형 + 손강영설 — 가난 속의 면학'),
        ('鑿壁引光 (착벽인광)', '광형착벽 — 환경의 한계를 깨다'),
        ('臥龍 (와룡)', '공명와룡 — 잠겨 있는 큰 인재'),
        ('三遷之敎 (삼천지교)', '맹모삼천 — 환경이 곧 교육'),
        ('風流 (풍류)', '맹가낙모 — 작은 일에 흔들리지 않는 자연스러움'),
        ('知音 (지음)', '백아절현 — 나를 알아주는 한 사람'),
        ('四知 (사지)', '양진사지 — 慎獨의 원형'),
        ('一室 (일실)', '진번일실 — 큰 뜻과 작은 일'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.52)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.2), Inches(0.45), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(3.2), Inches(0.45),
                    k, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.1), y, Inches(8.8), Inches(0.45),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC7)
def vii_spread(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '사자성어가 동아시아 공통 교양이 된 비밀',
              '몽구의 결정적 역할')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 한국·중국·일본·베트남이 공유하는 사자성어 — 그 통로가 몽구', {'font_size': 17, 'space_before': 6}),
        ('● 어린이가 외운 한 단위 = 평생의 교양 어휘', {'font_size': 17, 'space_before': 10}),
        ('● 형설지공·삼천지교·와룡·풍류·지음 — 1,300년 살아남은 4자', {'font_size': 17, 'space_before': 10}),
        ('● 「**짧고 강한 압축**」이 가장 오래 살아남는다 — 정보 시대에도 유효', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅷ. 한국·일본 수용 ==============
SEC8 = 'Ⅷ. 한국·일본 수용'

@S(SEC8)
def viii_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '한국에서의 수용',
              '천자문 → 명심보감 → 몽구 → 사서삼경')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.8), [
        ('● 고려·조선 사대부 자제의 기초 교재', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   학습 단계 — 천자문 → 명심보감 → **몽구** → 사서삼경', {'font_size': 14, 'color': SUB}),
        ('   몽구는 역사·인물 지식의 입문서 역할', {'font_size': 14, 'color': SUB}),
        ('● 조선 말기 — 홍익주(洪翼周)의 『**몽구주해**』', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   「자구는 간략하나 뜻이 깊어, 선비 가문의 매일 학습서로 이용」', {'font_size': 14, 'color': SUB}),
        ('● 속찬(續撰) 전통 — 한국사 인물을 다룬 「**○○몽구**」류 다수 편찬', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   동국몽구류 — 한국 학자들이 몽구 형식을 모방한 한국판', {'font_size': 14, 'color': SUB}),
    ])


@S(SEC8)
def viii_japan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '일본에서의 수용 — 동아시아 최대의 영향',
              '헤이안 귀족 사회의 절대 필독서')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.8), [
        ('● 헤이안(平安) 시대(794~1185)부터의 필수 학습서', {'font_size': 17, 'bold': True, 'color': ACCENT}),
        ('   878년 사다야스 친왕(9세)이 처음 학습한 기록이 남아 있음', {'font_size': 14, 'color': SUB}),
        ('● 「**권학원의 참새는 몽구를 지저귄다**」(勸學院の雀は蒙求を囀る)', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('   학교 처마의 참새조차 매일 듣고 따라 지저귈 정도였다는 속담', {'font_size': 14, 'color': SUB}),
        ('● 미요시 다메야스의 『**동몽송운(童蒙頌韻)**』 등 일본판 몽구류 출현', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 『**겐지모노가타리**』 등 헤이안 문학의 중국 고사 — 대부분 몽구가 매개', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC8)
def viii_xu_bzhu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '송 서자광의 『**몽구보주**』 — 결정적 보충',
              '압축된 4자를 풀어쓴 주석본')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 원문이 너무 압축적이어서 일화 배경을 모르면 이해 불가', {'font_size': 17, 'space_before': 6}),
        ('● 송 서자광(徐子光)이 각 4자마다 출전과 일화 전말을 상세히 주석', {'font_size': 17, 'space_before': 10}),
        ('● 오늘날 전해지는 몽구는 사실상 모두 「**보주본**」 계통', {'font_size': 17, 'space_before': 10}),
        ('● 압축의 매력 + 주석의 친절 — 두 층의 결합이 1,300년 생명력의 비결', {'font_size': 16, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC8)
def viii_legacy(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '후대 어린이 교재의 원형',
              '몽구가 만든 동몽서의 표준')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 천자문(6세기) → **몽구(8세기)** → 삼자경·백가성(13세기)', {'font_size': 17, 'space_before': 6}),
        ('● 몽구가 정착시킨 「**운문 + 인물 + 대구**」 — 후대 모든 동몽서가 계승', {'font_size': 17, 'space_before': 10}),
        ('● 동아시아 어린이 교육의 표준 형식이 됨', {'font_size': 17, 'space_before': 10}),
        ('● 「**한 권의 어린이 책이 동아시아 1,300년 교양을 만들었다**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


# ============== Ⅸ. 오늘 다시 펼치는 이유 ==============
SEC9 = 'Ⅸ. 오늘 다시 펼치는 이유'

@S(SEC9)
def ix_today1(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '1 — 「자기 4자」에 대한 자각',
              '나는 어떤 4자로 기억될 것인가')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 광형은 「鑿壁(벽을 뚫음)」으로 영원히 기억된다', {'font_size': 17, 'space_before': 6}),
        ('● 차윤은 「囊螢(반딧불 자루)」으로, 손강은 「映雪(눈빛)」으로', {'font_size': 17, 'space_before': 10}),
        ('● 평생의 수많은 행위 중 — 후세가 기억하는 것은 결정적 한 장면', {'font_size': 17, 'space_before': 10}),
        ('● 「**나는 어떤 4자로 기억될 것인가**」 — 어린이가 평생 묻는 자각', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 오늘의 행위 하나가 — 나의 4자가 될지 모른다', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today2(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '2 — 압축이 가장 강력한 교육',
              '본질만 남기면 가장 오래 기억된다')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 2,384자 — 600여 인물 — 1,200년 역사가 들어 있다', {'font_size': 17, 'space_before': 6}),
        ('● 오늘날 정보 과잉의 시대에도 — 압축의 원리는 그대로', {'font_size': 17, 'space_before': 10}),
        ('● 본질만 남기고 모든 잉여를 제거할 때 — 가장 오래 살아남는다', {'font_size': 17, 'space_before': 10}),
        ('● 「**핵심 메시지를 한 줄로**」 — 모든 커뮤니케이션의 정수', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ])


@S(SEC9)
def ix_today3(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '3 — 학습은 외우는 게 아니라 닮는 것',
              '몽구의 가장 깊은 교육 철학')
    add_paragraphs(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.5), [
        ('● 운율과 대구로 외우게 만든 것은 — 단순 암기 효율 때문이 아니다', {'font_size': 17, 'space_before': 6}),
        ('● 어린이가 평생 입에 올리며 — 인물의 행적이 자기 삶의 척도가 되도록', {'font_size': 17, 'space_before': 10}),
        ('● 역사를 「**외운다**」가 아니라 「**닮는다**」', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('● 모범의 내면화 — 이게 동양 인격 교육의 정수', {'font_size': 16, 'color': SUB, 'space_before': 10}),
    ])


# ============== Ⅹ. 마무리 ==============
SEC10 = 'Ⅹ. 마무리'

@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '몽구가 일러주는 7가지',
              '한 폭으로 정리')
    items = [
        '어린이가 스스로 찾아오는 학습 — 童蒙求我',
        '인물의 한 장면이 그를 영원히 정의 — 자기 4자에 대한 자각',
        '추상이 아닌 사례 — 구체적 인물 일화로 덕목을 배움',
        '두 인물 짝짓기 — 비교와 종합 사고력 형성',
        '운율과 대구 — 형식의 매력이 학습 동기를 만든다',
        '권선징악의 균형 — 모범과 반면교사를 함께 보인다',
        '압축이 가장 오래 살아남는다 — 본질만 남기는 교육',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.5 + i * 0.6)
        add_textbox(slide, Inches(0.9), y, Inches(0.6), Inches(0.5),
                    f'{i+1}.', font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.5),
                    txt, font_size=16, color=INK)


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                '몽구 — 어린이의 입과 가슴에 1,300년을 살린 600인의 한 장면',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '童 蒙 求 我',
                font_size=140, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
                '어린이가 나를 찾아온다 — 자발적 배움의 정수',
                font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.8), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
                '몽구가 묻는다 — 「당신은 어떤 4자로 기억될 것인가?」',
                font_size=15, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total = len(SLIDES)
for i, (fn, sec) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    fn(slide, i, total)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\몽구.pptx'
prs.save(out_path)
print(f'생성 완료: {out_path}  슬라이드 수: {total}')
