# -*- coding: utf-8 -*-
"""
사기(史記) 발표자료 — 망라적 83장 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 1. 표지 ==============
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8),
                '史 記', font_size=140, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                'Records of the Grand Historian · 사기',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.5),
                '사마천(司馬遷) · 130편 · 황제(黃帝)~한 무제 3,000년의 통사',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
                '본기 12 + 표 10 + 서 8 + 세가 30 + 열전 70 = 130편 · 52만 6,500자',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                '"史家之絕唱 無韻之離騷"',
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                '사가의 절창, 운율 없는 이소 — 노신(魯迅)',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 2. 목차 ==============
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.8), Inches(0.7),
                '목 차', font_size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.15), Inches(12.8))
    items = [
        [('Ⅰ', '개요 — 사기란 무엇인가'),
         ('Ⅱ', '사마천의 생애 — 궁형과 발분'),
         ('Ⅲ', '오체(五體) 구조 — 130편'),
         ('Ⅳ', '다섯 파격 — 격식을 깨다'),
         ('Ⅴ', '10대 주제'),
         ('Ⅵ', '핵심 인물 16인')],
        [('Ⅶ', '7대 명장면'),
         ('Ⅷ', '명구 10선'),
         ('Ⅸ', '사마천의 사상 — 삼구(三句)'),
         ('Ⅹ', '후대 영향'),
         ('Ⅺ', '현대적 의의'),
         ('Ⅻ', '마무리')],
    ]
    for col, group in enumerate(items):
        x = 0.7 + col * 6.4
        top = 1.6
        for num, title in group:
            add_textbox(slide, Inches(x), Inches(top), Inches(1.0), Inches(0.5),
                        num, font_size=17, bold=True, color=ACCENT)
            add_textbox(slide, Inches(x + 1.0), Inches(top), Inches(5.3), Inches(0.5),
                        title, font_size=17, color=INK)
            top += 0.7


# ============== Ⅰ. 개요 ==============
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '사기(史記)란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                '중국 최초의 기전체(紀傳體) 통사 · 24사(史)의 머리',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.5),
                '"인물"을 축으로 역사를 구성한 동양 역사 서술의 원형',
                font_size=15, color=ACCENT, align=PP_ALIGN.CENTER)
    nums = [('130', '편(篇)'), ('약 52만', '6,500자'), ('약 3,000', '년 통사'), ('14~15', '년 편찬')]
    for i, (n, lbl) in enumerate(nums):
        x = 0.6 + i * 3.05
        add_textbox(slide, Inches(x), Inches(4.0), Inches(2.9), Inches(1.0),
                    n, font_size=44, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(5.1), Inches(2.9), Inches(0.5),
                    lbl, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                '황제(黃帝)에서 한 무제 태초(太初) 연간까지 — 사실상 동양 문명사의 출발점',
                font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '"역사서이자 문학의 걸작" — 인물·일화·고사가 모두 여기서 시작',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_status(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '책의 위상 — 핵심 정보')
    rows = [
        ('정식 명칭', '태사공서(太史公書)·태사공기(太史公記)',  '위진(魏晉) 이후 "사기(史記)"로 통칭'),
        ('편자',     '사마천(司馬遷, BC 145?~86?)',           '한(漢) 무제 시대 태사령(太史令)'),
        ('편찬',     'BC 108 착수 → BC 91경 완성',              '약 14~15년 · 궁형(BC 99) 전후 작업'),
        ('분량',     '130편 · 52만 6,500자',                    '본기 12 + 표 10 + 서 8 + 세가 30 + 열전 70'),
        ('수록 기간', '약 3,000년',                              '전설의 황제(黃帝)~한 무제 태초 연간'),
        ('형식',     '기전체(紀傳體)',                          '인물 중심 — 사마천이 발명'),
        ('성격',     '24사(史)의 머리',                          '이후 모든 정사가 이 체제를 계승'),
        ('평가',     '"史家之絕唱, 無韻之離騷"',                 '— 노신(魯迅) — "역사이자 문학의 절정"'),
    ]
    top = 1.95
    for i, (tag, val, note) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.6), PALE)
        add_filled_rect(slide, Inches(2.85), Inches(top), Inches(10.0), Inches(0.6), bg)
        add_textbox(slide, Inches(0.55), Inches(top + 0.15), Inches(2.2), Inches(0.4),
                    tag, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.95), Inches(top + 0.05), Inches(4.5), Inches(0.5),
                    val, font_size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.5), Inches(top + 0.1), Inches(5.3), Inches(0.45),
                    note, font_size=12, color=SUB)
        top += 0.66


@S('Ⅰ. 개요')
def s_gijeon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '기전체(紀傳體) — 사마천의 발명')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '이전의 사서 — 편년체·국별체', font_size=18, bold=True, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('연대순(편년체)', {'font_size': 16, 'bold': True}),
        ('  춘추 · 좌전', {'font_size': 14, 'color': SUB}),
        ('', {'font_size': 8}),
        ('국가별(국별체)', {'font_size': 16, 'bold': True, 'space_before': 6}),
        ('  국어(國語) · 전국책(戰國策)', {'font_size': 14, 'color': SUB}),
        ('', {'font_size': 8}),
        ('한계', {'font_size': 14, 'color': ACCENT, 'space_before': 6}),
        ('"개인의 삶이 드러나지 않는다"', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '사기 — 기전체(紀傳體)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"인물"을 축으로 역사를 구성',
         {'font_size': 16, 'bold': True, 'color': ACCENT}),
        ('', {'font_size': 6}),
        ('• 본기(紀) — 제왕의 일생', {'font_size': 13, 'space_before': 6}),
        ('• 열전(傳) — 신하·인물의 일생', {'font_size': 13, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('"한 사람의 삶"이 역사의 단위',
         {'font_size': 14, 'bold': True, 'color': INK, 'space_before': 6}),
        ('"왕조의 흥망"이 아니라 "인간이 어떻게 살고 죽었는가"',
         {'font_size': 12, 'color': SUB}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '한서·후한서·24사 모두 이 체제를 계승 — 동양 역사 서술의 표준',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_nosin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '노신(魯迅)의 한 마디 — 사기의 본질')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('史 家 之 絕 唱   無 韻 之 離 騷',
         {'font_size': 32, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('사가지절창 · 무운지이소',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"사가(史家)의 절창(絕唱)이요',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 운율 없는 이소(離騷)"',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"역사이자 동시에 문학의 절정"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('사가로서는 가장 높은 노래(絕唱)이고',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('운율은 없으나 굴원의 이소(離騷)에 견줄 만큼 문학적이다',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅰ. 개요')
def s_eastasia(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '동아시아 한문 교양의 공통 기반',
              '"한·중·일 지식인이 공유하는 인물·고사·명구가 모두 사기에서 나왔다"')
    items = [
        ('중국',  '24사(史)의 머리 · 모든 정사의 표준',
         '당송팔대가의 고문 전범 · 삼국지연의·수호전의 인물 작법'),
        ('한국',  '김부식 『삼국사기』(1145) — 본기·연표·지·열전 4체 계승',
         '조선왕조실록 · 정조의 사기 강독 · 김원중 역(현대)'),
        ('일본',  '헤이안 시대 귀족 교양 · 에도 무사의 필독서',
         '시바 료타로(司馬遼太郎) 필명 · 라이 산요 『일본외사』'),
        ('동아시아 공통', '인물 일화 · 고사성어 · 명구의 공통 기반',
         '"한문 교양 = 사기를 읽었다는 뜻"'),
    ]
    top = 2.4
    for region, primary, secondary in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.95), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(2.5), Inches(0.5),
                    region, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(9.6), Inches(0.95), PALE)
        add_textbox(slide, Inches(3.35), Inches(top + 0.1), Inches(9.3), Inches(0.45),
                    primary, font_size=13, bold=True, color=INK)
        add_textbox(slide, Inches(3.35), Inches(top + 0.5), Inches(9.3), Inches(0.45),
                    secondary, font_size=12, color=SUB)
        top += 1.07


# ============== Ⅱ. 사마천 생애 ==============
@S('Ⅱ. 사마천')
def s_simajia(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 사마천', page, total)
    add_title(slide, '사마씨 — 주(周) 왕실부터 이어온 사관(史官) 가문')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '司\n馬\n遷', font_size=110, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('사마천(BC 145?~86?)', {'bold': True, 'font_size': 18, 'color': ACCENT}),
        ('  자(字) 자장(子長) · 한(漢) 무제 시대',
         {'font_size': 14}),
        ('', {'font_size': 6}),
        ('가문 — 사관(史官)의 세습',
         {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 6}),
        ('  주(周) 왕실 사관을 이어온 세습 가문',
         {'font_size': 14}),
        ('아버지 — 사마담(司馬談)',
         {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 6}),
        ('  한 무제 때 태사령(太史令)',
         {'font_size': 14}),
        ('  천문·역법·국가 제사·문서 기록 관장',
         {'font_size': 13, 'color': SUB}),
        ('  사기 편찬의 첫 불씨를 사마천에게 남김',
         {'font_size': 13, 'bold': True, 'color': INK}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.3)


@S('Ⅱ. 사마천')
def s_sajadam(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 사마천', page, total)
    add_title(slide, '아버지 사마담의 유언 (BC 110) — 사기의 첫 불씨')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('한 무제의 태산 봉선(封禪)에 참여 못 하고 낙양에서 임종 직전',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"내가 죽거든 너는 반드시 태사가 되어라',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 태사가 되거든 내가 논저(論著)하려 했던 바를 잊지 말라"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"공자가 춘추를 지은 뒤 400여 년이 지났고 역사의 기록이 끊어졌다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        (' 한이 흥한 이래 명주·현군·충신·의사들의 일을',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        (' 내가 논평하지 못한다면 천하의 문헌이 사라질까 두렵다"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('→ 사기의 동기 ① — 아버지의 유언을 잇는 효(孝)와 사관의 책임',
         {'font_size': 13, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅱ. 사마천')
def s_chunhalim(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 사마천', page, total)
    add_title(slide, '20세의 천하 순력(巡歷) — 사기의 동기 ②',
              '"문헌과 현장을 결합한 실증(實證)의 역사"')
    items = [
        ('회계(會稽)',     '우(禹) 임금의 무덤'),
        ('구의(九疑)',     '순(舜) 임금의 전설'),
        ('장강·회수·제노', '광활한 답사 — 풍토와 인심'),
        ('곡부(曲阜)',     '공자의 고향 — "高山仰止" 체험'),
        ('설(薛)',         '맹상군의 옛 자취'),
        ('오강(烏江)',     '항우가 자결한 자리'),
    ]
    top = 2.3
    for place, content in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(3.0), Inches(0.4),
                    place, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(9.1), Inches(0.6), PALE)
        add_textbox(slide, Inches(3.9), Inches(top + 0.13), Inches(8.8), Inches(0.4),
                    content, font_size=14, color=INK)
        top += 0.68
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '"후일 열전의 생생한 현장감은 이때의 답사에서 나왔다"',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 사마천')
def s_ireung(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 사마천', page, total)
    add_title(slide, '이릉(李陵)의 화(BC 99) — 궁형(宮刑)의 비극')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('BC 99년 — 명장 이광(李廣)의 손자 이릉이 흉노에 항복',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('조정의 모두가 이릉을 역적으로 매도',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('오직 사마천만이 무제 앞에서 이릉을 변호 — "중과부적이었습니다"',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('무제의 격노 — 사형 선고',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('한대 사형수의 두 길 — ① 돈 50만 전 또는 ② 궁형(宮刑, 거세형)',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('돈이 없었고 가족·친구 누구도 나서지 않음 → 궁형을 선택',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅱ. 사마천')
def s_boimanseo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 사마천', page, total)
    add_title(slide, '『보임안서(報任安書)』 — 발분저서(發憤著書)',
              '"문왕이 갇혀 주역을, 중니가 곤하여 춘추를…"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.6),
                '"人 固 有 一 死  或 重 於 泰 山  或 輕 於 鴻 毛"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.4),
                '인고유일사 혹중어태산 혹경어홍모',
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.55), Inches(11.7), Inches(1.3), [
        ('"사람은 누구나 한 번 죽는다',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        (' 그 죽음이 태산보다 무거울 수도 있고, 기러기 털보다 가벼울 수도 있다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"문왕이 갇혀 주역을 펼치고, 중니가 곤하여 춘추를 짓고,',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        (' 굴원이 쫓겨나 이소를 짓고, 좌구명이 실명하여 국어를 남기고,',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        (' 손자가 다리 잘리고 병법을, 한비자가 갇혀 세난·고분을 지었다"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('→ 발분저서(發憤著書) — 고난이 기록을 낳는다',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅱ. 사마천')
def s_hidden(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 사마천', page, total)
    add_title(slide, '"명산에 감춘 책" — 사후의 전승')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('완성 직후 사마천의 선언',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"藏 諸 名 山  傳 之 其 人"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"이 책을 명산(名山)에 감추고, 부본을 경사(京師)에 남겨 후세의 군자를 기다린다"',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('전승의 연쇄',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 사마천 사후 — 외손자 양운(楊惲)이 부본 보관',
         {'font_size': 14, 'space_before': 6}),
        ('• 선제(宣帝, BC 74~49) 시기 — 양운이 사기 공개',
         {'font_size': 14, 'space_before': 4}),
        ('• 반고 『한서』(1세기) — 사기의 무제 이전 부분을 흡수, 함께 보존',
         {'font_size': 14, 'space_before': 4}),
        ('• 송대 — 사기집해·색은·정의 "삼가주(三家注)" 합본',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅲ. 오체 구조 ==============
@S('Ⅲ. 오체')
def s_5che(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 오체', page, total)
    add_title(slide, '오체(五體) — 130편의 다섯 기둥')
    rows = [
        ('본기(本紀)', '12편', '제왕·실질 최고 권력자',  '연대기(紀)',  '황제~한 무제'),
        ('표(表)',     '10편', '시대 전체',              '연표·세계표', 'BC 2700~BC 100경'),
        ('서(書)',     '8편',  '제도·문물',              '분야사',       '고대~한대'),
        ('세가(世家)', '30편', '제후국·공신·특수 인물',   '가계 열전',    '주초~한 무제'),
        ('열전(列傳)', '70편', '개인·집단·이민족',        '인물 열전',    '은말~한 무제'),
    ]
    top = 2.0
    headers = ['부', '편수', '주인공', '형식', '시간 축']
    widths = [1.6, 1.2, 4.0, 2.7, 3.3]
    x = 0.5
    for i, (h, w) in enumerate(zip(headers, widths)):
        add_filled_rect(slide, Inches(x), Inches(top), Inches(w), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(x), Inches(top + 0.1), Inches(w), Inches(0.4),
                    h, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w + 0.05
    row_h = 0.65
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE
        x = 0.5
        for i, (val, w) in enumerate(zip(row, widths)):
            add_filled_rect(slide, Inches(x), Inches(y), Inches(w), Inches(row_h), bg)
            c = ACCENT if i == 0 else INK
            add_textbox(slide, Inches(x + 0.1), Inches(y + 0.15), Inches(w - 0.2),
                        Inches(0.4), val, font_size=13, color=c,
                        bold=(i == 0 or i == 1), align=PP_ALIGN.CENTER)
            x += w + 0.05
    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                '합계 130편 · 52만 6,500자 — 한 사건을 다섯 각도에서 비추는 입체 구조',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅲ. 오체')
def s_bongi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 오체', page, total)
    add_title(slide, '본기(本紀) 12편 — 제왕의 연대기')
    rows = [
        ('1',  '오제본기',       '황제·전욱·제곡·요·순 — 문명의 기원'),
        ('2',  '하본기',         '우 임금 — 치수와 선양→세습'),
        ('3',  '은본기',         '탕~주왕 — 천명(天命)의 이동'),
        ('4',  '주본기',         '문왕·무왕~동주 — 봉건의 흥망 800년'),
        ('5',  '진본기',         '진 목공~장양왕 — 변방의 강국화'),
        ('6',  '진시황본기',     '진시황·2세·자영 — 통일과 15년 멸망'),
        ('7',  '항우본기 ★',     '서초패왕 — 제왕이 아닌데도 본기 (파격)'),
        ('8',  '고조본기',       '한 태조 유방 — 정장(亭長)에서 황제로'),
        ('9',  '여태후본기 ★',   '여후의 수렴청정 — 여성을 본기에 (파격)'),
        ('10-12', '효문·효경·효무본기', '문제·경제·무제 — 문경의 치와 한 무제'),
    ]
    top = 1.95
    for num, name, desc in rows:
        is_special = '★' in name
        c = ACCENT if is_special else SUB
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(1.0), Inches(0.5), c)
        add_textbox(slide, Inches(0.5), Inches(top + 0.08), Inches(1.0), Inches(0.4),
                    num, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.6), Inches(top), Inches(3.5), Inches(0.5),
                        PALE if not is_special else RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(1.7), Inches(top + 0.08), Inches(3.4), Inches(0.4),
                    name, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.2), Inches(top), Inches(7.6), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.4), Inches(top + 0.08), Inches(7.3), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.55


@S('Ⅲ. 오체')
def s_pyo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 오체', page, total)
    add_title(slide, '표(表) 10편 — 시대의 골격', '"오늘날의 타임라인 차트"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('표(表)는 사기의 가장 독창적 발명 중 하나',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('편년체 서사에서 놓치기 쉬운',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"동시대 여러 국가·인물의 병렬 관계"를 도표로 정리',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('대표 표 — 한흥이래장상명신연표',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('한나라 역대 재상의 재임 기간과 사유를 한눈에',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"왜 왔다가 어떻게 떠났는가" — 한 도표로 1세기 정치사 파노라마',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('현대의 인포그래픽·타임라인의 동양 원형',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅲ. 오체')
def s_seo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 오체', page, total)
    add_title(slide, '서(書) 8편 — 문명의 제도 (분야사)')
    items = [
        ('예서(禮書)',      '예제의 역사',          '의전·사회 규범'),
        ('악서(樂書)',      '음악과 교화',          '문화 정책'),
        ('율서(律書)',      '병법과 음률',          '국방·기준 체계'),
        ('역서(曆書)',      '역법',                 '달력·시간'),
        ('천관서(天官書)',  '천문',                 '기상·천체'),
        ('봉선서(封禪書)',  '제왕의 태산 제사',     '종교 정책'),
        ('하거서(河渠書)',  '치수·운하',            '국토 개발'),
        ('평준서(平準書) ★','화폐·물가·재정',       '경제 정책 — 사기의 현대성'),
    ]
    top = 2.0
    for tag, content, modern in items:
        is_special = '★' in tag
        c = ACCENT if is_special else SUB
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(3.0), Inches(0.55), c)
        add_textbox(slide, Inches(0.5), Inches(top + 0.13), Inches(3.0), Inches(0.4),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_special else PALE
        add_filled_rect(slide, Inches(3.6), Inches(top), Inches(4.5), Inches(0.55), bg)
        add_textbox(slide, Inches(3.75), Inches(top + 0.13), Inches(4.3), Inches(0.4),
                    content, font_size=13, bold=True, color=INK)
        add_filled_rect(slide, Inches(8.2), Inches(top), Inches(4.6), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(8.35), Inches(top + 0.13), Inches(4.4), Inches(0.4),
                    modern, font_size=13, color=SUB)
        top += 0.62
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"평준서 — 국가 경제를 역사서 본문으로 다룬 최초의 시도"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅲ. 오체')
def s_sega(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 오체', page, total)
    add_title(slide, '세가(世家) 30편 — 제후국과 특별한 인물들')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('춘추전국의 주요 제후국 + 한 초 제후왕·공신의 전기',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('오·제·노·연·진·초·월·정·조·위·한 등',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(5.9), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(0.5), Inches(4.65), Inches(5.9), Inches(0.4),
                '파격 ① — 공자세가(孔子世家)', font_size=15, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(5.3), Inches(1.8), [
        ('포의(布衣)의 선비를 세가에', {'font_size': 14, 'bold': True}),
        ('"높은 산을 우러러본다(高山仰止)"',
         {'font_size': 13, 'color': SUB, 'space_before': 4}),
        ('사상가를 제후 반열에 — 전무후무',
         {'font_size': 13, 'color': ACCENT, 'bold': True, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(4.5), Inches(5.9), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(4.65), Inches(5.9), Inches(0.4),
                '파격 ② — 진섭세가(陳涉世家)', font_size=15, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(5.15), Inches(5.3), Inches(1.8), [
        ('진(秦)에 반기 든 빈농 진승(陳勝)', {'font_size': 14, 'bold': True}),
        ('"왕후장상에 어찌 씨가 있으랴',
         {'font_size': 13, 'color': SUB, 'space_before': 4}),
        ('(王侯將相寧有種乎)"',
         {'font_size': 13, 'color': ACCENT, 'bold': True}),
        ('역사를 바꾼 평민 봉기를 인정',
         {'font_size': 13, 'color': SUB, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅲ. 오체')
def s_yeoljeon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 오체', page, total)
    add_title(slide, '열전(列傳) 70편 — 인간 군상의 파노라마',
              '"사기의 백미 — 개인·집단·이민족을 아우른다"')
    items = [
        ('의리·지조형',     '백이·오자서·굴원·염파·인상여'),
        ('책략가·유세가',   '소진·장의·범저·채택'),
        ('전국 사공자',     '맹상군·평원군·신릉군·춘신군'),
        ('명장·병법가',     '손자·오기·백기·왕전·이광·위청·곽거병'),
        ('한 초 공신',       '소하·조참·장량·한신·경포·팽월'),
        ('사상가',           '노자·한비·맹자·순경'),
        ('자객 5인',         '조말·전제·예양·섭정·형가'),
        ('이민족',           '흉노·남월·동월·조선·서남이·대원'),
        ('서민 유형',        '순리·혹리·유협·골계·편작·화식'),
        ('총결',             '태사공자서(太史公自序) — 사마천의 자전'),
    ]
    top = 2.0
    for tag, content in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(3.5), Inches(0.4),
                    tag, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.2), Inches(top), Inches(8.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(4.4), Inches(top + 0.08), Inches(8.3), Inches(0.4),
                    content, font_size=13, color=INK)
        top += 0.5


# ============== Ⅳ. 다섯 파격 ==============
@S('Ⅳ. 다섯 파격')
def s_5pakyeok(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 다섯 파격', page, total)
    add_title(slide, '다섯 파격 — 사마천이 격식을 깨뜨린 다섯 지점')
    items = [
        ('1', '제왕 아닌 자를 본기에',         '7 항우본기',         '실질 권력 인정'),
        ('2', '여성을 본기에',                  '9 여태후본기',       '명분 아닌 실질'),
        ('3', '포의 학자를 세가에',             '17 공자세가',         '사상 권력 인정'),
        ('4', '평민 봉기 지도자를 세가에',      '18 진섭세가',         '역사 전환점 인정'),
        ('5', '상인·협객·점쟁이·해학가까지 열전에', '59~69 열전',      '인간 군상 전면화'),
    ]
    top = 2.4
    for num, title, where, why in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(5.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.22), Inches(4.9), Inches(0.5),
                    title, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.5), Inches(top), Inches(2.5), Inches(0.85),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(6.5), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    where, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(9.1), Inches(top), Inches(3.7), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(9.25), Inches(top + 0.22), Inches(3.5), Inches(0.5),
                    why, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        top += 0.95


@S('Ⅳ. 다섯 파격')
def s_pakyeok_essence(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 다섯 파격', page, total)
    add_title(slide, '파격의 공통 원리 — "신분(分)이 아니라 영향력(實)"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('사마천의 편제 원리',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('명분(名)이 아니라 실질(實)이 역사의 척도다',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"한 사람의 영향력이 역사에 남긴 자국을 본다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"成 一 家 之 言" — 성일가지언',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"한 집안(가문·학파)의 말을 이룬다"',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('단순 자료 편찬이 아니라',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('사마천 자신의 역사철학을 담은 저작',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅳ. 다섯 파격')
def s_jangpyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 다섯 파격', page, total)
    add_title(slide, '입체 구조의 예 — 장평대전(BC 260) 하나를 보면',
              '"한 사건이 다섯 차원에서 동시에 기록된다"')
    items = [
        ('본기',  '진본기·조세가 일부',     '국가 간 전쟁의 배경'),
        ('표',    '육국연표(六國年表)',     '그 해 여섯 나라가 무엇을 했나'),
        ('세가',  '조세가(趙世家)',         '조나라의 오판과 결과'),
        ('열전',  '백기왕전·염파인상여',    '장수들의 선택'),
        ('서',    '평준서(平準書)',         '40만을 묻은 뒤의 재정 파탄'),
    ]
    top = 2.5
    for tag, where, what in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.8), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(1.8), Inches(0.4),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.5), Inches(top), Inches(3.5), Inches(0.6), PALE)
        add_textbox(slide, Inches(2.6), Inches(top + 0.13), Inches(3.4), Inches(0.4),
                    where, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.1), Inches(top), Inches(6.7), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.3), Inches(top + 0.13), Inches(6.4), Inches(0.4),
                    what, font_size=13, color=INK)
        top += 0.7
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '국가·연표·분야·제후·개인 — 다섯 시선이 한 사건에 모인다',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅴ. 10대 주제 ==============
def make_theme_slide(num, title, original, principle, example, lesson):
    @S('Ⅴ. 10대 주제')
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, f'Ⅴ. 10대 주제 ({num}/10)', page, total)
        add_title(slide, f'주제 {num} — {title}', original)
        add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                    principle, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.8), PALE)
        add_textbox(slide, Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.4),
                    '대표 사례', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.1),
                    example, font_size=14, color=INK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.4),
                    '사기의 메시지', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.1),
                    lesson, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return renderer


make_theme_slide('1', '천도시비(天道是非)',
    '"천도(天道)는 선인을 돕는가 — 그러나 백이는 굶어 죽었다"',
    '선악과 운명의 불일치 — 사기 전체의 화두',
    '백이·숙제 수양산 아사 (1 백이열전)\n도척의 천수 · 굴원의 멱라 · 이광의 자결',
    '"답을 주는 책이 아니라 질문을 던지는 책"\n그들의 이름을 기록함으로써 천도를 대신한다')

make_theme_slide('2', '창업(創業) vs 수성(守成)',
    '"공수지세이야(攻守之勢異也)" — 가의 「과진론」',
    '얻는 덕과 지키는 덕은 다르다',
    '진시황 통일 → 15년 멸망 (6 진시황본기)\n한 고조 건국 → 여후의 전횡 (8·9)\n문경의 치 → 무제의 팽창',
    '"힘으로 얻고 덕으로 지킨다"\n정관정요의 창업수성론과 동일한 주제')

make_theme_slide('3', '용인(用人)',
    '"장량·소하·한신을 쓸 줄 알았기에 천하를 얻었다"',
    '사람을 쓸 줄 아는 자가 이긴다',
    '유방 vs 항우 — 삼걸 활용 vs 범증조차 못 씀\n진 목공·제 환공·진 효공·연 소왕의 인재 경영\n연 소왕의 "천금매골(千金買骨)" 일화',
    '"혼자 잘난 자가 아니라 남을 쓸 줄 아는 자"\n사기가 제시하는 리더십의 최종 답')

make_theme_slide('4', '공성신퇴(功成身退)',
    '"토사구팽(兔死狗烹)·조진궁장(鳥盡弓藏)"',
    '공을 이루고 물러설 줄 아는 자만이 살아남는다',
    '산 자 — 범려(월)·장량(한)·조참·장창\n죽은 자 — 문종(월)·백기(진)·상앙·오기·이사·한신·팽월·경포',
    '"공을 세우는 능력보다 공을 세운 후 처신이 더 어렵다"\n한신의 유언이 사기 전편을 흐르는 경고')

make_theme_slide('5', '법가의 이중주',
    '"법으로 세운 자는 법으로 망한다"',
    '법가는 강국을 만드나 자신을 보호하지 못한다',
    '상앙(위→진) · 오기(위→초) · 이사(초→진)\n한비(한→진) · 조조(한)\n개혁가 5인 모두 비참한 최후',
    '"법가의 효능을 인정하되 — 덕 없이 법만으로는\n자기도 보호 못 한다"는 아이러니')

make_theme_slide('6', '사위지기자사(士爲知己者死)',
    '"선비는 자기를 알아준 이를 위해 죽는다" — 예양',
    '한 번의 은혜에 목숨을 바치는 의리의 극단',
    '자객 5인 — 조말·전제·예양·섭정·형가\n평원군의 모수자천 · 위공자의 후영·주해\n한신의 표모(漂母) 일식천금',
    '사마천 자신이 이릉(자신을 몰랐던 자)을 변호하다\n궁형을 받은 체험의 문학적 승화')

make_theme_slide('7', '발분저서(發憤著書)',
    '"문왕이 갇혀 주역을, 중니가 곤하여 춘추를…"',
    '고난이 기록을 낳는다',
    '문왕(주역) · 중니(춘추) · 굴원(이소)\n좌구명(국어) · 손자(병법) · 여불위(여람)\n한비자(세난·고분) · 그리고 사마천(사기)',
    '"역사는 성공한 자가 아니라\n실패와 고난에서 펴 올린다"')

make_theme_slide('8', '현실주의 경제관',
    '"창고가 차야 예절을 안다(倉廩實則知禮節)" — 관중',
    '인간의 이(利) 추구를 인정하라',
    '관중 — "창고가 차야 예절"\n화식열전 — "천하가 오는 것은 모두 이익 때문"\n범려(도주공)·자공·백규·의돈·오씨나',
    '2,100년 전에 "경제가 역사의 엔진"임을 선언\n사기의 가장 현대적인 면모')

make_theme_slide('9', '이민족 서술의 객관성',
    '"흉노·조선·서남이도 한 사람의 역사"',
    '이민족을 야만이 아니라 자기 질서를 가진 사회로',
    '흉노열전 — 묵돌선우·평성 포위·화친·정벌\n조선열전 — 위만조선의 흥망 (한국 고대사 1차 사료)\n남월·동월·서남이·대원 — 장건의 서역 탐사',
    '이후 24사 이민족 서술 전통의 출발점\n중국 외부도 역사의 주체')

make_theme_slide('10', '자기 시대의 반(反)서사',
    '"승리자의 시대"인 무제의 공식 서사에 대한 행간의 비판',
    '사기는 무제 시대의 가장 집요한 반(反) 서사',
    '효무본기 — 봉선 의례의 허황됨\n봉선서 — 진시황·무제가 신선 쫓은 어리석음\n평준서 — 흉노 원정의 재정 파탄\n혹리열전 — 공포정치의 민낯',
    '"명산에 감춘 책" — 무제가 노할 수 있는 내용\n현대 중국사학자들의 공통 평가')


# ============== Ⅵ. 핵심 인물 16인 ==============
def make_person_slide(num, name_kor, name_han, era, role, story, lesson):
    @S(f'Ⅵ. 인물 ({num}/16)')
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, f'Ⅵ. 인물 ({num}/16)', page, total)
        add_title(slide, f'{name_kor}({name_han}) — {role}', era)
        add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.8), PALE)
        add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.4),
                    '인생·일화', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(2.95), Inches(11.7), Inches(2.0),
                    story, font_size=14, color=INK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.4),
                    '사마천의 평·교훈', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0),
                    lesson, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return renderer


make_person_slide('1', '백이·숙제', '伯夷·叔齊', '은말주초 · 고죽국 두 왕자',
    '의(義)의 화두 — 사기 전체의 시작',
    '왕위를 사양 → 주 무왕이 은 주왕 치려는 것을 만류\n"신하로 임금을 치는 것은 의가 아니다"\n수양산에 들어가 고사리를 캐 먹다 굶어 죽음',
    '"천도시야비야 — 천도는 옳은가 그른가"\n사기 전체를 관통하는 첫 화두')

make_person_slide('2', '관중·안영', '管仲·晏嬰', '제(齊) 환공·경공 시대',
    '명재상의 두 전형',
    '관중 — 환공을 도와 춘추 최초의 패업\n"관포지교(管鮑之交)" · "창고가 차야 예절을 안다"\n안영 — "남귤북지(南橘北枳)" · 키 작은 명재상',
    '"이런 재상을 만나는 것은 임금의 복"\n경제가 예의 토대임을 선언한 정치학')

make_person_slide('3', '오자서', '伍子胥', '춘추 말 · 초→오→월',
    '복수의 화신, 충신의 비극',
    '초 평왕이 아버지·형 처형 → 오로 망명\n수도 영(郢) 함락 후 평왕 시신을 300번 매질\n오왕 부차의 자결 명령 — "내 눈을 동문에 걸어라"',
    '"日暮途遠 倒行逆施" — 해는 저물고 길은 멀어\n사마천 자신(궁형)과 깊이 공명하는 인물')

make_person_slide('4', '공자', '孔子', 'BC 551~479 · 노(魯)',
    '세가에 오른 유일한 사상가',
    '포의(布衣)임에도 세가(世家)에 배치 — 전무후무한 파격\n사마천이 곡부에 가서 공자의 묘당·수레·의복을 보고\n"머물러 차마 떠나지 못하였다"',
    '"高山仰止 景行行止 — 높은 산을 우러러"\n권력이 아니라 사상의 영향력이 역사의 척도')

make_person_slide('5', '진시황', '秦始皇', 'BC 259~210 · 진(秦)',
    '통일 제국의 설계자',
    '6국 통일(BC 221) — 황제 칭호 창제\n군현제·문자·도량형·수레 궤간 통일\n분서갱유·아방궁·여산릉 → 사후 15년 멸망',
    '가의 「과진론」 — "인의를 베풀지 않았으니\n공수의 형세가 달라진 것이다"')

make_person_slide('6', '항우', '項羽', 'BC 232~202 · 초(楚)',
    '본기에 오른 비극의 영웅',
    '진을 멸하고 천하 제패 → 해하(垓下)에서 사면초가\n우희와 이별 — "力拔山兮氣蓋世"\n오강(烏江)에서 자결 (31세)',
    '"왜 본기에 넣었는가" — 실질 권력자였기에\n"근고 이래 미증유한 인물"의 비극미')

make_person_slide('7', '한 고조 유방', '漢高祖 劉邦', 'BC 256~195 · 한(漢) 태조',
    '정장(亭長)에서 황제로',
    '패현의 건달 출신 · 60세에 황제\n낙양 잔치의 자기 평가 — "장량·소하·한신은 인걸\n내가 그들을 쓸 수 있었기에 천하를 얻었다"',
    '"혼자 잘난 자가 아니라 남을 쓸 줄 아는 자"\n사기가 제시하는 리더십의 최종 답')

make_person_slide('8', '한초삼걸', '漢初三傑', '장량·소하·한신 — 유방의 세 인걸',
    '책략·내정·전쟁의 세 천재',
    '장량 — 박랑사 진시황 암살 시도 · 황석공의 태공병법\n소하 — 진의 관부에서 도서·율령·지도를 챙김\n한신 — 배수진·정형·해하 — 백전백승',
    '"세 인걸"의 협력이 한 천하를 만들었다\n그러나 한신은 처세를 몰랐다')

make_person_slide('9', '한신의 토사구팽', '韓信', '회음후(淮陰侯) · BC ?~196',
    '"鳥盡弓藏 兔死狗烹" — 사기 전편의 경고',
    '빨래터 노파의 밥(一飯千金) · 가랑이 밑(胯下之辱)\n소하월하추한신 → 대장군 임명\n배수진 · 정형 · 해하 → 모반의 누명 → 처형',
    '"교활한 토끼가 죽으면 사냥개도 삶긴다"\n전장의 천재가 정치의 생리를 모를 때')

make_person_slide('10', '여불위', '呂不韋', '진(秦) · BC ?~235',
    '"奇貨可居" — 상인에서 승상으로',
    '조나라에서 진 인질 자초(子楚)에게 투자\n조희(趙姬)를 자초에게 바침 → 그 아들이 진시황\n승상 · 문신후 · 『여씨춘추』 — 일자천금(一字千金)',
    '"최대의 투자는 사람에게 한다"\n그러나 권력 본체가 되는 순간 제거된다')

make_person_slide('11', '굴원', '屈原', 'BC 343?~278? · 초(楚)',
    '충신의 비극, 문학의 시조',
    '회왕에게 합종을 권했으나 추방 → 양왕에게 다시 쫓겨남\n「이소」·「천문」·「구가」 — 울분의 문학\n"擧世皆濁我獨淸 衆人皆醉我獨醒" → 멱라 자결',
    '"진흙에서 매미가 흐린 껍질을 벗듯\n그 지조는 해와 달과 빛을 다투어도 좋다"')

make_person_slide('12', '전국사공자', '戰國四公子', '맹상군·평원군·신릉군·춘신군',
    '용인(用人)의 네 전범',
    '맹상군 — 식객 3,000 · 계명구도(鷄鳴狗盜)\n평원군 — 모수자천 · 첩의 목을 베어 식객 회유\n신릉군 — 절부구조(竊符救趙) · 백정 주해\n춘신군 — 말년에 음모로 암살',
    '"용인의 최고 수준은 신분·학벌·외모를 보지 않는 것"\n사기가 네 사람의 대비로 입체화')

make_person_slide('13', '염파·인상여', '廉頗·藺相如', '조(趙) · 장상(將相)의 덕',
    '문경지교(刎頸之交)의 모범',
    '인상여 — 완벽귀조(完璧歸趙) · 민지지회(澠池之會)\n염파의 불만 → 인상여가 피해 다님\n"先國家之急而後私讎也" → 염파의 부형청죄',
    '"공(公)을 앞에 두고 사(私)를 뒤에 두는 신하"\n"자기 잘못을 인정하는 장군"')

make_person_slide('14', '상앙·이사', '商鞅·李斯', '법가의 두 거장 — 비극의 계열',
    '"법으로 세운 자는 법으로 망한다"',
    '상앙 — 진 효공의 변법 · 이목지신 · 거열형\n이사 — 진시황 6국 통일 · 분서갱유 건의\n2세 호해 옹립 · 조고의 모함으로 요참(腰斬)\n"황견지탄 — 너와 함께 토끼 잡으러 가고 싶구나"',
    '"법가는 강국을 만드나 자신을 보호하지 못한다"\n사마천의 반복되는 아이러니')

make_person_slide('15', '손무·오기·손빈', '孫武·吳起·孫臏', '병법의 3대 거장',
    '재능과 비극',
    '손무 — 손자병법 · 합려 앞에서 궁녀 두 명을 베다\n오기 — 살처구장(殺妻求將) · 도왕 사후 암살\n손빈 — 방연의 질투로 빈형(臏刑) · 마릉 전투로 복수',
    '"병법의 천재는 인간관계의 얕음 때문에 몰락"\n"실력보다 시기(猜忌)가 더 무서운 적"')

make_person_slide('16', '형가', '荊軻', '자객열전의 주인공',
    '"風 蕭 蕭 兮 易 水 寒"',
    '연 태자 단의 부탁 → 진시황 암살 시도\n번오기의 목 + 독이 발린 비수를 지도에 감춤\n"圖窮匕見" — 마지막 일격을 놓치고 참수',
    '"바람 소슬한데 역수는 차갑구나\n장사 한 번 가면 다시 돌아오지 못하리"')


# ============== Ⅶ. 7대 명장면 ==============
def make_scene_slide(num, title, scene, location, quote):
    @S(f'Ⅶ. 7대 명장면 ({num}/7)')
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, f'Ⅶ. 7대 명장면 ({num}/7)', page, total)
        add_title(slide, f'명장면 {num} — {title}', location)
        add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.7), PALE)
        add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.4),
                    '장면', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(2.9), Inches(11.7), Inches(2.0),
                    scene, font_size=15, color=INK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.9),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.4),
                    '명구·의미', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1.2),
                    quote, font_size=16, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return renderer


make_scene_slide('1', '요·순의 선양(禪讓)',
    '요(堯)가 천하를 자식이 아니라 어진 순(舜)에게 물려줌\n순도 다시 우(禹)에게 — 권력의 가장 이상적인 이전',
    '1 오제본기 · 사기 전체의 출발점',
    '"천하를 사사로이 하지 않는다(天下爲公)"\n— 권력의 이상형, 사기 전체가 변주하는 첫 모범')

make_scene_slide('2', '백이·숙제의 수양산 아사(餓死)',
    '주 무왕이 은 주왕을 치는 것을 만류 → 듣지 않자\n수양산에 들어가 고사리를 캐 먹다 굶어 죽음',
    '1 백이열전 · 열전의 첫 편',
    '"天 道 是 邪 非 邪 — 천도는 옳은가 그른가"\n— 사기 전체를 관통하는 화두')

make_scene_slide('3', '항우의 사면초가·패왕별희',
    '해하(垓下) 포위 · 사방에서 초나라 노래\n우희와 이별의 노래 → 우희 자결 → 오강 자결',
    '7 항우본기 · 영웅 비극의 정점',
    '"力拔山兮氣蓋世 — 힘은 산을 뽑고 기개는 세상을 덮건만\n時不利兮騅不逝 — 때가 이롭지 않으니 명마도 가지 않는구나"')

make_scene_slide('4', '유방의 자기 평가 — 삼걸을 썼기에 이겼다',
    '낙양 남궁 잔치 — "내가 천하를 얻은 까닭은?"\n"장량보다 못하나 그를 썼고\n소하·한신을 썼기에 천하를 얻었다"',
    '8 고조본기 · 리더십의 최종 답',
    '"夫運籌帷帳之中 決勝千里之外 吾不如子房…\n項羽有一范增而不能用 此其所以爲我擒也"')

make_scene_slide('5', '한신의 토사구팽',
    '연·제·조 멸하고 백만 대군의 장수\n그러나 유방의 의심을 사 여후에게 목 베임\n임종 직전 "교활한 토끼가 죽으면 사냥개가 삶긴다"',
    '32 회음후열전 · 공성신퇴 실패의 비극',
    '"鳥 盡 弓 藏  兔 死 狗 烹"\n— 사기 전편의 가장 강력한 경고')

make_scene_slide('6', '형가의 풍소소혜역수한',
    '진시황 암살을 떠나는 형가의 이별\n역수(易水) 강가에서 부른 노래\n진 왕궁의 "圖窮匕見" → 일격 실패 → 참수',
    '26 자객열전 · 사위지기자사의 극한',
    '"風 蕭 蕭 兮 易 水 寒\n壯 士 一 去 兮 不 復 還"')

make_scene_slide('7', '사마천의 궁형과 태사공자서',
    '이릉을 변호하다 궁형 → 죽음을 택하지 않고 사기를 완성\n70 태사공자서 — 자기 자신을 130번째 인물로 기록\n"명산에 감추고 후세의 군자를 기다린다"',
    '70 태사공자서 · 사기 전체의 서명(署名)',
    '"究 天 人 之 際   通 古 今 之 變   成 一 家 之 言"\n— 사기 최고 이념 · 사마천의 자기 선언')


# ============== Ⅷ. 명구 10선 ==============
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=42):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=20, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('Ⅷ. 명구 (1/10)',
    '究 天 人 之 際   通 古 今 之 變   成 一 家 之 言',
    '구천인지제 · 통고금지변 · 성일가지언',
    '천인을 궁구하고 고금에 통달하여 일가의 말을 이룬다',
    '보임안서 · 태사공자서 — 사기 최고 이념', hanmun_size=22), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (2/10)',
    '人 固 有 一 死\n或 重 於 泰 山   或 輕 於 鴻 毛',
    '인고유일사 · 혹중어태산 혹경어홍모',
    '사람은 한 번 죽는다. 태산보다 무거울 수도, 기러기 털보다 가벼울 수도 있다',
    '사마천 「보임안서」', hanmun_size=24), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (3/10)',
    '天 道 是 邪 非 邪',
    '천 도 시 야 비 야',
    '천도는 옳은가 그른가 — 사기 전체의 화두',
    '1 백이열전', hanmun_size=72), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (4/10)',
    '倉 廩 實 則 知 禮 節\n衣 食 足 則 知 榮 辱',
    '창름실즉지예절 · 의식족즉지영욕',
    '창고가 차야 예절을 알고, 의식이 족해야 영욕을 안다',
    '관중 → 2 관안열전 · 69 화식열전', hanmun_size=28), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (5/10)',
    '士 爲 知 己 者 死',
    '사 위 지 기 자 사',
    '선비는 자기를 알아준 자를 위해 죽는다',
    '예양 → 26 자객열전', hanmun_size=68), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (6/10)',
    '王 侯 將 相   寧 有 種 乎',
    '왕후장상 영유종호',
    '왕후장상이 어찌 씨가 있으랴 — 평민 봉기의 외침',
    '진승 → 18 진섭세가', hanmun_size=44), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (7/10)',
    '鳥 盡 弓 藏   兔 死 狗 烹',
    '조진궁장 · 토사구팽',
    '새가 다 잡히면 활을 감추고, 토끼가 죽으면 개를 삶는다',
    '한신 → 32 회음후열전', hanmun_size=44), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (8/10)',
    '力 拔 山 兮 氣 蓋 世\n時 不 利 兮 騅 不 逝',
    '역발산혜기개세 · 시불리혜추불서',
    '힘은 산을 뽑고 기개는 세상을 덮건만, 때가 이롭지 않으니 명마도 가지 않는구나',
    '항우 → 7 항우본기 (해하가)', hanmun_size=24), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (9/10)',
    '桃 李 不 言   下 自 成 蹊',
    '도리불언 · 하자성혜',
    '복숭아와 오얏은 말이 없어도 그 아래 저절로 길이 난다 — 덕이 있는 자에게는 사람이 모인다',
    '이광 → 49 이장군열전', hanmun_size=44), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (10/10)',
    '天 下 熙 熙   皆 爲 利 來\n天 下 壤 壤   皆 爲 利 往',
    '천하희희 개위리래 · 천하양양 개위리왕',
    '천하가 들썩이는 것은 모두 이익 때문 — 인간 경제의 본성',
    '69 화식열전', hanmun_size=22), 'Ⅷ. 명구'))


# ============== Ⅸ. 사마천의 사상 — 삼구(三句) ==============
@S('Ⅸ. 사상')
def s_3gu(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 사상', page, total)
    add_title(slide, '삼구(三句) — 사기의 최고 이념',
              '"보임안서·태사공자서에서 자기 사상을 한 줄로 선언"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.5), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
                '究 天 人 之 際',
                font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.7),
                '通 古 今 之 變   成 一 家 之 言',
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('구천인지제', '究天人之際',  '천(天)과 인(人)의 관계를 궁구한다'),
        ('통고금지변', '通古今之變',  '옛날과 지금의 변화를 꿰뚫는다'),
        ('성일가지언', '成一家之言',  '한 집안(가문·학파)의 말을 이룬다'),
    ]
    top = 4.0
    for eum, han, mean in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    han, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(2.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.22), Inches(2.0), Inches(0.5),
                    eum, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.3), Inches(top), Inches(7.5), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.5), Inches(top + 0.22), Inches(7.2), Inches(0.5),
                    mean, font_size=14, color=INK)
        top += 0.95


@S('Ⅸ. 사상')
def s_chuncheonin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 사상', page, total)
    add_title(slide, '① 究天人之際 — 하늘과 인간의 관계를 궁구한다')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"천도는 정의로운가?"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('왜 백이는 굶어 죽고 도척은 장수했는가?',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('이릉은 분전했는데 사마천은 궁형을 받았는가?',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('사마천의 응답',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"답을 주는 책이 아니라 질문을 던지는 책"',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"천도는 선인을 돕는다"는 상투적 답을 거부',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('대신 그들의 이름을 기록함으로써 천도를 대신한다',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅸ. 사상')
def s_tonggogeum(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 사상', page, total)
    add_title(slide, '② 通古今之變 — 옛날과 지금의 변화를 꿰뚫는다')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"제도는 끊임없이 변하되, 흥망에는 보편 패턴이 있다"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('선양 → 세습 (요·순·우 → 하)',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('봉건 → 군현 (주 → 진)',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('무위 → 유술 (문경 → 무제)',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('표(表) 10편이 이 주제의 구체적 실현',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('공시적(共時的) 시선으로 변화를 추적',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"같은 시대 여러 나라가 무엇을 했는가"의 동시적 파노라마',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('통감의 1,362년 통사와 일맥상통하는 사관',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅸ. 사상')
def s_taesagongwal(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 사상', page, total)
    add_title(slide, '③ 成一家之言 — 한 집안의 말을 이룬다 · 태사공왈',
              '"각 편 끝의 130편의 평(評)이 누적되어 사마천의 사상이 된다"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
                '"太 史 公 曰"  태사공이 말한다',
                font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.55), Inches(11.7), Inches(1.5), [
        ('자료 편찬이 아니라 사마천 자신의 역사철학을 담은 저작',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('포폄(褒貶)과 직필(直筆) — 황제도 거리낌이 없다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.45), Inches(11.7), Inches(1.6), [
        ('• 실패한 자를 감싸고 성공한 자를 의심한다',
         {'font_size': 14, 'space_before': 6}),
        ('• 황제를 향해서도 무제 봉선·흉노 정책을 비판',
         {'font_size': 14, 'space_before': 4}),
        ('• 자기 자신(궁형)도 비평 대상에 올린다',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅹ. 후대 영향 ==============
@S('Ⅹ. 후대 영향')
def s_24sa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '24사(史)의 머리 — 정사의 표준')
    items = [
        ('사기',         '사마천 · BC 91경',     '기전체 통사의 시조 — 130편'),
        ('한서',         '반고 · AD 82경',       '단대사(斷代史)의 시조 — 100편'),
        ('후한서',       '범엽 · 5세기',         '동한 13편 + 열전 80편'),
        ('삼국지',       '진수 · 3세기',         '위·촉·오 삼국'),
        ('진서·송서·…', '남북조 정사',           '시대마다 같은 체제'),
        ('당서·…',      '당대 이후',             '정사 표준 굳어짐'),
        ('명사',         '청 장정옥 · 1739',     '24사의 마지막'),
    ]
    top = 2.3
    for tag, author, char in items:
        c = ACCENT if '사기' in tag else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if '사기' in tag else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.55), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(3.8), Inches(0.55), bg)
        add_textbox(slide, Inches(3.35), Inches(top + 0.13), Inches(3.6), Inches(0.4),
                    author, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.1), Inches(top), Inches(5.7), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.25), Inches(top + 0.13), Inches(5.5), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.62
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '사기 → 한서 → 24사 — 2,000년 동양 정사 체제의 출발점',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅹ. 후대 영향')
def s_korea(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '한국 — 김부식 『삼국사기』(1145)의 정면 계승',
              '"사기 5체 → 삼국사기 4체의 직계 후예"')
    rows = [
        ('사기 (BC 91경)',          '본기 12 + 표 10 + 서 8 + 세가 30 + 열전 70 = 130편'),
        ('한서 (AD 82경)',          '본기 12 + 표 8 + 지 10 + 열전 70 = 100편'),
        ('삼국사기 (1145)',         '본기 28 + 연표 3 + 지 9 + 열전 10 = 50권'),
    ]
    top = 2.5
    for tag, content in rows:
        is_kr = '삼국' in tag
        c = ACCENT if is_kr else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_kr else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.75), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.18), Inches(3.5), Inches(0.4),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.2), Inches(top), Inches(8.6), Inches(0.75), bg)
        add_textbox(slide, Inches(4.35), Inches(top + 0.18), Inches(8.3), Inches(0.4),
                    content, font_size=13, bold=is_kr, color=INK)
        top += 0.85
    add_filled_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.6), [
        ('"본기·연표·지(志)·열전" 4체 구조 — 사기 5체에서 표·서를 합친 형태',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('한국 고대사 1차 사료 — 사기의 「조선열전」과 함께 동방 역사의 두 축',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅹ. 후대 영향')
def s_japan(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '일본·동아시아 — 사기 한자 교양의 확산')
    items = [
        ('헤이안 시대',  '794~1185',
         '귀족·승려의 한문 교양서 · 일본에 정착'),
        ('에도 시대',    '1603~1868',
         '무사 계급 교양의 핵심 · 막부의 사기 강독'),
        ('라이 산요',    '1780~1832',
         '『일본외사』 — 사기 체제를 모방한 일본 통사'),
        ('시바 료타로',  '1923~1996',
         '필명이 사마천에서 — "司馬(사마)에 미치지 못하는(遼) 太郞"'),
        ('베트남·류큐',  '한자 문화권',
         '동일하게 사기를 정사의 표준으로 학습'),
    ]
    top = 2.4
    for era, when, char in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.2), Inches(2.5), Inches(0.4),
                    era, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(2.5), Inches(0.8), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.2), Inches(2.5), Inches(0.4),
                    when, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.8), Inches(top), Inches(7.0), Inches(0.8),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.0), Inches(top + 0.2), Inches(6.7), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.9


@S('Ⅹ. 후대 영향')
def s_sagi_jachi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '사기 vs 자치통감 — 동양 통사의 두 축')
    rows = [
        ('편자·시기', '사마천 · BC 91경',          '사마광 · 1084'),
        ('형식',       '기전체(紀傳體) — 인물 중심',  '편년체(編年體) — 연대 중심'),
        ('주인공',    '한 사람의 일생',              '같은 시간의 사건들'),
        ('분량',       '130편 · 약 52만 자',          '294권 · 약 300만 자'),
        ('수록',       '약 3,000년 (황제~한 무제)',   '1,362년 (BC 403~959)'),
        ('관점',       '"인간이 어떻게 살았는가"',    '"왜 흥하고 망하는가"'),
        ('특징',       '개인의 입체 조명',            '거시적 흐름의 분석'),
    ]
    top = 1.95
    add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.2), Inches(0.55), SUB)
    add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(2.2), Inches(0.4),
                '항목', font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(2.75), Inches(top), Inches(5.0), Inches(0.55), ACCENT)
    add_textbox(slide, Inches(2.75), Inches(top + 0.1), Inches(5.0), Inches(0.4),
                '사기(史記)', font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(7.8), Inches(top), Inches(5.0), Inches(0.55), SUB)
    add_textbox(slide, Inches(7.8), Inches(top + 0.1), Inches(5.0), Inches(0.4),
                '자치통감', font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    row_h = 0.65
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(y), Inches(2.2), Inches(row_h), PALE)
        add_textbox(slide, Inches(0.55), Inches(y + 0.17), Inches(2.1), Inches(0.4),
                    row[0], font_size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.75), Inches(y), Inches(5.0), Inches(row_h), bg)
        add_textbox(slide, Inches(2.9), Inches(y + 0.17), Inches(4.7), Inches(0.4),
                    row[1], font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.8), Inches(y), Inches(5.0), Inches(row_h), bg)
        add_textbox(slide, Inches(7.95), Inches(y + 0.17), Inches(4.7), Inches(0.4),
                    row[2], font_size=13, color=INK, align=PP_ALIGN.CENTER)


@S('Ⅹ. 후대 영향')
def s_literature(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '문학적 영향 — 인물 묘사의 원형')
    items = [
        ('당송팔대가',     '한유·유종원·소식 등',  '사한문장(史漢文章) — 중국 고문의 절정'),
        ('삼국지연의',     '나관중 · 14세기',      '인물 묘사·고사·일화 작법이 사기에서'),
        ('수호전',         '시내암 · 14세기',      '108호걸의 입체 조명 — 열전 작법 계승'),
        ('서유기·금병매', '16세기',                '캐릭터 작법의 토대'),
        ('현대 전기 문학', 'biography 일반',       '모든 인물 평전의 원형'),
    ]
    top = 2.4
    for tag, author, char in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.8), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.15), Inches(2.8), Inches(0.4),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.5), Inches(top), Inches(3.0), Inches(0.7), PALE)
        add_textbox(slide, Inches(3.5), Inches(top + 0.15), Inches(3.0), Inches(0.4),
                    author, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.6), Inches(top), Inches(6.2), Inches(0.7),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.8), Inches(top + 0.15), Inches(6.0), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.8
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '"고사(故事)·일화·말 한마디로 인물을 그린다"',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '— 사기의 캐릭터 작법이 동아시아 모든 인물 서사의 원형',
                font_size=12, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅺ. 현대 의의 ==============
@S('Ⅺ. 현대 의의')
def s_modern_biography(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ① — 인물 중심 사관 · 전기 문학의 원형')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"한 사람의 삶이 역사의 단위"라는 사기의 사관',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('현대의 모든 전기(biography)·평전·인물 서사가 이 사관 위에',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 응용',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('• Walter Isaacson의 잡스·아인슈타인·다빈치 평전',
         {'font_size': 14, 'space_before': 6}),
        ('• Robert Caro의 LBJ·로버트 모지스 평전 — "한 사람의 일생으로 시대를 그린다"',
         {'font_size': 14, 'space_before': 4}),
        ('• 시바 료타로의 역사 소설 — 사기의 인물 구도를 직접 차용',
         {'font_size': 14, 'space_before': 4}),
        ('• "그 사람의 일화·말 한마디가 모든 것을 보여준다"',
         {'font_size': 13, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_modern_micro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ② — 미시사·사회사의 선구')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"황제만이 아니라 자객·유협·점쟁이·해학가까지 역사의 주체"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('60대 자객열전 · 64 유협열전 · 66 골계열전',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('67 일자열전 · 68 귀책열전 · 45 편작창공열전',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 사회사·미시사와의 일치',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• Annales 학파의 일상생활사 — 사기의 서민 열전과 통함',
         {'font_size': 14, 'space_before': 6}),
        ('• 카를로 긴즈부르그 『치즈와 구더기』 — 미시사의 선구',
         {'font_size': 14, 'space_before': 4}),
        ('• 사기는 2,100년 전에 이미 "평범한 사람의 역사"를 시도',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_modern_economy(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ③ — 화식열전 · 평준서 — 경제사의 효시')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
                '"天下熙熙 皆爲利來 天下壤壤 皆爲利往"',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.4),
                '천하희희 개위리래 · 천하양양 개위리왕 — 화식열전',
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.55), Inches(11.7), Inches(1.5), [
        ('69 화식열전 — 상인을 역사의 주체로',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('범려(도주공)·자공·백규·의돈·오씨나 — 상인 전기',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('평준서 — 한 무제 시대 재정의 날카로운 분석',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.45), Inches(11.7), Inches(1.5), [
        ('현대 의의',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"2,100년 전에 경제가 역사의 엔진임을 선언" — 사기의 가장 현대적인 면',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('근대 경제사의 효시로 재평가 — 막스 베버보다 2,000년 앞선 통찰',
         {'font_size': 13, 'color': INK, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_modern_minority(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ④ — 이민족 서술의 객관성 · 한국 고대사 1차 사료')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"이민족을 야만이 아니라 자기 질서를 가진 사회로"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('흉노·조선·남월·동월·서남이·대원 — 6개 이민족 열전',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"이후 24사 이민족 서술 전통이 여기서 시작"',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('한국에 대한 의미 — 「조선열전」',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 위만조선의 흥망 — 한국 고대사의 가장 오래된 1차 사료',
         {'font_size': 14, 'space_before': 6}),
        ('• 한 무제의 침공과 한사군 설치의 구체적 기록',
         {'font_size': 14, 'space_before': 4}),
        ('• 김부식 『삼국사기』가 사기의 형식을 빌려 우리 역사를 기록',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_modern_balbun(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ⑤ — 발분저서(發憤著書) · 고난과 창조')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"고난이 기록을 낳는다" — 사마천의 자기 선언',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"문왕이 갇혀 주역을, 중니가 곤하여 춘추를…"',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('역사는 성공한 자가 아니라 실패와 고난에서 펴 올린다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 적용 — Post-traumatic growth(외상 후 성장)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 빅터 프랭클 『죽음의 수용소에서』 — 발분저서의 현대 사례',
         {'font_size': 14, 'space_before': 6}),
        ('• 솔제니친 『이반 데니소비치의 하루』 — 굴라크의 발분',
         {'font_size': 14, 'space_before': 4}),
        ('• "고난이 가장 위대한 창조의 출발점"임을 사기가 증명',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_modern_what(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ⑥ — 사기가 우리에게 던지는 질문')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('사기를 다 읽고 나면',
         {'font_size': 18, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('독자는 결국 사마천과 마주 앉는다',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('그 앞에서 묻게 된다',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"나는 역사에 무엇으로 남을 것인가"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"내 죽음은 태산처럼 무거울 것인가',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        (' 기러기 털처럼 가벼울 것인가"',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('이 질문을 던지는 순간 — 사기는 2,100년 된 자료가 아니라 오늘의 거울이 된다',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


# ============== Ⅻ. 마무리 ==============
@S('Ⅻ. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 사기')
    add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(11.9), Inches(5.0), PALE)
    add_paragraphs(slide, Inches(1.1), Inches(2.2), Inches(11.1), Inches(4.7), [
        ('궁형(宮刑)을 당한 사관이 죽음을 택하는 대신',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('"말을 이어가기" 위해 14~15년에 걸쳐 쓴 책',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('황제(黃帝)에서 한 무제까지 3,000년을',
         {'font_size': 16, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('본기·표·서·세가·열전 5체의 입체 구조로 재구성',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('격식이 아닌 실질로 사람을 보고',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('(항우·여후·공자·진섭의 파격)',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('성공자가 아닌 실패자에 공감하며',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('(굴원·이광·한신)',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('"천도시야비야" — 부조리를 인정하면서도',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('의인의 이름을 남기는 것으로 천도를 대신한다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.2)


@S('Ⅻ. 마무리')
def s_3gu_final(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 마무리', page, total)
    add_title(slide, '사마천의 삼구(三句) — 사기의 영원한 이념')
    add_filled_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.9),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.6), [
        ('究 天 人 之 際',
         {'font_size': 30, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('하늘과 인간의 관계를 궁구한다',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 8}),
        ('通 古 今 之 變',
         {'font_size': 30, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('옛날과 지금의 변화를 꿰뚫는다',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 8}),
        ('成 一 家 之 言',
         {'font_size': 32, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('한 집안의 말을 이룬다',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 8}),
        ('— 사마천이 보임안서·태사공자서에서 직접 선언한 사기의 최고 이념 —',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 12}),
    ], line_spacing=1.2)


@S('Ⅻ. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.5),
                '太 史 公 曰',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '태 사 공 왈', font_size=24, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.7),
                '"태사공(太史公)이 말한다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '— 각 편 끝의 130편 사평(史評)이 누적되어 사마천의 사상이 된다 —',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '"나는 역사에 무엇으로 남을 것인가"',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '— 사기가 2,100년 건너 우리에게 던지는 질문 —',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '감사합니다', font_size=24, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\사기.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')