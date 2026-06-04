# -*- coding: utf-8 -*-
"""
주역(周易) 발표자료 — 망라적 76장 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 1. 표지 ==============
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8),
                '周 易', font_size=140, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                'The Book of Changes · 주역(周易)',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.5),
                '복희·문왕·주공·공자 — 三古三聖이 완성한 변화의 책',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
                '8괘 → 64괘 384효 · 십익(十翼) 7편 · 군경지수(群經之首)',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                '"一 陰 一 陽 之 謂 道"',
                font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                '한 번 음하고 한 번 양하는 것을 도(道)라 한다 — 계사전',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 2. 목차 ==============
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.8), Inches(0.7),
                '목 차', font_size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.15), Inches(12.8))
    items = [
        [('Ⅰ', '개요 — 주역이란 무엇인가'),
         ('Ⅱ', '6층위 구조 — 음양·8괘·64괘·6효'),
         ('Ⅲ', '십익(十翼) — 공자의 날개'),
         ('Ⅳ', '64괘의 큰 흐름 — 6막의 서사시'),
         ('Ⅴ', '10대 핵심 사상'),
         ('Ⅵ', '핵심 괘 9선')],
        [('Ⅶ', '64괘 주제별 클러스터'),
         ('Ⅷ', '계사전 12명구'),
         ('Ⅸ', '해석의 두 학파 — 상수·의리'),
         ('Ⅹ', '후대 영향 — 동아시아 3,000년'),
         ('Ⅺ', '현대적 의의'),
         ('Ⅻ', '마무리')],
    ]
    for col, group in enumerate(items):
        x = 0.7 + col * 6.4
        top = 1.6
        for num, title in group:
            add_textbox(slide, Inches(x), Inches(top), Inches(1.0), Inches(0.5),
                        num, font_size=17, bold=True, color=ACCENT)
            add_textbox(slide, Inches(x + 1.0), Inches(top), Inches(5.3), Inches(0.5),
                        title, font_size=17, color=INK)
            top += 0.7


# ============== Ⅰ. 개요 ==============
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '주역(周易)이란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                '복희·문왕·주공·공자가 완성한 "변화(易)"의 책',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.5),
                '"변화 속의 불변, 불변 속의 변화"를 음양 두 기호로 압축',
                font_size=15, color=ACCENT, align=PP_ALIGN.CENTER)
    nums = [('64', '괘(卦)'), ('384', '효(爻)'), ('7', '편 십익(十翼)'), ('3,000', '년의 책')]
    for i, (n, lbl) in enumerate(nums):
        x = 0.6 + i * 3.05
        add_textbox(slide, Inches(x), Inches(4.0), Inches(2.9), Inches(1.0),
                    n, font_size=48, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(5.15), Inches(2.9), Inches(0.5),
                    lbl, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                '점서(占書)에서 시작 → 공자의 십익으로 철학서로 승격',
                font_size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '"군경지수(群經之首) · 대도지원(大道之源)" — 모든 경전의 머리, 큰 도의 근원',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_status(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '책의 위상 — 군경지수(群經之首)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '한대 이후 오경(五經)의 첫머리 — 모든 경전의 머리',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('주역(周易)', '천지(天地)의 원리',     '"우주의 원형 구조"'),
        ('시(詩)',     '인사(人事)의 노래',     '문학·정서'),
        ('서(書)',     '인사(人事)의 기록',     '정치·역사'),
        ('예(禮)',     '인사(人事)의 규범',     '윤리·제도'),
        ('춘추(春秋)', '인사(人事)의 평가',     '역사 비평'),
    ]
    top = 3.4
    for tag, role, char in items:
        is_yi = '주역' in tag
        c = ACCENT if is_yi else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_yi else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.5), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(2.5), Inches(0.4),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(4.0), Inches(0.5), bg)
        add_textbox(slide, Inches(3.35), Inches(top + 0.08), Inches(3.8), Inches(0.4),
                    role, font_size=13, bold=is_yi, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.3), Inches(top), Inches(5.5), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.45), Inches(top + 0.08), Inches(5.3), Inches(0.4),
                    char, font_size=13, color=INK)
        top += 0.55
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"인사(人事)의 책들" 위에 "천지의 원리" 주역이 놓인다 — 메타 경전',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_jueyi_meaning(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '"주역"의 의미 — 周와 易의 풀이')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '周 — 주(周)의 두 해석', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('① 주(周)나라의 역', {'font_size': 15, 'bold': True}),
        ('   주 문왕·주공이 다듬음', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('② "두루 미친다(周, 遍也)"', {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('   만물에 빠짐없이 미치는 역', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('둘 다 틀리지 않다', {'font_size': 13, 'color': INK, 'bold': True, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '易 — 역(易)의 세 의미', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('① 변역(變易)', {'font_size': 15, 'bold': True}),
        ('   만물은 변한다', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('② 불역(不易)', {'font_size': 15, 'bold': True, 'space_before': 4}),
        ('   변화 속에 불변의 법칙', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('③ 간역(簡易)', {'font_size': 15, 'bold': True, 'space_before': 4}),
        ('   음양 두 기호로 간단히', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('세 의미가 동시에 작동', {'font_size': 13, 'color': ACCENT, 'bold': True, 'space_before': 6}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '실전된 連山易(하)·歸藏易(은) → 周易만 남음 — 三易 중 유일한 생존자',
                font_size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_3go3seong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '삼고삼성(三古三聖) — 4성인이 합작한 책',
              '"人更三聖 世歷三古" — 한서 예문지')
    items = [
        ('상고(上古)', '복희(伏羲)',  '팔괘(八卦)를 그어 천지 만물을 상징화'),
        ('중고(中古)', '문왕(文王)',  '유리(羑里) 옥중에서 팔괘를 겹쳐 64괘 · 괘사(卦辭)'),
        ('중고(中古)', '주공(周公)',  '각 괘 6효에 효사(爻辭)를 붙임'),
        ('하고(下古)', '공자(孔子)',  '십익(十翼) — 점서를 철학서로 승격'),
    ]
    top = 2.5
    for era, name, role in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.9), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(2.5), Inches(0.5),
                    era, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(2.8), Inches(0.9), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.27), Inches(2.8), Inches(0.5),
                    name, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.1), Inches(top), Inches(6.7), Inches(0.9),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.3), Inches(top + 0.27), Inches(6.4), Inches(0.5),
                    role, font_size=13, color=INK)
        top += 1.0
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"4성 공동 저술" 서사가 주역의 권위를 받친다 — 한 사람이 아닌 문명의 책',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_jeomseo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '점서(占書)에서 철학서로 — 공자의 십익')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '본래 — 점서(占書)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('거북점(卜) · 시초점(筮)', {'font_size': 15, 'bold': True}),
        ('', {'font_size': 6}),
        ('50개의 시초(蓍草) 줄기를 뽑아',
         {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ('괘를 얻고 길흉을 판단',
         {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('"실용적 점복서"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '공자 이후 — 철학서', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"君子居則觀其象而玩其辭', {'font_size': 13, 'bold': True}),
        (' 動則觀其變而玩其占"', {'font_size': 13, 'bold': True}),
        ('— 계사상', {'font_size': 11, 'color': SUB}),
        ('', {'font_size': 6}),
        ('"군자는 거하면 그 상을 보고', {'font_size': 12, 'color': SUB, 'space_before': 6}),
        (' 움직이면 그 변을 본다"', {'font_size': 12, 'color': SUB}),
        ('', {'font_size': 6}),
        ('"점치지 않고도 읽히는 책', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('(不占而已)"', {'font_size': 14, 'bold': True, 'color': ACCENT}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '"길흉을 묻는 책에서 어떻게 살 것인가를 묻는 책으로"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 6층위 구조 ==============
@S('Ⅱ. 6층위')
def s_6stack(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 6층위', page, total)
    add_title(slide, '6층위 구조 — 음양에서 십익까지')
    items = [
        ('①', '음양(陰陽)',       '최소 단위 — 양효(⚊) · 음효(⚋)'),
        ('②', '팔괘(八卦)',       '3효의 조합 — 2³ = 8가지 원형'),
        ('③', '64괘(大成卦)',    '팔괘의 쌍 — 8×8 = 64'),
        ('④', '6효의 자리·관계',  '位·中·正·應·比 — 상황과 관계의 시스템'),
        ('⑤', '괘사·효사',         '본문 — 길흉회린의 판단'),
        ('⑥', '십익(十翼)',       '공자의 7편 해설 — 점서를 철학으로'),
    ]
    top = 2.3
    for num, tag, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(0.7), Inches(0.4),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(3.5), Inches(0.6), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.13), Inches(3.4), Inches(0.4),
                    tag, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.0), Inches(top), Inches(7.8), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.2), Inches(top + 0.13), Inches(7.5), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.7
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"6층위가 동시에 작동" — 한 효를 읽을 때 모든 층위가 함께 의미를 만든다',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 6층위')
def s_eumyang(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 6층위', page, total)
    add_title(slide, '① 음양(陰陽) — 최소 단위')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '陽 — 양효 (⚊)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(3.1), Inches(5.9), Inches(0.6),
                '⚊',
                font_size=80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('이어진 선', {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('하늘 · 남성 · 강(剛)', {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('동(動) · 활(活)', {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('숫자로 — 九(9, 노양)', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '陰 — 음효 (⚋)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.1), Inches(5.9), Inches(0.6),
                '⚋',
                font_size=80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('끊어진 선', {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('땅 · 여성 · 유(柔)', {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('정(靜) · 수(受)', {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('숫자로 — 六(6, 노음)', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '"대립이 아니라 상호 생성" — 양 속에 음이 숨고 음 속에 양이 싹튼다',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 6층위')
def s_8gwae(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 6층위', page, total)
    add_title(slide, '② 팔괘(八卦) — 3효의 8가지 조합 · 자연의 8원형')
    items = [
        ('☰', '건(乾)', '天 하늘',  '健 굳셈',  '父'),
        ('☷', '곤(坤)', '地 땅',     '順 순함',  '母'),
        ('☳', '진(震)', '雷 우레',   '動 움직임','長男'),
        ('☴', '손(巽)', '風 바람',   '入 들어감','長女'),
        ('☵', '감(坎)', '水 물',     '陷 빠짐',  '中男'),
        ('☲', '이(離)', '火 불',     '明 밝음',  '中女'),
        ('☶', '간(艮)', '山 산',     '止 멈춤',  '少男'),
        ('☱', '태(兌)', '澤 못',     '說 기쁨',  '少女'),
    ]
    top = 2.2
    for i, (sym, name, nature, char, fam) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = top + row * 1.15
        add_filled_rect(slide, Inches(x), Inches(y), Inches(1.0), Inches(1.0), ACCENT)
        add_textbox(slide, Inches(x), Inches(y + 0.15), Inches(1.0), Inches(0.7),
                    sym, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 1.1), Inches(y), Inches(2.0), Inches(1.0), PALE)
        add_textbox(slide, Inches(x + 1.1), Inches(y + 0.3), Inches(2.0), Inches(0.5),
                    name, font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 3.2), Inches(y), Inches(2.7), Inches(1.0),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_paragraphs(slide, Inches(x + 3.3), Inches(y + 0.15), Inches(2.5),
                       Inches(0.8), [
                           (nature, {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER}),
                           (char, {'font_size': 12, 'color': SUB, 'align': PP_ALIGN.CENTER}),
                           (fam, {'font_size': 11, 'color': SUB, 'align': PP_ALIGN.CENTER}),
                       ], line_spacing=1.25)


@S('Ⅱ. 6층위')
def s_64gwae(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 6층위', page, total)
    add_title(slide, '③ 64괘(大成卦) — 팔괘의 쌍으로 이룬 6효')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('8 × 8 = 64괘 — 각 괘는 6효(상괘 3 + 하괘 3)',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('상괘(上卦) — 외(外) · 환경 · 결과 · 위 3효',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('하괘(下卦) — 내(內) · 개인 · 원인 · 아래 3효',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('예 — 11번 태괘(泰卦, 地天泰)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('상괘 ☷ 곤(땅) + 하괘 ☰ 건(하늘) — 땅이 위, 하늘이 아래',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('하늘의 기운은 위로, 땅의 기운은 아래로 → 두 기운이 가운데서 교류',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('→ 태평(泰平) · 소통',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"제자리"가 반드시 좋은 것이 아니다 — 교류·역전이 통(通)을 낳는다',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅱ. 6층위')
def s_6hyo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 6층위', page, total)
    add_title(slide, '④ 6효의 자리·관계 — 中·正·應·比의 시스템')
    rows = [
        ('상효(上爻)', '上九/上六', '은퇴자·초월자',     '지나침·노년'),
        ('오효(五爻)', '九五/六五', '★ 군주 자리',       '전성기 — "비룡재천"'),
        ('사효(四爻)', '九四/六四', '대신·측근',          '이행기'),
        ('삼효(三爻)', '九三/六三', '하급 지도자',        '위기·갈림길'),
        ('이효(二爻)', '九二/六二', '중견 관료',          '내면의 중심'),
        ('초효(初爻)', '初九/初六', '신참·평민',          '시작'),
    ]
    top = 2.0
    for tag, name, role, stage in rows:
        is_5 = '오효' in tag
        c = ACCENT if is_5 else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_5 else PALE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.0), Inches(0.5), c)
        add_textbox(slide, Inches(0.5), Inches(top + 0.08), Inches(2.0), Inches(0.4),
                    tag, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.6), Inches(top), Inches(2.0), Inches(0.5), bg)
        add_textbox(slide, Inches(2.6), Inches(top + 0.08), Inches(2.0), Inches(0.4),
                    name, font_size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.7), Inches(top), Inches(3.5), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.85), Inches(top + 0.08), Inches(3.3), Inches(0.4),
                    role, font_size=12, color=INK, bold=is_5)
        add_filled_rect(slide, Inches(8.3), Inches(top), Inches(4.5), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(8.45), Inches(top + 0.08), Inches(4.3), Inches(0.4),
                    stage, font_size=12, color=INK, bold=is_5)
        top += 0.55
    add_filled_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.4), [
        ('5가지 관계 — 位 · 中 · 正 · 應 · 比',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('當位(제자리) · 中(가운데) · 正(바름) · 應(호응) · 比(이웃)',
         {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"중정(中正)" — 가운데이면서 바른 자리, 주역 최고의 상태',
         {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅱ. 6층위')
def s_giljum(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 6층위', page, total)
    add_title(slide, '⑤ 길흉회린(吉凶悔吝) — 4단계 판단어')
    items = [
        ('吉', '길', '좋음',           ACCENT,
         '제자리에서 바르게 움직임 → 형통(亨通)'),
        ('凶', '흉', '나쁨',           SUB,
         '자리·때에 어긋난 행동 → 화(禍)'),
        ('悔', '회', '후회',           RGBColor(0xA0, 0x60, 0x40),
         '반성하면 길로 간다 — 흉의 회복 가능성'),
        ('吝', '린', '인색·막힘',      RGBColor(0x60, 0x40, 0xA0),
         '방치하면 흉으로 간다 — 길의 손실 위험'),
    ]
    top = 2.4
    for han, eum, mean, color, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.3), Inches(1.0), color)
        add_textbox(slide, Inches(0.6), Inches(top + 0.25), Inches(1.3), Inches(0.6),
                    han, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.0), Inches(top), Inches(2.5), Inches(1.0), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.15), Inches(2.5), Inches(0.4),
                    eum, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.0), Inches(top + 0.55), Inches(2.5), Inches(0.4),
                    mean, font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.6), Inches(top), Inches(8.2), Inches(1.0),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.8), Inches(top + 0.3), Inches(7.9), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 1.1
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"吉凶悔吝 生乎動者也" — 길흉회린은 모두 "움직임"에서 생긴다',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅲ. 십익 ==============
@S('Ⅲ. 십익')
def s_10ik_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 십익', page, total)
    add_title(slide, '십익(十翼) — 공자의 7편 10개 날개')
    rows = [
        ('1·2', '단전(彖傳) 상·하',     '괘사 해설',         '괘 전체의 덕성·시세 판단'),
        ('3·4', '상전(象傳) 상·하',     '괘상·효상 해설',     '"君子以~" 군자의 실천 규범'),
        ('5·6', '계사전(繫辭傳) 상·하', '주역 철학의 총론',   '★ 주역 사상의 정수'),
        ('7',   '문언전(文言傳)',       '건·곤 두 괘 해설',   '덕성 수양의 본보기'),
        ('8',   '설괘전(說卦傳)',       '팔괘의 상징 체계',   '상수(象數) 해석의 원천'),
        ('9',   '서괘전(序卦傳)',       '64괘의 배열 이유',   '괘의 연쇄 논리'),
        ('10',  '잡괘전(雜卦傳)',       '32쌍 괘의 대비',     '짝지어 파악하는 관점'),
    ]
    top = 1.95
    for num, name, role, char in rows:
        is_key = '계사' in name
        c = ACCENT if is_key else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_key else PALE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(1.0), Inches(0.55), c)
        add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(1.0), Inches(0.4),
                    num, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.6), Inches(top), Inches(3.5), Inches(0.55), bg)
        add_textbox(slide, Inches(1.7), Inches(top + 0.13), Inches(3.4), Inches(0.4),
                    name, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.2), Inches(top), Inches(2.8), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.35), Inches(top + 0.13), Inches(2.6), Inches(0.4),
                    role, font_size=12, color=INK)
        add_filled_rect(slide, Inches(8.1), Inches(top), Inches(4.7), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(8.25), Inches(top + 0.13), Inches(4.5), Inches(0.4),
                    char, font_size=12, color=INK, bold=is_key)
        top += 0.62
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"翼(날개)" — 경문 없이는 날지 못한다 · 주역을 철학으로 띄운 날개',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅲ. 십익')
def s_dan_sang(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 십익', page, total)
    add_title(slide, '단전(彖傳) · 상전(象傳) — 괘와 효의 해설')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '단전(彖傳) — 괘 전체의 판단', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"彖(단)"이란?', {'font_size': 14, 'bold': True, 'color': ACCENT}),
        ('  큰 돼지의 어금니 — "물어 끊는다"', {'font_size': 12, 'color': SUB}),
        ('  → "단정적으로 판단한다"', {'font_size': 12, 'color': SUB}),
        ('', {'font_size': 6}),
        ('역할', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('• 괘의 핵심 덕성', {'font_size': 13}),
        ('• 그 괘가 가리키는 시세(時勢)', {'font_size': 13}),
        ('• 괘사가 왜 그러한지 설명', {'font_size': 13}),
        ('', {'font_size': 6}),
        ('예 — 건괘 단전', {'font_size': 13, 'bold': True, 'color': INK, 'space_before': 6}),
        ('"大哉乾元 萬物資始" — 위대하다 건의 으뜸, 만물이 비롯한다',
         {'font_size': 11, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '상전(象傳) — 군자의 실천', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('대상(大象) + 소상(小象)', {'font_size': 14, 'bold': True, 'color': ACCENT}),
        ('  대상 — 괘 전체의 상', {'font_size': 12, 'color': SUB}),
        ('  소상 — 각 효의 상', {'font_size': 12, 'color': SUB}),
        ('', {'font_size': 6}),
        ('특징 — "君子以~" 형식', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('  "군자는 이로써 ~한다"', {'font_size': 13}),
        ('  자연 현상 → 군자의 실천 규범', {'font_size': 13, 'bold': True}),
        ('', {'font_size': 6}),
        ('예 — 건 상전', {'font_size': 13, 'bold': True, 'color': INK, 'space_before': 6}),
        ('"天行健 君子以自强不息"', {'font_size': 13, 'bold': True, 'color': ACCENT}),
        ('하늘의 운행이 강건하니 군자는 스스로 쉬지 않는다',
         {'font_size': 11, 'color': SUB}),
    ], line_spacing=1.3)


@S('Ⅲ. 십익')
def s_gyesa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 십익', page, total)
    add_title(slide, '계사전(繫辭傳) — 주역 철학의 정수',
              '"십익 중 가장 중요한 두 편 — 주역의 총론·우주론·인생론"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '"繫(계)" — 매다 · 辭(사) — 말씀 · "괘사·효사에 매단 큰 말씀"',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(3.4), [
        ('계사전이 다루는 5가지 큰 주제',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('① 도(道)·기(器)의 형이상학 — "形而上者謂之道"',
         {'font_size': 14, 'space_before': 6}),
        ('② 음양의 작동 원리 — "一陰一陽之謂道"',
         {'font_size': 14, 'space_before': 4}),
        ('③ 사물의 변화·생성 — "生生之謂易"',
         {'font_size': 14, 'space_before': 4}),
        ('④ 시(時)와 위(位)의 판단 — 군자의 처신',
         {'font_size': 14, 'space_before': 4}),
        ('⑤ 기미(幾微)와 점복의 의미 — "知幾其神乎"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"주역의 절반은 계사전에 있다"',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)


@S('Ⅲ. 십익')
def s_moonjeon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 십익', page, total)
    add_title(slide, '문언·설괘·서괘·잡괘전 — 나머지 네 편')
    items = [
        ('문언전(文言傳)',  '건·곤 두 괘의 집중 해설',
         '"여섯 효 모두에 사구(辭句)를 더함" — 덕성 수양의 본보기'),
        ('설괘전(說卦傳)',  '팔괘의 상징 체계',
         '☰ 천·☷ 지·☳ 뢰·☴ 풍 등 — 자연·동물·신체·방위 배당'),
        ('서괘전(序卦傳)',  '64괘의 배열 이유',
         '"왜 이 순서인가" — 괘의 연쇄 논리'),
        ('잡괘전(雜卦傳)',  '32쌍의 대비',
         '"건·곤" "태·비" "박·복" 짝지어 보는 관점 — 짧고 함축적'),
    ]
    top = 2.4
    for tag, role, char in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.95), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(3.0), Inches(0.5),
                    tag, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(9.1), Inches(0.95), PALE)
        add_textbox(slide, Inches(3.85), Inches(top + 0.1), Inches(8.9), Inches(0.45),
                    role, font_size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(3.85), Inches(top + 0.55), Inches(8.9), Inches(0.45),
                    char, font_size=12, color=INK)
        top += 1.07


# ============== Ⅳ. 64괘의 큰 흐름 ==============
@S('Ⅳ. 64괘 흐름')
def s_6mak(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 64괘 흐름', page, total)
    add_title(slide, '서괘전(序卦傳)의 6막 서사시 — "만물의 생성·발전·쇠퇴·재생"')
    items = [
        ('제1막', '1~10',   '천지의 개벽과 만물의 시작'),
        ('제2막', '11~30',  '흥성과 막힘의 순환 — 태↔비, 박↔복'),
        ('제3막', '31~40',  '인간관계의 전개 — 함·항·가인·규'),
        ('제4막', '41~50',  '성장과 변혁 — 손·익·혁·정'),
        ('제5막', '51~60',  '움직임과 멈춤의 교향 — 진↔간'),
        ('제6막', '61~64',  '완성과 미완성 — 기제·미제'),
    ]
    top = 2.4
    for tag, scope, content in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.15), Inches(2.0), Inches(0.4),
                    tag, font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.7), Inches(top), Inches(1.8), Inches(0.7), PALE)
        add_textbox(slide, Inches(2.7), Inches(top + 0.15), Inches(1.8), Inches(0.4),
                    scope, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.6), Inches(top), Inches(8.2), Inches(0.7),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.8), Inches(top + 0.15), Inches(7.9), Inches(0.4),
                    content, font_size=14, color=INK)
        top += 0.8
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '64괘는 일직선이 아니라 원(圓) — 미제(64) 뒤에 다시 건(1)이 시작',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 64괘 흐름')
def s_act_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 64괘 흐름', page, total)
    add_title(slide, '제1막 (1~10) — 천지의 개벽과 만물의 시작')
    items = [
        ('1·2', '건·곤(乾·坤)',       '하늘과 땅이 열린다 — 주역의 두 근원'),
        ('3',   '둔(屯)',              '천지가 처음 교류하니 만물이 비로소 생긴다'),
        ('4',   '몽(蒙)',              '생겨난 것은 어리다 — 가르침이 필요'),
        ('5',   '수(需)',              '어린 것은 기다려 자란다 — 먹이고 기름'),
        ('6',   '송(訟)',              '자라면 다툼(訟)이 생긴다'),
        ('7',   '사(師)',              '다툼이 커지면 군대(師)가 일어난다'),
        ('8',   '비(比)',              '전쟁 뒤에는 친목(比)이 필요'),
        ('9·10','소축·리(小畜·履)',   '작은 것을 쌓고 예로 밟아간다'),
    ]
    top = 2.3
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(1.0), Inches(0.4),
                    num, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(3.5), Inches(0.5), PALE)
        add_textbox(slide, Inches(1.8), Inches(top + 0.08), Inches(3.4), Inches(0.4),
                    name, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.3), Inches(top), Inches(7.5), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.45), Inches(top + 0.08), Inches(7.3), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.56
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"생겨남 → 가르침 → 먹임 → 다툼 → 질서"의 원형 — 세계의 발생학',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 64괘 흐름')
def s_act_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 64괘 흐름', page, total)
    add_title(slide, '제2막 (11~30) — 흥성과 막힘의 순환',
              '"태가 끝에 이르면 비가 온다(泰極否來)"')
    items = [
        ('11↔12',  '태(泰) ↔ 비(否)',    '태평과 막힘의 교체 — 주역 최대의 주제'),
        ('13~16',  '동인·대유·겸·예',     '풍요의 국면에서 지킬 덕목 (특히 겸은 6효 모두 길)'),
        ('17~20',  '수·고·임·관',         '따름·다스림의 단계'),
        ('21·22',  '서합·비(噬嗑·賁)',    '깨물어 합치고 꾸민다'),
        ('23↔24',  '박(剝) ↔ 복(復)',    '극도로 벗겨진 뒤 다시 돌아온다 — 동지(冬至)의 상'),
        ('25~28',  '무망·대축·이·대과',   '거짓 없음과 축적, 그리고 큰 지나침'),
        ('29·30',  '감·이(坎·離)',         '물과 불 — 상경의 끝, 천지의 두 작용'),
    ]
    top = 2.3
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.5), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(1.5), Inches(0.4),
                    num, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.2), Inches(top), Inches(3.5), Inches(0.6), PALE)
        add_textbox(slide, Inches(2.3), Inches(top + 0.13), Inches(3.4), Inches(0.4),
                    name, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.8), Inches(top), Inches(7.0), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.95), Inches(top + 0.13), Inches(6.8), Inches(0.4),
                    desc, font_size=12, color=INK)
        top += 0.66
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"만물은 극에 이르면 반전한다(物極必反)" — 그 반전을 미리 보는 눈이 지혜',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 64괘 흐름')
def s_act_3(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 64괘 흐름', page, total)
    add_title(slide, '제3막 (31~40) — 인간관계의 전개 (하경의 시작)',
              '"하경은 남녀의 감응으로 열린다"')
    items = [
        ('31', '함(咸)',           '감응 — 젊은 남녀의 사귐이 하경의 입구'),
        ('32', '항(恒)',           '감응이 지속되어 부부의 항구함이 됨'),
        ('33', '둔(遯)',           '물러날 때가 온다 — 은둔의 지혜'),
        ('34', '대장(大壯)',       '나아가는 큰 힘'),
        ('35', '진(晉)',           '나아가 밝아짐'),
        ('36', '명이(明夷)',       '밝음이 상함 — 어리석은 왕 아래 현자가 숨음'),
        ('37', '가인(家人)',       '집안의 도 — "여자 정이 안에, 남자 정이 밖에"'),
        ('38', '규(睽)',           '어긋남 — 가인이 흐트러지면 규가 된다'),
        ('39', '건(蹇)',           '절뚝거림 — 앞이 막혀 어려움'),
        ('40', '해(解)',           '어려움이 풀린다'),
    ]
    top = 2.0
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.8), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(0.8), Inches(0.4),
                    num, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.5), Inches(top), Inches(2.5), Inches(0.5), PALE)
        add_textbox(slide, Inches(1.6), Inches(top + 0.08), Inches(2.4), Inches(0.4),
                    name, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.1), Inches(top), Inches(8.7), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.25), Inches(top + 0.08), Inches(8.5), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.5
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '"인사(人事)의 세계" — 남녀·부부·군신·가족의 관계가 하경의 중심',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 64괘 흐름')
def s_act_4(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 64괘 흐름', page, total)
    add_title(slide, '제4막 (41~50) — 성장과 변혁',
              '"덜고 더하기, 결단과 조우, 그리고 혁명"')
    items = [
        ('41·42', '손·익(損·益)',  '덜어냄과 더함의 역설 — 손이 익이고 익이 손'),
        ('43',    '쾌(夬)',          '결단 — 5양효가 1음효를 몰아내는 형세'),
        ('44',    '구(姤)',          '조우 — 5양효 밑에 음효 하나가 싹튼다 (작은 악)'),
        ('45·46', '췌·승(萃·升)',  '모여들고 올라간다'),
        ('47',    '곤(困)',          '곤궁 — "곤이형(困而亨), 궁 중에서도 형통"'),
        ('48',    '정(井)',          '우물 — "바꿀 수 없는 것" · 변하지 않는 근원'),
        ('49',    '혁(革)',          '혁명 — "탕무혁명, 順乎天而應乎人"'),
        ('50',    '정(鼎)',          '혁명 후 새 질서의 건설 (鼎 = 세 발 솥)'),
    ]
    top = 2.3
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.3), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(1.3), Inches(0.4),
                    num, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.0), Inches(top), Inches(2.8), Inches(0.55), PALE)
        add_textbox(slide, Inches(2.1), Inches(top + 0.1), Inches(2.7), Inches(0.4),
                    name, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.9), Inches(top), Inches(7.9), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.05), Inches(top + 0.1), Inches(7.7), Inches(0.4),
                    desc, font_size=12, color=INK)
        top += 0.6


@S('Ⅳ. 64괘 흐름')
def s_act_5_6(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 64괘 흐름', page, total)
    add_title(slide, '제5·6막 (51~64) — 움직임과 멈춤 / 완성과 미완성',
              '"完成 뒤에 다시 未完이 온다"')
    items = [
        ('51·52', '진·간(震·艮)',     '움직임과 멈춤의 두 극단 — "艮其背"의 멈춤'),
        ('53·54', '점·귀매(漸·歸妹)', '점진적 결합 vs 성급한 결합'),
        ('55·56', '풍·여(豐·旅)',     '전성기(豐)와 나그네(旅) — "부귀해도 결국 떠난다"'),
        ('57·58', '손·태(巽·兌)',     '바람의 들어감과 못의 기쁨'),
        ('59·60', '환·절(渙·節)',     '흩어짐과 절제'),
        ('61·62', '중부·소과',         '속이 참됨과 작게 지나침'),
        ('63',    '기제(旣濟)',       '★ 이미 건넜다 — 6효 모두 정위 · 완성'),
        ('64',    '미제(未濟)',       '★ 아직 건너지 못함 — 6효 모두 부정위 · 미완'),
    ]
    top = 2.3
    for num, name, desc in items:
        is_end = (num == '64' or num == '63')
        c = ACCENT if is_end else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_end else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.3), Inches(0.55), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.1), Inches(1.3), Inches(0.4),
                    num, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.0), Inches(top), Inches(3.5), Inches(0.55), bg)
        add_textbox(slide, Inches(2.1), Inches(top + 0.1), Inches(3.4), Inches(0.4),
                    name, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.6), Inches(top), Inches(7.2), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.75), Inches(top + 0.1), Inches(7.0), Inches(0.4),
                    desc, font_size=12, color=INK, bold=is_end)
        top += 0.6
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '"왜 결말이 기제가 아니라 미제인가" — 역사는 멈추지 않는다 · 순환의 책',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅴ. 10대 핵심 사상 ==============
def make_thought_slide(num, title, original, principle, desc):
    @S('Ⅴ. 10대 사상')
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, f'Ⅴ. 10대 사상 ({num}/10)', page, total)
        add_title(slide, f'사상 {num} — {title}', original)
        add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9), ACCENT)
        add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
                    principle, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.6), PALE)
        add_textbox(slide, Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.4),
                    '풀이', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_paragraphs(slide, Inches(0.8), Inches(4.05), Inches(11.7), Inches(2.9),
                       desc, line_spacing=1.4, font_size=15)
    return renderer


make_thought_slide('1', '일음일양지도(一陰一陽之道)',
    '"한 번 음하고 한 번 양하는 것을 도(道)라 한다" — 계사상',
    '존재의 법칙 — 우주는 음양의 율동',
    [
        ('"낮과 밤 · 여름과 겨울 · 들숨과 날숨 · 성공과 실패 · 흥함과 망함"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('모두 음양의 한 리듬',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('고정된 상태는 없다 — 모든 것은 흐르고 교체된다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ])

make_thought_slide('2', '변역불변(變易不變)',
    '"變하는 것 속에 不變의 법칙이 있다"',
    '만물은 변하되, 변하는 방식 자체는 변하지 않는다',
    [
        ('계절이 매년 바뀌는 것 같지만',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('춘하추동의 순서는 불변',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('"변하는 패턴의 불변성"을 읽는 것이 지혜',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('변화를 두려워하지 않고 그 안의 질서를 본다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ])

make_thought_slide('3', '중(中)과 정(正) — 행위의 기준',
    '"中正" — 가운데이면서 바른 자리, 주역 최고의 상태',
    '치우치지 않고 제자리에 있는 덕',
    [
        ('中 — 치우치지 않음 (二爻 · 五爻)',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('正 — 제자리에 있음 (양효는 홀수, 음효는 짝수 자리)',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 6}),
        ('"中庸" 사상의 뿌리가 여기에',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('특히 九五(양효가 군주 자리에 중정) — 대부분 괘에서 가장 좋은 효',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ])

make_thought_slide('4', '시(時)와 위(位) — 때와 자리',
    '"같은 양효라도 어느 효냐, 어느 괘 속이냐에 따라 의미가 다르다"',
    '잘하기는 "절대 기준"이 아니라 "때와 자리에 맞느냐"의 문제',
    [
        ('한 효의 길흉은 두 가지에 달렸다',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('① 지금 어느 자리(位)에 있는가',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('② 지금이 어느 때(時)인가',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('"때와 자리"를 모르고 절대적 기준만 고집하면 — 모두 흉(凶)',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ])

make_thought_slide('5', '물극필반(物極必反)',
    '"窮則變 變則通 通則久" — 궁하면 변하고, 변하면 통하고, 통하면 오래간다',
    '극은 반드시 돌아선다 — "성공의 정점이 몰락의 시작"',
    [
        ('태 → 비 (태평이 끝나면 막힘)',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('박 → 복 (벗겨짐 뒤에 회복)',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('기제 → 미제 (완성 뒤에 다시 미완)',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('"극점을 경계하라"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('가득 차면 기울고, 기울면 다시 찬다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ])

make_thought_slide('6', '기미(幾微) — 조짐을 본다',
    '"幾者 動之微 吉之先見者也" — 기는 움직임의 미묘함, 길함이 먼저 나타나는 것',
    '일이 벌어지기 전 희미한 조짐을 읽는 능력',
    [
        ('"君子見幾而作 不俟終日"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('군자는 기미를 보고 움직이되 하루를 기다리지 않는다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('현대의 약신호 감지(weak signal) · 트렌드 예측의 원리',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"이미 벌어진 뒤"가 아니라 "벌어지기 전"의 사고법',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ])

make_thought_slide('7', '길흉회린 — 운명이 아니라 행위에 달림',
    '"吉凶悔吝 生乎動者也" — 길흉회린은 움직임에서 생긴다',
    '주역은 운명론이 아니라 실천철학',
    [
        ('같은 상황에서도 —',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 흉 → 회(반성) → 길로 간다',
         {'font_size': 15, 'space_before': 6}),
        ('• 길 → 린(교만) → 흉으로 간다',
         {'font_size': 15, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('"정해진 운명은 없다" — 행위가 길흉을 결정한다',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ])

make_thought_slide('8', '상(象)과 수(數) — 이미지와 숫자',
    '"聖人立象以盡意" — 성인이 상(象)을 세워 뜻을 다했다',
    '언어 이전의 직관적 통찰 — 이미지로 세계를 파악',
    [
        ('상(象) — 이미지(☰☷☳☴)로 세계의 원형을 파악',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('수(數) — 1~9의 조합으로 패턴의 수학적 구조',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('라이프니츠 — 이진법 발견 후 주역 괘를 보고 감탄',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('"동양에 이미 있었구나"',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ])

make_thought_slide('9', '천인합일(天人合一) — 자연과 인간의 상응',
    '"자연의 리듬이 곧 인간 삶의 리듬"',
    '자연을 관찰하는 것이 곧 자기 수양',
    [
        ('"天行健 君子以自强不息" (건 상전)',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('하늘의 운행이 강건하니 군자는 스스로 쉬지 않는다',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('"地勢坤 君子以厚德載物" (곤 상전)',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('땅의 형세가 곤이니 군자는 두터운 덕으로 만물을 싣는다',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ])

make_thought_slide('10', '적선지가(積善之家) — 윤리적 인과',
    '"積善之家 必有餘慶 積不善之家 必有餘殃" — 문언전',
    '냉정한 법칙 안의 윤리적 인과',
    [
        ('"선을 쌓은 집에는 반드시 넘치는 경사가 있고',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 불선을 쌓은 집에는 반드시 넘치는 재앙이 있다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('주역의 우주론은 냉정한 법칙이지만',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('그 법칙 안에 윤리적 인과가 함께 작동한다',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ])


# ============== Ⅵ. 핵심 괘 9선 ==============
def make_gwae_slide(num, han, eum, symbol, desc, principle, lesson):
    @S(f'Ⅵ. 핵심 괘 ({num}/9)')
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, f'Ⅵ. 핵심 괘 ({num}/9)', page, total)
        add_title(slide, f'{han}({eum}) — {desc}', principle)
        add_textbox(slide, Inches(0.6), Inches(2.3), Inches(3.2), Inches(3.0),
                    symbol, font_size=72, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(4.0), Inches(2.3), Inches(8.8), Inches(3.0), PALE)
        add_textbox(slide, Inches(4.0), Inches(2.45), Inches(8.8), Inches(0.4),
                    '핵심 가르침', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(4.2), Inches(2.95), Inches(8.4), Inches(2.2),
                    lesson, font_size=15, color=INK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    return renderer


make_gwae_slide('1', '건(乾)', '乾為天 · 1번', '☰\n☰',
    '하늘이 둘 — 모든 것의 근원, 강건한 양의 정수',
    '"天行健 君子以自强不息"',
    '하늘의 운행이 강건하니 군자는 스스로 강해지기를 쉬지 않는다\n\n'
    '용의 6단계 — 잠룡(初九) → 견룡(二) → 종일건건(三) → 혹약재연(四) → 비룡재천(五) → 항룡유회(上九)\n\n'
    '"자기 강화의 끝없는 노력" — 주역의 첫 가르침')

make_gwae_slide('2', '곤(坤)', '坤為地 · 2번', '☷\n☷',
    '땅이 둘 — 만물을 싣는 두터운 덕, 순응하는 음의 정수',
    '"地勢坤 君子以厚德載物"',
    '땅의 형세가 곤이니 군자는 두터운 덕으로 만물을 싣는다\n\n'
    '"承乘" — 받들고 따르는 덕 · "厚載" — 두텁게 싣는 덕\n\n'
    '"순응하는 힘이 가장 큰 힘이다" — 건괘와 짝을 이루는 주역의 두 기둥')

make_gwae_slide('3', '태↔비(泰↔否)', '11·12번 — 태평과 막힘', '☷\n☰  ↔  ☰\n☷',
    '땅이 위·하늘이 아래(태) ↔ 하늘이 위·땅이 아래(비)',
    '"泰極否來" — 태평이 끝에 이르면 막힘이 온다',
    '태(泰) — 두 기운이 가운데서 교류 → 태평·소통\n\n'
    '비(否) — 각자 제자리에 있어 보이지만 기운이 교류 안 함 → 막힘\n\n'
    '"제자리"가 반드시 좋은 것이 아니다 — 교류·역전이 통(通)을 낳는다\n주역 최대의 주제 — 흥성과 막힘의 영원한 교체')

make_gwae_slide('4', '겸(謙)', '15번 · 地山謙', '☷\n☶',
    '땅속에 산이 있다 — 큰 것을 안에 감추고 낮은 자리에',
    '★ 주역 64괘 중 유일하게 6효 모두 길(吉)',
    '"謙謙君子 卑以自牧也"\n겸손하고 겸손한 군자는 낮춤으로 자기를 친다\n\n'
    '"勞謙君子 萬民服也" — 수고하면서도 겸손하면 만민이 복종한다\n\n'
    '"풍요할수록 겸손하라" — 주역이 제시하는 최고의 덕\n— 산은 본래 높지만 땅 속에 감추니 더욱 두터워진다')

make_gwae_slide('5', '박↔복(剝↔復)', '23·24번 — 벗겨짐과 회복', '☶\n☷  ↔  ☷\n☳',
    '산이 땅 위에 거의 무너짐(박) ↔ 어둠 끝에 한 양효가 돌아옴(복)',
    '"七日來復" — 7일이면 돌아온다 · 동지(冬至)의 상',
    '박(剝) — 5음효 위에 1양효만 남은 위태로움\n복(復) — 5음효 아래에 1양효가 돌아옴 (동지의 첫 양)\n\n'
    '"不遠復" — 멀지 않은 곳에서 돌아온다 — 안회(顔回)의 덕\n\n'
    '"가장 어두운 순간에 빛이 시작된다"\n작은 잘못을 즉시 바로잡는 자기 회복력')

make_gwae_slide('6', '손↔익(損↔益)', '41·42번 — 덜어냄과 더함의 역설', '☶\n☱  ↔  ☴\n☳',
    '덜어내야 더해지고, 위에서 덜어 아래에 더해야 진짜 이익',
    '"損上益下 民說無疆" — 위를 덜어 아래에 더하면 백성이 끝없이 기뻐한다',
    '손(損) — 손하익상(損下益上, 아래를 덜어 위에 더함) → 백성의 고통\n익(益) — 손상익하(損上益下, 위를 덜어 아래에 더함) → 백성의 이익\n\n'
    '"손해가 반드시 손해가 아니고, 이익이 반드시 이익이 아니다"\n— 진짜 이익은 베푸는 데서, 진짜 손해는 쌓는 데서')

make_gwae_slide('7', '혁↔정(革↔鼎)', '49·50번 — 혁명과 새 질서', '☱\n☲  ↔  ☲\n☴',
    '못 가운데 불(혁) ↔ 나무 위에 불(정 = 세 발 솥)',
    '"湯武革命 順乎天而應乎人" — 천명에 순응하고 민심에 응하는 혁명',
    '혁(革) — 옛 것을 바꾼다 · "天地革而四時成"\n정(鼎) — 새 것을 세운다 · 세 발 솥에 새 음식\n\n'
    '"破舊立新" — 깨고 세우는 한 쌍\n\n'
    '혁명은 천명과 민심의 일치 — 함부로 일으키는 것이 아니다\n그러나 때가 되면 두려워 말고 단행하라')

make_gwae_slide('8', '기제↔미제(旣濟↔未濟)', '63·64번 — 완성과 미완성', '☵\n☲  ↔  ☲\n☵',
    '이미 건넜다(기제) — 6효 모두 정위 ↔ 아직 건너지 못함(미제) — 6효 모두 부정위',
    '★ 64괘의 결말이 "기제"가 아니라 "미제"인 까닭',
    '기제 — 모든 효가 제자리에 있는 완성 상태\n그러나 "완성 다음에 다시 미완이 온다"\n\n'
    '미제 — 아직 끝나지 않은 새로운 시작\n\n'
    '"완성이 끝이 아니라 새로운 출발"\n— 주역은 순환의 책 · 미제 다음에 다시 건괘가 시작')

make_gwae_slide('9', '핵심 9괘 통합', '주역의 골격', '☰☷\n泰否謙\n剝復損益\n革鼎旣未',
    '핵심 9괘로 보는 주역의 본체',
    '"음양 → 변역 → 중정 → 시위 → 순환"',
    '건·곤 — 우주의 두 근원 (양과 음)\n태·비 — 흥성과 막힘의 교체\n겸 — 풍요할수록 겸손\n박·복 — 극에 이르러 돌아옴\n손·익 — 덜어냄이 더함\n혁·정 — 깨고 세움\n기제·미제 — 완성 뒤의 새로운 시작')


# ============== Ⅶ. 주제별 클러스터 ==============
@S('Ⅶ. 주제별')
def s_cluster_start(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 주제별', page, total)
    add_title(slide, '클러스터 A·B — 시작과 끝 / 어려운 시기')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                'A. 시작과 끝의 괘', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('1 건(乾) · 2 곤(坤)', {'font_size': 15, 'bold': True}),
        ('  모든 것의 근원', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('63 기제(旣濟) · 64 미제(未濟)',
         {'font_size': 15, 'bold': True, 'space_before': 6}),
        ('  완성과 미완의 영원한 순환', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('"끝은 새로운 시작"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                'B. 어려운 시기에 읽는 괘', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('3 둔(屯) — 시작의 진통', {'font_size': 14, 'bold': True}),
        ('4 몽(蒙) — 무지의 극복', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('29 감(坎) — 반복되는 위험', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('39 건(蹇) — 절뚝거림', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('47 곤(困) — 곤궁', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('"困而亨" — 궁 중에서도 형통',
         {'font_size': 13, 'color': ACCENT, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)


@S('Ⅶ. 주제별')
def s_cluster_leader(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 주제별', page, total)
    add_title(slide, '클러스터 C — 리더십·군주의 도')
    items = [
        ('5',  '수(需)',   '기다림의 전략'),
        ('7',  '사(師)',   '"사정(師貞), 군대는 바름으로 움직인다"'),
        ('19', '임(臨)',   '다가가 이끌음'),
        ('20', '관(觀)',   '"관이불언 — 관찰하되 말하지 않는다"'),
        ('46', '승(升)',   '아래에서 위로 올라감'),
        ('49', '혁(革)',   '"湯武革命, 順乎天而應乎人" — 천명·민심의 일치'),
        ('50', '정(鼎)',   '혁명 후 새 질서의 건설'),
    ]
    top = 2.4
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(1.0), Inches(0.4),
                    num, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(2.8), Inches(0.5), PALE)
        add_textbox(slide, Inches(1.8), Inches(top + 0.08), Inches(2.7), Inches(0.4),
                    name, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.6), Inches(top), Inches(8.2), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.75), Inches(top + 0.08), Inches(8.0), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.58


@S('Ⅶ. 주제별')
def s_cluster_human(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 주제별', page, total)
    add_title(slide, '클러스터 D — 인간관계·가족의 도')
    items = [
        ('8',  '비(比)',     '친합 — 친구·동료의 사귐'),
        ('13', '동인(同人)', '타인과 함께함'),
        ('31', '함(咸)',     '남녀의 감응 — 하경의 첫 괘'),
        ('32', '항(恒)',     '부부의 항구함'),
        ('37', '가인(家人)', '가정의 도 — "여자 정이 안에, 남자 정이 밖에"'),
        ('38', '규(睽)',     '어긋남 — 가인이 흐트러지면 규'),
        ('54', '귀매(歸妹)', '혼사의 주의'),
    ]
    top = 2.4
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(1.0), Inches(0.4),
                    num, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(2.8), Inches(0.5), PALE)
        add_textbox(slide, Inches(1.8), Inches(top + 0.08), Inches(2.7), Inches(0.4),
                    name, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.6), Inches(top), Inches(8.2), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.75), Inches(top + 0.08), Inches(8.0), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.58


@S('Ⅶ. 주제별')
def s_cluster_self(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 주제별', page, total)
    add_title(slide, '클러스터 E — 자기 수양')
    items = [
        ('15', '겸(謙)',     '★ 주역 최고의 덕 · 6효 모두 길'),
        ('24', '복(復)',     '"不遠復" — 멀지 않은 곳에서 돌아온다 · 안회의 덕'),
        ('25', '무망(无妄)', '거짓 없음'),
        ('26', '대축(大畜)', '크게 쌓음 — 덕의 축적'),
        ('52', '간(艮)',     '그쳐야 할 곳에 그침 — "艮其背 不獲其身"'),
        ('61', '중부(中孚)', '속이 참됨 — "豚魚之孚" 신뢰가 돼지·물고기까지'),
        ('41', '손(損)',     '덜어냄 — 자기를 비우는 덕'),
    ]
    top = 2.4
    for num, name, desc in items:
        is_special = (num == '15')
        c = ACCENT if is_special else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_special else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(0.5), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(1.0), Inches(0.4),
                    num, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(2.8), Inches(0.5), bg)
        add_textbox(slide, Inches(1.8), Inches(top + 0.08), Inches(2.7), Inches(0.4),
                    name, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.6), Inches(top), Inches(8.2), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.75), Inches(top + 0.08), Inches(8.0), Inches(0.4),
                    desc, font_size=13, color=INK)
        top += 0.58


@S('Ⅶ. 주제별')
def s_cluster_warning(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 주제별', page, total)
    add_title(slide, '클러스터 F·G — 변화·전환 / 경계해야 할 괘')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                'F. 변화·전환의 괘', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('11 태 ↔ 12 비', {'font_size': 14, 'bold': True}),
        ('  태평과 막힘의 교체', {'font_size': 12, 'color': SUB}),
        ('23 박 → 24 복', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  벗겨짐 뒤의 회복', {'font_size': 12, 'color': SUB}),
        ('41 손 → 42 익', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  덜어냄과 더함의 역설', {'font_size': 12, 'color': SUB}),
        ('43 쾌 / 44 구', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  결단과 조우', {'font_size': 12, 'color': SUB}),
        ('49 혁 / 50 정', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  혁명과 건설', {'font_size': 12, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                'G. 경계해야 할 괘', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('6 송(訟)', {'font_size': 14, 'bold': True}),
        ('  다툼의 위험', {'font_size': 12, 'color': SUB}),
        ('28 대과(大過)', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  크게 지나침', {'font_size': 12, 'color': SUB}),
        ('38 규(睽)', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  어긋남', {'font_size': 12, 'color': SUB}),
        ('56 여(旅)', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  나그네 처지 — 객지에서의 조심', {'font_size': 12, 'color': SUB}),
        ('62 소과(小過)', {'font_size': 14, 'bold': True, 'space_before': 4}),
        ('  작게 지나침', {'font_size': 12, 'color': SUB}),
    ], line_spacing=1.3)


# ============== Ⅷ. 계사전 12명구 ==============
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=42):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=20, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('Ⅷ. 명구 (1/12)',
    '一 陰 一 陽 之 謂 道',
    '일음일양지위도',
    '한 번 음하고 한 번 양하는 것을 도(道)라 한다 — 주역 사상의 정수',
    '계사상', hanmun_size=58), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (2/12)',
    '生 生 之 謂 易',
    '생생지위역',
    '끊임없이 낳고 낳는 것을 역(易)이라 한다',
    '계사상', hanmun_size=80), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (3/12)',
    '形 而 上 者 謂 之 道\n形 而 下 者 謂 之 器',
    '형이상자위지도 · 형이하자위지기',
    '형상 위의 것을 도(道), 아래의 것을 기(器)라 한다 — 동양 형이상학의 출발점',
    '계사상', hanmun_size=26), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (4/12)',
    '窮 則 變   變 則 通   通 則 久',
    '궁즉변 · 변즉통 · 통즉구',
    '궁하면 변하고, 변하면 통하고, 통하면 오래간다 — 변화의 법칙',
    '계사하', hanmun_size=36), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (5/12)',
    '樂 天 知 命   故 不 憂',
    '낙천지명 · 고불우',
    '천명을 즐겨 알기에 근심하지 않는다',
    '계사상', hanmun_size=48), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (6/12)',
    '仁 者 見 之 謂 之 仁\n知 者 見 之 謂 之 知',
    '인자견지위지인 · 지자견지위지지',
    '인자(仁者)는 이를 보고 인이라 하고, 지자(知者)는 보고 지(知)라 한다',
    '계사상 — 시각의 다양성', hanmun_size=22), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (7/12)',
    '君 子 藏 器 於 身   待 時 而 動',
    '군자장기어신 · 대시이동',
    '군자는 그릇(才)을 몸에 감추고 때를 기다려 움직인다',
    '계사하', hanmun_size=32), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (8/12)',
    '見 幾 而 作   不 俟 終 日',
    '견기이작 · 불사종일',
    '기미를 보고 일어서되 하루를 기다리지 않는다 — 약신호 감지의 원리',
    '계사하', hanmun_size=38), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (9/12)',
    '天 行 健   君 子 以 自 强 不 息',
    '천행건 · 군자이자강불식',
    '하늘의 운행이 강건하니 군자는 스스로 강해지기를 쉬지 않는다',
    '건 상전(象傳)', hanmun_size=32), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (10/12)',
    '地 勢 坤   君 子 以 厚 德 載 物',
    '지세곤 · 군자이후덕재물',
    '땅의 형세가 곤이니 군자는 두터운 덕으로 만물을 싣는다',
    '곤 상전(象傳)', hanmun_size=32), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (11/12)',
    '積 善 之 家   必 有 餘 慶',
    '적선지가 · 필유여경',
    '선을 쌓은 집에는 반드시 넘치는 경사가 있다',
    '곤괘 문언전(文言傳)', hanmun_size=40), 'Ⅷ. 명구'))

SLIDES.append((make_quote_slide('Ⅷ. 명구 (12/12)',
    '聖 人 立 象 以 盡 意',
    '성인입상이진의',
    '성인은 상(象, 이미지)을 세워 뜻을 다하였다 — 언어 너머의 직관',
    '계사상', hanmun_size=46), 'Ⅷ. 명구'))


# ============== Ⅸ. 해석 학파 ==============
@S('Ⅸ. 해석 학파')
def s_sangsu(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 해석 학파', page, total)
    add_title(slide, '상수학파(象數學派) — "기호와 숫자의 학"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '관심 — 괘상·호괘·종괘·착괘·납갑·오행 배당 등 형식 구조의 수리적 법칙',
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.5), [
        ('대표 인물',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('한대 — 맹희·경방 · 송대 — 소옹(邵雍, 소강절)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"주역의 수학적·구조적 아름다움을 드러낸다"',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.6), [
        ('장점과 위험',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('• 주역의 수학적·구조적 아름다움을 드러냄',
         {'font_size': 14, 'space_before': 4}),
        ('• 위험: 지나친 도식화로 실제 인간 문제에서 멀어질 수 있음',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅸ. 해석 학파')
def s_uirei(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 해석 학파', page, total)
    add_title(slide, '의리학파(義理學派) — "뜻과 이치의 학"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '관심 — 괘사·효사에 담긴 도덕적·철학적 의미 · 인간 삶의 상황으로 읽기',
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.5), [
        ('대표 인물',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('위(魏)의 왕필(王弼) · 송의 정이(程頤)·주희(朱熹)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"괘를 인간 삶의 상황으로 읽기"',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.6), [
        ('표준 텍스트',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('• 정이 『이천역전(伊川易傳)』 · 주희 『주역본의(周易本義)』',
         {'font_size': 14, 'space_before': 4}),
        ('• 조선의 퇴계·율곡·다산 — 모두 의리학의 관점에서 주역을 읽음',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅸ. 해석 학파')
def s_two_eyes(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 해석 학파', page, total)
    add_title(slide, '두 학파의 통합 — "두 눈으로 보아야 입체가 된다"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"어느 한쪽을 배척할 필요는 없다"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('상수 — 구조의 눈 · 의리 — 뜻의 눈',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('두 눈으로 함께 보아야 주역이 입체가 된다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('한국 — 정약용 『주역사전(周易四箋)』',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('의리(義理) + 상수(象數) 종합 — 양 학파의 장점을 흡수',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('조선 후기 실학의 정점',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"주역은 만능 열쇠가 아니라 구체적 상황의 실천서"',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


# ============== Ⅹ. 후대 영향 ==============
@S('Ⅹ. 후대 영향')
def s_song_master(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '송대 — 정이·주희의 의리학 정립')
    items = [
        ('정이(程頤)',  '1033~1107',  '『이천역전(伊川易傳)』',
         '의리학의 표준 — 조선 유학의 기본 교재'),
        ('주희(朱熹)',  '1130~1200',  '『주역본의(周易本義)』',
         '"본의(本義)" — 의리와 점복의 통합'),
        ('소옹(邵雍)',  '1011~1077',  '『황극경세서(皇極經世書)』',
         '상수학의 정점 — "선천도(先天圖)" 발명'),
        ('주돈이(周敦頤)','1017~1073', '『태극도설(太極圖說)』',
         '주역에서 성리학으로 — "無極而太極"'),
    ]
    top = 2.4
    for name, era, work, char in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(1.0), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    name, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.65), Inches(2.5), Inches(0.3),
                    era, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(3.5), Inches(1.0), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.3), Inches(3.5), Inches(0.5),
                    work, font_size=13, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.8), Inches(top), Inches(6.0), Inches(1.0),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.0), Inches(top + 0.3), Inches(5.7), Inches(0.5),
                    char, font_size=12, color=INK)
        top += 1.13


@S('Ⅹ. 후대 영향')
def s_joseon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '조선 — 퇴계·율곡·다산의 주역 학습')
    items = [
        ('퇴계 이황',   '1501~1570',  '『계몽전의(啓蒙傳疑)』',
         '주희 『역학계몽』 주석 — 의리학의 깊이'),
        ('율곡 이이',   '1536~1584',  '『역수책(易數策)』',
         '학문 시험 답안 — 상수와 의리의 균형'),
        ('성호 이익',   '1681~1763',  '『역경질서(易經疾書)』',
         '근기학파의 주역 — 실용적 접근'),
        ('다산 정약용', '1762~1836',  '『주역사전(周易四箋)』 · 『역학서언』',
         '실학적 종합 — 의리+상수 통합의 정점'),
        ('해석',       '조선 5백 년',  '경연 + 사대부 필독서',
         '"수신·치국의 거울" — 정치·인생 결정의 기준'),
    ]
    top = 2.3
    for name, era, work, char in items:
        is_dasan = '다산' in name
        c = ACCENT if is_dasan else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_dasan else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.12), Inches(2.5), Inches(0.4),
                    name, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(2.5), Inches(0.3),
                    era, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(3.5), Inches(0.85), bg)
        add_textbox(slide, Inches(3.2), Inches(top + 0.25), Inches(3.5), Inches(0.4),
                    work, font_size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.8), Inches(top), Inches(6.0), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.0), Inches(top + 0.25), Inches(5.7), Inches(0.4),
                    char, font_size=12, color=INK)
        top += 0.95


@S('Ⅹ. 후대 영향')
def s_leibniz(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '라이프니츠와 이진법(1703) — 서양의 주역 발견')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('1703년 — 라이프니츠(G. W. Leibniz)가 이진법 발견',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('예수회 선교사 부베(Bouvet)가 라이프니츠에게',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('소옹의 선천도(先天圖) — 64괘 배열을 보냄',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"이미 동양에 있었구나!" — 라이프니츠의 감탄',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('양효(⚊) = 1, 음효(⚋) = 0 · 64괘 = 6비트 이진수',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"동양의 4,000년 된 이진 체계를 서양이 17세기에 발견"',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅹ. 후대 영향')
def s_modern_west(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '20세기 서양 — 융·DNA·시스템 이론과의 만남')
    items = [
        ('칼 융(Carl Jung)',     '1875~1961',
         '"동시성(Synchronicity)" 원리 — 주역 점복의 심리학적 해석'),
        ('빌헬름 역(Wilhelm I Ching)', '1923',
         '독일어 번역 + 융의 서문 — 서양 주역 수용의 결정적 텍스트'),
        ('DNA 64 코돈',           '1968',
         '단백질 합성 코돈 64개 = 64괘 — 우연한 일치인가?'),
        ('시스템 이론',           '20세기 후반',
         '피드백 순환·복잡계 — 주역의 변화 사상과 통함'),
        ('VUCA 시대',             '21세기',
         '불확실성을 다루는 책으로 재조명'),
    ]
    top = 2.4
    for name, era, char in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.15), Inches(3.0), Inches(0.4),
                    name, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(2.0), Inches(0.7), PALE)
        add_textbox(slide, Inches(3.7), Inches(top + 0.15), Inches(2.0), Inches(0.4),
                    era, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.8), Inches(top), Inches(7.0), Inches(0.7),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.95), Inches(top + 0.15), Inches(6.8), Inches(0.4),
                    char, font_size=12, color=INK)
        top += 0.8


@S('Ⅹ. 후대 영향')
def s_eastasia(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 후대 영향', page, total)
    add_title(slide, '동아시아 확산 — 문명의 공통 문법')
    items = [
        ('중국',     '주역은 "群經之首" — 모든 학문의 기초',
         '한대~청대 학자가 한 번은 주역에 손을 댐'),
        ('한국',     '신라 백제부터 사대부 필독서',
         '세종이 집현전에 주역 강독 명령 · 다산 『주역사전』'),
        ('일본',     '헤이안 시대부터 귀족 한문 교양',
         '에도 시대 주자학의 핵심 텍스트'),
        ('베트남',   '한자 문화권의 공통 경전',
         '리·쩐 왕조의 사대부 교양'),
        ('전 세계',  '20세기 이후 영어 번역 다수',
         'Wilhelm-Baynes 역이 서양 표준'),
    ]
    top = 2.4
    for region, char, where in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.0), Inches(0.5),
                    region, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.7), Inches(top), Inches(5.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(2.85), Inches(top + 0.22), Inches(4.8), Inches(0.5),
                    char, font_size=12, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.8), Inches(top), Inches(5.0), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.95), Inches(top + 0.22), Inches(4.8), Inches(0.5),
                    where, font_size=12, color=INK, align=PP_ALIGN.CENTER)
        top += 0.97


# ============== Ⅺ. 현대 의의 ==============
@S('Ⅺ. 현대 의의')
def s_vuca(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ① — 불확실성(VUCA)의 시대를 위한 책')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                'VUCA — Volatility · Uncertainty · Complexity · Ambiguity',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.5), [
        ('"변동성·불확실성·복잡성·모호성의 시대"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('현대 경영·정치·기술이 직면한 핵심 도전',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.6), [
        ('"주역은 3,000년 전부터 변화는 상수(常數)다 라는 전제 위에"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('변화를 두려워하지 않고 읽어내는 기술이 주역의 본령',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"變則通 通則久" — 변하면 통하고, 통하면 오래간다',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_timing(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ② — 타이밍(時) 감각의 사전')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"인생에서 가장 어려운 것은 \'무엇을\'이 아니라 \'언제\'"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('나아갈 때와 물러날 때',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('말할 때와 침묵할 때',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('주역 64괘는',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"인생의 64가지 전형적 상황을 미리 시뮬레이션한 때의 사전(時之典)"',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('• 시작 → 둔(3) / 준비기 → 수(5)',
         {'font_size': 13, 'space_before': 6}),
        ('• 관계 → 동인(13)·규(38) / 성공 뒤 → 겸(15)·박(23)',
         {'font_size': 13, 'space_before': 4}),
        ('• 변혁 → 혁(49) / 마무리 → 기제·미제',
         {'font_size': 13, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_leadership(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ③ — 조직 리더를 위한 7개 괘')
    items = [
        ('1 건',    '리더의 자강불식(自强不息)'),
        ('2 곤',    '포용과 실행 — 후덕재물'),
        ('7 사',    '군대의 운용 (조직 통솔)'),
        ('11·12 태비', '시대의 호흡 읽기'),
        ('15 겸',   '성공 뒤의 겸손'),
        ('49·50 혁정','변혁과 안정의 균형'),
        ('63·64 기제미제', '성과에 안주하지 않는 자세'),
    ]
    top = 2.3
    for tag, role in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(3.5), Inches(0.4),
                    tag, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.2), Inches(top), Inches(8.6), Inches(0.6), PALE)
        add_textbox(slide, Inches(4.4), Inches(top + 0.13), Inches(8.3), Inches(0.4),
                    role, font_size=14, color=INK)
        top += 0.65
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"리더의 7가지 시기 — 한 사람의 일생이 곧 7괘의 순환"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅺ. 현대 의의')
def s_mirror(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ④ — 점서를 넘어 "거울"로')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('공자의 선언 — "점치지 않고도 읽히는 책(不占而已)"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('한 괘를 읽는 일은',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('자신의 현재를 그 괘의 구조에 비추어 보는 일',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"점괘가 아니라 거울"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('미래를 예측하는 도구가 아니라',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('자기 위치를 객관화하는 거울',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"내가 지금 어디에 있고, 무엇을 향해 어떻게 움직이고 있는가"',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅺ. 현대 의의')
def s_system(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ⑤ — 시스템 사고·피드백 순환')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('20세기 시스템 사고와의 만남',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('피터 셍게(Peter Senge) 『The Fifth Discipline』 —',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"피드백 순환 (Feedback Loop) = 음양의 작동"',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('주역의 시스템적 통찰',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 강화 루프 — 양효 누적이 흥성을 만들지만 결국 反이 온다',
         {'font_size': 14, 'space_before': 6}),
        ('• 균형 루프 — 음·양 교체가 순환을 만든다 (태↔비)',
         {'font_size': 14, 'space_before': 4}),
        ('• 지연 효과 — 박(剝)에서 복(復)까지 7효(七日來復)',
         {'font_size': 14, 'space_before': 4}),
        ('• "복잡계의 작동을 64가지로 단순화한 모델"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅺ. 현대 의의')
def s_which_gwae(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 현대 의의', page, total)
    add_title(slide, '현대 ⑥ — "내가 지금 어떤 괘에 있는가"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('주역이 살아나는 순간 —',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('"내가 지금 어떤 괘에 있는가?"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('를 자기에게 물을 때',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('자기 진단 질문',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 지금 시(時)는? — 시작·발전·전성·쇠퇴·전환·재생 중 어디?',
         {'font_size': 13, 'space_before': 6}),
        ('• 지금 위(位)는? — 초·이·삼·사·오·상효 어느 자리?',
         {'font_size': 13, 'space_before': 4}),
        ('• 응(應)·비(比)는 어떤가? — 누구와 호응하고 이웃하는가?',
         {'font_size': 13, 'space_before': 4}),
        ('• 어떻게 움직여야 길(吉)인가? — 같은 자리도 행위로 달라진다',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅻ. 마무리 ==============
@S('Ⅻ. 마무리')
def s_one_page(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 마무리', page, total)
    add_title(slide, '한 장으로 보는 주역')
    add_filled_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.1), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.8), [
        ('太 極 (태극)',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('  ↓',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('陰 陽 (음양)',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('  ↓',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('四 象 (사상)',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('  ↓',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('八 卦 (팔괘) — ☰☱☲☳☴☵☶☷',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('  ↓',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('六 十 四 卦 (64괘) — 8 × 8 = 64',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('  ↓',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('十 翼 (십익)',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('— "一 陰 一 陽 之 謂 道" —',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.15)


@S('Ⅻ. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 주역')
    add_filled_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.7), [
        ('복희가 8괘를 긋고,',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('문왕·주공이 64괘 384효에 사(辭)를 매단 점서를',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('공자의 십익이 철학서로 승격시킨',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('"변화의 책(Book of Changes)" —',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('', {'font_size': 6}),
        ('"변하지 않는 것은 단 하나 — \'만물이 변한다\'는 사실"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('그 변화의 리듬을 먼저 읽는 자가 스스로를 구한다',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 6}),
        ('— 3,000년 동안 동아시아 모든 학문의 머리로 자리한 군경지수(群經之首) —',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
    ], line_spacing=1.2)


@S('Ⅻ. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
                '一 陰 一 陽 之 謂 道',
                font_size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '일 음 일 양 지 위 도', font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.1), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.7),
                '"한 번 음하고 한 번 양하는 것을 도(道)라 한다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
                '— 계사상 · 주역 철학의 정수, 동양 우주관의 한 문장 —',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '"내가 지금 어떤 괘에 있는가" — 이 질문이 주역을 살린다',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '감사합니다', font_size=24, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\주역_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')