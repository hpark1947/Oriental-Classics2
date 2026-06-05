# -*- coding: utf-8 -*-
"""
서경(書經) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '오경(五經)의 으뜸 · 동양 정치사상의 원천',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.6),
                '書 經',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '서 경  ·  상서(尙書)',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
                '天視自我民視 天聽自我民聽',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '우서·하서·상서·주서 — 요순부터 동주 초까지의 정치 문서',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 서경이란 무엇인가'),
        ('Ⅱ.', '전승 — 금문상서와 고문상서'),
        ('Ⅲ.', '육체(六體) — 여섯 가지 문체'),
        ('Ⅳ.', '우서(虞書) — 요순의 이상 정치'),
        ('Ⅴ.', '하서(夏書) — 우왕의 치수'),
        ('Ⅵ.', '상서(商書) — 은의 흥망'),
        ('Ⅶ.', '주서(周書) — 주의 건국과 제도'),
    ]
    items_right = [
        ('Ⅷ.', '명편 깊이 읽기'),
        ('Ⅸ.', '핵심 사상 — 천명·덕치·민본'),
        ('Ⅹ.', '명구 모음'),
        ('Ⅺ.', '후대 영향'),
        ('Ⅻ.', '현대적 의의'),
        ('XIII.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(0.9), y, Inches(0.8), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.7), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(7.2), y, Inches(0.9), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(8.1), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '서경(書經) — 중국 최초의 정치 산문집',
              '요순부터 동주 초까지, 「옛 문서」의 모음')
    add_paragraphs(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.5), [
        ('· 중국에서 가장 오래된 정치 문서·역사 기록 모음집',
         {'font_size': 18, 'space_before': 6}),
        ('· 「상서(尙書)」 — 「상(尙)」은 「위(上)」 곧 「옛것」 → 「옛 문서」',
         {'font_size': 18, 'space_before': 6}),
        ('· 한대(漢代)에 오경의 하나로 격상되어 「서경(書經)」으로 불림',
         {'font_size': 18, 'space_before': 6}),
        ('· 요(堯)·순(舜) 시대부터 하·상·주 삼대에 이르는 제왕의 말씀과 정치 문서',
         {'font_size': 18, 'space_before': 6}),
        ('· 훈시(誥)·서약(誓)·명령(命)·모의(謨)·법전(典) 등 정치 행위 그 자체의 기록',
         {'font_size': 18, 'space_before': 6}),
        ('· 동양 정치사상 — 천명·덕치·민본의 원천이 되는 책',
         {'font_size': 18, 'space_before': 6}),
    ])


@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '명칭의 변천 — 書에서 經으로')
    rows = [
        ('書 서',       '원래의 이름 — 「적은 것」, 곧 정치 문서·기록의 뜻'),
        ('尙書 상서',   '「옛 문서」라는 뜻으로 한대 이전까지의 표준 명칭'),
        ('書經 서경',   '한대 오경(五經)으로 격상되며 「경(經)」의 지위를 얻음'),
        ('六經 / 五經', '시·서·역·예·악·춘추 → 악이 빠진 오경 — 서는 그 한 축'),
    ]
    for i, (name, desc) in enumerate(rows):
        y = Inches(2.3 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.75), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.75),
                    name, font_size=20, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y + Inches(0.05), Inches(9.4), Inches(0.7),
                    desc, font_size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_era(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '시대와 규모')
    rows = [
        ('시대 폭',  '요순(전설) → 하 → 상(은) → 주 (BC 2300년경 ~ BC 8세기)'),
        ('수록 인물', '요·순·우·탕·이윤·반경·부열·문왕·무왕·주공·강왕 등'),
        ('현존 편수', '58편 — 금문 33편 + 고문(僞古文) 25편'),
        ('네 부분',   '우서 5 · 하서 4 · 상서 17 · 주서 32'),
        ('체재',      '왕의 훈시·맹세·명령·모의 — 1인칭 발화체 다수'),
        ('성격',      '단순한 역사서가 아닌 「정치 행위 그 자체」의 기록'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.55),
                    k, font_size=16, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.55),
                    v, font_size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_thought(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '핵심 사상 — 한 폭으로 보기')
    boxes = [
        ('天 命', '천명', '하늘의 뜻 — 덕 있는 자에게 천하를 맡긴다'),
        ('德 治', '덕치', '덕으로 다스리기 — 형벌보다 교화·위력보다 덕행'),
        ('民 本', '민본', '백성이 나라의 근본 — 民惟邦本 本固邦寧'),
        ('敬天勤民', '경천근민', '하늘을 공경하고 백성을 위해 부지런히 힘쓴다'),
    ]
    for i, (han, kor, desc) in enumerate(boxes):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(6.0), Inches(2.0), PALE)
        add_textbox(slide, x, y + Inches(0.25), Inches(6.0), Inches(0.7),
                    han, font_size=30, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(1.0), Inches(6.0), Inches(0.4),
                    kor, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.4), Inches(6.0), Inches(0.5),
                    desc, font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 전승 ==============
SEC2 = 'Ⅱ. 전승 — 금문·고문 상서'

@S(SEC2)
def ii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '서경의 험난한 전승사',
              '진시황의 분서 → 한대의 복원 → 청대의 위작 고증')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● BC 213 — 진시황의 분서갱유로 서경의 원본이 거의 소실',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 4}),
        ('● 한 초(漢初) — 노학자 복생(伏生, 90세)이 구술로 28편을 전함 → 금문상서',
         {'font_size': 18, 'space_before': 10}),
        ('● 한 무제 — 공자 옛집 벽에서 고문(古文)으로 쓰인 「공벽본」 발견 (16편 추가)',
         {'font_size': 18, 'space_before': 10}),
        ('● 동진(東晉) — 매색(梅賾)이 고문상서 25편을 추가로 헌상',
         {'font_size': 18, 'space_before': 10}),
        ('● 청대 염약거(閻若璩) — 매색본 25편이 위작임을 고증 → 「위고문(僞古文)」 확정',
         {'font_size': 18, 'space_before': 10}),
        ('● 그러나 사상적 가치는 「위(僞)」라도 부정되지 않음 — 「인심도심」 등은 그 안에 있음',
         {'font_size': 17, 'color': SUB, 'space_before': 10}),
    ])


@S(SEC2)
def ii_compare(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '금문상서 vs 고문상서 — 한눈에 비교')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.55), INK)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.5), Inches(0.55),
                '구분', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(4.2), Inches(2.3), Inches(4.2), Inches(0.55),
                '금문상서(今文)', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(8.5), Inches(2.3), Inches(4.2), Inches(0.55),
                '고문상서(古文)', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('전승 시기',  '한초 복생이 구전',          '공자 구택 발견·동진 매색 헌상'),
        ('문자',       '한대 예서(隸書)',           '선진 시대 고문(古文)'),
        ('편수',       '28편 (또는 29편)',          '추가 25편'),
        ('진위',       '대체로 진본으로 인정',      '청 염약거 「위고문」 고증'),
        ('현존',       '현행 58편 중 33편',          '현행 58편 중 25편'),
        ('대표 편',     '요전·우공·반경·홍범 등',     '대우모(「인심도심」)·태서·열명 등'),
    ]
    for i, (k, j, g) in enumerate(rows):
        y = Inches(2.85 + i * 0.6)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.6), bg)
        add_textbox(slide, Inches(0.7), y, Inches(3.5), Inches(0.6),
                    k, font_size=15, color=SUB, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.2), y, Inches(4.2), Inches(0.6),
                    j, font_size=14, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(8.5), y, Inches(4.2), Inches(0.6),
                    g, font_size=14, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅲ. 육체 ==============
SEC3 = 'Ⅲ. 육체(六體) — 여섯 가지 문체'

@S(SEC3)
def iii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '서경의 글은 여섯 갈래로 분류된다',
              '典 · 謨 · 訓 · 誥 · 誓 · 命')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('· 서경은 「문서」이지만 그 형식이 한 가지가 아니다',
         {'font_size': 18, 'space_before': 6}),
        ('· 「전(典)」은 법도, 「모(謨)」는 의논, 「훈(訓)」은 가르침의 기록',
         {'font_size': 18, 'space_before': 8}),
        ('· 「고(誥)」는 위에서 아래로 내리는 훈시, 「서(誓)」는 출정의 맹세, 「명(命)」은 임명·책봉의 명령',
         {'font_size': 18, 'space_before': 8}),
        ('· 같은 정치 문서라도 「누가 누구에게 무엇을 말하는가」에 따라 문체가 갈라진다',
         {'font_size': 18, 'space_before': 8}),
        ('· 후대 한문 산문 — 특히 詔·勅·誥·命 등 황실 문서 양식의 원형',
         {'font_size': 18, 'space_before': 8}),
    ])


@S(SEC3)
def iii_table(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '육체(六體) — 누가 누구에게')
    rows = [
        ('典 전', '제왕의 행적과 법도 기록',     '요전(堯典), 순전(舜典)'),
        ('謨 모', '군신 간의 모략·의논',         '대우모(大禹謨), 고요모(皐陶謨)'),
        ('訓 훈', '신하가 군주에게 올리는 가르침', '이훈(伊訓)'),
        ('誥 고', '군주가 신하·백성에게 내리는 훈시', '탕고(湯誥), 강고(康誥), 주고(酒誥), 대고(大誥)'),
        ('誓 서', '출전·출정 시의 서약',          '감서(甘誓), 탕서(湯誓), 목서(牧誓), 진서(秦誓)'),
        ('命 명', '군주가 신하에게 내리는 명령',  '열명(說命), 문후지명(文侯之命), 고명(顧命)'),
    ]
    for i, (name, role, exam) in enumerate(rows):
        y = Inches(2.3 + i * 0.78)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.5), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.65),
                    name, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, Inches(2.4), y, Inches(4.8), Inches(0.65), PALE)
        add_textbox(slide, Inches(2.5), y, Inches(4.7), Inches(0.65),
                    role, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.4), y, Inches(5.5), Inches(0.65),
                    exam, font_size=14, color=SUB, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅳ. 우서 ==============
SEC4 = 'Ⅳ. 우서(虞書)'

@S(SEC4)
def iv_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '우서(虞書) — 요순의 이상 정치 5편',
              '선양(禪讓)의 이상과 덕치의 원형')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「우(虞)」 — 순임금의 나라 — 요로부터 천하를 받음', {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● 요·순의 「선양(禪讓)」 — 천하를 사사로이 하지 않고 덕 있는 자에게 물려줌',
         {'font_size': 17, 'space_before': 10}),
        ('● 5편', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('     1. 요전(堯典) — 요임금의 성덕·치세, 순에게 선양', {'font_size': 16, 'space_before': 4}),
        ('     2. 순전(舜典) — 순임금의 등용·관직 설치, 「詩言志」의 출전', {'font_size': 16, 'space_before': 4}),
        ('     3. 대우모(大禹謨) — 우(禹)에 대한 순의 훈시, 「人心惟危 道心惟微」', {'font_size': 16, 'space_before': 4}),
        ('     4. 고요모(皐陶謨) — 고요의 정치론, 구덕(九德)', {'font_size': 16, 'space_before': 4}),
        ('     5. 익직(益稷) — 익과 후직의 보좌', {'font_size': 16, 'space_before': 4}),
    ])


@S(SEC4)
def iv_yaoshun(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '요(堯)·순(舜) — 동양 정치의 영원한 표상')
    cols = [
        ('堯 요', '요전(堯典)',
         '「克明俊德, 以親九族」\n큰 덕을 밝혀 구족을 친애.\n\n능력과 덕이 있는 자에게\n천하를 물려준 첫 임금.\n\n순을 시험하여 등용·선양.'),
        ('舜 순', '순전(舜典)',
         '「愼徽五典, 五典克從」\n오륜을 신중히 펴니 백성이 따름.\n\n역산에서 농사·하빈에서 도자기.\n효성으로 부모를 감동시킴.\n\n구관(九官)을 설치, 천하를 다스림.'),
    ]
    for i, (han, src, body) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), INK)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=28, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.85), Inches(5.9), Inches(0.4),
                    src, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.3), Inches(3.55), Inches(5.5), Inches(3.5),
                       [(body, {'font_size': 15, 'color': INK})], line_spacing=1.45)


# ============== Ⅴ. 하서 ==============
SEC5 = 'Ⅴ. 하서(夏書)'

@S(SEC5)
def v_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '하서(夏書) — 4편',
              '우왕(禹王)의 치수와 국토 정리의 기록')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1. 우공(禹貢) — 구주(九州)의 지리·공물·토질 — 중국 최초의 지리서',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 4}),
        ('     기·연·청·서·양·형·예·양·옹 — 아홉 주(州)로 천하를 나눠 다스림', {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 2. 감서(甘誓) — 유호씨 정벌을 앞두고 한 우왕의 군중 맹세',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     서경의 첫 「서(誓)」 — 후대 모든 군중 맹세의 원형', {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 3. 오자지가(五子之歌) — 태강(太康)의 실정을 한탄하는 다섯 형제의 노래',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     「民惟邦本 本固邦寧」 — 「백성은 나라의 근본」 민본주의의 절대 명구', {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 4. 윤정(胤征) — 천문관 의·화의 직무 태만에 대한 정벌',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 14}),
    ])


@S(SEC5)
def v_yugong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '우공(禹貢)과 구주(九州) — 중국 최초의 지리서',
              '치수로 국토를 정리한 우왕의 행정 청사진')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 우왕은 13년간 「세 번 자기 집 앞을 지나면서도 들어가지 않고」 치수에 전념',
         {'font_size': 18, 'space_before': 4}),
        ('● 황하의 범람을 다스리며 구주(九州)를 획정 — 중국 최초의 지리·행정 구획',
         {'font_size': 18, 'space_before': 10}),
        ('● 구주', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('     冀(기) · 兗(연) · 靑(청) · 徐(서) · 揚(양) · 荊(형) · 豫(예) · 梁(양) · 雍(옹)',
         {'font_size': 17, 'color': INK, 'space_before': 4, 'font_name': 'Batang'}),
        ('● 각 주의 토질·산천·공물(貢物)·교통로를 상세히 기록 — 행정과 지리학의 동시 효시',
         {'font_size': 17, 'space_before': 12}),
        ('● 「구주(九州)」는 이후 「중국 전체」를 가리키는 대명사로 자리잡음',
         {'font_size': 17, 'color': SUB, 'space_before': 8}),
    ])


# ============== Ⅵ. 상서 ==============
SEC6 = 'Ⅵ. 상서(商書) — 은의 흥망'

@S(SEC6)
def vi_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '상서(商書) — 17편',
              '탕왕의 건국 → 이윤·반경·부열 → 미자의 절망')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1. 탕서(湯誓) — 탕왕이 하의 걸왕을 정벌하기 전의 맹세 — 「予畏上帝, 不敢不正」',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 2. 이훈(伊訓) — 명재상 이윤이 어린 태갑 왕에게 올린 훈계',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 3. 태갑(太甲) 상·중·하 — 태갑의 유배와 회개·복위 — 「習與性成」',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 4. 함유일덕(咸有一德) — 군신이 모두 「한결같은 덕」을 가져야 함',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 5. 반경(盤庚) 상·중·하 — 도읍을 은(殷)으로 옮긴 반경의 천도 설득문',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 6. 열명(說命) 상·중·하 — 고종 무정이 부열(傅說)을 등용하고 정치를 논함',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 7. 미자(微子) — 은 말기 미자의 절망과 결단 — 망국 앞의 충신',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC6)
def vi_pangeng(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '반경(盤庚)의 천도 — 설득의 정치학',
              '은(殷)으로의 천도, 6번째 수도를 정하다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 잦은 홍수와 귀족의 사치로 흔들리던 상나라 — 반경은 결단을 내림',
         {'font_size': 17, 'space_before': 4}),
        ('● 백성과 귀족 모두가 천도를 꺼리자, 반경은 강압이 아닌 「설득」을 택함',
         {'font_size': 17, 'space_before': 10}),
        ('● 명구 — 「若網在綱, 有條而不紊」', {'font_size': 19, 'bold': True, 'color': ACCENT, 'space_before': 12, 'font_name': 'Batang'}),
        ('     그물이 벼리에 있어야 가지런하고 어지럽지 않은 것처럼, 정치도 강기(綱紀)가 서야 한다',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 「天命不易」 — 천명은 쉽지 않으니 백성을 두려워하라',
         {'font_size': 17, 'space_before': 12}),
        ('● 결국 은으로 천도 — 이후 「은(殷)」이 상나라의 별칭이 됨',
         {'font_size': 17, 'space_before': 8, 'color': SUB}),
    ])


# ============== Ⅶ. 주서 ==============
SEC7 = 'Ⅶ. 주서(周書)'

@S(SEC7)
def vii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '주서(周書) — 32편',
              '서경에서 가장 방대하고 가장 핵심적인 부분')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 건국 — 태서(상·중·하)·목서·무성', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 4}),
        ('     무왕의 은 정벌과 정복 — 「天視自我民視」 출전', {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 법전·홍범 — 홍범(洪範)', {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('     기자(箕子)가 무왕에게 전한 구주홍범(九疇) — 정치의 대법', {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 주공 섭정 — 금등·대고·강고·주고·재재·소고·낙고·다사·무일·군석',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('     성왕을 보필한 주공의 정치 문서 모음 — 주서의 정수', {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 성·강 시대 — 다방·입정·주관·고명·강왕지고',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('     주의 황금기 — 제도와 의례의 정비', {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 후기 — 여형·문후지명·진서',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'space_before': 10}),
        ('     형법(呂刑)·서주의 쇠퇴와 동주의 도래', {'font_size': 14, 'color': SUB, 'space_before': 4}),
    ])


@S(SEC7)
def vii_zhougong(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '주공(周公) — 주서의 주인공',
              '섭정 7년 — 어린 조카를 위해 천하를 다스리다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 무왕의 동생이자 성왕의 숙부 — 무왕 사후 어린 조카를 대신해 섭정',
         {'font_size': 17, 'space_before': 4}),
        ('● 「一沐三捉髮, 一飯三吐哺」 — 머리 감다 세 번, 밥 먹다 세 번 손님을 맞이함',
         {'font_size': 17, 'space_before': 10}),
        ('● 관숙·채숙의 반란을 진압(삼감의 난)하고 동방을 평정 — 「대고(大誥)」',
         {'font_size': 17, 'space_before': 10}),
        ('● 낙읍(雒邑) 건설 — 동도(東都) 설립으로 천하 통치의 중심 마련',
         {'font_size': 17, 'space_before': 10}),
        ('● 봉건제·종법제·예악을 정비 — 「제례작악(制禮作樂)」의 주인공',
         {'font_size': 17, 'space_before': 10}),
        ('● 7년 후 성왕에게 권력을 돌려주고 신하의 자리로 — 충성과 도덕의 표상',
         {'font_size': 17, 'space_before': 10}),
        ('● 공자가 가장 흠모한 사람 — 「久矣吾不復夢見周公」 (오랫동안 꿈에 주공을 못 뵈었구나)',
         {'font_size': 16, 'color': SUB, 'space_before': 12, 'font_name': 'Batang'}),
    ])


@S(SEC7)
def vii_hongfan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '홍범(洪範) — 구주(九疇), 정치의 큰 법',
              '기자(箕子)가 무왕에게 전한 통치의 9가지 강령')
    add_paragraphs(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(5.0), [
        ('● 1. 오행(五行) — 수·화·목·금·토 — 자연의 기본 원리', {'font_size': 15, 'space_before': 4}),
        ('● 2. 오사(五事) — 모(貌)·언(言)·시(視)·청(聽)·사(思) — 몸가짐의 다섯', {'font_size': 15, 'space_before': 6}),
        ('● 3. 팔정(八政) — 식·화·사·사공·사도·사구·빈·사 — 여덟 가지 정사', {'font_size': 15, 'space_before': 6}),
        ('● 4. 오기(五紀) — 세·월·일·성신·역수 — 시간과 천문', {'font_size': 15, 'space_before': 6}),
        ('● 5. 황극(皇極) — 군주가 세우는 큰 중도 — 구주의 중심', {'font_size': 15, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 6. 삼덕(三德) — 정직·강극·유극 — 통치의 세 결', {'font_size': 15, 'space_before': 6}),
        ('● 7. 계의(稽疑) — 거북점·시초로 의심을 풀다', {'font_size': 15, 'space_before': 6}),
        ('● 8. 서징(庶徵) — 비·바람 같은 징조로 정치를 점검', {'font_size': 15, 'space_before': 6}),
        ('● 9. 오복·육극(五福·六極) — 다섯 가지 복과 여섯 가지 화', {'font_size': 15, 'space_before': 6}),
        ('「홍범」은 동아시아 통치 철학의 최고(最古) 강령서 — 한대 이후 정치·천문·역의 종합 매뉴얼이 됨',
         {'font_size': 14, 'color': SUB, 'space_before': 12}),
    ])


# ============== Ⅷ. 명편 깊이 읽기 ==============
SEC8 = 'Ⅷ. 명편 깊이 읽기'

def make_chapter_slide(num, total, name_han, name_kor, source, original, modern, theme, point):
    @S(SEC8)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC8} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=28, bold=True, color=INK)
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    f'출전 — {source}', font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.4), [
            (original, {'font_size': 18, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 핵심 주제', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(theme, {'font_size': 16, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 음미할 지점', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 16, 'color': INK})], line_spacing=1.35)


CHAPTERS = [
    ('요전(堯典)', '요전', '우서 — 서경의 첫 편',
     '克明俊德 以親九族  百姓昭明 協和萬邦',
     '큰 덕을 밝혀 구족을 친애하니, 백성이 환히 밝아지고 만방이 화합하다',
     '한 사람의 덕이 가족 → 백성 → 천하로 동심원처럼 퍼지는 정치의 출발점.',
     '『대학』 「수신제가치국평천하」 — 동양 도덕정치의 모형은 모두 이 구절에서 비롯한다.'),
    ('순전(舜典)', '순전', '우서',
     '詩言志 歌永言  聲依永 律和聲',
     '시는 뜻을 말한 것, 노래는 말을 길게 펼친 것 / 소리는 길이에 의지하고, 가락은 소리에 화합한다',
     '동양 시론(詩論)의 원전 — 「시언지(詩言志)」가 처음 등장하는 구절.',
     '시경의 「대서」는 이 구절을 받아 동아시아 시 이론의 토대를 마련했다.'),
    ('대우모(大禹謨)', '대우모', '우서 — 16자심전(十六字心傳)',
     '人心惟危 道心惟微  惟精惟一 允執厥中',
     '인심은 위태롭고 도심은 은미하니 / 정밀하고 한결같이 그 중도를 잡으라',
     '동양 도통(道統)의 핵심 — 「16자심전」으로 불리는 유학의 최고 명제.',
     '주희가 『중용장구서』에서 이 구절을 도통 전수의 핵심으로 지목 — 성리학의 출발점.'),
    ('고요모(皐陶謨)', '고요모', '우서 — 구덕(九德)',
     '寬而栗 柔而立 愿而恭 亂而敬  擾而毅 直而溫 簡而廉 剛而塞 彊而義',
     '너그러우면서 위엄있고, 부드러우면서 굳세고, 성실하면서 공손하고...',
     '인재 등용의 기준 — 한 사람의 덕을 아홉 가지 모순적 균형으로 본다.',
     '편향이 아닌 「양극의 통합」이 덕 — 동양 인재론의 깊이를 보여주는 명제.'),
    ('우공(禹貢)', '우공', '하서 — 중국 최초의 지리서',
     '禹敷土 隨山刊木  奠高山大川',
     '우가 토지를 정돈하니 산을 따라 나무를 베고 / 큰 산과 큰 강을 정했다',
     '치수와 함께 진행된 국토 정리 — 구주(九州)의 획정.',
     '단순한 지리가 아닌 「정치 행정의 청사진」 — 한대 이후 모든 행정 구획의 모범.'),
    ('오자지가(五子之歌)', '오자지가', '하서',
     '民惟邦本 本固邦寧',
     '백성이 곧 나라의 근본이니, 근본이 굳어야 나라가 평안하다',
     '동양 민본주의의 절대 명제 — 민(民)을 근본으로 삼는 정치의 출발점.',
     '맹자의 「민귀군경(民貴君輕)」, 정약용의 민본 사상이 모두 이 구절을 잇는다.'),
    ('탕서(湯誓)', '탕서', '상서 — 탕왕의 정벌 맹세',
     '予畏上帝 不敢不正  夏氏有罪 予畏上帝',
     '나는 상제(上帝)를 두려워하여 바로잡지 않을 수 없으니 / 하씨가 죄가 있어 상제를 두려워한다',
     '폭군 정벌의 정당성을 천명(天命)으로 논증 — 후대 「방벌론」의 원형.',
     '나의 의지가 아닌 「하늘의 명」으로 정벌한다는 논리 — 맹자 역성혁명론의 출전.'),
    ('이훈(伊訓)', '이훈', '상서 — 이윤의 훈계',
     '與人不求備 檢身若不及',
     '남에게는 완전함을 구하지 않되, 자신을 살피기는 미치지 못한 듯이 하라',
     '재상 이윤이 어린 태갑 왕에게 올린 「세 가지 풍(三風)·열 가지 죄(十愆)」의 가르침.',
     '타인에겐 관대, 자신에겐 엄격 — 동양 자기수양의 핵심 원칙.'),
    ('태갑(太甲)', '태갑', '상서 — 「習與性成」의 출전',
     '習與性成  茲乃不義 習與性成',
     '습관이 성품을 이룬다 / 이런 의롭지 못함이 습관이 되어 성품이 된다',
     '태갑이 어려서 무도하다 이윤에 의해 동궁(桐宮)에 유배 → 3년 만에 회개·복위.',
     '「습관이 곧 운명」 — 자기 변화·교육의 가능성을 보여주는 위대한 회개의 서사.'),
    ('반경(盤庚)', '반경', '상서 — 천도(遷都)의 정치학',
     '若網在綱 有條而不紊',
     '그물이 벼리에 있어야 가닥가닥 어지럽지 않다',
     '천도를 반대하는 백성·귀족을 설득하는 반경의 정치적 수사.',
     '「강기(綱紀)」 — 정치는 큰 줄기가 잡혀야 비로소 가지가 흐트러지지 않는다.'),
    ('열명(說命)', '열명', '상서 — 고종이 부열을 등용함',
     '非知之艱 行之惟艱',
     '아는 것이 어려운 것이 아니라, 행하는 것이 어렵다',
     '재상 부열(傅說)이 고종 무정에게 올린 정치 격언.',
     '「지행(知行)」의 동양적 명제 — 왕양명의 「지행합일」, 정약용의 「행기지(行其知)」의 원천.'),
    ('태서(泰誓)', '태서', '주서 — 무왕의 은 정벌 맹세',
     '天視自我民視  天聽自我民聽',
     '하늘이 보는 것은 우리 백성이 보는 것에서, 하늘이 듣는 것은 우리 백성이 듣는 것에서',
     '「하늘의 뜻은 곧 백성의 뜻」 — 천명을 민의(民意)와 동일시한 동양 정치의 정점.',
     '맹자·동중서·왕양명까지 — 모든 민본·여론 사상의 절대 출전이 된 한 구절.'),
    ('홍범(洪範)', '홍범', '주서 — 기자가 무왕에게 전한 구주',
     '皇建其有極  斂時五福 用敷錫厥庶民',
     '임금이 그 큰 중도(皇極)를 세워, 오복을 모아 백성에게 베푼다',
     '「황극(皇極)」 — 9주(疇)의 중심에 놓인 「임금의 큰 중도」.',
     '동양 통치 철학의 최고 강령 — 정치·천문·윤리를 하나의 체계로 묶은 최고(最古)의 청사진.'),
    ('무일(無逸)', '무일', '주서 — 주공이 성왕에게',
     '先知稼穡之艱難 乃逸 則知小人之依',
     '먼저 농사의 어려움을 알고 나서 편안해야, 백성의 의지하는 바를 안다',
     '편안함(逸)에 빠지지 말라는 주공의 절절한 훈계 — 「농사의 고통을 알라」.',
     '치자(治者)는 피치자(被治者)의 삶을 직접 알아야 함 — 「민이 의지하는 바」를 아는 지도자의 조건.'),
    ('여형(呂刑)', '여형', '주서 — 형법의 신중함',
     '哀矜折獄  輕重諸罰有權',
     '슬퍼하고 가엾이 여기며 옥사를 결단하고 / 가볍고 무거운 형벌엔 헤아림이 있어야 한다',
     '형벌을 가볍게 하기보다 「신중하게」 적용 — 동양 형법 사상의 원전.',
     '「의심나면 가볍게」 「죄인이라도 가엾이」 — 인본주의 사법 정신의 출발점.'),
    ('진서(秦誓)', '진서', '주서 — 서경의 마지막 편',
     '人之有技 若己有之  人之彦聖 其心好之',
     '남에게 기예가 있으면 마치 자기가 가진 듯이 / 남이 아름답고 거룩하면 마음으로 좋아한다',
     '진 목공이 패전 후 자기를 반성하며 한 맹세 — 「겸손과 인재 포용」의 표본.',
     '『대학』이 통째로 인용 — 「사람을 알아보고 시기 없이 포용함」이 치자의 큰 그릇.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅸ. 핵심 사상 ==============
SEC9 = 'Ⅸ. 핵심 사상'

@S(SEC9)
def ix_tianming(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '천명(天命) — 하늘의 명',
              '왕조 교체의 정당성과 통치의 근거')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85),
                '皇 天 無 親, 惟 德 是 輔',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(2.0), Inches(3.2), Inches(9.3), Inches(0.4),
                '하늘은 친소(親疏)가 없으니, 오직 덕을 가진 자를 돕는다 — 채중지명(蔡仲之命)',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('· 하늘은 특정 왕조·혈통에 매이지 않는다 — 「天命靡常」',
         {'font_size': 18, 'space_before': 8}),
        ('· 덕이 있으면 천명을 받고, 덕을 잃으면 천명이 거두어간다',
         {'font_size': 18, 'space_before': 8}),
        ('· 「폭군 방벌」의 논리적 토대 — 하 걸·은 주의 천명 박탈',
         {'font_size': 18, 'space_before': 8}),
        ('· 후대 동양 정치사 — 모든 왕조 교체의 정당화 언어가 되다',
         {'font_size': 18, 'space_before': 8}),
    ])


@S(SEC9)
def ix_dezhi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '덕치(德治) — 덕으로 다스리기',
              '형벌이 아닌 교화, 위력이 아닌 덕행')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「克明德 慎罰」 — 덕을 밝히고 형벌을 신중히 한다 (강고)',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 형벌은 최후 수단 — 그 적용도 「슬퍼하고 가엾이 여기며(哀矜)」',
         {'font_size': 18, 'space_before': 12}),
        ('● 「만물에는 도(道)가 있다」 — 군주는 사물의 결을 따라 다스릴 뿐',
         {'font_size': 18, 'space_before': 10}),
        ('● 「正德 利用 厚生」 — 덕을 바르게, 쓰임을 이롭게, 삶을 두텁게 (대우모)',
         {'font_size': 18, 'space_before': 10, 'font_name': 'Batang'}),
        ('● 후대 — 공자의 「爲政以德」, 맹자의 「王道」가 모두 이 자리에서 자라남',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC9)
def ix_minben(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '민본(民本) — 백성이 나라의 근본',
              '서경이 동양 민본주의의 절대 원천이 된 까닭')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85),
                '民 惟 邦 本, 本 固 邦 寧',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(2.0), Inches(3.2), Inches(9.3), Inches(0.4),
                '백성이 곧 나라의 근본이니, 근본이 굳어야 나라가 평안하다 — 오자지가(五子之歌)',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('· 백성이 군주의 도구가 아니라, 군주가 백성을 위한 직책',
         {'font_size': 18, 'space_before': 8}),
        ('· 「天視自我民視 天聽自我民聽」 — 하늘의 뜻은 곧 백성의 뜻',
         {'font_size': 18, 'space_before': 8, 'font_name': 'Batang'}),
        ('· 맹자 「民貴君輕」, 황종희 『명이대방록』, 정약용 『목민심서』의 원천',
         {'font_size': 18, 'space_before': 8}),
        ('· 동양 정치사상 2500년의 가장 깊은 뿌리',
         {'font_size': 18, 'space_before': 8, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC9)
def ix_jingtian(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '경천근민(敬天勤民) — 군주의 자세',
              '하늘을 공경하고 백성을 위해 부지런히')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 경천(敬天) — 천명을 두려워하는 자세, 「天命不易」', {'font_size': 19, 'bold': True, 'color': ACCENT}),
        ('     군주의 권한이 자기 것이 아니라 「하늘로부터 빌린 것」이라는 자각',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 근민(勤民) — 백성을 위해 부지런함', {'font_size': 19, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     「先知稼穡之艱難」 — 농사의 어려움부터 알아야 한다 (무일)',
         {'font_size': 15, 'color': SUB, 'space_before': 4, 'font_name': 'Batang'}),
        ('● 안일(逸)의 경계 — 무일(無逸) 한 편이 모두 「안일하지 말라」',
         {'font_size': 18, 'space_before': 14}),
        ('● 「臨深履薄」 — 깊은 못 앞에, 살얼음 위에 선 듯한 신중함',
         {'font_size': 18, 'space_before': 10, 'font_name': 'Batang'}),
    ])


@S(SEC9)
def ix_renxin(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '인심·도심 — 16자심전(十六字心傳)',
              '대우모 — 동양 도통(道統)의 핵심')
    add_filled_rect(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.4),
                '人 心 惟 危  道 心 惟 微\n惟 精 惟 一  允 執 厥 中',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 「人心」 — 욕망에 흔들리는 일상의 마음 — 위태롭다(危)',
         {'font_size': 18, 'space_before': 8}),
        ('● 「道心」 — 도리에 따르는 마음 — 은미하여 잘 보이지 않는다(微)',
         {'font_size': 18, 'space_before': 8}),
        ('● 「惟精惟一」 — 정밀히 분별하고 한결같이 지킴',
         {'font_size': 18, 'space_before': 8}),
        ('● 「允執厥中」 — 진실로 그 중도(中)를 잡음',
         {'font_size': 18, 'space_before': 8}),
        ('● 주희가 『중용장구서』 머리에서 「道統의 비밀」로 풀이 — 성리학의 핵심 명제',
         {'font_size': 16, 'color': SUB, 'space_before': 12}),
    ])


# ============== Ⅹ. 명구 모음 ==============
SEC10 = 'Ⅹ. 명구 모음'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC10)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC10} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=32, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=17, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 17, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('오자지가(五子之歌)',  '民 惟 邦 本, 本 固 邦 寧',
     '백성이 나라의 근본이니, 근본이 굳어야 나라가 평안하다',
     '동양 민본주의의 절대 명제. 정약용 『목민심서』의 머리글, 한국 헌법 정신의 깊은 원천.'),
    ('태서 중(泰誓 中)', '天 視 自 我 民 視, 天 聽 自 我 民 聽',
     '하늘이 보는 것은 우리 백성이 보는 것에서, 하늘이 듣는 것은 우리 백성이 듣는 것에서',
     '천명을 민의(民意)와 동일시한 동양 정치사상 최고의 명제. 맹자·왕양명까지 끊임없이 인용.'),
    ('대우모(大禹謨)',  '人 心 惟 危, 道 心 惟 微  惟 精 惟 一, 允 執 厥 中',
     '인심은 위태롭고 도심은 은미하니, 정밀하고 한결같이 그 중도를 잡으라',
     '「16자심전」 — 주희가 『중용장구서』에서 도통 전수의 핵심으로 지목. 성리학의 토대.'),
    ('대우모(大禹謨)', '正 德 利 用 厚 生',
     '덕을 바르게 하고, 쓰임을 이롭게 하며, 삶을 두텁게 한다',
     '동양 정치의 세 큰 과제. 한국 「홍익인간」의 정신과도 통하는 통치의 세 축.'),
    ('이훈(伊訓)',  '與 人 不 求 備, 檢 身 若 不 及',
     '남에게는 완전함을 구하지 않되, 자신을 살피기는 미치지 못한 듯이 하라',
     '타인에겐 관대, 자신에겐 엄격 — 동양 자기수양의 결정적 한 줄.'),
    ('태갑(太甲)', '習 與 性 成',
     '습관이 성품을 이룬다',
     '습관 = 성품 = 운명. 어려서 잘못 든 태갑이 회개·복위했듯, 사람은 끝까지 변할 수 있다.'),
    ('열명(說命)', '非 知 之 艱, 行 之 惟 艱',
     '아는 것이 어려운 것이 아니라, 행하는 것이 어렵다',
     '「지행(知行)」의 동양적 명제. 왕양명의 「지행합일」, 정약용의 「행기지」 모두 여기서 출발.'),
    ('소민(小民) 인용·다방', '與 治 同 道 罔 不 興  與 亂 同 事 罔 不 亡',
     '다스림과 도(道)를 함께하면 흥하지 않을 수 없고, 어지러움과 일을 함께하면 망하지 않을 수 없다',
     '정치의 도리는 단순하다 — 어떤 결을 따르느냐가 흥망을 가른다.'),
    ('홍범(洪範)', '無 偏 無 黨, 王 道 蕩 蕩',
     '치우치지 않고 편당하지 않으니, 왕도가 평탄하구나',
     '동양 정치의 「공평무사」 이상. 「왕도」라는 말의 가장 오래된 출전 중 하나.'),
    ('무일(無逸)', '先 知 稼 穡 之 艱 難',
     '먼저 농사의 어려움을 알라',
     '주공이 어린 성왕에게 — 치자가 피치자의 삶을 모르면 「민이 의지하는 바」를 알 수 없다.'),
    ('진서(秦誓)', '人 之 有 技, 若 己 有 之',
     '남이 가진 재주를, 마치 자기 것인 듯이 (기뻐한다)',
     '인재 포용의 정수. 『대학』이 통째로 인용 — 사람을 알아보고 시기하지 않는 그릇이 큰 정치를 만든다.'),
    ('주관(周官)', '功 崇 惟 志, 業 廣 惟 勤',
     '공이 높은 것은 뜻이 있어서이고, 업이 넓은 것은 부지런해서이다',
     '큰 공과 큰 업은 결국 「뜻」과 「부지런함」 — 동양 일의 원리의 한 줄 격언.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅺ. 후대 영향 ==============
SEC11 = 'Ⅺ. 후대 영향'

@S(SEC11)
def xi_confucian(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '유가 경전의 핵심으로',
              '논어 · 맹자 · 대학 · 중용이 끊임없이 서경을 인용')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 『논어』 — 「子曰: 書云, 孝乎惟孝, 友于兄弟」 (효와 우애의 출전을 서경에 둠)',
         {'font_size': 17, 'space_before': 4}),
        ('● 『맹자』 — 「天視自我民視」를 인용해 역성혁명론 전개',
         {'font_size': 17, 'space_before': 8}),
        ('● 『대학』 — 「克明俊德」 (요전), 「日新 又日新」 (탕지반명), 「進賢退不肖」 (진서)를 통째 인용',
         {'font_size': 17, 'space_before': 8}),
        ('● 『중용』 — 「人心道心」(대우모)을 도통 전수의 비밀로 풀이',
         {'font_size': 17, 'space_before': 8}),
        ('● 『순자』·『묵자』·『여씨춘추』까지 — 선진 사상 거의 모두 서경을 인용 토대로 삼음',
         {'font_size': 17, 'space_before': 8}),
        ('● 서경 없이는 유학의 정치 철학과 도통론을 말할 수 없다',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC11)
def xi_song(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '송학(宋學)과 16자심전',
              '주희가 『중용장구서』에서 도통의 비밀로 풀다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 주희(朱熹) — 『중용장구서』 — 도통(道統)의 16자심전 정립',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 4}),
        ('     「人心惟危, 道心惟微, 惟精惟一, 允執厥中」 — 요·순·우·탕·문·무·주공·공자·맹자로 이어진 도의 비밀',
         {'font_size': 15, 'color': SUB, 'space_before': 4, 'font_name': 'Batang'}),
        ('● 채침(蔡沈) — 주희의 명을 받아 『서집전(書集傳)』 완성 (1209)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('     이후 동아시아 서경 강독의 표준 — 조선·일본·베트남 모두 채침본을 사용', {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 위(僞)고문 논쟁 — 청 염약거(閻若璩) 『상서고문소증(尚書古文疏證)』 (1745)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 12}),
        ('     고문 25편이 위작임을 정밀히 고증 — 그러나 사상적 가치는 살아남음', {'font_size': 15, 'color': SUB, 'space_before': 4}),
    ])


@S(SEC11)
def xi_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '한국과 일본에서의 서경')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 한국 — 삼국시대부터 수용, 조선조에 「치국의 교과서」로 자리잡음',
         {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('     · 세종·정조의 경연(經筵)에서 가장 자주 강독된 책', {'font_size': 15, 'space_before': 4}),
        ('     · 정약용 『상서고훈(尙書古訓)』·『매씨서평(梅氏書平)』 — 고문상서 비판의 정점', {'font_size': 15, 'space_before': 4}),
        ('     · 「민유방본」, 「천시자아민시」가 동학(東學)·민본 사상의 토양이 됨', {'font_size': 15, 'space_before': 4}),
        ('● 일본 — 헤이안 시대부터 한학자 강독, 에도 막부의 통치 이념 형성에 활용',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     · 야마자키 안사이(山崎闇齋)·아라이 하쿠세키(新井白石)의 서경 연구', {'font_size': 15, 'space_before': 4}),
        ('     · 「경천근민」의 이념이 막부 정치 윤리의 한 축', {'font_size': 15, 'space_before': 4}),
        ('● 베트남 — 한문 교양의 기초, 사서삼경 체계의 한 축',
         {'font_size': 18, 'space_before': 14}),
    ])


@S(SEC11)
def xi_modern(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '동아시아 한문학사에서의 위상')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 오경 중 「최고(最古)의 산문집」 — 시경이 운문이라면 서경은 산문의 원조',
         {'font_size': 18, 'space_before': 4}),
        ('● 한문 산문의 원형 — 後代 詔·策·誥·命 등 황실 문서 양식의 토대',
         {'font_size': 18, 'space_before': 8}),
        ('● 「3대(三代)의 정치」 = 「이상 정치」의 동의어가 됨 — 서경의 영향',
         {'font_size': 18, 'space_before': 8}),
        ('● 모든 동양 정치 어휘의 보고 — 천명·덕치·민본·왕도·중도(中)·황극(皇極) 등',
         {'font_size': 18, 'space_before': 8}),
        ('● 『사기』·『한서』·『자치통감』까지 — 모든 정사(正史)가 서경에서 출발',
         {'font_size': 18, 'space_before': 8}),
    ])


# ============== Ⅻ. 현대적 의의 ==============
SEC12 = 'Ⅻ. 현대적 의의'

@S(SEC12)
def xii_democracy(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '민주주의와 「민유방본」',
              '동양 민본주의는 현대 민주주의와 어떻게 만나는가')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 서양 민주주의 — 「민(民)이 곧 권력의 주체」 (popular sovereignty)',
         {'font_size': 18, 'space_before': 4}),
        ('● 동양 민본주의 — 「민(民)이 정치의 근본·목적」 (for the people)',
         {'font_size': 18, 'space_before': 10}),
        ('● 두 사상은 「주체」와 「목적」의 차이를 두고도 깊이 만남 — 모두 사람을 향한다',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 「天視自我民視」 = 「Vox populi, vox dei」 (백성의 소리가 곧 신의 소리)',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 21세기 정치 — 여론·선거·SNS 또한 「하늘의 시청(視聽)」이 발현되는 통로',
         {'font_size': 17, 'space_before': 12}),
    ])


@S(SEC12)
def xii_leadership(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '리더십의 5가지 원리 — 서경에서 길어낸 교훈')
    boxes = [
        ('1', '경천(敬天)',  '큰 책임 앞에서 두려워하라 — 권한은 빌린 것이다'),
        ('2', '근민(勤民)',  '편안함을 경계하라 — 안일이 곧 패망의 시작'),
        ('3', '경덕(敬德)',  '능력보다 덕 — 큰 자리는 큰 그릇이 채워야 한다'),
        ('4', '집중(執中)',  '극단이 아닌 중도 — 인심도심을 가려 잡으라'),
        ('5', '포현(包賢)',  '남의 재주를 자기 것처럼 — 시기 없이 인재를 안다'),
    ]
    for i, (num, name, desc) in enumerate(boxes):
        y = Inches(2.2 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.0), Inches(0.75), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.0), Inches(0.75),
                    num, font_size=24, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(1.9), y, Inches(2.5), Inches(0.75), PALE)
        add_textbox(slide, Inches(1.9), y, Inches(2.5), Inches(0.75),
                    name, font_size=17, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.6), y + Inches(0.05), Inches(8.3), Inches(0.7),
                    desc, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC12)
def xii_governance(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '정의로운 사법 — 「哀矜折獄」',
              '여형(呂刑) — 인본주의 형법의 원천')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「哀矜折獄」 — 슬퍼하고 가엾이 여기며 옥사를 결단하라',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 의심나는 죄는 가볍게, 의심나는 공은 무겁게 — 罪疑惟輕, 功疑惟重',
         {'font_size': 18, 'space_before': 12, 'font_name': 'Batang'}),
        ('● 형벌의 목적은 「벌(罰)」이 아니라 「교정(矯正)」',
         {'font_size': 18, 'space_before': 10}),
        ('● 현대 회복적 사법(restorative justice)·인본 사법의 동양적 원형',
         {'font_size': 17, 'space_before': 12, 'color': SUB}),
        ('● 「열 명의 죄인을 놓치는 것이, 한 사람의 무고한 자를 죽이는 것보다 낫다」',
         {'font_size': 17, 'space_before': 10}),
    ])


@S(SEC12)
def xii_today(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '서경이 오늘 우리에게 말하는 것')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 권력은 결국 「사람을 위한 것」이다 — 民惟邦本',
         {'font_size': 18, 'space_before': 6}),
        ('● 큰 자리는 큰 두려움이 따라야 한다 — 敬天',
         {'font_size': 18, 'space_before': 10}),
        ('● 편안함이 가장 위험하다 — 無逸',
         {'font_size': 18, 'space_before': 10}),
        ('● 사람의 마음은 늘 위태롭다 — 중도(中)를 잡지 못하면 흔들린다 — 允執厥中',
         {'font_size': 18, 'space_before': 10}),
        ('● 남의 재주를 자기 것처럼 기뻐할 수 있을 때, 비로소 큰 그릇이다 — 진서',
         {'font_size': 18, 'space_before': 10}),
        ('● 안다는 것이 어려운 게 아니라 행하는 것이 어렵다 — 非知之艱 行之惟艱',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


# ============== XIII. 마무리 ==============
SEC13 = 'XIII. 마무리'

@S(SEC13)
def xiii_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC13, n, t)
    add_title(slide, '서경, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 서경은 요순부터 동주 초까지, 약 1500년의 정치 문서 모음이다.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 우서 5 · 하서 4 · 상서 17 · 주서 32 — 사대(四代)의 말씀.',
         {'font_size': 18, 'space_before': 8}),
        ('● 典·謨·訓·誥·誓·命 — 여섯 갈래 문체로 정치 행위를 기록한다.',
         {'font_size': 18, 'space_before': 8}),
        ('● 천명(天命)·덕치(德治)·민본(民本)·경천근민 — 동양 정치의 네 기둥.',
         {'font_size': 18, 'space_before': 8}),
        ('● 16자심전 — 인심도심·惟精惟一·允執厥中 — 도통(道統)의 비밀.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「하늘의 뜻과 백성의 마음이 하나임을 일러주는 책」.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC13)
def xiii_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.0),
                '民 惟 邦 本',
                font_size=100, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.7),
                '本 固 邦 寧',
                font_size=60, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.5), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '백성이 곧 나라의 근본이니, 근본이 굳어야 나라가 평안하다',
                font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
                '— 서경, 하서 「오자지가(五子之歌)」',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                '書  經',
                font_size=22, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\서경.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
