# -*- coding: utf-8 -*-
"""
한비자 발표자료 — 망라적 102장 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 1. 표지 ==============
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8),
                '韓 非 子', font_size=104, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
                'Han Feizi · 한비자', font_size=22, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.3), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.55), Inches(12.3), Inches(0.5),
                '한비(韓非) 저 — 동양 정치철학 현실주의의 정점',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.4),
                '전국시대 말기 (BC 3세기) · 55편 약 10만 자 · 법(法)·술(術)·세(勢)의 완성',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.5),
                '"凡治天下, 必因人情"',
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
                '— 천하를 다스리려면 반드시 사람의 실정에 근거하라',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 2. 목차 ==============
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.8), Inches(0.7),
                '목 차', font_size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.15), Inches(12.8))
    items = [
        [('Ⅰ', '개요 — 한비자란 무엇인가'),
         ('Ⅱ', '한비의 생애 — 비극의 천재'),
         ('Ⅲ', '시대 배경 — 전국 말기'),
         ('Ⅳ', '법가 3선구 — 상앙·신불해·신도'),
         ('Ⅴ', '法 — 공개된 규칙'),
         ('Ⅵ', '術 — 군주의 은밀한 기법'),
         ('Ⅶ', '勢 — 지위의 권력'),
         ('Ⅷ', '三位一體 — 법·술·세의 통합'),
         ('Ⅸ', '인간관 — 이기심의 현실주의'),
         ('Ⅹ', '주도·유도·이병'),
         ('Ⅺ', '八姦 — 8가지 간사의 수법')],
        [('Ⅻ', '十過 — 10가지 치명적 실수'),
         ('ⅩⅢ', '고분·세난·화씨 — 개혁가의 비극'),
         ('ⅩⅣ', '亡徵 — 47가지 망국 징조'),
         ('ⅩⅤ', '비내 · 해로유로'),
         ('ⅩⅥ', '오두·현학 — 시대 비판'),
         ('ⅩⅦ', '유명 우화 10편'),
         ('ⅩⅧ', '10대 원리'),
         ('ⅩⅨ', '명구절 10선'),
         ('ⅩⅩ', '다른 사상과의 비교'),
         ('ⅩⅩⅠ', '현대적 의의'),
         ('ⅩⅩⅡ', '마무리')],
    ]
    for col, group in enumerate(items):
        x = 0.7 + col * 6.4
        top = 1.5
        for num, title in group:
            add_textbox(slide, Inches(x), Inches(top), Inches(1.0), Inches(0.4),
                        num, font_size=14, bold=True, color=ACCENT)
            add_textbox(slide, Inches(x + 1.0), Inches(top), Inches(5.3), Inches(0.4),
                        title, font_size=14, color=INK)
            top += 0.49


# ============== Ⅰ. 개요 ==============
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '한비자(韓非子)란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                '전국시대 말기 법가(法家) 사상의 집대성',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.5),
                '"인간은 이익으로 움직인다 — 그러므로 제도로 다스려라"',
                font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)
    nums = [('55', '편(篇)'), ('약 10만', '자(字)'), ('200+', '우화·일화'), ('1,700', '년 앞선 마키아벨리')]
    for i, (n, lbl) in enumerate(nums):
        x = 0.6 + i * 3.05
        add_textbox(slide, Inches(x), Inches(4.0), Inches(2.9), Inches(1.0),
                    n, font_size=46, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(5.1), Inches(2.9), Inches(0.5),
                    lbl, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                '동양 정치철학 현실주의의 정점 — 마키아벨리 『군주론』보다 1,700년 앞섬',
                font_size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '진(秦)의 천하 통일을 설계한 사상 — 중국 2,000년 제국 체제의 토대',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_status(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '책의 위상 — 핵심 정보')
    rows = [
        ('서명',  '한비자(韓非子)',                  '"한비 선생의 책"'),
        ('저자',  '한비(韓非, BC 280?~233)',         '한(韓)나라 공자(公子)'),
        ('시대',  '전국시대 말기',                    'BC 3세기 중반~말'),
        ('분량',  '55편 약 10만 자',                  '체계적 논설 + 200+ 우화'),
        ('학파',  '법가(法家)의 집대성',              '상앙·신불해·신도의 사상 통합'),
        ('성격',  '동양 현실주의 정치철학의 정점',    '마키아벨리 군주론보다 1,700년 앞섬'),
        ('영향',  '진(秦) 통일·중국 2,000년 제국',    '"외유내법(外儒內法)" 구조의 원형'),
    ]
    top = 2.15
    for i, (tag, val, note) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.65), PALE)
        add_filled_rect(slide, Inches(2.85), Inches(top), Inches(10.0), Inches(0.65), bg)
        add_textbox(slide, Inches(0.55), Inches(top + 0.18), Inches(2.2), Inches(0.4),
                    tag, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.95), Inches(top + 0.05), Inches(4.5), Inches(0.5),
                    val, font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.5), Inches(top + 0.08), Inches(5.3), Inches(0.5),
                    note, font_size=13, color=SUB)
        top += 0.7


@S('Ⅰ. 개요')
def s_idealism_failure(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '왜 한비자인가 — 이상주의의 실패에서 태어난 현실주의')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.45), Inches(11.3), Inches(1.8), [
        ('전국시대 250년의 끝자락 — 이상주의가 현실을 바꾸지 못한 시대',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 공자의 인(仁)은 제후들에게 외면당했고',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('• 묵자의 겸애(兼愛)는 실험에 그쳤으며',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('• 맹자의 왕도(王道)는 이상에 머물렀다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.6), Inches(4.5), Inches(11.9), Inches(2.4),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.9), Inches(4.65), Inches(11.3), Inches(2.2), [
        ('법가의 혁명성',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 이상이 아니라 현실에서 출발',
         {'font_size': 15, 'space_before': 4}),
        ('• 개인의 덕이 아니라 제도의 힘에 의존',
         {'font_size': 15, 'space_before': 4}),
        ('• 성인을 기다리지 않고 평범한 군주도 다스릴 수 있는 시스템 설계',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅰ. 개요')
def s_qin_love(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '진시황이 한비를 사랑한 이유')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('사마천 『사기』「노자한비열전」',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"진왕(秦王)이 『고분』·『오두』의 글을 보고 말하기를',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('「아! 내가 이 사람을 만나 사귈 수 있다면 죽어도 여한이 없겠다」"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('진시황은 한비의 책만 읽고 그를 꼭 만나고 싶어',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('한나라를 쳐서 한비를 인질로 데려오게 했다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 6}),
        ('→ 책 한 권이 한 사람의 운명과 한 제국의 운명을 바꾼 순간',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅰ. 개요')
def s_paradox_victory(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '한비의 역설적 승리 — "개인은 죽었으나 사상은 천하를 지배"')
    timeline = [
        ('BC 233', '한비 옥중 독살 — 이사의 모함',                       True),
        ('BC 221', '진시황이 6국 통일 — 한비 사후 12년',                 False),
        ('진(秦) 제도', '관료제·도량형·법령·중앙집권 — 한비가 설계한 법가 구현', False),
        ('한(漢) 이후', '외유내법(外儒內法) — 겉은 유교, 속은 법가 2,000년', False),
    ]
    top = 2.3
    for era, event, is_death in timeline:
        c = ACCENT if is_death else INK
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_death else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.95), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(3.0), Inches(0.5),
                    era, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.8), Inches(top), Inches(9.0), Inches(0.95), bg)
        add_textbox(slide, Inches(4.0), Inches(top + 0.27), Inches(8.7), Inches(0.5),
                    event, font_size=15, color=INK)
        top += 1.07
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"개혁가는 죽지만 제도는 남는다" — 한비 자신이 화씨편에서 예고한 운명',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 한비의 생애 ==============
@S('Ⅱ. 생애')
def s_han_prince(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 생애', page, total)
    add_title(slide, '한(韓)나라 공자 — 망해가는 조국의 왕자')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '韓\n非', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [
        ('출생', {'bold': True, 'font_size': 18, 'color': ACCENT}),
        ('  기원전 280년경 · 한나라 왕족 출생',
         {'font_size': 16}),
        ('한나라의 위치', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('  전국 7웅 중 가장 약한 나라',
         {'font_size': 16}),
        ('  위·진·초 사이에 끼여 끊임없이 영토 침탈',
         {'font_size': 16}),
        ('정치적 좌절', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('  조국을 구하려 한왕에게 법술을 여러 차례 건의',
         {'font_size': 16}),
        ('  한왕은 듣지 않았다 — 귀족·연줄 정치에 머묾',
         {'font_size': 16}),
        ('한비자의 정조(情調)', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 10}),
        ('  "조국을 구하려 했으나 구하지 못한 자의 절규"',
         {'font_size': 16, 'bold': True, 'color': INK}),
    ]
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5),
                   lines, line_spacing=1.35)


@S('Ⅱ. 생애')
def s_stutter(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 생애', page, total)
    add_title(slide, '말더듬(口吃)과 글의 천재 — 세난(說難)의 역설')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.45), Inches(11.3), Inches(1.8), [
        ('『사기』의 기록 — "韓非 口吃 — 말더듬이"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('말하는 데 어려움이 있었으나',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('글쓰기는 당대 최고였다',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.6), Inches(4.5), Inches(11.9), Inches(2.4),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.9), Inches(4.7), Inches(11.3), Inches(2.1), [
        ('깊은 역설 — 제12편 『세난(說難)』',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"설득의 어려움"을 분석한 명문 — 정작 저자 자신이 말을 더듬었다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 6}),
        ('"세상에서 가장 설득이 어려운 사람이 설득의 이론을 쓴다"',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('— 이것이 한비자의 출발점',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅱ. 생애')
def s_xunzi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 생애', page, total)
    add_title(slide, '순자(荀子) 문하 — 동학(同學) 이사(李斯)와 함께')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(1.0), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(11.9), Inches(0.6),
                '한비의 스승 — 순자(荀子, BC 313?~238)',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.55), Inches(11.5), Inches(3.4), [
        ('순자의 사상 토대 — 성악설(性惡說)', {'bold': True, 'font_size': 17, 'color': ACCENT}),
        ('  인간 본성은 악하니 예(禮)와 법(法)으로 교화해야 한다',
         {'font_size': 15, 'color': INK}),
        ('', {'font_size': 8}),
        ('한비의 발전 — 유교의 탯줄을 자르다',
         {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 10}),
        ('  순자의 "예(禮)"를 버리고 "법(法)"만을 취함 → 법가로 완전 이동',
         {'font_size': 15, 'color': INK}),
        ('', {'font_size': 8}),
        ('동문 동학 — 이사(李斯)',
         {'bold': True, 'font_size': 17, 'color': ACCENT, 'space_before': 10}),
        ('  훗날 진시황의 승상이 될 자 — 한비의 친구이자 결국 살해자',
         {'font_size': 15, 'color': INK}),
        ('  운명적 비극의 씨앗이 이미 학문 공동체 안에 있었다',
         {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)


@S('Ⅱ. 생애')
def s_qin_tragedy1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 생애', page, total)
    add_title(slide, '진(秦)나라에서의 비극 ① — 진왕의 감탄')
    timeline = [
        ('BC 234', '진(秦)이 한(韓)을 침공',                                  False),
        ('BC 234',  '한왕이 한비를 사절(使節)로 진에 파견',                   False),
        ('BC 233',  '진왕 정(政, 훗날 진시황)이 한비를 만남 — 감탄',          True),
        ('BC 233',  '진왕이 한비를 곁에 두려 함',                              True),
    ]
    top = 2.3
    for era, event, is_meet in timeline:
        c = ACCENT if is_meet else SUB
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.2), Inches(2.5), Inches(0.5),
                    era, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.3), Inches(top), Inches(9.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.5), Inches(top + 0.22), Inches(9.2), Inches(0.5),
                    event, font_size=15, color=INK)
        top += 0.95
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
                '진왕이 한비의 글을 읽은 직후 — 사상의 결정적 만남',
                font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 생애')
def s_qin_tragedy2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 생애', page, total)
    add_title(slide, '진(秦)나라에서의 비극 ② — 이사의 모함')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.6),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.3), [
        ('이사(李斯)의 모함',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"한비는 한나라 사람이다.',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 진이 쓰면 본심은 한을 위한 것이고,',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 쓰지 않고 돌려보내면 후환이 된다"',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('진왕은 일단 한비를 옥에 가둠',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.9), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.7), [
        ('이사가 몰래 독약을 보내 자살을 강요',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('한비는 진왕을 만나 변명하고자 했으나 기회를 얻지 못하고',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('기원전 233년 옥중에서 독살 — 향년 47세경',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅱ. 생애')
def s_sima_lament(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 생애', page, total)
    add_title(slide, '사마천의 탄식 — 비극의 완결')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('사마천 『사기』「노자한비열전」',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"비(非)는 『세난(說難)』을 지어 설득의 어려움을 깊이 알았으면서도',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('스스로는 그 화를 벗어나지 못했다"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('역설 ①',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"설득의 이론을 쓴 자가 자기 한 사람을 설득하지 못해 죽었다"',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('역설 ②',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('한비가 『고분』에 예고한 "법술지사가 기득권에 제거되는" 구조 그대로',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('자기 자신에게 실현되었다',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


# ============== Ⅲ. 시대 배경 ==============
@S('Ⅲ. 시대')
def s_warring_states(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 시대', page, total)
    add_title(slide, '전국시대(戰國時代, BC 403~221) — 250년의 격동')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(11.9), Inches(0.55),
                '전국 7웅 — 제(齊)·초(楚)·연(燕)·한(韓)·조(趙)·위(魏)·진(秦)',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.3), Inches(5.9), Inches(3.6), PALE)
    add_textbox(slide, Inches(0.6), Inches(3.45), Inches(5.9), Inches(0.5),
                '춘추 vs 전국 — 결정적 차이',
                font_size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(4.1), Inches(5.3), Inches(2.7), [
        ('• 춘추: 귀족 의례 전쟁', {'font_size': 14}),
        ('• 전국: 국가 총동원 체제', {'font_size': 14, 'bold': True, 'space_before': 6}),
        ('• 작은 나라가 큰 나라에 망함', {'font_size': 14, 'space_before': 6}),
        ('• 7웅 중 가장 약한 한나라', {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ('  — 한비의 절박함의 배경', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(3.3), Inches(5.9), Inches(3.6), PALE)
    add_textbox(slide, Inches(6.8), Inches(3.45), Inches(5.9), Inches(0.5),
                '한비의 시대 인식',
                font_size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(4.1), Inches(5.3), Inches(2.7), [
        ('"상고는 도덕으로 겨루고', {'font_size': 14}),
        (' 중고는 지모로 다투고', {'font_size': 14}),
        (' 당금(當今)은 힘으로 싸운다"', {'font_size': 14, 'bold': True, 'color': ACCENT}),
        ('', {'font_size': 6}),
        ('→ 옛 방법은 옛 시대용', {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ('   새 시대는 새 방법이 필요', {'font_size': 14, 'color': SUB}),
    ], line_spacing=1.3)


@S('Ⅲ. 시대')
def s_battle(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 시대', page, total)
    add_title(slide, '전쟁의 대형화 — 장평대전(長平之戰, BC 260)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(2.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(5.9), Inches(0.5),
                '춘추시대 전투', font_size=18, bold=True, color=SUB,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(3.1), Inches(5.3), Inches(1.6), [
        ('• 수천 명 규모', {'font_size': 16}),
        ('• 며칠 만에 끝', {'font_size': 16, 'space_before': 6}),
        ('• 귀족 의례 전쟁', {'font_size': 16, 'space_before': 6}),
        ('• 도덕적 명분 중시', {'font_size': 16, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.4)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.45), Inches(5.9), Inches(0.5),
                '전국시대 전투', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(3.1), Inches(5.3), Inches(1.6), [
        ('• 수십만 명 규모', {'font_size': 16, 'bold': True}),
        ('• 수년에 걸친 총력전', {'font_size': 16, 'bold': True, 'space_before': 6}),
        ('• 국가 총동원 체제', {'font_size': 16, 'bold': True, 'space_before': 6}),
        ('• 승부가 곧 멸망', {'font_size': 16, 'color': ACCENT, 'space_before': 6}),
    ], line_spacing=1.4)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('BC 260년 장평대전 — 진(秦) vs 조(趙)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('진(秦)이 조(趙)의 포로 40만 명을 생매장(坑殺)',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('"이 규모의 전쟁에서는 귀족 의례가 불가능하다 — 총동원만 가능하다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅲ. 시대')
def s_iron(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 시대', page, total)
    add_title(slide, '철기·인구 폭증·상앙의 변법 — 사회 토대의 격변')
    items = [
        ('철기 보급',         '청동기 → 철기 농기구 → 농업 생산성 폭증'),
        ('인구 폭발',         '한 변경에서 수만 → 수십만 인구의 도시'),
        ('상앙의 변법 (BC 356)','정전제 폐지 · 개간 장려 · 군공에 따른 신분 상승'),
        ('농전(農戰)의 총동원', '"모든 역량을 농업과 전쟁에 집중"의 법가 체제'),
        ('승자의 법칙',       '귀족 의례가 아니라 시스템 설계가 경쟁의 승자를 가름'),
    ]
    top = 2.5
    for tag, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.5), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.17), Inches(3.5), Inches(0.4),
                    tag, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.3), Inches(top), Inches(8.5), Inches(0.7), PALE)
        add_textbox(slide, Inches(4.5), Inches(top + 0.17), Inches(8.2), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.8
    add_textbox(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
                '→ 이런 토대 위에서만 한비자의 "제도 우선" 사상이 설득력을 가졌다',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅲ. 시대')
def s_baekga(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 시대', page, total)
    add_title(slide, '제자백가(諸子百家) — 사상의 폭발')
    schools = [
        ('유가(儒家)',   '공자·맹자·순자',   '인(仁)·예(禮)·덕치',     SUB),
        ('도가(道家)',   '노자·장자',         '무위자연(無爲自然)',     SUB),
        ('묵가(墨家)',   '묵자',              '겸애(兼愛)·비공(非攻)', SUB),
        ('명가(名家)',   '혜시·공손룡',       '명실(名實)의 변(辯)',    SUB),
        ('종횡가(縱橫家)','소진·장의',         '합종연횡의 외교술',      SUB),
        ('법가(法家)',   '관중·상앙·신불해·신도·한비', '법·술·세의 통치',  ACCENT),
    ]
    top = 2.2
    for name, figures, idea, color in schools:
        is_legalist = '법가' in name
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_legalist else PALE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(0.65), color)
        add_textbox(slide, Inches(0.5), Inches(top + 0.17), Inches(2.5), Inches(0.4),
                    name, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.05), Inches(top), Inches(4.5), Inches(0.65), bg)
        add_textbox(slide, Inches(3.2), Inches(top + 0.17), Inches(4.3), Inches(0.4),
                    figures, font_size=13, color=INK)
        add_filled_rect(slide, Inches(7.6), Inches(top), Inches(5.2), Inches(0.65), bg)
        add_textbox(slide, Inches(7.8), Inches(top + 0.17), Inches(4.9), Inches(0.4),
                    idea, font_size=13, color=INK)
        top += 0.7
    add_textbox(slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
                '제후들이 인재를 모으는 시대 — 천하를 놓고 사상이 경쟁',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅳ. 법가 3선구 ==============
@S('Ⅳ. 법가 3선구')
def s_3pioneers(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 법가 3선구', page, total)
    add_title(slide, '법가 3선구 — 한비 이전의 세 거장')
    pioneers = [
        ('상앙\n商鞅', '法', '법(法)', '명문화된 공개 규칙'),
        ('신불해\n申不害', '術', '술(術)', '군주의 은밀한 통치술'),
        ('신도\n愼到', '勢', '세(勢)', '지위에서 나오는 권력'),
    ]
    for i, (name, key, han, role) in enumerate(pioneers):
        x = 0.5 + i * 4.3
        add_filled_rect(slide, Inches(x), Inches(2.3), Inches(4.1), Inches(4.7), PALE)
        add_textbox(slide, Inches(x), Inches(2.5), Inches(4.1), Inches(1.3),
                    key, font_size=130, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.0), Inches(4.1), Inches(0.6),
                    name, font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 0.3), Inches(4.95), Inches(3.5), Inches(0.5),
                        ACCENT)
        add_textbox(slide, Inches(x + 0.3), Inches(5.05), Inches(3.5), Inches(0.4),
                    han, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.3), Inches(5.6), Inches(3.5), Inches(1.0),
                    role, font_size=14, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '한비는 세 사람의 정수를 합쳐 "법·술·세"를 통합 완성',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 법가 3선구')
def s_sangang(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 법가 3선구', page, total)
    add_title(slide, '상앙(商鞅) — 法의 화신', '진(秦)의 변법(變法) 주도자')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '업적 — BC 356년 변법(變法)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('• 법의 명문화·공개·일관 적용', {'font_size': 14}),
        ('• 정전제 폐지·개간 장려', {'font_size': 14, 'space_before': 6}),
        ('• 군공(軍功)에 따른 작위', {'font_size': 14, 'space_before': 6}),
        ('• 연좌제·도량형 통일', {'font_size': 14, 'space_before': 6}),
        ('• 진을 강국으로 만든 토대', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('저술: 『상군서(商君書)』', {'font_size': 13, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '한계와 비극', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('• 법은 세웠으나 신하 통제술(術) 부재',
         {'font_size': 14, 'bold': True}),
        ('• 효공 사후 → 거열형(車裂刑) 처형', {'font_size': 14, 'color': ACCENT, 'space_before': 6}),
        ('  자기가 세운 법에 의해 죽음', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 4}),
        ('한비의 진단', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('"법만으로는 부족하다 —', {'font_size': 14}),
        (' 술(術)이 있어야 법이 보존된다"', {'font_size': 14}),
    ], line_spacing=1.3)


@S('Ⅳ. 법가 3선구')
def s_shenbuhai(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 법가 3선구', page, total)
    add_title(slide, '신불해(申不害) — 術의 대가', '한(韓)나라 재상')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '업적 — 술(術)의 정립', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('• 형명참동(形名參同) 시스템', {'font_size': 14, 'bold': True}),
        ('• 군주의 의중 비노출', {'font_size': 14, 'space_before': 6}),
        ('• 여러 정보원 활용', {'font_size': 14, 'space_before': 6}),
        ('• 능력 중심 인사', {'font_size': 14, 'space_before': 6}),
        ('• 한 재상 시절 한이 강성', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('저술: 『신자(申子)』 — 전해지지 않음',
         {'font_size': 13, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '한계', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('• 법을 세우지 않음', {'font_size': 14, 'bold': True}),
        ('• 일관된 기준 부족', {'font_size': 14, 'space_before': 6}),
        ('• 신불해 사후 술이 전수되지 못함', {'font_size': 14, 'space_before': 6, 'color': ACCENT}),
        ('', {'font_size': 4}),
        ('한비의 진단', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('"술만으로는 부족하다 —', {'font_size': 14}),
        (' 법(法)이 뼈대가 되어야 한다"', {'font_size': 14}),
    ], line_spacing=1.3)


@S('Ⅳ. 법가 3선구')
def s_shendao(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 법가 3선구', page, total)
    add_title(slide, '신도(愼到) — 勢의 이론가')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '핵심 사상 — 勢의 절대성', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('"위치만 있으면 누구나 다스린다"',
         {'font_size': 14, 'bold': True}),
        ('', {'font_size': 6}),
        ('• 요(堯)가 필부였다면', {'font_size': 14, 'space_before': 6}),
        ('  세 사람도 다스리지 못함', {'font_size': 14, 'color': SUB}),
        ('• 걸(桀)이 천자였기에', {'font_size': 14, 'space_before': 6}),
        ('  천하를 어지럽힘', {'font_size': 14, 'color': SUB}),
        ('', {'font_size': 4}),
        ('저술: 『신자(愼子)』', {'font_size': 13, 'color': SUB, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '한계', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('• 세(勢)만으로는 부족', {'font_size': 14, 'bold': True}),
        ('• 법이 없으면 자의적 통치', {'font_size': 14, 'space_before': 6}),
        ('• 술이 없으면 통제 불가', {'font_size': 14, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('한비의 진단', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('"세에 법과 술이 결합해야', {'font_size': 14}),
        (' 비로소 완성된 통치"', {'font_size': 14}),
    ], line_spacing=1.3)


# ============== Ⅴ. 法 ==============
@S('Ⅴ. 法')
def s_law_def(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 法', page, total)
    add_title(slide, '法 — 공개된 규칙', '"법이란 도적(圖籍)에 명확히 기록하여 백성에게 공포한 것이다"')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '法', font_size=240, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(4.2), Inches(2.4), Inches(8.7), Inches(4.5), [
        ('원문 (제6편 유도)', {'bold': True, 'font_size': 16, 'color': ACCENT}),
        ('"法者 編著之圖籍,', {'font_size': 18, 'bold': True}),
        (' 設之於官府,',     {'font_size': 18, 'bold': True}),
        (' 而布之於百姓者也"', {'font_size': 18, 'bold': True}),
        ('', {'font_size': 6}),
        ('해석', {'bold': True, 'font_size': 16, 'color': ACCENT, 'space_before': 10}),
        ('법이란 도적(圖籍·문서)에 명확히 기록하여',
         {'font_size': 15}),
        ('관부에 비치하고',
         {'font_size': 15}),
        ('백성에게 공포(公布)한 것이다',
         {'font_size': 15, 'bold': True, 'color': INK}),
    ], line_spacing=1.3)


@S('Ⅴ. 法')
def s_law_4req(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 法', page, total)
    add_title(slide, '법의 4가지 요건')
    items = [
        ('1', '명문화(明文化)',  '반드시 글로 적혀 있어야 함',
         '구두 전승은 법이 아님'),
        ('2', '공개성(公開性)',  '백성이 모두 알 수 있도록 공포',
         '비밀법은 법이 아님'),
        ('3', '일관성(一貫性)',  '어제와 오늘, 이 사람과 저 사람에게 다르게 적용되지 않음',
         '상황·기분에 따른 변화 금지'),
        ('4', '평등성(平等性)',  '신분과 무관하게 동일 적용',
         '"刑過不避大臣, 賞善不遺匹夫"'),
    ]
    top = 2.3
    for num, title, desc, note in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.9), Inches(1.05), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.32), Inches(0.9), Inches(0.5),
                    num, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(2.6), Inches(1.05), PALE)
        add_textbox(slide, Inches(1.7), Inches(top + 0.3), Inches(2.6), Inches(0.5),
                    title, font_size=17, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.5), Inches(top), Inches(8.3), Inches(1.05),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.7), Inches(top + 0.13), Inches(8.0), Inches(0.45),
                    desc, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(4.7), Inches(top + 0.55), Inches(8.0), Inches(0.45),
                    note, font_size=12, color=SUB)
        top += 1.18


@S('Ⅴ. 法')
def s_imokjisin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 法', page, total)
    add_title(slide, '이목지신(移木之信) — 상앙의 변법 시작')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(3.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.5), Inches(11.3), Inches(2.7), [
        ('BC 356년 — 상앙이 진나라에서 변법을 시행하기 전',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('도성 남문에 나무 기둥을 세우고 포고:',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"이 기둥을 북문으로 옮기는 자에게 금 열 냥을 준다"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('아무도 믿지 않자 상금을 50냥으로 올림',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('한 사람이 옮기자 즉시 50냥을 지급',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.6), Inches(5.5), Inches(11.9), Inches(1.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.9), Inches(5.65), Inches(11.3), Inches(1.3), [
        ('"법이 반드시 실행된다(法必行)"는 신뢰 구축',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('법치의 첫 단계는 "약속을 지키는 권력자"임을 증명하는 것',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅴ. 法')
def s_law_vs_man(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 法', page, total)
    add_title(slide, '법치(法治) vs 인치(人治)')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(5.9), Inches(0.6), SUB)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(0.5),
                '인치(人治)', font_size=20, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(6.9), Inches(2.2), Inches(5.9), Inches(0.6), ACCENT)
    add_textbox(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(0.5),
                '법치(法治)', font_size=20, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    rows = [
        ('군주의 개인 능력에 의존',  '시스템에 의존'),
        ('기분에 따른 정책 변화',     '일관된 기준'),
        ('신하의 아첨 유발',           '능력 평가 가능'),
        ('카리스마 사라지면 붕괴',     '시스템이 계속 작동'),
        ('잡스 같은 천재 필요',        '평범한 후계자도 운영 가능'),
    ]
    top = 3.0
    for left, right in rows:
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(5.9), Inches(0.65), PALE)
        add_textbox(slide, Inches(0.7), Inches(top + 0.17), Inches(5.5), Inches(0.4),
                    left, font_size=14, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.9), Inches(top), Inches(5.9), Inches(0.65),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(7.1), Inches(top + 0.17), Inches(5.5), Inches(0.4),
                    right, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        top += 0.75
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '"國無常強 無常弱 奉法者強 則國強" — 법을 받드는 자가 강하면 나라가 강하다',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅵ. 術 ==============
@S('Ⅵ. 術')
def s_sul_def(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 術', page, total)
    add_title(slide, '術 — 군주의 은밀한 통치 기법')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '術', font_size=240, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(4.2), Inches(2.4), Inches(8.7), Inches(4.5), [
        ('원문 (제43편 정법)', {'bold': True, 'font_size': 16, 'color': ACCENT}),
        ('"術者', {'font_size': 17, 'bold': True}),
        ('  因任而授官,', {'font_size': 17, 'bold': True}),
        ('  循名而責實,', {'font_size': 17, 'bold': True}),
        ('  操殺生之柄,', {'font_size': 17, 'bold': True}),
        ('  課群臣之能者也"', {'font_size': 17, 'bold': True}),
        ('', {'font_size': 6}),
        ('해석', {'bold': True, 'font_size': 16, 'color': ACCENT, 'space_before': 8}),
        ('능력에 따라 관직을 주고,',
         {'font_size': 14}),
        ('이름(직책)에 따라 실적을 묻고,',
         {'font_size': 14}),
        ('생살의 권한을 쥐고,',
         {'font_size': 14}),
        ('신하의 능력을 검증한다',
         {'font_size': 14, 'bold': True, 'color': INK}),
    ], line_spacing=1.25)


@S('Ⅵ. 術')
def s_hyeongmyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 術', page, total)
    add_title(slide, '형명참동(形名參同) — 현대 KPI·OKR의 원조')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '循 名 責 實',
                font_size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '순명책실 — "이름(名)에 따라 실적(實)을 묻는다"',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    steps = [
        ('1', '약속(名)',     '신하가 "이 일을 맡으면 이렇게 하겠습니다"'),
        ('2', '실행(實)',     '신하가 일을 수행'),
        ('3', '대조(參)',     '결과(形)를 약속(名)과 대조'),
        ('4', '판정',         '일치→상 / 미달→벌 / 초과→벌(월권)'),
    ]
    top = 4.05
    for num, tag, desc in steps:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.65), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.12), Inches(0.7), Inches(0.4),
                    num, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.5), Inches(0.65), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    tag, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.0), Inches(top), Inches(8.8), Inches(0.65),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.2), Inches(top + 0.13), Inches(8.5), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.75


@S('Ⅵ. 術')
def s_excess_punished(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 術', page, total)
    add_title(slide, '"왜 초과 성취도 벌하는가" — 월권(越權)의 원리')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('한비의 유명한 일화 — 모자 담당과 옷 담당',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('한 군주가 술에 취해 잠들었다 — 추워 보였다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('모자 담당 신하가 친절히 군주에게 옷을 덮어주었다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('깨어난 군주는 모자 담당과 옷 담당을 모두 벌했다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('옷 담당 — 자기 일을 안 했으므로 벌',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('모자 담당 — "친절"이 아니라 "월권(越權)"이므로 벌',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"직분의 경계를 흐리면 조직이 무너진다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


# ============== Ⅶ. 勢 ==============
@S('Ⅶ. 勢')
def s_se_def(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 勢', page, total)
    add_title(slide, '勢 — 지위에서 나오는 권력')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '勢', font_size=240, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(4.2), Inches(2.4), Inches(8.7), Inches(4.5), [
        ('원문 (제40편 난세)', {'bold': True, 'font_size': 16, 'color': ACCENT}),
        ('"夫有材而無勢',  {'font_size': 22, 'bold': True}),
        (' 雖賢不能制不肖"', {'font_size': 22, 'bold': True}),
        ('', {'font_size': 6}),
        ('해석', {'bold': True, 'font_size': 16, 'color': ACCENT, 'space_before': 10}),
        ('재능이 있어도 세(勢)가 없으면',
         {'font_size': 16}),
        ('아무리 어질어도',
         {'font_size': 16}),
        ('불초한 자를 제어할 수 없다',
         {'font_size': 16, 'bold': True, 'color': INK}),
        ('', {'font_size': 6}),
        ('의미', {'bold': True, 'font_size': 16, 'color': ACCENT, 'space_before': 10}),
        ('지위·자리·상황·구조에서 나오는 힘',
         {'font_size': 15}),
    ], line_spacing=1.3)


@S('Ⅶ. 勢')
def s_yo_geol(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 勢', page, total)
    add_title(slide, '요(堯)와 걸(桀)의 비유 — 직위가 곧 힘이다')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '요(堯) — 성왕(聖王)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('만약 필부(匹夫)였다면?',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('세 사람도', {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('다스리지 못한다', {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('아무리 어질어도 세(勢) 없으면', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('어떤 권위도 없다', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '걸(桀) — 폭군', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.1), Inches(5.3), Inches(3.5), [
        ('만약 천자(天子)였다면?',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('천하를', {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('어지럽힐 수 있다', {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('아무리 어리석어도 세(勢) 있으면', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('힘을 행사한다', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅶ. 勢')
def s_se_3principles(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 勢', page, total)
    add_title(slide, '세(勢)의 3원칙')
    items = [
        ('1', '권력을 나누지 마라',  '한 자루의 칼(이병二柄)은 군주만이 쥔다'),
        ('2', '권위를 지켜라',        '군주가 가벼워 보이면 모든 것이 무너진다'),
        ('3', '구조를 설계하라',      '개인의 능력이 아니라 시스템에 의존'),
    ]
    top = 2.5
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.9), Inches(1.1), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.32), Inches(0.9), Inches(0.5),
                    num, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.7), Inches(top), Inches(3.5), Inches(1.1), PALE)
        add_textbox(slide, Inches(1.7), Inches(top + 0.35), Inches(3.5), Inches(0.5),
                    title, font_size=18, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.4), Inches(top), Inches(7.4), Inches(1.1),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.6), Inches(top + 0.35), Inches(7.1), Inches(0.5),
                    desc, font_size=15, color=INK)
        top += 1.3


# ============== Ⅷ. 三位一體 ==============
@S('Ⅷ. 三位一體')
def s_jeongbeop(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 三位一體', page, total)
    add_title(slide, '정법편(定法) — 상앙과 신불해의 한계 통합')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"申不害不擅其法 不能獨行其術',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 商鞅不用其術 不能長保其法"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('"신불해는 법을 겸하지 않아 술만으로 독행할 수 없었고',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        (' 상앙은 술을 쓰지 않아 법을 오래 보존할 수 없었다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('한비의 통합 — 법·술·세 三位一體',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('법: 뼈대 (공개적·일관적 기준)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('술: 운용 (은밀한 관리 기법)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('세: 동력 (지위에서 나오는 권위)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('법이 술의 남용을 방지 ↔ 술이 법의 경직성을 보완',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅷ. 三位一體')
def s_triangle(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 三位一體', page, total)
    add_title(slide, '법·술·세 삼각형 — 한비의 최종 종합')
    # 큰 삼각형 시각화 (3개 원으로)
    add_filled_rect(slide, Inches(5.5), Inches(2.3), Inches(2.3), Inches(1.6), ACCENT)
    add_textbox(slide, Inches(5.5), Inches(2.55), Inches(2.3), Inches(0.5),
                '法', font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.5),
                '규칙·뼈대', font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(1.5), Inches(4.4), Inches(2.3), Inches(1.6), ACCENT)
    add_textbox(slide, Inches(1.5), Inches(4.65), Inches(2.3), Inches(0.5),
                '術', font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(5.4), Inches(2.3), Inches(0.5),
                '기법·운용', font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(9.5), Inches(4.4), Inches(2.3), Inches(1.6), ACCENT)
    add_textbox(slide, Inches(9.5), Inches(4.65), Inches(2.3), Inches(0.5),
                '勢', font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.5), Inches(5.4), Inches(2.3), Inches(0.5),
                '권력·동력', font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    # 연결선
    add_textbox(slide, Inches(5.5), Inches(6.4), Inches(2.3), Inches(0.5),
                '통치 시스템',
                font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
                '세 가지가 모두 있어야 비로소 완전한 통치 — 한비의 최종 종합',
                font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅸ. 인간관 ==============
@S('Ⅸ. 인간관')
def s_human_nature(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 인간관', page, total)
    add_title(slide, '凡治天下 必因人情 — 인간관의 근본 명제')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '凡 治 天 下   必 因 人 情',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '범치천하 필인인정',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
                '— 제48편 팔경(八經)', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
                '"천하를 다스리려면 반드시 사람의 실정(人情)에 근거한다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.1), [
        ('사람은 좋아하고 싫어한다(有好惡) → 상벌(賞罰)을 쓸 수 있다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('인간은 이익(利)을 좋아하고, 해(害)를 싫어한다 — 이 본성은 변하지 않는다',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅸ. 인간관')
def s_cart_coffin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 인간관', page, total)
    add_title(slide, '수레 장인 vs 관 장인 — 한비 인간관의 상징')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.32), Inches(12.3), Inches(0.5),
                '輿人欲人之富貴, 匠人欲人之夭死',
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.1), Inches(5.9), Inches(2.8), PALE)
    add_textbox(slide, Inches(0.6), Inches(3.25), Inches(5.9), Inches(0.5),
                '수레 장인(輿人)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.85), Inches(5.3), Inches(1.9), [
        ('사람들이 부자이기를 바람',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('"사람이 귀해지지 않으면',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        (' 수레가 팔리지 않는다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(3.1), Inches(5.9), Inches(2.8),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(3.25), Inches(5.9), Inches(0.5),
                '관 장인(匠人)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.85), Inches(5.3), Inches(1.9), [
        ('사람들이 일찍 죽기를 바람',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('"사람이 죽지 않으면',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        (' 관이 팔리지 않는다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_textbox(slide, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.5),
                '수레 장인이 인자해서가 아니고, 관 장인이 잔인해서도 아니다',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '→ 애덤 스미스의 "푸줏간 주인의 이익"(1776)보다 2,000년 앞선 선언',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅸ. 인간관')
def s_parent_child(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 인간관', page, total)
    add_title(slide, '부모와 자식 간에도 이해타산 — 가장 충격적 관찰')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('제50편 현학 — 충격적 관찰',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"父母之於子也',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        (' 産男則相賀, 産女則殺之"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"부모가 자식에 대해, 아들을 낳으면 서로 축하하고',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        (' 딸을 낳으면 죽이기도 한다"',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"훗날의 편익을 생각하고 장기적 이익을 계산하기 때문이다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"가장 가까운 부모 자식 간에도 이해타산이 작동한다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 하물며 남남 사이에서 순수한 도덕만으로 질서를 유지하겠는가?"',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅸ. 인간관')
def s_ruler_no_exception(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 인간관', page, total)
    add_title(slide, '군주도 예외가 아니다 — 비내(備內)의 충격')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '人主之患在於信人   信人則制於人',
                font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.5),
                '인주지환재어신인 신인즉제어인',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.75), Inches(12.3), Inches(0.4),
                '— 제17편 비내', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.35), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.6),
                '"군주의 걱정은 사람을 믿는 데 있다 — 믿으면 사람에게 제어당한다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.5), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.65), Inches(11.3), Inches(1.3), [
        ('"군주가 죽으면 이익을 얻는 자가 누구인가?"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('태자(빨리 왕 되기) · 왕비(총애 상실 방지) · 권신(새 군주의 공신)',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('→ "모두를 의심하고, 제도로 통제하라"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅸ. 인간관')
def s_zero_trust(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 인간관', page, total)
    add_title(slide, '제로 트러스트(Zero Trust) 철학 — 2,300년 전의 선언')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"夫以妻之近與子之親而猶不可信',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 則其餘無可信者矣"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"아내의 가까움과 자식의 친밀함으로도 믿을 수 없다면',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        (' 나머지는 믿을 자가 없다"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 정보보안의 핵심 원칙',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"Never trust, always verify"',
         {'font_size': 22, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('— 결코 신뢰하지 말고 항상 검증하라 —',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('Google·Microsoft가 채택한 21세기 보안 패러다임 = 한비의 비내(備內)',
         {'font_size': 13, 'color': ACCENT, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('Ⅸ. 인간관')
def s_three_conclusions(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 인간관', page, total)
    add_title(slide, '한비 인간관의 3대 귀결')
    items = [
        ('1', '도덕 교화는 한계가 있다', '성인(聖人)을 기다리면 나라가 망한다'),
        ('2', '법과 제도로 이기심을 활용하라', '이익이 곧 공익이 되게 설계하라'),
        ('3', '견제와 균형은 의심에서 나온다', '아무도 믿지 말고 모두를 검증하라'),
    ]
    top = 2.4
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(1.2), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.35), Inches(1.0), Inches(0.5),
                    num, font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.8), Inches(top), Inches(11.0), Inches(1.2), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.2), Inches(10.6), Inches(0.5),
                    title, font_size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.0), Inches(top + 0.7), Inches(10.6), Inches(0.4),
                    desc, font_size=15, color=INK)
        top += 1.4
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '현대 민주주의의 삼권분립·기업의 내부통제·조직 거버넌스와 동일한 논리',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅹ. 주도·유도·이병 ==============
@S('Ⅹ. 주도·유도·이병')
def s_judo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 주도·유도·이병', page, total)
    add_title(slide, '주도(主道) — 군주의 도리 · 허정무위(虛靜無爲)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '제5편 주도(主道) — "속마음을 숨기고 법·제도로 자동 통치"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.8), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.6), [
        ('"是故 去智而有明',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 去賢而有功 去勇而有強"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"자기 지혜를 버려야 밝아지고, 자기 어짐을 버려야 공이 있고, 자기 용기를 버려야 강해진다"',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.45), Inches(11.7), Inches(1.5), [
        ('3대 핵심 개념',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('허정무위(虛靜無爲) — 마음을 비우고 표정·의견을 드러내지 않음',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('형명참동(形名參同) — 신하의 말(名)과 실적(形) 대조',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('무위이치(無爲而治) — 직접 하지 않으면서 이루어지게',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅹ. 주도·유도·이병')
def s_show_preference(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 주도·유도·이병', page, total)
    add_title(slide, '"군주가 좋아하는 것을 드러내면" — 관찰자 효과')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '君 見 其 所 欲   臣 自 將 雕 琢',
                font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '군견기소욕 신자장조탁',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
                '— 제5편 주도', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '"군주가 좋아하는 것을 드러내면 신하들은 스스로 자신을 꾸민다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.5), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.3), [
        ('현대 사례 — 아마존 제프 베조스의 원칙',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('회의에서 가장 나중에 발언 — CEO가 먼저 의견 내면 팀은 그에 맞춰 말한다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('관찰자의 선호를 드러내면 관찰 대상이 그에 맞게 위장된다',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅹ. 주도·유도·이병')
def s_yudo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 주도·유도·이병', page, total)
    add_title(slide, '유도(有度) — 법치(法治) > 인치(人治)')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '國 無 常 強   無 常 弱',
                font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.9),
                '奉 法 者 強   則 國 強',
                font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.5),
                '국무상강 무상약 · 봉법자강 즉국강',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.55), Inches(12.3), Inches(0.4),
                '— 제6편 유도', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.15), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.6),
                '"나라의 강약은 일정하지 않다. 법을 받드는 자가 강하면 나라가 강하다"',
                font_size=17, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.9), PALE)
    add_textbox(slide, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.5),
                '현대 — 잡스 사후에도 애플이 유지되는 이유는 시스템(법치)이 완성되어 있었기 때문',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅹ. 주도·유도·이병')
def s_equal_punish(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 주도·유도·이병', page, total)
    add_title(slide, '刑過不避大臣 賞善不遺匹夫 — 법 앞의 평등')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '刑 過 不 避 大 臣',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.9),
                '賞 善 不 遺 匹 夫',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                '형과불피대신 상선불유필부',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
                '— 제6편 유도', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.45), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.6),
                '"형벌은 대신(大臣)도 피하지 못하고, 상은 평민(匹夫)도 놓치지 않는다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                '→ 현대 법치국가·"법 앞의 평등"의 2,300년 전 선언',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅹ. 주도·유도·이병')
def s_ibyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 주도·유도·이병', page, total)
    add_title(slide, '이병(二柄) — 상벌의 두 자루')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '二 柄 者   刑 德 也',
                font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '이병자 형덕야 — "두 자루는 형(刑)과 덕(德)이다"',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
                '— 제7편 이병', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(4.7), Inches(5.9), Inches(2.3),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(0.6), Inches(4.85), Inches(5.9), Inches(0.5),
                '刑 — 형(刑)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(5.5), Inches(5.9), Inches(0.5),
                '"殺戮之謂刑"', font_size=18, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(6.1), Inches(5.9), Inches(0.7),
                '살육·처벌 — 두려움의 자루',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(6.8), Inches(4.7), Inches(5.9), Inches(2.3), PALE)
    add_textbox(slide, Inches(6.8), Inches(4.85), Inches(5.9), Inches(0.5),
                '德 — 덕(德)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(5.5), Inches(5.9), Inches(0.5),
                '"慶賞之謂德"', font_size=18, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(6.1), Inches(5.9), Inches(0.7),
                '경상(慶賞) — 동기 부여의 자루',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅹ. 주도·유도·이병')
def s_tiger_teeth(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 주도·유도·이병', page, total)
    add_title(slide, '"호랑이의 발톱과 이빨" — 상벌권을 넘기면 군주가 굴복한다')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.6),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.3), [
        ('"虎之所以能服狗者, 爪牙也.',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 使虎釋其爪牙而使狗用之,',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 則虎反服於狗矣"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"호랑이가 개를 굴복시키는 것은 발톱과 이빨 때문이다.',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 호랑이가 발톱과 이빨을 버리고 개에게 주면, 호랑이가 개에게 굴복한다"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.7), [
        ('"상벌권이 신하에게 넘어가는 순간 군주는 형식뿐"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('역사 — 조선 세도정치: 안동 김씨·풍양 조씨가 인사권·사법권 장악 → 왕은 허수아비',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('현대 — 중간관리자의 평가권 독점 → 충성이 CEO 아닌 중간관리자로 이동',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


# ============== Ⅺ. 八姦 ==============
@S('Ⅺ. 八姦')
def s_8gan_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 八姦', page, total)
    add_title(slide, '팔간(八姦) — 8가지 간사의 수법')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(11.9), Inches(0.55),
                '"신하가 군주를 속이고 권력을 빼앗는 여덟 가지 수법"',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('1', '동상(同牀)',  '후궁·왕비 이용'),
        ('2', '재방(在旁)',  '측근·환관 이용'),
        ('3', '부형(父兄)',  '왕족·외척 이용'),
        ('4', '양앙(養殃)',  '유흥·향락 이용'),
        ('5', '민맹(民萌)',  '백성의 인기 이용'),
        ('6', '유행(流行)',  '유세객·변론 이용'),
        ('7', '위강(威強)',  '무력·군사 이용'),
        ('8', '사방(四方)',  '외세 이용'),
    ]
    top = 3.3
    for i, (num, name, role) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.4
        y = top + row * 0.85
        add_filled_rect(slide, Inches(x), Inches(y), Inches(0.7), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(x), Inches(y + 0.15), Inches(0.7), Inches(0.4),
                    num, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 0.8), Inches(y), Inches(2.5), Inches(0.7), PALE)
        add_textbox(slide, Inches(x + 0.8), Inches(y + 0.18), Inches(2.5), Inches(0.4),
                    name, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 3.4), Inches(y), Inches(2.6), Inches(0.7),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(x + 3.5), Inches(y + 0.18), Inches(2.5), Inches(0.4),
                    role, font_size=13, color=INK)


@S('Ⅺ. 八姦')
def s_8gan_1_4(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 八姦', page, total)
    add_title(slide, '팔간 1~4 — 가까운 자들의 함정')
    items = [
        ('1', '동상(同牀)', '잠자리를 함께하는 자(후궁·왕비)',
         '현대 — CEO 배우자의 인사 간섭'),
        ('2', '재방(在旁)', '곁에 있는 자(측근·환관)',
         '현대 — 비서실의 비공식 권력'),
        ('3', '부형(父兄)', '아버지와 형(왕족·외척)',
         '현대 — 오너 일가의 전횡'),
        ('4', '양앙(養殃)', '재앙을 기르는 자(유흥·향락)',
         '현대 — 과도한 접대, 골프·요트 접대'),
    ]
    top = 2.3
    for num, han, desc, modern in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(1.05), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.32), Inches(0.7), Inches(0.5),
                    num, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.5), Inches(1.05), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.32), Inches(2.5), Inches(0.5),
                    han, font_size=18, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.0), Inches(top), Inches(8.8), Inches(1.05),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.2), Inches(top + 0.12), Inches(8.5), Inches(0.4),
                    desc, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(4.2), Inches(top + 0.55), Inches(8.5), Inches(0.4),
                    modern, font_size=12, color=SUB)
        top += 1.18


@S('Ⅺ. 八姦')
def s_8gan_5_8(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 八姦', page, total)
    add_title(slide, '팔간 5~8 — 외부와 결탁의 함정')
    items = [
        ('5', '민맹(民萌)', '백성의 인기를 이용',
         '현대 — 포퓰리즘·세력 형성'),
        ('6', '유행(流行)', '유세객·변론가의 말솜씨',
         '현대 — PR 회사·SNS 인플루언서'),
        ('7', '위강(威強)', '무력·군사·기술 독점',
         '현대 — 군사 쿠데타·기술 독점'),
        ('8', '사방(四方)', '외세(다른 나라)와 결탁',
         '현대 — 외부 투자자·외국 정부 결탁'),
    ]
    top = 2.3
    for num, han, desc, modern in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(1.05), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.32), Inches(0.7), Inches(0.5),
                    num, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.5), Inches(1.05), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.32), Inches(2.5), Inches(0.5),
                    han, font_size=18, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.0), Inches(top), Inches(8.8), Inches(1.05),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.2), Inches(top + 0.12), Inches(8.5), Inches(0.4),
                    desc, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(4.2), Inches(top + 0.55), Inches(8.5), Inches(0.4),
                    modern, font_size=12, color=SUB)
        top += 1.18


@S('Ⅺ. 八姦')
def s_8gan_modern(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅺ. 八姦', page, total)
    add_title(slide, '현대 조직의 팔간 체크리스트')
    items = [
        '경영진 가족이 인사에 간섭하는가? (동상·부형)',
        '비서실·보좌진이 공식 라인을 우회하는가? (재방)',
        '오너 혈족이 능력 무관 승진하는가? (부형)',
        '접대·향응이 의사결정에 영향을 주는가? (양앙)',
        '직원 인기를 무기로 세력을 형성하는 임원이 있는가? (민맹)',
        '외부 컨설턴트가 편향된 정보로 판단을 흔드는가? (유행)',
        '특정 팀이 기술·고객을 독점하여 협상력을 행사하는가? (위강)',
        '외부 투자자와 결탁해 내부 의사결정을 흔드는 임원이 있는가? (사방)',
    ]
    top = 2.3
    for item in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.5), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(0.5), Inches(0.4),
                    '☐', font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.2), Inches(top), Inches(11.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(1.35), Inches(top + 0.12), Inches(11.3), Inches(0.4),
                    item, font_size=14, color=INK)
        top += 0.57


# ============== Ⅻ. 十過 ==============
@S('Ⅻ. 十過')
def s_10gua_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 十過', page, total)
    add_title(slide, '십과(十過) — 군주의 10가지 치명적 실수')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(2.45), Inches(11.3), Inches(1.8), [
        ('제10편 십과 — 군주가 범하기 쉬운 10가지 결함',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('"개별로는 사소해 보이지만 누적되면 망국(亡國)"',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('각각이 망국에 이르는 구체적 역사 사례로 입증된다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.6), Inches(4.5), Inches(11.9), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_textbox(slide, Inches(0.6), Inches(4.65), Inches(11.9), Inches(0.5),
                '10가지 한눈에', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    items = ['① 작은 충성', '② 소탐대실', '③ 행동 괴이', '④ 음악·유흥', '⑤ 탐욕',
             '⑥ 여색', '⑦ 마이크로매니지', '⑧ 결단력 부족', '⑨ 한 신하 의존', '⑩ 경고 무시']
    for i, item in enumerate(items):
        col = i % 5
        row = i // 5
        x = 0.7 + col * 2.4
        y = 5.3 + row * 0.55
        add_filled_rect(slide, Inches(x), Inches(y), Inches(2.3), Inches(0.45), PALE)
        add_textbox(slide, Inches(x), Inches(y + 0.08), Inches(2.3), Inches(0.4),
                    item, font_size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)


@S('Ⅻ. 十過')
def s_10gua_1_5(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 十過', page, total)
    add_title(slide, '십과 1~5')
    items = [
        ('1', '작은 충성 → 큰 충성 해침',  '개인적 호의 > 국가적 충성'),
        ('2', '소탐대실(小貪大失)',         '눈앞 이익으로 큰 이익 상실'),
        ('3', '행동 괴이·남 업신여김',     '기행 → 거짓 보고 증가'),
        ('4', '음악·유흥 탐닉',             '향락 → 정사 방치'),
        ('5', '탐욕(貪欲)',                 '끝없는 욕심 → 민란'),
    ]
    top = 2.3
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(5.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.22), Inches(4.9), Inches(0.5),
                    name, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.5), Inches(top), Inches(6.3), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.7), Inches(top + 0.22), Inches(6.0), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 0.95


@S('Ⅻ. 十過')
def s_10gua_6_10(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 十過', page, total)
    add_title(slide, '십과 6~10')
    items = [
        ('6',  '여색 탐닉',         '감정적 집착 → 판단 흐림'),
        ('7',  '직접 일함',         '군주가 궁을 비우면 위험 — 마이크로매니지'),
        ('8',  '결단력 부족',       '결정 미룸 → 기회 상실'),
        ('9',  '한 신하 의존',      '한 사람 말만 따름'),
        ('10', '무시할 수 없는 것 무시', '경고 신호 외면'),
    ]
    top = 2.3
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(5.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.22), Inches(4.9), Inches(0.5),
                    name, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.5), Inches(top), Inches(6.3), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.7), Inches(top + 0.22), Inches(6.0), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 0.95


@S('Ⅻ. 十過')
def s_10gua_history(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 十過', page, total)
    add_title(slide, '십과의 대표 역사 사례')
    cases = [
        ('제2과 소탐대실', '우(虞)나라와 순망치한(脣亡齒寒)',
         '진(晉)이 우에게 보물 주고 길 빌려 괵(虢)을 친 후, 돌아오는 길에 우도 멸망'),
        ('제5과 탐욕',     '진시황의 대규모 토목공사',
         '백성 혹사 → 진 15년 만에 멸망'),
        ('제6과 여색',     '서주 유왕과 포사 — 봉화희제후(烽火戲諸侯)',
         '거짓 봉화로 제후들을 속이다가 진짜 침공 시 제후가 오지 않아 멸망'),
        ('제9과 한 신하 의존', '제 환공과 역아·수조·개방',
         '관중 사후 세 간신에 의존 → 환공 사후 유폐되어 아사'),
        ('제10과 방심',    '오 부차와 월 구천',
         '항복 받고 경계 해제 → 와신상담한 구천에게 멸망'),
    ]
    top = 2.2
    for tag, name, desc in cases:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(9.6), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.35), Inches(top + 0.08), Inches(9.3), Inches(0.45),
                    name, font_size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(3.35), Inches(top + 0.5), Inches(9.3), Inches(0.4),
                    desc, font_size=12, color=INK)
        top += 0.95


# ============== XⅢ. 고분·세난·화씨 ==============
@S('XⅢ. 고분·세난·화씨')
def s_gobun(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 고분·세난·화씨', page, total)
    add_title(slide, '고분(孤憤) — 외로운 분노')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(11.9), Inches(0.55),
                '제11편 고분 — "법술지사(法術之士)가 기득권에 의해 죽는 구조적 비극"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.3), Inches(5.9), Inches(3.6), PALE)
    add_textbox(slide, Inches(0.6), Inches(3.45), Inches(5.9), Inches(0.5),
                '법술지사(法術之士)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(4.1), Inches(5.3), Inches(2.6), [
        ('법과 술을 아는 개혁가', {'font_size': 16, 'bold': True}),
        ('', {'font_size': 6}),
        ('• 국가를 위한 합리적 제안', {'font_size': 14, 'space_before': 6}),
        ('• 기득권의 사익을 침해', {'font_size': 14, 'space_before': 6}),
        ('• 보호 세력이 없음', {'font_size': 14, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(3.3), Inches(5.9), Inches(3.6),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(3.45), Inches(5.9), Inches(0.5),
                '당여지신(當與之臣)', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(4.1), Inches(5.3), Inches(2.6), [
        ('사적 파벌의 기득권', {'font_size': 16, 'bold': True}),
        ('', {'font_size': 6}),
        ('• 인사권 장악', {'font_size': 14, 'space_before': 6}),
        ('• 당파 형성·집단 비방', {'font_size': 14, 'space_before': 6}),
        ('• 공적 사유화', {'font_size': 14, 'space_before': 6}),
    ], line_spacing=1.3)


@S('XⅢ. 고분·세난·화씨')
def s_gobun_prophecy(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 고분·세난·화씨', page, total)
    add_title(slide, '고분의 자기 예언 — 한비 자신의 운명')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('개혁가의 비극적 운명',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('침묵하면 — 나라가 망하는 것을 지켜본다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('말하면 — 기득권의 미움으로 죽는다',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('역사 사례 — 그리고 자기 자신',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('상앙 — 거열형(車裂刑) · 오기 — 화살에 맞음',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('그리고 한비 자신 — 옥중 독살(BC 233)',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"한비가 『고분』에 예고한 구조가 자기 자신에게 그대로 실현되었다"',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('XⅢ. 고분·세난·화씨')
def s_senan(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 고분·세난·화씨', page, total)
    add_title(slide, '세난(說難) — 설득의 어려움')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(11.9), Inches(0.55),
                '제12편 세난 — 말더듬 한비의 설득 심리학',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.3), Inches(11.9), Inches(1.8), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(3.45), Inches(11.3), Inches(1.6), [
        ('설득이 어려운 진짜 이유',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('지식·언변·용기가 부족해서가 아니다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('상대방의 속마음을 정확히 파악하고 그에 맞추어 말해야 하기 때문',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.6), Inches(5.3), Inches(11.9), Inches(1.8),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.9), Inches(5.45), Inches(11.3), Inches(1.6), [
        ('입장(Position) vs 이해관계(Interest) — 하버드 협상 이론의 원조',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('겉으로 명예를 말하지만 속으로 이익을 탐할 때 — 양쪽 다 위험',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('이 역설을 푸는 것이 설득의 기술',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('XⅢ. 고분·세난·화씨')
def s_yeoklin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 고분·세난·화씨', page, total)
    add_title(slide, '역린(逆鱗) — 용의 거꾸로 난 비늘')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '逆 鱗',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.4),
                '— 제12편 세난', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(4.45), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.3), [
        ('"용(龍)은 길들이면 탈 수도 있다',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 그러나 목 아래에 한 자 길이의 거꾸로 난 비늘이 있는데',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        (' 이것을 건드리는 자는 반드시 죽임을 당한다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('모든 권력자에게는 절대 건드려서는 안 되는 금기가 있다',
         {'font_size': 15, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('CEO의 과거 실패·정치인의 스캔들·리더의 콤플렉스 — 모두 역린',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('XⅢ. 고분·세난·화씨')
def s_hwa_ssi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 고분·세난·화씨', page, total)
    add_title(slide, '화씨지벽(和氏之璧) — 인정받지 못하는 보옥')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(3.2), [
        ('변화(卞和)의 우화',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('변화가 옥 원석을 세 번 바침',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('여왕(厲王): "이것은 돌이다" → 왼발을 자름',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('무왕(武王): 또 "돌이다" → 오른발을 자름',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('문왕(文王): 쪼개보니 천하의 명옥(名玉) → "화씨지벽(和氏之璧)"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(6.05), Inches(11.7), Inches(1.0), [
        ('변화의 눈물 — "발이 잘린 것이 슬픈 것이 아니라',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 보옥을 돌이라 하고 정직한 선비를 거짓말쟁이라 부르는 것이 슬프다"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('XⅢ. 고분·세난·화씨')
def s_reformers(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅢ. 고분·세난·화씨', page, total)
    add_title(slide, '개혁가 열전 — 죽었으나 제도는 살았다')
    items = [
        ('오기(吳起)',  '초나라 개혁',         '도왕 사후 화살에 맞음 — 자기 몸을 왕의 시신 위에 던져 복수'),
        ('상앙(商鞅)',  '진나라 변법',         '효공 사후 거열형(車裂刑) — 그러나 그의 법은 진 통일의 기반'),
        ('한비(韓非)',  '법가 사상 완성',      '이사에게 옥중 독살 — 그러나 사상이 진 제국의 설계도가 됨'),
        ('공통점',      '개혁의 본질적 특성',  '기존 이익의 재분배 → 누군가의 이익을 빼앗음 → 저항'),
    ]
    top = 2.5
    for name, role, fate in items:
        is_summary = '공통점' in name
        bg_left = SUB if is_summary else ACCENT
        bg_right = PALE if is_summary else RGBColor(0xFA, 0xE5, 0xE5)
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.9), bg_left)
        add_textbox(slide, Inches(0.6), Inches(top + 0.25), Inches(2.5), Inches(0.5),
                    name, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(2.5), Inches(0.9), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.25), Inches(2.5), Inches(0.5),
                    role, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.8), Inches(top), Inches(7.0), Inches(0.9), bg_right)
        add_textbox(slide, Inches(6.0), Inches(top + 0.25), Inches(6.7), Inches(0.5),
                    fate, font_size=13, color=INK)
        top += 1.0
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                '"개혁가는 죽지만 제도는 남는다" — 화씨편의 핵심 통찰',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== XⅣ. 망징 ==============
@S('XⅣ. 망징')
def s_mangjing_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 망징', page, total)
    add_title(slide, '망징(亡徵) — 47가지 망국 징조',
              '제15편 — 현대 조직 진단에도 그대로 적용되는 체크리스트')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                '47',
                font_size=200, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
                '가지 망국 징조',
                font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.6), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.4), [
        ('세 가지 범주', {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('① 군주의 개인적 결함 (15개)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('② 제도적·구조적 문제 (18개)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('③ 사회적 병폐 (14개)',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('XⅣ. 망징')
def s_mangjing_examples(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 망징', page, total)
    add_title(slide, '대표 망징 — 13가지')
    items = [
        '사치가 지나치고 국고가 빈약하면 망한다',
        '법령이 조석으로 바뀌면 망한다 (정책 불확실성)',
        '공적 없이 높은 작위를 받는 자가 많으면 망한다 (피터의 원리)',
        '관리가 사사로이 재물을 모으면 망한다 (부패)',
        '같은 죄에 다른 벌을 내리면 망한다 (법의 불공정)',
        '약속을 지키지 않아 이웃 나라의 신뢰를 잃으면 망한다',
        '상인이 존경받고 농민이 천대받으면 망한다',
        '사사로이 무력을 쓰는 자가 떠받들어지면 망한다',
        '말재주만 갈고 실질 공적 없는 자가 존경받으면 망한다',
        '외국의 힘을 빌려 국내 문제를 해결하면 망한다',
        '후계 구도가 불안정하면 망한다',
        '측근이 정치에 간섭하면 망한다',
        '군주가 간언을 듣지 않으면 망한다',
    ]
    top = 1.95
    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = top + row * 0.7
        add_filled_rect(slide, Inches(x), Inches(y), Inches(0.5), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(x), Inches(y + 0.12), Inches(0.5), Inches(0.4),
                    str(i + 1), font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 0.6), Inches(y), Inches(5.7), Inches(0.55), PALE)
        add_textbox(slide, Inches(x + 0.7), Inches(y + 0.12), Inches(5.5), Inches(0.4),
                    item, font_size=11, color=INK)


@S('XⅣ. 망징')
def s_mangjing_modern(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅣ. 망징', page, total)
    add_title(slide, '현대 조직 망징 체크리스트')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.4), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.1), [
        ('"망징이 하나라도 있으면 위태롭고, 여럿이 겹치면 망한다',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 그러나 이를 아는 자는 드물고, 안다 해도 고치는 자는 더 드물다"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    items = [
        '리더가 세부에만 간섭하는가?',
        '인사가 실력 아닌 관계에 의존하는가?',
        '정책이 너무 자주 바뀌는가?',
        '성과와 보상이 연동되지 않는가?',
        '비공식 권력이 공식 구조를 우회하는가?',
        '법 적용이 사람에 따라 달라지는가?',
        '기득권층이 의무를 회피하는가?',
        '내부 고발이 보호받지 않는가?',
    ]
    top = 3.85
    for item in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.5), Inches(0.4), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.05), Inches(0.5), Inches(0.3),
                    '☐', font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.15), Inches(top), Inches(11.65), Inches(0.4),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(1.3), Inches(top + 0.07), Inches(11.4), Inches(0.3),
                    item, font_size=13, color=INK)
        top += 0.42


# ============== XⅤ. 비내·해로유로 ==============
@S('XⅤ. 비내·해로유로')
def s_binae(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 비내·해로유로', page, total)
    add_title(slide, '비내(備內) — 내부 경계')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('제17편 비내 — "군주의 진짜 적은 외적이 아니라 내부"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"누가 군주의 죽음에서 이익을 얻는가?"',
         {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('• 태자 — 빨리 왕이 되고 싶다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('• 후궁 — 총애가 사라지면 모든 것을 잃는다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
        ('• 권신 — 새 군주를 세우면 공신이 된다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('제도로 막아라 — 4대 원칙',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('① 견제와 균형 ② 정보 다원화 ③ 법에 의한 통제 ④ 이해관계 재설계',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('현대 — 엔론·리먼브라더스 모두 내부 문제로 붕괴',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('XⅤ. 비내·해로유로')
def s_haero(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 비내·해로유로', page, total)
    add_title(slide, '해로(解老) — 노자의 법가적 재해석',
              '제20편 — 도가와 법가의 융합 텍스트')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(5.9), Inches(0.5),
                '도(道)의 재해석', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.05), Inches(5.3), Inches(3.6), [
        ('노자', {'bold': True, 'font_size': 16, 'color': SUB}),
        ('  우주의 근원 — 말할 수 없는 궁극',
         {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 4}),
        ('한비', {'bold': True, 'font_size': 16, 'color': ACCENT, 'space_before': 6}),
        ('  만물이 그렇게 되는 까닭', {'font_size': 14}),
        ('  — 이치(理) 그 자체', {'font_size': 14, 'bold': True}),
        ('', {'font_size': 4}),
        ('한비는 도를', {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ('"신비"에서 "합리"로 가져왔다', {'font_size': 13, 'color': SUB, 'bold': True}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.45), Inches(5.9), Inches(0.5),
                '무위(無爲)의 변환', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.05), Inches(5.3), Inches(3.6), [
        ('노자', {'bold': True, 'font_size': 16, 'color': SUB}),
        ('  자연에 맡기고 인위를 삼가라',
         {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 4}),
        ('한비', {'bold': True, 'font_size': 16, 'color': ACCENT, 'space_before': 6}),
        ('  법을 잘 만들어놓으면', {'font_size': 14}),
        ('  군주가 간섭하지 않아도', {'font_size': 14}),
        ('  나라가 저절로 돌아간다', {'font_size': 14, 'bold': True}),
        ('', {'font_size': 4}),
        ('현대 — 잘 설계된 시스템은', {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ('관리자 개입 없이 자동 작동', {'font_size': 13, 'color': SUB, 'bold': True}),
    ], line_spacing=1.3)


@S('XⅤ. 비내·해로유로')
def s_pyeonjak(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 비내·해로유로', page, total)
    add_title(slide, '편작(扁鵲)과 채환공(蔡桓公) — 조기 진단의 우화')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.2), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.9), [
        ('명의 편작이 채환공을 세 번 만나 경고:',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('1차 — "피부에 병이 있습니다" → 무시',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('2차 — "살 속에 병이 있습니다" → 무시',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('3차 — "장부(臟腑)에 병이 있습니다" → 무시',
         {'font_size': 16, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('4차 — 환공을 보자마자 편작은 그냥 돌아감',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('  "병이 골수까지 갔으니 이제 어찌할 수 없다"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('환공은 곧 죽었다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)
    add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.6),
                '"작은 경고를 무시하면 큰 재앙 — 조기 개입의 중요성"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '— 제21편 유로(喩老)',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S('XⅤ. 비내·해로유로')
def s_yuro_others(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅤ. 비내·해로유로', page, total)
    add_title(slide, '유로(喩老)의 다른 유명 우화')
    items = [
        ('천 길 제방과 개미구멍', '千里之堤 潰於蟻穴',
         '천 길 제방도 개미구멍 때문에 무너진다 — 나비효과·조기 경보'),
        ('천리마의 뼈',           '死馬骨 五百金',
         '죽은 말의 뼈도 오백 금에 산다 → 살아 있는 천리마가 몰려든다'),
        ('범려의 공성신퇴',       '功成身退',
         '범려는 월왕 구천을 도와 오를 멸하고 즉시 물러나 화를 면함'),
        ('화·복의 상호전환',     '禍福相依',
         '"화에는 복이 기대어 있고, 복에는 화가 숨어 있다"'),
    ]
    top = 2.4
    for name, han, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.8), Inches(0.95), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(2.8), Inches(0.5),
                    name, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.5), Inches(top), Inches(2.5), Inches(0.95), PALE)
        add_textbox(slide, Inches(3.5), Inches(top + 0.27), Inches(2.5), Inches(0.5),
                    han, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.1), Inches(top), Inches(6.7), Inches(0.95),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.3), Inches(top + 0.27), Inches(6.4), Inches(0.5),
                    desc, font_size=13, color=INK)
        top += 1.1


# ============== XⅥ. 시대비판 ==============
@S('XⅥ. 시대비판')
def s_5_worms(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 시대비판', page, total)
    add_title(slide, '오두(五蠹) — 다섯 좀벌레', '제49편 — 나라를 좀먹는 5가지')
    items = [
        ('1', '학자(學者)',         '유가 — 비현실적 고대 숭배, 법 약화',     '이상주의 경영자'),
        ('2', '언담자(言談者)',     '종횡가 — 무책임한 말꾼, 실행력 부재',   '실적 없는 경영 구루'),
        ('3', '대검자(帶劍者)',     '협객 — 사적 폭력, 법질서 파괴',         '온라인 자경단'),
        ('4', '환어자(患御者)',     '권문 기생자 — 의무 회피, 특권 악용',    '병역 면제·탈세 특권층'),
        ('5', '상공지민(商工之民)', '비생산적 부, 사치 조장',                '실물 없는 투기'),
    ]
    top = 2.4
    for num, name, desc, modern in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.2), Inches(0.7), Inches(0.4),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.6), Inches(0.8), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.2), Inches(2.6), Inches(0.4),
                    name, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.1), Inches(top), Inches(5.4), Inches(0.8),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(4.3), Inches(top + 0.2), Inches(5.2), Inches(0.4),
                    desc, font_size=13, color=INK)
        add_filled_rect(slide, Inches(9.6), Inches(top), Inches(3.2), Inches(0.8), PALE)
        add_textbox(slide, Inches(9.7), Inches(top + 0.2), Inches(3.0), Inches(0.4),
                    modern, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        top += 0.92


@S('XⅥ. 시대비판')
def s_sujudaeto(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 시대비판', page, total)
    add_title(slide, '수주대토(守株待兔) — 그루터기 지키며 토끼 기다리기')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '守 株 待 兔',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.4),
                '— 제49편 오두', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.35), Inches(11.7), Inches(1.7), [
        ('송나라 농부가 밭을 갈고 있는데 토끼가 달려와 그루터기에 부딪혀 죽었다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('농부는 이후 쟁기를 버리고 그루터기 옆에 앉아 토끼를 기다렸다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('토끼는 다시 오지 않았고, 농부는 송나라의 웃음거리가 됐다',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                '"지금 옛 선왕의 방법으로 당세를 다스리려는 자는, 모두 이 그루터기를 지키는 자다"',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('XⅥ. 시대비판')
def s_hyeonhak(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 시대비판', page, total)
    add_title(slide, '현학(顯學) — 유가와 묵가 비판')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(11.9), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(2.45), Inches(11.9), Inches(0.55),
                '"공자는 한 사람인데 가르침이 여덟 가지, 묵자는 한 사람인데 셋. 누가 진짜인가?"',
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.6), Inches(3.3), Inches(5.9), Inches(3.6), PALE)
    add_textbox(slide, Inches(0.6), Inches(3.45), Inches(5.9), Inches(0.5),
                '유가의 분열 — 8파', font_size=17, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(4.05), Inches(5.3), Inches(2.5),
                '자장파 · 자사파\n안씨파 · 맹씨파\n칠조씨파 · 중량씨파\n손씨파 · 악정씨파',
                font_size=15, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(6.8), Inches(3.3), Inches(5.9), Inches(3.6), PALE)
    add_textbox(slide, Inches(6.8), Inches(3.45), Inches(5.9), Inches(0.5),
                '묵가의 분열 — 3파', font_size=17, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.1), Inches(4.5), Inches(5.3), Inches(1.5),
                '상리씨\n상부씨\n등릉씨',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)


@S('XⅥ. 시대비판')
def s_unverifiable(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅥ. 시대비판', page, total)
    add_title(slide, '검증 불가능한 권위 인용의 문제')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"요순의 도를 따르라"고 하지만',
         {'font_size': 17, 'align': PP_ALIGN.CENTER}),
        ('요순은 죽은 지 수백 년 — 진위 확인 불가',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"죽은 사람의 뜻을 빌려 자신의 주장을 정당화"',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('= 지적 사기(知的 詐欺)',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('법치 vs 도덕 교화',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('도덕: 모호("인이란?")·주관적·검증 불가·느림',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('법: 명확(명문화)·보편적·검증 가능·즉시',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('현대 — "열정적으로 일하라"보다 명확한 KPI·성과급이 효과적',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


# ============== XⅦ. 유명 우화 10편 ==============
def make_fable_slide(section, idx, title_kor, title_han, eum, story, lesson, ref):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_title(slide, f'우화 {idx} — {title_kor}({title_han})',
                  f'{eum} · {ref}')
        add_filled_rect(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.7), PALE)
        add_textbox(slide, Inches(0.6), Inches(2.55), Inches(12.1), Inches(0.5),
                    '이야기', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(3.1), Inches(11.7), Inches(2.0),
                    story, font_size=15, color=INK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(0.6), Inches(5.45), Inches(12.1), Inches(0.5),
                    '교훈', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(5.95), Inches(11.7), Inches(1.0),
                    lesson, font_size=16, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return renderer


SLIDES.append((make_fable_slide('XⅦ. 우화 (1/10)', '①',
    '수주대토', '守株待兔', '수주대토',
    '송나라 농부가 그루터기에 부딪혀 죽은 토끼를 얻은 뒤\n쟁기를 버리고 그루터기 옆에 앉아 토끼를 기다림.\n토끼는 다시 오지 않았고, 농부는 송나라의 웃음거리가 됨.',
    '과거의 성공에 안주하지 마라 — 시대는 변한다',
    '제49편 오두'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (2/10)', '②',
    '모순', '矛盾', '모순',
    '초나라 상인이 "무엇이든 뚫는 창"과 "무엇이든 막는 방패"를\n동시에 팔자, 한 사람이 물었다.\n"그 창으로 그 방패를 찌르면 어떻게 되느냐?" — 상인은 답하지 못했다.',
    '논리적 일관성의 중요성 — 두 극단을 동시에 주장할 수 없다',
    '제36편 난일, 제40편 난세'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (3/10)', '③',
    '화씨지벽', '和氏之璧', '화씨지벽',
    '변화(卞和)가 옥 원석을 세 번 바쳤다.\n두 왕은 "돌이다" 하여 두 발을 자르고\n세 번째 왕이 쪼개보니 천하의 명옥(名玉) — 화씨지벽.',
    '진정한 가치는 알아보는 눈이 있어야 빛난다 — 개혁가의 비극',
    '제13편 화씨'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (4/10)', '④',
    '정인매리', '鄭人買履', '정인매리',
    '정나라 사람이 자기 발의 치수를 재놓고 시장에 갔다.\n치수를 집에 두고 온 것을 알고 집으로 돌아가 가져왔으나 시장이 파함.\n"직접 신어보면 될 것을" 하니 "치수는 믿어도 내 발은 못 믿는다"고.',
    '경직된 원칙주의의 어리석음 — 수단과 목적을 혼동하지 마라',
    '제32편 외저설 좌상'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (5/10)', '⑤',
    '악양식자', '樂羊食子', '악양식자',
    '위나라 장군 악양이 중산국을 칠 때, 중산왕이 악양의 아들을\n잡아 국을 끓여 보냄. 악양은 "군법을 어길 수 없다"며 그 국을 다 먹음.\n위 문후는 "자기 아들 고기를 먹은 자가 남의 고기도 먹겠다"며 경계.',
    '극단적 충성은 극단적 배신의 가능성이다 — 인간 본성의 양면',
    '제22편 세림 상'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (6/10)', '⑥',
    '오기연저', '吳起吮疽', '오기연저',
    '장군 오기가 병사의 종기를 입으로 빨아주자 그 병사의 어머니가 통곡.\n"전에 오 장군이 그 아버지 종기를 빨아주었는데\n아버지는 감격하여 결국 싸움에서 죽었다. 내 아들도 죽을 것이다"',
    '리더의 친밀한 행동에도 전략적 의도가 있을 수 있다 — 사랑도 무기',
    '제32편 외저설 좌상'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (7/10)', '⑦',
    '편작견환공', '扁鵲見桓公', '편작견환공',
    '편작이 환공의 병을 세 번 경고했으나 무시당함.\n피부 → 살 속 → 장부 → 골수에 이르자 편작은 도망.\n환공은 결국 죽음.',
    '작은 경고를 무시하면 큰 재앙 — 조기 개입의 중요성',
    '제21편 유로'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (8/10)', '⑧',
    '천리마지골', '千里馬之骨', '천리마지골',
    '연 소왕이 인재를 구하자 곽외가 조언:\n"죽은 천리마의 뼈도 5백 금에 산다는 말처럼\n먼저 저 같은 평범한 자를 후히 대우하소서" → 악의·추연 같은 인재가 몰려듦.',
    '진정성 있는 대우가 최고의 인재를 끌어온다',
    '제21편 유로'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (9/10)', '⑨',
    '삼인성호', '三人成虎', '삼인성호',
    '한 사람이 "시장에 호랑이가 있다"고 하면 믿지 않음.\n두 사람이 말하면 의심. 세 사람이 말하면 믿는다.',
    '거짓 정보도 반복되면 진실이 된다 — 한 경로 정보에 의존 말 것',
    '제30편 내저설 상'), 'XⅦ. 우화'))

SLIDES.append((make_fable_slide('XⅦ. 우화 (10/10)', '⑩',
    '초왕호세요', '楚王好細腰', '초왕호세요',
    '초 영왕이 가는 허리의 신하를 좋아하자\n조정의 모든 신하들이 굶어서 허리를 가늘게 만들어\n얼굴이 검게 되고 기둥을 붙잡고 걸어다녔다.',
    '위가 좋아하면 아래가 반드시 따른다 — 군주의 선호는 조직 전체를 왜곡',
    '제32편 외저설 좌상'), 'XⅦ. 우화'))


# ============== XⅧ. 10대 원리 ==============
@S('XⅧ. 10대 원리')
def s_10principles_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 10대 원리', page, total)
    add_title(slide, '한비자를 관통하는 10대 원리 ①')
    items = [
        ('1', '시대는 변한다',         '世異則事異 — 과거에 안주하지 말고 현재의 조건에 맞춰라'),
        ('2', '인간은 이익으로 움직인다','凡治天下 必因人情 — 악이 아니라 현실, 이익을 활용해 통제'),
        ('3', '제도가 개인을 초월한다', '法治 > 人治 — 좋은 시스템이 좋은 사람을 이긴다'),
        ('4', '상벌은 두 자루(二柄)',   '상과 벌은 군주만이 쥔다, 절대 나누지 마라'),
        ('5', '형명참동(循名責實)',     '말(名)과 실적(形)을 대조하여 평가 — 현대 KPI의 원조'),
    ]
    top = 2.4
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(4.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.22), Inches(4.4), Inches(0.5),
                    name, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.0), Inches(top), Inches(6.8), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.2), Inches(top + 0.22), Inches(6.5), Inches(0.5),
                    desc, font_size=13, color=INK)
        top += 0.95


@S('XⅧ. 10대 원리')
def s_10principles_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XⅧ. 10대 원리', page, total)
    add_title(slide, '한비자를 관통하는 10대 원리 ②')
    items = [
        ('6',  '뜻을 숨겨라',          '虛靜無爲 — 속마음을 드러내면 신하가 맞춰 위장한다'),
        ('7',  '법·술·세의 통합',       '三位一體 — 세 가지가 모두 있어야 완전한 통치'),
        ('8',  '내부의 적을 경계하라',  '備內 — 외적보다 위험한 것은 이해관계 있는 측근'),
        ('9',  '역린을 피하라',          '逆鱗 — 권력자의 심리적 지뢰를 알고 접근하라'),
        ('10', '무시할 수 없는 것 무시 말 것', '小過 積爲 大過 — 개미구멍으로 제방이 무너진다'),
    ]
    top = 2.4
    for num, name, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(4.5), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.5), Inches(top + 0.22), Inches(4.4), Inches(0.5),
                    name, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.0), Inches(top), Inches(6.8), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.2), Inches(top + 0.22), Inches(6.5), Inches(0.5),
                    desc, font_size=13, color=INK)
        top += 0.95


# ============== XⅨ. 명구절 ==============
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=42):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=22, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('XⅨ. 명구절 (1/10)',
    '凡 治 天 下   必 因 人 情',
    '범치천하 필인인정',
    '천하를 다스리려면 반드시 사람의 실정(人情)에 근거하라',
    '제48편 팔경', hanmun_size=42), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (2/10)',
    '國 無 常 強   奉 法 者 強   則 國 強',
    '국무상강 봉법자강 즉국강',
    '나라의 강약은 일정하지 않다 — 법을 받드는 자가 강하면 나라가 강하다',
    '제6편 유도', hanmun_size=32), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (3/10)',
    '刑 過 不 避 大 臣\n賞 善 不 遺 匹 夫',
    '형과불피대신 · 상선불유필부',
    '형벌은 대신도 피하지 못하고, 상은 평민도 놓치지 않는다',
    '제6편 유도', hanmun_size=30), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (4/10)',
    '二 柄 者   刑 德 也',
    '이병자 형덕야',
    '두 자루는 형(刑)과 덕(德)이다 — 군주만이 쥐는 절대 권능',
    '제7편 이병', hanmun_size=52), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (5/10)',
    '循 名 而 責 實',
    '순명책실',
    '이름(名)에 따라 실적(實)을 묻는다 — 형명참동의 원리',
    '제43편 정법', hanmun_size=72), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (6/10)',
    '人 主 之 患 在 於 信 人\n信 人 則 制 於 人',
    '인주지환재어신인 · 신인즉제어인',
    '군주의 걱정은 사람을 믿는 데 있다 — 믿으면 사람에게 제어당한다',
    '제17편 비내', hanmun_size=26), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (7/10)',
    '輿 人 欲 人 之 富 貴\n匠 人 欲 人 之 夭 死',
    '여인욕인지부귀 · 장인욕인지요사',
    '수레 장인은 사람들이 부자이기를, 관 장인은 일찍 죽기를 바란다',
    '제17편 비내', hanmun_size=26), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (8/10)',
    '其 喉 下 有 逆 鱗',
    '기후하유역린',
    '용의 목 아래에 거꾸로 난 비늘이 있다 — 권력자의 절대 금기',
    '제12편 세난', hanmun_size=52), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (9/10)',
    '守 株 待 兔',
    '수 주 대 토',
    '그루터기를 지키며 토끼를 기다린다 — 시대 변화를 모르는 어리석음',
    '제49편 오두', hanmun_size=104), 'XⅨ. 명구절'))

SLIDES.append((make_quote_slide('XⅨ. 명구절 (10/10)',
    '千 里 之 堤   潰 於 蟻 穴',
    '천리지제 궤어의혈',
    '천 리 제방도 개미구멍으로 무너진다 — 작은 신호가 큰 재앙으로',
    '제21편 유로', hanmun_size=42), 'XⅨ. 명구절'))


# ============== XX. 비교 ==============
@S('XX. 비교')
def s_vs_confucian(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XX. 비교', page, total)
    add_title(slide, '한비자 vs 유가 — 정면 대립')
    rows = [
        ('인간 본성',  '성선(공자·맹자) / 성악(순자)',  '이기(利己) — 선악 이전'),
        ('통치 방법',  '덕치(德治)·예치(禮治)',          '법치(法治)'),
        ('이상',       '요순의 도',                       '당세(當世)의 현실'),
        ('지식인',     '군자(君子)',                      '"좀벌레(蠹)"'),
        ('시간관',     '복고 — 옛날이 좋았다',            '역사 진화론'),
    ]
    top = 2.2
    add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.5), Inches(0.55), SUB)
    add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(2.5), Inches(0.4),
                '쟁점', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(3.05), Inches(top), Inches(4.8), Inches(0.55), SUB)
    add_textbox(slide, Inches(3.05), Inches(top + 0.1), Inches(4.8), Inches(0.4),
                '유가', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(7.9), Inches(top), Inches(4.9), Inches(0.55), ACCENT)
    add_textbox(slide, Inches(7.9), Inches(top + 0.1), Inches(4.9), Inches(0.4),
                '한비자', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    row_h = 0.8
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(y), Inches(2.5), Inches(row_h), PALE)
        add_textbox(slide, Inches(0.55), Inches(y + 0.22), Inches(2.4), Inches(0.4),
                    row[0], font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.05), Inches(y), Inches(4.8), Inches(row_h), bg)
        add_textbox(slide, Inches(3.2), Inches(y + 0.22), Inches(4.5), Inches(0.4),
                    row[1], font_size=14, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.9), Inches(y), Inches(4.9), Inches(row_h), bg)
        add_textbox(slide, Inches(8.05), Inches(y + 0.22), Inches(4.6), Inches(0.4),
                    row[2], font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '한비의 유가 비판: "요순이 있었다는 것도 증명할 수 없는데, 따라야 한다고?"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('XX. 비교')
def s_vs_taoism_mohism(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XX. 비교', page, total)
    add_title(slide, '한비자 vs 도가·묵가')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                'vs 도가(노자·장자)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('대립 지점', {'font_size': 15, 'bold': True, 'color': ACCENT}),
        ('도가: 무위자연(無爲自然)', {'font_size': 14}),
        ('법가: 제도·법치', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('공유 지점', {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('무위이치(無爲而治) — 도가의 "무위"를', {'font_size': 14}),
        ('한비는 "법을 세워놓고 간섭 안 함"으로 재해석', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('→ 한비는 노자를 법가의 선배로 재해석', {'font_size': 13, 'color': SUB, 'bold': True, 'space_before': 6}),
        ('  (제20·21편 해로·유로)', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                'vs 묵가(묵자)', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('묵가의 주장', {'font_size': 15, 'bold': True, 'color': ACCENT}),
        ('겸애(兼愛)·비공(非攻)', {'font_size': 14}),
        ('모든 사람을 평등하게 사랑하라', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('한비의 비판', {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('"가족조차 이해타산으로 움직이는데', {'font_size': 14}),
        (' 타인을 어찌 평등하게 사랑하겠는가?"', {'font_size': 14, 'bold': True}),
        ('', {'font_size': 6}),
        ('현학편에서 묵가의 내부 분열도 비판', {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ('  (3파로 분열)', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)


@S('XX. 비교')
def s_vs_machiavelli(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XX. 비교', page, total)
    add_title(slide, '한비자 vs 마키아벨리 군주론 — 1,700년의 시차')
    rows = [
        ('인간관',    '이익 추구',       '이기적·배은망덕'),
        ('도덕',      '제도로 대체',     '효용으로 판단'),
        ('통치',      '법·술·세',         '역량(virtù)·운(fortuna)'),
        ('배신 대응', '이해관계 분석',   '"두려워지는 것이 낫다"'),
        ('시대관',    '진화론',           '순환론'),
    ]
    top = 2.2
    add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.55), SUB)
    add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(2.3), Inches(0.4),
                '비교', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(2.85), Inches(top), Inches(5.0), Inches(0.55), ACCENT)
    add_textbox(slide, Inches(2.85), Inches(top + 0.1), Inches(5.0), Inches(0.4),
                '한비자 (BC 233경)', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(7.9), Inches(top), Inches(4.9), Inches(0.55), SUB)
    add_textbox(slide, Inches(7.9), Inches(top + 0.1), Inches(4.9), Inches(0.4),
                '마키아벨리 (AD 1513)', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    row_h = 0.7
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(y), Inches(2.3), Inches(row_h), PALE)
        add_textbox(slide, Inches(0.55), Inches(y + 0.18), Inches(2.2), Inches(0.4),
                    row[0], font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.85), Inches(y), Inches(5.0), Inches(row_h), bg)
        add_textbox(slide, Inches(3.0), Inches(y + 0.18), Inches(4.7), Inches(0.4),
                    row[1], font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.9), Inches(y), Inches(4.9), Inches(row_h), bg)
        add_textbox(slide, Inches(8.05), Inches(y + 0.18), Inches(4.6), Inches(0.4),
                    row[2], font_size=14, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                '동서양 현실주의 정치철학의 양대 산맥 — 한비가 1,700년 앞선 선구자',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== XXI. 현대 의의 ==============
@S('XXI. 현대 의의')
def s_modern_biz(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXI. 현대 의의', page, total)
    add_title(slide, '경영학 — 제도 설계의 원조')
    rows = [
        ('형명참동(形名參同)',  'KPI·OKR·MBO — 약속과 실적의 대조'),
        ('이병(二柄)',           '인사권·평가권의 경영진 집중'),
        ('법·술·세',             '공식 시스템 + 비공식 리더십 + 권한 구조'),
        ('팔간·십과',            '코퍼레이트 리스크 매트릭스'),
        ('망징 47',              '조직 진단 체크리스트'),
        ('비내(備內)',           '내부통제·감사위원회·내부고발자 보호'),
        ('주도(主道)',           '베조스의 "가장 나중에 발언" — 관찰자 효과'),
    ]
    top = 2.2
    for tag, modern in rows:
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(4.5), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.5), Inches(top + 0.15), Inches(4.5), Inches(0.4),
                    tag, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.1), Inches(top), Inches(7.7), Inches(0.6), PALE)
        add_textbox(slide, Inches(5.3), Inches(top + 0.15), Inches(7.4), Inches(0.4),
                    modern, font_size=14, color=INK)
        top += 0.65


@S('XXI. 현대 의의')
def s_modern_econ(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXI. 현대 의의', page, total)
    add_title(slide, '행동경제학·제도경제학과의 접점')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '행동경제학', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('한비의 명제', {'font_size': 15, 'bold': True, 'color': ACCENT}),
        ('인간은 이익으로 움직인다', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('행동경제학', {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('인간 = 합리적 이기심 +', {'font_size': 14}),
        ('편향·휴리스틱', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('리처드 탈러의 넛지(Nudge)', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('= 인센티브 설계로 행동 유도', {'font_size': 13, 'color': SUB}),
        ('= 한비의 "이익을 활용한 통제"', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '제도경제학', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('더글러스 노스 (Nobel 1993)', {'font_size': 14, 'bold': True}),
        ('"제도가 경제 성과를 결정한다"', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('한비의 법치(法治)', {'font_size': 15, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('= 노스의 제도(Institutions)', {'font_size': 14}),
        ('', {'font_size': 6}),
        ('"좋은 제도가 좋은 사람을 이긴다"', {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 6}),
        ('한비가 2,300년 전 선언', {'font_size': 13, 'color': SUB}),
    ], line_spacing=1.3)


@S('XXI. 현대 의의')
def s_zero_trust_modern(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXI. 현대 의의', page, total)
    add_title(slide, '제로 트러스트 보안 — 정보보안의 2,300년 전 원조')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.0),
                '人 主 之 患   在 於 信 人',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '인주지환 재어신인 — "군주의 걱정은 사람을 믿는 데 있다"',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.1), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.55), Inches(11.7), Inches(2.3), [
        ('Zero Trust — "Never trust, always verify"',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('Google·Microsoft가 채택한 21세기 정보보안 패러다임',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('= 비내(備內)의 정확한 현대적 구현',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('내부 사용자도, 신뢰받는 시스템도 매번 검증한다',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"가장 가까운 자조차 의심하라"는 한비의 정신과 동일',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('XXI. 현대 의의')
def s_modern_gov(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXI. 현대 의의', page, total)
    add_title(slide, '현대 거버넌스 — 견제와 균형의 원조')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"한비는 군주제 옹호자인데, 현대 민주주의와 어떻게 연결되는가?"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('한비의 핵심 — 권력자도 법의 지배를 받아야 한다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('권력 집중은 위험하다 → 견제 시스템 필요',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('현대 거버넌스 원리와의 일치',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 팔간·비내 → 삼권분립의 동양적 원조',
         {'font_size': 14, 'space_before': 6}),
        ('• 형과불피대신 → 법 앞의 평등',
         {'font_size': 14, 'space_before': 4}),
        ('• 망징 체크리스트 → 기업 거버넌스 진단',
         {'font_size': 14, 'space_before': 4}),
        ('• 형명참동 → 360도 평가·실적 검증',
         {'font_size': 14, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== XXII. 마무리 ==============
@S('XXII. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXII. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 한비자')
    add_filled_rect(slide, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.9), PALE)
    add_paragraphs(slide, Inches(1.1), Inches(2.3), Inches(11.1), Inches(4.6), [
        ('"인간은 이익으로 움직인다"는 냉정한 관찰에서 출발하여',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('성인을 기다리는 도덕 정치 대신',
         {'font_size': 17, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('평범한 군주도 다스릴 수 있는 제도 시스템을 설계한 책',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('법(法, 공개된 규칙) · 술(術, 은밀한 기법) · 세(勢, 지위의 권력)',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('의 삼각형이 그 뼈대',
         {'font_size': 16, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('이병(상·벌)을 군주가 독점하고',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('형명참동(약속과 실적 대조)으로 평가하며',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('허정무위(뜻을 숨김)로 조직을 왜곡 없이 관찰한다',
         {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('— 동양 정치철학 현실주의의 정점 —',
         {'font_size': 18, 'bold': True, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
    ], line_spacing=1.2)


@S('XXII. 마무리')
def s_prison_legacy(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'XXII. 마무리', page, total)
    add_title(slide, '한비의 옥중 유언 — 비극과 승리의 동시')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.3),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0), [
        ('BC 233년 진(秦) 옥중 — 47세경, 이사의 모함으로 독약 마심',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"수레 장인과 관 장인이 서로 다른 것을 바라는 이 세상에서',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 도덕 교화로는 질서가 서지 않고',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 성인을 기다리면 나라가 망하니',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 이익을 인정하고 법으로 다스려야 한다"',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.95), Inches(11.7), Inches(2.0), [
        ('한비는 죽었으나 — 12년 후 진시황이 6국 통일',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('그가 설계한 법가 시스템은',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('중국 2,000년 제국의 바탕이 되었다',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"외유내법(外儒內法)" — 겉은 유교, 속은 법가',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.25)


@S('XXII. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.5),
                '守 株 待 兔',
                font_size=130, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
                '수 주 대 토', font_size=24, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.7),
                '"그루터기를 지키며 토끼를 기다리지 마라"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '시대는 변한다 — 변화를 읽는 자가 이긴다',
                font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
                '— 제49편 오두 · 한비가 2,300년 전 남긴 마지막 메시지',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
                '감사합니다', font_size=26, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\한비자.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')