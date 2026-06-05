# -*- coding: utf-8 -*-
"""
중용 발표자료 — 망라적 77장 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 1. 표지 ==============
@S('표지')
def s_cover(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8),
                '中 庸', font_size=130, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                'The Doctrine of the Mean · 중용',
                font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.5), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.75), Inches(12.3), Inches(0.5),
                '자사(子思) 저 — 사서(四書)의 정점, 동아시아 제왕학의 심장',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
                '본래 『예기(禮記)』 제31편 · 33장 · 약 3,500자 · 공문전수지심법(孔門傳授之心法)',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '"誠者, 天之道也; 誠之者, 人之道也"',
                font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                '— 성(誠)은 하늘의 도이고, 성실하고자 함은 사람의 도다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 2. 목차 ==============
@S('목차')
def s_toc(slide, page, total):
    set_white_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.8), Inches(0.7),
                '목 차', font_size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.15), Inches(12.8))
    items = [
        [('Ⅰ', '개요 — 중용이란 무엇인가'),
         ('Ⅱ', '왜 왕의 책이었나 — 제왕학의 심장'),
         ('Ⅲ', '구조 — 33장의 3부 구성'),
         ('Ⅳ', '1장 (총론) — 중용의 형이상학적 뼈대'),
         ('Ⅴ', '시중(時中)과 군자 — 2~11장'),
         ('Ⅵ', '도의 보편성 — 12~19장'),
         ('Ⅶ', '20장 — 중용의 중심, 구경(九經)'),
         ('Ⅷ', '성(誠)의 형이상학 — 21~26장')],
        [('Ⅸ', '성인의 도 — 27~33장'),
         ('Ⅹ', '7대 주제어'),
         ('Ⅺ', '명구절 10선'),
         ('Ⅻ', '7대 핵심 메시지'),
         ('ⅩⅢ', '현대적 의의'),
         ('ⅩⅣ', '다른 고전과의 비교'),
         ('ⅩⅤ', '마무리')],
    ]
    for col, group in enumerate(items):
        x = 0.7 + col * 6.4
        top = 1.5
        for num, title in group:
            add_textbox(slide, Inches(x), Inches(top), Inches(1.0), Inches(0.4),
                        num, font_size=15, bold=True, color=ACCENT)
            add_textbox(slide, Inches(x + 1.0), Inches(top), Inches(5.3), Inches(0.4),
                        title, font_size=15, color=INK)
            top += 0.55


# ============== Ⅰ. 개요 ==============
@S('Ⅰ. 개요')
def s_what_is(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '중용(中庸)이란 무엇인가')
    add_filled_rect(slide, Inches(0.5), Inches(2.2), Inches(5.9), Inches(4.7), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(5.9), Inches(0.5),
                '中 — 중(中)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(3.0), Inches(5.3), Inches(3.7), [
        ('한쪽으로 치우치지 않음', {'font_size': 16, 'bold': True}),
        ('', {'font_size': 6}),
        ('지나침(過)과 모자람(不及)의 중간',
         {'font_size': 14, 'space_before': 8}),
        ('', {'font_size': 6}),
        ('단순한 산술적 중간이 아니라',
         {'font_size': 14, 'color': SUB, 'space_before': 8}),
        ('"그 상황에 가장 알맞은 지점"',
         {'font_size': 16, 'bold': True, 'color': ACCENT}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.2), Inches(5.9), Inches(4.7),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.4), Inches(5.9), Inches(0.5),
                '庸 — 용(庸)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(3.0), Inches(5.3), Inches(3.7), [
        ('두 뜻이 공존', {'font_size': 16, 'bold': True}),
        ('', {'font_size': 6}),
        ('① 평상(平常)', {'font_size': 14, 'space_before': 8, 'bold': True}),
        ('  일상에서 변하지 않음', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 4}),
        ('② 용구(用久)', {'font_size': 14, 'space_before': 6, 'bold': True}),
        ('  쓰임이 오래감', {'font_size': 13, 'color': SUB}),
        ('', {'font_size': 6}),
        ('"평소에도 변함없이 지킬 수 있는 도"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 10}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '"적당히 절충"이 아니라 "가장 알맞은 한 점을 정확히 찾아 거기 머무는 것"',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅰ. 개요')
def s_status(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '책의 위상 — 핵심 정보')
    rows = [
        ('서명',    '중용(中庸)',                       '"치우치지 않고 평소에도 지키는 도"'),
        ('저자',    '자사(子思, BC 483?~402?)',         '공자의 손자 · 공급(孔伋)'),
        ('학통',    '공자 → 증자 → 자사 → 맹자',         '유학 정통 계보의 중심축'),
        ('원출처',  '『예기(禮記)』 제31편',              '본래 독립 책이 아닌 한 편'),
        ('분량',    '33장 · 약 3,500자',                  '간결하지만 가장 심오'),
        ('주희의 평가','공문전수지심법(孔門傳授之心法)',  '공자 문하가 대대로 전수한 마음의 법'),
        ('사서의 위치','대학·논어·맹자 다음 — 정점',     '"가장 나중에 읽어야 할 책"'),
    ]
    top = 2.15
    for i, (tag, val, note) in enumerate(rows):
        bg = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.3), Inches(0.65), PALE)
        add_filled_rect(slide, Inches(2.85), Inches(top), Inches(10.0), Inches(0.65), bg)
        add_textbox(slide, Inches(0.55), Inches(top + 0.18), Inches(2.2), Inches(0.4),
                    tag, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.95), Inches(top + 0.05), Inches(4.5), Inches(0.5),
                    val, font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.5), Inches(top + 0.08), Inches(5.3), Inches(0.5),
                    note, font_size=13, color=SUB)
        top += 0.7


@S('Ⅰ. 개요')
def s_jasa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '자사(子思, BC 483?~402?) — 공자의 손자')
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(3.5),
                '子\n思', font_size=120, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(4.2), Inches(2.3), Inches(8.7), Inches(4.5), [
        ('이름과 신분', {'bold': True, 'font_size': 18, 'color': ACCENT}),
        ('  본명 공급(孔伋), 자(字) 자사(子思)',
         {'font_size': 15}),
        ('  공자의 손자 — 공자의 외동아들 백어(伯魚)의 아들',
         {'font_size': 15}),
        ('학통', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 8}),
        ('  공자 → 증자(曾子) → 자사 → 맹자',
         {'font_size': 15, 'bold': True}),
        ('  유학 정통 학통의 중심에 위치',
         {'font_size': 13, 'color': SUB}),
        ('저작 동기', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 8}),
        ('  주희의 진단 — "도학(道學)이 실전될 것을 근심하여 지음"',
         {'font_size': 14}),
        ('  16자 심법(十六字心法)을 계승·체계화',
         {'font_size': 14, 'color': SUB}),
        ('사상사적 위치', {'bold': True, 'font_size': 18, 'color': ACCENT, 'space_before': 8}),
        ('  공자의 인(仁)과 맹자의 성선(性善)을 연결하는 다리',
         {'font_size': 14}),
    ], line_spacing=1.3)


@S('Ⅰ. 개요')
def s_jeongja_zhuxi(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '정자(程子)·주희(朱熹)의 평가')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('정이천(程伊川)의 한 마디',
         {'font_size': 16, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"공문전수지심법(孔門傳授之心法)"',
         {'font_size': 24, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('— 공자 문하가 대대로 전수한 마음의 법 —',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        ('"펼치면 우주 여섯 방위에 가득차고, 거두면 은밀한 곳에 숨는다"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('(放之則彌六合, 卷之則退藏於密)',
         {'font_size': 12, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('주희 『중용장구』 서문',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('"중용은 자사가 도학이 실전될 것을 근심하여 지은 책이다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 요·순의 「인심유위, 도심유미」 16자 심법을 계승한 것"',
         {'font_size': 14, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅰ. 개요')
def s_reverse_order(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅰ. 개요', page, total)
    add_title(slide, '사서(四書)의 역순 — 가장 마지막에 읽는 정점',
              '"중용은 사서의 종착지 · 가장 심오하기에 마지막에 둔다"')
    books = [
        ('대학(大學)',  '規模 — 규모',  '전체의 지도',     '먼저 읽어 큰 그림'),
        ('논어(論語)',  '根本 — 근본',  '인(仁)의 뿌리',   '두 번째로 근본 세움'),
        ('맹자(孟子)',  '發揮 — 발휘',  '활달한 논변',     '세 번째로 발휘'),
        ('중용(中庸)',  '微妙 — 미묘',  '형이상학·심법',   '마지막 — 정점'),
    ]
    top = 2.4
    for i, (name, role, desc, when) in enumerate(books):
        is_last = (i == 3)
        c = ACCENT if is_last else SUB
        bg = RGBColor(0xFA, 0xE5, 0xE5) if is_last else PALE
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.95), c)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(2.5), Inches(0.5),
                    name, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(2.5), Inches(0.95), PALE)
        add_textbox(slide, Inches(3.2), Inches(top + 0.27), Inches(2.5), Inches(0.5),
                    role, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.8), Inches(top), Inches(3.5), Inches(0.95), bg)
        add_textbox(slide, Inches(5.9), Inches(top + 0.27), Inches(3.3), Inches(0.5),
                    desc, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(9.4), Inches(top), Inches(3.4), Inches(0.95), bg)
        add_textbox(slide, Inches(9.5), Inches(top + 0.27), Inches(3.2), Inches(0.5),
                    when, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        top += 1.05
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                '대학이 "사업계획서"라면 중용은 "존재론" — 가장 나중에 가장 깊이',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 왜 왕의 책인가 ==============
@S('Ⅱ. 왕의 책')
def s_king_book(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 왕의 책', page, total)
    add_title(slide, '동아시아 제왕학의 심장 — 왜 왕이 꼭 읽었나')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '중용 = 통치의 형이상학 + 판단의 기준 + 자기감찰의 도구 + 감정 관리 철학',
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('16자 심법',        '요·순·우에게 전수된 왕권의 가장 깊은 뿌리'),
        ('구경(九經)',        '천하국가를 다스리는 9가지 법도 — 20장의 핵심'),
        ('성자천지도',        '왕권의 정당성을 지탱하는 최종 근거'),
        ('시중(時中)',        '제왕의 판단력의 다른 이름'),
        ('신독(愼獨)',        '구중궁궐의 "홀로 있음"을 감찰'),
        ('중화(中和)',        '왕의 감정 통치의 철학'),
    ]
    top = 3.25
    for tag, desc in items:
        add_filled_rect(slide, Inches(0.5), Inches(top), Inches(3.0), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.5), Inches(top + 0.13), Inches(3.0), Inches(0.4),
                    tag, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.6), Inches(top), Inches(9.2), Inches(0.55),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.8), Inches(top + 0.13), Inches(8.9), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.62


@S('Ⅱ. 왕의 책')
def s_16char(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 왕의 책', page, total)
    add_title(slide, '16자 심법(十六字心法) — 왕권의 뿌리')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.5), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
                '人 心 惟 危   道 心 惟 微',
                font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.6),
                '惟 精 惟 一   允 執 厥 中',
                font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
                '인심유위 도심유미 · 유정유일 윤집궐중',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"인심(人心)은 위태롭고 도심(道心)은 미묘하니',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 오직 정밀(惟精)하게 하고 오직 한결같이(惟一) 하여',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 진실로 그 중(中)을 잡으라(允執厥中)"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('요(堯) → 순(舜) → 우(禹)로 전수된 제왕의 16자',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('중용은 이 "중(中)"의 철학적 해명서',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('Ⅱ. 왕의 책')
def s_human_dao_mind(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 왕의 책', page, total)
    add_title(slide, '人心 vs 道心 — 왕이 매 순간 선택해야 하는 두 마음')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '人 心 — 인심', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('욕망에 이끌리는 마음', {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"惟 危" — 늘 위태롭다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 6}),
        ('• 권력의 유혹', {'font_size': 13, 'space_before': 8}),
        ('• 측근의 아첨', {'font_size': 13, 'space_before': 4}),
        ('• 감각의 쾌락', {'font_size': 13, 'space_before': 4}),
        ('• 분노·욕심', {'font_size': 13, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '道 心 — 도심', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('하늘의 이치를 따르는 마음', {'font_size': 17, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"惟 微" — 희미하고 잡기 어렵다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 6}),
        ('• 도덕적 판단', {'font_size': 13, 'space_before': 8}),
        ('• 공정성', {'font_size': 13, 'space_before': 4}),
        ('• 백성의 안위', {'font_size': 13, 'space_before': 4}),
        ('• 천명의 인식', {'font_size': 13, 'space_before': 4}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '"왕이 중(中)을 잃으면 나라가 기운다" — 인심·도심 사이에서 매 순간 시중',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅱ. 왕의 책')
def s_9governance(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 왕의 책', page, total)
    add_title(slide, '구경(九經) — 천하국가를 다스리는 9가지 법도',
              '— 중용 20장, 노 애공(哀公)의 정치 물음에 대한 공자의 답')
    items = [
        ('1', '修身', '수신',   '자기 몸을 닦음'),
        ('2', '尊賢', '존현',   '어진 이를 존경'),
        ('3', '親親', '친친',   '친족을 친하게'),
        ('4', '敬大臣', '경대신', '대신을 공경'),
        ('5', '體群臣', '체군신', '군신을 몸처럼'),
        ('6', '子庶民', '자서민', '서민을 자식처럼'),
        ('7', '來百工', '래백공', '장인을 오게 함'),
        ('8', '柔遠人', '유원인', '먼 곳 사람을 부드럽게'),
        ('9', '懷諸侯', '회제후', '제후를 품음'),
    ]
    top = 2.3
    for i, (num, han, eum, kor) in enumerate(items):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.3
        y = top + row * 1.45
        add_filled_rect(slide, Inches(x), Inches(y), Inches(0.6), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(x), Inches(y + 0.1), Inches(0.6), Inches(0.4),
                    num, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x + 0.7), Inches(y), Inches(3.4), Inches(0.55), PALE)
        add_textbox(slide, Inches(x + 0.7), Inches(y + 0.1), Inches(3.4), Inches(0.4),
                    f'{han} ({eum})', font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(x), Inches(y + 0.6), Inches(4.1), Inches(0.75),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.78), Inches(3.8), Inches(0.4),
                    kor, font_size=13, color=INK, align=PP_ALIGN.CENTER)


@S('Ⅱ. 왕의 책')
def s_seong_cheondo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 왕의 책', page, total)
    add_title(slide, '誠者 天之道 — 왕권의 정당성을 지탱하는 최종 근거')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '誠 者   天 之 道 也',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.9),
                '誠 之 者   人 之 道 也',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
                '성자 천지도야 · 성지자 인지도야',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
                '— 중용 20장', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.3), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.6), Inches(5.55), Inches(11.9), Inches(1.5), PALE)
    add_paragraphs(slide, Inches(0.9), Inches(5.7), Inches(11.3), Inches(1.3), [
        ('"왕의 권위는 천명(天命)에서 나온다"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('하늘의 본성이 "성(誠)"이라면, 왕도 "성"을 근본으로',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('왕이 거짓되면 천명을 잃고, 천명을 잃으면 혁명이 일어난다',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅱ. 왕의 책')
def s_sijung_sindok_jungwa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 왕의 책', page, total)
    add_title(slide, '시중·신독·중화 — 왕의 3대 자기 점검')
    items = [
        ('時中', '시중',  '"군자이시중" — 군자는 때에 맞게 중을 쓴다',
         '제왕의 판단력 — 같은 일도 때마다 다른 정답'),
        ('愼獨', '신독',  '"군자신기독" — 홀로 있음을 삼간다',
         '왕의 자기감찰 — 침전(寢殿)의 모습이 진짜 왕'),
        ('中和', '중화',  '"치중화 천지위언" — 중화를 이루면 천지가 제자리',
         '왕의 감정 통치 — 절도에 맞게 발하라'),
    ]
    top = 2.4
    for han, eum, principle, application in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.5), Inches(1.3), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.35), Inches(1.5), Inches(0.6),
                    han, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.2), Inches(top), Inches(10.6), Inches(1.3), PALE)
        add_textbox(slide, Inches(2.4), Inches(top + 0.15), Inches(10.2), Inches(0.4),
                    eum, font_size=12, color=SUB)
        add_textbox(slide, Inches(2.4), Inches(top + 0.45), Inches(10.2), Inches(0.4),
                    principle, font_size=14, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.4), Inches(top + 0.85), Inches(10.2), Inches(0.4),
                    application, font_size=13, color=INK)
        top += 1.45


@S('Ⅱ. 왕의 책')
def s_history(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅱ. 왕의 책', page, total)
    add_title(slide, '실제 교육 역사 — 중국·한국·일본 왕실')
    items = [
        ('중국',  '청 강희제 — 『일강중용해의(日講中庸解義)』 매일 강독',
         '당 태종 공영달의 『예기정의』, 송 주희 이후 과거 표준'),
        ('한국 — 경연(經筵)',  '조선 왕실의 경연에서 가장 많이 강독된 책 중 하나',
         '논어·맹자와 함께 핵심 교재'),
        ('한국 — 학자',  '퇴계 『성학십도』 6·9도 (심통성정·경재잠)',
         '율곡 『성학집요』 — 구경(九經)을 왕도의 골격으로'),
        ('한국 — 정조',  '18세 때 깊이 읽고 『중용강의』 친저',
         '정약용 『중용강의보』 — 실학의 중용 해석'),
        ('일본',  '도쿠가와 이에야스 — 하야시 라잔에게 강독',
         '막부 말기 요시다 쇼인이 제자들에게 반복 강의'),
    ]
    top = 2.4
    for region, primary, secondary in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    region, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(9.6), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.4), Inches(top + 0.08), Inches(9.3), Inches(0.4),
                    primary, font_size=13, bold=True, color=INK)
        add_textbox(slide, Inches(3.4), Inches(top + 0.5), Inches(9.3), Inches(0.4),
                    secondary, font_size=12, color=SUB)
        top += 0.95


# ============== Ⅲ. 구조 ==============
@S('Ⅲ. 구조')
def s_structure(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 구조', page, total)
    add_title(slide, '33장 + 3부 구조 — 주희 『중용장구』')
    items = [
        ('1부',  '제 1 장',         '강령(綱領, 총론)',
         '"천명지위성"에서 "치중화"까지 — 중용의 형이상학적 뼈대'),
        ('2부',  '제 2~20장',      '본론 ① — 중용의 실천, 도, 정치',
         '공자의 어록 — 중용이 현실에 펼쳐지는 모습'),
        ('3부',  '제 21~33장',     '본론 ② — 성(誠)의 형이상학',
         '자사의 논설 — 중용의 궁극적 정점'),
    ]
    top = 2.5
    for part, scope, role, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.5), Inches(1.2), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.4), Inches(1.5), Inches(0.5),
                    part, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.2), Inches(top), Inches(2.5), Inches(1.2), PALE)
        add_textbox(slide, Inches(2.2), Inches(top + 0.4), Inches(2.5), Inches(0.5),
                    scope, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.8), Inches(top), Inches(8.0), Inches(1.2),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.0), Inches(top + 0.18), Inches(7.7), Inches(0.45),
                    role, font_size=15, bold=True, color=INK)
        add_textbox(slide, Inches(5.0), Inches(top + 0.65), Inches(7.7), Inches(0.45),
                    desc, font_size=13, color=SUB)
        top += 1.4


@S('Ⅲ. 구조')
def s_key_chapters(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 구조', page, total)
    add_title(slide, '핵심 장별 지도 — 한눈에 보는 중용')
    rows = [
        ('1',    '천명·중화·신독',        '중용의 형이상학적 뼈대'),
        ('2',    '군자시중',              '중용은 "때의 중(時中)"'),
        ('3',    '중용난능',              '중용은 가장 어려운 덕'),
        ('6',    '집기양단 용기중',        '순 임금의 지혜'),
        ('13',   '도불원인',              '도는 사람에게서 멀지 않다'),
        ('20',   '구경 · 성자천지도',      '중용의 중심 — 가장 긴 장'),
        ('22',   '지성진성',              '천지와 더불어 셋이 됨(天地參)'),
        ('26',   '지성무식',              '지극한 성은 쉬지 않는다'),
        ('33',   '무성무취',              '결론 — 하늘의 일은 소리·냄새 없다'),
    ]
    top = 2.0
    for ch, name, desc in rows:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.9), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(0.9), Inches(0.4),
                    ch + '장', font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.6), Inches(top), Inches(3.5), Inches(0.5), PALE)
        add_textbox(slide, Inches(1.7), Inches(top + 0.08), Inches(3.4), Inches(0.4),
                    name, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.2), Inches(top), Inches(7.6), Inches(0.5),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.4), Inches(top + 0.08), Inches(7.4), Inches(0.4),
                    desc, font_size=14, color=INK)
        top += 0.55


@S('Ⅲ. 구조')
def s_flow(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅲ. 구조', page, total)
    add_title(slide, '강령 → 실천 → 형이상학 — 3단 흐름')
    flow = [
        ('1장',         '강령',         '천명지위성 · 중화 · 신독',  ACCENT),
        ('2~11장',      '시중·군자',     '군자시중 · 순임금 · 안회',  RGBColor(0xA0, 0x40, 0x40)),
        ('12~19장',     '도의 보편성',   '도불원인 · 소위소행 · 귀신',RGBColor(0x70, 0x40, 0x60)),
        ('20장',        '정치의 중심',   '구경 · 성자천지도 · 5단계',  ACCENT),
        ('21~26장',     '성의 형이상학', '지성진성 · 천지참 · 지성무식', RGBColor(0xA0, 0x40, 0x40)),
        ('27~33장',     '성인의 도',     '왕도의 완성 · 무성무취',     SUB),
    ]
    top = 2.3
    for scope, tag, content, color in flow:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(0.65), color)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(2.0), Inches(0.4),
                    scope, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.75), Inches(top), Inches(2.5), Inches(0.65), PALE)
        add_textbox(slide, Inches(2.75), Inches(top + 0.13), Inches(2.5), Inches(0.4),
                    tag, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(5.4), Inches(top), Inches(7.4), Inches(0.65),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.6), Inches(top + 0.13), Inches(7.2), Inches(0.4),
                    content, font_size=13, color=INK)
        top += 0.75


# ============== Ⅳ. 1장 (총론) ==============
@S('Ⅳ. 1장')
def s_chap1_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 1장', page, total)
    add_title(slide, '1장 — 중용의 형이상학적 뼈대 (강령)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('단 한 장에 중용 전체의 뼈대가 압축',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 천명지위성 — 4단계 인과 (天→性→道→敎)',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('• 신독(愼獨) — 막현호은 막현호미',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('• 중화(中和) — 미발의 중, 이발의 화',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        ('• 치중화(致中和) — 천지위언 만물육언',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('수미상관(首尾相關)의 구조',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('1장 — "막현호은(莫見乎隱)" — 숨은 것보다 더 드러난 것 없다',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('33장 — "무성무취(無聲無臭)" — 하늘의 일은 소리·냄새가 없다',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"드러나지 않음"의 두 끝이 책 전체를 둘러싼다',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅳ. 1장')
def s_chunmyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 1장', page, total)
    add_title(slide, '天命之謂性 — 유학 형이상학의 출발점')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '天 命 之 謂 性',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.9),
                '率 性 之 謂 道   修 道 之 謂 敎',
                font_size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
                '천명지위성 · 솔성지위도 · 수도지위교',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    chain = ['天', '→', '性', '→', '道', '→', '敎']
    kor = ['하늘', '', '본성', '', '도', '', '교육']
    box_w = 1.4
    arrow_w = 0.4
    total_w = box_w * 4 + arrow_w * 3
    start_x = (13.333 - total_w) / 2
    cur_x = start_x
    for i, item in enumerate(chain):
        is_arrow = item == '→'
        w = arrow_w if is_arrow else box_w
        if not is_arrow:
            add_filled_rect(slide, Inches(cur_x), Inches(5.0), Inches(box_w),
                            Inches(1.2), PALE)
            add_textbox(slide, Inches(cur_x), Inches(5.15), Inches(box_w), Inches(0.6),
                        item, font_size=36, bold=True, color=ACCENT,
                        align=PP_ALIGN.CENTER)
            label = kor[i]
            add_textbox(slide, Inches(cur_x), Inches(5.75), Inches(box_w), Inches(0.4),
                        label, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        else:
            add_textbox(slide, Inches(cur_x), Inches(5.4), Inches(arrow_w), Inches(0.5),
                        '→', font_size=22, color=SUB, align=PP_ALIGN.CENTER)
        cur_x += w
    add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
                '"인간의 선함은 하늘의 명(命)이다" — 맹자 성선설의 뿌리',
                font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 1장')
def s_gyesin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 1장', page, total)
    add_title(slide, '戒愼恐懼 — 보이지 않는 곳에서도 경계하라')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.6),
                '"君子戒愼乎其所不睹 恐懼乎其所不聞"',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
                '군자계신호기소부도 공구호기소불문',
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(3.9), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '"군자는 남이 보지 않는 곳에서도 경계하고',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
                ' 남이 듣지 않는 곳에서도 두려워한다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.4), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.2), [
        ('남이 안 볼 때의 행동이 진짜 인격',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('왕에게 — 구중궁궐의 침전에서의 모습이 진짜 왕',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅳ. 1장')
def s_makhyeon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 1장', page, total)
    add_title(slide, '莫見乎隱 莫顯乎微 — 숨은 것이 가장 드러난다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '莫 見 乎 隱   莫 顯 乎 微',
                font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '막현호은 막현호미',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.0), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '"숨은 것보다 더 드러난 것이 없고',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.6),
                ' 미세한 것보다 더 분명한 것이 없다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.4), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.2), [
        ('"故君子愼其獨也" — 그러므로 군자는 홀로 있음을 삼간다',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('성리학 수양법의 핵심 — 퇴계 "경(敬)"·율곡 "성(誠)"의 뿌리',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅳ. 1장')
def s_jungwa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 1장', page, total)
    add_title(slide, '中和 — 미발(未發)의 중, 이발(已發)의 화')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '中 — 중(中)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('喜怒哀樂之 未發',
         {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('희노애락이 아직 발하지 않음',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"감정이 일어나기 전의 고요"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 6}),
        ('"천하의 큰 근본(天下之大本)"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '和 — 화(和)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('發而皆 中節',
         {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('발하여 모두 절도에 맞음',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('"일어난 감정이 딱 맞게 발현"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 6}),
        ('"천하의 통달하는 도(天下之達道)"',
         {'font_size': 14, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '"감정을 억누르라"가 아니라 "발할 때 절도(節)에 맞게 발하라"',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅳ. 1장')
def s_chijungwa(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅳ. 1장', page, total)
    add_title(slide, '致中和 天地位焉 萬物育焉 — 중화의 우주적 의미')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '致 中 和',
                font_size=66, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.7),
                '天 地 位 焉   萬 物 育 焉',
                font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                '치중화 · 천지위언 만물육언',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.95), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.6),
                '"중화를 이루면 천지가 제자리를 잡고',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.6),
                ' 만물이 자라난다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '→ 개인의 마음이 우주와 연결되는 동양적 우주론',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅴ. 시중과 군자 ==============
@S('Ⅴ. 시중과 군자')
def s_gunja_sijung(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 시중과 군자', page, total)
    add_title(slide, '君子時中 — 군자는 때에 맞게 중을 쓴다')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '君 子 之 中 庸 也',
                font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.9),
                '君 子 而 時 中',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.5),
                '군자지중용야 · 군자이시중',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4),
                '— 중용 2장', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.2), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.4), [
        ('"군자의 중용이란 군자이면서 때에 맞게 중을 쓰는 것"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('→ 중은 고정된 지점이 아니다 — 매번 새로 발견해야 하는 "그때의 정답"',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅴ. 시중과 군자')
def s_minseon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 시중과 군자', page, total)
    add_title(slide, '中庸之爲德也 其至矣乎 — 중용의 덕은 지극하다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '中 庸 之 爲 德 也',
                font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.9),
                '其 至 矣 乎   民 鮮 久 矣',
                font_size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
                '중용지위덕야 · 기지의호 민선구의',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
                '— 공자의 말, 중용 3장 인용', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(5.3), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.6),
                '"중용의 덕됨이 지극하구나!',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.6),
                ' 백성들 중에 이를 오래 지키는 자가 드물다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
                '쉬워 보이지만 가장 어려운 덕 — 중용의 역설',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S('Ⅴ. 시중과 군자')
def s_difficult(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 시중과 군자', page, total)
    add_title(slide, '中庸難能 — 중용은 가장 어려운 덕')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('공자의 안타까운 진단',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"백성들 중에 중용을 오래 지키는 자가 드물다(民鮮久矣)"',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('"천하의 나라도 균등히 다스릴 수 있고, 작록도 사양할 수 있고',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 흰 칼날도 밟을 수 있으나, 중용은 능하기 어렵다(中庸不可能也)"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('왜 가장 어려운가?',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('• 극단은 쉽다 — 분명한 선택', {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 중간(中)은 어렵다 — 매번 새로 찾아야 하므로', {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('• 지속(庸)은 더 어렵다 — 한 번이 아니라 평생이므로', {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅴ. 시중과 군자')
def s_sun(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 시중과 군자', page, total)
    add_title(slide, '순(舜) 임금의 지혜 — 執其兩端 用其中於民')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.5),
                '"舜其大知也與!" — 순 임금은 참으로 큰 지혜로운 분',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('好問',     '호문',     '묻기를 좋아함'),
        ('好察邇言', '호찰이언', '가까운 말(서민·일상의 말)을 살피기 좋아함'),
        ('隱惡揚善', '은악양선', '악은 덮고 선은 드러냄'),
        ('執其兩端', '집기양단', '극단의 양쪽을 모두 파악'),
        ('用其中於民', '용기중어민', '그 중(中)을 백성에게 적용'),
    ]
    top = 3.25
    for han, eum, kor in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(3.0), Inches(0.6), PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.08), Inches(3.0), Inches(0.4),
                    han, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(2.5), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.7), Inches(top + 0.14), Inches(2.5), Inches(0.4),
                    eum, font_size=13, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.3), Inches(top), Inches(6.5), Inches(0.6),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(6.5), Inches(top + 0.14), Inches(6.2), Inches(0.4),
                    kor, font_size=13, color=INK)
        top += 0.7
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '"극단을 이해하지 못하면 중을 알 수 없다" — 왕의 사고법',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅴ. 시중과 군자')
def s_anhoe(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅴ. 시중과 군자', page, total)
    add_title(slide, '안회(顔回)의 실천 — 拳拳服膺')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('"回之爲人也 擇乎中庸',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 得一善 則拳拳服膺 而弗失之矣"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"안회의 사람됨은 중용을 택하여',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 하나의 선을 얻으면 받들어 가슴에 품고 잃지 않았다"',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"권권복응(拳拳服膺)" — 받들어 가슴에 품다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('한 번 얻은 좋은 것을 평생 놓지 않는 안회의 지속력',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"중용은 천재의 일이 아니라 끈기의 일"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


# ============== Ⅵ. 도의 보편성 ==============
@S('Ⅵ. 도의 보편성')
def s_dobulwon(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 도의 보편성', page, total)
    add_title(slide, '道不遠人 — 도는 사람에게서 멀지 않다 (13장)')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '道 不 遠 人',
                font_size=80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '도불원인',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
                '— 중용 13장 (공자)', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.55), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.6),
                '"도는 사람에게서 멀지 않다',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.6),
                ' 사람이 도를 행하면서 사람을 멀리 한다면 도라 할 수 없다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.8), PALE)
    add_textbox(slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
                '"도는 산속 수도가 아니라 부모·형제·부부·붕우·군신의 일상 안에"',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅵ. 도의 보편성')
def s_so_so(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 도의 보편성', page, total)
    add_title(slide, '素位而行 — 자기 위치에서 행한다 (14장)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '"君子素其位而行 不願乎其外"',
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
                '군자소기위이행 불원호기외',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.4),
                '— 중용 14장', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.3), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.6),
                '"군자는 자기 자리(位)에 맞게 행하고',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
                ' 그 밖의 것을 원하지 않는다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(6.1), Inches(11.9), Inches(1.0), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.8), [
        ('"부귀에 처하면 부귀에 맞게, 빈천에 처하면 빈천에 맞게',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
        (' 환난에 처하면 환난에 맞게 — 어느 자리에서나 흔들리지 않는다"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅵ. 도의 보편성')
def s_guisin(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅵ. 도의 보편성', page, total)
    add_title(slide, '鬼神之德 — 은미함과 드러남의 통일 (16장)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), [
        ('"視之而弗見 聽之而弗聞',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 體物而不可遺"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"보아도 보이지 않고, 들어도 들리지 않으나',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 만물에 깃들어 빠뜨릴 수 없다"',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)
    add_filled_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.8), [
        ('"夫微之顯 誠之不可揜 如此夫"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"은미함이 드러나는 것 — 성(誠)을 가릴 수 없음이 이와 같다"',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('가장 미세한 것이 가장 분명하게 드러난다',
         {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


# ============== Ⅶ. 20장 — 중용의 중심 ==============
@S('Ⅶ. 20장')
def s_chap20_intro(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '20장 — 중용의 중심 · 가장 긴 장')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
                '哀公問政 — 노(魯) 애공이 공자에게 정치를 묻다',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.7), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.5), Inches(11.7), Inches(3.4), [
        ('20장이 담고 있는 4가지 핵심',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('① 구경(九經) — 천하국가 통치의 9가지 법도',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('② 박학·심문·신사·명변·독행 — 학습의 5단계',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('③ 인일능지기백지 — 천 번의 반복의 정신 (그릿)',
         {'font_size': 16, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('④ 성자천지도 — 성(誠) 개념의 결정적 도입',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('"왕들이 줄을 치며 읽은 장" — 책 전체의 중심',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
    ], line_spacing=1.3)


@S('Ⅶ. 20장')
def s_9governance_overview(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '九經 개관 — "凡爲天下國家有九經"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
                '"凡爲天下國家有九經"',
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.4),
                '범위천하국가유구경 — "무릇 천하국가를 다스림에 아홉 가지 법이 있다"',
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        ('① 修身', '수신'),  ('② 尊賢', '존현'),  ('③ 親親', '친친'),
        ('④ 敬大臣', '경대신'), ('⑤ 體群臣', '체군신'), ('⑥ 子庶民', '자서민'),
        ('⑦ 來百工', '래백공'), ('⑧ 柔遠人', '유원인'), ('⑨ 懷諸侯', '회제후'),
    ]
    top = 3.7
    for i, (han, eum) in enumerate(items):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.3
        y = top + row * 1.1
        add_filled_rect(slide, Inches(x), Inches(y), Inches(4.1), Inches(0.95), PALE)
        add_textbox(slide, Inches(x), Inches(y + 0.15), Inches(4.1), Inches(0.4),
                    han, font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(y + 0.6), Inches(4.1), Inches(0.4),
                    eum, font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅶ. 20장')
def s_9governance_1_5(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '九經 1~5 — 자기에서 군신까지')
    items = [
        ('1', '修身',   '수신',     '자기 몸을 닦음',           '리더의 자기 관리'),
        ('2', '尊賢',   '존현',     '어진 이를 존경',           '탁월한 인재의 우대'),
        ('3', '親親',   '친친',     '친족을 친하게',            '가족·측근 질서'),
        ('4', '敬大臣', '경대신',   '대신을 공경',              '핵심 임원의 권위 존중'),
        ('5', '體群臣', '체군신',   '군신을 몸처럼',            '조직원과 일체감'),
    ]
    top = 2.3
    for num, han, eum, classical, modern in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.0), Inches(0.85), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.1), Inches(2.0), Inches(0.4),
                    han, font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.4), Inches(top + 0.5), Inches(2.0), Inches(0.4),
                    eum, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.5), Inches(top), Inches(4.0), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.7), Inches(top + 0.22), Inches(3.7), Inches(0.5),
                    classical, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.6), Inches(top), Inches(5.2), Inches(0.85),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(7.8), Inches(top + 0.22), Inches(4.9), Inches(0.5),
                    modern, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        top += 0.95


@S('Ⅶ. 20장')
def s_9governance_6_9(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '九經 6~9 — 백성에서 제후까지')
    items = [
        ('6', '子庶民', '자서민',   '서민을 자식처럼',          '일반 구성원·고객을 자녀처럼'),
        ('7', '來百工', '래백공',   '장인을 오게 함',           '각 분야 전문가의 영입'),
        ('8', '柔遠人', '유원인',   '먼 곳 사람을 부드럽게',    '외부·국제 관계 관리'),
        ('9', '懷諸侯', '회제후',   '제후를 품음',              '동맹·협력 파트너 관리'),
    ]
    top = 2.4
    for num, han, eum, classical, modern in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.95), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.27), Inches(0.7), Inches(0.5),
                    num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.0), Inches(0.95), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.15), Inches(2.0), Inches(0.4),
                    han, font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.4), Inches(top + 0.55), Inches(2.0), Inches(0.4),
                    eum, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.5), Inches(top), Inches(4.0), Inches(0.95),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.7), Inches(top + 0.3), Inches(3.7), Inches(0.5),
                    classical, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.6), Inches(top), Inches(5.2), Inches(0.95),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(7.8), Inches(top + 0.3), Inches(4.9), Inches(0.5),
                    modern, font_size=13, color=INK, align=PP_ALIGN.CENTER)
        top += 1.05


@S('Ⅶ. 20장')
def s_9governance_effect(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '九經의 효과 — 각 항목의 결과')
    items = [
        ('修身', '도(道)가 선다'),
        ('尊賢', '유혹되지 않는다'),
        ('親親', '친족 갈등이 없다'),
        ('敬大臣', '현혹되지 않는다'),
        ('體群臣', '선비가 보답한다'),
        ('子庶民', '백성이 부지런해진다'),
        ('來百工', '재용이 족해진다'),
        ('柔遠人', '사방이 귀의한다'),
        ('懷諸侯', '천하가 두려워한다'),
    ]
    top = 2.0
    for i, (han, effect) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = top + row * 0.65 + (0.35 if i == 8 else 0)
        # 마지막 항목은 가운데 정렬
        if i == 8:
            x = (13.333 - 6.0) / 2
        add_filled_rect(slide, Inches(x), Inches(y), Inches(2.0), Inches(0.55), ACCENT)
        add_textbox(slide, Inches(x), Inches(y + 0.13), Inches(2.0), Inches(0.4),
                    han, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 2.05), Inches(y + 0.13), Inches(0.6), Inches(0.4),
                    '→', font_size=18, color=ACCENT, align=PP_ALIGN.CENTER)
        w = 4.0 if i != 8 else 3.4
        add_filled_rect(slide, Inches(x + 2.65), Inches(y), Inches(w), Inches(0.55), PALE)
        add_textbox(slide, Inches(x + 2.8), Inches(y + 0.13), Inches(w - 0.2), Inches(0.4),
                    effect, font_size=13, color=INK)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                '왕이 이 9가지를 체크리스트로 삼으면 그 치세가 흥한다',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅶ. 20장')
def s_seong_cheondo_full(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '誠者 天之道 — 중용 전체를 뒤바꾸는 명제')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '誠 者   天 之 道 也',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.9),
                '誠 之 者   人 之 道 也',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
                '성자 천지도야 · 성지자 인지도야',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.9), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
                '"성(誠)은 하늘의 도이고, 성실하고자 함(誠之)은 사람의 도다"',
                font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.1), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.9), [
        ('성(誠) = 단순한 성실성이 아니라 "존재의 진정성"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('하늘의 도는 저절로 참되고 / 사람의 도는 참되고자 노력하는 것',
         {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅶ. 20장')
def s_5stages(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '학습의 5단계 — 博學·審問·愼思·明辨·篤行')
    items = [
        ('1', '博 學', '박학', '널리 배운다',          '양의 축적 — T자형 학습'),
        ('2', '審 問', '심문', '자세히 묻는다',        '비판적 사고 — Critical Thinking'),
        ('3', '愼 思', '신사', '신중하게 생각한다',    '깊은 사유 — Deep Work'),
        ('4', '明 辨', '명변', '밝게 분별한다',        '판단력 — Judgement'),
        ('5', '篤 行', '독행', '도탑게 실천한다',      '실천력 — Execution'),
    ]
    top = 2.3
    for num, han, eum, classical, modern in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(0.7), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.18), Inches(0.7), Inches(0.4),
                    num, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.4), Inches(top), Inches(2.2), Inches(0.8), PALE)
        add_textbox(slide, Inches(1.4), Inches(top + 0.05), Inches(2.2), Inches(0.4),
                    han, font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.4), Inches(top + 0.5), Inches(2.2), Inches(0.4),
                    eum, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.7), Inches(top), Inches(3.5), Inches(0.8),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(3.85), Inches(top + 0.2), Inches(3.3), Inches(0.4),
                    classical, font_size=14, bold=True, color=INK)
        add_filled_rect(slide, Inches(7.3), Inches(top), Inches(5.5), Inches(0.8),
                        RGBColor(0xFA, 0xE5, 0xE5))
        add_textbox(slide, Inches(7.5), Inches(top + 0.2), Inches(5.2), Inches(0.4),
                    modern, font_size=13, color=ACCENT, bold=True)
        top += 0.9


@S('Ⅶ. 20장')
def s_inbaek(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '人一能之 己百之 — 그릿(Grit)의 동양 원형')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '人 一 能 之   己 百 之',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.9),
                '人 十 能 之   己 千 之',
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
                '인일능지 기백지 · 인십능지 기천지',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.85), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.6),
                '"남이 한 번에 능하면 나는 백 번을 하고',
                font_size=18, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.6),
                ' 남이 열 번에 능하면 나는 천 번을 한다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '"果能此道矣 雖愚必明 雖柔必強" — 어리석어도 반드시 밝아지고, 유약해도 반드시 강해진다',
                font_size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
                '→ Angela Duckworth의 Grit 이론·1만 시간의 법칙의 원형 — 조선 선비의 좌우명',
                font_size=12, color=SUB, align=PP_ALIGN.CENTER)


@S('Ⅶ. 20장')
def s_taeksun(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅶ. 20장', page, total)
    add_title(slide, '擇善而固執之 — 현실적 수양의 방법')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '擇 善 而 固 執 之',
                font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
                '택선이고집지',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
                '— 중용 20장', font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.65), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.6),
                '"선을 택하여 굳게 붙드는 것"',
                font_size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(5.95), Inches(11.7), Inches(1.1), [
        ('성인(聖人) — 힘쓰지 않아도 중에 맞는다 (不勉而中)',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('보통 사람 — 선을 택해 굳게 붙드는 것이 수양의 길',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅷ. 성의 형이상학 ==============
@S('Ⅷ. 성의 형이상학')
def s_jaseong_jamyeong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 성의 형이상학', page, total)
    add_title(slide, '自誠明 自明誠 — 성(誠)과 명(明)의 두 방향 (21장)')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '自 誠 明 — 자성명', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"성(誠)으로부터 밝아지는 것"',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('이것을 "성(性)"이라 한다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 6}),
        ('성인(聖人)의 길', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('타고난 본성에서', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('저절로 통달',     {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '自 明 誠 — 자명성', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('"밝음으로부터 성에 이르는 것"',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('이것을 "교(敎)"라 한다',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('', {'font_size': 6}),
        ('보통 사람의 길', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('배움과 명변을 통해', {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('성(誠)에 이름',       {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)


@S('Ⅷ. 성의 형이상학')
def s_jiseong_jinseong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 성의 형이상학', page, total)
    add_title(slide, '至誠盡性 — 자기→사람→만물→천지 (22장)')
    chain = [
        ('盡 其 性',       '자기 성을 다함'),
        ('盡 人 之 性',     '사람의 성을 다함'),
        ('盡 物 之 性',     '만물의 성을 다함'),
        ('贊 天 地 化 育',  '천지의 화육을 도움'),
        ('與 天 地 參',     '천지와 더불어 셋(參)'),
    ]
    top = 2.3
    for han, kor in chain:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(5.5), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.13), Inches(5.5), Inches(0.4),
                    han, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(6.3), Inches(top + 0.18), Inches(0.4), Inches(0.4),
                    '→', font_size=18, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(6.9), Inches(top), Inches(5.9), Inches(0.7), PALE)
        add_textbox(slide, Inches(7.1), Inches(top + 0.18), Inches(5.6), Inches(0.4),
                    kor, font_size=15, color=INK)
        top += 0.8
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '"唯天下至誠爲能盡其性" — 천하의 지성(至誠)만이 자기 성을 다할 수 있다',
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅷ. 성의 형이상학')
def s_cheonjisam(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 성의 형이상학', page, total)
    add_title(slide, '與天地參 — 인간이 천지의 동반자가 된다')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '與 天 地 參',
                font_size=110, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                '여천지참 — "천지와 더불어 셋(三)이 된다"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.3), Inches(2.3), color=RULE, weight=1.5)
    add_filled_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.4), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.75), Inches(11.7), Inches(2.2), [
        ('중용의 절정 — 인간의 존재론적 위상 선언',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('천(天) — 하늘 · 지(地) — 땅 · 인(人) — 인간',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('인간은 천지의 피조물이 아니라 동반자(三)',
         {'font_size': 17, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('서양 종교의 인간관과 결정적으로 다른 지점',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('Ⅷ. 성의 형이상학')
def s_kichachigok(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 성의 형이상학', page, total)
    add_title(slide, '其次致曲 — 작은 것부터 성(誠)을 이룬다 (23장)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '"其次致曲 曲能有誠 誠則形…"',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
                '기차치곡 곡능유성 성즉형…',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    chain = [
        ('致 曲', '곡진히 함', '작은 것에 정성을 다함'),
        ('有 誠', '성이 있음', '진실함이 생김'),
        ('形',    '형',       '겉으로 드러남'),
        ('著',    '저',       '뚜렷해짐'),
        ('明',    '명',       '밝아짐'),
        ('動·變·化', '동·변·화', '움직이고 변하고 화함'),
    ]
    top = 3.9
    for han, eum, desc in chain:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.0), Inches(0.45), PALE)
        add_textbox(slide, Inches(0.6), Inches(top + 0.07), Inches(2.0), Inches(0.3),
                    han, font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.7), Inches(top), Inches(2.0), Inches(0.45),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(2.7), Inches(top + 0.08), Inches(2.0), Inches(0.3),
                    eum, font_size=12, color=SUB, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(4.8), Inches(top), Inches(8.0), Inches(0.45),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(5.0), Inches(top + 0.08), Inches(7.7), Inches(0.3),
                    desc, font_size=13, color=INK)
        top += 0.52


@S('Ⅷ. 성의 형이상학')
def s_seongja_jaseong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 성의 형이상학', page, total)
    add_title(slide, '誠者自成 — 성(誠)은 스스로 이루는 것 (25장)')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '誠 者 自 成 也',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                '성자자성야',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.0), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '"성은 스스로(自) 이루는 것이고',
                font_size=20, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.6),
                ' 도는 스스로 행하는 것이다"',
                font_size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.3), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(5.95), Inches(11.3), Inches(1.1), [
        ('"誠者 物之終始 不誠無物"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('성은 사물의 처음과 끝 — 성이 없으면 사물이 없다',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅷ. 성의 형이상학')
def s_jiseong_musik(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅷ. 성의 형이상학', page, total)
    add_title(slide, '至誠無息 — 지극한 성은 쉬지 않는다 (26장)')
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.9),
                '至 誠 無 息',
                font_size=104, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                '지성무식',
                font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
    chain = [
        ('無 息',  '쉬지 않음'),
        ('久',     '오래감'),
        ('徵',     '드러남'),
        ('悠 遠',  '유원(멀고 깊음)'),
        ('博 厚',  '박후(넓고 두터움)'),
        ('高 明',  '고명(높고 밝음)'),
    ]
    cur_x = 0.5
    box_w = 2.0
    arrow_w = 0.13
    for i, (han, kor) in enumerate(chain):
        add_filled_rect(slide, Inches(cur_x), Inches(4.7), Inches(box_w),
                        Inches(0.55), PALE)
        add_textbox(slide, Inches(cur_x), Inches(4.78), Inches(box_w), Inches(0.4),
                    han, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(cur_x), Inches(5.3), Inches(box_w),
                        Inches(0.5), RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(cur_x), Inches(5.38), Inches(box_w), Inches(0.4),
                    kor, font_size=11, color=INK, align=PP_ALIGN.CENTER)
        if i < len(chain) - 1:
            add_textbox(slide, Inches(cur_x + box_w), Inches(4.85), Inches(arrow_w),
                        Inches(0.4), '→', font_size=14, color=ACCENT,
                        align=PP_ALIGN.CENTER)
        cur_x += box_w + arrow_w
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
                '"쉬지 않으면 오래가고, 오래가면 드러나고, 드러나면 유원하고…"',
                font_size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                '"지속의 힘이 공간을 확장한다" — 우주를 움직이는 원리는 폭발이 아니라 지속',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅸ. 결론 ==============
@S('Ⅸ. 결론')
def s_sage_way(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 결론', page, total)
    add_title(slide, '성인(聖人)의 도 — 27~32장')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.43), Inches(12.3), Inches(0.55),
                '"大哉聖人之道! 洋洋乎發育萬物 峻極于天"',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
                '대재성인지도! 양양호발육만물 준극우천',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.8), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(3.95), Inches(11.7), Inches(1.6), [
        ('"위대하구나, 성인의 도!',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        (' 양양히 만물을 발육시키고 높이 하늘에 닿는다"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.3),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(5.95), Inches(11.7), Inches(1.1), [
        ('"優優大哉! 禮儀三百 威儀三千"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('우우대재 — 예의(禮儀) 300, 위의(威儀) 3,000 — 성인의 광대한 도(道)',
         {'font_size': 13, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅸ. 결론')
def s_cheondo_indo(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 결론', page, total)
    add_title(slide, '天道와 人道의 합치 — 왕도의 완성')
    add_filled_rect(slide, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(0.5),
                '天道 — 하늘의 도', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.9), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('誠者 天之道', {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 저절로 참됨', {'font_size': 14, 'space_before': 8}),
        ('• 쉬지 않음(無息)', {'font_size': 14, 'space_before': 4}),
        ('• 소리·냄새 없음(無聲無臭)', {'font_size': 14, 'space_before': 4}),
        ('• 만물을 화육함', {'font_size': 14, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(0.5),
                '人道 — 사람의 도', font_size=20, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.5), [
        ('誠之者 人之道', {'font_size': 18, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 8}),
        ('• 참되고자 노력', {'font_size': 14, 'space_before': 8}),
        ('• 선을 택해 굳게 붙듦', {'font_size': 14, 'space_before': 4}),
        ('• 박학·심문·신사·명변·독행', {'font_size': 14, 'space_before': 4}),
        ('• 人一能之 己百之', {'font_size': 14, 'space_before': 4}),
    ], line_spacing=1.3)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                '두 길이 한 곳에서 만난다 — "與天地參"의 인간 — 왕이 이르러야 할 경지',
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('Ⅸ. 결론')
def s_museong(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅸ. 결론', page, total)
    add_title(slide, '上天之載 無聲無臭 — 33장 결론, 중용의 대미')
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
                '上 天 之 載',
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.9),
                '無 聲 無 臭',
                font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.5),
                '상천지재 무성무취 — "지극하도다(至矣)!"',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.8), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
                '"상천(上天)의 일은 소리도 없고 냄새도 없다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.7), Inches(5.95), Inches(11.9), Inches(1.2), PALE)
    add_paragraphs(slide, Inches(1.0), Inches(6.1), Inches(11.3), Inches(1.0), [
        ('"가장 궁극의 도는 드러나지 않는다 — 요란한 것은 진짜가 아니다"',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('1장 "막현호은"과 33장 "무성무취" — 수미상관의 완성',
         {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅹ. 7대 주제어 ==============
@S('Ⅹ. 7대 주제어')
def s_keywords_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 7대 주제어', page, total)
    add_title(slide, '7대 주제어 ① — 中(중) · 庸(용)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.5),
                '中 — 중', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(5.9), Inches(1.0),
                '中', font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"그때 그 자리의 정답"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 치우치지 않음', {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 기하학적 중간이 아니라', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('  상황·관계·시대의 적정점', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 살아 움직이는 판단', {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.9), Inches(0.5),
                '庸 — 용', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.9), Inches(3.1), Inches(5.9), Inches(1.0),
                '庸', font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"평상시에도 흔들리지 않음"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 특별한 때의 지혜가 아니라', {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('  일상에서 지켜지는 도', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 매일의 식사·말·옷차림에서', {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('  이어지는 품격', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• "지속 가능한 덕"', {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅹ. 7대 주제어')
def s_keywords_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 7대 주제어', page, total)
    add_title(slide, '7대 주제어 ② — 性(성) · 誠(성)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.5),
                '性 — 성(本性)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(5.9), Inches(1.0),
                '性', font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"하늘이 준 본래의 나"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 天命之謂性', {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 후천적 학습이 아니라', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('  하늘에서 부여받음', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• "이미 있는 것의 발현"', {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.9), Inches(0.5),
                '誠 — 성(진실)', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.9), Inches(3.1), Inches(5.9), Inches(1.0),
                '誠', font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"거짓 없는 존재의 방식"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 誠者 天之道', {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 인간이 하늘과 연결되는 지점', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 자기 본성이 삶으로 흘러나옴', {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('• Authenticity의 동양 원형', {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅹ. 7대 주제어')
def s_keywords_3(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 7대 주제어', page, total)
    add_title(slide, '7대 주제어 ③ — 愼獨(신독) · 時中(시중)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.5),
                '愼 獨 — 신독', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(5.9), Inches(1.0),
                '愼獨', font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"아무도 보지 않을 때가 진짜"',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 莫見乎隱 (1장)', {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 드러난 행위가 아니라', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('  숨은 순간의 태도가 인격', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 왕권을 지키는 내적 장치', {'font_size': 13, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.9), Inches(0.5),
                '時 中 — 시중', font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.9), Inches(3.1), Inches(5.9), Inches(1.0),
                '時中', font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"때를 아는 판단"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 君子而時中 (2장)', {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 중은 때에 따라 달라진다', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 매번 다시 발견하는 중', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 리더의 판단력', {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


@S('Ⅹ. 7대 주제어')
def s_keywords_4(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅹ. 7대 주제어', page, total)
    add_title(slide, '7대 주제어 ④ — 至誠無息(지성무식) · 中和(중화)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(5.9), Inches(4.5), PALE)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.5),
                '至誠無息 — 지속의 힘', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(5.9), Inches(1.0),
                '無息', font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.8), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"멈추지 않음"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 26장 — 至誠無息', {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 짧은 열정이 아니라', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('  긴 지속이 우주의 원리', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 人一能之 己百之', {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(6.9), Inches(2.3), Inches(5.9), Inches(4.5),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_textbox(slide, Inches(6.9), Inches(2.5), Inches(5.9), Inches(0.5),
                '中和 — 감정의 절도', font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.9), Inches(3.1), Inches(5.9), Inches(1.0),
                '中和', font_size=58, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.2), Inches(4.4), Inches(5.3), Inches(2.4), [
        ('"감정 통치의 철학"',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 6}),
        ('• 미발의 中 + 이발의 和', {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('• 억누름이 아니라', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('  절도(節)에 맞게 발함', {'font_size': 13, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('• 致中和 → 천지위·만물육', {'font_size': 13, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)


# ============== Ⅺ. 명구절 ==============
def make_quote_slide(section, hanmun, eum, mean, ref, *, hanmun_size=42):
    def renderer(slide, page, total):
        set_white_background(slide)
        add_page_header(slide, section, page, total)
        add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4),
                    hanmun, font_size=hanmun_size, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                    eum, font_size=20, color=SUB, align=PP_ALIGN.CENTER)
        add_rule(slide, Inches(5.5), Inches(4.4), Inches(2.3), color=RULE, weight=1.5)
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                    mean, font_size=22, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    f'— {ref}', font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    return renderer


SLIDES.append((make_quote_slide('Ⅺ. 명구절 (1/10)',
    '天 命 之 謂 性\n率 性 之 謂 道   修 道 之 謂 敎',
    '천명지위성 · 솔성지위도 수도지위교',
    '하늘이 명한 것을 성이라, 성을 따르는 것을 도라, 도를 닦는 것을 교라 한다',
    '1장 (강령)', hanmun_size=26), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (2/10)',
    '致 中 和   天 地 位 焉   萬 物 育 焉',
    '치중화 · 천지위언 만물육언',
    '중화를 이루면 천지가 제자리를 잡고 만물이 자라난다',
    '1장', hanmun_size=28), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (3/10)',
    '君 子 之 中 庸 也   君 子 而 時 中',
    '군자지중용야 · 군자이시중',
    '군자의 중용이란 군자이면서 때에 맞게 중을 쓰는 것이다',
    '2장 — 시중(時中)', hanmun_size=28), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (4/10)',
    '道 不 遠 人',
    '도 불 원 인',
    '도는 사람에게서 멀지 않다 — 일상의 관계 안에 있다',
    '13장', hanmun_size=92), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (5/10)',
    '執 其 兩 端   用 其 中 於 民',
    '집기양단 용기중어민',
    '양쪽 끝을 잡아 그 중(中)을 백성에게 쓴다 — 순(舜) 임금의 지혜',
    '6장', hanmun_size=32), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (6/10)',
    '誠 者   天 之 道 也\n誠 之 者   人 之 道 也',
    '성자 천지도야 · 성지자 인지도야',
    '성(誠)은 하늘의 도이고, 성실하고자 함은 사람의 도다',
    '20장 — 성(誠)의 결정적 선언', hanmun_size=26), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (7/10)',
    '博 學 之   審 問 之\n愼 思 之   明 辨 之   篤 行 之',
    '박학지 심문지 · 신사지 명변지 독행지',
    '널리 배우고, 자세히 묻고, 신중히 생각하고, 밝게 분별하고, 도탑게 행하라',
    '20장 — 학습의 5단계', hanmun_size=22), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (8/10)',
    '人 一 能 之   己 百 之\n人 十 能 之   己 千 之',
    '인일능지 기백지 · 인십능지 기천지',
    '남이 한 번에 능하면 나는 백 번을, 남이 열 번에 능하면 나는 천 번을 한다',
    '20장 — 그릿(Grit)의 동양 원형', hanmun_size=22), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (9/10)',
    '至 誠 無 息',
    '지 성 무 식',
    '지극한 성(誠)은 쉼이 없다 — 지속의 힘이 우주의 원리',
    '26장', hanmun_size=110), 'Ⅺ. 명구절'))

SLIDES.append((make_quote_slide('Ⅺ. 명구절 (10/10)',
    '上 天 之 載   無 聲 無 臭',
    '상천지재 무성무취',
    '하늘의 일은 소리도 없고 냄새도 없다 — "지극하도다!"',
    '33장 — 결론', hanmun_size=38), 'Ⅺ. 명구절'))


# ============== Ⅻ. 7대 메시지 ==============
@S('Ⅻ. 7대 메시지')
def s_messages_1(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 7대 메시지', page, total)
    add_title(slide, '7대 핵심 메시지 ① ~ ②')
    items = [
        ('1', '당신의 본성은 하늘이 내린 선(善)이다',
         '천명지위성 — 인간은 결핍이 아니라 이미 하늘의 선을 받고 태어났다'),
        ('2', '중(中)은 정답이 아니라 과정이다',
         '계산으로 나오는 답이 아니라 매 상황에서 새로 발견해야 하는 한 지점'),
    ]
    top = 2.7
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(1.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(1.0), Inches(0.5),
                    num, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.8), Inches(top), Inches(11.0), Inches(1.5), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.3), Inches(10.6), Inches(0.5),
                    title, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.0), Inches(top + 0.9), Inches(10.6), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 1.7


@S('Ⅻ. 7대 메시지')
def s_messages_2(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 7대 메시지', page, total)
    add_title(slide, '7대 핵심 메시지 ③ ~ ⑤')
    items = [
        ('3', '감정을 억누르지 말고 절도에 맞춰라',
         '중화(中和) — 분노를 지우지 말고 옳은 분노를 정확한 대상에게'),
        ('4', '남이 안 보이는 곳이 당신이다',
         '신독(愼獨) — 무대 위가 아니라 분장실의 모습이 진짜 당신'),
        ('5', '진정성(誠)은 모든 것의 근본이다',
         '誠者物之終始 不誠無物 — 성이 없으면 사물 자체가 없다'),
    ]
    top = 2.5
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(1.3), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.42), Inches(1.0), Inches(0.5),
                    num, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.8), Inches(top), Inches(11.0), Inches(1.3), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.22), Inches(10.6), Inches(0.5),
                    title, font_size=19, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.0), Inches(top + 0.75), Inches(10.6), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 1.5


@S('Ⅻ. 7대 메시지')
def s_messages_3(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'Ⅻ. 7대 메시지', page, total)
    add_title(slide, '7대 핵심 메시지 ⑥ ~ ⑦')
    items = [
        ('6', '지속이 강함이다',
         '至誠無息 — 한 번의 천재성보다 천 번의 반복이 강하다'),
        ('7', '가장 지극한 도는 드러나지 않는다',
         '上天之載 無聲無臭 — 시끄러운 선은 얕고, 조용한 덕이 깊다'),
    ]
    top = 2.7
    for num, title, desc in items:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(1.0), Inches(1.5), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.5), Inches(1.0), Inches(0.5),
                    num, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(1.8), Inches(top), Inches(11.0), Inches(1.5), PALE)
        add_textbox(slide, Inches(2.0), Inches(top + 0.3), Inches(10.6), Inches(0.5),
                    title, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(2.0), Inches(top + 0.9), Inches(10.6), Inches(0.5),
                    desc, font_size=14, color=INK)
        top += 1.7


# ============== ⅩⅢ. 현대 의의 ==============
@S('ⅩⅢ. 현대 의의')
def s_modern_sijung(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅢ. 현대 의의', page, total)
    add_title(slide, '현대 ① — 극단의 시대에 "시중(時中)"을 다시 배우다')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('극단화(Polarization) 시대',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('SNS · 정치 · 문화 모든 영역이 양극단으로 치닫는 시대',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"중도"는 무기력한 타협으로 오해받는다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('중용의 답 — "執其兩端, 用其中"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"중간에서 무난하게 살자"가 아니라',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('"양쪽 끝을 다 이해한 뒤 이 상황의 최선을 찾아라"',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('= 가장 치밀한 사고력을 요구하는 진짜 지성의 훈련',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('ⅩⅢ. 현대 의의')
def s_modern_authenticity(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅢ. 현대 의의', page, total)
    add_title(slide, '현대 ② — 진정성(Authenticity)의 동양적 뿌리 "성(誠)"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('현대 리더십·심리학이 재발견한 가치 — Authenticity',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('내면의 느낌과 외면의 행동',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('혼자일 때의 나와 남들 앞의 나가 일치할 때',
         {'font_size': 15, 'bold': True, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('중용은 2,500년 전에 이미 선언',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"성(誠)은 하늘의 도이다(誠者天之道)"',
         {'font_size': 17, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 4}),
        ('일관성·진정성·신뢰는 21세기의 가장 희소한 자원',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('기업·브랜드·개인 평판의 축 = 중용의 성(誠)',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.25)


@S('ⅩⅢ. 현대 의의')
def s_modern_eq(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅢ. 현대 의의', page, total)
    add_title(slide, '현대 ③ — 감정 관리(EQ)의 최고 수준 "중화(中和)"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"감정을 억누르라"가 아니라 "절도에 맞게 표현하라"',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('분노할 자리에 분노하고, 기뻐할 자리에 기뻐하되',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('과하지도 모자라지도 않게',
         {'font_size': 16, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('다니엘 골먼 "감성지능(EQ)" — 정확히 이 주제',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('• 조직 리더십 · 협상 · 부부 관계 · 자녀 양육의 공통 기술',
         {'font_size': 14, 'space_before': 6}),
        ('• 중용의 "중절(中節)" = 현대의 emotion regulation 이론',
         {'font_size': 14, 'space_before': 4}),
        ('• 편도체 과잉 활성의 절제 = "발이개중절(發而皆中節)"',
         {'font_size': 14, 'bold': True, 'color': ACCENT, 'space_before': 4}),
    ], line_spacing=1.3)


@S('ⅩⅢ. 현대 의의')
def s_modern_metacog(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅢ. 현대 의의', page, total)
    add_title(slide, '현대 ④ — 메타인지(Metacognition)의 고전 "신독(愼獨)"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('"자기를 감찰하는 자기"',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('AI 시대의 생존 능력 = 메타인지(Metacognition)',
         {'font_size': 16, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('= 중용의 "신독(愼獨)"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                    RGBColor(0xFA, 0xFA, 0xFA))
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('SNS·원격 근무·개인 사업의 시대',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"아무도 나를 감독하지 않는 시간"이 길어졌다',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('이 시간 동안 자기 감찰이 없으면 누구도 나를 지켜줄 수 없다',
         {'font_size': 14, 'color': SUB, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('신독 = 현대인의 가장 절실한 습관',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


@S('ⅩⅢ. 현대 의의')
def s_modern_grit(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅢ. 현대 의의', page, total)
    add_title(slide, '현대 ⑤ — Grit(끈기)의 원형 "人一能之 己百之"')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.7), [
        ('Angela Duckworth의 그릿(Grit) 이론',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"열정과 끈기가 재능보다 성취를 결정한다"',
         {'font_size': 15, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('= 중용 20장의 "人一能之 己百之"의 실증 연구',
         {'font_size': 16, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.3)
    add_filled_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), PALE)
    add_paragraphs(slide, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), [
        ('"果能此道矣 雖愚必明 雖柔必強"',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('"진실로 이 도를 행하면 어리석어도 반드시 밝아지고',
         {'font_size': 14, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        (' 유약해도 반드시 강해진다"',
         {'font_size': 14, 'bold': True, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('1만 시간의 법칙·일본 직인(職人) 정신 — 모두 이 한 구절의 변주',
         {'font_size': 13, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
    ], line_spacing=1.25)


# ============== ⅩⅣ. 비교 ==============
@S('ⅩⅣ. 비교')
def s_compare_daehak(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅣ. 비교', page, total)
    add_title(slide, '대학 vs 중용 — 사서의 두 축')
    rows = [
        ('성격',       '"무엇을 어떻게 할 것인가"',   '"왜 그래야 하는가"'),
        ('비유',       '사업계획서',                  '존재론'),
        ('주제',       '수기치인(修己治人)',          '천명·성·도·교'),
        ('출발점',     '명명덕 — 내 안의 등불',       '천명지위성 — 하늘의 명'),
        ('핵심 개념',  '삼강령 + 팔조목',              '중·용·성·신독·시중'),
        ('주희의 평',  '規模 (전체의 틀)',             '微妙 (가장 심오)'),
        ('읽는 순서',  '사서의 첫 번째',               '사서의 마지막'),
    ]
    top = 2.0
    add_filled_rect(slide, Inches(0.5), Inches(top), Inches(2.2), Inches(0.55), SUB)
    add_textbox(slide, Inches(0.5), Inches(top + 0.1), Inches(2.2), Inches(0.4),
                '항목', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(2.85), Inches(top), Inches(5.0), Inches(0.55), SUB)
    add_textbox(slide, Inches(2.85), Inches(top + 0.1), Inches(5.0), Inches(0.4),
                '대학(大學)', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(7.9), Inches(top), Inches(4.9), Inches(0.55), ACCENT)
    add_textbox(slide, Inches(7.9), Inches(top + 0.1), Inches(4.9), Inches(0.4),
                '중용(中庸)', font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    row_h = 0.62
    for r_idx, row in enumerate(rows):
        y = top + 0.55 + r_idx * row_h
        bg = RGBColor(0xFA, 0xFA, 0xFA) if r_idx % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.5), Inches(y), Inches(2.2), Inches(row_h), PALE)
        add_textbox(slide, Inches(0.55), Inches(y + 0.15), Inches(2.1), Inches(0.4),
                    row[0], font_size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(2.85), Inches(y), Inches(5.0), Inches(row_h), bg)
        add_textbox(slide, Inches(2.95), Inches(y + 0.15), Inches(4.8), Inches(0.4),
                    row[1], font_size=13, color=INK, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(7.9), Inches(y), Inches(4.9), Inches(row_h), bg)
        add_textbox(slide, Inches(8.0), Inches(y + 0.15), Inches(4.7), Inches(0.4),
                    row[2], font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


@S('ⅩⅣ. 비교')
def s_compare_modern(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅣ. 비교', page, total)
    add_title(slide, '중용과 현대 도서·이론')
    rows = [
        ('Daniel Goleman',   'Emotional Intelligence',     '중화(中和)·중절(中節)의 현대판'),
        ('Angela Duckworth', 'Grit (그릿)',                 '人一能之 己百之의 실증 연구'),
        ('Stephen Covey',    '7 Habits of Highly Effective People',
         '신독·성의·지성무식의 현대 경영학 버전'),
        ('Cal Newport',      'Deep Work',                   '"신사(愼思)"의 21세기 재발견'),
        ('Brené Brown',      'Daring Greatly',              '성(誠, authenticity)의 현대 심리학'),
    ]
    top = 2.4
    for author, book, link in rows:
        add_filled_rect(slide, Inches(0.6), Inches(top), Inches(2.5), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.6), Inches(top + 0.22), Inches(2.5), Inches(0.5),
                    author, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_filled_rect(slide, Inches(3.2), Inches(top), Inches(3.8), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.35), Inches(top + 0.22), Inches(3.6), Inches(0.5),
                    book, font_size=13, bold=True, color=ACCENT)
        add_filled_rect(slide, Inches(7.1), Inches(top), Inches(5.7), Inches(0.85),
                        RGBColor(0xFA, 0xFA, 0xFA))
        add_textbox(slide, Inches(7.25), Inches(top + 0.22), Inches(5.4), Inches(0.5),
                    link, font_size=13, color=INK)
        top += 0.97


# ============== ⅩⅤ. 마무리 ==============
@S('ⅩⅤ. 마무리')
def s_one_page(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅤ. 마무리', page, total)
    add_title(slide, '한 장으로 보는 중용')
    add_filled_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.15), Inches(12.3), Inches(0.5),
                '1부 · 1장 — 강령',
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.7),
                '天命之謂性 · 中和 · 致中和 天地位焉 萬物育焉',
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.7),
                    RGBColor(0xA0, 0x40, 0x40))
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.5),
                '2부 · 2~20장 — 실천 · 정치',
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.6),
                '君子時中 · 執其兩端 · 道不遠人',
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5),
                '20장 — 九經 · 博學審問慎思明辨篤行 · 人一能之己百之 · 誠者天之道',
                font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5),
                    RGBColor(0x70, 0x40, 0x60))
    add_textbox(slide, Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.5),
                '3부 · 21~33장 — 성(誠)의 형이상학',
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
                '至誠盡性 · 與天地參 · 至誠無息',
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                '33장 결론 — 上天之載 無聲無臭 · 至矣',
                font_size=14, color=WHITE, align=PP_ALIGN.CENTER)


@S('ⅩⅤ. 마무리')
def s_summary(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅤ. 마무리', page, total)
    add_title(slide, '한 문장으로 정리하는 중용')
    add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(11.9), Inches(5.0), PALE)
    add_paragraphs(slide, Inches(1.1), Inches(2.2), Inches(11.1), Inches(4.7), [
        ('"당신의 본성은 하늘이 주었다(天命之謂性)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        (' 그 본성을 거스르지 않고',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        (' 매 상황마다 한쪽으로 치우치지 않으며(時中)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        (' 아무도 보지 않을 때도 스스로를 삼가고(愼獨)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        (' 내면이 외면과 어긋나지 않게 참되라(誠)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 4}),
        ('', {'font_size': 4}),
        (' 이를 멈추지 말라(無息)"',
         {'font_size': 20, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 8}),
        ('', {'font_size': 6}),
        ('당신의 덕이 하늘에 이를 때, 그것은 소리도 냄새도 없이',
         {'font_size': 15, 'color': SUB, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('이미 세상을 움직이고 있을 것이다(無聲無臭)',
         {'font_size': 15, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
    ], line_spacing=1.2)


@S('ⅩⅤ. 마무리')
def s_compass(slide, page, total):
    set_white_background(slide)
    add_page_header(slide, 'ⅩⅤ. 마무리', page, total)
    add_title(slide, '맺음 — 나침반의 자침(磁針)')
    add_filled_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.6),
                    RGBColor(0xFA, 0xE5, 0xE5))
    add_paragraphs(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(4.3), [
        ('주희가 중용을 사서의 마지막에 두었던 이유',
         {'font_size': 19, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER}),
        ('', {'font_size': 4}),
        ('대학이 지도(地圖)라면',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 10}),
        ('중용은 그 지도가 가리키는 나침반의 자침(磁針)',
         {'font_size': 20, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 8}),
        ('어디를 가든',
         {'font_size': 16, 'color': INK, 'align': PP_ALIGN.CENTER, 'space_before': 14}),
        ('이 자침 하나만 있으면',
         {'font_size': 18, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER}),
        ('길을 잃지 않는다',
         {'font_size': 22, 'bold': True, 'color': ACCENT, 'align': PP_ALIGN.CENTER, 'space_before': 6}),
        ('', {'font_size': 6}),
        ('— 극단의 시대에 진정성으로 자기 길을 지속하는 사람 —',
         {'font_size': 14, 'color': SUB, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 12}),
    ], line_spacing=1.25)


@S('ⅩⅤ. 마무리')
def s_closing(slide, page, total):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), RULE)
    add_textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.5),
                '無 聲 無 臭',
                font_size=120, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                '무 성 무 취', font_size=24, color=SUB, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=RULE, weight=1.5)
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.7),
                '"하늘의 일은 소리도 없고 냄새도 없다"',
                font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '— 33장 결론 · 중용 전체의 대미 ·  "지극하도다(至矣)!"',
                font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '정말 깊은 것은 조용하다 — 자랑할 필요 없는 것이 진짜다',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
                '감사합니다', font_size=26, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), RULE)


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)

out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\중용.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')