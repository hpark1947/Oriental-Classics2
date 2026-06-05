# -*- coding: utf-8 -*-
"""
맹자(孟子) 발표자료 — 전면 보강판 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
7편 14장 각 1장 깊이 읽기 · 맹자 생애·맹모 일화 · 사상 7기둥 · 후대 영향 망라
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '아성(亞聖)의 책 · 공자의 인(仁)을 체계화한 사서(四書)의 둘째',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '孟 子',
                font_size=120, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '맹 자',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
                '民 爲 貴  社 稷 次 之  君 爲 輕  — 백성이 가장 귀하다',
                font_size=17, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '맹가(孟軻, BC 372?~289?) · 전국시대 중기 · 7편 14장 · 약 35,000자',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '성선설 · 사단 · 왕도정치 · 호연지기 · 민본 · 인의(仁義)',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 맹자란 무엇인가'),
        ('Ⅱ.', '맹자의 생애 — 어머니와 한 철학자'),
        ('Ⅲ.', '제자와 논적들 · 사기 열전'),
        ('Ⅳ.', '7편 14장의 구조 · 오륜(五倫)'),
        ('Ⅴ.', '7편 14장 깊이 읽기'),
        ('Ⅵ.', '핵심 사상 9기둥'),
    ]
    items_right = [
        ('Ⅶ.', '맹자의 유명한 비유들'),
        ('Ⅷ.', '정전제 · 항산항심의 경제'),
        ('Ⅸ.', '명구 18선'),
        ('Ⅹ.', '동아시아 수용 · 사단칠정 논쟁'),
        ('Ⅺ.', '오늘 우리에게 · 마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 맹자')
    rows = [
        ('서명',  '맹자(孟子) — 「맹 선생의 책」'),
        ('저자',  '맹가(孟軻, BC 372?~289?) · 자(字) 자여(子輿) · 제자 만장·공손추와 공저'),
        ('시대',  '전국시대 중기 (BC 4세기) — 양혜왕·제선왕·등문공 시대'),
        ('분량',  '7편 14장 (상·하) · 260여 장 · 약 35,000자'),
        ('학파',  '유가(儒家) — 공자의 적통(嫡統)'),
        ('도통',  '공자 → 증자 → 자사 → 맹자 (한유 정립)'),
        ('위상',  '사서(四書)의 둘째 — 「공자의 인을 체계화한 책」'),
        ('존칭',  '아성(亞聖) — 공자 다음가는 성인 (송 주희 확립)'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.1), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_systematic(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '왜 맹자인가 — 유교의 철학적 체계화',
              '논어가 단편적 어록이라면, 맹자는 체계적 논변서')
    cols = [
        ('論 語', '논어',
         '단편적 어록\n질문 → 짧은 답\n\n「述而不作」\n— 전할 뿐 짓지 않음\n\n소박하고 실천적',
         INK),
        ('孟 子', '맹자',
         '체계적 논변서\n긴 대화 → 비유 → 논증\n\n양주·묵자·고자 비판\n\n철학적 체계화',
         ACCENT),
    ]
    for i, (han, kor, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    kor, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.5), Inches(3.6), Inches(4.9), Inches(3.5),
                       [(body, {'font_size': 15, 'color': INK})], line_spacing=1.5)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '맹자가 체계화한 것 — 성선설 · 사단 · 왕도 vs 패도 · 호연지기 · 인의(仁義) 병칭',
                font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


@S(SEC1)
def i_canon(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '경전이 되기까지 — 1,500년의 여정',
              '제자백가의 하나에서 사서(四書)의 둘째까지')
    rows = [
        ('전국시대 (BC 3C)', '제자백가 중 하나',  '묵가·법가·도가와 같은 급'),
        ('한(漢) (BC 2C~AD 2C)', '조기(趙岐)가 주석',  '경전화 시작'),
        ('당(唐) 8C', '한유(韓愈) 도통론',     '「공자 → 증자 → 자사 → 맹자」의 적통 선언'),
        ('북송 11C', '왕안석이 과거에 포함',    '준경전 지위'),
        ('남송 12C', '주희 『사서집주』',         '사서(四書)로 확정 · 아성(亞聖)으로 격상'),
        ('원(元) 14C', '과거 시험 필수',         '동아시아 지식인의 필독서'),
    ]
    for i, (era, who, what) in enumerate(rows):
        y = Inches(2.4 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.4), Inches(0.55),
                    era, font_size=13, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.3), y, Inches(3.8), Inches(0.55),
                    who, font_size=14, color=INK, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.3), y, Inches(5.6), Inches(0.55),
                    what, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
                '맹자는 1,500년을 기다려 경전이 된 책 — 송 주희에 의해 「아성(亞聖)」으로 최종 확립',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC1)
def i_thought(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '핵심 사상 한 폭으로')
    boxes = [
        ('性 善 說', '성선설', '인간 본성은 선하다 — 물이 아래로 흐르듯'),
        ('四 端', '사단', '측은·수오·사양·시비 — 사람의 네 가지 마음'),
        ('王 道', '왕도', '덕(德)으로 다스림 — 패도(力)의 반대'),
        ('民 本', '민본', '백성이 가장 귀하다 — 民貴君輕'),
        ('浩 然 之 氣', '호연지기', '정직으로 기른 천지 가득한 기운'),
        ('仁 義', '인의', '공자의 인(仁) → 인의(仁義) 병칭으로 확장'),
    ]
    for i, (han, kor, desc) in enumerate(boxes):
        col, row = i % 3, i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(3.9), Inches(2.0), PALE)
        add_textbox(slide, x, y + Inches(0.2), Inches(3.9), Inches(0.7),
                    han, font_size=22, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(0.95), Inches(3.9), Inches(0.4),
                    kor, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.5), Inches(0.5),
                    desc, font_size=12, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 맹자의 생애 ==============
SEC2 = 'Ⅱ. 맹자의 생애'

@S(SEC2)
def ii_profile(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '맹가(孟軻) — 한 어머니가 키운 철학자',
              '추(鄒)나라에서 태어나 천하를 떠돌다')
    rows = [
        ('이름',     '맹가(孟軻), 자(字) 자여(子輿) 또는 자거(子車)'),
        ('출생',     'BC 372년경 추(鄒)나라 (현 산동성 추성시 — 공자 고향 곡부와 인접)'),
        ('사망',     'BC 289년경 (향년 약 83세)'),
        ('부친',     '맹자 3세경 일찍 사망 — 홀어머니 밑에서 성장'),
        ('모친',     '장씨(仉氏) — 동양 최고의 「교육 어머니」'),
        ('스승',     '자사(子思)의 문인에게 수학 — 「사숙(私淑)」'),
        ('활동',     '40여 년 천하 유세 — 양혜왕·제선왕·등문공 등 군주 설득'),
        ('말년',     '70세 전후 귀향 → 만장·공손추 등 제자와 7편 저술'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.0), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.9), y, Inches(10.1), Inches(0.5),
                    v, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC2)
def ii_mengmu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '맹모(孟母) 두 일화 — 동양 교육사의 전설',
              '맹모삼천지교 · 맹모단기지교')
    cols = [
        ('孟 母 三 遷', '맹모삼천지교',
         '묘지 → 시장 → 서당 근처\n세 번 이사한 어머니\n\n「이제 비로소 아이를\n기를 곳을 얻었구나」\n\n환경이 사람을 만든다.\n학군의 동양적 원조.\n루소 「에밀」보다 2,100년 앞선\n환경 교육론.'),
        ('孟 母 斷 機', '맹모단기지교',
         '학업 중 집에 돌아온 맹자.\n어머니가 베틀 실을 끊음\n\n「짜던 베를 끊으면\n쓸모없는 천이 되듯,\n학문을 중도에 그만두면\n아무것도 이루지 못한다」\n\n포기하지 말라 — 그릿(grit) 이론의\n동양적 원체험.'),
    ]
    for i, (han, kor, body) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), ACCENT)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=28, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    kor, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.4), Inches(3.6), Inches(5.1), Inches(3.5),
                       [(body, {'font_size': 13, 'color': INK})], line_spacing=1.5)


@S(SEC2)
def ii_tour(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '40년의 천하 유세 — 왕도(王道)를 설파하다',
              '공자의 14년 천하주유에 필적하는 긴 설득 여정')
    rows = [
        ('梁(魏) 양(위)', '양혜왕(梁惠王)',
         '맹자의 첫 만남 — 「왕께서 하필 이익을 말씀하십니까? 오직 인의(仁義)만 있을 뿐입니다」 / 등용 실패'),
        ('齊 제',         '제선왕(齊宣王)',
         '직하학궁의 군주 · 객경(客卿) 대우 / 「여민동락(與民同樂)」의 대화 — 군주의 약점을 왕도의 도구로 전환'),
        ('滕 등',         '등문공(滕文公)',
         '소국 군주 · 정전제(井田制) 자문 / 「소국이라도 왕도를 실현할 수 있다」는 맹자 주장의 실험장'),
        ('宋 송 · 鄒 추',  '여러 군주',
         '짧은 방문과 고향 재방문 — 어느 제후도 맹자의 도를 채택하지 않음'),
    ]
    for i, (country, ruler, story) in enumerate(rows):
        y = Inches(2.4 + i * 1.0)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(2.4), Inches(0.85),
                    country, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, Inches(3.3), y, Inches(2.2), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.3), y, Inches(2.2), Inches(0.85),
                    ruler, font_size=14, color=ACCENT, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.7), y + Inches(0.05), Inches(7.2), Inches(0.8),
                    story, font_size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
                '「500년 만에 반드시 왕자(王者)가 나타난다 — 나 이외에 누가 이 일을 하겠는가?」 — 맹자의 사명감',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅲ. 제자와 논적들 ==============
SEC3 = 'Ⅲ. 제자와 논적들'

@S(SEC3)
def iii_disciples(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '핵심 제자 — 만장·공손추와 책의 편집자들')
    rows = [
        ('萬章 만장',  '제5편의 주인공',
         '맹자의 가장 중요한 제자 · 고대 성현(요·순·우·탕)의 행적에 관한 질문자\n「요임금이 천하를 순에게 준 것이 사실입니까?」 → 「天與之 — 하늘이 준 것이다」 / 천명관(天命觀)의 핵심'),
        ('公孫丑 공손추', '제2편의 주인공',
         '제(齊)나라 출신 · 호연지기·부동심(不動心)에 대한 긴 문답의 상대\n「선생님은 어떻게 그처럼 부동심이십니까?」 → 「我善養吾浩然之氣」 — 호연지기 대화'),
        ('樂正子 악정자', '정치 실천의 제자',
         '등(滕)나라에서 관직 / 맹자의 정치 이상을 실제로 시도한 제자'),
        ('桃應·陳臻 등', '윤리·처세 질문자',
         '도응 — 「순 임금의 아버지가 살인하면?」이라는 유명 윤리 딜레마 제기 / 진진 — 맹자의 처세에 관한 질문'),
    ]
    for i, (name, label, body) in enumerate(rows):
        y = Inches(2.3 + i * 1.15)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(1.0), ACCENT)
        add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(2.4), Inches(0.5),
                    name, font_size=17, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), y + Inches(0.55), Inches(2.4), Inches(0.4),
                    label, font_size=12, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(3.3), y + Inches(0.1), Inches(9.6), Inches(0.9),
                       [(body, {'font_size': 12, 'color': INK})], line_spacing=1.4)


@S(SEC3)
def iii_simaqian(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '사기(史記)가 본 맹자 — 사마천의 「맹자열전」',
              '냉정한 역사가가 그린 한 좌절한 사상가의 초상')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 『사기·맹자순경열전』 — 맹자와 순자를 한데 묶어 다룸',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 사마천의 평 — 「迂遠而闊於事情」 — 우원(迂遠)하고 일의 정세에 어둡다',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
        ('     · 당시 군주의 눈으로 본 맹자 — 「현실 정치엔 통하지 않는 이상주의자」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 그러나 사마천 자신은 맹자를 깊이 존경 — 「余讀孟子書 …  廢書而嘆」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     · 「맹자 책을 읽으며 책을 덮고 탄식하지 않을 수 없었다」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 묘사한 맹자의 모습 — 「자사 문하에서 배우고 도가 통했으나 시대가 받아들이지 않음」',
         {'font_size': 17, 'space_before': 12}),
        ('● 「만장 등과 함께 시·서·중니의 뜻을 풀이하여 맹자 7편을 지었다」',
         {'font_size': 16, 'space_before': 10, 'color': SUB, 'bold': True}),
        ('● 좌절한 사상가가 7편의 책으로 1,500년 후 「아성」이 되는 역설 — 사기가 남긴 가장 깊은 메시지',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC3)
def iii_opponents(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '맹자의 논적들 — 「양묵을 막는 것이 내 사명」',
              '고자·묵가·양주·허행 등 당대 학파와의 정면 충돌')
    rows = [
        ('告子 고자',  '인성론 최대 논적',
         '성무선악설(性無善惡說) — 「본성은 선도 악도 아니다」\n물이 동·서로 흐르듯 외부 환경이 결정 / 맹자의 「성선설」 정립의 계기'),
        ('墨家 묵가',  '겸애(兼愛) 논쟁',
         '묵자 — 「모든 사람을 똑같이 사랑하라」\n맹자의 비판 — 「아비를 무시하는 것이다(無父) · 금수와 다를 바 없다!」'),
        ('楊朱 양주',  '극단적 이기주의',
         '「拔一毛利天下 而不爲也」 — 천하를 위해 머리털 하나 뽑지 않겠다\n맹자의 비판 — 「임금이 없는 것이다(無君) · 금수와 다를 바 없다!」'),
        ('許行 허행',  '농가 논쟁',
         '「군주도 직접 농사를 지어야 한다」\n맹자의 답 — 「백공(百工)의 일은 본래 분업」 / 애덤 스미스보다 2,000년 앞선 분업론'),
    ]
    for i, (name, label, body) in enumerate(rows):
        y = Inches(2.3 + i * 1.15)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(1.0), INK)
        add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(2.4), Inches(0.5),
                    name, font_size=17, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), y + Inches(0.55), Inches(2.4), Inches(0.4),
                    label, font_size=12, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(3.3), y + Inches(0.1), Inches(9.6), Inches(0.9),
                       [(body, {'font_size': 12, 'color': INK})], line_spacing=1.4)


# ============== Ⅳ. 7편 14장 구조 ==============
SEC4 = 'Ⅳ. 7편 14장의 구조'

@S(SEC4)
def iv_table(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '7편 14장 한 폭으로 — 각 편은 상·하로 나뉨',
              '편명은 대부분 첫 등장 인물에서 따옴')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.45), INK)
    headers = [('편', 0.7), ('편명', 3.5), ('상·하', 1.3), ('핵심 주제', 6.5)]
    x = Inches(0.7)
    for label, w in headers:
        add_textbox(slide, x, Inches(2.2), Inches(w), Inches(0.45),
                    label, font_size=13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(w)
    rows = [
        ('1', '梁惠王 양혜왕', '上·下', '왕도와 패도 · 「왕하필왈리」 · 여민동락'),
        ('2', '公孫丑 공손추', '上·下', '호연지기 · 부동심 · 사단(四端)'),
        ('3', '滕文公 등문공', '上·下', '정전제 · 분업론 · 「대장부론」'),
        ('4', '離婁 이루', '上·下', '예와 인 · 자기 반성 · 「반구저기」'),
        ('5', '萬章 만장', '上·下', '요순의 천하 양도 · 천명론'),
        ('6', '告子 고자', '上·下', '성선설 vs 성무선악 · 인성 논쟁'),
        ('7', '盡心 진심', '上·下', '진심·지성·지천 · 민귀군경 · 마무리'),
    ]
    for i, (no, name, ud, desc) in enumerate(rows):
        y = Inches(2.7 + i * 0.6)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.55), bg)
        add_textbox(slide, Inches(0.7), y, Inches(0.7), Inches(0.55),
                    no, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.4), y, Inches(3.5), Inches(0.55),
                    name, font_size=15, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.9), y, Inches(1.3), Inches(0.55),
                    ud, font_size=13, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.3), y, Inches(6.4), Inches(0.55),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_olun(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '오륜(五倫) — 맹자가 정립한 동양 인간관계의 다섯 기둥',
              '등문공 상편 — 「人之有道也  飽食煖衣 逸居而無敎  則近於禽獸」')
    rows = [
        ('父 子 有 親', '부자유친', '부모와 자식 사이엔 친(親)함이 있다',  '효(孝)의 원리'),
        ('君 臣 有 義', '군신유의', '임금과 신하 사이엔 의(義)가 있다',     '충(忠)과 의의 통합'),
        ('夫 婦 有 別', '부부유별', '부부 사이엔 분별(別)이 있다',         '역할의 구분과 존중'),
        ('長 幼 有 序', '장유유서', '어른과 아이 사이엔 차례(序)가 있다',   '제(悌)의 원리'),
        ('朋 友 有 信', '붕우유신', '벗 사이엔 믿음(信)이 있다',           '신(信)의 원리'),
    ]
    for i, (han, kor, body, label) in enumerate(rows):
        y = Inches(2.3 + i * 0.85)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.7), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.7),
                    han, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, Inches(3.5), y, Inches(1.8), Inches(0.7), PALE)
        add_textbox(slide, Inches(3.5), y, Inches(1.8), Inches(0.7),
                    kor, font_size=14, color=SUB, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.5), y + Inches(0.05), Inches(5.0), Inches(0.6),
                    body, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(10.7), y + Inches(0.05), Inches(2.2), Inches(0.6),
                    label, font_size=13, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
                '소학(小學)·삼강행실도·동몽선습 — 동양 윤리 교육의 뼈대가 된 다섯 기둥',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅴ. 7편 14장 깊이 읽기 ==============
SEC5 = 'Ⅴ. 7편 14장 깊이 읽기'

def make_chapter_slide(num, total, name_han, name_kor, headline,
                        original, modern, theme, point):
    @S(SEC5)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC5} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=26, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    headline, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.4), [
            (original, {'font_size': 18, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 편의 핵심', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(theme, {'font_size': 14, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 오늘에의 적용', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 15, 'color': INK})], line_spacing=1.35)


CHAPTERS = [
    ('梁惠王 上 양혜왕 상', '양혜왕 상', '제1편 상 — 왕도와 패도 · 맹자의 첫 발화',
     '王 何 必 曰 利  亦 有 仁 義 而 已 矣',
     '왕께서 하필 이익을 말씀하십니까? 오직 인의(仁義)만 있을 뿐입니다',
     '맹자가 양혜왕을 만난 역사적 첫 대화 / 「五十步百步」 — 도토리 키재기 / 「不嗜殺人者能一之」 — 사람 죽이기를 좋아하지 않는 자가 천하를 통일 / 「樂民之樂 民亦樂其樂」 — 여민동락의 시작',
     '경쟁이 「이익(利)」으로만 흐르면 결국 모두가 진다 — 의(義)가 토대일 때만 지속. 비즈니스 윤리·ESG의 동양적 원천.'),
    ('梁惠王 下 양혜왕 하', '양혜왕 하', '제1편 하 — 여민동락과 폭군 방벌',
     '聞 誅 一 夫 紂 矣  未 聞 弑 君 也',
     '한 사내(紂)를 죽였다 들었지, 임금을 시해했다 듣지 못했다',
     '제선왕과의 여민동락 대화 / 「과인은 음악을 좋아한다」 → 「백성과 함께 즐기면 왕도가 가깝다」 / 폭군 방벌론 — 주왕은 더 이상 임금이 아닌 일개 「부(夫)」 / 혁명권의 사상적 근거',
     '리더의 자격은 직책이 아닌 행위 — 잘못된 권력은 정당성을 잃는다. 「리더십의 합법성」의 동양 원형.'),
    ('公孫丑 上 공손추 상', '공손추 상', '제2편 상 — 호연지기와 사단(四端)',
     '惻 隱 之 心  仁 之 端 也',
     '측은히 여기는 마음이 인(仁)의 단서이다',
     '맹자 핵심 사상의 정점 — 사단(四端) / 측은(仁)·수오(義)·사양(禮)·시비(智) / 「우물에 빠지려는 아이」 — 모든 사람에게 측은지심이 있다는 증명 / 호연지기 — 「我善養吾浩然之氣」',
     '도덕은 외부의 명령이 아닌 내면의 발견 — 누구나 가진 「네 가지 마음」을 길러내는 것이 인격 완성.'),
    ('公孫丑 下 공손추 하', '공손추 하', '제2편 하 — 천시·지리·인화',
     '天 時 不 如 地 利  地 利 不 如 人 和',
     '하늘의 때는 땅의 이로움만 못하고, 땅의 이로움은 사람의 화합만 못하다',
     '승부의 3단계 — 천시·지리·인화 / 「得道者多助 失道者寡助」 — 도를 얻은 자는 도움이 많고, 잃은 자는 적다 / 인화(人和)가 가장 결정적 / 맹자의 사명감 — 「500년 만에 왕자가 난다」',
     '경쟁의 가장 큰 자산은 「사람의 마음」 — 시기·자원보다 신뢰가 우선. 조직 문화의 동양적 정의.'),
    ('滕文公 上 등문공 상', '등문공 상', '제3편 상 — 정전제와 분업론',
     '勞 心 者 治 人  勞 力 者 治 於 人',
     '마음을 쓰는 자는 남을 다스리고, 힘을 쓰는 자는 남에게 다스려진다',
     '등문공에게 정전제(井田制) 자문 — 토지를 우물 정(井)자로 9등분 / 허행과의 분업 논쟁 — 「백공의 일은 본래 분업」 / 사회의 분업 원리 · 노심자(勞心者)와 노력자(勞力者)',
     '사회는 분업으로 작동 — 머리와 손이 서로 다른 일을 하되 결합한다. 애덤 스미스보다 2,000년 앞선 분업론.'),
    ('滕文公 下 등문공 하', '등문공 하', '제3편 하 — 대장부론과 양묵 비판',
     '富 貴 不 能 淫  貧 賤 不 能 移  威 武 不 能 屈  此 之 謂 大 丈 夫',
     '부귀가 흔들지 못하고, 빈천이 옮기지 못하며, 위세가 굽히지 못하는 자 — 이것이 대장부',
     '맹자의 「대장부론(大丈夫論)」 — 동양 군자상의 결정판 / 「養浩然之氣」의 인격 정의 / 양주(楊朱)·묵적(墨翟) 비판 — 「양묵을 막는 것이 내 사명」 / 「無父無君 是禽獸也」',
     '진짜 인격은 외부 조건에 흔들리지 않는 자 — 부귀·빈천·위세 모두 인격의 기준이 못 된다.'),
    ('離婁 上 이루 상', '이루 상', '제4편 상 — 인의 정치와 자기 반성',
     '愛 人 不 親  反 其 仁   治 人 不 治  反 其 智',
     '남을 사랑해도 친해지지 않으면 자기 인(仁)을 돌아보고, 다스려도 다스려지지 않으면 자기 지혜를 돌아보라',
     '「反求諸己(반구저기)」 — 모든 문제의 원인을 자기에게서 찾는다 / 「離婁之明 公輸子之巧 不以規矩 不能成方員」 — 시력 좋은 이루도 컴퍼스 없이 원을 못 그린다 / 「規矩」 — 표준의 중요성',
     '결과가 안 나올 때 「남 탓」이 아니라 「내 탓」 — 가장 강력한 자기 책임의 원칙. 외부 표준(規矩)의 가치.'),
    ('離婁 下 이루 하', '이루 하', '제4편 하 — 군신 관계의 상호성',
     '君 之 視 臣 如 手 足  則 臣 視 君 如 腹 心',
     '임금이 신하를 손발같이 보면, 신하는 임금을 배와 심장같이 본다',
     '군신 관계의 상호성 — 일방적 충성이 아닌 「반응」의 윤리 / 「대인 — 적자지심(大人者 不失其赤子之心)」 — 큰 사람은 어린아이 마음을 잃지 않는다 / 「人之異於禽獸者幾希」 — 사람과 짐승의 차이는 미세하다',
     '관계는 항상 「상호적」 — 윗사람이 어떻게 대하느냐가 아랫사람의 태도를 결정. 일방주의 리더십에 대한 경고.'),
    ('萬章 上 만장 상', '만장 상', '제5편 상 — 요순의 천하 양도와 천명',
     '天 不 言  以 行 與 事 示 之 而 已 矣',
     '하늘은 말하지 않는다 — 행위와 일로 보여줄 뿐이다',
     '만장의 「요임금이 천하를 순에게 준 것이 사실?」 질문 / 「天與之 — 하늘이 준 것이다」 / 「天視自我民視 天聽自我民聽」 (서경 인용) / 천명(天命)은 곧 민의(民意)',
     '권력의 정당성은 결국 「백성의 마음」 — 하늘의 뜻이 백성의 시청에 드러난다는 동양 민본주의의 가장 깊은 진술.'),
    ('萬章 下 만장 하', '만장 하', '제5편 하 — 성인의 네 가지 풍모',
     '伯 夷 聖 之 淸 者 也   伊 尹 聖 之 任 者 也   柳 下 惠 聖 之 和 者 也   孔 子 聖 之 時 者 也',
     '백이는 청(淸)의 성인, 이윤은 임(任)의 성인, 유하혜는 화(和)의 성인, 공자는 시(時)의 성인',
     '성인의 네 가지 유형론 — 청(맑음)·임(책임)·화(조화)·시(때) / 공자가 가장 위 — 「시중(時中)」의 인격 / 「集大成」 — 모든 것을 모은 완성',
     '인격에도 다양한 모범이 있다 — 어떤 모범을 택하느냐는 자기 길의 선택. 「때를 아는 자」가 가장 큰 인격.'),
    ('告子 上 고자 상', '고자 상', '제6편 상 — 성선설의 정점 · 인성 논쟁',
     '人 性 之 善 也  猶 水 之 就 下 也',
     '사람의 본성이 선한 것은 물이 아래로 흐르는 것과 같다',
     '고자와의 4대 논쟁 — 본성·비유·식색·의(義) / 「水信無分於東西 無分於上下乎?」 — 물에 동서 구분은 없어도 상하 구분은 있다 / 「우산(牛山)의 비유」 — 본래 선한 산도 도끼질로 황폐해질 수 있다',
     '인간 본성에 대한 가장 낙관적인 답 — 「누구나 선하다, 환경이 가린다」. 교육·재활·인본주의 사상의 동양 원천.'),
    ('告子 下 고자 하', '고자 하', '제6편 하 — 천장강대임(天將降大任)',
     '天 將 降 大 任 於 是 人 也  必 先 苦 其 心 志  勞 其 筋 骨',
     '하늘이 큰 책임을 이 사람에게 내리려 할 때, 반드시 먼저 그 마음과 뜻을 괴롭히고 근육과 뼈를 수고롭게 한다',
     '시련의 의미 — 「큰일을 맡기려는 자에게 먼저 시련을 준다」 / 순·부열·교격·관중·손숙오·백리해 등 「시련을 거친 인물들」의 예 / 「生於憂患 死於安樂」 — 우환에서 살고 안락에서 죽는다',
     '시련을 「의미」로 전환하는 가장 오래된 명제 — 어려움이 곧 단련의 과정. 현대 회복탄력성·그릿 이론의 동양 원천.'),
    ('盡心 上 진심 상', '진심 상', '제7편 상 — 진심·지성·지천',
     '盡 其 心 者  知 其 性 也   知 其 性  則 知 天 矣',
     '마음을 다하는 자는 본성을 알고, 본성을 알면 하늘을 안다',
     '맹자 수양론의 정점 — 「盡心 → 知性 → 知天」 / 「萬物皆備於我矣」 — 만물이 모두 나에게 갖추어져 있다 / 「君子有三樂」 — 군자의 세 즐거움 (부모형제·앙불괴어천·천하영재교지)',
     '내면 깊이 들어가는 것이 곧 우주에 닿는 길 — 동양 명상·내면 수양의 가장 정교한 정의. 자기 내부에 모든 답이 있다.'),
    ('盡心 下 진심 하', '진심 하', '제7편 하 — 민귀군경과 맹자의 마지막',
     '民 爲 貴  社 稷 次 之  君 爲 輕',
     '백성이 가장 귀하고, 사직(나라)이 그 다음이며, 임금은 가장 가볍다',
     '맹자의 마지막 편 / 「民爲貴 社稷次之 君爲輕」 — 동양 민본주의의 절대 명제 / 「盡信書 不如無書」 — 책을 다 믿느니 차라리 없는 게 낫다 / 「五百年必有王者興」 — 500년에 반드시 왕자가 난다',
     '권력의 위계를 뒤집은 동양 정치사상의 정점 — 백성이 가장 위, 임금이 가장 아래. 모든 민주주의 사상의 동양 원천.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅵ. 핵심 사상 9기둥 ==============
SEC6 = 'Ⅵ. 핵심 사상 9기둥'

def make_concept_slide(num, total, han, kor, source, principle, today):
    @S(SEC6)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC6} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{han}  ({kor})',
                    font_size=30, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    source, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.7), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.6),
                       [(principle, {'font_size': 14, 'color': INK})], line_spacing=1.45)
        add_textbox(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.4),
                    '◆ 오늘에의 함의', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.3), Inches(12.0), Inches(3.0),
                       [(today, {'font_size': 15, 'color': INK})], line_spacing=1.5)


CONCEPTS = [
    ('性 善 說', '성선설', '고자편 · 맹자 사상의 출발점',
     '「人性之善也 猶水之就下也」 — 사람의 본성이 선한 것은 물이 아래로 흐르는 것과 같다\n물은 아래로 흐르는 것이 본성 — 위로 튀게 할 수 있으나 그것은 외부 힘\n악행은 본성이 아니라 외부 환경·욕심의 탓',
     '인간에 대한 가장 낙관적인 답 — 누구나 선하다, 환경이 가린다.\n교육·재활·인본주의 사상의 동양 원천.\n순자의 성악설과 정면 대립 — 동아시아 인성론의 두 축.'),
    ('四 端', '사단', '공손추 상편 · 인간 본성의 네 단서',
     '· 측은지심(惻隱之心) — 인(仁)의 단서 — 우물에 빠지려는 아이\n· 수오지심(羞惡之心) — 의(義)의 단서 — 부끄러워하는 마음\n· 사양지심(辭讓之心) — 예(禮)의 단서 — 양보하는 마음\n· 시비지심(是非之心) — 지(智)의 단서 — 옳고 그름을 가리는 마음',
     '도덕은 외부 명령이 아닌 내면의 발견.\n현대 도덕심리학(조나단 하이트 등)의 「도덕 직관」 이론의 2,300년 선행.\n「누구나 가진 네 가지 마음」 — 인간 존엄의 근거.'),
    ('王 道 / 覇 道', '왕도·패도', '양혜왕편 · 정치의 두 길',
     '· 왕도(王道) — 덕(德)으로 다스림 · 「以德服人 心悅而誠服也」\n· 패도(覇道) — 힘(力)으로 다스림 · 「以力服人者 非心服也」\n왕도는 마음으로 복종 · 패도는 두려움으로 복종 — 지속 가능성이 다르다',
     '권력의 두 종류 — 외부 강제 vs 내부 동의.\n현대 「소프트 파워 vs 하드 파워」(조지프 나이)의 동양 원천.\n조직 리더십도 같다 — 강제로 시킨 일과 마음으로 한 일의 차이.'),
    ('民 本 / 民 貴 君 輕', '민본·민귀군경', '진심 하편 · 동양 정치사상의 정점',
     '「民爲貴 社稷次之 君爲輕」 — 백성이 가장 귀하고, 사직이 다음, 임금이 가장 가볍다\n권력의 위계를 뒤집는 혁명적 명제\n「天視自我民視 天聽自我民聽」(서경 인용) — 하늘의 뜻은 곧 백성의 뜻',
     '동양 민본주의의 가장 강력한 표현.\n현대 민주주의·국민주권론과 직접 통한다.\n조선 정약용 『목민심서』, 동학(東學), 현대 한국 헌법 정신의 사상적 원천.'),
    ('與 民 同 樂', '여민동락', '양혜왕 하편 · 군주와 백성의 함께함',
     '제선왕과의 대화 — 「과인은 음악을 좋아한다」 → 「백성과 함께 즐기시면 왕도가 가깝다」\n혼자 즐기는 즐거움은 작고, 함께 즐기는 즐거움이 크다\n「樂民之樂 民亦樂其樂」 — 백성의 즐거움을 즐기면 백성도 그 즐거움을 즐긴다',
     '리더십의 가장 따뜻한 정의 — 「부하의 즐거움을 같이 즐길 수 있는가」.\n현대 조직 문화·복지 경영의 동양 원형.\n불평등이 깊어진 오늘 가장 절실한 메시지.'),
    ('浩 然 之 氣', '호연지기', '공손추 상편 · 맹자의 수양론 정점',
     '공손추: 「선생님은 어떻게 부동심이십니까?」 → 맹자: 「我善養吾浩然之氣」\n「至大至剛 以直養而無害 則塞於天地之間」 — 지극히 크고 강하여 정직으로 기르면 천지 사이에 가득\n「集義所生」 — 의(義)를 쌓아서 생긴다 — 한 번에 안 되고 누적이 필요',
     '동양 수양론의 가장 웅장한 단어 — 「천지에 가득한 기운」.\n현대의 「플로우(flow)」·「존재감(presence)」의 동양적 원형.\n그러나 비결은 평범 — 매일의 의(義)를 쌓을 것.'),
    ('仁 義', '인의', '맹자 전체 · 공자의 인을 인의 병칭으로',
     '공자는 인(仁)을 중심 — 맹자는 인의(仁義)를 병칭\n「居仁由義」 — 인에 머물고 의로 행한다\n「仁 人之安宅也  義 人之正路也」 — 인은 사람의 편안한 집, 의는 바른 길\n「不仁不智 無禮無義 人役也」 — 인의예지가 없으면 남에게 부려질 뿐',
     '맹자가 공자에 더한 결정적 한 축 — 「의(義)」.\n인이 사랑이라면, 의는 그 사랑의 방향성·올바름.\n맹자 이후 동양 윤리의 표준 어휘 — 「인의예지」 사덕(四德).'),
    ('恒 産 恒 心', '항산항심', '양혜왕 상편 · 경제와 도덕의 관계',
     '「無恒産而有恒心者 惟士爲能 若民則無恒産 因無恒心」\n항산(恒産, 안정된 생업)이 없으면 항심(恒心, 변함없는 마음)이 없다\n선비는 항산 없이도 항심을 지킬 수 있으나, 백성은 항산이 없으면 도덕을 지킬 수 없다\n→ 정전제(井田制) 제안의 사상적 토대',
     '도덕보다 「생계의 안정」이 먼저 — 「의식주가 족해야 예의를 안다」(管子)의 맹자적 정식.\n현대 복지국가·기본소득·경제 정책의 동양 원천.\n「가난이 도덕의 문제가 아니라 정치의 문제」라는 통찰.'),
    ('良 知 良 能', '양지양능', '진심 상편 · 왕양명 심학(心學)의 출전',
     '「人之所不學而能者 其良能也  所不慮而知者 其良知也」\n사람이 배우지 않고도 할 수 있는 것이 양능, 생각하지 않고도 아는 것이 양지\n어린아이가 부모를 사랑하고 형을 공경할 줄 아는 것 — 본래적 도덕 능력의 증명\n→ 명대 왕양명이 「致良知」로 발전시킨 심학의 원천',
     '도덕 능력은 외부 학습이 아닌 「본래 갖춘 것」 — 양지·양능.\n왕양명 「致良知」, 한국 양명학(정제두), 일본 양명학(나카에 도쥬)의 출발점.\n현대 도덕 직관·아이의 도덕성 연구의 동양적 정식.'),
]

for i, c in enumerate(CONCEPTS, 1):
    make_concept_slide(i, len(CONCEPTS), *c)


# ============== Ⅶ. 맹자의 유명한 비유 ==============
SEC7 = 'Ⅶ. 맹자의 유명한 비유'

@S(SEC7)
def vii_metaphors(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '맹자의 비유 — 추상을 일상으로 옮기는 천재',
              '오십보백보·연목구어·송인양묘·물의 비유 등')
    rows = [
        ('五 十 步 百 步', '오십보백보',
         '전쟁터에서 50보 도망친 자가 100보 도망친 자를 비웃을 수 있는가?\n「정도의 차이일 뿐 본질은 같다」 — 양혜왕에게 자기 정치의 한계를 깨우치는 비유'),
        ('緣 木 求 魚', '연목구어',
         '나무에 올라가 물고기를 구하는 격이다 — 잘못된 방법으로 목적을 추구\n제선왕에게 「패도」로는 왕도의 효과를 못 얻는다고 비판'),
        ('揠 苗 助 長', '알묘조장 (송인양묘)',
         '송나라 사람이 벼 자라기를 도우려 벼를 잡아당겨 모두 시들게 함\n호연지기는 단번에 안 되고 누적이 필요 — 인위적 조장의 어리석음 경고'),
        ('水 之 就 下', '물의 비유',
         '사람의 본성은 물이 아래로 흐르는 것과 같다 — 性善의 비유\n외부 힘으로 위로 튀게 할 수 있으나 본성은 아래로'),
        ('牛 山 之 木', '우산지목',
         '우산(牛山)의 나무가 본래 무성했으나 도끼질로 황폐해진 것 — 본성이 환경으로 가려지는 비유\n「夜氣」 — 밤의 맑은 기운으로도 회복 가능'),
        ('反 求 諸 己', '반구저기',
         '활을 쏘아 맞지 않으면 이긴 자를 원망 말고 자기를 돌아보라\n모든 문제의 원인을 자기에게서 찾는 동양 자기 책임의 원리'),
    ]
    for i, (han, kor, body) in enumerate(rows):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 1.55)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(1.35), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.0), Inches(0.45),
                    han, font_size=16, bold=True, color=ACCENT, font_name='Batang')
        add_textbox(slide, x + Inches(3.3), y + Inches(0.05), Inches(2.4), Inches(0.45),
                    kor, font_size=12, color=SUB, bold=True,
                    align=PP_ALIGN.RIGHT)
        add_paragraphs(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.75),
                       [(body, {'font_size': 11, 'color': INK})], line_spacing=1.3)


@S(SEC7)
def vii_daejangbu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '대장부론(大丈夫論) — 등문공 하편',
              '동양 군자상의 결정판')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(2.0), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.6),
                '富 貴 不 能 淫',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(2.9), Inches(12.0), Inches(0.6),
                '貧 賤 不 能 移',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.6),
                '威 武 不 能 屈  此 之 謂 大 丈 夫',
                font_size=24, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.5), Inches(12.0), Inches(2.5), [
        ('● 富貴不能淫(부귀불능음) — 부귀가 (인격을) 흔들지 못한다',
         {'font_size': 17, 'space_before': 4, 'font_name': 'Batang'}),
        ('● 貧賤不能移(빈천불능이) — 빈천이 (인격을) 옮기지 못한다',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
        ('● 威武不能屈(위무불능굴) — 위세와 무력이 (인격을) 굽히지 못한다',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
        ('● 이것이 곧 「大丈夫」 — 외부 조건이 흔들지 못하는 인격',
         {'font_size': 17, 'space_before': 14, 'bold': True, 'color': ACCENT}),
        ('● 안중근·윤봉길·김구·이순신 — 모두 이 구절을 좌우명 삼았던 한국의 대장부들',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC7)
def vii_trial(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '천장강대임(天將降大任) — 고자 하편',
              '시련을 의미로 전환하는 가장 오래된 명제')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(2.2), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.6),
                '天 將 降 大 任 於 是 人 也',
                font_size=22, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(2.9), Inches(12.0), Inches(0.6),
                '必 先 苦 其 心 志   勞 其 筋 骨',
                font_size=20, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.6),
                '餓 其 體 膚   空 乏 其 身   行 拂 亂 其 所 爲',
                font_size=20, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(0.5),
                '所 以 動 心 忍 性  曾 益 其 所 不 能',
                font_size=18, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.5), [
        ('● 「하늘이 큰 책임을 이 사람에게 내리려 할 때, 반드시 먼저 그 마음과 뜻을 괴롭히고 근육과 뼈를 수고롭게 하며...」',
         {'font_size': 14, 'space_before': 4, 'color': SUB}),
        ('● 「마음을 흔들고 본성을 참게 하여, 일찍이 할 수 없던 것을 늘리기 위함이다」',
         {'font_size': 14, 'space_before': 6, 'color': SUB}),
        ('● 「生於憂患 死於安樂」 — 우환에서 살고 안락에서 죽는다',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 회복탄력성(resilience)·그릿(grit)·anti-fragile의 2,300년 선행 명제',
         {'font_size': 15, 'space_before': 8, 'color': SUB}),
    ])


# ============== Ⅷ. 정전제와 항산항심의 경제 ==============
SEC_ECON = 'Ⅷ. 정전제와 경제사상'

@S(SEC_ECON)
def econ_jeongjeon(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_ECON, n, t)
    add_title(slide, '정전제(井田制) — 맹자가 그린 이상적 토지 제도',
              '등문공 상편 · 토지를 우물 정(井)자로 9등분')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(6.0), Inches(4.6), PALE)
    cells = [
        ('私田', '私田', '私田'),
        ('私田', '公田', '私田'),
        ('私田', '私田', '私田'),
    ]
    for r in range(3):
        for c in range(3):
            x = Inches(0.9 + c * 1.85)
            y = Inches(2.4 + r * 1.45)
            is_center = (r == 1 and c == 1)
            color = ACCENT if is_center else INK
            add_filled_rect(slide, x, y, Inches(1.7), Inches(1.3), color)
            label = cells[r][c]
            add_textbox(slide, x, y + Inches(0.2), Inches(1.7), Inches(0.5),
                        label, font_size=18, bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER, font_name='Batang')
            if is_center:
                add_textbox(slide, x, y + Inches(0.7), Inches(1.7), Inches(0.4),
                            '공동 경작', font_size=11, color=RULE,
                            align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(4.5), [
        ('● 토지를 우물 정(井)자로 9등분',
         {'font_size': 16, 'bold': True, 'color': ACCENT}),
        ('● 8가구가 둘레 8칸을 「사전(私田)」으로 경작',
         {'font_size': 14, 'space_before': 10}),
        ('● 가운데 1칸은 「공전(公田)」 — 8가구가 함께 경작 → 세금으로 납입',
         {'font_size': 14, 'space_before': 8}),
        ('● 「方里而井 井九百畝 其中爲公田 八家皆私百畝」',
         {'font_size': 13, 'space_before': 10, 'color': SUB, 'font_name': 'Batang'}),
        ('     1리(里)에 정(井) 하나 — 900묘 · 가운데 공전, 8가구가 사전 100묘씩',
         {'font_size': 12, 'color': SUB, 'space_before': 4}),
        ('● 9분의 1만 세금 — 동양에서 가장 오래된 「적정 조세」 사상',
         {'font_size': 14, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 협력 노동·공동체 의식·균등 토지 — 동양 평등주의의 원천',
         {'font_size': 13, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC_ECON)
def econ_hangsan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_ECON, n, t)
    add_title(slide, '항산항심(恒産恒心) — 경제와 도덕의 관계',
              '양혜왕 상편 · 「항산 없이 항심 없다」')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.8), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.7),
                '無 恒 産 而 有 恒 心 者  惟 士 爲 能',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(12.0), Inches(0.7),
                '若 民 則 無 恒 産  因 無 恒 心',
                font_size=22, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                '「항산 없이도 항심을 지킬 수 있는 자는 오직 선비뿐 · 백성은 항산 없으면 항심도 없다」',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.3), Inches(12.0), Inches(3.0), [
        ('● 항산(恒産) — 안정된 생업(토지·일자리·소득)',
         {'font_size': 17, 'space_before': 6}),
        ('● 항심(恒心) — 변함없는 도덕적 마음',
         {'font_size': 17, 'space_before': 10}),
        ('● 도덕을 가르치기 전에 먼저 생계를 안정시켜야 한다 — 「先養後敎」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 백성의 잘못은 백성의 잘못이 아니라 「제도의 잘못」 — 정치의 1차 책임',
         {'font_size': 17, 'space_before': 10}),
        ('● 현대 복지국가·기본소득·최저임금·노동권 등 사회 경제 정책의 동양적 원천',
         {'font_size': 15, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


# ============== Ⅸ. 명구 18선 ==============
SEC8 = 'Ⅸ. 명구 18선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC8)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC8} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=16, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 16, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('양혜왕 상편 · 맹자의 첫 발화', '王 何 必 曰 利  亦 有 仁 義 而 已 矣',
     '왕께서 하필 이익을 말씀하십니까? 오직 인의(仁義)만 있을 뿐입니다',
     '맹자가 양혜왕을 만난 역사적 첫 대화. 맹자 7편 260장의 모든 논의는 이 한 문장의 전개다.'),
    ('양혜왕 상편', '五 十 步 笑 百 步',
     '오십 보 도망친 자가 백 보 도망친 자를 비웃는다',
     '도토리 키재기. 정도의 차이일 뿐 본질은 같다 — 양혜왕에게 자기 정치의 한계를 깨우치는 가장 유명한 비유.'),
    ('양혜왕 하편 · 제선왕 대화', '樂 民 之 樂 者  民 亦 樂 其 樂',
     '백성의 즐거움을 즐기는 자는, 백성도 그의 즐거움을 즐긴다',
     '여민동락(與民同樂)의 정의. 리더십의 가장 따뜻한 한 줄 — 「부하의 즐거움을 같이 즐길 수 있는가」.'),
    ('양혜왕 하편 · 폭군 방벌론', '聞 誅 一 夫 紂 矣  未 聞 弑 君 也',
     '한 사내(紂)를 죽였다 들었지, 임금을 시해했다 듣지 못했다',
     '리더의 자격은 직책이 아닌 행위 — 잘못된 권력은 정당성을 잃는다. 동양 혁명권의 가장 강한 진술.'),
    ('공손추 상편 · 사단(四端)', '惻 隱 之 心  仁 之 端 也',
     '측은히 여기는 마음이 인(仁)의 단서이다',
     '맹자 핵심 사상의 정점 — 사단(四端). 「누구나 가진 네 가지 마음」 — 도덕은 외부 명령이 아닌 내면의 발견.'),
    ('공손추 상편 · 호연지기', '我 善 養 吾 浩 然 之 氣',
     '나는 나의 호연지기를 잘 기른다',
     '동양 수양론의 가장 웅장한 단어. 「지극히 크고 강하여 정직으로 기르면 천지 사이에 가득」 — 그러나 비결은 매일의 의를 쌓는 평범함.'),
    ('공손추 하편', '天 時 不 如 地 利  地 利 不 如 人 和',
     '하늘의 때는 땅의 이로움만 못하고, 땅의 이로움은 사람의 화합만 못하다',
     '승부의 3단계 — 천시·지리·인화. 가장 결정적인 것은 「사람의 마음」. 조직 문화의 동양적 정의.'),
    ('등문공 상편 · 분업론', '勞 心 者 治 人  勞 力 者 治 於 人',
     '마음을 쓰는 자는 남을 다스리고, 힘을 쓰는 자는 남에게 다스려진다',
     '사회 분업론의 동양적 원형. 애덤 스미스보다 2,000년 앞선 분업 사상. 머리와 손이 서로 다른 일을 하되 결합한다.'),
    ('등문공 하편 · 대장부론', '富 貴 不 能 淫  貧 賤 不 能 移  威 武 不 能 屈',
     '부귀가 흔들지 못하고, 빈천이 옮기지 못하며, 위세가 굽히지 못한다',
     '동양 군자상의 결정판. 안중근·윤봉길·김구·이순신이 모두 좌우명 삼았던 한국 독립운동의 정신적 기둥.'),
    ('이루 상편', '反 求 諸 己',
     '돌이켜 자기에게서 구한다',
     '활을 쏘아 맞지 않으면 이긴 자를 원망 말고 자기를 돌아보라 — 가장 강력한 자기 책임의 원칙. 4글자에 압축된 인격의 핵심.'),
    ('이루 하편 · 군신관계', '君 之 視 臣 如 手 足  則 臣 視 君 如 腹 心',
     '임금이 신하를 손발같이 보면, 신하는 임금을 배와 심장같이 본다',
     '관계의 상호성. 일방적 충성이 아닌 「반응」의 윤리 — 일방주의 리더십에 대한 가장 오래된 경고.'),
    ('만장 상편 · 천명론', '天 不 言  以 行 與 事 示 之 而 已 矣',
     '하늘은 말하지 않는다 — 행위와 일로 보여줄 뿐이다',
     '천명(天命)은 곧 민의(民意). 권력의 정당성은 결국 「백성의 마음」 — 「天視自我民視 天聽自我民聽」의 짝.'),
    ('고자 상편 · 성선설', '人 性 之 善 也  猶 水 之 就 下 也',
     '사람의 본성이 선한 것은 물이 아래로 흐르는 것과 같다',
     '성선설의 가장 시적인 표현. 인간에 대한 가장 낙관적인 답 — 「누구나 선하다, 환경이 가린다」.'),
    ('고자 하편 · 시련론', '生 於 憂 患  死 於 安 樂',
     '우환에서 살고, 안락에서 죽는다',
     '시련을 의미로 전환하는 한 줄. 회복탄력성(resilience)·그릿(grit)·anti-fragile의 2,300년 선행 명제.'),
    ('진심 상편', '萬 物 皆 備 於 我 矣',
     '만물이 모두 나에게 갖추어져 있다',
     '맹자 수양론의 정점. 내면 깊이 들어가는 것이 곧 우주에 닿는 길 — 자기 안에 모든 답이 있다는 동양 내면주의의 정수.'),
    ('진심 하편 · 민귀군경', '民 爲 貴   社 稷 次 之   君 爲 輕',
     '백성이 가장 귀하고, 사직이 다음, 임금이 가장 가볍다',
     '맹자의 마지막 정수. 동양 민본주의의 절대 명제. 정약용 『목민심서』·동학·현대 한국 헌법 정신의 사상적 원천.'),
    ('양혜왕 상편 · 항산항심', '無 恒 産 者  無 恒 心',
     '항산(안정된 생업)이 없는 자에겐 항심(변함없는 마음)이 없다',
     '경제와 도덕의 관계를 정의한 한 줄. 「의식주가 족해야 예의를 안다」의 맹자적 정식 — 현대 복지·기본소득 사상의 동양 원천.'),
    ('진심 상편 · 양지양능', '不 學 而 能 者  良 能 也   不 慮 而 知 者  良 知 也',
     '배우지 않고도 할 수 있는 것이 양능, 생각하지 않고도 아는 것이 양지',
     '왕양명 「치양지(致良知)」 심학(心學)의 출전. 도덕 능력은 본래 갖춘 것 — 어린아이의 사랑·공경이 그 증거.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅹ. 동아시아 수용 ==============
SEC9 = 'Ⅹ. 동아시아 수용'

@S(SEC9)
def ix_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '한국 — 조선 민본주의의 사상적 기둥',
              '퇴계·율곡부터 정약용·동학·독립운동까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 고려 — 안향이 주자학과 함께 도입, 『맹자』가 사서로 들어옴',
         {'font_size': 17, 'space_before': 4}),
        ('● 조선 — 사서 학습의 두 번째 책 (대학 → 논어 → 맹자 → 중용)',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 퇴계 이황 — 『맹자석의』 / 율곡 이이 — 『맹자석의』',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 정약용 — 『맹자요의』 — 가장 방대한 조선 시대 맹자 주석서',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 「民爲貴」 — 조선 민본주의의 사상적 원천 → 정약용 『목민심서』',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 동학(東學) — 「人乃天 — 사람이 곧 하늘」의 사상적 토대',
         {'font_size': 17, 'space_before': 10}),
        ('● 독립운동 — 안중근·윤봉길·김구의 좌우명 「富貴不能淫·貧賤不能移·威武不能屈」',
         {'font_size': 16, 'space_before': 10, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC9)
def ix_sachil(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '사단칠정(四端七情) 논쟁 — 조선 성리학의 정점',
              '맹자의 「사단」이 한국 철학사에서 200년 논쟁을 낳다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 사단(四端) — 맹자 「측은·수오·사양·시비」 / 칠정(七情) — 예기 「희·노·애·구·애·오·욕」',
         {'font_size': 17, 'space_before': 4}),
        ('● 「둘은 어떻게 다른가?」가 조선 성리학의 200년 논쟁',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 1차 논쟁 — 퇴계 이황 vs 고봉 기대승 (1559~1566, 7년)',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 퇴계 — 「사단은 理의 발(理發), 칠정은 氣의 발(氣發)」 — 이기호발설(理氣互發說)',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 고봉 — 「사단과 칠정은 모두 정(情), 사단은 칠정 중 善한 부분」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 2차 논쟁 — 율곡 이이 vs 우계 성혼 (1572~1573)',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 율곡 — 「氣發理乘一途說」 — 발하는 것은 氣뿐, 理는 氣를 타고 따른다',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 맹자의 한 단어가 한국 철학사 최대 논쟁의 씨앗이 된 역사적 사례',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC9)
def ix_japan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '일본 — 메이지 유신과 맹자',
              '에도의 비판적 독해부터 메이지 「혁명」의 사상적 무기까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 에도 막부 — 처음엔 「방벌론」 때문에 위험시 → 점차 학자들이 수용',
         {'font_size': 17, 'space_before': 4}),
        ('● 이토 진사이(伊藤仁齋) — 『맹자고의』 — 일본 고학파(古學派)의 정수',
         {'font_size': 17, 'space_before': 10}),
        ('● 오규 소라이(荻生徂徠) — 맹자에 비판적이었으나 깊이 연구',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 요시다 쇼인(吉田松陰) — 메이지 유신 사상의 토대로 맹자 활용',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 「민귀군경」·「방벌론」을 막부 타도의 사상적 무기로',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 메이지 이후 — 「자유민권운동」의 사상적 자원으로 다시 부상',
         {'font_size': 17, 'space_before': 10}),
        ('● 시부사와 에이이치 — 『논어와 주판』에 맹자의 의(義)도 통합',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC9)
def ix_west(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '서양 — Mencius로 알려진 맹자',
              '계몽주의자들이 발견한 동양 인본주의의 정점')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 17세기 — 예수회 선교사들이 「Mencius」로 라틴어 명명',
         {'font_size': 17, 'space_before': 4}),
        ('● 1711 — 노엘(François Noël) 라틴어 번역 『Sinensis Imperii Libri Classici Sex』',
         {'font_size': 17, 'space_before': 10}),
        ('● 계몽주의자들의 발견 — 볼테르·라이프니츠가 맹자의 민본·인본을 격찬',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 19~20세기 — 영어 번역(James Legge, 1861) — 동양 정치사상의 표준 텍스트',
         {'font_size': 17, 'space_before': 10}),
        ('● 「성선설」 — 루소·간디 등 서구 인본주의자들의 동양적 짝',
         {'font_size': 17, 'space_before': 10}),
        ('● 현대 — 「민귀군경」이 동아시아 민주주의 사상의 동양 원천으로 재조명',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 영문판이 동양 정치사상서 중 가장 자주 인용 — Confucius 다음의 「Mencius」',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


# ============== Ⅺ. 오늘 우리에게 + 마무리 ==============
SEC10 = 'Ⅺ. 오늘 우리에게'

@S(SEC10)
def x_today(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '맹자가 오늘 우리에게 일러주는 12가지')
    items = [
        '1. 이익(利) 이전에 의(義)를 묻는다 — 王何必曰利',
        '2. 사람은 본래 선하다 — 性善',
        '3. 누구나 「네 가지 마음」을 가졌다 — 四端',
        '4. 백성이 가장 귀하다 — 民貴君輕',
        '5. 부귀·빈천·위세가 흔들지 못하는 자 — 大丈夫',
        '6. 천지에 가득한 기운 — 浩然之氣',
        '7. 같이 즐기는 즐거움이 크다 — 與民同樂',
        '8. 모든 문제의 답은 자기에게 — 反求諸己',
        '9. 시련은 큰 책임의 준비 — 天將降大任',
        '10. 우환에서 살고 안락에서 죽는다 — 生於憂患 死於安樂',
        '11. 만물이 나에게 갖추어져 있다 — 萬物皆備於我',
        '12. 손발같이 보면 배와 심장같이 본다 — 관계의 상호성',
    ]
    for i, txt in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 0.7)
        add_textbox(slide, x, y, Inches(6.0), Inches(0.6),
                    txt, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '맹자, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 맹가(BC 372?~289?)와 제자 만장·공손추가 함께 엮은 7편 14장 약 35,000자.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 사서의 둘째 · 아성(亞聖)의 책 — 공자의 인(仁)을 체계화.',
         {'font_size': 18, 'space_before': 8}),
        ('● 7대 사상 — 성선설 · 사단 · 왕도 · 민본 · 여민동락 · 호연지기 · 인의.',
         {'font_size': 18, 'space_before': 8}),
        ('● 맹모삼천·맹모단기 — 동양 교육의 원형. 40년 천하 유세 — 좌절과 사명감.',
         {'font_size': 18, 'space_before': 8}),
        ('● 오십보백보·연목구어·반구저기·민귀군경 — 한국어에 살아 있는 맹자.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「인간은 본래 선하다」고 가장 단호히 말한 책.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
                '民 爲 貴',
                font_size=84, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.7),
                '社 稷 次 之',
                font_size=44, bold=True, color=SUB,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.7),
                '君 爲 輕',
                font_size=44, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.5), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '백 성 이  가 장  귀 하 다',
                font_size=22, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '사 직 이  그  다 음 이 요  ·  임 금 은  가 장  가 볍 다',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
                '— 맹자 제7편 진심(盡心) 하, 마지막 가까이',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '孟  子',
                font_size=20, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\맹자.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
