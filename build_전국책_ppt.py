# -*- coding: utf-8 -*-
"""
전국책(戰國策) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '시대 이름을 만든 책 · 동아시아 변론·외교·인간학의 원천',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '戰 國 策',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '전 국 책',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '세 치 혀가 백만 군대보다 강하다',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '유향(劉向, 전한) 편 · 12국 33편 · 약 460장 · 12만 자',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '소진·장의·맹상군·신릉군·노중련·형가 — 격변기 250년의 변사(辯士)들',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 책 이름이 시대 이름이 되다'),
        ('Ⅱ.', '시대 배경 — 전국시대 250년'),
        ('Ⅲ.', '책의 정체 — 유향의 편집'),
        ('Ⅳ.', '12국 33편의 구성'),
        ('Ⅴ.', '두 주인공 — 소진과 장의'),
        ('Ⅵ.', '전국 사군자 — 식객의 정치'),
    ]
    items_right = [
        ('Ⅶ.', '의인들 — 노중련·예양·형가'),
        ('Ⅷ.', '책사들 — 범저·상앙'),
        ('Ⅸ.', '종횡가의 사상'),
        ('Ⅹ.', '명구와 고사성어 16선'),
        ('Ⅺ.', '후대 영향과 한국'),
        ('Ⅻ.', '마무리 — 5가지 메시지'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_what(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '전국책(戰國策) — 어떤 책인가',
              '전국시대의 책략집 · 동아시아 산문의 영원한 모범')
    add_paragraphs(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.5), [
        ('· 전국시대 250년의 외교·변론·일화를 12개 나라별로 정리한 역사 산문집',
         {'font_size': 18, 'space_before': 6}),
        ('· 전한(前漢)의 학자 유향(劉向, BC 77?~6)이 흩어진 자료를 모아 편집',
         {'font_size': 18, 'space_before': 6}),
        ('· 12국 33편 · 약 460장 · 총 약 12만 자',
         {'font_size': 18, 'space_before': 6}),
        ('· 종횡가(縱橫家) — 소진·장의의 합종연횡 외교 책략의 결정판',
         {'font_size': 18, 'space_before': 6}),
        ('· 사마천 『사기』 전국시대 서술의 1차 근거 — 사기보다 먼저 그 자료가 있었다',
         {'font_size': 18, 'space_before': 6}),
        ('· 「전국시대(戰國時代)」라는 시대 이름 자체가 이 책에서 비롯됨',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한 책의 이름이 한 시대의 이름이 되다',
              '유향의 작명 — 「戰 國 의 策」')
    add_filled_rect(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.1), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.1),
                '戰 國 = 전쟁하는 나라들      策 = 책략·계책',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('● 유향이 정리하기 전 — 이 시대를 부르는 통일된 이름이 없었다',
         {'font_size': 18, 'space_before': 6}),
        ('● 사마천 『사기』조차 이 시대를 따로 명명하지 않았다',
         {'font_size': 18, 'space_before': 10, 'color': SUB}),
        ('● 유향이 「戰國의 策」 — 「전쟁하는 나라들의 책략집」으로 이름 붙임',
         {'font_size': 18, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 후세 학자들이 이 이름을 따라 시대 자체를 「전국시대」라 부르게 됨',
         {'font_size': 18, 'space_before': 10}),
        ('● 「한 책의 이름이 한 시대의 이름이 된」 드문 사례',
         {'font_size': 17, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC1)
def i_status(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 전국책')
    rows = [
        ('편자',     '유향(劉向, BC 77?~6) — 전한 종실(宗親)·학자'),
        ('시대 범위',  '전국시대 (BC 476경 ~ BC 222, 약 250년)'),
        ('구성',     '12국 33편(권) · 약 460장 · 약 12만 자'),
        ('대상 12국', '동주·서주·진·제·초·조·위·한·연·송·위(衛)·중산'),
        ('성격',     '유세가·종횡가의 책략·언설·일화 모음 — 역사 산문집'),
        ('사료 위상',  '전국시대의 거의 유일한 1차 사료군'),
        ('문학 위상',  '중국 산문 문학의 가장 완성된 모범 — 변론체의 정점'),
        ('출토 검증',  '1973년 마왕퇴 한묘 출토 백서로 일부 내용 검증됨'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅱ. 시대 배경 ==============
SEC2 = 'Ⅱ. 시대 배경'

@S(SEC2)
def ii_era(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '전국시대 — 중국 역사상 가장 격렬했던 250년',
              'BC 453 삼가분진 → BC 221 진의 천하 통일')
    cols = [
        ('春秋 춘추',
         'BC 770 ~ 476\n\n주(周) 왕실 명목상 인정\n100여 개 제후국\n\n「패자(覇者)」의 시대\n동맹과 회맹'),
        ('戰國 전국',
         'BC 476 ~ 221\n\n주 왕실 무시\n7개 강국으로 정리\n\n「왕(王)·제(帝)」의 시대\n천하 통일을 노림'),
    ]
    for i, (han, body) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), ACCENT if i == 1 else INK)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(1.0),
                    han, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_paragraphs(slide, x + Inches(0.5), Inches(3.5), Inches(4.9), Inches(3.5),
                       [(body, {'font_size': 16, 'color': INK})], line_spacing=1.5)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '시작 — 三家分晉(삼가분진) : 진(晉)이 한·위·조 세 가문에 분할 → 「신하가 군주를 삼키는 시대」의 신호탄',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC2)
def ii_seven(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '전국 칠웅(七雄) — 250년을 호령한 일곱 강국')
    rows = [
        ('秦 진', '서쪽 (관중)',   '함곡관 천연 요새 · 상앙 변법 · 최후의 승자', 'BC 221 천하 통일'),
        ('楚 초', '남쪽',         '양자강 유역 거대 영토 · 초사(楚辭) 문화',    'BC 223 진에 멸망'),
        ('齊 제', '동쪽 (산동)',  '바다·소금·철의 부 · 직하학궁의 학문 중심',  'BC 221 진에 항복'),
        ('燕 연', '북쪽 (북경)',  '약소국이나 의리의 나라 · 형가의 진왕 암살', 'BC 222 진에 멸망'),
        ('趙 조', '중북부',       '호복기사(胡服騎射) 군사 개혁',             'BC 228 진에 멸망'),
        ('魏 위', '중원',         '전국 초기 최강국 → 점차 쇠퇴',             'BC 225 진에 멸망'),
        ('韓 한', '중원 (정주)',  '가장 작은 강국 · 한비자의 고향',           'BC 230 가장 먼저 멸망'),
    ]
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.5), INK)
    headers = [('나라', 1.2), ('위치', 2.0), ('특징', 6.3), ('운명', 2.5)]
    x = Inches(0.7)
    for label, w in headers:
        add_textbox(slide, x, Inches(2.2), Inches(w), Inches(0.5),
                    label, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(w)
    for i, (han, loc, feat, fate) in enumerate(rows):
        y = Inches(2.7 + i * 0.6)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.6), bg)
        add_textbox(slide, Inches(0.7), y, Inches(1.2), Inches(0.6),
                    han, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(1.9), y, Inches(2.0), Inches(0.6),
                    loc, font_size=12, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.9), y, Inches(6.3), Inches(0.6),
                    feat, font_size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(10.2), y, Inches(2.5), Inches(0.6),
                    fate, font_size=12, color=ACCENT, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC2)
def ii_society(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '전국시대의 사회상 — 「능력만이 길이다」',
              '출신과 신분이 무너진 시대 · 평민의 천하 등용')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 춘추까지 — 귀족(士) 계급만이 정치에 참여',
         {'font_size': 17, 'space_before': 4}),
        ('● 전국시대 — 평민·외국 출신이라도 능력만 있으면 재상이 될 수 있다',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 대표 인물', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     · 소진 — 시골 출신, 형수에게 멸시받던 가난한 청년 → 6국의 재상',
         {'font_size': 15, 'space_before': 4}),
        ('     · 장의 — 빈털터리 떠돌이 → 진나라의 재상',
         {'font_size': 15, 'space_before': 4}),
        ('     · 범저 — 위에서 매 맞아 죽다 살아남은 자 → 진의 재상',
         {'font_size': 15, 'space_before': 4}),
        ('     · 상앙 — 위 출신 → 진 변법을 주도해 진을 강국으로',
         {'font_size': 15, 'space_before': 4}),
        ('● 「실력 있는 자가 천하를 얻는다」는 새로운 사회 윤리의 등장',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


# ============== Ⅲ. 책의 정체 ==============
SEC3 = 'Ⅲ. 책의 정체'

@S(SEC3)
def iii_liu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '편자 유향(劉向, BC 77?~6)',
              '한 황실의 종친 · 천록각 도서관의 학자')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 한 고조 유방의 동생 유교(劉交)의 4대손 — 황실 종친',
         {'font_size': 17, 'space_before': 4}),
        ('● 한 성제(成帝)의 명을 받아 천록각(天祿閣) 황실 도서관의 옛 전적 정리',
         {'font_size': 17, 'space_before': 10}),
        ('● 그 결과물 — 『별록(別錄)』 (한대까지의 모든 책의 해제집)',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 함께 편찬한 책들', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 『전국책(戰國策)』 — 전국시대 책략집',
         {'font_size': 15, 'space_before': 4}),
        ('     · 『설원(說苑)』 — 고사 모음집',
         {'font_size': 15, 'space_before': 4}),
        ('     · 『신서(新序)』 — 정치 일화집',
         {'font_size': 15, 'space_before': 4}),
        ('     · 『열녀전(列女傳)』 — 여성 전기집',
         {'font_size': 15, 'space_before': 4}),
        ('● 단지 자료를 모은 사람이 아닌 — 「옛 자료로 한대 사회에 교훈을 주려 한 학자」',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC3)
def iii_six(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '여섯 가지 이름으로 흩어졌던 자료',
              '유향이 천록각에서 발견한 6종의 단편')
    rows = [
        ('國 策 국책',  '나라별 책략집'),
        ('國 事 국사',  '나라의 일'),
        ('短 長 단장',  '짧고 긴 글 (변론의 기교)'),
        ('事 語 사어',  '사건과 말'),
        ('長 書 장서',  '긴 문서'),
        ('脩 書 수서',  '다듬은 문서'),
    ]
    for i, (han, desc) in enumerate(rows):
        y = Inches(2.4 + i * 0.65)
        add_filled_rect(slide, Inches(0.7), y, Inches(4.0), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(4.0), Inches(0.55),
                    han, font_size=17, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.9), y + Inches(0.05), Inches(8.0), Inches(0.5),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
                '유향이 중복을 정리하고 나라별·연대순으로 재편집 → 33편의 『전국책』으로 통합',
                font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅳ. 12국 33편 ==============
SEC4 = 'Ⅳ. 12국 33편'

@S(SEC4)
def iv_table(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '12국 33편의 구성',
              '동주·서주에서 시작해 중산국까지 — 한 시대의 모든 나라')
    rows = [
        ('권1',     '東周策 동주책', '1편', '망해가는 천자의 외교'),
        ('권2',     '西周策 서주책', '1편', '서주의 외교 — 천자국의 몰락'),
        ('권3~7',   '秦策 진책',   '5편', '진의 변법·외교·통일의 길'),
        ('권8~13',  '齊策 제책',   '6편', '맹상군·노중련·직하학궁 — 가장 많은 분량'),
        ('권14~17', '楚策 초책',   '4편', '남방 거대국·굴원의 비극'),
        ('권18~21', '趙策 조책',   '4편', '호복기사·인상여·평원군'),
        ('권22~25', '魏策 위책',   '4편', '문후의 전성기와 신릉군의 호부 사건'),
        ('권26~28', '韓策 한책',   '3편', '약소국의 외교·한비자의 고향'),
        ('권29~31', '燕策 연책',   '3편', '의리의 나라·형가의 진왕 암살'),
        ('권32',    '宋·衛策',    '1편', '약소국의 마지막 이야기'),
        ('권33',    '中山策 중산책', '1편', '사라진 작은 나라의 기록'),
    ]
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.45), INK)
    headers = [('권', 1.0), ('편명', 2.4), ('편수', 0.8), ('핵심 내용', 7.8)]
    x = Inches(0.7)
    for label, w in headers:
        add_textbox(slide, x, Inches(2.2), Inches(w), Inches(0.45),
                    label, font_size=13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(w)
    for i, (vol, name, num, desc) in enumerate(rows):
        y = Inches(2.65 + i * 0.42)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.42), bg)
        add_textbox(slide, Inches(0.7), y, Inches(1.0), Inches(0.42),
                    vol, font_size=12, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.7), y, Inches(2.4), Inches(0.42),
                    name, font_size=13, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.1), y, Inches(0.8), Inches(0.42),
                    num, font_size=12, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.0), y, Inches(7.7), Inches(0.42),
                    desc, font_size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_message(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '편 구성이 보여주는 것',
              '큰 나라부터 작은 나라까지 · 흥에서 망까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 12국의 분량 비대칭 — 제(6편) > 진(5편) > 초·조·위(4편) > 한·연(3편) > 동주·서주·송·중산(1편)',
         {'font_size': 17, 'space_before': 6}),
        ('● 분량이 곧 그 나라의 외교·문화·일화의 풍부함을 반영',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 제책 — 사군자 맹상군·직하학궁·노중련 → 「전국시대 지성사의 보고」',
         {'font_size': 17, 'space_before': 10}),
        ('● 진책 — 변법·외교·통일의 길 → 「변방국이 천하를 통일하는 250년의 과정」',
         {'font_size': 17, 'space_before': 10}),
        ('● 동주·서주·송·중산 — 작은 나라들의 외교 → 「약소국의 생존 모범」',
         {'font_size': 17, 'space_before': 10}),
        ('● 한·위·조 — 진과 가장 가까이 있던 나라들 → 「가장 먼저 멸망한 자들의 기록」',
         {'font_size': 17, 'space_before': 10}),
    ])


# ============== Ⅴ. 두 주인공 ==============
SEC5 = 'Ⅴ. 두 주인공'

@S(SEC5)
def v_meet(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '소진과 장의 — 같은 스승, 정반대의 길',
              '귀곡선생(鬼谷先生) 문하의 두 동기')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 전국시대의 전설적 스승 귀곡선생(鬼谷先生) — 종횡가의 비조(鼻祖)',
         {'font_size': 17, 'space_before': 4}),
        ('● 종횡술(縱橫術) — 외교와 변론의 기술 — 두 사람이 함께 배움',
         {'font_size': 17, 'space_before': 10}),
        ('● 그러나 두 사람이 간 길은 정반대였다', {'font_size': 18, 'space_before': 14, 'bold': True, 'color': ACCENT}),
        ('     · 소진 → 합종(合縱) — 6국이 진에 맞서 동맹', {'font_size': 16, 'space_before': 6}),
        ('     · 장의 → 연횡(連橫) — 진이 6국을 차례로 분열시킴', {'font_size': 16, 'space_before': 4}),
        ('● 두 사람의 대립은 「약자들의 연합 vs 강자 중심 분열」이라는 외교 전략의 두 원형',
         {'font_size': 16, 'space_before': 14, 'color': SUB, 'bold': True}),
        ('● 이 대립은 오늘날까지도 모든 외교의 기본 패턴 (NATO vs 양자 동맹 등)',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC5)
def v_sojin(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '소진(蘇秦) — 6국의 재상 인(印)을 동시에 차다',
              '낙양 평민에서 6국 합종의 영웅으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 낙양 출신 가난한 청년 — 진에 가서 유세했으나 받아들여지지 않음',
         {'font_size': 17, 'space_before': 4}),
        ('● 빈털터리로 귀향 — 가족과 형수에게조차 멸시 당함',
         {'font_size': 17, 'space_before': 8, 'color': SUB}),
        ('● 『음부경(陰符經)』 1년간 다시 공부 — 졸리면 송곳으로 허벅지를 찔러 깨움',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     「자고현량(刺股懸樑)」 — 머리를 대들보에 매고 송곳으로 허벅지를 찌르며 공부',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 연·조·한·위·제·초 — 6국 군주를 차례로 만나며 합종(合縱) 책략 설파',
         {'font_size': 17, 'space_before': 12}),
        ('● 합종 성공 — 6국 재상의 인을 동시에 차다 (중국 역사상 전무후무)',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 진은 이후 15년간 함곡관 밖으로 군대를 내지 못했다',
         {'font_size': 17, 'space_before': 10}),
        ('● 후에 제나라에서 자객의 칼에 — 자기 시신을 미끼로 자객을 잡아낸 마지막 책략',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC5)
def v_zhangyi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '장의(張儀) — 「내 혀가 아직 있느냐」',
              '모욕에서 시작된 연횡(連橫)의 사기꾼 외교가')
    add_filled_rect(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.5), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(0.7),
                '吾 舌 尙 在',
                font_size=32, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.8),
                '내 혀가 아직 있다',
                font_size=20, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(3.0), [
        ('● 초나라 재상 연회 — 벽옥(璧玉)이 사라지자 의심받아 매를 맞고 쫓겨남',
         {'font_size': 16, 'space_before': 6}),
        ('● 피투성이가 된 그가 아내에게 — 「내 혀가 아직 있느냐?」 → 「있습니다」 → 「그럼 됐다」',
         {'font_size': 16, 'space_before': 8, 'bold': True, 'color': ACCENT}),
        ('● 진 혜문왕에게 등용 — 연횡(連橫) 책략으로 6국 합종을 차례로 깨뜨림',
         {'font_size': 16, 'space_before': 10}),
        ('● 초 회왕을 속여 「600리 땅을 주마」 → 실제로는 6리만 — 외교 사기의 정점',
         {'font_size': 16, 'space_before': 8, 'color': SUB}),
        ('● 결과 — 6국 합종 붕괴 → 진의 통일 토대 마련',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC5)
def v_compare(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '소진 vs 장의 — 합종(合縱) vs 연횡(連橫)')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.5), INK)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(2.6), Inches(0.5),
                '항목', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(3.3), Inches(2.2), Inches(4.7), Inches(0.5),
                '소진(蘇秦) — 합종', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(8.0), Inches(2.2), Inches(4.7), Inches(0.5),
                '장의(張儀) — 연횡', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('의미',  '남북으로 합친다 (縱)',     '동서로 잇는다 (橫)'),
        ('전략',  '6국 동맹 → 진에 대항',     '진 + 1국씩 → 6국 분열'),
        ('누구를 위해', '6국 (약자들)',         '진 (강자)'),
        ('출신',  '낙양 평민',              '위나라 평민'),
        ('결과',  '15년간 진을 묶음',        '6국 합종 붕괴'),
        ('역사적 의미', '「약자 연합」의 원형', '「강자 중심 분열책」의 원형'),
    ]
    for i, (k, s, j) in enumerate(rows):
        y = Inches(2.75 + i * 0.6)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.6), bg)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.6),
                    k, font_size=14, bold=True, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.3), y, Inches(4.7), Inches(0.6),
                    s, font_size=14, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(8.0), y, Inches(4.7), Inches(0.6),
                    j, font_size=14, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅵ. 전국 사군자 ==============
SEC6 = 'Ⅵ. 전국 사군자'

@S(SEC6)
def vi_four(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '전국 사군자(四君子) — 식객 3,000명을 거느린 4대 공자',
              '「천하의 인재가 모이는 작은 정부」')
    rows = [
        ('孟嘗君 맹상군', '제(齊)', '3,000명', '계명구도 · 교토삼굴 · 풍훤의 의리'),
        ('平原君 평원군', '조(趙)', '수천 명', '모수자천 · 한단 포위에서 초에 구원'),
        ('信陵君 신릉군', '위(魏)', '수천 명', '호부(虎符) 탈취로 조나라 구함'),
        ('春申君 춘신군', '초(楚)', '3,000명', '평원군의 사신을 식객 화려함으로 압도'),
    ]
    for i, (name, country, num, episode) in enumerate(rows):
        y = Inches(2.4 + i * 1.05)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.0), Inches(0.85), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(3.0), Inches(0.85),
                    name, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_filled_rect(slide, Inches(3.9), y, Inches(1.4), Inches(0.85), PALE)
        add_textbox(slide, Inches(3.9), y, Inches(1.4), Inches(0.85),
                    country, font_size=15, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(5.5), y + Inches(0.1), Inches(1.6), Inches(0.65),
                    num, font_size=14, color=SUB, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.3), y + Inches(0.1), Inches(5.6), Inches(0.65),
                    episode, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC6)
def vi_meng(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '맹상군과 계명구도(鷄鳴狗盜)',
              '제책 — 「천한 재주가 결정적 순간에 사람을 살린다」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 맹상군이 진(秦) 소왕에게 잡혀 죽을 위기',
         {'font_size': 18, 'space_before': 4}),
        ('● 식객 중 한 사람 — 개 도둑 흉내로 호백구(여우 가죽옷)를 훔쳐 진왕 애첩에게 바침',
         {'font_size': 17, 'space_before': 12}),
        ('● 풀려나 도망 → 함곡관에 도착했으나 새벽이 되지 않아 문이 닫혀 있음',
         {'font_size': 17, 'space_before': 10}),
        ('● 다른 식객 — 닭 울음을 흉내내자 마을 닭이 함께 울어 문이 열림 → 탈출 성공',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 「계명구도(鷄鳴狗盜)」 — 닭 울음과 개 도둑의 잔재주 — 그러나 살린 것은 사람',
         {'font_size': 17, 'space_before': 12, 'color': SUB, 'bold': True}),
        ('● 후일 풍훤(馮諼)의 「교토삼굴(狡兎三窟)」 — 영악한 토끼는 세 굴을 판다',
         {'font_size': 16, 'space_before': 12, 'font_name': 'Batang'}),
    ])


@S(SEC6)
def vi_shin(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC6, n, t)
    add_title(slide, '신릉군과 호부(虎符) 탈취',
              '위책 — 친구를 위해 자기 나라의 법을 어긴 의리의 결단')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 진(秦)이 조(趙)의 한단(邯鄲)을 포위 — 평원군이 친척인 신릉군에게 구원 요청',
         {'font_size': 17, 'space_before': 4}),
        ('● 위 안리왕이 응하지 않자, 신릉군은 결단',
         {'font_size': 17, 'space_before': 10}),
        ('● 식객 후영(侯嬴)의 책략 — 왕의 침실에서 호부(虎符, 군대 동원의 부절)를 훔침',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 협조하지 않는 장수 진비(晉鄙)는 식객 주해(朱亥)가 철퇴로 죽임',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 군대를 이끌고 한단으로 — 조나라를 구하고 진의 포위를 풀어냄',
         {'font_size': 17, 'space_before': 10}),
        ('● 친구를 위해 자기 나라의 법을 어긴 「의리(義)의 정점」',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 「선비는 자기를 알아주는 자를 위해 죽는다」 — 식객 정치의 가장 극적인 장면',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅶ. 의인들 ==============
SEC7 = 'Ⅶ. 의인들'

@S(SEC7)
def vii_lu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '노중련(魯仲連) — 의(義)의 화신',
              '제책 — 「의불제진(義不帝秦)」, 차라리 동해에 빠져 죽겠다')
    add_filled_rect(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5),
                '義 不 帝 秦',
                font_size=44, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(3.0), [
        ('● 진(秦)이 조(趙)의 한단을 포위 — 위(魏) 사신이 「진을 황제로 인정하면 풀어준다」 권유',
         {'font_size': 16, 'space_before': 6}),
        ('● 노중련의 일갈 — 「차라리 동해에 빠져 죽을지언정 진의 백성은 되지 않으리라」',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 분위기가 바뀌어 조는 항복 거부 → 신릉군의 호부 군대가 와서 포위가 풀림',
         {'font_size': 16, 'space_before': 10}),
        ('● 조에서 봉토를 주려 하자 모두 사양 — 「남의 환난을 풀어주고도 취하지 않는 것이 선비」',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 전국시대 가장 순수한 의(義)의 화신 — 도덕적 인격이 가장 완성된 인물',
         {'font_size': 15, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


@S(SEC7)
def vii_yu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '예양(豫讓) — 사위지기자사(士爲知己者死)',
              '조책 — 「선비는 자기를 알아주는 자를 위해 죽는다」')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5),
                '士 爲 知 己 者 死  女 爲 悅 己 者 容',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(3.0), [
        ('● 주군 지백(智伯)이 조양자(趙襄子)에게 살해당함 — 예양은 복수를 결심',
         {'font_size': 16, 'space_before': 6}),
        ('● 자기 몸에 옻칠을 해 문둥이로 변장, 숯을 삼켜 벙어리가 되어 시장에서 구걸',
         {'font_size': 16, 'space_before': 10}),
        ('● 조양자의 행차에 칼을 들고 다리 밑에 숨음 → 발각',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 「선비는 자기를 알아주는 자를 위해 죽는다」 — 잡힌 그가 한 말',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 처형 전 조양자의 옷을 빌려 세 번 찔러 명목상 복수 후 자결',
         {'font_size': 16, 'space_before': 10}),
        ('● 신의(信義)의 극치 — 후세에 무수히 인용된 명문',
         {'font_size': 15, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


@S(SEC7)
def vii_jingke(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '형가(荊軻) — 진왕 암살의 비극',
              '연책 마지막 · 전국시대 250년의 클라이맥스')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.7), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.7),
                '風 蕭 蕭 兮 易 水 寒',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(12.0), Inches(0.7),
                '壯 士 一 去 兮 不 復 還',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(3.0), [
        ('● 「바람은 쓸쓸하고 역수는 차구나, 장사 한 번 가면 다시 돌아오지 못하리」',
         {'font_size': 16, 'space_before': 6}),
        ('● 연 태자 단(丹)이 보낸 자객 — 진왕(시황) 암살을 위해 떠나는 날의 노래',
         {'font_size': 16, 'space_before': 10}),
        ('● 친구 고점리(高漸離)가 축(筑, 악기)을 켜고 형가가 노래하던 역수(易水) 강가',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 진궁에서 지도를 펴는 순간 비수로 진왕을 찔렀으나 실패 — 웃으며 죽음',
         {'font_size': 16, 'space_before': 10}),
        ('● 5년 후 연 멸망 — 그러나 「이루지 못할 것을 알면서도 의리를 위해 가는 자」의 모범',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 동아시아 비장미(悲壯美)의 최고 명문구',
         {'font_size': 14, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅷ. 책사들 ==============
SEC8 = 'Ⅷ. 책사들'

@S(SEC8)
def viii_fan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '범저(范雎) — 원교근공(遠交近攻)',
              '진책 — 한 구절이 천하 통일의 청사진이 되다')
    add_filled_rect(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.4), PALE)
    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.4),
                '遠 交 近 攻',
                font_size=48, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(3.5), [
        ('● 위(魏) 출신 — 재상 위제의 모함으로 매 맞아 거의 죽고 변소에 버려져 오줌 세례',
         {'font_size': 16, 'space_before': 4}),
        ('● 이름을 장록(張祿)으로 바꾸고 진(秦)으로 도망 → 진 소왕(昭王)에게 등용',
         {'font_size': 16, 'space_before': 8}),
        ('● 책략 — 「멀리 있는 나라(齊·楚)와 사귀고, 가까운 나라(韓·魏)부터 친다」',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 「한 치를 얻으면 왕의 한 치, 한 자를 얻으면 왕의 한 자」 — 누적의 전략',
         {'font_size': 15, 'space_before': 8, 'color': SUB, 'font_name': 'Batang'}),
        ('● 진은 이 책략대로 한 → 조 → 위 → 초 → 연 → 제 순서로 통일 (BC 230~221)',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 한 책의 한 구절이 250년의 역사를 결정 — 「언어의 가장 큰 위력」',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC8)
def viii_yan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '연 소왕과 곽외(郭隗) — 천금매골(千金買骨)',
              '연책 — 「큰일은 작고 가까운 것에서 시작된다」')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 연(燕)이 제(齊)에 큰 패배를 당함 — 연 소왕이 인재를 모으려 함',
         {'font_size': 17, 'space_before': 4}),
        ('● 신하 곽외(郭隗)의 비유 — 「천금으로 죽은 천리마의 뼈를 산 임금」',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 「임금이 죽은 말의 뼈도 천금에 사니, 산 천리마는 더 비싸겠다」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 「오시지 않을 줄 알았던 천리마들이 1년 안에 셋이 왔다」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 「선시어외(先始於隗)」 — 「먼저 저(곽외)부터 후하게 대접하십시오」',
         {'font_size': 17, 'space_before': 14, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 효과 — 악의(樂毅)·추연(鄒衍)·극신(劇辛) 등이 모여 연이 일시 강국으로',
         {'font_size': 17, 'space_before': 12}),
        ('● 「인재 영입은 가장 가까운 데서, 가장 큰 성의로」',
         {'font_size': 16, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


# ============== Ⅸ. 종횡가의 사상 ==============
SEC9 = 'Ⅸ. 종횡가의 사상'

@S(SEC9)
def ix_who(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '종횡가(縱橫家) — 외교와 변론의 전문가',
              '백가(百家) 중 외교·설득을 전문으로 한 학파')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 이름의 유래 — 합종(合縱) + 연횡(連橫)의 두 글자',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 핵심 가치', {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 현실주의 — 도덕 명분보다 이해관계 분석',
         {'font_size': 16, 'space_before': 6}),
        ('     · 세(勢)의 통찰 — 세력 균형과 흐름을 읽는 능력',
         {'font_size': 16, 'space_before': 4}),
        ('     · 변론과 설득 — 말 한 마디로 천하를 움직임',
         {'font_size': 16, 'space_before': 4}),
        ('     · 유연성 — 상황에 따라 책략을 바꿈',
         {'font_size': 16, 'space_before': 4}),
        ('     · 공리주의 — 결과가 모든 것을 정당화',
         {'font_size': 16, 'space_before': 4}),
        ('● 비조(鼻祖) — 귀곡선생(鬼谷先生) · 대표 인물 — 소진·장의·범저',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
    ])


@S(SEC9)
def ix_judge(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '후대의 두 평가 — 폄하와 옹호',
              '도덕가의 비판 vs 문장가의 옹호')
    cols = [
        ('정자·주자 (송 성리학)',
         '「권모술수의 책」\n\n명분과 도덕보다\n술수를 가르친다는 비판\n\n도덕적 평가의 거부'),
        ('소동파 (송 대문호)',
         '「문장 그 자체로 영원한 모범」\n\n도덕적 평가와 별개로\n글솜씨와 인간 통찰은\n후세 어떤 책도 따라가지 못함'),
    ]
    for i, (title, body) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        color = INK if i == 0 else ACCENT
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(1.0),
                    title, font_size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.5), Inches(3.6), Inches(4.9), Inches(3.5),
                       [(body, {'font_size': 15, 'color': INK, 'align': PP_ALIGN.CENTER})],
                       line_spacing=1.6)
    add_textbox(slide, Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.4),
                '두 평가는 모두 옳다 — 도덕서로 읽으면 위험, 「인간과 권력의 거울」로 읽으면 무한한 통찰',
                font_size=14, color=SUB, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅹ. 명구·고사성어 16선 ==============
SEC10 = 'Ⅹ. 명구·고사성어 16선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC10)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC10} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=36, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=17, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 16, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('초책 — 진진(陳軫)이 소양(昭陽)에게', '蛇 足',
     '뱀에 다리를 그리다',
     '술 한 잔을 두고 뱀 그리기 시합 — 가장 먼저 그린 자가 「발도 그릴 수 있다」 자랑하다 술을 빼앗긴 일화. 「과욕·군더더기로 일을 망친다」.'),
    ('연책 — 소대(蘇代)가 조 혜문왕에게', '漁 父 之 利',
     '어부의 이익',
     '조개와 도요새가 다투는 사이 어부가 잡아간 비유. 조·연이 다투지 말 것을 권유. 「둘이 다투면 제3자가 이득을 본다」 — 외교의 영원한 경고.'),
    ('초책 — 강을(江乙)이 초 선왕에게', '狐 假 虎 威',
     '여우가 호랑이의 위세를 빌리다',
     '호랑이가 여우를 잡자 여우가 「내가 백수의 왕이다」 — 호랑이 뒤에 따라가니 모두 도망. 「남의 권세를 빌려 위세 부리는 자」의 원형.'),
    ('위책 — 방총(龐葱)이 위 혜왕에게', '三 人 成 虎',
     '세 사람이면 호랑이가 생긴다',
     '시장에 호랑이가 있다 한 사람이 말하면 안 믿고, 두 사람이면 의심, 세 사람이면 믿는다. 「거짓도 여럿이 말하면 사실이 된다」 — 가장 강력한 여론·소문의 비유.'),
    ('위책 — 계량(季梁)이 위 안리왕에게', '南 轅 北 轍',
     '수레는 남으로 가는데 끌채는 북으로',
     '초나라로 간다면서 북쪽으로 가는 자 — 「말이 빠르다·노자가 많다·마부가 능하다」 자랑할수록 더 멀어진다. 「방향이 틀린 노력의 어리석음」.'),
    ('제책 — 맹상군의 식객', '鷄 鳴 狗 盜',
     '닭 울음과 개 도둑',
     '맹상군이 진에서 탈출할 때 호백구를 훔친 자와 닭 울음 흉내내는 자가 살림. 「천한 재주도 결정적 순간엔 쓸모」 — 「인재는 평가의 그릇에 따라 다른 모양」.'),
    ('제책 — 풍훤(馮諼)의 책략', '狡 兎 三 窟',
     '영악한 토끼는 세 굴을 판다',
     '맹상군의 식객 풍훤이 봉지의 빚 문서를 불태우고, 위에 추천하고, 종묘를 세우게 한 「세 굴」. 「위험에 대비한 다중 안전장치」 — 전국시대 리스크 관리의 정수.'),
    ('조책 — 평원군에게 자청한 모수', '毛 遂 自 薦',
     '모수가 스스로 자기를 천거하다',
     '한단 포위 시 평원군이 초로 구원을 청할 사절 20명 중 하나가 부족 — 식객 모수가 「저를 데려가십시오」. 결국 그가 가장 큰 공을 세움. 「실력자는 자기 입을 연다」.'),
    ('조책 — 같은 일화에서', '囊 中 之 錐',
     '주머니 속의 송곳',
     '평원군이 「뛰어난 자는 주머니 속 송곳처럼 저절로 드러난다」 하자 모수가 「저는 주머니에 들지 못했을 뿐」. 「진짜 실력은 끝내 드러난다」.'),
    ('연책 — 곽외의 비유', '千 金 買 骨',
     '천금으로 말의 뼈를 사다',
     '천리마를 구하던 임금이 죽은 천리마의 뼈를 천금에 샀더니 1년 안에 산 천리마 셋이 왔다는 비유. 「인재 영입은 큰 성의에서 시작」.'),
    ('연책 — 곽외 일화의 결어', '先 始 於 隗',
     '먼저 저(곽외)부터 시작하소서',
     '「임금이 천하의 인재를 부르고자 하시면, 먼저 저부터 후하게 대접하십시오. 그러면 저보다 뛰어난 자들이 천리를 마다하지 않고 올 것입니다」. 「큰 일은 가까운 작은 일부터」.'),
    ('진책 — 범저의 책략', '遠 交 近 攻',
     '먼 곳과 사귀고 가까운 곳을 친다',
     '한 책의 한 구절이 250년 역사의 청사진. 진의 통일 순서(한→조→위→초→연→제)가 모두 이 한 줄에서 나옴. 「외교 우선순위 설정의 영원한 원칙」.'),
    ('소진과 장의의 전체 책략', '合 縱 連 橫',
     '합종과 연횡',
     '소진의 합종(6국 동맹) vs 장의의 연횡(진+1국씩 분열). 「약자 연합 vs 강자 중심 분열」의 영원한 외교 패턴 — 오늘 NATO·양자 동맹 논쟁의 원형.'),
    ('연책 — 노중련의 한 마디', '義 不 帝 秦',
     '의로움으로 진을 황제로 인정하지 않는다',
     '「차라리 동해에 빠져 죽을지언정 진의 백성은 되지 않으리」 — 의(義)의 가장 순수한 표현. 도덕적 거부의 결정적 한 줄.'),
    ('연책 — 형가의 노래', '風 蕭 蕭 兮 易 水 寒  壯 士 一 去 兮 不 復 還',
     '바람은 쓸쓸하고 역수는 차구나, 장사 한 번 가면 다시 돌아오지 못하리',
     '진왕 암살을 위해 떠나는 형가의 비장한 결의. 동아시아 비장미(悲壯美)의 최고봉. 「이루지 못할 것을 알고도 가는 자」의 모범.'),
    ('조책 — 예양의 명언', '士 爲 知 己 者 死',
     '선비는 자기를 알아주는 자를 위해 죽는다',
     '예양이 지백을 위해 복수에 나서며 한 말. 「자기를 알아주는 사람」의 무게 — 동양 의리(義理) 사상의 가장 짧은 핵심 명제. 후대 무수한 인용의 출발점.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅺ. 후대 영향 ==============
SEC11 = 'Ⅺ. 후대 영향과 한국'

@S(SEC11)
def xi_sima(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '『사기』·당송팔대가 — 산문의 영원한 교과서',
              '사마천이 옮겨 쓴 책 · 당송 명문장의 원천')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 사마천 『사기』 — 전국시대 부분이 사실상 『전국책』의 재구성',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('     · 「소진열전」·「장의열전」·「범저채택열전」·「맹상군열전」·「자객열전」 등',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 한대 산문가들 — 변론체의 모범으로 학습',
         {'font_size': 17, 'space_before': 12}),
        ('● 당송팔대가 — 한유(韓愈)·유종원(柳宗元)·소식(蘇軾)·구양수 등 모두 전국책 문체에서 영향',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 명청 시대 — 모든 문장 학습의 필독서',
         {'font_size': 17, 'space_before': 10}),
        ('● 비유와 우화의 보고 — 16개 사자성어가 일상 언어에 살아 있음',
         {'font_size': 17, 'space_before': 10}),
        ('● 2,000년 동안 동아시아 산문의 영원한 교과서',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC11)
def xi_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '한국에서의 수용 — 조선 외교의 교본',
              '명·청 사대 외교 · 통신사 외교 · 변사의 학습서')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 조선 사대부의 필독서 — 외교 문서 작성과 변론의 교본',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 광해군 시기 — 명·청 사이의 외교 줄타기 = 「합종연횡의 한국적 적용」',
         {'font_size': 17, 'space_before': 12}),
        ('● 임진왜란 — 명나라 원군 청병 외교 = 「약소국 외교의 모범」 따름',
         {'font_size': 17, 'space_before': 10}),
        ('● 정약용 등 실학자 — 『전국책』을 깊이 학습',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 한국어 번역본', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 임동석 (동서문화사) — 학술 역주 / 진기환 (명문당) — 본격 완역',
         {'font_size': 14, 'space_before': 4}),
        ('     · 신동준 (인간사랑) — 일반 독자용 풀이',
         {'font_size': 14, 'space_before': 4}),
        ('● 우리는 알게 모르게 매일 전국책의 표현을 쓰고 있다 — 일상 언어 속에 살아 있는 책',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC11)
def xi_value(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '사료적 가치와 한계',
              '역사서이자 동시에 문학서 — 사실과 우화의 혼합')
    cols = [
        ('가치',
         '전국시대 250년의 거의 유일한 1차 사료\n\n외교·전쟁·인물·풍속·언어 망라\n\n사대부의 사고방식 직접 보존\n\n1973년 마왕퇴 백서로 일부 검증', ACCENT),
        ('한계',
         '윤색·과장 — 유세가의 자기 자랑\n\n연대·인물 착오 — 사기와 불일치\n\n문학적 가공 — 실화/창작 경계\n\n어부지리 등 일부는 우화로 보아야', INK),
    ]
    for i, (title, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(1.0),
                    title, font_size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.4), Inches(3.6), Inches(5.1), Inches(3.5),
                       [(body, {'font_size': 14, 'color': INK})], line_spacing=1.5)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '학자들은 『전국책』을 「역사서이자 동시에 문학서」로 본다 — 그 모호함이 책의 깊이를 만든다',
                font_size=14, color=SUB, bold=True, align=PP_ALIGN.CENTER)


# ============== Ⅻ. 마무리 ==============
SEC12 = 'Ⅻ. 마무리'

@S(SEC12)
def xii_messages(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '전국책이 일러주는 5가지 메시지')
    items = [
        ('1', '세 치 혀가 백만 군대보다 강하다',
         '소진의 합종이 진을 15년간 함곡관에 묶었다 — 한 사람의 말이 백만 군대보다 강한 시대'),
        ('2', '능력은 출신과 신분을 넘는다',
         '소진·장의·범저·이사 — 모두 평민·외국 출신. 출신을 벗어난 능력 사회의 첫 모범'),
        ('3', '외교는 도덕이 아니라 이해관계다',
         '영원한 동맹은 없다 — 오늘의 적이 내일의 동지. 종횡가의 현실주의 외교는 오늘도 유효'),
        ('4', '위기는 인격을 드러낸다',
         '형가의 결의·노중련의 의·신릉군의 우정·예양의 신의 — 평상시엔 안 보이던 인격들'),
        ('5', '역사는 인물의 선택이 만든다',
         '진의 통일은 우연이 아니다 — 상앙·장의·범저·이사의 선택이 쌓인 결과. 6국 멸망도 그들의 선택'),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(2.3 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(0.7), Inches(0.8), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(0.7), Inches(0.8),
                    num, font_size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(1.7), y + Inches(0.05), Inches(11.2), Inches(0.8), [
            (title, {'font_size': 16, 'bold': True, 'color': ACCENT}),
            (desc,  {'font_size': 13, 'color': INK, 'space_before': 4}),
        ], line_spacing=1.25)


@S(SEC12)
def xii_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '전국책, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 유향이 12국 33편으로 정리한 전국시대 250년의 변사·책사들의 책략집.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 책 이름이 시대 이름이 된 「유일한 사례」 — 전국시대(戰國時代).',
         {'font_size': 18, 'space_before': 8}),
        ('● 두 주인공 소진과 장의 — 합종(合縱)과 연횡(連橫)의 영원한 외교 패턴.',
         {'font_size': 18, 'space_before': 8}),
        ('● 사군자·노중련·예양·형가·범저 — 격동기 인간 군상의 모든 모습.',
         {'font_size': 18, 'space_before': 8}),
        ('● 사족·어부지리·호가호위·삼인성호·계명구도·교토삼굴·원교근공 — 일상에 살아 있는 책.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「세 치 혀가 천하를 움직인다」는 명제의 책.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC12)
def xii_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.1),
                '風 蕭 蕭 兮 易 水 寒',
                font_size=66, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.0),
                '壯 士 一 去 兮 不 復 還',
                font_size=66, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.4), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.6),
                '바람은 쓸쓸하고 역수는 차구나',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
                '장사 한 번 가면 다시는 돌아오지 못하리',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
                '— 형가(荊軻), 역수가(易水歌) · 전국책 연책',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '戰  國  策',
                font_size=18, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\전국책.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
