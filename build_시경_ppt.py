# -*- coding: utf-8 -*-
"""
시경(詩經) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '오경(五經)의 으뜸 · 동아시아 시가의 원류',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.6),
                '詩 經',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
                '시 경',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
                '思無邪 — 생각에 사악함이 없다',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '풍(風) 160 · 아(雅) 105 · 송(頌) 40 — 총 305편',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 시경이란 무엇인가'),
        ('Ⅱ.', '육의(六義) — 풍·아·송과 부·비·흥'),
        ('Ⅲ.', '국풍(國風) — 15국의 민요'),
        ('Ⅳ.', '아(雅) — 소아·대아'),
        ('Ⅴ.', '송(頌) — 종묘의 노래'),
        ('Ⅵ.', '명편 깊이 읽기'),
    ]
    items_right = [
        ('Ⅶ.', '시경의 사상'),
        ('Ⅷ.', '공자의 시교(詩敎)'),
        ('Ⅸ.', '명구 모음'),
        ('Ⅹ.', '해석 학파와 후대 영향'),
        ('Ⅺ.', '현대적 의의'),
        ('Ⅻ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '시경(詩經) — 중국 최초의 시가 총집',
              '약 500년의 노래, 305편으로 모이다')
    add_paragraphs(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.5), [
        ('· 중국에서 가장 오래된 시가 모음집, 동아시아 시가 문학의 원류',
         {'font_size': 18, 'space_before': 6}),
        ('· 서주(西周) 초기 ~ 춘추(春秋) 중기의 노래 305편을 수록 (약 500년)',
         {'font_size': 18, 'space_before': 6}),
        ('· 원래 「시(詩)」, 「시삼백(詩三百)」이라 부르다 한대(漢代)에 경전(經)이 되어 「시경」으로 격상',
         {'font_size': 18, 'space_before': 6}),
        ('· 공자가 3,000여 편에서 305편을 가려 정리했다는 전승 — 「산시(刪詩)」',
         {'font_size': 18, 'space_before': 6}),
        ('· 백성의 노래 · 궁중의 음악 · 종묘의 찬가까지 — 사회 전 영역의 정서를 수렴',
         {'font_size': 18, 'space_before': 6}),
        ('· 사실주의 문학의 시조이자, 인격 수양·언어 교육의 핵심 교재',
         {'font_size': 18, 'space_before': 6}),
    ])


@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '명칭의 변천 — 詩에서 經으로')
    rows = [
        ('시(詩)',     '원래의 이름. 단지 「노래·시가」를 뜻함'),
        ('시삼백(詩三百)', '공자가 「시 삼백 편(詩三百)」으로 부른 데서 통칭이 됨'),
        ('모시(毛詩)',  '한대(漢代) 모형(毛亨)·모장(毛萇)의 전수본이 표준이 되어 붙은 이름'),
        ('시경(詩經)',  '한대에 오경(五經)으로 격상되면서 정식 경전의 지위를 얻음'),
    ]
    for i, (name, desc) in enumerate(rows):
        y = Inches(2.3 + i * 0.95)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(0.75), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.4), Inches(0.75),
                    name, font_size=20, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.3), y + Inches(0.05), Inches(9.5), Inches(0.7),
                    desc, font_size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_compile(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '편찬 — 누가, 언제, 어떻게')
    add_paragraphs(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.5), [
        ('● 채시(采詩) — 천자의 사신이 각지를 돌며 민요를 수집했다는 제도', {'font_size': 19, 'bold': True, 'color': ACCENT}),
        ('     「왕이 풍속을 살피고 정치의 득실을 알기 위해 사방의 시를 모았다」', {'font_size': 17, 'color': SUB, 'space_before': 4}),
        ('● 헌시(獻詩) — 대부·공경이 풍자·풍간을 위해 시를 지어 임금에게 바침', {'font_size': 19, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     아(雅)·송(頌)의 일부가 이 경로로 형성', {'font_size': 17, 'color': SUB, 'space_before': 4}),
        ('● 산시(刪詩) — 공자가 3,000여 편에서 305편을 가렸다는 전승', {'font_size': 19, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     『사기·공자세가』의 기록 — 다만 학계 내 논쟁은 있음', {'font_size': 17, 'color': SUB, 'space_before': 4}),
        ('● 「예의에 합당하지 않은 것을 빼고 305편을 남겼다」 — 공자 평어', {'font_size': 17, 'color': INK, 'space_before': 14}),
    ])


@S(SEC1)
def i_era(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '시대와 규모')
    rows = [
        ('시대',   'BC 11세기 ~ BC 6세기경 (서주 초기 ~ 춘추 중기)'),
        ('규모',   '총 305편 (시삼백) — 가사만 남고 곡조는 실전'),
        ('구성',   '풍(風) 160편 + 아(雅) 105편 + 송(頌) 40편'),
        ('편명만 남은 시', '소아 6편 — 가사를 잃고 제목만 전해짐 (생시笙詩)'),
        ('형식',   '주로 4언(四言)을 기본으로, 1~9자가 섞임'),
        ('주제',   '연애·노동·전쟁·제사·풍자·송축까지 — 삶의 전 영역'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.55), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.55),
                    k, font_size=16, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.55),
                    v, font_size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_thought(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '핵심 사상 — 한 폭으로 보기')
    boxes = [
        ('詩 言 志', '시언지', '시는 마음에 품은 뜻을 노래한 것'),
        ('思 無 邪', '사무사', '시 305편을 한 마디로 — 생각에 사악함이 없다'),
        ('溫 柔 敦 厚', '온유돈후', '시교의 풍모 — 부드럽고 두터운 인격'),
        ('興 觀 群 怨', '흥관군원', '시의 네 가지 효용 — 일으키고·살피고·모으고·풀어내다'),
    ]
    for i, (han, kor, desc) in enumerate(boxes):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(6.0), Inches(2.0), PALE)
        add_textbox(slide, x, y + Inches(0.25), Inches(6.0), Inches(0.7),
                    han, font_size=30, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.0), Inches(6.0), Inches(0.4),
                    kor, font_size=15, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.4), Inches(6.0), Inches(0.5),
                    desc, font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 육의 ==============
SEC2 = 'Ⅱ. 육의(六義)'

@S(SEC2)
def ii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '육의(六義) — 시경을 이해하는 여섯 가지 틀',
              '내용 분류 셋(풍·아·송) + 표현 기법 셋(부·비·흥)')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.0), [
        ('「대서(大序)」 — 시에는 육의가 있으니, 풍(風)·부(賦)·비(比)·흥(興)·아(雅)·송(頌)이다',
         {'font_size': 19, 'bold': True, 'color': ACCENT}),
        ('', {'space_before': 6}),
        ('· 내용에 따른 세 분류 — 풍(風) · 아(雅) · 송(頌)', {'font_size': 18, 'space_before': 8}),
        ('· 표현 기법의 세 가지 — 부(賦) · 비(比) · 흥(興)', {'font_size': 18, 'space_before': 6}),
        ('· 시경 305편은 이 여섯 범주의 조합으로 모두 해명된다', {'font_size': 18, 'space_before': 6}),
        ('· 후대 한문 시가 비평의 출발점이자 표준이 됨', {'font_size': 18, 'space_before': 6}),
    ])


@S(SEC2)
def ii_content(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '내용 분류 — 풍(風) · 아(雅) · 송(頌)')
    cols = [
        ('風 풍', '160편', '각 지방의 민요·백성의 노래',
         '15국의 민간 가요. 일상·연애·노동의 정서가 가장 풍부.\n시경 중 문학적 가치가 가장 높은 부분.'),
        ('雅 아', '105편', '궁중 연회·조회의 노래',
         '소아 74편 + 대아 31편.\n귀족의 정서와 정치적 송축·풍자가 어우러진다.'),
        ('頌 송', '40편', '종묘 제사의 노래',
         '주송 31편 + 노송 4편 + 상송 5편.\n조상의 덕을 찬양하는 엄숙한 제례악.'),
    ]
    for i, (han, num, label, desc) in enumerate(cols):
        x = Inches(0.7 + i * 4.2)
        add_filled_rect(slide, x, Inches(2.3), Inches(3.9), Inches(0.9), INK)
        add_textbox(slide, x, Inches(2.3), Inches(3.9), Inches(0.55),
                    han, font_size=26, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, Inches(2.78), Inches(3.9), Inches(0.4),
                    num, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, x, Inches(3.2), Inches(3.9), Inches(0.55), PALE)
        add_textbox(slide, x, Inches(3.2), Inches(3.9), Inches(0.55),
                    label, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.2), Inches(4.0), Inches(3.5), Inches(3.0),
                       [(desc, {'font_size': 14, 'color': INK})])


@S(SEC2)
def ii_tech(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '표현 기법 — 부(賦) · 비(比) · 흥(興)')
    rows = [
        ('賦 부', '직접 서술',
         '있는 그대로 펼쳐 말함 (敷陳其事 而直言之)',
         '예) 「부이(芣苢)」 — 질경이를 캐는 동작을 그대로 나열'),
        ('比 비', '비유 · 대조',
         '다른 사물에 빗대어 말함 (以彼物比此物也)',
         '예) 「종사(螽斯)」 — 메뚜기 떼의 번성으로 자손 번영을 비유'),
        ('興 흥', '연상 · 유발',
         '먼저 다른 사물을 들어 시상을 일으킴 (先言他物 以引起所詠之詞)',
         '예) 「관저(關雎)」 — 물수리의 울음에서 군자의 그리움을 일으킴'),
    ]
    for i, (han, label, principle, example) in enumerate(rows):
        y = Inches(2.3 + i * 1.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.5), Inches(1.35), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.5), Inches(0.8),
                    han, font_size=28, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.7), y + Inches(0.85), Inches(1.5), Inches(0.4),
                    label, font_size=13, color=PALE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(2.4), y + Inches(0.1), Inches(10.5), Inches(1.2), [
            (principle, {'font_size': 17, 'bold': True, 'color': INK}),
            (example, {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.25)


@S(SEC2)
def ii_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '육의 한 폭 정리')
    add_filled_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.55), INK)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(2.0), Inches(0.55),
                '구분', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(2.7), Inches(2.3), Inches(2.0), Inches(0.55),
                '명칭', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(4.7), Inches(2.3), Inches(8.0), Inches(0.55),
                '핵심 설명', font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('내용', '風 풍', '각 지방의 민요'),
        ('내용', '雅 아', '궁중의 공식 음악'),
        ('내용', '頌 송', '종묘 제사의 찬가'),
        ('기법', '賦 부', '직접 서술 — 사실을 그대로 풀어 말함'),
        ('기법', '比 비', '비유 — 다른 사물에 빗대어 말함'),
        ('기법', '興 흥', '연상 — 다른 사물을 들어 시상을 일으킴'),
    ]
    for i, (cat, name, desc) in enumerate(rows):
        y = Inches(2.85 + i * 0.55)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.55), bg)
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.55),
                    cat, font_size=15, color=SUB, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.7), y, Inches(2.0), Inches(0.55),
                    name, font_size=16, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.9), y, Inches(7.8), Inches(0.55),
                    desc, font_size=15, color=INK,
                    anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅲ. 국풍 ==============
SEC3 = 'Ⅲ. 국풍(國風)'

@S(SEC3)
def iii_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '국풍(國風) — 15국의 민요 160편',
              '시경의 핵심이자 문학적 가치의 정수')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('· 황하 유역 각 제후국의 민간 가요 — 「풍속·정치의 거울」',
         {'font_size': 18, 'space_before': 6}),
        ('· 부역·전쟁의 고통, 연애·이별의 정서, 풍자와 한탄까지 — 백성의 삶을 가장 진솔하게 담음',
         {'font_size': 18, 'space_before': 6}),
        ('· 시경 305편 중 가장 많은 160편, 문학적 가치가 가장 높이 평가되는 영역',
         {'font_size': 18, 'space_before': 6}),
        ('· 「풍자(諷)」, 「풍속(俗)」, 「가곡의 풍격(風)」 세 뜻이 결합된 명칭',
         {'font_size': 18, 'space_before': 6}),
        ('· 정치 풍자의 정신은 후대 「악부(樂府)」 시·당대 신악부 운동의 원천이 됨',
         {'font_size': 18, 'space_before': 6}),
    ])


@S(SEC3)
def iii_south(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '남방의 노래 — 주남 · 소남 · 빈풍',
              '주나라 발상지와 문왕 덕화권의 정서')
    rows = [
        ('1. 주남(周南)', '11편', '주나라 남쪽(남양) · 문왕의 덕화', '부부의 도리, 가정의 화목'),
        ('2. 소남(召南)', '14편', '소공의 봉지(남방)', '여성의 덕, 교화의 풍속'),
        ('15. 빈풍(豳風)', '7편',  '빈(豳) 지역 · 주족의 발상지', '농사·세시풍속의 백과사전'),
    ]
    for i, (name, num, region, theme) in enumerate(rows):
        y = Inches(2.5 + i * 1.4)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.8), Inches(1.2), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.8), Inches(0.65),
                    name, font_size=18, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.7), y + Inches(0.7), Inches(2.8), Inches(0.4),
                    num, font_size=14, color=SUB,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(3.7), y + Inches(0.15), Inches(9.0), Inches(1.0), [
            (region, {'font_size': 16, 'bold': True, 'color': INK}),
            (theme,  {'font_size': 15, 'color': SUB, 'space_before': 6}),
        ])


@S(SEC3)
def iii_central(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '중원의 노래 — 패·용·위·왕·정·제',
              '구은(舊殷) 유민과 왕실 변천의 음성')
    rows = [
        ('3. 패풍(邶風)', '19편', '이별과 그리움, 사회 비판'),
        ('4. 용풍(鄘風)', '10편', '풍속과 남녀 사랑'),
        ('5. 위풍(衛風)', '10편', '사랑과 이별, 정치 풍자'),
        ('6. 왕풍(王風)', '10편', '동주(東周) 왕실 쇠락의 비애'),
        ('7. 정풍(鄭風)', '21편', '남녀 사랑의 노래가 가장 풍부'),
        ('8. 제풍(齊風)', '11편', '호방한 기풍, 사냥의 노래'),
    ]
    for i, (name, num, theme) in enumerate(rows):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 1.3)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(1.1), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.5), Inches(0.45),
                    name, font_size=17, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(4.5), y + Inches(0.05), Inches(1.2), Inches(0.45),
                    num, font_size=14, color=SUB, align=PP_ALIGN.RIGHT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.5),
                    theme, font_size=14, color=INK)


@S(SEC3)
def iii_periphery(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '변경과 소국의 노래 — 위·당·진·진·회·조',
              '변경의 무용(武勇)과 약소국의 비애')
    rows = [
        ('9. 위풍(魏風)',  '7편',  '백성의 고단한 삶, 정치 비판'),
        ('10. 당풍(唐風)', '12편', '우국지정, 계절의 정서'),
        ('11. 진풍(秦風)', '10편', '상무 정신, 전쟁의 노래'),
        ('12. 진풍(陳風)', '10편', '남녀 사랑, 제사와 풍속'),
        ('13. 회풍(檜風)', '4편',  '소국의 비애와 망국의 한'),
        ('14. 조풍(曹風)', '4편',  '현실 비판과 풍자'),
    ]
    for i, (name, num, theme) in enumerate(rows):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 1.3)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(1.1), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.5), Inches(0.45),
                    name, font_size=17, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(4.5), y + Inches(0.05), Inches(1.2), Inches(0.45),
                    num, font_size=14, color=SUB, align=PP_ALIGN.RIGHT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.5),
                    theme, font_size=14, color=INK)


# ============== Ⅳ. 아 ==============
SEC4 = 'Ⅳ. 아(雅)'

@S(SEC4)
def iv_xiaoya(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '소아(小雅) — 궁중 연회의 노래 74편',
              '귀족과 관리의 정서, 민간가요의 일부까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 시대 — 서주(西周) 중기 ~ 동주(東周) 초기',
         {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● 성격 — 작은 연회(잔치)의 음악, 일부 정치 풍자·사회 비판도 포함',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 8}),
        ('● 대표작', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 8}),
        ('     · 녹명(鹿鳴) — 군신 화합의 연회시', {'font_size': 16, 'space_before': 4}),
        ('     · 사모(四牡) — 공무에 지친 사신의 향수', {'font_size': 16, 'space_before': 4}),
        ('     · 채미(采薇) — 변경 수자리 병사의 귀향', {'font_size': 16, 'space_before': 4}),
        ('     · 학명(鶴鳴) — 「他山之石 可以攻玉」의 출전', {'font_size': 16, 'space_before': 4}),
        ('     · 사간(斯干) — 왕궁 신축과 후사 축원', {'font_size': 16, 'space_before': 4}),
        ('● 6편은 가사를 잃고 제목만 전함 — 「생시(笙詩)」', {'font_size': 16, 'color': SUB, 'space_before': 8}),
    ])


@S(SEC4)
def iv_daya(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '대아(大雅) — 조회와 제사의 노래 31편',
              '주나라 건국 서사시와 정치적 훈계')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 시대 — 서주 초기·중기가 중심 — 가장 격조 높은 영역',
         {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● 성격 — 조회·대제(大祭) 등 큰 행사의 정악(正樂)',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 8}),
        ('● 대표작', {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 8}),
        ('     · 생민(生民) — 주의 시조 후직(后稷) 탄생 신화', {'font_size': 16, 'space_before': 4}),
        ('     · 공류(公劉) — 빈(豳) 정착의 영웅 서사', {'font_size': 16, 'space_before': 4}),
        ('     · 면(緜)     — 고공단보의 기산 이주', {'font_size': 16, 'space_before': 4}),
        ('     · 황의(皇矣) — 문왕의 천명 수임', {'font_size': 16, 'space_before': 4}),
        ('     · 문왕(文王) — 「周雖舊邦 其命維新」의 출전', {'font_size': 16, 'space_before': 4}),
        ('     · 탕(蕩)     — 폭정에 대한 격렬한 풍간', {'font_size': 16, 'space_before': 4}),
        ('● 다섯 편의 사시(史詩)는 주나라 흥기를 노래한 서사시 연작', {'font_size': 16, 'color': SUB, 'space_before': 8}),
    ])


# ============== Ⅴ. 송 ==============
SEC5 = 'Ⅴ. 송(頌)'

@S(SEC5)
def v_open(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '송(頌) — 종묘 제사의 노래 40편',
              '시경에서 가장 오래되고 가장 엄숙한 영역')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「송(頌)」 = 용(容) — 「제사 때의 춤사위·악곡 전체」를 가리키는 옛 글자',
         {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('● 종묘에서 조상의 덕을 찬양하는 제례악 — 길이는 짧고 한 장(章)으로 끝나는 시가 많음',
         {'font_size': 18, 'space_before': 8}),
        ('● 시경 중 가장 오래된 층 — 주송(周頌)은 서주 초기 작품으로 추정',
         {'font_size': 18, 'space_before': 6}),
        ('● 운(韻)이 거의 없고, 4언의 정형성을 따르지 않는 자유로운 형식이 특징',
         {'font_size': 18, 'space_before': 6}),
        ('● 문학적 미감보다 의례적 엄숙함 · 신성성 · 정통성의 표현이 우선',
         {'font_size': 18, 'space_before': 6}),
    ])


@S(SEC5)
def v_three(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC5, n, t)
    add_title(slide, '주송·노송·상송')
    cols = [
        ('周頌 주송', '31편',
         '주나라 종묘 제사악\n• 시경에서 가장 오래된 시\n• 청묘(清廟)·열문(烈文) 등\n• 형식 자유, 운 거의 없음'),
        ('魯頌 노송', '4편',
         '노나라 종묘 제사 · 희공 찬양\n• 「공자의 고향」 노나라의 송\n• 비궁(閟宮) 등 — 길이 가장 김\n• 형식은 풍·아에 가까워짐'),
        ('商頌 상송', '5편',
         '은(상)의 후예 송(宋)의 제사악\n• 「나(那)」, 「열조(烈祖)」, 「현조(玄鳥)」 등\n• 은의 시조 설(契) 탄생 신화\n• 가장 신화적·웅혼한 분위기'),
    ]
    for i, (han, num, body) in enumerate(cols):
        x = Inches(0.7 + i * 4.2)
        add_filled_rect(slide, x, Inches(2.3), Inches(3.9), Inches(1.0), INK)
        add_textbox(slide, x, Inches(2.3), Inches(3.9), Inches(0.6),
                    han, font_size=24, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, Inches(2.85), Inches(3.9), Inches(0.4),
                    num, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.2), Inches(3.6), Inches(3.5), Inches(3.5),
                       [(body, {'font_size': 15, 'color': INK})], line_spacing=1.35)


# ============== Ⅵ. 명편 깊이 읽기 ==============
SEC6 = 'Ⅵ. 명편 깊이 읽기'

def make_poem_slide(poem_num, total_poem, name_han, name_kor, source, original, modern, theme, point):
    @S(SEC6)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC6} ({poem_num}/{total_poem})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=28, bold=True, color=INK)
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    f'출전 — {source}', font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.4), [
            (original, {'font_size': 18, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 14, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 핵심 주제', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(theme, {'font_size': 16, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 음미할 지점', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 16, 'color': INK})], line_spacing=1.35)


POEMS = [
    ('관저(關雎)', '관저', '국풍·주남 — 시경 첫 편',
     '關關雎鳩 在河之洲  窈窕淑女 君子好逑',
     '관관 우는 물수리는 강가 모래톱에 있고 / 아리따운 숙녀는 군자의 좋은 배필이로다',
     '강가의 물수리에서 시상을 일으켜, 덕 있는 여인을 사모하는 군자의 마음을 그린다.',
     '공자가 「樂而不淫 哀而不傷」(즐거우되 음란하지 않고, 슬프되 상하지 않는다)이라 평한 시경의 표상.'),
    ('도요(桃夭)', '도요', '국풍·주남',
     '桃之夭夭 灼灼其華  之子于歸 宜其室家',
     '복숭아 어여뻐 그 꽃이 활짝 빛나네 / 이 아가씨 시집가니 그 집안에 어울리리라',
     '활짝 핀 복숭아꽃에 신부를 빗대어 가정의 번창을 축원하는 혼인의 노래.',
     '「桃之夭夭」는 동아시아 혼인 축가의 원형 — 꽃→열매→무성한 잎으로 점층 전개.'),
    ('부이(芣苢)', '부이', '국풍·주남',
     '采采芣苢 薄言采之  采采芣苢 薄言有之',
     '질경이를 캐세 질경이를 캐세 / 질경이를 캐세 질경이를 따세',
     '동사만 바꾸는 순수 반복 — 여인들이 함께 나물 캐는 노동의 리듬과 즐거움.',
     '시경에서 가장 단순하고도 완벽한 형식미. 「부(賦)」의 극치이자 노동요의 원형.'),
    ('한광(漢廣)', '한광', '국풍·주남',
     '南有喬木 不可休息  漢有游女 不可求思',
     '남쪽 높은 나무엔 쉴 수 없고 / 한수의 그녀를 구할 수 없네',
     '한수(漢水)를 사이에 둔 이루지 못한 사랑의 안타까움.',
     '「不可」의 반복으로 좌절감을 강조하면서도 끝내 희망을 놓지 못하는 깊은 서정.'),
    ('권이(卷耳)', '권이', '국풍·주남',
     '采采卷耳 不盈頃筐  嗟我懷人 寘彼周行',
     '도꼬마리 캐고 캐도 광주리에 차지 않네 / 그리운 이를 두고 길가에 광주리 놓네',
     '시점을 교차하며 떨어져 있는 부부의 상호 그리움을 노래.',
     '여인 → 남편으로 시점이 전환되는 독특한 구성. 「不盈頃筐」으로 그리움을 간접화.'),
    ('갈담(葛覃)', '갈담', '국풍·주남',
     '葛之覃兮 施于中谷  維葉萋萋 黃鳥于飛',
     '칡넝쿨이 골짜기에 뻗고 / 잎이 무성한데 꾀꼬리가 날아오네',
     '칡을 베어 옷감을 짜고 친정으로 가는 효성스러운 여인의 덕행.',
     '봄 풍경에서 노동·효성으로 시상이 흘러가는 흥(興)의 모범 사례.'),
    ('칠월(七月)', '칠월', '국풍·빈풍 — 시경 최장편',
     '七月流火 九月授衣  一之日觱發 二之日栗烈',
     '칠월에 큰불별 서쪽으로 흐르고 구월에 옷감 나눠준다 / 십일월 삭풍, 십이월 매서운 추위',
     '농민의 한 해 노동을 월별로 그린 시경의 백과사전 — 농사·양잠·사냥·세시풍속 망라.',
     '8장의 대편. 하력과 주력을 혼용. 사료 가치·민속학적 보고로 그 자체가 한 권의 책.'),
    ('치효(鴟鴞)', '치효', '국풍·빈풍',
     '迨天之未陰雨 徹彼桑土  綢繆牖戶 今女下民 或敢侮予',
     '비 오기 전에 뽕뿌리 캐어 / 창과 문 단단히 동여매니 아랫것들이 감히 업신여기랴',
     '주공이 둥지를 지키는 새에 자신을 비유 — 위기 앞의 결의와 사전 대비.',
     '「未雨綢繆(미우주무)」 — 「비 오기 전에 대비한다」 고사성어의 출전.'),
    ('동산(東山)', '동산', '국풍·빈풍',
     '我徂東山 慆慆不歸  我來自東 零雨其濛',
     '동산으로 갔다가 오래 돌아오지 못했네 / 동쪽에서 돌아오니 가랑비가 부슬부슬',
     '주공의 동정에서 돌아오는 병사의 귀향길 — 황량한 빈집과 평화의 갈망.',
     '동아시아 반전시(反戰詩)의 원형. 가랑비·거미줄·빈집의 이미지가 깊은 여운을 남김.'),
    ('녹명(鹿鳴)', '녹명', '소아 첫 편',
     '呦呦鹿鳴 食野之苹  我有嘉賓 鼓瑟吹笙',
     '어여삐 우는 사슴은 들판의 마름을 먹고 / 내게 귀한 손님 있어 비파 타고 생황 부네',
     '군신이 연회에서 화합하며 정치의 도리를 나누는 모습.',
     '소아의 첫 시이자 가장 격조 높은 연회시. 후대 향응시(饗宴詩)의 표준이 됨.'),
    ('채미(采薇)', '채미', '소아',
     '昔我往矣 楊柳依依  今我來思 雨雪霏霏',
     '내가 떠날 때엔 버드나무 늘어졌더니 / 지금 돌아오니 진눈깨비 흩날리네',
     '변경에서 돌아오는 수자리 병사의 깊은 비애와 향수.',
     '「楊柳依依」·「雨雪霏霏」 — 시경에서 가장 자주 인용되는 시구. 시간과 정서의 대비가 절묘.'),
    ('학명(鶴鳴)', '학명', '소아',
     '它山之石 可以攻玉',
     '다른 산의 돌이라도 옥을 다듬을 수 있다',
     '겉으론 보잘것없어 보이는 사물도 큰 쓸모를 가질 수 있음.',
     '「他山之石(타산지석)」의 출전 — 자기 발전·인재 등용의 비유로 2,500년간 인용.'),
    ('생민(生民)', '생민', '대아 — 주나라 시조 신화',
     '厥初生民 時維姜嫄  生民如何 克禋克祀',
     '백성의 시초는 강원이라 / 백성이 어떻게 났는가, 정성껏 제사를 올리고',
     '주나라 시조 후직(后稷)의 탄생과 농경의 시작을 노래한 서사시.',
     '주의 천명관과 농경 문화의 신화적 표상 — 대아 5편의 「주민족 서사시」 중 하나.'),
    ('문왕(文王)', '문왕', '대아 — 주의 천명',
     '周雖舊邦 其命維新  侯文王 於昭于天',
     '주는 비록 옛 나라이나 그 명은 새로워라 / 문왕이여, 하늘에 빛나도다',
     '문왕의 덕과 주나라의 천명을 노래.',
     '「周雖舊邦 其命維新」 — 『대학』과 후대 개혁 사상에서 끊임없이 인용되는 구절.'),
    ('청묘(清廟)', '청묘', '주송 첫 편',
     '於穆清廟 肅雝顯相  濟濟多士 秉文之德',
     '아, 깊고 맑은 종묘여, 엄숙하고 화락한 보필이여 / 많은 선비가 문왕의 덕을 받드네',
     '종묘 제사의 엄숙한 의식 — 문왕의 덕을 기리는 송의 표상.',
     '주송의 첫 시이자 송의 전형. 운(韻) 없는 자유로운 형식, 의례의 깊이가 압도적.'),
]

for i, p in enumerate(POEMS, 1):
    make_poem_slide(i, len(POEMS), *p)


# ============== Ⅶ. 시경의 사상 ==============
SEC7 = 'Ⅶ. 시경의 사상'

@S(SEC7)
def vii_yanji(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '시언지(詩言志) — 시는 뜻을 말한 것',
              '동아시아 문학론의 출발점')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85),
                '詩 言 志, 歌 永 言',
                font_size=30, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(2.0), Inches(3.2), Inches(9.3), Inches(0.4),
                '시는 뜻을 말한 것이요, 노래는 그 말을 길게 펼친 것이다 — 『서경·요전(舜典)』',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('· 시의 본질은 「뜻(志)」의 표현 — 외부 묘사가 아니라 내심의 토로',
         {'font_size': 18, 'space_before': 8}),
        ('· 「志」는 마음에 머문 것 — 그것이 말로 나오면 시가 되고, 길게 펼쳐지면 노래가 된다',
         {'font_size': 18, 'space_before': 8}),
        ('· 「대서(大序)」 — 「情動於中而形於言, 言之不足故嗟嘆之」 (마음이 움직이매 말로 형성되고, 말이 부족하매 탄식한다)',
         {'font_size': 17, 'color': SUB, 'space_before': 8}),
        ('· 후대 중국·한국·일본 시론(詩論)의 출발점 — 정(情)과 지(志)를 토대로 한 시 이해',
         {'font_size': 18, 'space_before': 8}),
    ])


@S(SEC7)
def vii_wuxie(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '사무사(思無邪) — 생각에 사악함이 없다',
              '공자의 시경 한 줄 평')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85),
                '詩 三 百, 一 言 以 蔽 之, 曰 思 無 邪',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(2.0), Inches(3.2), Inches(9.3), Inches(0.4),
                '시 삼백 편을 한마디로 줄이면, 생각에 사악함이 없다 — 『논어·위정』',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('· 「思無邪」 자체는 노송 「경(駉)」편의 한 구절을 빌려온 것',
         {'font_size': 18, 'space_before': 8}),
        ('· 305편 모두가 인간 정서의 정직한 표현 — 거짓된 마음이 없다는 평어',
         {'font_size': 18, 'space_before': 8}),
        ('· 사랑·이별·분노·풍자도 「邪」가 아니라 본래 마음의 결을 따른 진실',
         {'font_size': 18, 'space_before': 8}),
        ('· 시경을 도덕 교과서가 아니라 「진실의 노래책」으로 자리매김한 공자의 핵심 명제',
         {'font_size': 18, 'space_before': 8}),
    ])


@S(SEC7)
def vii_wenrou(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '온유돈후(溫柔敦厚) — 시교의 풍모',
              '시가 빚어내는 인격의 결')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(0.85),
                '溫 柔 敦 厚, 詩 教 也',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(2.0), Inches(3.2), Inches(9.3), Inches(0.4),
                '온화하고 부드러우며 도탑고 두터운 것이 시의 가르침이다 — 『예기·경해』',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(3.0), [
        ('· 溫(온) — 따뜻함 · 柔(유) — 부드러움 · 敦(돈) — 도타움 · 厚(후) — 두터움',
         {'font_size': 18, 'space_before': 8}),
        ('· 시는 격렬함보다 절제 — 풍자도 곧장 찌르지 않고 빗대어 일러줌',
         {'font_size': 18, 'space_before': 8}),
        ('· 「怨而不怒, 哀而不傷」 — 원망하되 노하지 않고, 슬프되 상하지 않는다',
         {'font_size': 18, 'space_before': 8}),
        ('· 시를 익힌 사람의 인격적 결 — 격정에 휩쓸리지 않는 절제와 두터움',
         {'font_size': 18, 'space_before': 8}),
    ])


@S(SEC7)
def vii_xingguan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '흥관군원(興觀群怨) — 시의 네 가지 효용',
              '공자의 시 효용론')
    add_textbox(slide, Inches(0.7), Inches(2.1), Inches(12.0), Inches(0.4),
                '『논어·양화』 — 「詩 可以興 可以觀 可以群 可以怨」',
                font_size=15, color=SUB, font_name='Batang')
    cols = [
        ('興 흥', '일으키다',
         '시는 사람의 마음을 일깨우고 시상을 일으킨다.\n정서를 자극하고 감수성을 깨움.'),
        ('觀 관', '살피다',
         '시는 풍속과 정치의 득실을 살피게 한다.\n사회 현실을 비추는 거울.'),
        ('群 군', '모이다',
         '시는 사람들이 함께 모여 어울리게 한다.\n공동체적 정서의 결속.'),
        ('怨 원', '풀어내다',
         '시는 원망과 비판을 표현할 수 있게 한다.\n억눌린 정서의 정당한 출구.'),
    ]
    for i, (han, kor, body) in enumerate(cols):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.7 + row * 2.1)
        add_filled_rect(slide, x, y, Inches(5.9), Inches(1.85), PALE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(1.6), Inches(0.55),
                    han, font_size=24, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + Inches(1.8), y + Inches(0.15), Inches(3.5), Inches(0.55),
                    kor, font_size=14, color=SUB, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(1.0),
                       [(body, {'font_size': 14, 'color': INK})], line_spacing=1.3)


@S(SEC7)
def vii_meishi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '미자(美刺) — 찬미와 풍자',
              '시의 정치적 기능')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「대서(大序)」 — 「上以風化下, 下以風刺上」', {'font_size': 19, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('     위로는 시로써 아래를 교화하고, 아래로는 시로써 위를 풍자한다',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 「美」 — 임금과 신하의 덕을 찬미함', {'font_size': 19, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     송(頌)이 그 전형 — 종묘의 찬가, 조상 덕업의 송축', {'font_size': 16, 'space_before': 4}),
        ('● 「刺」 — 정치·풍속의 잘못을 풍자·풍간', {'font_size': 19, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     국풍·아의 다수 — 폭정·부역·실정에 대한 백성과 사대부의 항의', {'font_size': 16, 'space_before': 4}),
        ('● 풍자도 직설이 아닌 빗대어 — 「主文而譎諫」 (문(文)을 주로 하고 완곡히 간언함)',
         {'font_size': 16, 'color': SUB, 'space_before': 14}),
        ('● 후대 한대 악부 · 당대 신악부 운동의 정신적 원천이 됨',
         {'font_size': 16, 'color': SUB, 'space_before': 6}),
    ])


# ============== Ⅷ. 공자의 시교 ==============
SEC8 = 'Ⅷ. 공자의 시교(詩敎)'

@S(SEC8)
def viii_value(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '왜 공자는 시를 가르쳤는가',
              '인격 수양과 언어 교육의 핵심 교재')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 시는 감정을 일깨우는 가장 빠른 길 — 「흥(興)」', {'font_size': 18, 'space_before': 6}),
        ('● 시는 풍속을 살피는 가장 정확한 거울 — 「관(觀)」', {'font_size': 18, 'space_before': 6}),
        ('● 시는 외교의 언어 — 사신은 시구로 뜻을 전했다 (賦詩言志)',
         {'font_size': 18, 'space_before': 6}),
        ('● 시는 박물(博物)의 교과서 — 새·짐승·초목의 이름을 다 알게 됨',
         {'font_size': 18, 'space_before': 6}),
        ('● 시는 효(孝)와 충(忠)의 가르침 — 가까이는 어버이를, 멀리는 임금을 섬김',
         {'font_size': 18, 'space_before': 6}),
        ('● 시는 인격의 결을 빚는 가장 부드러운 길 — 「溫柔敦厚」',
         {'font_size': 18, 'space_before': 6}),
    ])


@S(SEC8)
def viii_buxueshi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '不學詩 無以言 — 시를 배우지 않으면 말할 수 없다',
              '아들 백어(伯魚)에게 — 『논어·계씨』')
    add_filled_rect(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(1.0), PALE)
    add_textbox(slide, Inches(2.0), Inches(2.3), Inches(9.3), Inches(1.0),
                '不 學 詩, 無 以 言',
                font_size=36, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(3.5), [
        ('● 공자가 아들에게 — 「시를 배우지 않으면 더불어 말할 수 없다」',
         {'font_size': 18, 'space_before': 8}),
        ('● 시경의 구절은 당시 외교·정치의 공통 언어 — 시를 모르면 말의 무게가 없다',
         {'font_size': 18, 'space_before': 8}),
        ('● 시를 익혀야 말이 점잖아지고, 비유가 풍부해지고, 정서가 깊어진다',
         {'font_size': 18, 'space_before': 8}),
        ('● 「興於詩, 立於禮, 成於樂」 — 시에서 일으키고, 예에서 서고, 악에서 완성한다',
         {'font_size': 18, 'color': SUB, 'space_before': 14, 'font_name': 'Batang'}),
    ])


@S(SEC8)
def viii_comments(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '공자의 시 평론 — 세 구절')
    blocks = [
        ('關雎', '관저편에 대하여',
         '「樂而不淫, 哀而不傷」\n즐거우되 음란하지 않고, 슬프되 상하지 않는다'),
        ('詩三百', '시 전체에 대하여',
         '「一言以蔽之, 曰: 思無邪」\n한마디로 줄이면 — 생각에 사악함이 없다'),
        ('興觀群怨', '시의 효용에 대하여',
         '「詩可以興, 可以觀, 可以群, 可以怨」\n시는 흥을 일으키고, 살피고, 모으고, 풀어낸다'),
    ]
    for i, (han, label, body) in enumerate(blocks):
        y = Inches(2.3 + i * 1.65)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.2), Inches(1.45), ACCENT)
        add_textbox(slide, Inches(0.7), y + Inches(0.15), Inches(2.2), Inches(0.6),
                    han, font_size=24, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, Inches(0.7), y + Inches(0.85), Inches(2.2), Inches(0.45),
                    label, font_size=12, color=PALE, align=PP_ALIGN.CENTER)
        add_paragraphs(slide, Inches(3.1), y + Inches(0.15), Inches(10.0), Inches(1.3),
                       [(body, {'font_size': 16, 'color': INK, 'font_name': 'Batang'})],
                       line_spacing=1.4)


# ============== Ⅸ. 명구 모음 ==============
SEC9 = 'Ⅸ. 명구 모음'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC9)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC9} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=34, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=17, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 17, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('주남·관저',  '窈 窕 淑 女, 君 子 好 逑',
     '아리따운 숙녀는 군자의 좋은 배필이로다',
     '시경의 첫 구절. 후대의 모든 동아시아 연애시·혼인시가 그 그늘에서 자라났다고 해도 과언이 아니다.'),
    ('주남·관저(공자 평)', '樂 而 不 淫, 哀 而 不 傷',
     '즐거우되 음란하지 않고, 슬프되 상하지 않는다',
     '공자가 관저편을 평한 말이자, 동아시아 미학의 절제(中和) 원칙. 격정도 절제 안에서 빛난다.'),
    ('주남·도요',  '桃 之 夭 夭, 灼 灼 其 華',
     '복숭아 어여뻐 그 꽃이 활짝 빛나도다',
     '혼인 축가의 원형. 「夭夭」, 「灼灼」 같은 첩어가 자아내는 음악성이 한자 시가의 한 출발점이 됨.'),
    ('위풍·목과', '投 我 以 木 桃, 報 之 以 瓊 瑤',
     '내게 모과를 던지매, 나는 옥으로 갚노라',
     '받은 작은 정을 큰 마음으로 되갚는다 — 동아시아 의리·예물 교환 윤리의 시구적 원형.'),
    ('소아·학명', '它 山 之 石, 可 以 攻 玉',
     '다른 산의 돌이라도 옥을 다듬을 수 있다',
     '타산지석(他山之石) — 보잘것없는 것에서도 큰 쓸모를 찾는 지혜. 2,500년간 가장 자주 인용된 시구.'),
    ('소아·차할', '高 山 仰 止, 景 行 行 止',
     '높은 산은 우러르고, 큰길은 따라 걷는다',
     '사마천이 『사기·공자세가』 끝에서 인용 — 공자에 대한 흠모를 표한 명구. 인격의 표상에 대한 경모.'),
    ('소아·소민', '戰 戰 兢 兢, 如 臨 深 淵, 如 履 薄 冰',
     '전전긍긍 — 깊은 못 앞에 선 듯, 살얼음을 밟듯',
     '신중함의 극한. 「전전긍긍」, 「여리박빙」 두 사자성어의 출전. 책임자의 자세를 상징.'),
    ('대아·탕', '靡 不 有 初, 鮮 克 有 終',
     '시작이 없는 일은 없으나, 끝을 잘 맺는 일은 드물다',
     '성공의 어려움은 시작에 있지 않고 마지막에 있다. 「초심 → 유종의 미」 윤리의 원천.'),
    ('대아·문왕', '周 雖 舊 邦, 其 命 維 新',
     '주는 옛 나라이나 그 명은 새로워라',
     '『대학』과 후대 개혁 사상의 단골 인용. 오래된 나라일수록 끊임없는 갱신이 명운을 좌우한다.'),
    ('빈풍·칠월', '七 月 流 火, 九 月 授 衣',
     '칠월에 큰불별 서쪽으로 흐르고, 구월에 옷감을 나눠준다',
     '시경에서 가장 자주 인용되는 농가월령. 자연의 순환과 인간의 노동을 잇는 농경 문명의 표지.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅹ. 해석 학파와 후대 영향 ==============
SEC10 = 'Ⅹ. 해석 학파와 후대 영향'

@S(SEC10)
def x_moshi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '모시(毛詩)와 한대(漢代) 4가시',
              '한대에 오경의 하나가 되다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 한대에 시경은 「제(齊)·노(魯)·한(韓)·모(毛)」 네 학파로 전수됨',
         {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('     · 제시(齊詩) — 원고생(轅固生) 계열', {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('     · 노시(魯詩) — 신배(申培) 계열', {'font_size': 15, 'color': SUB, 'space_before': 2}),
        ('     · 한시(韓詩) — 한영(韓嬰) 계열 — 『한시외전』만 남음', {'font_size': 15, 'color': SUB, 'space_before': 2}),
        ('     · 모시(毛詩) — 모형(毛亨)·모장(毛萇) 계열', {'font_size': 15, 'color': SUB, 'space_before': 2}),
        ('● 동한(東漢) 정현(鄭玄)의 「전(箋)」으로 모시가 표준이 됨',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     이후 시경 = 모시(毛詩) — 다른 세 가지는 점차 실전(失傳)', {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 「대서(大序)」, 「소서(小序)」 — 각 시의 배경과 주제를 정치적 풍자로 해석',
         {'font_size': 18, 'space_before': 14}),
    ])


@S(SEC10)
def x_zhuxi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '주희(朱熹)의 『시집전(詩集傳)』',
              '송대 — 시경 해석의 대전환')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1177년 — 주희가 모시의 정치적 강독을 비판하며 새로운 주석을 완성',
         {'font_size': 18, 'space_before': 6}),
        ('● 「국풍은 본디 백성의 노래 — 정치 풍자가 아닌 「민간의 정(情)」으로 읽어야 한다」',
         {'font_size': 18, 'color': ACCENT, 'space_before': 8}),
        ('● 「음시(淫詩)」 논쟁 — 정풍·위풍의 연애시를 「음란한 시」로 규정하여 풍자가 아닌 직접 표현으로 봄',
         {'font_size': 17, 'space_before': 8}),
        ('● 「부(賦)·비(比)·흥(興)」 표기를 시마다 명기 — 표현 기법 분석의 정착',
         {'font_size': 17, 'space_before': 8}),
        ('● 이후 조선·일본의 시경 강독은 『시집전』이 표준 — 「모시 → 주자」로 권위 이동',
         {'font_size': 17, 'space_before': 8}),
        ('● 청대 고증학(顧炎武·戴震)은 모시·주자 모두를 비판하며 원전 복원 시도',
         {'font_size': 17, 'color': SUB, 'space_before': 8}),
    ])


@S(SEC10)
def x_chinese(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '중국 시가에 미친 영향',
              '시경 → 초사 → 한부 → 악부 → 당시')
    rows = [
        ('초사(楚辭)',   '굴원 — 남방의 무가(巫歌) 전통이 시경의 북방 4언을 잇는 새 길', '6언·산문조의 서정시'),
        ('한부(漢賦)',   '시경의 「부(賦)」 기법이 독자 장르로 발전',                  '대규모 사물 묘사'),
        ('악부(樂府)',   '한대 — 시경 채시(采詩) 정신의 재현, 민가 수집',              '풍자·서사의 부활'),
        ('건안 5언',    '조조·조비·조식 — 시경·악부의 정신을 5언으로 계승',           '서정과 비분(悲憤)'),
        ('당시(唐詩)',   '두보의 사회시는 풍의 풍자, 왕유의 자연시는 흥의 정조',         '시 황금기의 양 축'),
        ('신악부 운동', '백거이 — 「文章合為時而著, 歌詩合為事而作」',                  '시경 정신의 부활'),
    ]
    for i, (name, content, sub) in enumerate(rows):
        y = Inches(2.3 + i * 0.78)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.2), Inches(0.65), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.2), Inches(0.65),
                    name, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, Inches(3.1), y + Inches(0.05), Inches(7.5), Inches(0.6), [
            (content, {'font_size': 14, 'color': INK}),
        ])
        add_textbox(slide, Inches(10.7), y + Inches(0.05), Inches(2.5), Inches(0.6),
                    sub, font_size=13, color=SUB, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC10)
def x_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '한국과 일본에서의 시경')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 한국 — 삼국시대부터 수용, 조선조에 「사서삼경」의 핵심으로 자리잡음',
         {'font_size': 18, 'bold': True, 'color': ACCENT}),
        ('     · 과거(科擧) 필독 — 사대부의 기본 교양', {'font_size': 16, 'space_before': 4}),
        ('     · 정약용 『시경강의(詩經講義)』, 이익·박지원 등 실학자들의 새로운 독해', {'font_size': 16, 'space_before': 4}),
        ('     · 시조·가사 문학에 「온유돈후」·「사무사」의 미학이 깊이 스며듦', {'font_size': 16, 'space_before': 4}),
        ('● 일본 — 헤이안 시대부터 한학자들이 강독, 에도 시대 본격 연구',
         {'font_size': 18, 'bold': True, 'color': ACCENT, 'space_before': 14}),
        ('     · 이토 진사이(伊藤仁齋)·오규 소라이(荻生徂徠) — 모시 비판과 원전 회복', {'font_size': 16, 'space_before': 4}),
        ('     · 『만요슈(萬葉集)』의 서민가요 정신과의 비교 연구', {'font_size': 16, 'space_before': 4}),
        ('● 베트남 — 한문 교양의 기초로 「四書三經」 체계에 편입됨',
         {'font_size': 18, 'space_before': 14}),
    ])


@S(SEC10)
def x_position(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '동아시아 한문학사에서의 위상')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「오경(五經)」 중 유일한 시가집 — 詩·書·易·禮·春秋 중 문학 영역의 대표',
         {'font_size': 18, 'space_before': 6}),
        ('● 「사가(史家)」가 본 시경 — 사마천 『사기』는 시경 구절을 곳곳에 인용',
         {'font_size': 18, 'space_before': 8}),
        ('● 「논자(論者)」가 본 시경 — 『논어』·『맹자』·『순자』가 끊임없이 시구를 인용해 사상을 받쳐 줌',
         {'font_size': 18, 'space_before': 8}),
        ('● 「외교의 언어」 — 춘추시대 사신이 시경 구절을 인용해 뜻을 전했다(賦詩)',
         {'font_size': 18, 'space_before': 8}),
        ('● 「언어의 보고」 — 후대 사자성어·고사·어휘의 가장 풍부한 출전',
         {'font_size': 18, 'space_before': 8}),
    ])


# ============== Ⅺ. 현대적 의의 ==============
SEC11 = 'Ⅺ. 현대적 의의'

@S(SEC11)
def xi_realism(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '사실주의 문학의 원류',
              '백성의 말로 백성의 삶을 적다')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('· 「부역의 고단함」·「전쟁의 참혹함」·「폭정에 대한 항의」를 가감 없이 노래',
         {'font_size': 18, 'space_before': 8}),
        ('· 「벌단(伐檀)」 — 「不稼不穡, 胡取禾三百廛兮」 일하지 않는 자가 어찌 곡식을 거두는가 — 사회 비판의 모범',
         {'font_size': 17, 'space_before': 8, 'font_name': 'Batang'}),
        ('· 「석서(碩鼠)」 — 「큰 쥐야, 큰 쥐야, 내 곡식을 먹지 말라」 — 폭정을 쥐에 비유한 직접 풍자',
         {'font_size': 17, 'space_before': 8}),
        ('· 「소민(小民)」의 시선 — 통치자의 기록인 사서(史書)가 채우지 못하는 영역의 보존',
         {'font_size': 17, 'space_before': 8}),
        ('· 두보·백거이의 사회시, 정약용 「애절양」 등의 정신적 원조',
         {'font_size': 17, 'space_before': 8}),
    ])


@S(SEC11)
def xi_folklore(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '민속·역사 자료로서의 가치',
              '시경은 곧 BC 11~6세기의 동영상')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('· 농경 — 「칠월」 한 편이 농민의 1년 노동·세시풍속의 백과사전',
         {'font_size': 18, 'space_before': 6}),
        ('· 혼인·연애 — 정풍·위풍의 사랑 노래가 고대 남녀 교제의 실상을 보여줌',
         {'font_size': 18, 'space_before': 6}),
        ('· 전쟁 — 진풍·동산·채미 등은 군역의 고통과 향수의 1차 기록',
         {'font_size': 18, 'space_before': 6}),
        ('· 의례 — 송(頌)은 종묘 제사의 절차와 음악을 보존한 유일한 문헌',
         {'font_size': 18, 'space_before': 6}),
        ('· 동식물 — 새·짐승·물고기·풀·나무의 이름이 가장 풍부 — 박물학의 보고',
         {'font_size': 18, 'space_before': 6}),
        ('· 사서(『상서』·『춘추』)가 「대인의 기록」이라면, 시경은 「소인의 기록」',
         {'font_size': 17, 'color': SUB, 'space_before': 14}),
    ])


@S(SEC11)
def xi_aesthetics(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '동아시아 미학의 뿌리',
              '절제 · 중화 · 함축 · 흥(興)의 미학')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 절제 — 「樂而不淫, 哀而不傷」 — 격정도 일정한 결을 따라 흐른다',
         {'font_size': 18, 'space_before': 6}),
        ('● 중화(中和) — 강하지도 약하지도 않은 균형의 미',
         {'font_size': 18, 'space_before': 8}),
        ('● 함축 — 직설보다 비유, 비유보다 연상(興)의 깊이',
         {'font_size': 18, 'space_before': 8}),
        ('● 일상의 발견 — 새의 울음·풀의 흔들림·물의 흐름에서 시작되는 시상',
         {'font_size': 18, 'space_before': 8}),
        ('● 한국 시조·일본 와카·중국 절구·율시 — 모두 시경이 마련한 미학의 변주',
         {'font_size': 18, 'space_before': 8}),
    ])


@S(SEC11)
def xi_labor(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '일상과 노동에 대한 시선',
              '시경이 우리에게 일러주는 것')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('· 시경의 시인은 영웅이 아니다 — 나물 캐는 여인, 군역 가는 남자, 베 짜는 아내',
         {'font_size': 18, 'space_before': 6}),
        ('· 노동의 리듬이 곧 시의 리듬이 된다 — 「부이(芣苢)」가 보여주듯',
         {'font_size': 18, 'space_before': 8}),
        ('· 「내 일상」을 노래하는 권리 — 시는 권력의 전유물이 아니다',
         {'font_size': 18, 'space_before': 8}),
        ('· 작은 슬픔·작은 기쁨이 그대로 시가 되는 세계 — 현대 서정시의 원형',
         {'font_size': 18, 'space_before': 8}),
        ('· 시경은 「큰 것을 말하지 않고도 깊은 것을 말하는 법」을 보여준다',
         {'font_size': 18, 'space_before': 8, 'bold': True, 'color': ACCENT}),
    ])


# ============== Ⅻ. 마무리 ==============
SEC12 = 'Ⅻ. 마무리'

@S(SEC12)
def xii_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC12, n, t)
    add_title(slide, '시경, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 시경은 BC 11세기~6세기, 약 500년의 노래 305편이다.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 풍(160)·아(105)·송(40) — 백성의 노래·궁중의 음악·종묘의 찬가.',
         {'font_size': 18, 'space_before': 8}),
        ('● 부(賦)·비(比)·흥(興) — 직설·비유·연상의 세 갈래.',
         {'font_size': 18, 'space_before': 8}),
        ('● 시언지(詩言志)·사무사(思無邪)·온유돈후(溫柔敦厚)·흥관군원(興觀群怨).',
         {'font_size': 18, 'space_before': 8}),
        ('● 동아시아 시가·미학·언어·외교·교양의 공동 토대.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「사람의 마음을 가장 정직하게 노래한 책」.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC12)
def xii_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.0),
                '思 無 邪',
                font_size=120, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.0), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '생각에  사악함이  없다',
                font_size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
                '— 공자, 시 삼백을 한 마디로 평하다 (『논어·위정』)',
                font_size=16, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                '詩  經',
                font_size=22, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\시경.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
