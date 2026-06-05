# -*- coding: utf-8 -*-
"""
손자병법(孫子兵法) 발표자료 — 전면 보강판 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
13편 각 1장 깊이 읽기 · 손무 생애 · 무경칠서 위상 · 후대 영향 망라
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '2,500년의 병법 · 세계 3대 병서의 최고봉 · 무경칠서의 첫머리',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '孫 子 兵 法',
                font_size=92, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '손 자 병 법',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '不 戰 而 屈 人 之 兵 — 싸우지 않고 이기는 것이 최선',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '손무(孫武) · 춘추 말 BC 6세기 · 13편 약 6,000자',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '병법서이자 인생서 · 동아시아·서양 군사학의 영원한 교과서',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 6,000자의 역설'),
        ('Ⅱ.', '손무와 시대 — 춘추 말의 오나라'),
        ('Ⅲ.', '위상과 판본 · 은작산 죽간'),
        ('Ⅳ.', '13편의 구조 · 풍림화산'),
        ('Ⅴ.', '13편 각 편 깊이 읽기'),
        ('Ⅵ.', '핵심 사상 6기둥'),
    ]
    items_right = [
        ('Ⅶ.', '명구 16선'),
        ('Ⅷ.', '짝 책들 — 손빈·오자·삼십육계'),
        ('Ⅸ.', '역사 적용 — 적벽·이순신·마오'),
        ('Ⅹ.', '동서양 후대 영향 · 현대 경영'),
        ('Ⅺ.', '오늘 우리에게 · 마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.6 + i * 0.65)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 손자병법')
    rows = [
        ('서명',  '손자병법(孫子兵法) · 손자(孫子) · 오손자(吳孫子)'),
        ('저자',  '손무(孫武), 자(字) 장경(長卿) — 제(齊)나라 출신'),
        ('시대',  '춘추시대 말기 (BC 6세기 말 ~ 5세기 초)'),
        ('분량',  '13편 · 약 6,000자 — 동양 고전 중 가장 압축된 단일 저작'),
        ('성격',  '동양 최고(最古)·최고(最高)의 군사 전략서 · 보편적 지혜의 보고'),
        ('위상',  '세계 3대 병서 (손자병법 · 클라우제비츠 전쟁론 · 미야모토 오륜서)'),
        ('주석사', '조조(曹操) 최초 주석 · 송 무경칠서 첫머리 · 1972 은작산 죽간'),
        ('적용 영역', '전쟁 · 외교 · 경영 · 협상 · 스포츠 · 일상'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_brevity(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '왜 6,000자인가 — 간결함의 힘',
              '단행본 한 챕터 분량이 2,500년 동안 살아남은 이유')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 13편 · 6,000자 — 오늘날 단행본 한 챕터 분량',
         {'font_size': 18, 'space_before': 4}),
        ('● 그러나 각 문장은 극도로 압축된 명제 — 한 구절이 한 원리, 한 편이 한 체계',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 『논어』가 공자의 단편적 어록 모음이라면,',
         {'font_size': 17, 'space_before': 12}),
        ('     손자병법은 한 인물이 체계적으로 설계한 단일 저작',
         {'font_size': 17, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 이 점에서 동양 고전 중 독특한 위치를 차지',
         {'font_size': 17, 'space_before': 12}),
        ('● 압축의 미학 — 읽는 자가 자기 상황에 비추어 해석할 수 있는 「열린 텍스트」',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
        ('● 2,500년 동안 끊임없이 새 의미를 낳을 수 있었던 이유',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC1)
def i_paradox(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '위대한 역설 — 「병법서가 싸움을 부정한다」',
              '서양 클라우제비츠 『전쟁론』과 결정적으로 다른 지점')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.9), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.9),
                '不 戰 而 屈 人 之 兵  善 之 善 者 也',
                font_size=24, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.1), Inches(12.0), Inches(0.4),
                '싸우지 않고 적의 군대를 굴복시키는 것이 최선이다 — 제3편 모공',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(3.0), [
        ('● 제1편 시계 — 「兵者 國之大事 不可不察也」',
         {'font_size': 16, 'space_before': 4, 'font_name': 'Batang'}),
        ('     전쟁은 중대사이니 살피지 않을 수 없다',
         {'font_size': 13, 'color': SUB, 'space_before': 2}),
        ('● 제3편 모공 — 「不戰而屈人之兵 善之善者也」',
         {'font_size': 16, 'space_before': 10, 'font_name': 'Batang'}),
        ('     싸우지 않고 이기는 것이 최선',
         {'font_size': 13, 'color': SUB, 'space_before': 2}),
        ('● 제12편 화공 — 「亡國不可以復存 死者不可以復生」',
         {'font_size': 16, 'space_before': 10, 'font_name': 'Batang'}),
        ('     망한 나라는 되살릴 수 없고, 죽은 자는 다시 살릴 수 없다',
         {'font_size': 13, 'color': SUB, 'space_before': 2}),
        ('● 「전쟁을 준비하되 전쟁을 넘어서라」 — 2,500년간 살아남은 이유',
         {'font_size': 17, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC1)
def i_keyideas(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '핵심 사상 한 폭으로')
    boxes = [
        ('不 戰 勝', '부전승', '싸우지 않고 이긴다'),
        ('知 彼 知 己', '지피지기', '적을 알고 나를 안다'),
        ('五 事 七 計', '오사칠계', '도·천·지·장·법 + 7가지 비교'),
        ('奇 正 / 虛 實', '기정·허실', '정공과 기습, 빈 곳과 찬 곳'),
        ('勢', '세', '활시위의 에너지 — 형세의 운용'),
        ('詭 道', '궤도', '전쟁은 본질적으로 속임수'),
    ]
    for i, (han, kor, desc) in enumerate(boxes):
        col, row = i % 3, i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.3 + row * 2.3)
        add_filled_rect(slide, x, y, Inches(3.9), Inches(2.0), PALE)
        add_textbox(slide, x, y + Inches(0.2), Inches(3.9), Inches(0.7),
                    han, font_size=22, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, font_name='Batang')
        add_textbox(slide, x, y + Inches(0.95), Inches(3.9), Inches(0.4),
                    kor, font_size=14, color=SUB, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.5), Inches(0.5),
                    desc, font_size=13, color=INK, align=PP_ALIGN.CENTER)


# ============== Ⅱ. 손무와 시대 ==============
SEC2 = 'Ⅱ. 손무와 시대'

@S(SEC2)
def ii_sunwu(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '손무(孫武) — 제나라 출신의 오나라 장군',
              '제나라 내분을 피해 오로 망명 · 오자서의 추천')
    rows = [
        ('이름',     '손무(孫武), 자(字) 장경(長卿)'),
        ('출생지',   '제(齊)나라 (현 산동성)'),
        ('가문',     '제나라 명문 전씨(田氏) 가문 출신'),
        ('활동 시기', 'BC 544경 ~ BC 496경'),
        ('활동 국가', '오(吳)나라'),
        ('주군',     '오왕 합려(闔閭)'),
        ('동료',     '오자서(伍子胥) — 손무를 합려에게 추천'),
        ('대표 업적', '13편 저술 · 오의 초(楚) 정벌 (수도 영郢 함락)'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC2)
def ii_concubines(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '후궁 180명 훈련 일화 — 사기에 전하는 손무의 출발',
              '『사기·손자오기열전』의 가장 유명한 일화')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 합려가 손무의 실력을 시험하고자 후궁 180명을 주며 훈련시켜 보라 함',
         {'font_size': 17, 'space_before': 4}),
        ('● 손무는 두 부대로 나누고 합려의 총희(寵姬) 두 명을 대장으로 삼음',
         {'font_size': 17, 'space_before': 10}),
        ('● 명령을 내렸으나 궁녀들은 웃기만 함 — 한 번 용서',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 다시 명령했으나 여전히 웃음 → 총희 두 명을 즉시 참수',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 합려가 만류하자 — 「君命有所不受 — 장군이 전장에 있을 때는 임금의 명도 받지 않는다」',
         {'font_size': 16, 'space_before': 10, 'font_name': 'Batang'}),
        ('● 이후 궁녀들은 완벽하게 명령에 따랐다',
         {'font_size': 17, 'space_before': 10}),
        ('● 이 일화의 상징 — 엄격한 군법·실전에서 증명된 병법가·실력 평가가 감정보다 우선',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC2)
def ii_era(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '춘추 말기 — 손자병법이 나온 시대',
              '전쟁 양상의 대형화 · 현실주의 전쟁관의 등장')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 주(周) 왕실 권위 완전 상실 — 「춘추 5패」의 시대',
         {'font_size': 17, 'space_before': 4}),
        ('● 전쟁 양상의 변화 — 귀족 의례적 전투 → 대규모 동원 전쟁',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     철기 보급 · 보병 중심 편제 확대',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 오월(吳越)의 부상 — 남방 변방국이 중원을 뒤흔드는 시기',
         {'font_size': 17, 'space_before': 10}),
        ('     합려·부차·구천·오자서·범려 — 모두 이 시대의 인물',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 사상의 태동 — 공자(BC 551~479)가 거의 동시대',
         {'font_size': 17, 'space_before': 10}),
        ('● 손자병법은 「전쟁이 대형화·전문화되던 시점」의 최초 체계적 병법서',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 귀족 의례의 전쟁관을 버리고 「현실주의적·합리주의적」 전쟁관을 제시 — 획기적',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


# ============== Ⅲ. 무경칠서·판본 ==============
SEC3 = 'Ⅲ. 위상과 판본'

@S(SEC3)
def iii_wujing(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '무경칠서(武經七書)의 첫머리',
              '송 신종(1043) 공식 선정 · 동아시아 군사학의 정전 7권 중 으뜸')
    rows = [
        ('孫子兵法 손자병법',  '손무',          '전쟁의 원리 — 부전승의 철학', True),
        ('吳子兵法 오자병법',  '오기',          '군사 행정 · 장수와 사기'),
        ('司馬法 사마법',     '전양저',         '군례(軍禮) · 군의 윤리'),
        ('尉繚子 울료자',     '울료',          '법치 군사론 · 엄정한 군기'),
        ('李衛公問對 이위공문대', '이정',        '당대 전략 문답 · 정·기의 변용'),
        ('六韜 육도',         '태공망',         '종합 통치 · 인재 · 전략 · 전술'),
        ('三略 삼략',         '황석공',         '리더십 · 강유 조화 · 정당한 전쟁'),
    ]
    for i, row in enumerate(rows):
        name, author, role = row[0], row[1], row[2]
        highlight = len(row) > 3 and row[3]
        y = Inches(2.3 + i * 0.6)
        bg = ACCENT if highlight else (PALE if i % 2 == 1 else WHITE)
        text_color = WHITE if highlight else INK
        add_filled_rect(slide, Inches(0.7), y, Inches(3.7), Inches(0.55), bg)
        add_textbox(slide, Inches(0.7), y, Inches(3.7), Inches(0.55),
                    name, font_size=15, bold=True, color=text_color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.6), y, Inches(2.2), Inches(0.55),
                    author, font_size=14, color=ACCENT if highlight else SUB, bold=highlight,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(7.0), y, Inches(5.9), Inches(0.55),
                    role, font_size=14, color=ACCENT if highlight else INK, bold=highlight,
                    anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
                '송대부터 무과(武科) 필수 과목 — 동아시아 무관의 첫 교과서',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)


@S(SEC3)
def iii_versions(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '판본의 역사 — 1972년 은작산 죽간의 충격',
              '조조의 최초 주석부터 한묘 죽간 출토까지')
    rows = [
        ('춘추 말', '손무가 13편 저술, 오왕 합려에게 헌상'),
        ('전국시대', '후손 손빈(孫臏)이 별도의 『손빈병법』 저술'),
        ('한(漢)대', '장량·한신이 정리 · 『사기』 손자오기열전 수록'),
        ('삼국시대', '조조(曹操)가 최초 주석서 『위무주손자(魏武注孫子)』 — 현존 가장 권위 있는 주해'),
        ('송(宋)대', '무경칠서의 첫머리로 편입 — 무과 필수 과목'),
        ('1972년', '산동 은작산(銀雀山) 한묘에서 죽간본 출토 — 진본 입증'),
    ]
    for i, (era, desc) in enumerate(rows):
        y = Inches(2.4 + i * 0.75)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.4), Inches(0.6), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.4), Inches(0.6),
                    era, font_size=14, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.3), y + Inches(0.05), Inches(9.6), Inches(0.55),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '은작산 죽간 — 『손자병법』과 『손빈병법』이 별개의 책임이 처음 명확히 입증된 결정적 발견',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅳ. 13편의 구조 ==============
SEC4 = 'Ⅳ. 13편의 구조'

@S(SEC4)
def iv_flow(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '13편의 흐름 — 「계산 → 운용 → 변화 → 정보」',
              '한 인물이 설계한 단일 체계의 논리적 전개')
    rows = [
        ('단계 1',  '계산',     '1~3편 시계·작전·모공', '전쟁의 전제 — 「하기 전에 다 끝난다」'),
        ('단계 2',  '운용',     '4~6편 군형·병세·허실', '힘의 운용 — 형(形)·세(勢)·허실(虛實)'),
        ('단계 3',  '기동',     '7~8편 군쟁·구변',     '전장의 움직임 — 우직지계와 변통'),
        ('단계 4',  '환경',     '9~11편 행군·지형·구지', '환경의 활용 — 행군·지형·9가지 상황'),
        ('단계 5',  '특수·정보', '12~13편 화공·용간',    '특수 전술과 정보 — 결정적 비대칭'),
    ]
    for i, (step, kind, vol, desc) in enumerate(rows):
        y = Inches(2.4 + i * 0.9)
        add_filled_rect(slide, Inches(0.7), y, Inches(1.2), Inches(0.75), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(1.2), Inches(0.75),
                    step, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, Inches(2.0), y, Inches(1.4), Inches(0.75), PALE)
        add_textbox(slide, Inches(2.0), y, Inches(1.4), Inches(0.75),
                    kind, font_size=15, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y + Inches(0.05), Inches(3.5), Inches(0.7),
                    vol, font_size=13, bold=True, color=ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(7.1), y + Inches(0.05), Inches(5.8), Inches(0.7),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_table(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '13편 한 폭으로')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.45), INK)
    headers = [('편', 0.8), ('편명', 3.3), ('핵심 주제', 7.9)]
    x = Inches(0.7)
    for label, w in headers:
        add_textbox(slide, x, Inches(2.2), Inches(w), Inches(0.45),
                    label, font_size=13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(w)
    rows = [
        ('1', '始計 시계',  '전쟁의 기본 계산 — 오사칠계(五事七計)'),
        ('2', '作戰 작전',  '전쟁의 비용과 속전속결 — 졸속(拙速)'),
        ('3', '謀攻 모공',  '전략적 승리 — 부전승, 지피지기 백전불태'),
        ('4', '軍形 군형',  '공수의 형세 — 선승이후구전(先勝而後求戰)'),
        ('5', '兵勢 병세',  '기정(奇正)의 운용 · 세(勢)의 활용'),
        ('6', '虛實 허실',  '허와 실의 운용 · 주도권 장악 · 「피실격허」'),
        ('7', '軍爭 군쟁',  '우직지계(迂直之計) · 풍림화산(風林火山)'),
        ('8', '九變 구변',  '아홉 가지 변통 · 장수의 다섯 위험'),
        ('9', '行軍 행군',  '행군과 진영 배치 · 적정(敵情) 관찰'),
        ('10', '地形 지형', '여섯 가지 지형 · 장수의 리더십'),
        ('11', '九地 구지', '아홉 가지 전장 상황 · 사지즉전(死地則戰)'),
        ('12', '火攻 화공', '화공의 다섯 유형 · 전쟁의 신중함'),
        ('13', '用間 용간', '정보전 · 다섯 종류의 간첩 — 「先知」'),
    ]
    for i, (no, name, desc) in enumerate(rows):
        y = Inches(2.65 + i * 0.37)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.37), bg)
        add_textbox(slide, Inches(0.7), y, Inches(0.8), Inches(0.37),
                    no, font_size=13, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.5), y, Inches(3.3), Inches(0.37),
                    name, font_size=13, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(4.9), y, Inches(7.7), Inches(0.37),
                    desc, font_size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC4)
def iv_fenglinhuoshan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '풍림화산(風林火山) — 군쟁편의 4가지 행군 자세',
              '일본 다케다 신겐 깃발의 원전 · 손자병법 가장 시적인 문장')
    cols = [
        ('風', '풍', '其 疾 如 風',  '빠를 때는 바람처럼', '신속·기습'),
        ('林', '림', '其 徐 如 林',  '느릴 때는 숲처럼',   '정숙·은밀'),
        ('火', '화', '侵 掠 如 火',  '칠 때는 불처럼',     '맹렬·과감'),
        ('山', '산', '不 動 如 山',  '멈출 때는 산처럼',   '부동·견고'),
    ]
    for i, (han, kor, original, modern, label) in enumerate(cols):
        x = Inches(0.7 + i * 3.1)
        add_filled_rect(slide, x, Inches(2.3), Inches(2.9), Inches(2.2), ACCENT)
        add_textbox(slide, x, Inches(2.4), Inches(2.9), Inches(1.4),
                    han, font_size=120, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(3.95), Inches(2.9), Inches(0.4),
                    kor, font_size=15, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, x, Inches(4.6), Inches(2.9), Inches(0.6), PALE)
        add_textbox(slide, x, Inches(4.6), Inches(2.9), Inches(0.6),
                    original, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(5.3), Inches(2.9), Inches(0.4),
                    modern, font_size=13, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, Inches(5.8), Inches(2.9), Inches(0.4),
                    label, font_size=13, color=ACCENT, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4),
                '4가지 속도와 자세를 자유롭게 전환하는 유연성 — 다케다 신겐의 「風林火山」 깃발',
                font_size=14, color=SUB, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(6.9), Inches(12.0), Inches(0.4),
                '계속 — 「難 知 如 陰  動 如 雷 霆」 (알기 어려움은 그늘 같이, 움직임은 천둥 같이)',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅴ. 13편 깊이 읽기 ==============
SEC5 = 'Ⅴ. 13편 깊이 읽기'

def make_chapter_slide(num, total, name_han, name_kor, headline,
                        original, modern, theme, point):
    @S(SEC5)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC5} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=26, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    headline, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.4), [
            (original, {'font_size': 18, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 핵심 가르침', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(theme, {'font_size': 14, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 오늘에의 적용', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 15, 'color': INK})], line_spacing=1.35)


CHAPTERS = [
    ('始計 시계', '시계', '제1편 — 전쟁의 기본 계산 · 오사칠계(五事七計)',
     '兵 者  國 之 大 事  死 生 之 地  存 亡 之 道  不 可 不 察 也',
     '전쟁은 나라의 중대사 / 죽고 사는 터전, 존속하고 망하는 길이니, 살피지 않을 수 없다',
     '오사 — 道(도)·天(천)·地(지)·將(장)·法(법) / 칠계 — 일곱 가지 비교 분석 / 「兵者 詭道也」 — 전쟁은 본질적으로 속임수 / 「廟算多者勝」 — 묘산(廟算)에서 이기는 자가 실제로 이긴다',
     '모든 전략적 결정의 출발점은 신중한 분석. 전쟁 전 「조정의 계산」이 곧 현대의 사업 계획·시장 분석.'),
    ('作戰 작전', '작전', '제2편 — 전쟁의 비용과 속전속결',
     '兵 聞 拙 速  未 睹 巧 之 久 也',
     '전쟁에서 서툴더라도 빠른 것은 들었어도, 교묘하면서 오래 끄는 것은 본 적이 없다',
     '졸속(拙速)의 원리 — 길게 끄는 완벽함보다 짧은 부족함이 낫다 / 전쟁 비용의 엄청남 — 하루 천금(日費千金) / 「智將務食於敵」 — 적의 식량으로 먹인다 — 현지 조달',
     '완벽 강박을 버리고 빠른 실행 — 「good enough, but fast」. 보급선이 길수록 비용이 폭증.'),
    ('謀攻 모공', '모공', '제3편 — 부전승의 철학 · 지피지기',
     '不 戰 而 屈 人 之 兵  善 之 善 者 也    知 彼 知 己  百 戰 不 殆',
     '싸우지 않고 굴복시키는 것이 최선 / 적을 알고 나를 알면 백 번 싸워도 위태롭지 않다',
     '벌모(伐謀) → 벌교(伐交) → 벌병(伐兵) → 벌성(攻城)의 우선순위 / 「上兵伐謀」 — 최상의 병법은 적의 계책을 깬다 / 「知彼知己」 — 손자병법 전체에서 가장 유명한 명제',
     '경쟁 이전에 정보·관계·전략으로 이긴다. 정면 충돌은 마지막 수단. 자기 인식과 상대 분석의 짝.'),
    ('軍形 군형', '군형', '제4편 — 공수의 형세 · 선승이후구전',
     '勝 兵 先 勝 而 後 求 戰  敗 兵 先 戰 而 後 求 勝',
     '이기는 군대는 먼저 이기고 나서 싸우고, 지는 군대는 먼저 싸우고 나서 이기려 한다',
     '「先勝而後求戰」 — 이긴 다음 싸운다 / 「不可勝在己 可勝在敵」 — 지지 않는 것은 나에게, 이기는 것은 적에게 달림 / 「自保而全勝」 — 자기를 보전하면서 완전한 승리',
     '준비된 자에게 승부는 이미 결정. 「운이 좋아 이긴 것」은 진짜 승리가 아니다. 평소 축적이 모든 것.'),
    ('兵勢 병세', '병세', '제5편 — 기정(奇正)과 세(勢)',
     '凡 戰 者  以 正 合  以 奇 勝',
     '대저 전쟁이란 정공(正)으로 맞붙고, 기습(奇)으로 이긴다',
     '기정상생(奇正相生) — 정과 기는 서로 끝없이 바뀐다 / 「激水之疾 至於漂石者 勢也」 — 격류가 돌을 띄우는 것이 곧 세 / 「善戰者 求之於勢 不責於人」 — 잘 싸우는 자는 사람을 탓하지 않고 세를 구한다',
     '정공으로 안정, 기습으로 승부. 「에너지의 운용」이 전부 — 형세를 만든 자가 이긴다.'),
    ('虛實 허실', '허실', '제6편 — 허와 실 · 주도권의 비밀',
     '避 實 而 擊 虛   兵 形 象 水',
     '실(實)을 피하고 허(虛)를 친다 / 군대의 형세는 물과 같다',
     '「致人而不致於人」 — 남을 끌고 다니되 끌려다니지 않는다 — 주도권 / 「兵形象水」 — 물이 지형에 따르듯 적정(敵情)에 따라 변한다 / 적의 강한 곳을 피하고 약한 곳을 친다',
     '경쟁에서 「내가 강한 곳, 적이 약한 곳」으로 승부를 옮기는 능력. 주도권을 잡고 놓지 않는 자가 이긴다.'),
    ('軍爭 군쟁', '군쟁', '제7편 — 우직지계 · 풍림화산',
     '其 疾 如 風  其 徐 如 林  侵 掠 如 火  不 動 如 山',
     '빠르기는 바람, 느림은 숲, 침략은 불, 움직이지 않음은 산같이',
     '우직지계(迂直之計) — 돌아가는 길이 빠른 길 / 풍림화산(風林火山) — 4가지 행군 자세 — 일본 다케다 신겐의 「풍림화산」 깃발의 원전 / 「以迂爲直 以患爲利」 — 우회를 직진으로, 어려움을 이점으로',
     '직진이 항상 빠른 길이 아니다 — 돌아가는 지혜. 4가지 속도와 자세를 자유롭게 전환할 수 있는 유연성.'),
    ('九變 구변', '구변', '제8편 — 아홉 가지 변통 · 장수의 다섯 위험',
     '將 有 五 危 — 必 死 · 必 生 · 忿 速 · 廉 潔 · 愛 民',
     '장수의 다섯 위험 — 죽기만 하려는 것·살기만 하려는 것·성내고 급한 것·청렴함에 매이는 것·백성을 사랑하기에 매이는 것',
     '구변(九變) — 9가지 변통의 원리 / 장수의 다섯 위험 — 극단의 덕도 결정적 약점 / 「君命有所不受」 — 임금의 명도 거역할 수 있다',
     '한 덕목에 고착되면 그것이 곧 약점. 청렴·애민·용감조차 극단이면 위험. 상황에 따른 변통이 곧 능력.'),
    ('行軍 행군', '행군', '제9편 — 행군과 진영 · 적정 관찰',
     '辭 卑 而 益 備 者  進 也   辭 强 而 進 驅 者  退 也',
     '말이 겸손하면서 준비가 많으면 진격할 것이고, 말이 강하면서 전진하면 물러나려는 것이다',
     '32가지 적정 관찰법 — 적의 말·행동·연기·새의 움직임 등으로 의도 파악 / 「兵非益多」 — 군대가 많은 것이 능사가 아니다 / 「令之以文 齊之以武」 — 文으로 명령하고 武로 가지런히 한다',
     '관찰의 힘 — 표면이 아닌 패턴을 읽는 안목. 협상장의 미세한 신호 읽기, 데이터 행간 읽기의 원형.'),
    ('地形 지형', '지형', '제10편 — 여섯 가지 지형과 장수의 리더십',
     '視 卒 如 嬰 兒  故 可 與 之 赴 深 谿   視 卒 如 愛 子  故 可 與 之 俱 死',
     '병사를 갓난아기처럼 보면 깊은 골짜기에도 같이 갈 수 있고, 사랑하는 자식처럼 보면 함께 죽을 수 있다',
     '여섯 가지 지형 — 통(通)·괘(掛)·지(支)·애(隘)·험(險)·원(遠) / 「知彼知己 知天知地 勝乃可全」 — 적·나·하늘·땅을 다 알아야 완전한 승리 / 장수의 인격 — 병사를 자기 몸으로 본다',
     '환경(地)에 대한 분석 + 부하에 대한 사랑 — 리더십의 양 축. 「내 사람」으로 만드는 능력.'),
    ('九地 구지', '구지', '제11편 — 아홉 가지 전장 상황 · 사지즉전',
     '投 之 亡 地 然 後 存  陷 之 死 地 然 後 生',
     '망할 곳에 던진 뒤에야 살아나고, 죽을 곳에 빠뜨린 뒤에야 살아난다',
     '구지(九地) — 산지·경지·쟁지·교지·구지·중지·비지·위지·사지 / 「사지즉전(死地則戰)」 — 죽을 곳에서는 결사적으로 싸운다 / 배수진의 사상적 원형',
     '편안한 곳에서 안주하면 무너지고, 절박한 곳에서 결사적이면 살아난다 — 위기의 역설.'),
    ('火攻 화공', '화공', '제12편 — 화공의 다섯 유형 · 전쟁의 신중함',
     '亡 國 不 可 以 復 存  死 者 不 可 以 復 生   故 明 主 慎 之  良 將 警 之',
     '망한 나라는 되살릴 수 없고, 죽은 자는 다시 살릴 수 없다 / 그러므로 명군은 신중하고 양장은 경계한다',
     '화공오법(火攻五法) — 인(人)·적(積)·치(輜)·고(庫)·대(隊)의 다섯 유형 / 그러나 마지막 메시지는 「전쟁의 비가역성」 / 「非利不動 非得不用 非危不戰」 — 이익이 없으면 움직이지 말라',
     '극단적 수단도 알아두되, 절대 가볍게 쓰지 말라. 전쟁(경쟁)은 되돌릴 수 없는 결과를 낳는다는 경고.'),
    ('用間 용간', '용간', '제13편 — 정보전 · 다섯 종류의 간첩',
     '故 明 君 賢 將  所 以 動 而 勝 人  成 功 出 於 衆 者  先 知 也',
     '명군과 어진 장수가 움직여 이기고 무리에 뛰어난 공을 세우는 까닭은 「먼저 아는 것(先知)」 때문이다',
     '오간(五間) — 향간(鄕間)·내간(內間)·반간(反間)·사간(死間)·생간(生間) / 「선지(先知)」 — 먼저 아는 자가 이긴다 / 정보의 가치 — 천금을 아끼지 말라',
     '정보는 가장 적은 비용으로 가장 큰 효과를 낳는 무기. 현대의 인텔리전스·시장 정보·HR이 모두 「용간편」의 후예.'),
]

for i, p in enumerate(CHAPTERS, 1):
    make_chapter_slide(i, len(CHAPTERS), *p)


# ============== Ⅵ. 핵심 사상 6기둥 ==============
SEC6 = 'Ⅵ. 핵심 사상 6기둥'

def make_concept_slide(num, total, han, kor, source, principle, today):
    @S(SEC6)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC6} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{han}  ({kor})',
                    font_size=30, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    source, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.7), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.6),
                       [(principle, {'font_size': 14, 'color': INK})], line_spacing=1.45)
        add_textbox(slide, Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.4),
                    '◆ 오늘에의 함의', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.3), Inches(12.0), Inches(3.0),
                       [(today, {'font_size': 15, 'color': INK})], line_spacing=1.5)


CONCEPTS = [
    ('不 戰 勝', '부전승', '제3편 모공 — 손자병법의 정점',
     '· 「不戰而屈人之兵 善之善者也」 — 싸우지 않고 굴복시키는 것이 최선\n· 우선순위 — 벌모(伐謀) → 벌교(伐交) → 벌병(伐兵) → 벌성(攻城)\n· 「百戰百勝 非善之善者也」 — 백전백승도 최선이 아니다',
     '경쟁 이전에 전략으로, 충돌 이전에 협상으로 이긴다.\n현대 협상학·블루오션 전략·연합 외교의 동양적 원천.\n「싸우지 않고 이기는 자가 가장 강한 자」.'),
    ('知 彼 知 己', '지피지기', '제3편 모공 — 가장 자주 인용되는 명제',
     '· 「知彼知己 百戰不殆」 — 적을 알고 나를 알면 백 번 싸워도 위태롭지 않다\n· 「不知彼而知己 一勝一負」 — 적을 모르고 나만 알면 반은 이기고 반은 진다\n· 「不知彼不知己 每戰必敗」 — 둘 다 모르면 매번 진다',
     '자기 인식(知己)이 적 분석(知彼)에 선행한다 — 「나를 모르면 적도 모른다」.\n현대 SWOT 분석·자기 진단·데이터 분석의 2,500년 선행 원리.'),
    ('五 事 七 計', '오사칠계', '제1편 시계 — 전쟁 분석의 다섯 축',
     '오사(五事) — 道(도)·天(천)·地(지)·將(장)·法(법)\n  도(道) — 명분 · 천(天) — 시기 · 지(地) — 환경 · 장(將) — 리더 · 법(法) — 시스템\n칠계(七計) — 7가지 구체 비교 항목',
     '5가지 차원의 종합 분석 — 「대의 + 타이밍 + 환경 + 리더 + 시스템」.\n현대 경영의 PEST/5 forces 분석과 직접 통한다.\n특히 「도(道)」가 첫 번째 — 명분 없는 전쟁은 시작부터 진다.'),
    ('奇 正 / 虛 實', '기정·허실', '제5편 병세 · 제6편 허실',
     '「以正合 以奇勝」 — 정공으로 맞붙고, 기습으로 이긴다\n기정상생(奇正相生) — 정과 기는 끝없이 바뀐다\n「避實而擊虛」 — 실을 피하고 허를 친다\n「兵形象水」 — 물처럼 형태를 바꾼다',
     '정공(기본 역량)이 있어야 기습(혁신)이 통한다.\n경쟁에서 「내가 강한 곳, 상대가 약한 곳」으로 옮기는 능력.\n블루오션이 곧 「허(虛)」, 레드오션이 곧 「실(實)」.'),
    ('勢', '세', '제5편 병세 — 활시위의 에너지',
     '「激水之疾 至於漂石者 勢也」 — 격류가 돌을 띄우는 것이 세\n「善戰者 求之於勢 不責於人」 — 잘 싸우는 자는 사람을 탓하지 않고 세를 구한다\n「勢如張弩 節如發機」 — 활시위처럼 모으고 방아쇠처럼 풀어낸다',
     '개인의 노력보다 「형세의 운용」이 우선.\n트렌드·조류·모멘텀을 만드는 자가 이긴다.\n현대 마케팅의 momentum, 정치의 기세, 스타트업의 PMF — 모두 「세」의 변주.'),
    ('詭 道', '궤도', '제1편 시계 — 전쟁의 본질',
     '「兵者 詭道也」 — 전쟁은 본질적으로 속임수\n「能而示之不能 用而示之不用」 — 능력을 감추고 의도를 가린다\n「近而示之遠 遠而示之近」 — 가까이서 멀리 보이게, 멀리서 가까이 보이게',
     '경쟁에서 「투명성」은 미덕이 아니다 — 정보 비대칭이 곧 무기.\n포커·협상·시장 전략 — 「내 패를 보여주지 않는 것」이 기본.\n현대 정보전·심리전의 원형.'),
]

for i, c in enumerate(CONCEPTS, 1):
    make_concept_slide(i, len(CONCEPTS), *c)


# ============== Ⅶ. 명구 16선 ==============
SEC7 = 'Ⅶ. 명구 16선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC7)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC7} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=16, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 16, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('제1편 시계', '兵 者  國 之 大 事  死 生 之 地  存 亡 之 道  不 可 不 察 也',
     '전쟁은 나라의 중대사 · 죽고 사는 터전, 존속하고 망하는 길이니, 살피지 않을 수 없다',
     '손자병법의 첫 문장. 신중함이 모든 전략의 출발이라는 선언. 전쟁(경쟁)은 가벼이 시작할 일이 아니다.'),
    ('제1편 시계', '兵 者  詭 道 也',
     '전쟁은 본질적으로 속임수이다',
     '정직함이 미덕이 아닌 영역. 정보의 비대칭이 곧 우위 — 포커·협상·시장의 기본 원리.'),
    ('제2편 작전', '兵 聞 拙 速  未 睹 巧 之 久 也',
     '서툴더라도 빠른 것은 들었어도, 교묘하면서 오래 끄는 것은 본 적이 없다',
     '졸속(拙速)의 원리. 완벽 강박을 버리고 빠른 실행. 「good enough, but fast」의 가장 짧은 원전.'),
    ('제3편 모공', '不 戰 而 屈 人 之 兵  善 之 善 者 也',
     '싸우지 않고 적의 군대를 굴복시키는 것이 최선 중의 최선',
     '손자병법 전체의 정점. 「최고의 승리는 전쟁 자체가 없는 것」 — 협상·외교·연합의 원천 명제.'),
    ('제3편 모공', '知 彼 知 己  百 戰 不 殆',
     '적을 알고 나를 알면 백 번 싸워도 위태롭지 않다',
     '손자병법에서 가장 자주 인용되는 한 줄. 「위태롭지 않다」는 표현이 핵심 — 백전백승이 아니라 「불태(不殆)」.'),
    ('제3편 모공', '上 兵 伐 謀  其 次 伐 交  其 次 伐 兵  其 下 攻 城',
     '최상은 적의 계책을 깨고, 그 다음은 외교를 깨며, 그 다음이 군대 · 가장 하책이 성을 친다',
     '경쟁의 우선순위. 직접 충돌은 마지막. 머리 → 관계 → 무력 → 정면 충돌 — 모든 영역에 적용되는 위계.'),
    ('제4편 군형', '勝 兵 先 勝 而 後 求 戰',
     '이기는 군대는 먼저 이기고 나서 싸운다',
     '준비된 자에게 승부는 이미 결정 — 「선승이후구전」. 「운으로 이긴 승리」는 진짜 승리가 아니다.'),
    ('제5편 병세', '以 正 合  以 奇 勝',
     '정공으로 맞붙고, 기습으로 이긴다',
     '기정(奇正)의 원리. 정공(기본 역량)이 있어야 기습(혁신)이 통한다 — 둘은 끊임없이 바뀐다.'),
    ('제5편 병세', '激 水 之 疾  至 於 漂 石 者  勢 也',
     '격류가 돌을 띄우는 것이 곧 세(勢)이다',
     '개인의 힘이 아니라 「세」의 운용. 격류는 본래 무거운 돌도 띄운다 — 형세를 만드는 자가 이긴다.'),
    ('제6편 허실', '兵 形 象 水',
     '군대의 형세는 물과 같다',
     '물이 지형에 따라 흐르듯, 전략은 적정(敵情)에 따라 변한다. 손자병법의 가장 시적인 비유.'),
    ('제6편 허실', '致 人 而 不 致 於 人',
     '남을 끌고 다니되, 끌려다니지 않는다',
     '주도권의 정의. 모든 경쟁의 본질은 「누가 누구를 끌고 다니는가」 — 게임의 룰을 정하는 자가 이긴다.'),
    ('제7편 군쟁', '其 疾 如 風  其 徐 如 林  侵 掠 如 火  不 動 如 山',
     '빠를 때는 바람처럼, 느릴 때는 숲처럼, 칠 때는 불처럼, 멈출 때는 산처럼',
     '풍림화산(風林火山). 일본 다케다 신겐의 깃발 문장. 4가지 속도·자세를 자유로이 전환하는 유연성의 표상.'),
    ('제7편 군쟁', '以 迂 爲 直  以 患 爲 利',
     '돌아가는 것을 직진으로 삼고, 어려움을 이점으로 삼는다',
     '우직지계(迂直之計). 직진이 항상 빠른 길이 아니다 — 돌아가는 지혜. 어려움 속에서 기회를 보는 안목.'),
    ('제8편 구변', '君 命 有 所 不 受',
     '임금의 명령도 받지 않을 수 있다',
     '장수의 전권 위임. 손무가 후궁 일화에서 실제 적용한 원칙. 현장의 판단이 본부의 명령보다 우선할 때가 있다.'),
    ('제11편 구지', '投 之 亡 地 然 後 存  陷 之 死 地 然 後 生',
     '망할 곳에 던진 뒤에야 살아나고, 죽을 곳에 빠뜨린 뒤에야 살아난다',
     '사지즉전(死地則戰). 배수진의 사상적 원형. 절박함이 곧 동력 — 위기의 역설.'),
    ('제13편 용간', '故 明 君 賢 將  所 以 動 而 勝 人  成 功 出 於 衆 者  先 知 也',
     '명군과 어진 장수가 움직여 이기는 까닭은, 「먼저 아는 것(先知)」 때문이다',
     '정보의 가치. 가장 적은 비용으로 가장 큰 효과 — 현대 인텔리전스·시장 정보·HR의 2,500년 선행 명제.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅷ. 짝 책들 ==============
SEC_PAIR = 'Ⅷ. 짝 책들'

@S(SEC_PAIR)
def pair_sunbin(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_PAIR, n, t)
    add_title(slide, '손빈병법(孫臏兵法) — 2,000년 만에 부활한 책',
              '손무의 5대손 손빈 · 1972 은작산 죽간 출토')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 손빈(孫臏, BC 382?~316?) — 손무의 5대손 · 제(齊)나라 군사(軍師)',
         {'font_size': 17, 'space_before': 4}),
        ('● 방연(龐涓)의 모함으로 빈형(臏刑, 슬개골 도려내는 형벌) — 그 후 「손빈」으로 불림',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
        ('● 두 전설적 전투', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 계릉(桂陵) 전투 (BC 353) — 「위위구조(圍魏救趙)」의 원형',
         {'font_size': 14, 'space_before': 4}),
        ('     · 마릉(馬陵) 전투 (BC 341) — 방연을 죽음으로 몬 「감조계(減竈計)」',
         {'font_size': 14, 'space_before': 4}),
        ('● 송대 이후 실전 → 1972 산둥 은작산 한묘 죽간 출토로 부활',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     『손자병법』과 『손빈병법』이 별개의 책임이 죽간 출토로 확정',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 손자가 「전쟁의 원리」라면 손빈은 「전국시대 실전 사례 기록」',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC_PAIR)
def pair_wuzi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_PAIR, n, t)
    add_title(slide, '오자병법(吳子兵法) — 손자의 짝',
              '오기(吳起, BC 440~381) · 6편 · 군사 행정의 매뉴얼')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 오기 — 위(衛) 출신 · 노(魯)·위(魏)·초(楚) 3국에서 활약한 명장',
         {'font_size': 17, 'space_before': 4}),
        ('● 위 문후(文侯)·무후(武侯)를 도와 위나라를 전국 초기 최강국으로 만듦',
         {'font_size': 17, 'space_before': 10}),
        ('● 6편 구성 — 도국(圖國)·요적(料敵)·치병(治兵)·논장(論將)·응변(應變)·여사(勵士)',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 손자와의 차이 — 손자가 「전쟁의 본질」이라면 오자는 「군사 행정」',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 손자 — 압축·이론적 / 오자 — 구체적·실무적',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 손자 — 부전승의 철학 / 오자 — 장수와 병사의 사기 진작',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 손자와 함께 「손오병법(孫吳兵法)」으로 묶여 동양 병법의 양대 축',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
        ('● 오기의 최후 — 초나라에서 변법 추진 후 귀족들의 활에 죽음',
         {'font_size': 15, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC_PAIR)
def pair_36(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_PAIR, n, t)
    add_title(slide, '삼십육계(三十六計) — 손자의 응용 매뉴얼',
              '명말청초 작자 미상 · 6계 × 6류 = 36계의 사자성어')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 「三十六計 走爲上」 — 서른여섯 계책 중 달아나는 것이 최고의 계책',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 6류 × 6계 구조', {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('     · 승전계(勝戰計) — 우세할 때의 계책 (만천과해·위위구조·차도살인…)',
         {'font_size': 13, 'space_before': 6}),
        ('     · 적전계(敵戰計) — 대등할 때 (무중생유·암도진창·격안관화…)',
         {'font_size': 13, 'space_before': 4}),
        ('     · 공전계(攻戰計) — 공격 (타초경사·차시환혼·조호리산…)',
         {'font_size': 13, 'space_before': 4}),
        ('     · 혼전계(混戰計) — 혼란기 (부저추신·혼수모어·금선탈각…)',
         {'font_size': 13, 'space_before': 4}),
        ('     · 병전계(并戰計) — 합병 (투량환주·지상매괴·반객위주…)',
         {'font_size': 13, 'space_before': 4}),
        ('     · 패전계(敗戰計) — 패할 때 (미인계·공성계·반간계·고육계·연환계·주위상)',
         {'font_size': 13, 'space_before': 4}),
        ('● 손자가 「원리」라면 삼십육계는 「구체적 36개 계책 매뉴얼」 — 모공편 「벌모」의 응용',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC_PAIR)
def pair_compare(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_PAIR, n, t)
    add_title(slide, '동양 병서의 큰 지도 — 손자를 중심으로',
              '오자병법·삼십육계·육도삼략과의 관계')
    rows = [
        ('孫子兵法 손자병법',  '손무 · 13편 6,000자',  '전쟁의 본질 · 부전승의 철학', '본 PPT의 주제'),
        ('吳子兵法 오자병법',  '오기 · 6편',           '군사 행정 · 장수와 사기 · 실전 매뉴얼', '손자의 짝'),
        ('三十六計 삼십육계', '저자·시대 불명 · 36계', '구체적 계책 매뉴얼 · 모공의 「벌모」 사례집', '손자의 응용'),
        ('六韜 三略 육도삼략', '강태공·황석공',         '제왕학의 병서 · 통치와 결합', '손자의 보완'),
        ('司馬法 사마법',     '전양저 · 5편',          '군의 윤리(軍禮)', '무경칠서 동기'),
        ('尉繚子 울료자',     '울료 · 24편',          '법치 군사론 · 엄정한 군기', '무경칠서 동기'),
        ('李衛公問對 이위공문대', '이정 · 3권',        '당대 전략 문답 · 정·기의 변용', '손자의 해설'),
    ]
    for i, row in enumerate(rows):
        name, info, role, rel = row
        y = Inches(2.4 + i * 0.6)
        is_sunzi = i == 0
        bg = ACCENT if is_sunzi else (PALE if i % 2 == 0 else WHITE)
        color = WHITE if is_sunzi else INK
        add_filled_rect(slide, Inches(0.7), y, Inches(3.0), Inches(0.55), bg)
        add_textbox(slide, Inches(0.7), y, Inches(3.0), Inches(0.55),
                    name, font_size=14, bold=True, color=color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(3.9), y, Inches(2.4), Inches(0.55),
                    info, font_size=12, color=SUB, bold=is_sunzi,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(6.5), y, Inches(4.0), Inches(0.55),
                    role, font_size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(10.6), y, Inches(2.3), Inches(0.55),
                    rel, font_size=12, color=ACCENT, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC_PAIR)
def pair_pair(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_PAIR, n, t)
    add_title(slide, '손자와 짝 — 같이 읽으면 좋은 책',
              '「장수의 눈」과 「왕의 눈」을 함께 가지는 법')
    cols = [
        ('孫 子', '장수의 눈',
         '6,000자의 압축\n전장의 정밀한 승부 기술\n\n전술과 운용 중심\n「부전승」의 철학\n\n좁고 깊은 책',
         ACCENT),
        ('六 韜 三 略', '왕의 눈',
         '수만 자의 풍부\n천하 창건의 정치\n\n통치와 인재 중심\n「천하공유」의 철학\n\n넓고 풍부한 책',
         INK),
    ]
    for i, (han, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=30, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    label, font_size=15, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.5), Inches(3.6), Inches(4.9), Inches(3.5),
                       [(body, {'font_size': 15, 'color': INK})], line_spacing=1.5)
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                '한신(韓信)이 손자병법의 장수, 장량(張良)이 삼략의 책사 — 한(漢) 건국의 두 축',
                font_size=13, color=SUB, align=PP_ALIGN.CENTER)


# ============== Ⅸ. 역사 적용 사례 ==============
SEC_HIST = 'Ⅸ. 역사 적용 사례'

@S(SEC_HIST)
def hist_samguk(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_HIST, n, t)
    add_title(slide, '삼국지(三國志) — 손자병법의 살아 있는 무대',
              '조조·제갈량·주유 — 모두 손자병법의 대가')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 조조(曹操) — 현존 가장 권위 있는 손자 주석서 『위무주손자』 저술',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('     · 「내가 평생 사랑한 책은 손자병법뿐이다」 — 조조의 자평',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 적벽대전(208) — 손자병법 6편 허실편 「피실격허」의 정수',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 조조의 강군에 맞선 손권·유비 연합의 화공(火攻) — 12편 화공편 응용',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 제갈량의 「초선차전(草船借箭)」 — 13편 용간편 정보전',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 제갈량(諸葛亮) — 5편 병세편 「세(勢)」의 운용가',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 공성계(空城計) — 삼십육계 32번째 계, 손자 「궤도(詭道)」의 극치',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 「삼국지 = 손자병법의 살아 있는 케이스북」 — 6,000자의 모든 원리가 적용됨',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC_HIST)
def hist_yi(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_HIST, n, t)
    add_title(slide, '충무공 이순신 — 손자병법의 한국 화신',
              '23전 23승 · 명량·한산·노량의 모든 승부에 손자의 그림자')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 『난중일기』 곳곳에 손자의 정신이 흐른다 — 「知彼知己」·「先勝而後求戰」',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 한산도 대첩(1592) — 「유인 + 학익진」 = 5편 병세편 「以正合 以奇勝」',
         {'font_size': 17, 'space_before': 10}),
        ('     · 정공으로 만나고 기습으로 이긴다 — 학익진(鶴翼陣)의 완성',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 명량 해전(1597) — 「13척으로 133척을 이긴」 기적',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 11편 구지편 「死地則戰」 — 「必生則死 必死則生」의 정수',
         {'font_size': 14, 'color': SUB, 'space_before': 4, 'font_name': 'Batang'}),
        ('     · 지형(좁은 명량 해협) + 시간(조류) + 사기 = 손자 「오사」 완벽 적용',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 「閑山島月明夜上戍樓」 — 손자가 가르친 「先知」와 「致人而不致於人」의 한국적 표현',
         {'font_size': 15, 'space_before': 12, 'color': SUB, 'bold': True, 'font_name': 'Batang'}),
        ('● 충무공은 「손자병법을 가장 완벽하게 실전한 동양 장군」으로 평가됨',
         {'font_size': 16, 'space_before': 10, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC_HIST)
def hist_mao(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC_HIST, n, t)
    add_title(slide, '마오쩌둥 · 호치민 — 게릴라전의 손자',
              '약자가 강자를 이긴 20세기 두 전쟁의 사상적 토대')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 마오쩌둥(毛澤東) — 손자병법의 가장 위대한 20세기 응용가',
         {'font_size': 18, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 16자 전략 — 「敵進我退 敵駐我擾 敵疲我打 敵退我追」',
         {'font_size': 17, 'space_before': 10, 'font_name': 'Batang'}),
        ('     · 적이 진격하면 물러나고, 머무르면 괴롭히고, 지치면 치고, 물러나면 쫓는다',
         {'font_size': 13, 'color': SUB, 'space_before': 4}),
        ('     · 6편 허실편 「致人而不致於人」 + 7편 군쟁편 「以迂爲直」의 응용',
         {'font_size': 13, 'color': SUB, 'space_before': 4}),
        ('● 『지구전론(持久戰論)』(1938) — 손자 2편 작전편 「졸속」의 변용',
         {'font_size': 17, 'space_before': 12}),
        ('● 호치민 — 인도차이나 전쟁·베트남 전쟁에서 손자병법 직접 인용',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 「以迂爲直」 — 정글의 우회 작전 / 「죽도록 작은 비용으로 큰 효과」',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 20세기 양대 게릴라전 — 「약자가 강자를 이긴」 손자의 결정적 응용 사례',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


# ============== Ⅹ. 동서양 후대 영향 ==============
SEC8 = 'Ⅹ. 동서양 후대 영향'

@S(SEC8)
def viii_east(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '동아시아의 손자병법',
              '중국·한국·일본 — 2,500년의 영향사')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 중국 — 조조·제갈량·이정·악비 등 역대 명장의 필독서',
         {'font_size': 17, 'space_before': 4}),
        ('     · 조조 『위무주손자』 — 현존 가장 권위 있는 주석',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 송 무경칠서의 첫머리 — 무과 시험 필수 과목',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 한국 — 충무공 이순신의 애독서 · 조선 무과의 강경 과목',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 『난중일기』 곳곳에 「知彼知己」·「先勝而後求戰」의 정신',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 일본 — 다케다 신겐의 「풍림화산(風林火山)」 깃발 (제7편 군쟁 인용)',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 도쿠가와 이에야스 평생 옆에 두고 읽음',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('     · 메이지 이후 — 일본 육사·해사의 정식 교재',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 베트남 — 호치민이 인용한 「以迂爲直」의 게릴라 전략',
         {'font_size': 17, 'space_before': 10}),
    ])


@S(SEC8)
def viii_west(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '서양으로 건너간 손자병법',
              '1772년 아미오의 프랑스어 번역 → 나폴레옹 → 슈워츠코프')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1772 — 예수회 선교사 아미오(J.J.M. Amiot)의 프랑스어 첫 번역',
         {'font_size': 17, 'space_before': 4}),
        ('     · 「Art Militaire des Chinois」 — 중국인의 전쟁술',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 나폴레옹 — 손자병법을 깊이 연구했다는 기록 · 일부 전기에 언급',
         {'font_size': 17, 'space_before': 10}),
        ('● 20세기 — 영어 번역(Lionel Giles, 1910)으로 본격 확산',
         {'font_size': 17, 'space_before': 10}),
        ('● 1991 걸프전 — 슈워츠코프(Norman Schwarzkopf) 사령관',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     · 「내 책상엔 늘 손자병법이 있었다」 — 우회 기동(left hook)이 「以迂爲直」 응용',
         {'font_size': 14, 'color': SUB, 'space_before': 4}),
        ('● 마오쩌둥의 게릴라 전술 · 미 해병대·서점가 베스트셀러',
         {'font_size': 17, 'space_before': 10}),
        ('● 영문판이 동양 고전 중 가장 많이 팔린 책 — 서양에서 가장 사랑받는 동양 병서',
         {'font_size': 16, 'space_before': 12, 'color': ACCENT, 'bold': True}),
    ])


@S(SEC8)
def viii_modern(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC8, n, t)
    add_title(slide, '현대 응용 — 군사를 넘어 경영·외교·일상으로')
    rows = [
        ('경영 전략',  '마이클 포터의 경쟁 전략 · 블루오션 · 차별화 — 「避實擊虛」의 직역'),
        ('협상학',     '하버드 협상론 · BATNA · 정보 비대칭 — 「兵者詭道」의 응용'),
        ('마케팅',    '포지셔닝 · 모멘텀 · 차별화 — 「勢」와 「奇正」의 활용'),
        ('스포츠',     '농구·축구의 전술 다양화 — 풍림화산의 4가지 페이스'),
        ('정보·HR',   '인텔리전스 · 인재 분석 · 인사 평가 — 「先知」와 「五間」'),
        ('자기 관리',  '「知己」가 곧 자기 인식 · 「先勝而後求戰」이 곧 준비된 삶'),
        ('IT·스타트업', '아마존 베조스 등 — 「拙速」 = 「Move fast, ship fast」'),
        ('일본 기업', '소니·도요타의 「변통」 — 「九變」과 「兵形象水」의 응용'),
    ]
    for i, (cat, desc) in enumerate(rows):
        y = Inches(2.3 + i * 0.55)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.8), Inches(0.45), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(2.8), Inches(0.45),
                    cat, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.7), y + Inches(0.02), Inches(9.2), Inches(0.42),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅺ. 오늘 우리에게 + 마무리 ==============
SEC10 = 'Ⅺ. 오늘 우리에게'

@S(SEC10)
def x_today(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '손자병법이 오늘 우리에게 일러주는 10가지')
    items = [
        '1. 신중하라 — 큰 결정 앞에서는 「오사칠계」로 다섯 차원을 분석',
        '2. 싸우지 않고 이기는 길을 먼저 찾아라 — 부전승의 우선순위',
        '3. 적보다 자기를 먼저 알라 — 「知己」가 「知彼」에 선행',
        '4. 빠른 부족이 늦은 완벽보다 낫다 — 졸속의 원리',
        '5. 준비된 자에게 승부는 이미 결정 — 「先勝而後求戰」',
        '6. 정공과 기습을 자유로이 전환하라 — 기정의 운용',
        '7. 강한 곳을 피하고 약한 곳을 쳐라 — 「避實擊虛」',
        '8. 직진이 항상 빠른 길이 아니다 — 우직지계',
        '9. 한 덕목에 고착되면 그것이 약점 — 장수의 다섯 위험',
        '10. 정보가 가장 적은 비용으로 가장 큰 효과 — 「先知」',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.3 + i * 0.45)
        add_textbox(slide, Inches(0.7), y, Inches(12.0), Inches(0.4),
                    txt, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC10)
def x_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '손자병법, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 손무가 BC 6세기 말 오나라에서 합려에게 바친 13편 6,000자.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 동양 최고(最古)·최고(最高)의 병서, 무경칠서의 첫머리.',
         {'font_size': 18, 'space_before': 8}),
        ('● 6대 핵심 사상 — 부전승·지피지기·오사칠계·기정허실·세·궤도.',
         {'font_size': 18, 'space_before': 8}),
        ('● 풍림화산·우직지계·졸속·선승이후구전·사지즉전·선지 — 일상에 살아 있는 표현.',
         {'font_size': 18, 'space_before': 8}),
        ('● 조조의 주석 · 이순신의 애독서 · 다케다 신겐의 깃발 · 슈워츠코프의 책상.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「병법서이면서 싸우지 말라는」 위대한 역설의 책.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC10)
def x_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
                '不 戰 而 屈 人 之 兵',
                font_size=66, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.9),
                '善 之 善 者 也',
                font_size=60, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(4.3), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
                '싸 우 지  않 고  적 의  군 대 를  굴 복 시 키 는  것 이',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
                '최 선  중 의  최 선 이 다',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                '— 손자병법 제3편 모공(謀攻)',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '孫  子  兵  法',
                font_size=20, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\손자병법.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
