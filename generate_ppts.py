"""
4대 고전 PPT 제작 스크립트
논어, 맹자, 순자, 손자병법 각각의 발표자료를 PowerPoint로 생성한다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 디자인 상수 ──
BG_DARK = RGBColor(0x1B, 0x2A, 0x4A)
BG_LIGHTER = RGBColor(0x22, 0x35, 0x5C)
BG_ACCENT = RGBColor(0x19, 0x3D, 0x6B)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
TABLE_BG = RGBColor(0x2A, 0x3F, 0x6A)
TABLE_HEADER_BG = RGBColor(0x14, 0x24, 0x3E)
SUBTITLE_COLOR = RGBColor(0xA0, 0xB8, 0xD8)

FONT_KR = "맑은 고딕"
FONT_EN = "Arial"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_KR, line_spacing=1.3):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name=FONT_KR,
                  space_before=0, space_after=4, line_spacing=None):
    """텍스트프레임에 새 문단 추가"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return p


def add_decorated_line(slide, left, top, width, color=GOLD, height=Pt(2)):
    """장식선 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def make_title_slide(prs, title_kr, title_hanja, author, era, subtitle=""):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # 상단 장식선
    add_decorated_line(slide, Inches(1), Inches(1.2), Inches(11.3), GOLD, Pt(3))

    # 한자 제목
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
                title_hanja, font_size=52, color=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 한글 제목
    add_textbox(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1),
                title_kr, font_size=36, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 하단 장식선
    add_decorated_line(slide, Inches(4), Inches(3.7), Inches(5.3), GOLD, Pt(2))

    # 저자 & 시대
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.7),
                f"저자: {author}  |  시대: {era}", font_size=22, color=SUBTITLE_COLOR,
                alignment=PP_ALIGN.CENTER)

    if subtitle:
        add_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.7),
                    subtitle, font_size=18, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # 하단 장식
    add_decorated_line(slide, Inches(1), Inches(6.3), Inches(11.3), GOLD, Pt(3))


def make_toc_slide(prs, items):
    """목차 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
                "목차", font_size=36, color=GOLD, bold=True)
    add_decorated_line(slide, Inches(0.8), Inches(1.1), Inches(4), GOLD)

    tf = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(5.5),
                     "", font_size=20, color=WHITE)
    tf.paragraphs[0].text = ""
    for i, item in enumerate(items, 1):
        p = add_paragraph(tf, f"{i:02d}   {item}", font_size=20, color=WHITE,
                          space_before=4, space_after=6)


def make_section_slide(prs, title, content_lines, bg=BG_DARK):
    """일반 내용 슬라이드 (제목 + 줄 단위 본문)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                title, font_size=32, color=GOLD, bold=True)
    add_decorated_line(slide, Inches(0.8), Inches(1.0), Inches(5), GOLD)

    tf = add_textbox(slide, Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.8),
                     "", font_size=18, color=WHITE, line_spacing=1.0)
    tf.paragraphs[0].text = ""

    for line in content_lines:
        if isinstance(line, tuple):
            text, size, clr, bld = line
        else:
            text, size, clr, bld = line, 18, WHITE, False
        add_paragraph(tf, text, font_size=size, color=clr, bold=bld,
                      space_after=4)
    return slide


def make_quote_slide(prs, title, quotes):
    """명구절 슬라이드 — quotes: [(원문, 독음, 해석, 출처), ...]"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                title, font_size=32, color=GOLD, bold=True)
    add_decorated_line(slide, Inches(0.8), Inches(1.0), Inches(5), GOLD)

    tf = add_textbox(slide, Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.8),
                     "", font_size=16, color=WHITE, line_spacing=1.0)
    tf.paragraphs[0].text = ""

    for orig, reading, meaning, source in quotes:
        add_paragraph(tf, orig, font_size=17, color=GOLD, bold=True, space_before=8)
        add_paragraph(tf, reading, font_size=14, color=LIGHT_GRAY, space_after=2)
        add_paragraph(tf, f"→ {meaning}", font_size=15, color=WHITE, space_after=2)
        add_paragraph(tf, f"  — {source}", font_size=13, color=SUBTITLE_COLOR, space_after=6)


def make_table_slide(prs, title, headers, rows, col_widths=None):
    """표 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                title, font_size=32, color=GOLD, bold=True)
    add_decorated_line(slide, Inches(0.8), Inches(1.0), Inches(5), GOLD)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(11.5)
    tbl_left = Inches(0.9)
    tbl_top = Inches(1.4)
    tbl_height = Inches(5.5)

    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # 헤더
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = GOLD
            p.font.name = FONT_KR
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 데이터
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE
                p.font.name = FONT_KR
                p.alignment = PP_ALIGN.LEFT if j > 0 else PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_BG if i % 2 == 0 else BG_LIGHTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide


def make_comparison_slide(prs, title, content_lines):
    """비교 슬라이드 — section_slide와 동일하나 배경 약간 다름"""
    return make_section_slide(prs, title, content_lines, bg=BG_ACCENT)


def make_closing_slide(prs, title_kr, one_line_summary, key_message):
    """마무리 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_decorated_line(slide, Inches(1), Inches(1.5), Inches(11.3), GOLD, Pt(3))

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
                title_kr, font_size=40, color=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.5),
                one_line_summary, font_size=20, color=WHITE,
                alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    add_decorated_line(slide, Inches(4), Inches(4.8), Inches(5.3), GOLD, Pt(2))

    add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(1),
                key_message, font_size=18, color=SUBTITLE_COLOR,
                alignment=PP_ALIGN.CENTER)

    add_decorated_line(slide, Inches(1), Inches(6.3), Inches(11.3), GOLD, Pt(3))


# ════════════════════════════════════════════════
# 논어 PPT
# ════════════════════════════════════════════════
def create_lunyu_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. 표지
    make_title_slide(prs, "논어", "論語 (The Analects)", "공자(孔子)와 제자들",
                     "춘추시대 (BC 551~479)",
                     "사서(四書)의 으뜸 · 20편 482장 · 유교 사상의 핵심")

    # 2. 목차
    make_toc_slide(prs, [
        "개요 — 논어란 무엇인가",
        "구성 — 20편의 구조",
        "핵심 사상 ① 인(仁)",
        "핵심 사상 ② 예(禮)와 정치",
        "핵심 사상 ③ 군자(君子)",
        "명구절 (상)",
        "명구절 (하)",
        "구조적 특징",
        "현대적 의의",
        "다른 고전과의 비교",
        "마무리"
    ])

    # 3. 개요
    make_section_slide(prs, "논어란 무엇인가", [
        ("정의: 공자와 제자들의 언행을 기록한 유교 핵심 경전", 20, WHITE, True),
        ("  • 총 20편, 482장, 약 600여 문장", 18, WHITE, False),
        ("  • '論(논의하다) + 語(말씀)' = 토론하여 정리한 말씀", 18, WHITE, False),
        ("", 10, WHITE, False),
        ("공자(孔子, BC 551~479)", 20, GOLD, True),
        ("  • 이름: 공구(孔丘), 자: 중니(仲尼)", 18, WHITE, False),
        ("  • 노(魯)나라 출신 · 정치인, 사상가, 교육자", 18, WHITE, False),
        ("  • 제자 약 3,000명, 뛰어난 자 72명(칠십이현)", 18, WHITE, False),
        ("  • 14년간 천하 주유 후 고향에서 후학 양성", 18, WHITE, False),
        ("", 10, WHITE, False),
        ("편찬 과정", 20, GOLD, True),
        ("  • 1차: 직계 제자(중궁, 자유, 자하) 주도", 18, WHITE, False),
        ("  • 2차: 증자 사후 보충 → 3차: 전국시대 추가", 18, WHITE, False),
        ("  • 전한 말 장우(張禹) 편집본 기반, 후한에서 현재 형태 확정", 18, WHITE, False),
    ])

    # 4. 구성
    make_table_slide(prs, "논어의 구성 — 20편",
        ["구분", "편명", "핵심 주제"],
        [
            ["상론 1", "학이(學而)", "배움·효제·수신 — 논어의 총론"],
            ["상론 2", "위정(爲政)", "덕치·효도·학문의 단계"],
            ["상론 3", "팔일(八佾)", "예악의 본질, 인과 예의 관계"],
            ["상론 4", "이인(里仁)", "인(仁)의 내면, 충서(忠恕)"],
            ["상론 5~6", "공야장·옹야", "인물 평가, 문질빈빈"],
            ["상론 7~10", "술이~향당", "공자의 학문관·일상·절조"],
            ["하론 11~12", "선진·안연", "과유불급, 극기복례"],
            ["하론 13~15", "자로~위령공", "정명·정치·군자론 총결산"],
            ["하론 16~18", "계씨~미자", "수양·예악·은일과 입세"],
            ["하론 19~20", "자장·요왈", "제자 전승, 성왕의 이상"],
        ],
        col_widths=[1.5, 3.0, 7.0]
    )

    # 5. 핵심 사상 ① 인(仁)
    make_section_slide(prs, "핵심 사상 ① — 인(仁): 사람다움의 최고 덕목", [
        ("仁 — 사람(人) 둘이 함께하는 모양: 사람과 사람의 바람직한 관계", 20, GOLD, True),
        ("", 8, WHITE, False),
        ("克己復禮爲仁 — 자기를 이기고 예로 돌아감이 인이다 (안연 1장)", 17, WHITE, False),
        ("仁者愛人 — 인이란 사람을 사랑하는 것이다 (안연 22장)", 17, WHITE, False),
        ("己欲立而立人 己欲達而達人 — 자기가 서면 남을 세워준다 (옹야 30장)", 17, WHITE, False),
        ("己所不欲 勿施於人 — 원하지 않는 것을 남에게 하지 말라 (위령공 24장)", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("인의 실천 방법 — 충서(忠恕)", 20, GOLD, True),
        ("  • 충(忠): 자기 마음을 다하는 것 (盡己)", 18, WHITE, False),
        ("  • 서(恕): 자기를 미루어 남에게 미치는 것 (推己及人)", 18, WHITE, False),
        ("  • '夫子之道 忠恕而已矣' — 선생님의 도는 충서일 뿐이다 (이인 15장)", 17, LIGHT_GRAY, False),
        ("", 8, WHITE, False),
        ("인의 출발점 — 효(孝)", 20, GOLD, True),
        ("  • 孝弟也者 其爲仁之本與 — 효제는 인의 근본이다 (학이 2장)", 17, WHITE, False),
    ])

    # 6. 핵심 사상 ② 예와 정치
    make_section_slide(prs, "핵심 사상 ② — 예(禮)와 덕치(德治)", [
        ("예(禮) — 인(仁)의 외적 표현", 20, GOLD, True),
        ("  • 종교 제례에서 일상 행위 규범으로 확장", 18, WHITE, False),
        ("  • 人而不仁 如禮何 — 어질지 못하면 예를 어찌하겠는가 (팔일 3장)", 17, WHITE, False),
        ("  • 인 없는 예 = 빈 껍데기 / 예 없는 인 = 실현 불가", 17, LIGHT_GRAY, False),
        ("", 8, WHITE, False),
        ("덕치(德治) — 덕으로 다스림", 20, GOLD, True),
        ("  • 爲政以德 譬如北辰 — 덕으로 정치하면 북극성과 같다 (위정 1장)", 17, WHITE, False),
        ("  • 其身正 不令而行 — 자신이 바르면 명령 없이도 행해진다 (자로 6장)", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("정명(正名) — 이름을 바르게 함", 20, GOLD, True),
        ("  • 名不正則言不順 — 명분이 바르지 않으면 일이 이루어지지 못한다 (자로 3장)", 17, WHITE, False),
        ("  • 君君臣臣父父子子 — 각자 자기 역할에 맞게 (안연 11장)", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("교육 철학", 20, GOLD, True),
        ("  • 學而不思則罔 思而不學則殆 — 배움과 사유의 균형 (위정 15장)", 17, WHITE, False),
        ("  • 有敎無類 — 가르침에 차별이 없다 (위령공 39장)", 17, WHITE, False),
    ])

    # 7. 핵심 사상 ③ 군자
    make_section_slide(prs, "핵심 사상 ③ — 군자(君子): 이상적 인간상", [
        ("공자가 귀족 계층 개념을 '도덕적으로 수양된 이상적 인간'으로 재정의", 19, GOLD, True),
        ("", 8, WHITE, False),
        ("군자 vs 소인 대비", 20, GOLD, True),
        ("  • 和而不同 / 同而不和 — 조화롭되 같지 않다 (자로 23장)", 17, WHITE, False),
        ("  • 周而不比 / 比而不周 — 두루 사귀되 편당하지 않는다 (위정 14장)", 17, WHITE, False),
        ("  • 坦蕩蕩 / 長戚戚 — 마음이 넓다 / 항상 근심한다 (술이 37장)", 17, WHITE, False),
        ("  • 喩於義 / 喩於利 — 의에 밝다 / 이에 밝다 (이인 16장)", 17, WHITE, False),
        ("  • 求諸己 / 求諸人 — 자신에게서 구한다 / 남에게서 구한다 (위령공 21장)", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("군자의 세 덕목 — 삼달덕(三達德)", 20, GOLD, True),
        ("  • 知者不惑 — 지혜로운 자는 미혹되지 않는다", 18, WHITE, False),
        ("  • 仁者不憂 — 어진 자는 근심하지 않는다", 18, WHITE, False),
        ("  • 勇者不懼 — 용감한 자는 두려워하지 않는다 (자한 29장)", 18, WHITE, False),
    ])

    # 8. 명구절 (상)
    make_quote_slide(prs, "명구절 (상)", [
        ("學而時習之 不亦說乎", "학이시습지 불역열호",
         "배우고 때때로 익히면 기쁘지 아니한가", "학이편 1장"),
        ("溫故而知新 可以爲師矣", "온고이지신 가이위사의",
         "옛것을 익히고 새것을 알면 스승이 될 만하다", "위정편 11장"),
        ("朝聞道 夕死可矣", "조문도 석사가의",
         "아침에 도를 들으면 저녁에 죽어도 좋다", "이인편 8장"),
        ("三人行 必有我師焉", "삼인행 필유아사언",
         "세 사람이 가면 반드시 나의 스승이 있다", "술이편 22장"),
        ("歲寒然後 知松柏之後彫也", "세한연후 지송백지후조야",
         "추워진 뒤에야 소나무의 절개를 안다", "자한편 28장"),
    ])

    # 9. 명구절 (하)
    make_quote_slide(prs, "명구절 (하)", [
        ("克己復禮爲仁", "극기복례위인",
         "자기를 이기고 예로 돌아감이 인이다", "안연편 1장"),
        ("己所不欲 勿施於人", "기소불욕 물시어인",
         "자기가 원하지 않는 것을 남에게 하지 마라", "위령공편 24장"),
        ("君子和而不同 小人同而不和", "군자화이부동 소인동이불화",
         "군자는 조화롭되 같지 않고, 소인은 같되 조화롭지 않다", "자로편 23장"),
        ("過猶不及", "과유불급",
         "지나친 것은 못 미치는 것과 같다", "선진편 16장"),
        ("志士仁人 有殺身以成仁", "지사인인 유살신이성인",
         "뜻있는 선비는 목숨을 바쳐 인을 이루는 일이 있다", "위령공편 9장"),
    ])

    # 10. 구조적 특징
    make_section_slide(prs, "논어의 구조적 특징", [
        ("문체의 다양성", 20, GOLD, True),
        ("  • 子曰(자왈) — 공자의 직접 발언 (가장 기본 형식)", 17, WHITE, False),
        ("  • 문답체 — 제자의 질문 + 공자의 답 (안연, 자로편 등)", 17, WHITE, False),
        ("  • 서사체 — 사건·일화 서술 (미자, 향당편)", 17, WHITE, False),
        ("  • 제자 발언 — 19편 자장은 공자 발언이 없는 유일한 편", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("인(仁)의 등장 패턴", 20, GOLD, True),
        ("  • 논어 전체에 仁은 109회 등장", 17, WHITE, False),
        ("  • 안연편·이인편에 집중 / 계씨편에서는 단 한 번도 등장하지 않음", 17, WHITE, False),
        ("  • 제자마다 다른 인의 정의 → 인재시교(因材施敎)의 증거", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("수미상관(首尾相關)", 20, GOLD, True),
        ("  • 학이편 1장(시작): '學而時習之... 不亦君子乎'", 17, WHITE, False),
        ("  • 요왈편 3장(결론): '不知命 無以爲君子也'", 17, WHITE, False),
        ("  • 배움에서 출발 → 군자로 귀결하는 닫힌 구조", 17, LIGHT_GRAY, False),
    ])

    # 11. 현대적 의의
    make_section_slide(prs, "현대적 의의", [
        ("화이부동(和而不同) — 다원주의 사회의 원리", 20, GOLD, True),
        ("  서로 다른 종교·문화·가치관이 공존하는 다양성 속 조화", 18, WHITE, False),
        ("", 6, WHITE, False),
        ("기소불욕 물시어인(己所不欲 勿施於人) — 보편 윤리", 20, GOLD, True),
        ("  동서양을 관통하는 황금률. 세계인권선언, 기업 윤리, AI 윤리에 적용", 18, WHITE, False),
        ("", 6, WHITE, False),
        ("학이불사즉망(學而不思則罔) — 비판적 사고", 20, GOLD, True),
        ("  정보 과잉 시대에 단순 암기가 아닌 비판적 사고와 성찰의 중요성", 18, WHITE, False),
        ("", 6, WHITE, False),
        ("덕치·솔선수범 — 리더십", 20, GOLD, True),
        ("  서번트 리더십, 윤리적 리더십과 상통. 수기치인(修己治人)의 정신", 18, WHITE, False),
        ("", 6, WHITE, False),
        ("인(仁) — 인공지능 시대의 인문학", 20, GOLD, True),
        ("  기술만으로 해결할 수 없는 '사람다움'의 방향을 제시", 18, WHITE, False),
    ])

    # 12. 비교
    make_table_slide(prs, "논어와 다른 고전의 비교",
        ["비교 항목", "논어", "맹자", "순자"],
        [
            ["시대", "춘추 말기", "전국 중기", "전국 말기"],
            ["형식", "어록·문답", "대화·논변", "논설문"],
            ["인성론", "성상근(性相近)", "성선설(性善說)", "성악설(性惡說)"],
            ["핵심 덕목", "인(仁)", "인의(仁義)", "예(禮)"],
            ["수양 방법", "학·극기복례", "존심양성·확충", "화성기위·적(積)"],
            ["천(天)의 성격", "도덕적 천", "의지적 천", "자연적 천"],
            ["정치론", "덕치·정명", "왕도·역성혁명", "왕도·예법 병용"],
        ],
        col_widths=[2.5, 3.0, 3.0, 3.0]
    )

    # 13. 마무리
    make_closing_slide(prs, "논어를 한 문장으로",
        "배움(學)에서 출발하여, 인(仁)을 핵심으로, 예(禮)를 형식으로,\n"
        "충서(忠恕)를 방법으로, 군자(君子)를 이상으로 삼아,\n"
        "수기(修己)에서 치인(治人)으로 나아가는\n"
        "— 사람다움의 총체적 설계도.",
        "「半部論語治天下」 — 논어 반 권으로 천하를 다스린다")

    path = "Works/논어/논어_발표자료.pptx"
    prs.save(path)
    return path


# ════════════════════════════════════════════════
# 맹자 PPT
# ════════════════════════════════════════════════
def create_mengzi_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. 표지
    make_title_slide(prs, "맹자", "孟子 (Mencius)", "맹가(孟軻), 아성(亞聖)",
                     "전국시대 (BC 372~289경)",
                     "사서(四書) · 7편 14장 · 성선설과 왕도정치의 체계")

    # 2. 목차
    make_toc_slide(prs, [
        "개요 — 맹자란 무엇인가",
        "구성 — 7편 14장의 구조",
        "핵심 사상 ① 성선설과 사단",
        "핵심 사상 ② 왕도정치와 민본",
        "핵심 사상 ③ 호연지기와 수양론",
        "명구절 (상)",
        "명구절 (하)",
        "구조적 특징",
        "현대적 의의",
        "다른 고전과의 비교",
        "마무리"
    ])

    # 3. 개요
    make_section_slide(prs, "맹자란 무엇인가", [
        ("정의: 전국시대 유학자 맹자의 사상과 언행을 담은 유교 경전", 20, WHITE, True),
        ("  • 7편 14장(상·하), 260여 장, 약 35,000여 자", 18, WHITE, False),
        ("  • 논어가 간결한 어록이라면, 맹자는 체계적 논변과 비유", 18, WHITE, False),
        ("", 8, WHITE, False),
        ("맹자(孟軻, BC 372~289경)", 20, GOLD, True),
        ("  • 추(鄒)나라 출신, 자사(子思)의 문인에게 수학", 18, WHITE, False),
        ("  • 공자→증자→자사→맹자의 도통(道統) 계보", 18, WHITE, False),
        ("  • 맹모삼천지교(孟母三遷之敎) — 교육 환경의 중요성", 18, WHITE, False),
        ("  • 제·위·송·등 여러 나라를 유세하며 왕도정치 설파", 18, WHITE, False),
        ("  • 존칭: 아성(亞聖) — 공자 다음의 성인", 18, WHITE, False),
        ("", 8, WHITE, False),
        ("편찬 과정", 20, GOLD, True),
        ("  • 자술설: 맹자 본인이 만장·공손추 등과 직접 편찬", 18, WHITE, False),
        ("  • 후한 조기(趙岐)의 주석으로 7편 14장 체제 확정", 18, WHITE, False),
    ])

    # 4. 구성
    make_table_slide(prs, "맹자의 구성 — 7편 14장",
        ["편", "편명", "장 수", "핵심 주제"],
        [
            ["1", "양혜왕(梁惠王)", "상7+하16", "왕도정치, 인정, 의리 분별"],
            ["2", "공손추(公孫丑)", "상9+하14", "호연지기, 사단, 왕패 구분"],
            ["3", "등문공(滕文公)", "상5+하10", "정전제, 오륜, 이단 비판"],
            ["4", "이루(離婁)", "상28+하33", "인의 실천, 수신, 군신 관계"],
            ["5", "만장(萬章)", "상9+하9", "성인 행적, 경전 해석"],
            ["6", "고자(告子)", "상20+하16", "성선설 논변, 의(義)의 내재성"],
            ["7", "진심(盡心)", "상46+하38", "수양론 총결, 사상의 집대성"],
        ],
        col_widths=[1.0, 3.0, 2.0, 5.5]
    )

    # 5. 핵심 사상 ① 성선설과 사단
    make_section_slide(prs, "핵심 사상 ① — 성선설(性善說)과 사단(四端)", [
        ("人性之善也 猶水之就下也 — 인간 본성이 선한 것은 물이 아래로 흐르는 것과 같다", 17, GOLD, True),
        ("", 6, WHITE, False),
        ("성선설 — 인간 본성은 선하다", 20, GOLD, True),
        ("  • 악(惡)은 본성이 아니라 환경의 영향이나 본성을 기르지 못한 결과", 17, WHITE, False),
        ("  • 수양의 목표: 잃어버린 본성을 되찾는 것 (구방심, 求放心)", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("사단(四端) — 선한 본성의 네 가지 싹", 20, GOLD, True),
        ("  • 측은지심(惻隱之心) → 인(仁): 남의 고통을 차마 보지 못하는 마음", 17, WHITE, False),
        ("  • 수오지심(羞惡之心) → 의(義): 부끄러워하고 미워하는 마음", 17, WHITE, False),
        ("  • 사양지심(辭讓之心) → 예(禮): 남에게 양보하는 마음", 17, WHITE, False),
        ("  • 시비지심(是非之心) → 지(智): 옳고 그름을 가리는 마음", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("  '우물에 빠지려는 아이(孺子入井)' — 누구나 측은지심을 느낌 → 성선의 증거", 17, LIGHT_GRAY, False),
        ("  사단은 완성된 덕이 아니라 싹(端) → 기르고 확충해야(擴充) 완전한 덕이 됨", 17, LIGHT_GRAY, False),
    ])

    # 6. 핵심 사상 ② 왕도정치와 민본
    make_section_slide(prs, "핵심 사상 ② — 왕도정치와 민본사상", [
        ("왕도정치(王道政治) — 덕으로 다스리는 이상 정치", 20, GOLD, True),
        ("  • 왕도(덕·인) vs 패도(힘·이익): 진심 복종 vs 겉으로만 복종", 17, WHITE, False),
        ("  • 항산항심(恒産恒心): 안정된 생업이 있어야 안정된 마음이 있다", 17, WHITE, False),
        ("  • 정전제(井田制), 세금 경감, 교육 진흥", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("민본사상(民本思想) — 백성이 근본이다", 20, GOLD, True),
        ("  • 民爲貴 社稷次之 君爲輕 — 백성이 가장 귀하다 (진심 하)", 17, WHITE, False),
        ("  • 군주는 백성을 위해 존재 / 백성이 군주를 위해 존재하는 것이 아님", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("역성혁명(易姓革命)", 20, GOLD, True),
        ("  • 덕을 잃은 폭군은 더 이상 군주가 아닌 일부(一夫)에 불과", 17, WHITE, False),
        ("  • '聞誅一夫紂矣 未聞弑君也' — 필부 주를 죽였다 들었지 임금을 시해했다 듣지 못했다", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("의리론(義利論)", 20, GOLD, True),
        ("  • 王何必曰利 亦有仁義而已矣 — 하필 이익을 말씀하십니까? (양혜왕 상 1장)", 17, WHITE, False),
    ])

    # 7. 핵심 사상 ③ 호연지기
    make_section_slide(prs, "핵심 사상 ③ — 호연지기와 수양론", [
        ("호연지기(浩然之氣) — 도덕적 용기의 원천", 20, GOLD, True),
        ("  其爲氣也 至大至剛 以直養而無害 則塞于天地之間", 17, LIGHT_GRAY, False),
        ("  → 그 기운은 지극히 크고 강하니, 정직함으로 기르면 천지 사이에 가득 찬다", 17, WHITE, False),
        ("  • 도덕적 실천이 쌓여 생기는 정신적 에너지", 17, WHITE, False),
        ("  • 억지로 조장하면 안 됨 — 조장(助長) 우화의 교훈", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("수양의 단계", 20, GOLD, True),
        ("  1. 구방심(求放心) — 잃어버린 마음 되찾기", 17, WHITE, False),
        ("  2. 존심양성(存心養性) — 마음 보존, 본성 기르기", 17, WHITE, False),
        ("  3. 과욕(寡欲) — 욕심 줄여 본성 드러내기", 17, WHITE, False),
        ("  4. 양기(養氣) — 호연지기 기르기", 17, WHITE, False),
        ("  5. 진심지성(盡心知性) — 마음 다하여 하늘을 아는 경지", 17, WHITE, False),
        ("", 8, WHITE, False),
        ("  盡其心者 知其性也 知其性 則知天矣", 17, GOLD, False),
        ("  → 마음을 다하면 본성을 알고, 본성을 알면 하늘을 안다 (진심 상 1장)", 17, WHITE, False),
    ])

    # 8. 명구절 (상)
    make_quote_slide(prs, "명구절 (상)", [
        ("王何必曰利 亦有仁義而已矣", "왕하필왈리 역유인의이이의",
         "왕께서 하필 이익을 말씀하십니까? 오직 인의가 있을 뿐입니다", "양혜왕 상 1장"),
        ("惻隱之心 仁之端也", "측은지심 인지단야",
         "측은히 여기는 마음은 인(仁)의 싹이다", "공손추 상 6장"),
        ("天時不如地利 地利不如人和", "천시불여지리 지리불여인화",
         "하늘의 때보다 땅의 이로움, 땅보다 사람의 화합이 중요하다", "공손추 하 1장"),
        ("民爲貴 社稷次之 君爲輕", "민위귀 사직차지 군위경",
         "백성이 가장 귀하고, 사직이 다음이며, 임금이 가장 가볍다", "진심 하 14장"),
        ("五十步百步", "오십보백보",
         "정도의 차이일 뿐 본질은 같다", "양혜왕 상 3장"),
    ])

    # 9. 명구절 (하)
    make_quote_slide(prs, "명구절 (하)", [
        ("天將降大任於是人也 必先苦其心志", "천장강대임어시인야 필선고기심지",
         "하늘이 큰 일을 맡기려면 먼저 그 마음을 괴롭힌다", "고자 하 15장"),
        ("生於憂患而死於安樂也", "생우환이사우안락야",
         "근심 속에서 살아나고 안락 속에서 죽는다", "고자 하 15장"),
        ("捨生而取義者也", "사생이취의자야",
         "삶을 버리고 의로움을 취한다", "고자 상 10장"),
        ("盡其心者 知其性也 知其性 則知天矣", "진기심자 지기성야 지기성 즉지천의",
         "마음을 다하면 본성을 알고, 본성을 알면 하늘을 안다", "진심 상 1장"),
        ("自暴者 不可與有言也 自棄者 不可與有爲也", "자포자 불가여유언야 자기자 불가여유위야",
         "스스로 해치는 자, 스스로 버리는 자와는 함께할 수 없다", "이루 상 10장"),
    ])

    # 10. 구조적 특징
    make_section_slide(prs, "맹자의 구조적 특징", [
        ("웅변적 논변", 20, GOLD, True),
        ("  • 논리적이면서도 감정에 호소하는 설득의 기술", 17, WHITE, False),
        ("  • 양혜왕·제선왕과의 정치적 대화가 대표적", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("비유와 우화의 대가", 20, GOLD, True),
        ("  • 오십보백보 — 본질적 차이 없는 차이를 꼬집음", 17, WHITE, False),
        ("  • 연목구어 — 잘못된 방법으로는 목적 달성 불가", 17, WHITE, False),
        ("  • 조장(助長) — 억지로 성장을 재촉하면 망함", 17, WHITE, False),
        ("  • 우산지목(牛山之木) — 본성을 잃어버리는 과정", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("편별 흐름", 20, GOLD, True),
        ("  양혜왕(정치 출발) → 공손추(내면의 힘) → 등문공(사회 제도)", 17, WHITE, False),
        ("  → 이루(인의 실천) → 만장(역사적 검증) → 고자(인성론 논변)", 17, WHITE, False),
        ("  → 진심(수양론 완성과 종합)", 17, WHITE, False),
    ])

    # 11. 현대적 의의
    make_section_slide(prs, "현대적 의의", [
        ("인권과 민주주의", 20, GOLD, True),
        ("  • 민위귀(民爲貴) → 국민주권·민주주의의 선구적 이념", 18, WHITE, False),
        ("  • 역성혁명론 → 저항권(right of resistance)의 동아시아적 기원", 18, WHITE, False),
        ("", 6, WHITE, False),
        ("교육과 인격 형성", 20, GOLD, True),
        ("  • 성선설 → 긍정 심리학·인본주의 교육과 연결", 18, WHITE, False),
        ("  • 조장(助長) 경계 → 과도한 교육 압박의 위험성", 18, WHITE, False),
        ("", 6, WHITE, False),
        ("리더십과 경영", 20, GOLD, True),
        ("  • 왕도 vs 패도 → 덕에 의한 리더십 vs 힘에 의한 리더십", 18, WHITE, False),
        ("  • 항산항심 → 경제적 안정이 사회 안정의 기반", 18, WHITE, False),
        ("", 6, WHITE, False),
        ("윤리와 가치관", 20, GOLD, True),
        ("  • 의리론(義利之辨) → 이익보다 원칙을 우선하는 기업 윤리", 18, WHITE, False),
        ("  • 자포자기 경계 → 자기효능감, 성장 마인드셋", 18, WHITE, False),
    ])

    # 12. 비교
    make_table_slide(prs, "맹자와 다른 고전의 비교",
        ["항목", "공자(논어)", "맹자", "순자"],
        [
            ["핵심 덕목", "인(仁)", "인의(仁義)", "예(禮)"],
            ["인성론", "직접 언급 적음", "성선설 체계화", "성악설"],
            ["정치사상", "덕치(德治)", "왕도·민본·혁명론", "왕도·예법 병용"],
            ["수양론", "극기복례", "존심양성·구방심", "화성기위"],
            ["문체", "간결한 어록", "웅변적 논변", "체계적 논설문"],
            ["대화 상대", "주로 제자들", "제후·학자", "직하학궁 학자들"],
            ["존칭", "지성(至聖)", "아성(亞聖)", "후성(後聖)"],
        ],
        col_widths=[2.0, 3.2, 3.2, 3.1]
    )

    # 13. 마무리
    make_closing_slide(prs, "맹자를 한 문장으로",
        "인간이 태어나면서부터 선한 본성을 지니고 있으며(성선설),\n"
        "이 싹을 기르고 확충하는 것이 수양이요,\n"
        "덕으로 다스리는 왕도정치와 백성이 가장 귀하다는 민본사상을 핵심으로,\n"
        "의로움을 이익보다 앞세우는 도덕적 용기의 철학.",
        "「民爲貴 社稷次之 君爲輕」 — 백성이 가장 귀하다")

    path = "Works/맹자/맹자_발표자료.pptx"
    prs.save(path)
    return path


# ════════════════════════════════════════════════
# 순자 PPT
# ════════════════════════════════════════════════
def create_xunzi_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. 표지
    make_title_slide(prs, "순자", "荀子 (Xunzi)", "순황(荀況), 순경(荀卿)",
                     "전국시대 말기 (BC 313~238경)",
                     "32편 · 성악설 · 예론 · 천론 · 유가 철학의 체계적 논증")

    # 2. 목차
    make_toc_slide(prs, [
        "개요 — 순자란 무엇인가",
        "구성 — 32편의 구조",
        "핵심 사상 ① 성악설",
        "핵심 사상 ② 예론과 천론",
        "핵심 사상 ③ 수양론과 정명론",
        "명구절 (상)",
        "명구절 (하)",
        "구조적 특징",
        "현대적 의의",
        "다른 고전과의 비교",
        "마무리"
    ])

    # 3. 개요
    make_section_slide(prs, "순자란 무엇인가", [
        ("정의: 전국시대 말기 유학자 순황의 사상을 담은 32편의 유가 철학서", 20, WHITE, True),
        ("  • 논어·맹자가 어록 중심이라면, 순자는 최초의 논설문 형식", 18, WHITE, False),
        ("  • 선진 유학의 3대 저작 (논어·맹자·순자)", 18, WHITE, False),
        ("", 8, WHITE, False),
        ("순황(荀況, BC 313~238경)", 20, GOLD, True),
        ("  • 조(趙)나라 출신, 제나라 직하학궁에서 세 차례 제주(학장) 역임", 18, WHITE, False),
        ("  • 초나라 난릉령 역임 후 저술에 전념", 18, WHITE, False),
        ("  • 제자: 이사(李斯), 한비자(韓非子) — 법가에 결정적 영향", 18, WHITE, False),
        ("", 8, WHITE, False),
        ("사상사적 위치", 20, GOLD, True),
        ("  • 맹자의 성선설에 대해 성악설 주장 → 유가 내부 논쟁 촉발", 18, WHITE, False),
        ("  • 예(禮)를 사회 질서의 근본 원리로 체계화", 18, WHITE, False),
        ("  • 하늘(天)을 자연 현상으로 규정 → 합리주의 전통 수립", 18, WHITE, False),
    ])

    # 4. 구성
    make_table_slide(prs, "순자의 구성 — 32편",
        ["분류", "편", "핵심 주제"],
        [
            ["핵심 사상편", "1~9편", "권학(학문론), 수신, 비상, 유효, 왕제 등"],
            ["경세 실천편", "10~16편", "부국(경제), 왕패, 군도, 의병(군사론) 등"],
            ["철학 심화편", "17~23편", "천론, 예론, 악론, 해폐(인식론), 정명, 성악"],
            ["부록·잡편", "24~32편", "군자, 성상(운문), 격언, 공자 일화 등"],
        ],
        col_widths=[2.5, 2.5, 6.5]
    )

    # 5. 핵심 사상 ① 성악설
    make_section_slide(prs, "핵심 사상 ① — 성악설(性惡說)", [
        ("人之性惡 其善者僞也 — 사람의 본성은 악하며, 선한 것은 인위(僞)이다", 18, GOLD, True),
        ("", 6, WHITE, False),
        ("성(性)과 위(僞)의 구분", 20, GOLD, True),
        ("  • 성(性): 타고난 본성 — 이기적 욕망, 감각적 욕구 (자연적·본능적)", 17, WHITE, False),
        ("  • 위(僞): 인위적 노력 — 예의, 사양 (의식적·후천적)", 17, WHITE, False),
        ("  • 僞는 거짓이 아니라 '人+爲(사람의 행위)' = 인위적 노력", 17, LIGHT_GRAY, False),
        ("", 6, WHITE, False),
        ("화성기위(化性起僞) — 순자 사상의 핵심 공식", 20, GOLD, True),
        ("  성(性) → 위(僞) → 예의(禮義) → 법도(法度)", 18, WHITE, True),
        ("  • 성인은 본성을 변화시켜 인위를 일으키고", 17, WHITE, False),
        ("  • 인위가 일어나 예의가 생기며, 예의가 법도를 제정", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("교육 낙관론", 20, GOLD, True),
        ("  • 塗之人可以爲禹 — 길가의 보통 사람도 우(禹) 임금이 될 수 있다", 17, WHITE, False),
        ("  • 본성이 악하더라도 노력으로 성인이 될 수 있다는 평등적 수양론", 17, WHITE, False),
    ])

    # 6. 핵심 사상 ② 예론과 천론
    make_section_slide(prs, "핵심 사상 ② — 예론(禮論)과 천론(天論)", [
        ("예론 — 예는 사회 질서의 근본 원리", 20, GOLD, True),
        ("  • 예의 기원: 욕망 → 충족 불가 → 추구 → 한도 없으면 다툼 → 예가 필요", 17, WHITE, False),
        ("  • 양(養): 욕망을 적절히 충족 / 별(別): 사회적 분별 / 문(文): 문화적 형식", 17, WHITE, False),
        ("  • 禮者 養也 — 예란 기름(충족)이다. 욕망의 억압이 아니라 절도 있는 만족", 17, WHITE, False),
        ("  • 禮者 人道之極也 — 예는 인간 도의 극치", 17, LIGHT_GRAY, False),
        ("", 8, WHITE, False),
        ("천론 — 하늘은 자연이며 인간과 분리된다", 20, GOLD, True),
        ("  天行有常 不爲堯存 不爲桀亡", 17, GOLD, False),
        ("  → 하늘의 운행에는 법칙이 있으니, 요 때문에 존재하지도 걸 때문에 없어지지도 않는다", 17, WHITE, False),
        ("", 4, WHITE, False),
        ("  制天命而用之 — 천명을 제어하여 이용하라", 17, GOLD, False),
        ("  → 자연에 대한 수동적 경외가 아니라 적극적 활용", 17, WHITE, False),
        ("", 4, WHITE, False),
        ("  • 중국 사상사에서 가장 혁신적인 합리주의적 자연관", 17, LIGHT_GRAY, False),
    ])

    # 7. 핵심 사상 ③ 수양론과 정명론
    make_section_slide(prs, "핵심 사상 ③ — 수양론과 정명론", [
        ("적(積)의 사상 — 축적의 힘", 20, GOLD, True),
        ("  積土成山 風雨興焉 — 흙을 쌓아 산을 이루면 바람과 비가 일어난다", 17, WHITE, False),
        ("  • 한 번의 깨달음이 아니라 끊임없는 노력의 누적이 인격을 형성", 17, WHITE, False),
        ("  • 不積蹞步 無以至千里 — 반 걸음을 쌓지 않으면 천 리에 이를 수 없다", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("인식론 — 허일이정(虛壹而靜)", 20, GOLD, True),
        ("  • 허(虛): 마음을 비움 — 기존 지식에 얽매이지 않는 개방성", 17, WHITE, False),
        ("  • 일(壹): 집중 — 분산되지 않는 전일적 주의력", 17, WHITE, False),
        ("  • 정(靜): 고요함 — 감정에 흔들리지 않는 안정", 17, WHITE, False),
        ("  • 대청명(大清明)의 경지 → 사물을 있는 그대로 파악", 17, LIGHT_GRAY, False),
        ("", 6, WHITE, False),
        ("정명론(正名論) — 이름과 실재의 일치", 20, GOLD, True),
        ("  • 名無固宜 約之以命 — 이름에는 고정된 적합함이 없으며 약속으로 정해진다", 17, WHITE, False),
        ("  • 약정속성(約定俗成) — 현대 언어 철학의 규약주의와 유사", 17, WHITE, False),
    ])

    # 8. 명구절 (상)
    make_quote_slide(prs, "명구절 (상)", [
        ("學不可以已", "학불가이이",
         "배움은 그만둘 수 없다", "권학 제1편 첫 문장"),
        ("靑取之於藍而靑於藍", "청취지어람이청어람",
         "푸른색은 쪽에서 취하였으나 쪽보다 푸르다 (청출어람)", "권학 제1편"),
        ("積土成山 風雨興焉", "적토성산 풍우흥언",
         "흙을 쌓아 산을 이루면 바람과 비가 일어난다", "권학 제1편"),
        ("天行有常 不爲堯存 不爲桀亡", "천행유상 불위요존 불위걸망",
         "하늘의 운행에는 법칙이 있으니 요·걸과 무관하다", "천론 제17편"),
        ("君者舟也 庶人者水也 水則載舟 水則覆舟", "군자주야 서인자수야 수즉재주 수즉복주",
         "군주는 배, 백성은 물 — 물은 배를 띄우기도 뒤집기도 한다", "왕제 제9편"),
    ])

    # 9. 명구절 (하)
    make_quote_slide(prs, "명구절 (하)", [
        ("人之性惡 其善者僞也", "인지성악 기선자위야",
         "사람의 본성은 악하며, 선한 것은 인위(僞)이다", "성악 제23편"),
        ("禮起於何也... 求而無度量分界 則不能不爭", "예기어하야... 구이무도량분계 즉불능부쟁",
         "예는 욕망 조절의 필요에서 일어났다", "예론 제19편"),
        ("虛壹而靜 謂之大清明", "허일이정 위지대청명",
         "비우고 집중하며 고요히 하면 대청명에 이른다", "해폐 제21편"),
        ("塗之人可以爲禹", "도지인가이위우",
         "길가의 보통 사람도 우 임금이 될 수 있다", "성악 제23편"),
        ("木受繩則直 金就礪則利", "목수승즉직 금취려즉리",
         "나무는 먹줄을 받으면 곧아지고 쇠는 숫돌에 갈면 날카로워진다", "권학 제1편"),
    ])

    # 10. 구조적 특징
    make_section_slide(prs, "순자의 구조적 특징", [
        ("유가 최초의 체계적 논설문", 20, GOLD, True),
        ("  • 어록·대화 중심의 논어·맹자와 달리 비유→논증→결론 구조", 17, WHITE, False),
        ("  • 중국 산문 발전사에서 중요한 전환점", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("비유의 풍부함", 20, GOLD, True),
        ("  • 적토성산 — 축적의 힘", 17, WHITE, False),
        ("  • 청출어람 — 학문의 발전", 17, WHITE, False),
        ("  • 봉생마중 — 환경의 영향", 17, WHITE, False),
        ("  • 재주복주 — 민심의 결정적 중요성", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("32편의 논리 구조", 20, GOLD, True),
        ("  학문론(1편) → 수양론(2~4편) → 비판론(5~6편) → 유학론(7~8편)", 17, WHITE, False),
        ("  → 정치론(9~16편) → 철학론(17~23편) → 부록(24~32편)", 17, WHITE, False),
        ("  • 학문에서 시작 → 수양·비판 → 실천·정치 → 철학적 심화", 17, LIGHT_GRAY, False),
    ])

    # 11. 현대적 의의
    make_section_slide(prs, "현대적 의의", [
        ("성악설 → 현실주의적 인간관", 20, GOLD, True),
        ("  경제학·정치학의 합리적 행위자 전제와 유사", 18, WHITE, False),
        ("", 4, WHITE, False),
        ("예론(제도 설계) → 제도주의(institutionalism)", 20, GOLD, True),
        ("  개인의 덕성보다 제도적 장치의 중요성 강조", 18, WHITE, False),
        ("", 4, WHITE, False),
        ("천론(합리주의) → 과학적 세계관", 20, GOLD, True),
        ("  미신 비판, 자연에 대한 능동적·실용적 태도", 18, WHITE, False),
        ("", 4, WHITE, False),
        ("해폐(인지 편향) → 인지심리학·비판적 사고", 20, GOLD, True),
        ("  한쪽에 가려져 큰 이치에 어두운 것을 경계", 18, WHITE, False),
        ("", 4, WHITE, False),
        ("적(積)의 사상 → 성장 마인드셋·습관 형성", 20, GOLD, True),
        ("  작은 축적이 위대한 성취로 이어지는 점진적 발전론", 18, WHITE, False),
        ("", 4, WHITE, False),
        ("정명론(약정속성) → 언어 철학·기호학", 20, GOLD, True),
        ("  이름의 사회적 약정 성격에 대한 선구적 통찰", 18, WHITE, False),
    ])

    # 12. 비교
    make_table_slide(prs, "순자와 다른 사상가의 비교",
        ["항목", "공자(논어)", "맹자", "순자"],
        [
            ["인성론", "성상근(性相近)", "성선설", "성악설"],
            ["핵심 덕목", "인(仁)", "인의(仁義)", "예(禮)"],
            ["수양 방법", "학·극기복례", "확충(사단)", "화성기위(적)"],
            ["천관(天觀)", "도덕적 천", "의지적 천", "자연적 천"],
            ["정치론", "덕치·정명", "왕도·역성혁명", "왕도·예법 병용"],
            ["문체", "어록", "대화·논변", "논설문"],
            ["후대 영향", "유가 정통", "성리학", "법가·예학"],
        ],
        col_widths=[2.0, 3.2, 3.2, 3.1]
    )

    # 13. 마무리
    make_closing_slide(prs, "순자를 한 문장으로",
        "인간의 본성이 악하다고 진단하면서도,\n"
        "학문과 수양을 통해 누구나 성인이 될 수 있다고 확신한 사상가.\n"
        "하늘을 자연 법칙으로, 예를 문명의 근본으로, 이름을 사회적 약정으로 —\n"
        "합리주의와 제도주의의 선구자.",
        "「塗之人可以爲禹」 — 길가의 보통 사람도 우 임금이 될 수 있다")

    path = "Works/순자/순자_발표자료.pptx"
    prs.save(path)
    return path


# ════════════════════════════════════════════════
# 손자병법 PPT
# ════════════════════════════════════════════════
def create_sunzi_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. 표지
    make_title_slide(prs, "손자병법", "孫子兵法 (The Art of War)", "손무(孫武)",
                     "춘추시대 말기 (BC 6~5세기)",
                     "13편 약 6,000자 · 동양 최고(最古)의 군사 전략서")

    # 2. 목차
    make_toc_slide(prs, [
        "개요 — 손자병법이란 무엇인가",
        "구성 — 13편의 구조",
        "핵심 사상 ① 부전승과 오사칠계",
        "핵심 사상 ② 허실·기정·세",
        "핵심 사상 ③ 지피지기와 변화의 철학",
        "명구절 (상)",
        "명구절 (하)",
        "구조적 특징",
        "현대적 의의",
        "다른 사상과의 비교",
        "마무리"
    ])

    # 3. 개요
    make_section_slide(prs, "손자병법이란 무엇인가", [
        ("정의: 동양에서 가장 오래되고 가장 영향력 있는 군사 전략서", 20, WHITE, True),
        ("  • 13편, 약 6,000자의 간결한 철학적 서술", 18, WHITE, False),
        ("  • 전쟁 기술서를 넘어 전략적 사고의 교과서", 18, WHITE, False),
        ("", 8, WHITE, False),
        ("손무(孫武)", 20, GOLD, True),
        ("  • 제(齊)나라 출신, 오(吳)나라 장군", 18, WHITE, False),
        ("  • 오왕 합려(闔閭)에게 병법 13편을 올려 장군 임명", 18, WHITE, False),
        ("  • 초(楚)나라 수도 영(郢) 함락 등 혁혁한 전공", 18, WHITE, False),
        ("", 8, WHITE, False),
        ("판본의 역사", 20, GOLD, True),
        ("  • 조조(曹操): 최초의 주석서 편찬", 18, WHITE, False),
        ("  • 1972년 은작산(銀雀山) 한묘 죽간본 출토 → 연구의 전환점", 18, WHITE, False),
        ("  • 18세기 유럽 전파 → 20세기 서구 군사학교 필독서", 18, WHITE, False),
    ])

    # 4. 구성
    make_table_slide(prs, "손자병법의 구성 — 13편",
        ["구분", "편", "편명", "핵심 주제"],
        [
            ["전략론", "1", "시계(始計)", "전쟁의 기본 계산, 오사칠계"],
            ["전략론", "2", "작전(作戰)", "전쟁 비용과 속전속결"],
            ["전략론", "3", "모공(謀攻)", "부전승, 지피지기"],
            ["전략론", "4~5", "군형·병세", "선승이후구전, 기정, 세(勢)"],
            ["전략론", "6", "허실(虛實)", "주도권 장악, 물의 비유"],
            ["전술론", "7~8", "군쟁·구변", "풍림화산, 변통, 장수의 위험"],
            ["전술론", "9~11", "행군~구지", "지형 활용, 사지즉전"],
            ["특수전", "12~13", "화공·용간", "화공, 정보전, 다섯 간첩"],
        ],
        col_widths=[1.8, 0.8, 2.8, 6.1]
    )

    # 5. 핵심 사상 ① 부전승과 오사칠계
    make_section_slide(prs, "핵심 사상 ① — 부전승(不戰勝)과 오사칠계", [
        ("百戰百勝 非善之善者也 不戰而屈人之兵 善之善者也", 17, GOLD, True),
        ("→ 백전백승은 최선이 아니다. 싸우지 않고 적을 굴복시키는 것이 최선이다", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("승리의 등급", 20, GOLD, True),
        ("  • 최상: 벌모(伐謀) — 적의 전략을 무력화", 17, WHITE, False),
        ("  • 차선: 벌교(伐交) — 적의 동맹을 와해", 17, WHITE, False),
        ("  • 그 다음: 벌병(伐兵) — 적의 군대를 공격", 17, WHITE, False),
        ("  • 최하: 공성(攻城) — 적의 성을 공격", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("오사(五事) — 전쟁의 다섯 근본 요소", 20, GOLD, True),
        ("  • 도(道): 군주와 백성의 일체감, 대의명분", 17, WHITE, False),
        ("  • 천(天): 기후, 계절, 타이밍", 17, WHITE, False),
        ("  • 지(地): 지형, 거리, 경쟁 환경", 17, WHITE, False),
        ("  • 장(將): 지·신·인·용·엄 — 리더십의 자질", 17, WHITE, False),
        ("  • 법(法): 편제, 규율, 보급 — 조직 시스템", 17, WHITE, False),
    ])

    # 6. 핵심 사상 ② 허실·기정·세
    make_section_slide(prs, "핵심 사상 ② — 허실(虛實)·기정(奇正)·세(勢)", [
        ("허실(虛實) — 적의 빈곳을 치고, 내 빈곳을 감춰라", 20, GOLD, True),
        ("  • 避實而擊虛 — 실한 곳을 피하고 허한 곳을 치라", 17, WHITE, False),
        ("  • 致人而不致於人 — 적을 끌어오지, 끌려가지 않는다", 17, WHITE, False),
        ("  • 핵심은 주도권(initiative)의 장악", 17, LIGHT_GRAY, False),
        ("", 6, WHITE, False),
        ("기정(奇正) — 정공법과 기습의 조화", 20, GOLD, True),
        ("  • 以正合 以奇勝 — 정으로 맞서고, 기로 이긴다", 17, WHITE, False),
        ("  • 기정상생(奇正相生) — 끝없는 순환, 무한한 변화", 17, WHITE, False),
        ("  • 정(正): 정규적, 적을 묶어두는 주력", 17, WHITE, False),
        ("  • 기(奇): 비정규적, 결정적 승리의 별동대", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("세(勢) — 형세를 만들어 압도하라", 20, GOLD, True),
        ("  • 激水之疾 至於漂石者 勢也 — 급류가 돌을 떠내려 보내는 것이 세(勢)", 17, WHITE, False),
        ("  • 求之於勢 不責於人 — 세에서 승리를 구하지, 개인에게 책임 돌리지 않는다", 17, WHITE, False),
        ("  • 좋은 리더는 시스템을 설계하여 평범한 사람도 탁월한 성과를 내게 한다", 17, LIGHT_GRAY, False),
    ])

    # 7. 핵심 사상 ③ 지피지기와 변화
    make_section_slide(prs, "핵심 사상 ③ — 지피지기와 변화의 철학", [
        ("知彼知己 百戰不殆 — 적을 알고 나를 알면 백 번 싸워도 위태롭지 않다", 18, GOLD, True),
        ("", 4, WHITE, False),
        ("지(知)의 세 수준", 20, GOLD, True),
        ("  • 지피지기(知彼知己) → 백전불태 — 백 번 싸워도 위태롭지 않다", 17, WHITE, False),
        ("  • 부지피이지기 → 일승일부 — 한 번 이기고 한 번 진다", 17, WHITE, False),
        ("  • 부지피부지기 → 매전필태 — 매번 위태롭다", 17, WHITE, False),
        ("  ※ '백전백승'이 아니라 '백전불태' — 화려한 승리보다 안전을 추구", 17, LIGHT_GRAY, False),
        ("", 6, WHITE, False),
        ("물의 비유 — 유연한 변화의 철학", 20, GOLD, True),
        ("  兵形象水 — 군대의 형태는 물과 같다", 17, WHITE, False),
        ("  • 높은 곳을 피하고 낮은 곳으로 → 강한 곳 피하고 약한 곳 공격", 17, WHITE, False),
        ("  • 지형에 따라 형태가 변한다 → 상황에 따라 전략을 바꾼다", 17, WHITE, False),
        ("  • 兵無常勢 水無常形 — 일정한 형태가 없는 것이 최고", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("선승이후구전(先勝而後求戰) — 준비의 철학", 20, GOLD, True),
        ("  이기는 군대는 먼저 이겨놓고 싸우고, 지는 군대는 먼저 싸운 뒤 이기려 한다", 17, WHITE, False),
    ])

    # 8. 명구절 (상)
    make_quote_slide(prs, "명구절 (상)", [
        ("兵者 國之大事 死生之地 存亡之道 不可不察也", "병자 국지대사 사생지지 존망지도 불가불찰야",
         "전쟁은 나라의 중대한 일이니 살피지 않을 수 없다", "시계 제1편"),
        ("百戰百勝 非善之善者也 不戰而屈人之兵 善之善者也",
         "백전백승 비선지선자야 부전이굴인지병 선지선자야",
         "백전백승은 최선이 아니다. 싸우지 않고 이기는 것이 최선이다", "모공 제3편"),
        ("知彼知己 百戰不殆", "지피지기 백전불태",
         "적을 알고 나를 알면 백 번 싸워도 위태롭지 않다", "모공 제3편"),
        ("勝兵先勝而後求戰", "승병선승이후구전",
         "이기는 군대는 먼저 이겨놓고 싸운다", "군형 제4편"),
        ("凡戰者 以正合 以奇勝", "범전자 이정합 이기승",
         "전투는 정(正)으로 맞서고 기(奇)로 이긴다", "병세 제5편"),
    ])

    # 9. 명구절 (하)
    make_quote_slide(prs, "명구절 (하)", [
        ("故善戰者 致人而不致於人", "고선전자 치인이부치어인",
         "전쟁을 잘하는 자는 적을 끌어오지, 끌려가지 않는다", "허실 제6편"),
        ("其疾如風 其徐如林 侵掠如火 不動如山",
         "기질여풍 기서여림 침략여화 부동여산",
         "바람처럼 빠르고, 숲처럼 고요하며, 불처럼 공격하고, 산처럼 움직이지 않는다", "군쟁 제7편"),
        ("投之亡地然後存 陷之死地然後生", "투지망지연후존 함지사지연후생",
         "망할 곳에 던져야 살고, 죽을 곳에 빠뜨려야 산다", "구지 제11편"),
        ("主不可以怒而興師 將不可以慍而致戰",
         "주불가이노이흥사 장불가이온이치전",
         "분노로 군대를 일으켜서는 안 되고, 격분하여 싸워서는 안 된다", "화공 제12편"),
        ("明君賢將 所以動而勝人者 先知也", "명군현장 소이동이승인자 선지야",
         "명군과 현장이 움직이면 이기는 까닭은 미리 아는 것(先知) 때문이다", "용간 제13편"),
    ])

    # 10. 구조적 특징
    make_section_slide(prs, "손자병법의 구조적 특징", [
        ("13편의 논리 구조", 20, GOLD, True),
        ("  • 전략론(1~6편) → 전술론(7~11편) → 특수전법(12~13편)", 17, WHITE, False),
        ("  • 총론 → 원칙 → 실전 → 정보의 체계적 전개", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("간결함의 미학", 20, GOLD, True),
        ("  • 총 약 6,000자 — 극도로 압축된 문장", 17, WHITE, False),
        ("  • 한 문장이 하나의 원칙을 담고 있어 독립적으로 인용 가능", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("비유의 힘", 20, GOLD, True),
        ("  • 물(水) — 유연한 변화와 적응의 상징", 17, WHITE, False),
        ("  • 급류(激水) — 세(勢)의 위력", 17, WHITE, False),
        ("  • 상산의 뱀(率然) — 유기적 조직의 이상", 17, WHITE, False),
        ("  • 풍림화산 — 상황별 대응의 모범", 17, WHITE, False),
        ("", 6, WHITE, False),
        ("이중 구조: 전쟁론 + 반전론", 20, GOLD, True),
        ("  • '싸우지 않고 이기라'는 부전승 + '분노로 싸우지 마라'는 경고", 17, WHITE, False),
        ("  • 가장 뛰어난 병서이면서 동시에 가장 신중한 반전 메시지", 17, LIGHT_GRAY, False),
    ])

    # 11. 현대적 의의
    make_section_slide(prs, "현대적 의의", [
        ("경영 전략", 20, GOLD, True),
        ("  • 부전승 → 블루오션 전략, 경쟁 회피", 18, WHITE, False),
        ("  • 속전속결 → MVP, 린 스타트업", 18, WHITE, False),
        ("  • 지피지기 → 시장 조사, SWOT 분석", 18, WHITE, False),
        ("  • 허실 → 니치 마켓 공략, 선택과 집중", 18, WHITE, False),
        ("", 4, WHITE, False),
        ("리더십", 20, GOLD, True),
        ("  • 세(勢)에서 승리를 구하라 → 시스템 설계의 중요성", 18, WHITE, False),
        ("  • 장수의 다섯 덕목(지신인용엄) → 균형 잡힌 리더십", 18, WHITE, False),
        ("", 4, WHITE, False),
        ("일상의 지혜", 20, GOLD, True),
        ("  • 선승이후구전 → 준비 없이 도전하지 마라", 18, WHITE, False),
        ("  • 노이흥사 금지 → 감정에 의한 중요 결정을 피하라", 18, WHITE, False),
        ("  • 물의 비유 → 고정된 패턴에 얽매이지 마라", 18, WHITE, False),
    ])

    # 12. 비교
    make_table_slide(prs, "손자병법과 다른 사상의 비교",
        ["비교 대상", "공통점", "차이점"],
        [
            ["클라우제비츠(전쟁론)", "전쟁을 정치의 연장으로 봄", "손자: 부전승 / 클라우제비츠: 결전 중시"],
            ["마키아벨리(군주론)", "현실주의적 권력관", "손자: 도(道) 중시 / 마키아벨리: 권모술수"],
            ["노자(도덕경)", "유연함과 물의 비유", "노자: 무위 / 손자: 적극적 전략"],
            ["공자(논어)", "인(仁)의 중시", "공자: 도덕 우선 / 손자: 전략적 효용"],
            ["순자", "현실주의적 인간관", "순자: 제도 설계 / 손자: 전장 설계"],
        ],
        col_widths=[3.0, 3.5, 5.0]
    )

    # 13. 마무리
    make_closing_slide(prs, "손자병법을 한 문장으로",
        "전쟁의 가부를 신중히 판단하고(오사칠계),\n"
        "가능하면 싸우지 않고 이기며(부전승),\n"
        "반드시 싸워야 한다면 먼저 이길 조건을 갖추고(선승이후구전),\n"
        "이 모든 것의 기반은 정보(지피지기)와 유연한 변화(병무상세)이다.",
        "「不戰而屈人之兵 善之善者也」 — 싸우지 않고 이기는 것이 최선이다")

    path = "Works/손자병법/손자병법_발표자료.pptx"
    prs.save(path)
    return path


# ════════════════════════════════════════════════
# 메인 실행
# ════════════════════════════════════════════════
if __name__ == "__main__":
    os.chdir(os.path.dirname(os.path.abspath(__file__)))

    print("=" * 60)
    print("  4대 고전 PPT 생성 시작")
    print("=" * 60)

    results = []

    print("\n[1/4] 논어 PPT 생성 중...")
    p = create_lunyu_ppt()
    results.append(p)
    print(f"  [OK] {p}")

    print("\n[2/4] 맹자 PPT 생성 중...")
    p = create_mengzi_ppt()
    results.append(p)
    print(f"  [OK] {p}")

    print("\n[3/4] 순자 PPT 생성 중...")
    p = create_xunzi_ppt()
    results.append(p)
    print(f"  [OK] {p}")

    print("\n[4/4] 손자병법 PPT 생성 중...")
    p = create_sunzi_ppt()
    results.append(p)
    print(f"  [OK] {p}")

    print("\n" + "=" * 60)
    print("  생성 완료! 검증 중...")
    print("=" * 60)

    for path in results:
        if os.path.exists(path):
            size_kb = os.path.getsize(path) / 1024
            prs = Presentation(path)
            n_slides = len(prs.slides)
            print(f"  {path}")
            print(f"    슬라이드: {n_slides}장 | 파일 크기: {size_kb:.1f} KB")
        else:
            print(f"  [FAIL] {path} -- not found!")

    print("\n" + "=" * 60)
    print("  모든 작업 완료")
    print("=" * 60)
