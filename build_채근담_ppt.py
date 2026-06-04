# -*- coding: utf-8 -*-
"""
채근담(菜根譚) 발표자료 — 망라적 PPT
원칙: 흰 배경 · 한 슬라이드 한 주제 · 매수 무제약
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK    = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
SUB    = RGBColor(0x55, 0x60, 0x70)
RULE   = RGBColor(0xC8, 0xA2, 0x5B)
PALE   = RGBColor(0xF5, 0xEE, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *,
                   font_size=18, color=INK, font_name='맑은 고딕',
                   line_spacing=1.35, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font_name', font_name)
        run.font.size = Pt(opts.get('font_size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
    return tb


def add_rule(slide, left, top, width, color=RULE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_page_header(slide, section, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.3),
                section, font_size=11, color=SUB, bold=True)
    add_textbox(slide, Inches(12.0), Inches(0.3), Inches(1.3), Inches(0.3),
                f'{page_num} / {total}', font_size=10, color=SUB,
                align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(0.5), Inches(0.65), Inches(12.8))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.7),
                title, font_size=30, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.8), Inches(0.4),
                    subtitle, font_size=14, color=SUB)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDES = []
def S(section):
    def deco(fn):
        SLIDES.append((fn, section))
        return fn
    return deco


# ============== 표지 ==============
@S('표지')
def slide_cover(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
                '나물 뿌리를 씹는 자라야 백 가지 일을 이룬다 · 동양 최고의 인생 지침서',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.6),
                '菜 根 譚',
                font_size=110, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
                '채 근 담',
                font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(4.5), Inches(5.1), Inches(4.3))
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                '眞 味 只 是 淡 — 참된 맛은 담백할 뿐이다',
                font_size=18, color=SUB, align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
                '홍자성(洪自誠, 환초도인) · 명말 만력 연간 · 전집 225조 + 후집 135조 = 360조',
                font_size=16, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                '유(儒)·불(佛)·도(道) 삼교 융합 — 입세(入世)와 출세(出世)의 지혜',
                font_size=14, color=INK, align=PP_ALIGN.CENTER)


# ============== 목차 ==============
@S('목차')
def slide_toc(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, '목차', n, t)
    add_title(slide, '목 차', 'Contents')
    items_left = [
        ('Ⅰ.', '개요 — 채근담은 어떤 책인가'),
        ('Ⅱ.', '저자 홍자성과 명 말 시대'),
        ('Ⅲ.', '삼교합일 — 유·불·도의 융합'),
        ('Ⅳ.', '전집·후집 구조 — 입세와 출세'),
        ('Ⅴ.', '전집 5편 깊이 읽기'),
        ('Ⅵ.', '후집 5편 깊이 읽기'),
    ]
    items_right = [
        ('Ⅶ.', '책을 관통하는 10대 주제'),
        ('Ⅷ.', '명구 16선'),
        ('Ⅸ.', '동아시아의 채근담 — 일본과 한국'),
        ('Ⅹ.', '오늘 다시 펼치는 이유'),
        ('Ⅺ.', '마무리'),
    ]
    for i, (num, txt) in enumerate(items_left):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(0.9), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.6), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)
    for i, (num, txt) in enumerate(items_right):
        y = Inches(2.4 + i * 0.55)
        add_textbox(slide, Inches(7.2), y, Inches(0.7), Inches(0.5),
                    num, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.9), y, Inches(5.2), Inches(0.5),
                    txt, font_size=17, color=INK)


# ============== Ⅰ. 개요 ==============
SEC1 = 'Ⅰ. 개요'

@S(SEC1)
def i_name(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '제목 「채근(菜根)」의 뜻',
              '거칠고 질긴 나물 뿌리를 씹는 자라야 백 가지 일을 이룬다')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5),
                '人 常 咬 得 菜 根  則 百 事 可 做',
                font_size=26, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.4),
                '사람이 늘 나물 뿌리를 씹을 수 있다면, 모든 일을 이룰 수 있다 — 송 왕신민(汪信民)의 말',
                font_size=14, color=SUB, align=PP_ALIGN.CENTER)
    add_paragraphs(slide, Inches(0.7), Inches(4.4), Inches(12.0), Inches(3.0), [
        ('● 채근(菜根) — 거칠고 질긴 나물 뿌리 · 기름지지 않으나 씹을수록 참맛',
         {'font_size': 17, 'space_before': 8}),
        ('● 상징 — 인생의 쓴맛, 담박한 삶, 견뎌냄의 미덕',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 「나물 뿌리를 씹듯 인생의 고단함을 견딜 수 있는 자라야 큰일을 이룬다」',
         {'font_size': 17, 'space_before': 10}),
        ('● 책의 근본 이미지 — 진하고 기름진 삶이 아니라 「담박(淡泊)한 삶」',
         {'font_size': 17, 'space_before': 10, 'color': SUB, 'bold': True}),
    ])


@S(SEC1)
def i_overview(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '한눈에 보는 채근담')
    rows = [
        ('저자',     '홍자성(洪自誠) — 호 환초도인(還初道人), 자 자성'),
        ('시대',     '명(明) 말 만력(萬曆) 연간 (16세기 말~17세기 초)'),
        ('분량',     '전집 225조 + 후집 135조 = 총 360조'),
        ('형식',     '짧은 격언·경구(警句) · 대구(對句)와 리듬을 갖춘 변려풍 산문'),
        ('사상',     '유(儒)·불(佛)·도(道) 삼교 융합'),
        ('성격',     '처세·수양·은일(隱逸)의 지혜 — 일상의 한 줄 명상록'),
        ('일본 위상',  '에도 시대 이래 가장 많이 읽힌 중국 수양서 — 「일본의 논어」'),
        ('한국 수용',  '조선 후기부터 사대부·선비의 애독서 — 자기 계발서의 원형'),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.3 + i * 0.6)
        add_filled_rect(slide, Inches(0.7), y, Inches(2.6), Inches(0.5), PALE)
        add_textbox(slide, Inches(0.7), y, Inches(2.6), Inches(0.5),
                    k, font_size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.5), y, Inches(9.5), Inches(0.5),
                    v, font_size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC1)
def i_bestseller(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC1, n, t)
    add_title(slide, '왜 동아시아의 베스트셀러가 되었는가',
              '세 가지 이유 — 융합·생활·격언 형식')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 1. 유불도 삼교의 지혜가 한 책에',
         {'font_size': 19, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('     어느 종교·사상 배경의 독자도 받아들일 수 있는 보편성',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 2. 생활 밀착형',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
        ('     관념이 아니라 매일매일의 처세와 마음을 다룬다',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 3. 짧은 격언 형식',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
        ('     어디서 펼쳐도 한 조(條)가 독립적 가르침 — 명상록(Meditations)적 독법이 가능',
         {'font_size': 15, 'color': SUB, 'space_before': 4}),
        ('● 360개의 독립된 「진주알」이 실에 꿰여 있는 형태 — 어디서부터 읽어도 되는 책',
         {'font_size': 16, 'space_before': 14, 'color': ACCENT, 'bold': True}),
    ])


# ============== Ⅱ. 저자와 시대 ==============
SEC2 = 'Ⅱ. 저자와 시대'

@S(SEC2)
def ii_author(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '홍자성(洪自誠) — 은둔한 지식인',
              '벼슬에서 물러나 산림에 거처한 한 사대부')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 홍응명(洪應明) — 자 자성(自誠), 호 환초도인(還初道人)',
         {'font_size': 18, 'space_before': 4}),
        ('● 명 말 만력 연간(1573~1620)에 살았던 인물 — 생애는 상세히 전하지 않음',
         {'font_size': 17, 'space_before': 10}),
        ('● 호 「환초(還初)」 — 「처음으로 돌아간다(復歸於初)」 · 노장(老莊)적 지향',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
        ('● 우익(于孔兼)의 서문 — 「세상을 초탈해 산림에 거처한 인물」로 소개',
         {'font_size': 17, 'space_before': 12}),
        ('● 채근담은 그의 개인적 좌절과 관조의 산물',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 명 말의 혼란 속 한 사대부가 자기가 찾은 삶의 처방을 360조에 남긴 것',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC2)
def ii_era(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC2, n, t)
    add_title(slide, '명 말(明末)이라는 시대 — 채근담을 낳은 토양',
              '정치의 혼란 · 양명학의 폭발 · 선불교의 부흥')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 정치의 혼란 — 만력제의 30년 파업, 환관 전횡, 동림당 사화',
         {'font_size': 17, 'space_before': 4, 'bold': True, 'color': ACCENT}),
        ('● 사상의 지각 변동 — 양명학의 「심학(心學)」 시대',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('     「心卽理」 — 내면·양지(良知)로의 관심이 폭증',
         {'font_size': 14, 'color': SUB, 'space_before': 4, 'font_name': 'Batang'}),
        ('● 불교·도교의 부흥 — 선불교(禪佛敎)·전진교(全眞敎)의 수행이 사대부 사이에 광범위',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 출판 문화의 폭발 — 통속 격언집·처세서·선서(善書)가 대량 유통',
         {'font_size': 17, 'space_before': 10}),
        ('● 채근담은 이 시대의 결정적 산물 — 양명학·선·노장·유교가 한 사대부 안에서 융합',
         {'font_size': 16, 'space_before': 14, 'color': ACCENT, 'bold': True}),
    ])


# ============== Ⅲ. 삼교합일 ==============
SEC3 = 'Ⅲ. 삼교합일'

@S(SEC3)
def iii_three(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '삼교합일(三敎合一) — 채근담의 뼈대',
              '유·불·도가 한 책 안에서 자연스럽게 녹아들다')
    cols = [
        ('儒 유', '입세(入世)의 뼈대',
         '수신제가(修身齊家)\n중용(中庸)\n인의예지(仁義禮智)\n\n德者才之主\n— 덕이 재주의 주인\n\n주로 전집 225조의 영역',
         ACCENT),
        ('道 도', '담박(淡泊)의 색조',
         '무위자연(無爲自然)\n지족상락(知足常樂)\n공성신퇴(功成身退)\n\n眞味只是淡\n— 참맛은 담백할 뿐\n\n전·후집 전체에 흐름',
         INK),
        ('佛 불', '해탈의 심층',
         '무상(無常) · 공(空)\n방하착(放下着)\n관심(觀心)\n\n寵辱不驚\n— 영예·치욕에 놀라지 않음\n\n후집 135조에 두드러짐',
         SUB),
    ]
    for i, (han, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 4.2)
        add_filled_rect(slide, x, Inches(2.3), Inches(3.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(3.9), Inches(0.6),
                    han, font_size=28, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.85), Inches(3.9), Inches(0.4),
                    label, font_size=13, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.3), Inches(3.6), Inches(3.5), Inches(3.5),
                       [(body, {'font_size': 13, 'color': INK})], line_spacing=1.5)


@S(SEC3)
def iii_one(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC3, n, t)
    add_title(slide, '한 조(條) 안의 삼교 — 채근담의 작법',
              '세 사상이 한 문장 안에서 자연스럽게 겹친다')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.8), PALE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.7),
                '靜 中 靜 非 眞 靜',
                font_size=28, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(12.0), Inches(0.7),
                '動 處 靜 得 來  纔 是 性 天 之 眞 境',
                font_size=22, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
    add_paragraphs(slide, Inches(0.7), Inches(4.2), Inches(12.0), Inches(3.0), [
        ('● 「고요한 가운데의 고요는 참된 고요가 아니다. 움직임 속에서 고요를 얻어야 본성의 참 경지」',
         {'font_size': 15, 'space_before': 4, 'color': SUB}),
        ('● 유(儒) — 현실(動)을 떠나지 않고 지켜야 한다는 입세의 태도',
         {'font_size': 17, 'space_before': 12}),
        ('● 도(道) — 고요(靜)를 본성으로 삼는 무위',
         {'font_size': 17, 'space_before': 10}),
        ('● 불(佛) — 움직이는 현상 속에서 본성(性天)을 본다는 선적(禪的) 통찰',
         {'font_size': 17, 'space_before': 10}),
        ('● 이 한 문장에 세 사상이 녹아 있다 — 이것이 채근담의 방식',
         {'font_size': 16, 'space_before': 12, 'bold': True, 'color': ACCENT}),
    ])


# ============== Ⅳ. 전후집 구조 ==============
SEC4 = 'Ⅳ. 전후집 구조'

@S(SEC4)
def iv_two(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '전집 vs 후집 — 입세(入世)와 출세(出世)',
              '세상 속에서 살아가는 법 + 세상을 한 걸음 떨어져 바라보는 여유')
    cols = [
        ('前 集', '전집 225조',
         '入 世 — 입세의 지혜\n\n· 세상 속에서 살아가는\n  처세와 수양\n\n· 인간관계 · 도덕\n  자기 수양 · 권력과 부\n\n· 현실적 삶의 지침',
         ACCENT),
        ('後 集', '후집 135조',
         '出 世 — 출세의 지혜\n\n· 세상을 초월하는\n  한적한 삶\n\n· 자연 · 마음의 평정\n  인생의 본질 · 무욕\n\n· 정신적 삶의 지침',
         INK),
    ]
    for i, (han, label, body, color) in enumerate(cols):
        x = Inches(0.7 + i * 6.3)
        add_filled_rect(slide, x, Inches(2.3), Inches(5.9), Inches(1.0), color)
        add_textbox(slide, x, Inches(2.3), Inches(5.9), Inches(0.6),
                    han, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x, Inches(2.95), Inches(5.9), Inches(0.4),
                    label, font_size=14, color=RULE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraphs(slide, x + Inches(0.5), Inches(3.6), Inches(4.9), Inches(3.5),
                       [(body, {'font_size': 14, 'color': INK})], line_spacing=1.5)


@S(SEC4)
def iv_ten(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC4, n, t)
    add_title(slide, '10편의 주제 분류 (본 자료 재편집)',
              '전집 5편(입세) + 후집 5편(출세)')
    add_filled_rect(slide, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.5), INK)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(1.6), Inches(0.5),
                '구분', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(2.3), Inches(2.2), Inches(3.0), Inches(0.5),
                '편명', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(5.3), Inches(2.2), Inches(7.4), Inches(0.5),
                '핵심 내용', font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ('전집 1', '修養과 德行 수양과 덕행', '도덕적 자기 수양, 덕의 실천'),
        ('전집 2', '處世와 人間關係 처세와 인간관계', '세상을 살아가는 지혜, 대인관계'),
        ('전집 3', '權力과 富貴 권력과 부귀', '권세와 재물에 대한 올바른 태도'),
        ('전집 4', '治心 마음 다스리기', '욕망 절제, 내면의 평화'),
        ('전집 5', '敎育과 學問 교육과 학문', '배움의 자세, 가르침의 도리'),
        ('후집 1', '閑居 한적한 삶', '자연 속의 한가로운 삶의 기쁨'),
        ('후집 2', '自然과 風流 자연과 풍류', '산수자연과의 교감, 풍류의 멋'),
        ('후집 3', '人生 인생의 본질', '삶과 죽음, 세월의 무상함'),
        ('후집 4', '無慾 무욕과 해탈', '집착을 버리고 자유로운 경지'),
        ('후집 5', '覺悟 깨달음과 귀결', '인생의 궁극적 깨달음, 돌아감'),
    ]
    for i, (cat, name, desc) in enumerate(rows):
        y = Inches(2.7 + i * 0.42)
        bg = PALE if i % 2 == 0 else WHITE
        add_filled_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.42), bg)
        is_qian = cat.startswith('전집')
        cat_color = ACCENT if is_qian else INK
        add_textbox(slide, Inches(0.7), y, Inches(1.6), Inches(0.42),
                    cat, font_size=13, bold=True, color=cat_color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.3), y, Inches(3.0), Inches(0.42),
                    name, font_size=13, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(5.4), y, Inches(7.3), Inches(0.42),
                    desc, font_size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅴ. 전집 5편 깊이 읽기 ==============
SEC5 = 'Ⅴ. 전집 5편 깊이 읽기'

def make_chapter_slide(section, num, total, name_han, name_kor, headline,
                        original, modern, theme, point):
    @S(section)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{section} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.8), Inches(0.6),
                    f'{name_han}  ({name_kor})',
                    font_size=26, bold=True, color=INK, font_name='Batang')
        add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.8), Inches(0.4),
                    headline, font_size=14, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.5), PALE)
        add_paragraphs(slide, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.4), [
            (original, {'font_size': 17, 'bold': True, 'color': ACCENT, 'font_name': 'Batang'}),
            (modern,   {'font_size': 13, 'color': SUB, 'space_before': 6}),
        ], line_spacing=1.3)
        add_textbox(slide, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
                    '◆ 중심 주제', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
                       [(theme, {'font_size': 15, 'color': INK})], line_spacing=1.35)
        add_textbox(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.4),
                    '◆ 압축 메시지', font_size=14, bold=True, color=ACCENT)
        add_paragraphs(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.4),
                       [(point, {'font_size': 15, 'color': INK})], line_spacing=1.35)


# 전집 5편
QIAN = [
    ('修養과 德行', '전집 1편 — 수양과 덕행',
     '세상에 나서기 전에 먼저 자기를 닦는 것이 모든 것의 출발',
     '醲 肥 辛 甘 非 眞 味  眞 味 只 是 淡',
     '진하고 달콤한 것은 참맛이 아니다. 참된 맛은 담백할 뿐이다',
     '담박(淡泊)·겸손·신독(愼獨)·균형·점진적 축적·비움 — 화려함이 아니라 담백함이 본체.',
     '「지극한 사람은 평범해 보인다」 — 화려함이 아닌 담백함이, 탁월함의 과시가 아닌 꾸준한 축적이 수양의 본체.'),
    ('處世와 人間關係', '전집 2편 — 처세와 인간관계',
     '남에게 관대하고 자기에게 엄격하며, 언제나 한 걸음의 여지를 남긴다',
     '路 徑 窄 處 留 一 步 與 人 行  滋 味 濃 的 減 三 分 讓 人 嗜',
     '좁은 길에선 한 걸음 양보하고, 진한 맛은 삼분을 덜어 나누라 — 이것이 세상 사는 가장 편안한 법',
     '양보(讓)·관용·신뢰·말의 신중·여지(餘地) — 「不責人小過, 不發人陰私, 不念人舊惡」.',
     '「양보는 손해가 아니라 가장 편안한 길」 — 제로섬에서 벗어나 여지를 남길 줄 아는 자가 지속 가능한 관계를 얻는다.'),
    ('權力과 富貴', '전집 3편 — 권력과 부귀',
     '도덕에 뿌리내린 부귀만이 오래가며, 절정에서 물러설 줄 알아야 화를 면한다',
     '富 貴 名 譽 自 道 德 來 者  如 山 林 中 花  舒 徐 繁 衍',
     '도덕에서 온 부귀는 산속 꽃과 같아 느리되 오래 번성한다 / 권력으로 얻은 것은 병 속 꽃이라 곧 시든다',
     '근원(本源)·분수·지족(知足)·퇴보(退步)·불탐(不貪) — 「進步處便思退步」.',
     '「부귀의 뿌리를 보라」 — 같은 열매라도 뿌리가 다르면 수명이 다르다. 도덕에 뿌리박은 부귀만이 산속 꽃처럼 번성한다.'),
    ('治心 마음 다스리기', '전집 4편 — 마음 다스리기',
     '욕망과 분노를 초기에 다스리고, 분주함 속에서도 한가한 마음을 지킨다',
     '心 體 光 明  暗 室 中 有 靑 天   念 頭 暗 昧  白 日 下 有 厲 鬼',
     '마음이 밝으면 어두운 방에도 푸른 하늘이 있다 / 마음이 어두우면 대낮에도 무서운 귀신이 있다',
     '절제·관조(觀)·비움·신독·명상 — 「怒時暫緩須臾, 便不傷手足之情」.',
     '「세상은 내 마음이 만드는 거울」 — 외부 조건이 아니라 마음 상태가 나의 세계를 결정한다.'),
    ('敎育과 學問', '전집 5편 — 교육과 학문',
     '참된 학문은 지식 축적이 아니라 인격의 완성이며, 반드시 실천으로 연결',
     '德 者 才 之 主  才 者 德 之 奴',
     '덕이 재주의 주인이고, 재주는 덕의 종이다',
     '인격 수양·실천(踐)·겸손·질문(疑)·항심(恒心) — 「讀書不見聖賢, 爲鉛槧傭」.',
     '「지식은 도구, 인격이 목적」 — 책 1만 권을 읽고도 덕이 없으면 활자의 하인. 크게 의심하는 자만이 크게 깨닫는다.'),
]

for i, p in enumerate(QIAN, 1):
    make_chapter_slide(SEC5, i, len(QIAN), *p)


# ============== Ⅵ. 후집 5편 깊이 읽기 ==============
SEC6 = 'Ⅵ. 후집 5편 깊이 읽기'

HOU = [
    ('閑居 한적한 삶', '후집 1편 — 한적한 삶',
     '은둔이란 산속이 아니라, 세상 한복판에서도 마음이 한가로운 상태',
     '出 世 之 道  卽 在 涉 世 中   了 心 之 功  卽 在 盡 心 內',
     '세상을 벗어나는 길은 세상 속에 있고, 마음을 깨치는 공부는 마음을 다하는 안에 있다',
     '한(閑)·담박·여유·자족·관조 — 「閒中不放過, 忙處有受用」.',
     '「한가함은 장소가 아니라 상태」 — 시끄러운 저잣거리에서도 마음이 한가하면 그것이 산림이다.'),
    ('自然과 風流', '후집 2편 — 자연과 풍류',
     '자연은 말 없이 가르친다 — 꽃이 반쯤 핀 순간, 달이 차기 전의 조각에 참된 미가 있다',
     '花 看 半 開  酒 飮 微 醉   此 中 大 有 佳 趣',
     '꽃은 반쯤 피었을 때, 술은 살짝 취했을 때 — 그 안에 큰 멋이 있다',
     '반개(半開)의 미학·무흔(無痕)·여백·자연의 스승됨 — 「風來疏竹, 風過而竹不留聲」.',
     '「최고는 반쯤에 있다」 — 가득 찬 것은 기울기 직전. 반쯤의 아름다움, 흔적 없는 경지가 자연이 가르치는 풍류.'),
    ('人生 인생의 본질', '후집 3편 — 인생의 본질',
     '모든 것은 변하고 순환한다 — 유한을 자각할수록 현재가 소중해진다',
     '天 地 有 萬 古  此 身 不 再 得   人 生 只 百 年  此 日 最 易 過',
     '천지는 영원하나 이 몸은 다시 얻을 수 없다 / 인생은 백 년뿐이요 오늘은 가장 쉽게 간다',
     '무상(無常)·순환·유한·절정과 추락·현재 — 「衰颯的景象 就在盛滿中」.',
     '「영원을 붙잡으려 말고 오늘을 살아라」 — 번성 속에 쇠락의 씨앗이 있고, 쇠락 속에 재생의 기틀이 있다.'),
    ('無慾 무욕과 해탈', '후집 4편 — 무욕과 해탈',
     '욕망을 덜어낼수록 자유로워진다 — 감소(減)는 소극이 아니라 최고의 적극',
     '貪 得 者  身 富 而 心 貧   知 足 者  身 貧 而 心 富',
     '탐하는 자는 몸이 부유해도 마음이 가난하고, 만족하는 자는 몸이 가난해도 마음이 부유하다',
     '감(減)·지족·무집착·소박·비움 — 「減省一分, 便超脫一分」.',
     '「빼기가 곧 채움」 — 줄일수록 자유롭고, 비울수록 풍요롭다. 무욕은 무기력이 아니라 가장 강한 자유.'),
    ('覺悟 깨달음과 귀결', '후집 5편 — 깨달음과 귀결',
     '모든 수양의 종착지는 본래의 순수한 자리로 돌아옴 — 還初(환초)',
     '寵 辱 不 驚  閑 看 庭 前 花 開 花 落   去 留 無 意  漫 隨 天 外 雲 卷 雲 舒',
     '영예·치욕에 놀라지 않고 뜰 앞 꽃이 피고 짐을 한가히 보며 / 떠남과 머묾에 뜻이 없어 하늘 밖 구름의 말리고 펴짐을 그저 따른다',
     '본래성(本來)·졸(拙)·평정·회귀(還初)·삼세의 초월 — 「文以拙進, 道以拙成」.',
     '「배움의 끝은 배우기 전의 자리」 — 졸박(拙朴)한 처음으로 돌아오는 것이 깨달음. 저자의 호 「환초」가 그 정신.'),
]

for i, p in enumerate(HOU, 1):
    make_chapter_slide(SEC6, i, len(HOU), *p)


# ============== Ⅶ. 10대 주제 ==============
SEC7 = 'Ⅶ. 채근담의 10대 주제'

@S(SEC7)
def vii_ten(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC7, n, t)
    add_title(slide, '360조를 관통하는 10가지 주제')
    items = [
        ('1', '淡泊 담박', '참된 맛은 담백 — 「眞味只是淡」', ACCENT),
        ('2', '中庸 중용', '청렴하되 포용 · 인자하되 결단 — 양극의 균형', INK),
        ('3', '愼獨 신독', '혼자 있을 때의 나가 진짜 — 「冥冥中에서 먼저」', ACCENT),
        ('4', '退一步 퇴일보', '물러섬은 지혜 — 「進步의 반대는 좌초」', INK),
        ('5', '餘地 여지', '80%만 채우고 20%는 비움 — 「달은 차면 기운다」', ACCENT),
        ('6', '知足 지족', '만족하면 신선의 세계 — 같은 현실이 천국·지옥', INK),
        ('7', '澹心 담심', '발분(發憤)보다 담심 — 가라앉은 마음으로 사물을 본다', ACCENT),
        ('8', '觀心 관심', '밤이 깊을 때 자기 마음을 본다 — 「夜深人靜 獨坐觀心」', INK),
        ('9', '無痕 무흔', '바람·기러기 지나간 자리에 흔적 없음 — 집착 없는 마음', ACCENT),
        ('10', '還初 환초', '처음으로 돌아가기 — 모든 수양의 귀결점', INK),
    ]
    for i, (num, name, desc, color) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(2.3 + row * 0.95)
        add_filled_rect(slide, x, y, Inches(0.7), Inches(0.8), color)
        add_textbox(slide, x, y, Inches(0.7), Inches(0.8),
                    num, font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(slide, x + Inches(0.8), y, Inches(1.7), Inches(0.8), PALE)
        add_textbox(slide, x + Inches(0.8), y, Inches(1.7), Inches(0.8),
                    name, font_size=15, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, x + Inches(2.7), y + Inches(0.05), Inches(3.2), Inches(0.7),
                    desc, font_size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅷ. 명구 16선 ==============
SEC8 = 'Ⅷ. 명구 16선'

def make_quote_slide(num, total, source, han, kor, comment):
    @S(SEC8)
    def _slide(slide, n, t):
        set_white_background(slide)
        add_page_header(slide, f'{SEC8} ({num}/{total})', n, t)
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.8), Inches(0.45),
                    source, font_size=15, color=SUB)
        add_filled_rect(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.6), PALE)
        add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.8),
                    han, font_size=26, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Batang')
        add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                    kor, font_size=16, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rule(slide, Inches(0.7), Inches(4.5), Inches(12.0))
        add_paragraphs(slide, Inches(0.7), Inches(4.8), Inches(12.0), Inches(2.4),
                       [(comment, {'font_size': 16, 'color': INK})],
                       line_spacing=1.5)


QUOTES = [
    ('전집 수양', '棲 守 道 德 者 寂 寞 一 時  依 阿 權 勢 者 凄 涼 萬 古',
     '도덕을 지키면 한때 적막하지만, 권세에 아부하면 만고에 처량하다',
     '한때의 적막을 견디는 자가 영원의 영광을 얻고, 한때의 위세를 좇는 자가 영원의 처량함을 얻는다. 시야의 길이가 인격을 결정한다.'),
    ('전집 수양', '醲 肥 辛 甘 非 眞 味  眞 味 只 是 淡',
     '진하고 달콤한 것은 참맛이 아니다. 참된 맛은 담백할 뿐이다',
     '채근담 전체의 심미(審美) 원리. 화려함이 아니라 담백함이 본체 — 진한 음식이 곧 질리듯, 진한 인생도 곧 무너진다.'),
    ('전집 수양', '神 奇 卓 異 非 至 人  至 人 只 是 常',
     '신기한 것이 지인(至人)이 아니다. 지인은 오직 평범할 뿐이다',
     '진짜 큰 사람은 평범해 보인다. 「대교약졸(大巧若拙)」 — 큰 솜씨는 어리석어 보인다는 노자의 사상과 통한다.'),
    ('전집 처세', '路 徑 窄 處 留 一 步 與 人 行  滋 味 濃 的 減 三 分 讓 人 嗜',
     '좁은 길에선 한 걸음 양보하고, 진한 맛은 삼분을 덜어 나누라',
     '이것이 「세상 사는 가장 편안한 법(涉世一極安樂法)」. 제로섬에서 벗어나는 자가 가장 평화롭다.'),
    ('전집 처세', '不 責 人 小 過  不 發 人 陰 私  不 念 人 舊 惡',
     '남의 작은 허물을 꾸짖지 말고, 은밀한 사생활을 들추지 말며, 지난 잘못을 기억하지 말라',
     '관계의 세 가지 절제 — 작은 허물은 흘리고, 사생활은 묻어두고, 지난 일은 잊는다. 이 셋을 지키면 평생 적이 없다.'),
    ('전집 권력', '富 貴 名 譽 自 道 德 來 者  如 山 林 中 花  自 是 舒 徐 繁 衍',
     '도덕에서 온 부귀는 산속 꽃과 같아 느리되 오래 번성한다',
     '뿌리의 비유 — 도덕의 부귀는 산속 꽃(영원), 공업의 부귀는 화분의 꽃(이동 시 사라짐), 권력의 부귀는 병 속 꽃(곧 시듦).'),
    ('전집 권력', '進 步 處 便 思 退 步  庶 免 觸 藩 之 禍',
     '나아가는 곳에서 물러설 때를 생각하면, 화를 면한다',
     '진보의 정점에서 퇴보를 생각하라 — 노자 「공성신퇴(功成身退)」의 채근담 버전. 가장 높을 때 가장 위험하다.'),
    ('전집 마음', '心 體 光 明  暗 室 中 有 靑 天   念 頭 暗 昧  白 日 下 有 厲 鬼',
     '마음이 밝으면 어두운 방에도 푸른 하늘이 있고, 마음이 어두우면 대낮에도 무서운 귀신이 있다',
     '세상은 마음의 거울 — 외부 조건이 아니라 마음 상태가 나의 세계를 결정한다. 같은 방, 같은 시각이 천국이 되고 지옥이 된다.'),
    ('전집 학문', '德 者 才 之 主  才 者 德 之 奴',
     '덕이 재주의 주인이고, 재주는 덕의 종이다',
     '재능과 인격의 위계 — 덕 없는 재주는 흉기. 덕은 방향을, 재주는 추진력을 준다. 방향 없는 추진력은 재앙.'),
    ('후집 한거', '閒 中 不 放 過  忙 處 有 受 用',
     '한가할 때 허투루 보내지 않으면, 바쁠 때 쓸 것이 있다',
     '한가함의 정의 전환 — 빈 시간이 아니라 「축적의 시간」. 위기에 쓰일 자원은 평화 때에 모은다.'),
    ('후집 자연', '花 看 半 開  酒 飮 微 醉   此 中 大 有 佳 趣',
     '꽃은 반쯤 피었을 때, 술은 살짝 취했을 때 — 그 안에 큰 멋이 있다',
     '반개(半開)의 미학 — 가득 찬 것은 기울기 직전이다. 최고의 순간은 「조금 못 미친」 자리에 있다. 동양 절제미의 정수.'),
    ('후집 자연', '風 來 疏 竹  風 過 而 竹 不 留 聲   雁 度 寒 潭  雁 去 而 潭 不 留 影',
     '바람이 성긴 대숲에 불어도 떠난 뒤엔 소리를 남기지 않고, 기러기가 찬 못을 건너도 그림자를 남기지 않는다',
     '무흔(無痕)의 미학 — 일이 오면 드러내고 일이 가면 비운다. 「君子事來而心始現, 事去而心隨空」. 집착 없는 마음.'),
    ('후집 인생', '天 地 有 萬 古  此 身 不 再 得   人 生 只 百 年  此 日 最 易 過',
     '천지는 영원하나 이 몸은 다시 얻을 수 없다 / 인생은 백 년뿐이요 오늘은 가장 쉽게 간다',
     '유한의 자각이 현재를 살게 한다 — 「오늘은 가장 쉽게 가는 것」. 영원을 붙잡으려 말고 지금 이 자리를 살아라.'),
    ('후집 무욕', '貪 得 者  身 富 而 心 貧   知 足 者  身 貧 而 心 富',
     '탐하는 자는 몸이 부유해도 마음이 가난하고, 만족하는 자는 몸이 가난해도 마음이 부유하다',
     '부의 두 종류 — 몸의 부와 마음의 부. 둘은 종종 반대로 간다. 「知足者仙境, 不知足者凡境」.'),
    ('후집 무욕', '減 省 一 分  便 超 脫 一 分',
     '한 가지를 줄이면 한 걸음 초탈한다',
     '감(減)의 적극성 — 빼기가 곧 채움이고, 줄임이 곧 자유. 무욕은 무기력이 아니라 가장 강한 적극성.'),
    ('후집 깨달음', '寵 辱 不 驚  閑 看 庭 前 花 開 花 落',
     '영예와 치욕에 놀라지 않고, 뜰 앞의 꽃이 피고 짐을 한가히 본다',
     '채근담 마지막의 정수 — 「去留無意, 漫隨天外雲卷雲舒」와 짝. 영예·치욕·득실·거취에 흔들리지 않는 평정의 절정.'),
]

for i, q in enumerate(QUOTES, 1):
    make_quote_slide(i, len(QUOTES), *q)


# ============== Ⅸ. 동아시아의 채근담 ==============
SEC9 = 'Ⅸ. 동아시아의 채근담'

@S(SEC9)
def ix_japan(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '일본 — 「일본의 논어」',
              '에도 시대 이래 가장 많이 읽힌 중국 수양서')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 에도 시대 — 일본 사대부·승려의 필독서로 자리잡음',
         {'font_size': 18, 'space_before': 4}),
        ('● 메이지~쇼와 — 정치인·경영인·작가의 좌우명 자료',
         {'font_size': 18, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 무라카미 하루키부터 마쓰시타 고노스케까지 — 「내가 늘 가까이 두는 책」',
         {'font_size': 17, 'space_before': 12}),
        ('● 가와카미 데쓰타로(川上鐵太郞)의 1885년 일본어 역주가 결정적 보급',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 「論語·孟子보다 채근담이 더 일본인의 일상에 가까이 있다」 — 일본 동양학계의 평',
         {'font_size': 17, 'space_before': 12, 'bold': True, 'color': ACCENT}),
        ('● 일본의 「채근담 효과」 — 중국에서보다 일본에서 더 사랑받은 책',
         {'font_size': 16, 'space_before': 10, 'color': SUB}),
    ])


@S(SEC9)
def ix_korea(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC9, n, t)
    add_title(slide, '한국 — 조선 후기의 애독서',
              '사대부의 처세서에서 현대 자기 계발서의 원형까지')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 조선 후기부터 사대부·선비의 애독서로 정착',
         {'font_size': 18, 'space_before': 4}),
        ('● 한문 원문은 짧고 대구로 되어 있어 사대부 교양에 적합',
         {'font_size': 17, 'space_before': 10}),
        ('● 정조·이덕무·박지원 등 18세기 지식인 — 채근담 구절 인용 다수',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 일제 강점기 — 윤석중·정인보·홍명희 등 문인의 좌우명 자료',
         {'font_size': 17, 'space_before': 10}),
        ('● 1960년대 이후 — 한국 자기계발서·교양서의 가장 자주 인용되는 동양 고전',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 다양한 번역본 — 임어당·조지훈·홍자성 원전·생활 한문 시리즈',
         {'font_size': 16, 'space_before': 12, 'color': SUB}),
    ])


# ============== Ⅹ. 오늘 다시 펼치는 이유 ==============
SEC10 = 'Ⅹ. 오늘 다시 펼치는 이유'

@S(SEC10)
def x_lack(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '오늘의 결핍 지도 — 채근담이 채울 수 있는 자리')
    items = [
        ('과잉의 시대',         '진한 자극·진한 소비·진한 정보 — 「眞味只是淡」이 절실한 시대'),
        ('비교 피로',           '지족(知足)편 「貪得者身富而心貧」 — 비교에서 자기로의 시선 전환'),
        ('가득 채움 강박',       '여지(餘地) 「事事留個有餘」 — 80%만 채우고 20%를 비우는 지혜'),
        ('관계의 마찰',         '처세편 「路徑窄處 留一步」 — 양보가 곧 가장 편안한 길'),
        ('주의력 분산',         '관심(觀心) 「夜深人靜 獨坐觀心」 — 마음을 들여다보는 시간'),
        ('절정에서 내려옴 못함',  '권력편 「進步處便思退步」 — 잘 나갈 때 물러설 줄 아는 안목'),
    ]
    for i, (cat, desc) in enumerate(items):
        y = Inches(2.4 + i * 0.7)
        add_filled_rect(slide, Inches(0.7), y, Inches(3.5), Inches(0.6), ACCENT)
        add_textbox(slide, Inches(0.7), y, Inches(3.5), Inches(0.6),
                    cat, font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.4), y + Inches(0.05), Inches(8.5), Inches(0.55),
                    desc, font_size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)


@S(SEC10)
def x_howread(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '채근담 오늘 읽는 법 — 「한 조 한 조의 명상록」',
              '360개의 진주알을 마음의 결에 따라 한 알씩')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.5), [
        ('● 통독하지 말고 한 조(條)씩 — 마치 명상록(Meditations)을 펼치듯',
         {'font_size': 17, 'space_before': 4}),
        ('● 오늘의 마음 상태에 따라 다가오는 조가 다르다 — 그것이 채근담의 묘미',
         {'font_size': 17, 'space_before': 10, 'bold': True, 'color': ACCENT}),
        ('● 마음에 닿는 한 줄을 적어두고 일주일 — 「實感(실감)」이 와야 진짜 읽은 것',
         {'font_size': 17, 'space_before': 10}),
        ('● 전집은 일상에서 부딪치는 일들 옆에 두고 → 후집은 조용한 시간에 한 줄씩',
         {'font_size': 17, 'space_before': 10, 'color': SUB}),
        ('● 자녀와 함께 — 「오늘 우리 가족엔 어떤 한 줄?」 일주일에 한 조',
         {'font_size': 16, 'space_before': 12}),
        ('● 책 이름 그대로 — 「나물 뿌리를 씹는 마음」으로, 천천히 음미하며',
         {'font_size': 16, 'space_before': 12, 'color': SUB, 'bold': True}),
    ])


@S(SEC10)
def x_seven(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC10, n, t)
    add_title(slide, '채근담이 오늘 우리에게 일러주는 7가지')
    items = [
        '1. 참된 맛은 담백 — 진하게가 아니라 담박하게 살아라',
        '2. 양보가 가장 편안한 길 — 좁은 길에선 한 걸음 비킨다',
        '3. 부귀의 뿌리를 보라 — 도덕에 뿌리박은 부귀만 오래간다',
        '4. 세상은 마음의 거울 — 마음이 밝으면 어두운 방도 푸른 하늘',
        '5. 최고는 반쯤에 있다 — 꽃은 반개, 술은 미취',
        '6. 일이 가면 마음도 따라 빈다 — 바람·기러기처럼 흔적 없이',
        '7. 처음으로 돌아가라 — 모든 수양의 끝은 還初(환초)',
    ]
    for i, txt in enumerate(items):
        y = Inches(2.4 + i * 0.65)
        add_textbox(slide, Inches(0.7), y, Inches(12.0), Inches(0.55),
                    txt, font_size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ============== Ⅺ. 마무리 ==============
SEC11 = 'Ⅺ. 마무리'

@S(SEC11)
def xi_summary(slide, n, t):
    set_white_background(slide)
    add_page_header(slide, SEC11, n, t)
    add_title(slide, '채근담, 한 폭으로')
    add_paragraphs(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.8), [
        ('● 채근담은 명 말 홍자성이 360조로 엮은 처세·수양·은일의 격언집.',
         {'font_size': 18, 'space_before': 6, 'bold': True, 'color': ACCENT}),
        ('● 전집 225조(입세) + 후집 135조(출세) — 두 세계를 한 책에 담음.',
         {'font_size': 18, 'space_before': 8}),
        ('● 유·불·도 삼교가 한 사대부의 경험 속에 녹은 「생활의 융합」.',
         {'font_size': 18, 'space_before': 8}),
        ('● 10대 주제 — 담박·중용·신독·퇴일보·여지·지족·담심·관심·무흔·환초.',
         {'font_size': 18, 'space_before': 8}),
        ('● 일본의 「논어」 · 조선의 애독서 · 동아시아 자기 계발서의 원형.',
         {'font_size': 18, 'space_before': 8}),
        ('● 그리고 무엇보다, 「나물 뿌리를 씹는 자라야 백 가지 일을 이룬다」는 책.',
         {'font_size': 19, 'space_before': 14, 'bold': True, 'color': ACCENT}),
    ])


@S(SEC11)
def xi_final(slide, n, t):
    set_white_background(slide)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), INK)
    add_filled_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), INK)
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.9),
                '寵 辱 不 驚',
                font_size=70, bold=True, color=INK,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.7),
                '閑 看 庭 前 花 開 花 落',
                font_size=44, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name='Batang')
    add_rule(slide, Inches(4.0), Inches(3.9), Inches(5.3))
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.6),
                '영 예 와  치 욕 에  놀 라 지  않 고',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
                '뜰  앞 의  꽃 이  피 고  짐 을  한 가 히  본 다',
                font_size=20, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.5),
                '— 채근담 후집 · 깨달음의 마지막 자리',
                font_size=15, color=SUB, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                '菜  根  譚',
                font_size=22, color=INK, align=PP_ALIGN.CENTER, font_name='Batang')


# ============== 빌드 ==============
total_pages = len(SLIDES)
print(f'Total {total_pages} slides — building ...')
for i, (renderer, section) in enumerate(SLIDES, 1):
    slide = prs.slides.add_slide(blank)
    renderer(slide, i, total_pages)
out_path = r'C:\Users\박호군\ClaudeProjects\Oriental-Classics2\발표자료\채근담_발표자료.pptx'
prs.save(out_path)
print(f'Saved: {out_path}  ({total_pages} slides)')
